"""第1章 机器学习概述 - 扩展讲义PPT生成脚本 (详细版) [简体中文]"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色调色板 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    """以卡片形式绘制表格的辅助函数"""
    total_w = sum(col_widths)
    # 表头
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    # 行
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    """章节分隔幻灯片"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), accent)
    add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), accent)
    add_text(s, Inches(0), Inches(2.0), prs.slide_width, Inches(0.5),
             f"SECTION {section_num}", font_size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, Inches(5.5), Inches(2.7), Inches(2.3), accent)
    add_text(s, Inches(0), Inches(3.0), prs.slide_width, Inches(1.0),
             title, font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.5),
             subtitle, font_size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    return s


# ============================================================
# 幻灯片1：封面
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "机器学习 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "第1章：机器学习概述", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "Introduction to Machine Learning", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.3), Inches(11), Inches(0.8),
         "[ 扩展详细版 ]", font_size=20, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
         "核心关键词：ML定义 · 学习类型 · 偏差-方差权衡 · NFL定理 · 交叉验证 · ML流水线",
         font_size=14, color=DARK_GRAY)
add_shape(s, Inches(0), Inches(6.3), prs.slide_width, Inches(1.2), RGBColor(0x0A, 0x0A, 0x1A))
add_text(s, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
         "2026年第1学期", font_size=14, color=DARK_GRAY)

# ============================================================
# 幻灯片2：目录 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "目录 (Contents) - 1/2")
toc1 = [
    ("01", "什么是机器学习 - 定义、传统编程 vs ML、案例", ACCENT_BLUE),
    ("02", "机器学习的历史 - 里程碑、AI寒冬与复兴", ACCENT_CYAN),
    ("03", "学习类型 - 监督学习、无监督学习、强化学习详解", ACCENT_GREEN),
    ("04", "ML vs DL vs AI - 层次结构、比较、选择指南", ACCENT_PURPLE),
    ("05", "偏差-方差权衡 - 概念、公式推导、各模型特性", ACCENT_ORANGE),
    ("06", "No Free Lunch定理 - 数学定义、假设、启示", ACCENT_RED),
    ("07", "模型评估 - Hold-Out、K-Fold、Stratified、LOO比较", ACCENT_BLUE),
]
for i, (num, title, color) in enumerate(toc1):
    y = Inches(2.0) + Inches(0.7) * i
    add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.5), color, radius=True)
    add_text(s, Inches(1.2), y + Inches(0.05), Inches(0.55), Inches(0.4), num,
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), y + Inches(0.05), Inches(10), Inches(0.4), title,
             font_size=16, color=WHITE)

# ============================================================
# 幻灯片3：目录 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "目录 (Contents) - 2/2")
toc2 = [
    ("08", "ML流水线与数据质量 - 8个步骤、特征工程、数据泄露", ACCENT_CYAN),
    ("09", "5篇核心论文综述 - Domingos、Mitchell、Wolpert、Vapnik、ESL", ACCENT_GREEN),
    ("10", "实验1：偏差-方差权衡实现（代码详解）", ACCENT_BLUE),
    ("11", "实验2：NFL定理演示（代码详解）", ACCENT_ORANGE),
    ("12", "实验3：交叉验证比较（代码详解）", ACCENT_RED),
    ("13", "深入应用案例 - 医疗、金融、制造、NLP", ACCENT_PURPLE),
    ("14", "核心总结、公式整理、10道复习题", ACCENT_CYAN),
]
for i, (num, title, color) in enumerate(toc2):
    y = Inches(2.0) + Inches(0.7) * i
    add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.5), color, radius=True)
    add_text(s, Inches(1.2), y + Inches(0.05), Inches(0.55), Inches(0.4), num,
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), y + Inches(0.05), Inches(10), Inches(0.4), title,
             font_size=16, color=WHITE)

# ============================================================
# 幻灯片4：学习目标
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "学习目标 (Learning Objectives)")
objectives = [
    "能够解释机器学习的定义（Mitchell的T-P-E）及其与传统编程的区别",
    "能够从数据/目标/反馈的角度比较监督学习、无监督学习和强化学习的区别",
    "能够推导偏差-方差分解公式，并解释各项的含义",
    "能够阐述No Free Lunch定理的数学意义和实际启示",
    "能够比较4种交叉验证策略（Hold-Out、K-Fold、Stratified、LOO），并提出各场景的选择标准",
    "能够理解ML流水线8个步骤，并解释数据泄露（Data Leakage）防范原则",
    "能够总结5篇核心论文（Domingos、Mitchell、Wolpert、Vapnik、ESL）的主要贡献",
    "能够理解并运行3个实验代码（偏差-方差、NFL、交叉验证）",
]
for i, obj in enumerate(objectives):
    y = Inches(2.0) + Inches(0.62) * i
    colors = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE,
              ACCENT_RED, ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN]
    add_shape(s, Inches(0.8), y + Inches(0.08), Inches(0.35), Inches(0.35), colors[i], radius=True)
    add_text(s, Inches(0.8), y + Inches(0.08), Inches(0.35), Inches(0.35),
             str(i+1), font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.4), y, Inches(11), Inches(0.55),
             obj, font_size=15, color=LIGHT_GRAY)

# ============================================================
# SECTION 01: 什么是机器学习
# ============================================================
section_divider("什么是机器学习？", "Machine Learning Definition & Fundamentals", "01", ACCENT_BLUE)

# ============================================================
# 幻灯片6：机器学习定义 - Samuel & Mitchell
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "机器学习的定义", "Two Foundational Definitions")

add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Arthur Samuel (1959)", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.2),
         "\"Machine Learning is a field of study that gives\ncomputers the ability to learn without being\nexplicitly programmed.\"",
         font_size=14, color=LIGHT_GRAY)
add_text(s, Inches(0.9), Inches(3.7), Inches(5.2), Inches(0.4),
         "→ 赋予计算机无需显式编程即可学习能力的学科", font_size=12, color=DARK_GRAY)

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "Tom Mitchell (1997) - 形式化定义", font_size=16, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.2),
         "\"A computer program is said to learn from\nexperience E with respect to some class of tasks T\nand performance measure P, if its performance at\ntasks in T, as measured by P, improves with experience E.\"",
         font_size=13, color=LIGHT_GRAY)

# T-P-E 详细
labels = [
    ("T (Task)", "要执行的任务", "邮件垃圾分类", "做什么？", ACCENT_BLUE),
    ("P (Performance)", "性能衡量指标", "正确分类的邮件比例", "做得多好？", ACCENT_GREEN),
    ("E (Experience)", "用于学习的数据", "用户标注的邮件", "用什么学习？", ACCENT_ORANGE),
]
for i, (name, meaning, example, question, color) in enumerate(labels):
    x = Inches(0.6) + Inches(4.1) * i
    add_shape(s, x, Inches(4.6), Inches(3.8), Inches(2.7), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(4.7), Inches(3.4), Inches(0.4),
             name, font_size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, x + Inches(0.3), Inches(5.15), Inches(3.2), color)
    add_text(s, x + Inches(0.2), Inches(5.3), Inches(3.4), Inches(0.3),
             question, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(5.7), Inches(3.4), Inches(0.3),
             meaning, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(6.1), Inches(3.4), Inches(0.8),
             f"例：{example}", font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片7：传统编程 vs 机器学习
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "传统编程 vs 机器学习", "范式转换")

add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.8), Inches(2.3), Inches(5.4), Inches(0.5),
         "传统编程", font_size=22, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(1.0), Inches(3.0), Inches(5.0), Inches(3.8), [
    "输入：数据 + 规则 (Rules)",
    "输出：结果 (Answers)",
    "人工直接编写规则",
    "",
    "示例：基于if-else的垃圾邮件过滤器",
    '  if "中奖" in email: return "垃圾邮件"',
    '  if "免费" in email: return "垃圾邮件"',
    "",
    "局限：难以应对新模式",
    "      规则越复杂，维护越困难",
], font_size=14, color=LIGHT_GRAY)

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.0), Inches(2.3), Inches(5.4), Inches(0.5),
         "机器学习", font_size=22, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(7.2), Inches(3.0), Inches(5.0), Inches(3.8), [
    "输入：数据 + 标签 (Labels)",
    "输出：规则/模型 (Model)",
    "从数据中自动学习规则",
    "",
    "示例：通过数万封邮件学习",
    "  → 自动发现模式",
    "  → 也能检测新型垃圾邮件",
    "",
    "优点：能适应环境变化",
    "      能发现人类无法识别的模式",
], font_size=14, color=LIGHT_GRAY)

add_shape(s, Inches(6.0), Inches(4.0), Inches(1.2), Inches(0.8), ACCENT_BLUE, radius=True)
add_text(s, Inches(6.0), Inches(4.05), Inches(1.2), Inches(0.7),
         "VS", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片8：日常机器学习案例 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "日常生活中的机器学习案例 (1/2)")

# 推荐系统
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "1. 推荐系统 (Recommendation System)", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "Netflix、YouTube、Spotify：分析观看/收听记录 → 个性化推荐",
    "协同过滤 (Collaborative Filtering) + 基于内容的过滤",
    "Netflix：推荐引导约80%的观看时间",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 自动驾驶
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "2. 自动驾驶 (Autonomous Driving)", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "特斯拉、Waymo：摄像头、激光雷达、雷达传感器数据实时处理",
    "基于CNN的目标检测 + 基于强化学习的路径规划",
    "通过数百万公里行驶数据学习，识别行人/车辆/交通灯",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 医疗诊断
add_shape(s, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(5.2), Inches(0.4),
         "3. 医疗诊断 (Medical Diagnosis)", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(5.4), Inches(5.2), Inches(1.5), [
    "从X光、CT、MRI中自动检测肿瘤/肺炎/骨折",
    "AlphaFold：解决蛋白质结构预测问题 → 生物学革新",
    "乳腺癌诊断：部分ML模型准确率高于专科医生",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 核心要点
add_shape(s, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(7.1), Inches(4.9), Inches(5.2), Inches(0.4),
         "核心要点", font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(5.4), Inches(5.2), Inches(1.5), [
    "ML已应用于日常生活的方方面面",
    "多种数据类型：文本、图像、语音、传感器",
    "核心：从数据中自动学习模式进行预测/决策",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# 幻灯片9：日常机器学习案例 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "日常生活中的机器学习案例 (2/2)")

add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "4. 语音识别 (Speech Recognition)", font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "Siri、Google Assistant、Alexa：语音→文本(STT) + 意图理解(NLU)",
    "基于深度学习的RNN、Transformer模型（数万小时语音训练）",
    "OpenAI Whisper：在多语言语音识别中表现优异",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "5. 欺诈检测 (Fraud Detection)", font_size=16, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "信用卡盗刷、保险欺诈、洗钱实时检测",
    "学习正常模式 → 检测异常交易(anomaly)",
    "无监督学习(异常值检测) + 监督学习(分类)并行，每秒分析数千笔",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 为什么需要ML
add_shape(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.4),
         "为什么需要机器学习 (Mitchell, 1997)", font_size=18, color=ACCENT_ORANGE, bold=True)
reasons = [
    ("自动发现复杂规则", "人工难以逐一定义规则的问题（图像识别、语音识别）", ACCENT_BLUE),
    ("适应环境变化", "新数据到来时重新训练模型以应对变化", ACCENT_GREEN),
    ("从海量数据中发现模式", "大数据时代人工分析已不可能 → 自动化分析", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(reasons):
    x = Inches(0.8) + Inches(4.0) * i
    add_shape(s, x, Inches(5.5), Inches(3.7), Inches(1.5), RGBColor(0x2D, 0x2D, 0x45), color, radius=True)
    add_text(s, x + Inches(0.15), Inches(5.55), Inches(3.4), Inches(0.4),
             f"{i+1}. {title}", font_size=13, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(5.95), Inches(3.4), Inches(0.8),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# SECTION 02: 机器学习的历史
# ============================================================
section_divider("机器学习的历史", "History of Machine Learning", "02", ACCENT_CYAN)

# ============================================================
# 幻灯片11：历史时间线 (1/2) - 1943~1997
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "机器学习主要里程碑 (1943~1997)")

milestones1 = [
    ("1943", "McCulloch & Pitts", "人工神经元模型", "首个数学神经元模型\n神经网络的理论基础", ACCENT_BLUE),
    ("1950", "Alan Turing", "Computing Machinery\nand Intelligence", "提出图灵测试\n\"机器能思考吗？\"", ACCENT_CYAN),
    ("1957", "Rosenblatt", "感知机", "首个可学习的\n神经网络模型", ACCENT_GREEN),
    ("1959", "Arthur Samuel", "跳棋程序", "\"Machine Learning\"\n术语首次使用", ACCENT_GREEN),
    ("1969", "Minsky & Papert", "Perceptrons", "证明感知机局限(XOR)\n→ 引发第一次AI寒冬", ACCENT_RED),
    ("1986", "Rumelhart et al.", "反向传播算法", "多层神经网络可训练\n→ 神经网络复兴", ACCENT_ORANGE),
    ("1995", "Vapnik", "SVM", "基于统计学习理论的\n强大分类器", ACCENT_PURPLE),
    ("1997", "Mitchell / Wolpert", "ML教材 / NFL定理", "确立T-P-E定义\n否定通用最优算法", ACCENT_BLUE),
]
for i, (year, author, event, desc, color) in enumerate(milestones1):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(3.15) * col
    y = Inches(2.2) + Inches(2.7) * row
    add_shape(s, x, y, Inches(2.9), Inches(2.4), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.7), Inches(0.35), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.7), Inches(0.35),
             year, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.9), y + Inches(0.1), Inches(1.9), Inches(0.35),
             author, font_size=11, color=color, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(0.35),
             event, font_size=13, color=WHITE, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(1.2),
             desc, font_size=11, color=LIGHT_GRAY)

# ============================================================
# 幻灯片12：历史时间线 (2/2) - 2006~至今
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "机器学习主要里程碑 (2006~至今)")

milestones2 = [
    ("2006", "Hinton", "Deep Belief Networks", "深度学习复兴信号弹\n\"深度学习\"术语普及", ACCENT_BLUE),
    ("2012", "Krizhevsky et al.", "AlexNet (ImageNet)", "CNN性能爆发式提升\n深度学习热潮开始", ACCENT_CYAN),
    ("2014", "Goodfellow", "GAN", "生成对抗网络\n生成模型的革新", ACCENT_GREEN),
    ("2016", "DeepMind", "AlphaGo vs 李世石", "向全世界证明了\n强化学习的潜力", ACCENT_ORANGE),
    ("2017", "Vaswani et al.", "Transformer", "Attention Is All You Need\nNLP革命的开端", ACCENT_RED),
    ("2018~", "Google, OpenAI", "BERT、GPT系列", "大规模语言模型(LLM)\n时代开启", ACCENT_PURPLE),
    ("2022~", "OpenAI", "ChatGPT、GPT-4", "生成式AI的普及\nAGI讨论活跃", ACCENT_BLUE),
]
for i, (year, author, event, desc, color) in enumerate(milestones2):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(3.15) * col
    y = Inches(2.2) + Inches(2.7) * row
    add_shape(s, x, y, Inches(2.9), Inches(2.4), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.7), Inches(0.35), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.7), Inches(0.35),
             year, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.9), y + Inches(0.1), Inches(1.9), Inches(0.35),
             author, font_size=11, color=color, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(0.35),
             event, font_size=13, color=WHITE, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(1.2),
             desc, font_size=11, color=LIGHT_GRAY)

# ============================================================
# 幻灯片13：AI寒冬与复兴
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "AI寒冬与复兴", "AI Winters and Revivals")

phases = [
    ("第一次AI寒冬\n(1970年代)", "证明感知机的局限\n(XOR问题)\n→ 神经网络研究停滞\nMinsky & Papert(1969)",
     ACCENT_RED, Inches(0.5)),
    ("复兴\n(1980年代)", "发现反向传播算法\n→ 多层神经网络可训练\nRumelhart, Hinton,\nWilliams (1986)",
     ACCENT_GREEN, Inches(3.3)),
    ("第二次AI寒冬\n(1990年代初)", "神经网络训练困难\n(梯度消失)\n→ SVM、集成方法成为主流\nVapnik的SVM(1995)",
     ACCENT_ORANGE, Inches(6.1)),
    ("深度学习复兴\n(2006~2012)", "GPU计算 + 大数据\n+ ReLU激活函数\n→ 深度学习爆发式增长\nAlexNet(2012)",
     ACCENT_CYAN, Inches(8.9)),
]
# 时间线连接线
add_shape(s, Inches(0.8), Inches(3.7), Inches(11.7), Pt(4), ACCENT_BLUE)

for title, desc, color, x in phases:
    add_shape(s, x, Inches(2.2), Inches(2.5), Inches(1.3), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(2.3), Inches(2.2), Inches(1.1),
             title, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    # 圆点
    add_shape(s, x + Inches(1.1), Inches(3.55), Inches(0.3), Inches(0.3), color, radius=True)
    # 说明
    add_shape(s, x, Inches(4.1), Inches(2.5), Inches(2.5), CARD_BG, radius=True)
    add_text(s, x + Inches(0.15), Inches(4.2), Inches(2.2), Inches(2.2),
             desc, font_size=12, color=LIGHT_GRAY)

add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
         "启示：ML的发展并非直线前进。发现理论局限 → 停滞 → 新突破 → 复兴的循环",
         font_size=14, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 03: 学习类型
# ============================================================
section_divider("学习类型", "Supervised / Unsupervised / Reinforcement Learning", "03", ACCENT_GREEN)

# ============================================================
# 幻灯片15：监督学习详解
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "监督学习 (Supervised Learning)", "输入X与标签y同时给定的学习")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.35),
         "定义：在输入数据X与标签y同时给定的条件下，学习映射函数 f: X → y 的方法",
         font_size=16, color=WHITE, bold=True)
add_text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.35),
         "核心：朝着减小预测值与实际值差异（误差）的方向反复学习 → 回归（连续值）+ 分类（离散值）",
         font_size=14, color=LIGHT_GRAY)

# 回归
add_card(s, Inches(0.6), Inches(3.5), Inches(5.8), Inches(1.8),
         "回归 (Regression) - 连续值预测", [
             "输出值y为连续数值：房价、气温、销售额预测",
             "代表算法：线性回归、Ridge、Lasso、RF、GBM",
             "评估指标：MSE、RMSE、MAE、R²",
         ], ACCENT_GREEN, ACCENT_GREEN)
# 分类
add_card(s, Inches(6.8), Inches(3.5), Inches(5.8), Inches(1.8),
         "分类 (Classification) - 类别预测", [
             "输出值y为离散类别：垃圾邮件/正常、恶性/良性",
             "二分类 / 多分类 / 多标签分类",
             "评估指标：Accuracy、Precision、Recall、F1、AUC",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 行业案例表格
add_text(s, Inches(0.6), Inches(5.5), Inches(5), Inches(0.4),
         "监督学习行业应用案例", font_size=15, color=WHITE, bold=True)
add_table_slide(s,
    ["行业", "问题", "输入(X)", "输出(y)", "算法"],
    [
        ["金融", "信用评分预测", "收入、债务、交易记录", "信用等级(A~F)", "XGBoost"],
        ["制造", "产品缺陷预测", "传感器、温度、压力", "良品/不良品", "Random Forest"],
        ["营销", "客户流失预测", "使用模式、支付", "流失/留存", "Gradient Boosting"],
    ],
    Inches(0.6), Inches(5.9), [2.0, 1.8, 2.5, 2.5, 2.4], font_size=12, header_font_size=13, row_height=0.45)


# ============================================================
# PLACEHOLDER for remaining slides - will be appended
# ============================================================

# ============================================================
# 幻灯片16：无监督学习详解
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "无监督学习 (Unsupervised Learning)", "无标签情况下发现数据结构")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.7), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5),
         "定义：仅给定输入数据X，在没有标签y的情况下发现数据本身结构和模式的学习方法",
         font_size=16, color=WHITE, bold=True)

tasks = [
    ("聚类 (Clustering)", "将相似数据分组", "K-Means、DBSCAN\n层次聚类",
     "客户细分\n基因聚类", ACCENT_BLUE),
    ("降维\n(Dimensionality Reduction)", "高维 → 低维转换", "PCA、t-SNE\nAutoencoder",
     "可视化\n去噪", ACCENT_CYAN),
    ("异常检测\n(Anomaly Detection)", "检测偏离正常模式的\n数据", "Isolation Forest\nOne-Class SVM",
     "欺诈检测\n网络异常", ACCENT_ORANGE),
]
for i, (title, desc, algos, cases, color) in enumerate(tasks):
    x = Inches(0.5) + Inches(4.15) * i
    add_shape(s, x, Inches(3.2), Inches(3.9), Inches(4.0), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(3.3), Inches(3.5), Inches(0.6),
             title, font_size=15, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.9), Inches(3.5), Inches(0.5),
             desc, font_size=13, color=WHITE)
    add_text(s, x + Inches(0.2), Inches(4.5), Inches(3.5), Inches(0.3),
             "算法：", font_size=12, color=ACCENT_CYAN, bold=True)
    add_text(s, x + Inches(0.2), Inches(4.8), Inches(3.5), Inches(0.8),
             algos, font_size=12, color=LIGHT_GRAY)
    add_text(s, x + Inches(0.2), Inches(5.6), Inches(3.5), Inches(0.3),
             "应用案例：", font_size=12, color=ACCENT_GREEN, bold=True)
    add_text(s, x + Inches(0.2), Inches(5.9), Inches(3.5), Inches(0.8),
             cases, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 幻灯片17：强化学习详解
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "强化学习 (Reinforcement Learning)", "智能体与环境交互以最大化奖励")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.7), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5),
         "定义：智能体(Agent)与环境(Environment)交互，学习使奖励(Reward)最大化的行为策略(Policy)",
         font_size=16, color=WHITE, bold=True)

# 构成要素
components = [
    ("智能体 (Agent)", "学习并采取行动的主体", ACCENT_BLUE),
    ("环境 (Environment)", "智能体交互的世界", ACCENT_CYAN),
    ("状态 (State)", "当前环境的情况", ACCENT_GREEN),
    ("动作 (Action)", "智能体可采取的操作", ACCENT_ORANGE),
    ("奖励 (Reward)", "对行为的反馈信号", ACCENT_RED),
]
for i, (name, desc, color) in enumerate(components):
    x = Inches(0.4) + Inches(2.55) * i
    add_shape(s, x, Inches(3.2), Inches(2.35), Inches(1.2), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.1), Inches(3.3), Inches(2.15), Inches(0.4),
             name, font_size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(3.7), Inches(2.15), Inches(0.5),
             desc, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 交互图
add_shape(s, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(4.8), Inches(5.2), Inches(0.4),
         "强化学习交互循环", font_size=15, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(0.9), Inches(5.3), Inches(5.2), Inches(1.5),
         "智能体 --[动作(Action)]--> 环境\n环境 --[状态(State), 奖励(Reward)]--> 智能体\n\n反复进行此过程学习最优策略(Policy)\n目标：最大化累积奖励的期望值",
         font_size=14, color=LIGHT_GRAY)

# 行业案例
add_shape(s, Inches(6.8), Inches(4.7), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(4.8), Inches(5.2), Inches(0.4),
         "强化学习行业案例", font_size=15, color=ACCENT_CYAN, bold=True)
add_table_slide(s,
    ["领域", "奖励设计", "算法"],
    [
        ["游戏AI (AlphaGo)", "胜利+1, 失败-1", "MCTS+策略神经网络"],
        ["机械臂控制", "抓取目标成功/失败", "PPO, SAC"],
        ["数据中心冷却", "节能量", "DQN (DeepMind)"],
    ],
    Inches(6.8), Inches(5.3), [2.1, 1.9, 1.8], font_size=11, header_font_size=12, row_height=0.4)

# ============================================================
# 幻灯片18：三种学习类型比较
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "三种学习类型比较总结")

add_table_slide(s,
    ["类别", "监督学习", "无监督学习", "强化学习"],
    [
        ["数据", "输入 + 标签", "仅输入", "状态 + 奖励"],
        ["目标", "预测 / 分类", "发现结构", "奖励最大化"],
        ["反馈", "标签", "无", "奖励信号"],
        ["代表问题", "分类、回归", "聚类、降维", "游戏、机器人控制"],
        ["代表算法", "RF, XGB, SVM, LR", "K-Means, PCA", "Q-Learning, PPO"],
        ["难度", "相对较易", "中等", "相对较难"],
    ],
    Inches(0.8), Inches(2.2), [2.5, 2.8, 2.8, 2.8], row_height=0.55)

add_shape(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.3), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(1.1), Inches(5.9), Inches(11.1), Inches(0.4),
         "实务选择指南", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(1.1), Inches(6.3), Inches(11.1), Inches(0.7), [
    "有标签 → 监督学习（分类/回归）  |  无标签发现模式 → 无监督学习  |  序贯决策 → 强化学习",
    "实务中最常使用：监督学习（数据充足则性能高）。无监督学习常用于探索/预处理阶段",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# SECTION 04: ML vs DL vs AI
# ============================================================
section_divider("ML vs DL vs AI", "Hierarchy and Comparison", "04", ACCENT_PURPLE)

# ============================================================
# 幻灯片20：ML vs DL vs AI层次结构
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "AI ⊃ ML ⊃ DL 层次结构")

# 同心圆层次
circles = [
    (Inches(1.0), Inches(2.0), Inches(6.0), Inches(5.2), RGBColor(0x1E, 0x3A, 0x5F), ACCENT_BLUE),
    (Inches(1.8), Inches(2.6), Inches(4.4), Inches(4.0), RGBColor(0x1E, 0x4D, 0x3A), ACCENT_GREEN),
    (Inches(2.6), Inches(3.2), Inches(2.8), Inches(2.6), RGBColor(0x3A, 0x1E, 0x5F), ACCENT_PURPLE),
]
for x, y, w, h, fill, border in circles:
    add_shape(s, x, y, w, h, fill, border, radius=True)

add_text(s, Inches(1.3), Inches(2.2), Inches(3), Inches(0.4),
         "人工智能 (AI)", font_size=18, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(2.1), Inches(2.8), Inches(3), Inches(0.4),
         "机器学习 (ML)", font_size=16, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(2.9), Inches(4.2), Inches(2.2), Inches(0.4),
         "深度学习 (DL)", font_size=15, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)

# 说明卡片
descs = [
    ("AI (Artificial Intelligence)", "模仿人类智能的所有技术的总称\n包括基于规则、专家系统、搜索等\n最广泛的概念", ACCENT_BLUE),
    ("ML (Machine Learning)", "AI的子领域，专注于数据驱动学习\n不使用显式规则，从数据中自动学习模式\nSVM、RF、XGBoost等", ACCENT_GREEN),
    ("DL (Deep Learning)", "ML的子领域，利用深层神经网络\n通过多层非线性变换学习复杂表示\nCNN、RNN、Transformer", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(descs):
    y = Inches(2.0) + Inches(1.8) * i
    add_shape(s, Inches(7.5), y, Inches(5.2), Inches(1.5), CARD_BG, color, radius=True)
    add_text(s, Inches(7.7), y + Inches(0.1), Inches(4.8), Inches(0.35),
             title, font_size=14, color=color, bold=True)
    add_text(s, Inches(7.7), y + Inches(0.45), Inches(4.8), Inches(0.9),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 幻灯片21：ML vs DL详细比较
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "传统ML vs 深度学习详细比较")

add_table_slide(s,
    ["特性", "传统机器学习", "深度学习"],
    [
        ["特征工程", "人工手动设计", "模型自动学习"],
        ["数据需求量", "少量数据也可以", "需要大量数据"],
        ["计算资源", "CPU通常足够", "GPU/TPU必需"],
        ["可解释性", "相对较高", "黑箱倾向"],
        ["性能上限", "数据增加时性能饱和", "随数据增加持续提升"],
        ["适合的数据", "结构化数据（表格）", "非结构化（图像、文本、语音）"],
        ["代表算法", "SVM, RF, XGBoost", "CNN, RNN, Transformer"],
    ],
    Inches(0.8), Inches(2.2), [3.0, 4.3, 4.3], row_height=0.5, font_size=14)

# 选择指南
add_shape(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.35),
         "何时使用什么？(Domingos, 2012)", font_size=14, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7), [
    "结构化数据 + 中小规模 → XGBoost、LightGBM、RF  |  非结构化（图像、文本）→ CNN、Transformer",
    "可解释性重要（医疗、金融）→ 传统ML或XAI  |  数据很少时 → 传统ML + 领域知识特征",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# SECTION 05: 偏差-方差权衡
# ============================================================
section_divider("偏差-方差权衡", "Bias-Variance Tradeoff", "05", ACCENT_ORANGE)

# ============================================================
# 幻灯片23：偏差与方差概念
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "偏差(Bias)与方差(Variance)概念")

# 偏差
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.5),
         "偏差 (Bias)", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.9), Inches(5.2), Inches(2.0), [
    "模型预测值与真实值之间的系统性误差",
    "衡量模型捕捉训练数据模式的能力",
    "",
    "高偏差 → 模型过于简单 → 欠拟合 (Underfitting)",
    "低偏差 → 能很好地捕捉数据模式",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 方差
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.8), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.5),
         "方差 (Variance)", font_size=22, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.9), Inches(5.2), Inches(2.0), [
    "在不同训练数据集上预测值变化多少",
    "衡量模型预测的稳定性",
    "",
    "高方差 → 对训练数据敏感 → 过拟合 (Overfitting)",
    "低方差 → 即使数据改变也能稳定预测",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 靶心比喻
add_text(s, Inches(0.6), Inches(5.3), Inches(5), Inches(0.4),
         "直观类比：打靶（靶心 = 真实值）", font_size=15, color=WHITE, bold=True)
targets = [
    ("低偏差\n低方差", "理想 (Best)", ACCENT_GREEN),
    ("高偏差\n低方差", "欠拟合", ACCENT_ORANGE),
    ("低偏差\n高方差", "过拟合", ACCENT_RED),
    ("高偏差\n高方差", "最差 (Worst)", DARK_GRAY),
]
for i, (label, status, color) in enumerate(targets):
    x = Inches(0.6) + Inches(3.1) * i
    add_shape(s, x, Inches(5.8), Inches(2.8), Inches(1.4), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(5.9), Inches(1.3), Inches(1.2),
             label, font_size=12, color=LIGHT_GRAY)
    add_text(s, x + Inches(1.4), Inches(6.1), Inches(1.3), Inches(0.8),
             status, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片24：公式推导 (1/2) - 问题设定
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "偏差-方差分解：公式推导 (1/2)", "Problem Setup & Derivation")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "第1步：数据生成模型假设", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.5),
         "y = f(x) + e,    e ~ N(0, s²)\n其中f(x)为真实函数(ground truth)，e为不可约噪声(irreducible error)",
         font_size=15, color=LIGHT_GRAY)

add_shape(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(0.4),
         "第2步：期望预测误差 (Expected Prediction Error) 定义", font_size=16, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.5),
         "用训练数据集D训练的模型预测：f_hat_D(x)\nEPE(x) = E_D[(y - f_hat_D(x))²]  =  E_D[(f(x) + e - f_hat_D(x))²]",
         font_size=15, color=LIGHT_GRAY)

add_shape(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4),
         "第3步：引入平均预测", font_size=16, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.3),
         "f_bar(x) = E_D[f_hat_D(x)]  ← 多个数据集的模型预测期望值（平均预测）\n\n"
         "将f_hat_D(x)分解如下：\n"
         "f_hat_D(x) = f_bar(x) + (f_hat_D(x) - f_bar(x))\n"
         "              ↑ 平均预测    ↑ 偏离平均的偏差",
         font_size=14, color=LIGHT_GRAY)

# ============================================================
# 幻灯片25：公式推导 (2/2) - 最终分解
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "偏差-方差分解：公式推导 (2/2)", "Final Decomposition")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "第4步：展开与整理", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.0),
         "EPE(x) = E[(f(x) + e - f_hat(x))²]\n"
         "       = E[(f(x) - f_bar(x))² + (f_hat(x) - f_bar(x))² + e² + 交叉项]\n"
         "交叉项因e与f_hat独立且E[e]=0，全部为0",
         font_size=14, color=LIGHT_GRAY)

# 最终公式 - 大框
add_shape(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(1.5), RGBColor(0x15, 0x15, 0x30), ACCENT_CYAN, radius=True)
add_text(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.6),
         "EPE = Bias² + Variance + s²  (Irreducible Error)", font_size=28, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(0.5),
         "总误差 = 模型假设造成的系统性误差 + 数据变动造成的不稳定性 + 数据自身噪声",
         font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 各项说明表格
add_table_slide(s,
    ["项", "公式", "含义", "减小方法"],
    [
        ["Bias²", "(f(x) - f_bar(x))²", "模型假设造成的系统性误差", "增加模型复杂度"],
        ["Variance", "E[(f_hat - f_bar)²]", "数据变动造成的预测不稳定", "降低模型复杂度、集成"],
        ["s²（不可约）", "s²", "数据自身噪声", "无法减小"],
    ],
    Inches(0.6), Inches(5.8), [1.8, 2.8, 3.5, 3.0], row_height=0.5, font_size=13)

# ============================================================
# 幻灯片26：图表解读 + 各模型特性
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "偏差-方差权衡图 & 各模型特性")

# 图表说明（左侧）
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "图表解读", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(4.0), [
    "X轴：模型复杂度（多项式次数、树深度、层数）",
    "Y轴：误差 (Error)",
    "",
    "当模型复杂度增加时：",
    "  → 偏差(Bias)减小（可捕捉复杂模式）",
    "  → 方差(Variance)增大（对数据更敏感）",
    "",
    "最优模型复杂度：",
    "  → Bias² + Variance之和最小的点",
    "  → 在此点泛化性能最大",
    "",
    "左侧（低复杂度）= 欠拟合区域",
    "右侧（高复杂度）= 过拟合区域",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 各模型特性（右侧）
add_text(s, Inches(6.8), Inches(2.2), Inches(5.5), Inches(0.4),
         "各模型偏差-方差特性", font_size=16, color=ACCENT_CYAN, bold=True)
add_table_slide(s,
    ["模型", "偏差", "方差", "特征"],
    [
        ["线性回归", "高", "低", "简单、稳定"],
        ["KNN (K大)", "高", "低", "过度平滑"],
        ["KNN (K小)", "低", "高", "对噪声敏感"],
        ["决策树（深）", "低", "高", "过拟合倾向"],
        ["随机森林", "低", "适当", "集成→方差↓"],
        ["Boosting", "低", "适当~高", "减小偏差的优势"],
    ],
    Inches(6.8), Inches(2.7), [1.7, 0.8, 1.0, 1.6], row_height=0.48, font_size=12, header_font_size=13)

add_shape(s, Inches(6.8), Inches(6.1), Inches(5.5), Inches(0.8), CARD_BG, radius=True)
add_text(s, Inches(7.0), Inches(6.15), Inches(5.1), Inches(0.65),
         "ESL (Hastie et al., 2009) Ch2:\n最小二乘法（高偏差、低方差）vs k-NN（低偏差、高方差）比较",
         font_size=12, color=DARK_GRAY)

# ============================================================
# SECTION 06: NFL定理
# ============================================================
section_divider("No Free Lunch定理", "Wolpert & Macready (1997)", "06", ACCENT_RED)

# ============================================================
# 幻灯片28：NFL定理 - 内容与数学定义
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "No Free Lunch定理", "核心内容与数学定义")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "核心主张 (Wolpert & Macready, 1997)", font_size=16, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.6),
         "对所有可能的问题（目标函数）取平均时，任何两个优化（或学习）算法都表现出相同的性能。\n"
         "→ 不存在\"最佳算法\"。必须始终选择适合问题的算法。",
         font_size=15, color=LIGHT_GRAY)

# 数学定义
add_shape(s, Inches(0.6), Inches(3.8), Inches(12.1), Inches(2.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.4),
         "数学定义", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(4.3), Inches(5.5), Inches(1.8), [
    "X：候选解(candidate solution)的集合",
    "Y：代价值(cost value)的集合",
    "f: X → Y：目标函数",
    "a：优化算法",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(6.5), Inches(4.3), Inches(5.8), Inches(0.8), RGBColor(0x15, 0x15, 0x30))
add_text(s, Inches(6.5), Inches(4.35), Inches(5.8), Inches(0.7),
         "Sf P(d_m^y | f, m, a1) = Sf P(d_m^y | f, m, a2)",
         font_size=16, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.5), Inches(5.15), Inches(5.8), Inches(0.6),
         "对所有函数f求和时，对所有算法对\n(a1, a2)上述等式成立",
         font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 直观含义
add_shape(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7), CARD_BG, ACCENT_ORANGE, radius=True)
add_bullet_list(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6), [
    "若a1在某问题上优于a2，则必定存在其他问题中a2优于a1  |  随机搜索平均而言也与精密算法相同",
], font_size=13, color=LIGHT_GRAY)

# ============================================================
# 幻灯片29：NFL假设与启示
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "NFL定理：核心假设与实务启示")

# 假设
add_text(s, Inches(0.6), Inches(2.1), Inches(5), Inches(0.4),
         "核心假设与局限", font_size=16, color=WHITE, bold=True)
add_table_slide(s,
    ["假设", "含义", "现实解读"],
    [
        ["对所有函数均匀分布", "同等考虑所有可能的f", "现实问题具有特定结构"],
        ["不重复访问已评估点", "不对同一点评估两次", "大多数算法满足"],
        ["代价函数确定性", "f为deterministic", "概率问题需另行分析"],
    ],
    Inches(0.6), Inches(2.5), [3.0, 3.5, 4.5], row_height=0.5, font_size=13)

add_shape(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.6), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(4.25), Inches(11.5), Inches(0.5),
         "注意：NFL常被误解为\"所有算法都相同\"，但现实问题具有结构(structure)，特定算法更有优势！",
         font_size=14, color=ACCENT_ORANGE, bold=True)

# 启示
add_text(s, Inches(0.6), Inches(5.0), Inches(5), Inches(0.4),
         "5条实务启示", font_size=16, color=WHITE, bold=True)
implications = [
    ("1", "不存在\"最佳算法\"", "始终选择适合问题的算法", ACCENT_BLUE),
    ("2", "利用先验知识", "问题结构（归纳偏置）是关键", ACCENT_GREEN),
    ("3", "尝试多种算法", "不要依赖单一算法，要比较", ACCENT_ORANGE),
    ("4", "基准测试的局限", "特定性能 ≠ 所有问题优秀", ACCENT_RED),
    ("5", "AutoML的依据", "自动算法选择/调参的合理性", ACCENT_PURPLE),
]
for i, (num, title, desc, color) in enumerate(implications):
    x = Inches(0.4) + Inches(2.55) * i
    add_shape(s, x, Inches(5.5), Inches(2.35), Inches(1.8), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), Inches(5.6), Inches(0.35), Inches(0.35), color, radius=True)
    add_text(s, x + Inches(0.1), Inches(5.6), Inches(0.35), Inches(0.35),
             num, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.55), Inches(5.6), Inches(1.7), Inches(0.4),
             title, font_size=12, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(6.1), Inches(2.1), Inches(0.9),
             desc, font_size=11, color=LIGHT_GRAY)

# ============================================================
# SECTION 07: 模型评估 - 交叉验证
# ============================================================
section_divider("模型评估：交叉验证", "Cross-Validation Strategies", "07", ACCENT_BLUE)

# ============================================================
# 幻灯片31：为什么要划分数据 + Hold-Out
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "为什么要划分数据？& Hold-Out")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "泛化性能 (Generalization Performance) 评估的核心原则", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.5),
         "训练数据与评估数据必须分开。\n用训练数据评估无法判断是否过拟合 → 产生乐观偏差 (Optimistic Bias)",
         font_size=14, color=LIGHT_GRAY)

# Hold-Out
add_shape(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(3.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(0.4),
         "Hold-Out（简单划分）", font_size=18, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(0.9), Inches(4.3), Inches(5.5), Inches(0.4),
         "最简单的方法：将数据仅一次划分为训练/测试集（例：80:20）", font_size=14, color=WHITE)

add_bullet_list(s, Inches(0.9), Inches(4.8), Inches(5.5), Inches(2.0), [
    "优点：",
    "  • 计算成本最低",
    "  • 实现简单",
    "缺点：",
    "  • 划分不同结果差异大（高方差）",
    "  • 数据使用不高效",
    "使用时机：数据量很大时、快速原型开发",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# Hold-Out可视化
add_text(s, Inches(7.0), Inches(4.5), Inches(5.5), Inches(0.3),
         "Hold-Out划分示例：", font_size=13, color=WHITE, bold=True)
add_shape(s, Inches(7.0), Inches(4.9), Inches(4.0), Inches(0.4), ACCENT_BLUE)
add_text(s, Inches(7.0), Inches(4.9), Inches(4.0), Inches(0.4),
         "训练数据 (80%)", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
add_shape(s, Inches(11.0), Inches(4.9), Inches(1.0), Inches(0.4), ACCENT_RED)
add_text(s, Inches(11.0), Inches(4.9), Inches(1.0), Inches(0.4),
         "测试", font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(7.0), Inches(5.6), Inches(5.5), Inches(0.3),
         "不同划分 → 不同结果：", font_size=13, color=ACCENT_RED, bold=True)
add_shape(s, Inches(7.0), Inches(6.0), Inches(3.0), Inches(0.3), ACCENT_BLUE)
add_shape(s, Inches(10.0), Inches(6.0), Inches(2.0), Inches(0.3), ACCENT_RED)
add_shape(s, Inches(7.0), Inches(6.4), Inches(4.5), Inches(0.3), ACCENT_BLUE)
add_shape(s, Inches(11.5), Inches(6.4), Inches(0.5), Inches(0.3), ACCENT_RED)

# ============================================================
# 幻灯片32：K-Fold & Stratified K-Fold
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "K-Fold & Stratified K-Fold 交叉验证")

# K-Fold
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "K-Fold交叉验证", font_size=18, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(0.5),
         "将数据分为K个折叠，评估K次\n最终性能 = (性能1+性能2+...+性能K) / K", font_size=13, color=WHITE)

# K-Fold可视化
for fold_i in range(5):
    y = Inches(3.6) + Inches(0.28) * fold_i
    add_text(s, Inches(0.9), y, Inches(1.0), Inches(0.25),
             f"Fold {fold_i+1}:", font_size=10, color=DARK_GRAY)
    for block_j in range(5):
        x = Inches(2.0) + Inches(0.8) * block_j
        is_val = (block_j == fold_i)
        c = ACCENT_RED if is_val else RGBColor(0x2A, 0x4A, 0x6A)
        add_shape(s, x, y, Inches(0.7), Inches(0.22), c)

add_bullet_list(s, Inches(0.9), Inches(5.2), Inches(5.2), Inches(1.5), [
    "优点：所有数据都用于训练/评估",
    "缺点：类别不平衡时各折叠比例可能不同",
    "K=5或K=10为标准",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# Stratified K-Fold
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "Stratified K-Fold交叉验证", font_size=18, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(0.5),
         "K-Fold + 各折叠中保持与原始数据相同的类别比例", font_size=13, color=WHITE)

add_bullet_list(s, Inches(7.1), Inches(3.5), Inches(5.2), Inches(3.0), [
    "优点：",
    "  • 在类别不平衡数据中稳定",
    "  • 最广泛推荐的CV策略",
    "  • sklearn的默认CV策略",
    "",
    "缺点：",
    "  • 不能直接用于回归问题（连续值）",
    "",
    "使用时机：",
    "  • 分类问题，特别是类别不平衡",
    "  • 一般最推荐使用",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 幻灯片33：LOO + 策略比较
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "Leave-One-Out & 策略比较表")

# LOO
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.2), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Leave-One-Out (LOO)", font_size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "每次只测试1个样本（K=N，极端K-Fold）",
    "优点：偏差最低（几乎用全部数据训练）",
    "缺点：计算成本非常高（训练N次），方差可能较高",
    "使用：数据很少时（N < 50）",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# K值权衡
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.2), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "K值对偏差-方差权衡的影响", font_size=15, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "K小（K=2）：训练数据50% → 高偏差、低方差",
    "K大（K=N, LOO）：几乎全部训练 → 低偏差、高方差",
    "实务最优点：K = 5 ~ 10 最普遍推荐",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 比较表格
add_text(s, Inches(0.6), Inches(4.6), Inches(5), Inches(0.4),
         "四种策略比较", font_size=16, color=WHITE, bold=True)
add_table_slide(s,
    ["策略", "训练数据", "评估次数", "偏差", "方差", "计算成本"],
    [
        ["Hold-Out", "80%", "1次", "高", "高", "非常低"],
        ["5-Fold", "80%（反复）", "5次", "中等", "中等", "中等"],
        ["10-Fold", "90%（反复）", "10次", "低", "中~高", "中~高"],
        ["LOO", "(N-1)/N", "N次", "最低", "高", "非常高"],
    ],
    Inches(0.6), Inches(5.0), [1.8, 1.5, 1.3, 1.3, 1.3, 1.5], row_height=0.45, font_size=12, header_font_size=13)

# 按数据规模推荐
add_shape(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3), CARD_BG)
# 此位置不可见，省略

# ============================================================
# 幻灯片34：按数据规模推荐 + 评估指标
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "按数据规模推荐策略 & 评估指标概述")

# 推荐策略
recs = [
    ("大规模 (>10,000)", "Hold-Out或3-Fold", "数据充足，计算效率优先", ACCENT_BLUE),
    ("中等规模 (100~10,000)", "Stratified 5或10-Fold", "稳定估计 + 适当成本", ACCENT_GREEN),
    ("小规模 (<100)", "LOO或Repeated K-Fold", "最大限度利用所有数据", ACCENT_ORANGE),
]
for i, (size, strat, reason, color) in enumerate(recs):
    x = Inches(0.5) + Inches(4.15) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(1.6), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.35),
             size, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.7), Inches(3.5), Inches(0.35),
             strat, font_size=14, color=WHITE, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.1), Inches(3.5), Inches(0.5),
             reason, font_size=12, color=DARK_GRAY)

# 评估指标
add_text(s, Inches(0.6), Inches(4.1), Inches(5), Inches(0.4),
         "主要评估指标", font_size=16, color=WHITE, bold=True)

# 回归
add_card(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.7),
         "回归评估指标", [
             "RMSE = sqrt(mean((y_pred - y_true)²)) → 对大误差敏感",
             "MAE = mean(|y_pred - y_true|) → 对异常值稳健",
             "R² = 1 - SS_res/SS_tot → 解释力（越接近1越好）",
             "",
             "RMSE vs MAE：异常值多时优选MAE",
             "R²可能 < 0（比平均值还差的模型）",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 分类
add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.7),
         "分类评估指标", [
             "Accuracy = (TP+TN)/(TP+TN+FP+FN) → 注意不平衡！",
             "Precision = TP/(TP+FP) → 阳性预测的准确率",
             "Recall = TP/(TP+FN) → 实际阳性检出率（医疗核心）",
             "F1 = 2*(Prec*Rec)/(Prec+Rec) → 调和平均",
             "AUC-ROC = ROC曲线下面积 (0~1)",
             "  1.0=完美, 0.5=随机 → 阈值无关的综合指标",
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# SECTION 08: ML流水线
# ============================================================
section_divider("ML流水线与数据质量", "Machine Learning Pipeline & Data Quality", "08", ACCENT_CYAN)

# ============================================================
# 幻灯片36：ML流水线8个步骤
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "机器学习流水线8个步骤", "End-to-End ML Workflow")

steps = [
    ("1", "问题定义", "回归/分类决定\n设定成功标准\n业务→ML转换", ACCENT_BLUE),
    ("2", "数据收集", "DB、API、爬虫\n公开数据集\n数据获取", ACCENT_CYAN),
    ("3", "数据\n预处理", "缺失值、异常值\n编码、缩放\n原始→模型用加工", ACCENT_GREEN),
    ("4", "EDA", "分布、相关性\n可视化\n了解数据特征", ACCENT_GREEN),
    ("5", "建模", "算法选择\n训练\n多模型比较", ACCENT_ORANGE),
    ("6", "评估", "交叉验证\n测试集评估\n泛化性能测量", ACCENT_ORANGE),
    ("7", "HP调参", "Grid/Random\nBayesian Search\n最优参数", ACCENT_RED),
    ("8", "部署", "REST API\n监控\n重训练流水线", ACCENT_PURPLE),
]
for i, (num, name, desc, color) in enumerate(steps):
    x = Inches(0.3) + Inches(1.6) * i
    if i < len(steps) - 1:
        add_shape(s, x + Inches(1.45), Inches(3.5), Inches(0.25), Inches(0.15), DARK_GRAY)
    add_shape(s, x, Inches(2.3), Inches(1.4), Inches(2.5), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.5), Inches(2.1), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.5), Inches(2.1), Inches(0.4), Inches(0.4),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), Inches(2.55), Inches(1.3), Inches(0.6),
             name, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), Inches(3.2), Inches(1.3), Inches(1.4),
             desc, font_size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 底部核心
add_shape(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(2.2), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.4),
         "\"Garbage In, Garbage Out\" - 数据质量的重要性", font_size=16, color=ACCENT_ORANGE, bold=True)
add_table_slide(s,
    ["预处理步骤", "说明", "方法"],
    [
        ["缺失值处理", "处理缺失数据", "删除、均值/中位数填充、插值"],
        ["异常值处理", "异常极端值", "基于IQR、Z-score检测"],
        ["编码", "类别型→数值转换", "独热编码、标签编码"],
        ["缩放", "统一变量尺度/范围", "StandardScaler、MinMaxScaler"],
    ],
    Inches(0.8), Inches(5.7), [2.5, 3.0, 5.5], row_height=0.4, font_size=12, header_font_size=13)

# ============================================================
# 幻灯片37：特征工程 & 数据泄露
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "特征工程 & 数据泄露防范")

# 特征工程
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Domingos经验7: Feature Engineering Is the Key", font_size=14, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "特征设计比算法选择影响更大",
    "基于领域知识构建特征是关键",
    "例：房价预测",
    "  简单：面积、房间数、位置",
    "  工程化：每房面积、地铁站距离、学区等级",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 数据泄露
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "Data Leakage（数据泄露）防范", font_size=14, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "训练过程中测试数据信息流入模型的现象",
    "→ 过度乐观的性能估计",
    "",
    "防范原则：",
    "  1. 缩放：训练集fit_transform，测试集仅transform",
    "  2. 在交叉验证循环内进行预处理（利用Pipeline）",
    "  3. 时间序列：不将未来数据用于训练",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# 过拟合 vs 欠拟合
add_shape(s, Inches(0.6), Inches(4.8), Inches(3.9), Inches(2.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.8), Inches(4.9), Inches(3.5), Inches(0.4),
         "欠拟合 (Underfitting)", font_size=15, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(0.8), Inches(5.3), Inches(3.5), Inches(1.8), [
    "训练/测试均低性能",
    "模型过于简单，高偏差",
    "解决：使用复杂模型、增加训练时间",
    "      特征工程、减少正则化",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(4.8), Inches(4.8), Inches(3.5), Inches(2.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(5.0), Inches(4.9), Inches(3.1), Inches(0.4),
         "最优模型 (Good Fit)", font_size=15, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(5.0), Inches(5.3), Inches(3.1), Inches(1.8), [
    "训练/测试均高性能",
    "Bias² + Variance 最小",
    "泛化性能最大",
    "通过交叉验证确认",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(8.6), Inches(4.8), Inches(3.9), Inches(2.5), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(8.8), Inches(4.9), Inches(3.5), Inches(0.4),
         "过拟合 (Overfitting)", font_size=15, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(8.8), Inches(5.3), Inches(3.5), Inches(1.8), [
    "训练高、测试低",
    "模型过于复杂，高方差",
    "解决：更多数据、正则化(L1/L2)",
    "      Dropout、早停、特征选择",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# SECTION 09: 核心论文综述
# ============================================================
section_divider("5篇核心论文综述", "Key Paper Reviews", "09", ACCENT_GREEN)

# ============================================================
# 幻灯片39：Domingos (2012) - 12条经验 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "Domingos (2012)：12条核心经验 (1/2)",
             "A Few Useful Things to Know About Machine Learning")

add_shape(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.6), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.5),
         "Pedro Domingos | Communications of the ACM, Vol. 55, No. 10, pp. 78-87 | ML实务12条经验",
         font_size=13, color=LIGHT_GRAY)

lessons1 = [
    ("1", "Representation + Evaluation\n+ Optimization", "所有学习算法的\n三个组成部分", ACCENT_BLUE),
    ("2", "It's Generalization\nThat Counts", "泛化性能是\n核心目标", ACCENT_BLUE),
    ("3", "Data Alone Is\nNot Enough", "先验假设（归纳偏置）\n是必需的", ACCENT_CYAN),
    ("4", "Overfitting Has\nMany Faces", "过拟合以各种\n形式出现", ACCENT_CYAN),
    ("5", "Intuition Fails in\nHigh Dimensions", "在高维空间中直觉\n失效（维度灾难）", ACCENT_GREEN),
    ("6", "Theoretical Guarantees\nAre Not What They Seem", "理论保证\n可能与实际不同", ACCENT_GREEN),
]
for i, (num, text, desc, color) in enumerate(lessons1):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.9) + Inches(2.3) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.0), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.6), y + Inches(0.15), Inches(3.2), Inches(0.8),
             text, font_size=13, color=WHITE)
    add_accent_line(s, x + Inches(0.2), y + Inches(1.0), Inches(3.0), color)
    add_text(s, x + Inches(0.2), y + Inches(1.15), Inches(3.5), Inches(0.7),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 幻灯片40：Domingos (2012) - 12条经验 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "Domingos (2012)：12条核心经验 (2/2)")

lessons2 = [
    ("7", "Feature Engineering\nIs the Key", "特征设计比\n算法选择更重要", ACCENT_ORANGE),
    ("8", "More Data Beats a\nCleverer Algorithm", "数据量比\n算法精巧更重要", ACCENT_ORANGE),
    ("9", "Learn Many Models,\nNot Just One", "集成方法\n优于单一模型", ACCENT_RED),
    ("10", "Simplicity Does Not\nImply Accuracy", "简单模型不一定\n更准确", ACCENT_RED),
    ("11", "Representable Does Not\nImply Learnable", "可表示 ≠ 可学习\n（学习难度另论）", ACCENT_PURPLE),
    ("12", "Correlation Does Not\nImply Causation", "相关 ≠ 因果\n（解读需注意）", ACCENT_PURPLE),
]
for i, (num, text, desc, color) in enumerate(lessons2):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.2) + Inches(2.3) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.0), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.6), y + Inches(0.15), Inches(3.2), Inches(0.8),
             text, font_size=13, color=WHITE)
    add_accent_line(s, x + Inches(0.2), y + Inches(1.0), Inches(3.0), color)
    add_text(s, x + Inches(0.2), y + Inches(1.15), Inches(3.5), Inches(0.7),
             desc, font_size=12, color=LIGHT_GRAY)

# 核心总结
add_shape(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(6.83), Inches(11.5), Inches(0.4),
         "实务核心：泛化为目标 | 特征设计是关键 | 大量数据 > 精巧算法 | 利用集成 | 相关≠因果",
         font_size=13, color=ACCENT_CYAN, bold=True)

# ============================================================
# 幻灯片41：Mitchell、Wolpert、Vapnik、ESL论文
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "核心论文综述：Mitchell、Wolpert、Vapnik、ESL")

papers = [
    ("Mitchell (1997)", "Machine Learning (教材)", [
        "确立ML的形式化定义(T-P-E)",
        "概念学习：假设空间H、版本空间",
        "归纳偏置(Inductive Bias)：",
        "  所有学习算法为了泛化",
        "  需要先验假设 → 与NFL直接相关",
    ], ACCENT_BLUE),
    ("Wolpert & Macready (1997)", "NFL Theorems for Optimization", [
        "证明通用最优算法不可能",
        "对所有函数取平均则所有算法相同",
        "\"对这个特定问题哪个算法",
        "  最合适？\"才是正确的问题",
    ], ACCENT_RED),
    ("Vapnik (1995)", "Statistical Learning Theory", [
        "经验风险: R_emp = (1/n)SUM L(f(xi),yi)",
        "真实风险: R(f) = INT L(f(x),y) dP(x,y)",
        "VC维: 可shatter的最大点数",
        "SRM: R(f) <= R_emp + Phi(VCdim, n)",
        "核心：有限数据 → 泛化条件？",
    ], ACCENT_GREEN),
    ("Hastie et al. (2009)", "ESL Ch1-2", [
        "EPE(f) = E[(Y-f(X))²]",
        "最优预测器：f*(x) = E[Y|X=x]",
        "Bias-Variance Decomposition",
        "最小二乘法(高偏差,低方差) vs",
        "k-NN(低偏差,高方差) 对比说明",
    ], ACCENT_PURPLE),
]
for i, (author, title, items, color) in enumerate(papers):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(6.3) * col
    y = Inches(2.2) + Inches(2.6) * row
    add_shape(s, x, y, Inches(6.0), Inches(2.3), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.35),
             author, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(5.6), Inches(0.3),
             title, font_size=12, color=WHITE)
    add_bullet_list(s, x + Inches(0.2), y + Inches(0.8), Inches(5.6), Inches(1.3),
                    items, font_size=11, color=LIGHT_GRAY, spacing=Pt(2))

# ============================================================
# SECTION 10: 实验1 - 偏差-方差权衡
# ============================================================
section_divider("实验1：偏差-方差权衡", "01_bias_variance_tradeoff.py", "10", ACCENT_BLUE)

# ============================================================
# 幻灯片43：实验1概述 + 核心代码
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "实验1：偏差-方差权衡实现", "01_bias_variance_tradeoff.py")

add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(1.5), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "实验目标", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(0.8), [
    "用不同复杂度的多项式回归模型可视化偏差²、方差、总误差",
    "直观确认欠拟合/过拟合区域",
    "实验验证偏差-方差分解公式",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(1.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "实验设置", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(0.8), [
    "真实函数：f(x) = sin(1.5*pi*x)",
    "噪声：e ~ N(0, 0.3²)，200个数据集重复",
    "多项式次数1~15变化模型复杂度",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 核心代码1：数据生成
add_text(s, Inches(0.6), Inches(4.0), Inches(5), Inches(0.3),
         "核心代码1：真实函数 & 数据生成", font_size=13, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8), [
    "def true_function(x):",
    "    return np.sin(1.5 * np.pi * x)",
    "",
    "def generate_data(n_samples=30, noise_std=0.3):",
    "    X = np.sort(np.random.uniform(0, 1, n_samples))",
    "    y = true_function(X) + \\",
    "        np.random.normal(0, noise_std, n_samples)",
    "    return X, y",
    "",
    "# y = f(x) + epsilon",
    "# epsilon ~ N(0, sigma^2)  <- 不可约误差",
], font_size=11)

# 核心代码2：模型训练
add_text(s, Inches(6.8), Inches(4.0), Inches(5), Inches(0.3),
         "核心代码2：多项式回归模型训练", font_size=13, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.8), [
    "from sklearn.pipeline import make_pipeline",
    "from sklearn.preprocessing import PolynomialFeatures",
    "from sklearn.linear_model import LinearRegression",
    "",
    "model = make_pipeline(",
    "    PolynomialFeatures(degree, include_bias=True),",
    "    LinearRegression()",
    ")",
    "model.fit(X_train.reshape(-1, 1), y_train)",
    "predictions[i] = model.predict(X_test.reshape(-1,1))",
], font_size=11)

# ============================================================
# 幻灯片44：实验1 - 偏差方差计算代码
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "实验1：偏差-方差计算核心逻辑")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(4.8), [
    "def compute_bias_variance(degrees, n_datasets=200, n_samples=30, noise_std=0.3):",
    "    X_test = np.linspace(0, 1, 100)   # 固定测试点",
    "    y_true = true_function(X_test)      # 真实函数值",
    "    noise_var = noise_std ** 2          # 不可约误差 = 0.09",
    "",
    "    for degree in degrees:",
    "        predictions = np.zeros((n_datasets, 100))  # 200个数据集 x 100测试点",
    "",
    "        for i in range(n_datasets):            # 重复200次",
    "            X_train, y_train = generate_data()  # 每次生成新数据集",
    "            model = make_pipeline(PolynomialFeatures(degree), LinearRegression())",
    "            model.fit(X_train.reshape(-1,1), y_train)",
    "            predictions[i,:] = model.predict(X_test.reshape(-1,1))",
    "",
    "        # ============ 核心：偏差-方差分解计算 ============",
    "        mean_pred = predictions.mean(axis=0)            # f_bar(x) = E[f_hat(x)]",
    "        bias_sq = np.mean((mean_pred - y_true) ** 2)    # Bias² = (f(x) - f_bar(x))²",
    "        var = np.mean(predictions.var(axis=0))           # Var = E[(f_hat - f_bar)²]",
    "        total = bias_sq + var + noise_var                # Total = Bias² + Var + sigma²",
], font_size=11)

add_shape(s, Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3), CARD_BG)
add_text(s, Inches(0.9), Inches(7.1), Inches(11.5), Inches(0.3),
         "核心：通过200个数据集反复训练 → 实验计算预测的平均（偏差）和变动（方差）",
         font_size=12, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 11: 实验2 - NFL定理演示
# ============================================================
section_divider("实验2：NFL定理演示", "02_no_free_lunch_demo.py", "11", ACCENT_ORANGE)

# ============================================================
# 幻灯片46：实验2概述 + 核心代码
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "11", "实验2：NFL定理演示", "02_no_free_lunch_demo.py")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "实验目标：对4种数据结构应用5种算法 → 实验确认不存在单一最佳算法",
         font_size=15, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(0.4),
         "通过Stratified 5-Fold交叉验证进行公平比较", font_size=14, color=LIGHT_GRAY)

# 数据集说明
ds_items = [
    ("线性可分 (Linear)", "线性决策边界\n→ 对线性模型有利", ACCENT_BLUE),
    ("圆形 (Circles)", "同心圆结构\n→ 对RBF SVM有利", ACCENT_CYAN),
    ("XOR模式", "XOR逻辑运算\n→ 对决策树有利", ACCENT_GREEN),
    ("半月 (Moons)", "半月形状\n→ 对k-NN有利", ACCENT_ORANGE),
]
add_text(s, Inches(0.6), Inches(3.6), Inches(3), Inches(0.3),
         "4种数据集：", font_size=14, color=WHITE, bold=True)
for i, (name, desc, color) in enumerate(ds_items):
    x = Inches(0.5) + Inches(3.15) * i
    add_shape(s, x, Inches(3.95), Inches(2.9), Inches(1.3), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(4.0), Inches(2.6), Inches(0.35),
             name, font_size=13, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(4.4), Inches(2.6), Inches(0.7),
             desc, font_size=11, color=LIGHT_GRAY)

# 算法 + 核心代码
add_text(s, Inches(0.6), Inches(5.5), Inches(5), Inches(0.3),
         "5种算法（各自不同的归纳偏置）：", font_size=14, color=WHITE, bold=True)
add_code_block(s, Inches(0.6), Inches(5.85), Inches(5.8), Inches(1.4), [
    "LogisticRegression()     # 线性决策边界",
    "KNeighborsClassifier(5)  # 局部相似性",
    "SVC(kernel='rbf')        # 最大化间隔 + 非线性",
    "DecisionTreeClassifier() # 轴对齐分割",
    "RandomForestClassifier() # 集成 + 分割",
], font_size=11)

add_code_block(s, Inches(6.8), Inches(5.85), Inches(5.8), Inches(1.4), [
    "# 公平比较: Stratified 5-Fold CV",
    "cv = StratifiedKFold(n_splits=5, shuffle=True)",
    "scores = cross_val_score(clf, X_scaled, y,",
    "                         cv=cv, scoring='accuracy')",
    "# 核心结论：没有在所有数据集中排名第一的算法！",
], font_size=11)

# ============================================================
# SECTION 12: 实验3 - 交叉验证比较
# ============================================================
section_divider("实验3：交叉验证比较", "03_cross_validation_demo.py", "12", ACCENT_RED)

# ============================================================
# 幻灯片48：实验3概述 + 核心代码
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "12", "实验3：交叉验证策略比较", "03_cross_validation_demo.py")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.0), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.35),
         "实验目标：Hold-Out、K-Fold、Stratified K-Fold、LOO比较 | 50次重复测量性能估计方差 | 确认K值权衡",
         font_size=14, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.3),
         "数据：make_classification(150样本) + Iris数据集", font_size=13, color=LIGHT_GRAY)

# 方差比较代码
add_text(s, Inches(0.6), Inches(3.4), Inches(6), Inches(0.3),
         "核心代码1：性能估计方差比较（50次重复）", font_size=13, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(0.6), Inches(3.7), Inches(5.8), Inches(3.0), [
    "n_repeats = 50",
    "",
    "# Hold-Out: 每次划分结果不同",
    "for i in range(n_repeats):",
    "    X_tr, X_te, y_tr, y_te = train_test_split(",
    "        X, y, test_size=0.2, random_state=i)",
    "    model.fit(X_tr, y_tr)",
    "    holdout_scores.append(model.score(X_te, y_te))",
    "",
    "# Stratified K-Fold: 最稳定",
    "for i in range(n_repeats):",
    "    cv = StratifiedKFold(n_splits=5, shuffle=True,",
    "                         random_state=i)",
    "    scores = cross_val_score(model, X, y, cv=cv)",
    "    skfold_scores.append(scores.mean())",
], font_size=10)

# K值权衡代码
add_text(s, Inches(6.8), Inches(3.4), Inches(6), Inches(0.3),
         "核心代码2：K值权衡", font_size=13, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(6.8), Inches(3.7), Inches(5.8), Inches(3.0), [
    "k_values = [2, 3, 5, 7, 10, 15, 20, 50, 100]",
    "",
    "for k in k_values:",
    "    repeat_means = []",
    "    for i in range(30):  # 重复30次",
    "        cv = KFold(n_splits=k, shuffle=True,",
    "                   random_state=i)",
    "        scores = cross_val_score(model, X, y, cv=cv)",
    "        repeat_means.append(scores.mean())",
    "    means.append(np.mean(repeat_means))",
    "    stds.append(np.std(repeat_means))",
    "",
    "# 核心观察:",
    "# K小 → 高偏差, 低方差",
    "# K大 → 低偏差, 高方差",
], font_size=10)

add_shape(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(6.93), Inches(11.5), Inches(0.3),
         "核心结论：Hold-Out方差最大 | Stratified K-Fold最稳定 | K=5~10实务推荐",
         font_size=13, color=ACCENT_GREEN, bold=True)

# ============================================================
# SECTION 13: 深入应用案例
# ============================================================
section_divider("深入应用案例", "Healthcare, Finance, Manufacturing, NLP", "13", ACCENT_PURPLE)

# ============================================================
# 幻灯片50：应用案例 - 医疗 & 金融
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "13", "深入应用案例：医疗 & 金融")

# 医疗
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "医疗 (Healthcare) - 乳腺癌诊断", font_size=16, color=ACCENT_GREEN, bold=True)
add_table_slide(s,
    ["项目", "内容"],
    [
        ["数据", "Wisconsin Breast Cancer（569样本，30个特征）"],
        ["类别", "恶性37.3% / 良性62.7%"],
        ["核心指标", "召回率(Recall) - 不遗漏癌症患者"],
        ["性能", "仅逻辑回归准确率97~98%"],
    ],
    Inches(0.8), Inches(2.8), [1.5, 3.8], row_height=0.4, font_size=12, header_font_size=13)

add_bullet_list(s, Inches(0.9), Inches(4.6), Inches(5.0), Inches(2.0), [
    "医疗AI核心考虑事项：",
    "  • 最小化FN（假阴性）事关生命",
    "  • 重视召回率和F1-Score而非准确率",
    "  • 模型可解释性必需（需要医生理解）",
    "  • 满足监管要求（FDA批准等）",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# 金融
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "金融 (Finance) - 欺诈检测", font_size=16, color=ACCENT_ORANGE, bold=True)
add_table_slide(s,
    ["项目", "内容"],
    [
        ["问题", "极端类别不平衡（欺诈 < 1%）"],
        ["核心课题", "最小化FP + 实际欺诈Recall"],
        ["主要技术", "Isolation Forest、SMOTE、集成"],
        ["特殊点", "实时处理 + 模式持续变化"],
    ],
    Inches(7.0), Inches(2.8), [1.5, 3.8], row_height=0.4, font_size=12, header_font_size=13)

add_bullet_list(s, Inches(7.1), Inches(4.6), Inches(5.0), Inches(2.0), [
    "金融ML核心考虑事项：",
    "  • 时间序列：Time Series Split，防止未来数据泄露",
    "  • 类别不平衡：SMOTE、Cost-sensitive Learning",
    "  • 可解释性：监管机构需要模型说明",
    "  • 需要实时推理性能",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 幻灯片51：应用案例 - 制造 & NLP
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "13", "深入应用案例：制造 & NLP")

# 制造
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "制造 (Manufacturing) - 预测性维护", font_size=16, color=ACCENT_BLUE, bold=True)
add_table_slide(s,
    ["项目", "内容"],
    [
        ["问题", "预测设备故障时间/可能性"],
        ["数据", "IoT传感器（振动、温度、压力、噪声）"],
        ["技术", "RF、LSTM、Survival Analysis"],
        ["效果", "非计划停机减少30~50%"],
    ],
    Inches(0.8), Inches(2.8), [1.5, 3.8], row_height=0.4, font_size=12, header_font_size=13)

add_bullet_list(s, Inches(0.9), Inches(4.6), Inches(5.0), Inches(2.0), [
    "制造ML核心：",
    "  • 传感器数据的时间序列特性",
    "  • 故障数据非常稀少（不平衡）",
    "  • 通过Edge Computing实时推理",
    "  • 领域专家（工艺工程师）协作必需",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# NLP
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "NLP - 情感分析", font_size=16, color=ACCENT_PURPLE, bold=True)
add_table_slide(s,
    ["项目", "内容"],
    [
        ["问题", "文本分类（正面/负面/中性）"],
        ["数据", "评论、社交媒体、新闻"],
        ["传统方法", "TF-IDF + SVM/逻辑回归"],
        ["现代方法", "BERT、GPT Fine-tuning"],
    ],
    Inches(7.0), Inches(2.8), [1.5, 3.8], row_height=0.4, font_size=12, header_font_size=13)

add_bullet_list(s, Inches(7.1), Inches(4.6), Inches(5.0), Inches(2.0), [
    "NLP ML核心：",
    "  • 预处理：分词、停用词去除、标准化",
    "  • 特征：BoW、TF-IDF、Word2Vec、BERT Embedding",
    "  • 迁移学习：预训练模型Fine-tuning",
    "  • 中文分词处理的特殊性",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# SECTION 14: 核心总结 + 复习题
# ============================================================
section_divider("核心总结 & 复习题", "Summary & Review Questions", "14", ACCENT_CYAN)

# ============================================================
# 幻灯片53：核心总结
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "14", "核心总结 (Key Takeaways)")

summaries = [
    ("ML定义", "从数据中自动学习模式\nMitchell T-P-E框架", ACCENT_BLUE),
    ("学习类型", "监督(分类/回归) · 无监督 · 强化\nAI ⊃ ML ⊃ DL", ACCENT_CYAN),
    ("偏差-方差", "EPE = Bias² + Var + σ²\n模型复杂度最优平衡点", ACCENT_GREEN),
    ("NFL定理", "不存在通用最佳算法\n根据问题选择是关键", ACCENT_ORANGE),
    ("交叉验证", "Stratified K-Fold (K=5~10)\n根据数据规模选择", ACCENT_RED),
    ("ML流水线", "问题定义→收集→预处理→EDA\n→建模→评估→调参→部署", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(summaries):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.2) + Inches(2.5) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.1), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
             title, font_size=17, color=color, bold=True)
    add_accent_line(s, x + Inches(0.2), y + Inches(0.6), Inches(2.0), color)
    add_text(s, x + Inches(0.2), y + Inches(0.75), Inches(3.5), Inches(1.1),
             desc, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 幻灯片54：核心公式整理
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "14", "核心公式整理")

formulas = [
    ("偏差-方差分解", "EPE = Bias² + Variance + σ²", "总误差的三个来源"),
    ("偏差", "Bias(x) = f(x) - E[f_hat(x)]", "模型假设造成的系统性误差"),
    ("方差", "Var(x) = E[(f_hat(x) - E[f_hat(x)])²]", "数据变动造成的预测不稳定"),
    ("NFL定理", "Sf P(d|f,m,a1) = Sf P(d|f,m,a2)", "所有算法的平均性能相同"),
    ("结构风险最小化", "R(f) <= R_emp(f) + Phi(VCdim, n)", "泛化误差上界（Vapnik）"),
    ("最优预测器（回归）", "f*(x) = E[Y|X=x]", "最小化MSE的预测（ESL）"),
]
for i, (name, formula, meaning) in enumerate(formulas):
    y = Inches(2.2) + Inches(0.85) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    add_shape(s, Inches(0.6), y, Inches(12.1), Inches(0.75), bg, radius=True)
    add_text(s, Inches(0.9), y + Inches(0.05), Inches(2.5), Inches(0.65),
             name, font_size=14, color=ACCENT_CYAN, bold=True)
    add_text(s, Inches(3.5), y + Inches(0.05), Inches(4.5), Inches(0.65),
             formula, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.2), y + Inches(0.05), Inches(4.3), Inches(0.65),
             meaning, font_size=13, color=DARK_GRAY)

# ============================================================
# 幻灯片55：复习题 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "14", "复习题 (1/2)")

questions1 = [
    "Q1. 请将Mitchell的T-P-E定义应用到垃圾邮件过滤器进行说明。",
    "Q2. 请从数据、目标、反馈角度比较监督/无监督/强化学习。",
    "Q3. 请解释偏差-方差分解公式的各项，\n      并阐述与模型复杂度的关系。",
    "Q4. 请阐述NFL定理的核心及3条实务启示。",
    "Q5. 请比较Hold-Out、K-Fold、Stratified K-Fold、LOO的\n      优缺点。",
]
colors5 = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
for i, q in enumerate(questions1):
    y = Inches(2.2) + Inches(1.0) * i
    add_shape(s, Inches(0.8), y, Inches(11.5), Inches(0.85), CARD_BG, colors5[i], radius=True)
    add_shape(s, Inches(0.8), y + Inches(0.15), Inches(0.08), Inches(0.5), colors5[i])
    add_text(s, Inches(1.2), y + Inches(0.1), Inches(10.8), Inches(0.7),
             q, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 幻灯片56：复习题 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "14", "复习题 (2/2)")

questions2 = [
    "Q6. 请各列举3种以上过拟合和欠拟合的原因及解决方法。\n      请结合Domingos的\"Overfitting Has Many Faces\"进行说明。",
    "Q7. 请解释什么是VC维，以及二维线性分类器的\n      VC维为3的原因。",
    "Q8. 请解释ML流水线8个步骤和Data Leakage防范原则。",
    "Q9. 请通过ESL的最小二乘法 vs k-NN比较\n      来解释偏差-方差权衡。",
    "Q10. 请结合实务案例说明Domingos的\n       \"More Data Beats a Cleverer Algorithm\"和\"Feature Engineering Is the Key\"。",
]
colors5b = [ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE]
for i, q in enumerate(questions2):
    y = Inches(2.2) + Inches(1.0) * i
    add_shape(s, Inches(0.8), y, Inches(11.5), Inches(0.85), CARD_BG, colors5b[i], radius=True)
    add_shape(s, Inches(0.8), y + Inches(0.15), Inches(0.08), Inches(0.5), colors5b[i])
    add_text(s, Inches(1.2), y + Inches(0.1), Inches(10.8), Inches(0.7),
             q, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 幻灯片57：Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.5), prs.slide_width, Inches(1.0),
         "Thank You", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), prs.slide_width, Inches(0.6),
         "第1章：机器学习概述（扩展详细版）", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(4.2), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(4.5), prs.slide_width, Inches(0.5),
         "下一章：第2章 - Python基础", font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = os.path.join(os.path.dirname(__file__), "1장_머신러닝_개요_강의PPT_확장_chz.pptx")
prs.save(output_path)
print(f"PPT generated: {output_path}")
print(f"Total slides: {len(prs.slides)}")
