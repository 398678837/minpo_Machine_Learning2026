"""1장 머신러닝 개요 - 강의 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    card = add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))


def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)


# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
# 상단 장식 라인
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
# 중앙 정보
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "기계학습 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "1장: 머신러닝 개요", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "Introduction to Machine Learning", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
         "핵심 키워드: ML 정의 · 학습 유형 · 편향-분산 트레이드오프 · NFL 정리 · 교차검증 · ML 파이프라인",
         font_size=14, color=DARK_GRAY)
# 하단 정보
add_shape(s, Inches(0), Inches(6.3), prs.slide_width, Inches(1.2), RGBColor(0x0A, 0x0A, 0x1A))
add_text(s, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
         "2026년 1학기", font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 2: 목차
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (Contents)")

toc = [
    ("01", "머신러닝이란", ACCENT_BLUE),
    ("02", "머신러닝의 역사", ACCENT_CYAN),
    ("03", "학습 유형: 지도 / 비지도 / 강화학습", ACCENT_GREEN),
    ("04", "ML vs DL vs AI", ACCENT_PURPLE),
    ("05", "편향-분산 트레이드오프", ACCENT_ORANGE),
    ("06", "No Free Lunch 정리", ACCENT_RED),
    ("07", "모델 평가: 교차검증 전략", ACCENT_BLUE),
    ("08", "ML 파이프라인과 데이터 품질", ACCENT_CYAN),
    ("09", "핵심 논문 리뷰 (5편)", ACCENT_GREEN),
    ("10", "실습 소개 & 핵심 요약", ACCENT_PURPLE),
]
for i, (num, title, color) in enumerate(toc):
    y = Inches(2.0) + Inches(0.5) * i
    add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.38), color, radius=True)
    add_text(s, Inches(1.2), y, Inches(0.55), Inches(0.38), num,
             font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), y, Inches(8), Inches(0.38), title,
             font_size=17, color=WHITE)

# ============================================================
# 슬라이드 3: 머신러닝이란 - 정의
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "머신러닝이란?", "Machine Learning Definition")

# Samuel 정의
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(1.6), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Arthur Samuel (1959)", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(1.0),
         "\"컴퓨터가 명시적으로 프로그래밍하지 않아도\n학습할 수 있는 능력을 부여하는 학문 분야\"",
         font_size=15, color=LIGHT_GRAY)

# Mitchell 정의
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(1.6), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "Tom Mitchell (1997)", font_size=14, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(1.0),
         "\"프로그램이 작업 T에 대해 경험 E로부터 학습하며,\n성능 P가 E를 통해 향상되면 학습한다\"",
         font_size=15, color=LIGHT_GRAY)

# T-P-E 테이블
labels = [("T (Task)", "수행 작업", "이메일 스팸 분류", ACCENT_BLUE),
          ("P (Performance)", "성능 지표", "정확히 분류된 비율", ACCENT_GREEN),
          ("E (Experience)", "학습 경험", "레이블링된 이메일", ACCENT_ORANGE)]
for i, (name, meaning, example, color) in enumerate(labels):
    x = Inches(0.6) + Inches(4.1) * i
    add_shape(s, x, Inches(4.3), Inches(3.8), Inches(2.2), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(4.4), Inches(3.4), Inches(0.4),
             name, font_size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, x + Inches(0.3), Inches(4.85), Inches(3.2), color)
    add_text(s, x + Inches(0.2), Inches(5.0), Inches(3.4), Inches(0.4),
             meaning, font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(5.5), Inches(3.4), Inches(0.6),
             f"예: {example}", font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 4: 전통적 프로그래밍 vs 머신러닝
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "전통적 프로그래밍 vs 머신러닝")

# 왼쪽: 전통적
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.8), Inches(2.3), Inches(5.4), Inches(0.5),
         "전통적 프로그래밍", font_size=22, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(1.0), Inches(3.0), Inches(5.0), Inches(3.5), [
    "입력: 데이터 + 규칙 (Rules)",
    "출력: 결과 (Answers)",
    "사람이 규칙을 직접 코딩",
    "",
    "예시: if-else 기반 스팸 필터",
    "  → 특정 키워드가 포함되면 스팸으로 분류",
    "  → 새로운 패턴에 대응 어려움",
], font_size=15, color=LIGHT_GRAY)

# 오른쪽: 머신러닝
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.0), Inches(2.3), Inches(5.4), Inches(0.5),
         "머신러닝", font_size=22, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(7.2), Inches(3.0), Inches(5.0), Inches(3.5), [
    "입력: 데이터 + 정답 (Labels)",
    "출력: 규칙/모델 (Model)",
    "데이터로부터 규칙을 자동 학습",
    "",
    "예시: 수만 건의 이메일로 학습한 스팸 필터",
    "  → 패턴을 자동으로 발견",
    "  → 새로운 패턴에도 적응 가능",
], font_size=15, color=LIGHT_GRAY)

# 중앙 VS
add_shape(s, Inches(6.0), Inches(4.0), Inches(1.2), Inches(0.8), ACCENT_BLUE, radius=True)
add_text(s, Inches(6.0), Inches(4.05), Inches(1.2), Inches(0.7),
         "VS", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 5: 일상 속 머신러닝 사례
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "일상 속 머신러닝 사례 5가지")

cases = [
    ("추천 시스템", "넷플릭스, 유튜브\n시청 이력 기반\n개인화 추천", ACCENT_BLUE),
    ("자율주행", "테슬라, 웨이모\nCNN 기반 객체 탐지\n실시간 경로 계획", ACCENT_CYAN),
    ("의료 진단", "X-ray, CT, MRI\n종양·폐렴 탐지\nAlphaFold", ACCENT_GREEN),
    ("음성 인식", "시리, 알렉사\nSTT + NLU\nTransformer 기반", ACCENT_PURPLE),
    ("사기 탐지", "신용카드 부정사용\n이상 거래 감지\n실시간 분석", ACCENT_RED),
]
for i, (title, desc, color) in enumerate(cases):
    x = Inches(0.4) + Inches(2.5) * i
    add_shape(s, x, Inches(2.2), Inches(2.3), Inches(4.3), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(2.3), Inches(0.6), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(2.3), Inches(0.5),
             title, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.1), Inches(1.9), Inches(3.0),
             desc, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 6: 머신러닝의 역사
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "머신러닝의 역사", "Major Milestones")

milestones = [
    ("1943", "McCulloch & Pitts\n인공 뉴런 모델"),
    ("1957", "Rosenblatt\n퍼셉트론"),
    ("1959", "Arthur Samuel\nML 용어 최초 사용"),
    ("1986", "Rumelhart et al.\n역전파 알고리즘"),
    ("1995", "Vapnik\nSVM"),
    ("1997", "NFL 정리\nMitchell 교재"),
    ("2012", "AlexNet\n딥러닝 붐"),
    ("2017", "Transformer\nAttention Is All You Need"),
    ("2022~", "ChatGPT\n생성형 AI"),
]
# 타임라인 선
add_shape(s, Inches(0.8), Inches(3.6), Inches(11.7), Pt(3), ACCENT_BLUE)
for i, (year, desc) in enumerate(milestones):
    x = Inches(0.6) + Inches(1.33) * i
    # 원
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), Inches(3.45), Inches(0.35), Inches(0.35)) if False else \
          add_shape(s, x + Inches(0.35), Inches(3.45), Inches(0.25), Inches(0.25), ACCENT_CYAN, radius=True)
    add_text(s, x, Inches(2.5), Inches(1.2), Inches(0.4),
             year, font_size=13, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x - Inches(0.1), Inches(4.0), Inches(1.4), Inches(1.5),
             desc, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# AI 겨울
add_shape(s, Inches(0.6), Inches(5.7), Inches(12), Inches(1.3), CARD_BG, radius=True)
add_text(s, Inches(0.9), Inches(5.8), Inches(11.4), Inches(0.4),
         "AI 겨울과 부활", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(0.9), Inches(6.2), Inches(11.4), Inches(0.7),
         "1차 AI 겨울 (1970s): XOR 문제 → 부활 (1986): 역전파 → 2차 AI 겨울 (1990s 초): SVM 주류 → 딥러닝 부활 (2012~): GPU + 빅데이터 + ReLU",
         font_size=13, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 7: 학습 유형 - 지도학습
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "지도학습 (Supervised Learning)")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.3), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "정의: 입력 X와 정답 레이블 y가 함께 주어진 상태에서 매핑 함수 f: X → y를 학습", font_size=16, color=WHITE)
add_text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(0.5),
         "핵심: 예측값과 실제값의 차이(오차)를 줄이는 방향으로 학습  |  회귀 (연속값) + 분류 (이산값)", font_size=14, color=LIGHT_GRAY)

# 회귀 카드
add_card(s, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.2),
         "회귀 (Regression) - 연속값 예측", [
             "출력이 연속적 수치: 주택 가격, 기온, 매출",
             "대표 알고리즘: 선형회귀, 릿지, 라쏘, RF, GBM",
             "평가 지표: MSE, RMSE, MAE, R²",
             "예: 면적·위치·방 수 → 350,000달러",
         ], ACCENT_GREEN, ACCENT_GREEN)

# 분류 카드
add_card(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.2),
         "분류 (Classification) - 범주 예측", [
             "출력이 이산적 카테고리: 스팸/정상, 악성/양성",
             "이진분류 / 다중분류 / 다중레이블분류",
             "평가 지표: Accuracy, Precision, Recall, F1, AUC",
             "예: 이메일 내용 → 스팸 or 정상",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# 슬라이드 8: 비지도학습 & 강화학습
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "비지도학습 & 강화학습")

# 비지도학습
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.5),
         "비지도학습 (Unsupervised Learning)", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.9), Inches(5.2), Inches(4.0), [
    "정답 레이블 없이 데이터 구조 발견",
    "",
    "군집화: K-Means, DBSCAN, 계층적 군집화",
    "차원 축소: PCA, t-SNE, Autoencoder",
    "이상치 탐지: Isolation Forest",
    "",
    "활용: 고객 세분화, 유전자 군집화, 시각화",
], font_size=14, color=LIGHT_GRAY)

# 강화학습
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.5),
         "강화학습 (Reinforcement Learning)", font_size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.9), Inches(5.2), Inches(4.0), [
    "에이전트가 환경과 상호작용하며 보상 최대화",
    "",
    "에이전트 →[행동]→ 환경",
    "환경 →[상태, 보상]→ 에이전트",
    "",
    "알고리즘: Q-Learning, DQN, PPO, SAC",
    "활용: AlphaGo, 로봇 제어, 자율주행",
], font_size=14, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 9: 세 가지 학습 유형 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "세 가지 학습 유형 비교 요약")

headers = ["구분", "지도학습", "비지도학습", "강화학습"]
rows = [
    ["데이터", "입력 + 정답", "입력만", "상태 + 보상"],
    ["목표", "예측 / 분류", "구조 발견", "보상 최대화"],
    ["피드백", "정답 레이블", "없음", "보상 신호"],
    ["대표 문제", "분류, 회귀", "군집화, 차원축소", "게임, 로봇 제어"],
    ["난이도", "상대적 쉬움", "중간", "상대적 어려움"],
]
header_colors = [DARK_GRAY, ACCENT_BLUE, ACCENT_GREEN, ACCENT_PURPLE]

# 헤더 행
for j, (h, hc) in enumerate(zip(headers, header_colors)):
    x = Inches(1.0) + Inches(2.8) * j
    add_shape(s, x, Inches(2.3), Inches(2.6), Inches(0.6), hc)
    add_text(s, x, Inches(2.3), Inches(2.6), Inches(0.6),
             h, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 데이터 행
for i, row in enumerate(rows):
    y = Inches(3.0) + Inches(0.65) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(1.0) + Inches(2.8) * j
        add_shape(s, x, y, Inches(2.6), Inches(0.55), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y, Inches(2.6), Inches(0.55),
                 cell, font_size=14, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 10: ML vs DL vs AI
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "ML vs DL vs AI 관계", "계층 구조와 비교")

# 동심원 계층
circles = [
    (Inches(1.8), Inches(2.3), Inches(5.0), Inches(4.8), RGBColor(0x1E, 0x3A, 0x5F), "인공지능 (AI)", ACCENT_BLUE),
    (Inches(2.5), Inches(2.8), Inches(3.6), Inches(3.6), RGBColor(0x1E, 0x4D, 0x3A), "머신러닝 (ML)", ACCENT_GREEN),
    (Inches(3.2), Inches(3.3), Inches(2.2), Inches(2.4), RGBColor(0x3A, 0x1E, 0x5F), "딥러닝 (DL)", ACCENT_PURPLE),
]
for x, y, w, h, color, label, text_color in circles:
    shape = add_shape(s, x, y, w, h, color, text_color, radius=True)
    shape.fill.fore_color.rgb = color

for x, y, w, h, color, label, text_color in circles:
    if label == "인공지능 (AI)":
        add_text(s, x + Inches(0.3), y + Inches(0.2), w - Inches(0.6), Inches(0.4),
                 label, font_size=16, color=text_color, bold=True)
    elif label == "머신러닝 (ML)":
        add_text(s, x + Inches(0.3), y + Inches(0.1), w - Inches(0.6), Inches(0.4),
                 label, font_size=15, color=text_color, bold=True)
    else:
        add_text(s, x + Inches(0.2), Inches(4.2), w - Inches(0.4), Inches(0.4),
                 label, font_size=14, color=text_color, bold=True, align=PP_ALIGN.CENTER)

# 비교 테이블
comp_items = [
    ("피처 엔지니어링", "수동 설계", "자동 학습"),
    ("데이터 요구량", "적음~중간", "대량 필요"),
    ("계산 자원", "CPU 가능", "GPU/TPU 필수"),
    ("해석 가능성", "높음", "블랙박스"),
    ("적합 데이터", "정형 (테이블)", "비정형 (이미지, 텍스트)"),
]

add_shape(s, Inches(7.2), Inches(2.3), Inches(5.5), Inches(0.55), ACCENT_BLUE)
add_text(s, Inches(7.2), Inches(2.3), Inches(2.0), Inches(0.55),
         "특성", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(9.2), Inches(2.3), Inches(1.8), Inches(0.55),
         "전통 ML", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(11.0), Inches(2.3), Inches(1.7), Inches(0.55),
         "딥러닝", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, (feat, ml, dl) in enumerate(comp_items):
    y = Inches(2.95) + Inches(0.55) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    add_shape(s, Inches(7.2), y, Inches(5.5), Inches(0.5), bg)
    add_text(s, Inches(7.2), y, Inches(2.0), Inches(0.5),
             feat, font_size=13, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.2), y, Inches(1.8), Inches(0.5),
             ml, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(11.0), y, Inches(1.7), Inches(0.5),
             dl, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 하단 가이드
add_shape(s, Inches(7.2), Inches(5.8), Inches(5.5), Inches(1.3), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.5), Inches(5.9), Inches(5.0), Inches(0.3),
         "언제 무엇을 사용할 것인가? (Domingos, 2012)", font_size=13, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(7.5), Inches(6.2), Inches(5.0), Inches(0.9), [
    "정형 데이터 + 소~중규모 → XGBoost, LightGBM, RF",
    "비정형 데이터 (이미지, 텍스트) → CNN, Transformer",
    "해석 가능성 중요 (의료, 금융) → 전통 ML or XAI",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 11: 편향-분산 트레이드오프 (1)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "편향-분산 트레이드오프", "Bias-Variance Tradeoff")

# 편향 카드
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.5),
         "편향 (Bias)", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "모델 예측값과 실제 정답 사이의 체계적 오차",
    "높은 편향 → 모델이 너무 단순 → 과소적합",
    "낮은 편향 → 데이터 패턴을 잘 포착",
], font_size=14, color=LIGHT_GRAY)

# 분산 카드
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.5),
         "분산 (Variance)", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "다른 학습 데이터셋에서 예측값의 변동 정도",
    "높은 분산 → 학습 데이터에 민감 → 과적합",
    "낮은 분산 → 데이터 바뀌어도 안정적 예측",
], font_size=14, color=LIGHT_GRAY)

# 수식
add_shape(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(0.5),
         "EPE = Bias² + Variance + σ²  (비가약 오류)", font_size=24, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.4),
         "총 오차 = 모델 가정에 의한 체계적 오차  +  데이터 변동에 의한 불안정성  +  데이터 자체 노이즈 (줄일 수 없음)",
         font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 과녁 비유
targets = [
    ("낮은 편향\n낮은 분산", "이상적", ACCENT_GREEN),
    ("높은 편향\n낮은 분산", "과소적합", ACCENT_ORANGE),
    ("낮은 편향\n높은 분산", "과적합", ACCENT_RED),
    ("높은 편향\n높은 분산", "최악", DARK_GRAY),
]
add_text(s, Inches(0.6), Inches(6.2), Inches(3), Inches(0.4),
         "과녁 비유:", font_size=14, color=WHITE, bold=True)
for i, (label, status, color) in enumerate(targets):
    x = Inches(0.6) + Inches(3.1) * i
    add_shape(s, x, Inches(6.55), Inches(2.8), Inches(0.7), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.1), Inches(6.55), Inches(1.5), Inches(0.7),
             label, font_size=10, color=LIGHT_GRAY)
    add_text(s, x + Inches(1.5), Inches(6.55), Inches(1.2), Inches(0.7),
             status, font_size=13, color=color, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 12: 편향-분산 - 모델별 특성
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "모델별 편향-분산 특성")

models = [
    ("선형 회귀", "높음", "낮음", "단순하지만 안정적"),
    ("KNN (K 큼)", "높음", "낮음", "과도하게 평활화"),
    ("KNN (K 작음)", "낮음", "높음", "노이즈에 민감"),
    ("결정 트리 (깊음)", "낮음", "높음", "과적합 경향"),
    ("랜덤 포레스트", "낮음", "적절", "앙상블로 분산 감소"),
    ("부스팅", "낮음", "적절~높음", "편향 감소에 강점"),
]

h_items = ["모델", "편향", "분산", "특징"]
for j, h in enumerate(h_items):
    x = Inches(1.0) + Inches(2.8) * j
    add_shape(s, x, Inches(2.3), Inches(2.6), Inches(0.55), ACCENT_BLUE)
    add_text(s, x, Inches(2.3), Inches(2.6), Inches(0.55),
             h, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(models):
    y = Inches(2.95) + Inches(0.55) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(1.0) + Inches(2.8) * j
        add_shape(s, x, y, Inches(2.6), Inches(0.5), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        if j == 1 and cell == "높음":
            fc = ACCENT_RED
        elif j == 1 and cell == "낮음":
            fc = ACCENT_GREEN
        elif j == 2 and cell == "높음":
            fc = ACCENT_RED
        elif j == 2 and cell in ("낮음", "적절"):
            fc = ACCENT_GREEN
        add_text(s, x, y, Inches(2.6), Inches(0.5),
                 cell, font_size=13, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 하단 인용
add_shape(s, Inches(1.0), Inches(6.3), Inches(11.2), Inches(0.7), CARD_BG, radius=True)
add_text(s, Inches(1.3), Inches(6.35), Inches(10.6), Inches(0.6),
         "ESL (Hastie et al., 2009) Ch2: 최소제곱법 (높은 편향, 낮은 분산) vs k-NN (낮은 편향, 높은 분산) 비교를 통한 직관적 설명",
         font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 13: NFL 정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "No Free Lunch 정리", "Wolpert & Macready (1997)")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.5), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5),
         "핵심 주장", font_size=16, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(0.8),
         "모든 가능한 문제(목적 함수)에 대해 평균을 내면, 어떤 두 최적화(또는 학습) 알고리즘이든 동일한 성능을 보인다.\n→ \"최고의 알고리즘\"은 존재하지 않는다. 항상 문제에 맞는 알고리즘을 선택해야 한다.",
         font_size=15, color=LIGHT_GRAY)

# 수식
add_shape(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.8), RGBColor(0x15, 0x15, 0x30))
add_text(s, Inches(0.6), Inches(4.05), Inches(12.1), Inches(0.7),
         "Σf P(d_m^y | f, m, a₁) = Σf P(d_m^y | f, m, a₂)    ← 모든 알고리즘 쌍 (a₁, a₂)에 대해 성립",
         font_size=16, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# 시사점
implications = [
    ("1", "\"최고의 알고리즘\"은 없다", "항상 문제에 맞는 알고리즘 선택 필요", ACCENT_BLUE),
    ("2", "사전 지식 활용", "문제 구조(귀납적 편향)가 알고리즘 선택의 핵심", ACCENT_GREEN),
    ("3", "다양한 알고리즘 시도", "하나에 의존하지 말고 여러 알고리즘 비교", ACCENT_ORANGE),
    ("4", "AutoML의 근거", "자동 알고리즘 선택/하이퍼파라미터 튜닝의 당위성", ACCENT_PURPLE),
]
add_text(s, Inches(0.6), Inches(5.0), Inches(5), Inches(0.4),
         "실무적 시사점", font_size=16, color=WHITE, bold=True)
for i, (num, title, desc, color) in enumerate(implications):
    x = Inches(0.6) + Inches(3.1) * i
    add_shape(s, x, Inches(5.5), Inches(2.8), Inches(1.6), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(5.6), Inches(2.4), Inches(0.4),
             f"{num}. {title}", font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(6.05), Inches(2.4), Inches(0.8),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 14: 교차검증 전략
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "모델 평가: 교차검증 전략", "Cross-Validation Strategies")

strategies = [
    ("Hold-Out", "한 번만 분할 (80:20)\n장점: 빠름\n단점: 분할에 따라 결과 변동", ACCENT_ORANGE),
    ("K-Fold", "K개 폴드로 나눠 K번 평가\n모든 데이터 활용\nK=5 또는 K=10 표준", ACCENT_BLUE),
    ("Stratified K-Fold", "클래스 비율 유지하며 K-Fold\n불균형 데이터에 안정적\nsklearn 기본 CV 전략", ACCENT_GREEN),
    ("LOO (Leave-One-Out)", "K=N (한 번에 1개만 테스트)\n편향 가장 낮음\n데이터 적을 때 (N<50)", ACCENT_PURPLE),
]
for i, (name, desc, color) in enumerate(strategies):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(2.2), Inches(2.9), Inches(3.0), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(2.9), Inches(0.6), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(2.9), Inches(0.5),
             name, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.0), Inches(2.5), Inches(2.0),
             desc, font_size=13, color=LIGHT_GRAY)

# K-Fold 시각화
add_shape(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.6), CARD_BG, radius=True)
add_text(s, Inches(0.9), Inches(5.55), Inches(3), Inches(0.4),
         "5-Fold 교차검증 예시:", font_size=14, color=WHITE, bold=True)
fold_colors = [ACCENT_RED, ACCENT_BLUE, ACCENT_BLUE, ACCENT_BLUE, ACCENT_BLUE]
for fold_i in range(5):
    y = Inches(5.95) + Inches(0.22) * fold_i
    for block_j in range(5):
        x = Inches(4.0) + Inches(1.4) * block_j
        is_val = (block_j == fold_i)
        c = ACCENT_RED if is_val else RGBColor(0x2A, 0x4A, 0x6A)
        add_shape(s, x, y, Inches(1.3), Inches(0.18), c)
    add_text(s, Inches(11.2), y, Inches(1.5), Inches(0.2),
             f"→ 성능{fold_i + 1}", font_size=10, color=DARK_GRAY)

# ============================================================
# 슬라이드 15: 데이터 크기별 권장 전략
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "교차검증: 전략 비교 & 권장 가이드")

# 비교 테이블
th = ["전략", "학습 데이터", "평가 횟수", "편향", "분산", "계산 비용"]
trows = [
    ["Hold-Out", "80%", "1회", "높음", "높음", "매우 낮음"],
    ["5-Fold", "80% (반복)", "5회", "중간", "중간", "중간"],
    ["10-Fold", "90% (반복)", "10회", "낮음", "중간~높음", "중간~높음"],
    ["LOO", "(N-1)/N", "N회", "가장 낮음", "높음", "매우 높음"],
]
for j, h in enumerate(th):
    x = Inches(0.5) + Inches(2.05) * j
    add_shape(s, x, Inches(2.2), Inches(1.95), Inches(0.5), ACCENT_BLUE)
    add_text(s, x, Inches(2.2), Inches(1.95), Inches(0.5),
             h, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(trows):
    y = Inches(2.8) + Inches(0.5) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.5) + Inches(2.05) * j
        add_shape(s, x, y, Inches(1.95), Inches(0.45), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y, Inches(1.95), Inches(0.45),
                 cell, font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 데이터 크기별 권장
recs = [
    ("대용량 (>10,000)", "Hold-Out 또는 3-Fold", "충분한 데이터, 계산 효율성", ACCENT_BLUE),
    ("중간 (100~10,000)", "Stratified 5 or 10-Fold", "안정적 추정, 적절한 비용", ACCENT_GREEN),
    ("소용량 (<100)", "LOO 또는 Repeated K-Fold", "모든 데이터 최대 활용", ACCENT_ORANGE),
]
add_text(s, Inches(0.5), Inches(5.1), Inches(5), Inches(0.4),
         "데이터 크기별 권장 전략", font_size=16, color=WHITE, bold=True)
for i, (size, strategy, reason, color) in enumerate(recs):
    x = Inches(0.5) + Inches(4.15) * i
    add_shape(s, x, Inches(5.6), Inches(3.9), Inches(1.5), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(5.7), Inches(3.5), Inches(0.4),
             size, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(6.1), Inches(3.5), Inches(0.4),
             strategy, font_size=14, color=WHITE, bold=True)
    add_text(s, x + Inches(0.2), Inches(6.5), Inches(3.5), Inches(0.4),
             reason, font_size=12, color=DARK_GRAY)

# ============================================================
# 슬라이드 16: ML 파이프라인
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "ML 파이프라인", "Machine Learning Pipeline")

steps = [
    ("1", "문제\n정의", "회귀/분류 결정\n성공 기준 설정", ACCENT_BLUE),
    ("2", "데이터\n수집", "DB, API, 크롤링\n공개 데이터셋", ACCENT_CYAN),
    ("3", "데이터\n전처리", "결측치, 이상치\n인코딩, 스케일링", ACCENT_GREEN),
    ("4", "EDA", "분포, 상관관계\n시각화", ACCENT_GREEN),
    ("5", "모델링", "알고리즘 선택\n학습", ACCENT_ORANGE),
    ("6", "평가", "교차검증\n테스트셋 평가", ACCENT_ORANGE),
    ("7", "HP\n튜닝", "Grid/Random\nBayesian", ACCENT_RED),
    ("8", "배포", "REST API\n모니터링", ACCENT_PURPLE),
]
for i, (num, name, desc, color) in enumerate(steps):
    x = Inches(0.3) + Inches(1.6) * i
    # 화살표 (마지막 제외)
    if i < len(steps) - 1:
        add_shape(s, x + Inches(1.45), Inches(3.3), Inches(0.25), Inches(0.15), DARK_GRAY)

    add_shape(s, x, Inches(2.3), Inches(1.4), Inches(2.0), CARD_BG, color, radius=True)
    # 넘버 원
    add_shape(s, x + Inches(0.5), Inches(2.1), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.5), Inches(2.1), Inches(0.4), Inches(0.4),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(2.55), Inches(1.2), Inches(0.7),
             name, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(3.3), Inches(1.2), Inches(0.9),
             desc, font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 데이터 품질
add_shape(s, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.2), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(5.2), Inches(0.4),
         "\"Garbage In, Garbage Out\"", font_size=16, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(5.3), Inches(5.2), Inches(1.5), [
    "결측치 처리: 삭제, 평균/중앙값 대체, 보간",
    "이상치 처리: IQR, Z-score 기반 탐지",
    "인코딩: 원핫, 레이블 인코딩",
    "스케일링: StandardScaler, MinMaxScaler",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 피처 엔지니어링
add_shape(s, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.2), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(4.9), Inches(5.2), Inches(0.4),
         "Feature Engineering Is the Key (Domingos, 2012)", font_size=14, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(5.3), Inches(5.2), Inches(1.5), [
    "알고리즘 선택보다 피처 설계가 더 큰 영향",
    "도메인 지식 기반 피처 구성이 핵심",
    "Data Leakage 방지: Pipeline 활용",
    "스케일링은 학습셋 fit, 테스트셋 transform만",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 17: 과적합 vs 과소적합
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "과적합 (Overfitting) vs 과소적합 (Underfitting)")

# 과소적합
add_shape(s, Inches(0.6), Inches(2.2), Inches(3.8), Inches(4.8), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.6), Inches(2.3), Inches(3.8), Inches(0.5),
         "과소적합", font_size=20, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(0.9), Inches(2.9), Inches(3.2), Inches(3.8), [
    "학습/테스트 모두 낮은 성능",
    "모델이 너무 단순",
    "높은 편향, 낮은 분산",
    "",
    "해결:",
    "• 복잡한 모델 사용",
    "• 학습 시간 늘리기",
    "• 피처 엔지니어링",
    "• 정규화 강도 줄이기",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 적절
add_shape(s, Inches(4.75), Inches(2.2), Inches(3.8), Inches(4.8), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(4.75), Inches(2.3), Inches(3.8), Inches(0.5),
         "최적 모델", font_size=20, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(5.05), Inches(2.9), Inches(3.2), Inches(3.8), [
    "학습/테스트 모두 높은 성능",
    "적절한 모델 복잡도",
    "편향과 분산의 균형",
    "",
    "목표:",
    "• Bias² + Variance 최소",
    "• 일반화 성능 최대",
    "• 교차검증으로 확인",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 과적합
add_shape(s, Inches(8.9), Inches(2.2), Inches(3.8), Inches(4.8), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(8.9), Inches(2.3), Inches(3.8), Inches(0.5),
         "과적합", font_size=20, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(9.2), Inches(2.9), Inches(3.2), Inches(3.8), [
    "학습 높음, 테스트 낮음",
    "모델이 너무 복잡",
    "낮은 편향, 높은 분산",
    "",
    "해결:",
    "• 더 많은 데이터 수집",
    "• 정규화 (L1/L2)",
    "• 드롭아웃, 조기 종료",
    "• 특성 선택",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 18: 핵심 논문 리뷰
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "핵심 논문 리뷰 (5편)", "Key Paper Reviews")

papers = [
    ("Domingos (2012)", "A Few Useful Things to Know\nAbout Machine Learning",
     "ML 실무 12가지 교훈\n일반화, 피처, 데이터, 앙상블", ACCENT_BLUE),
    ("Mitchell (1997)", "Machine Learning (교재)",
     "ML의 형식적 정의 (T-P-E)\n귀납적 편향, 개념 학습", ACCENT_CYAN),
    ("Wolpert (1997)", "No Free Lunch Theorems\nfor Optimization",
     "보편적 최적 알고리즘\n불가능성 증명", ACCENT_RED),
    ("Vapnik (1995)", "The Nature of Statistical\nLearning Theory",
     "VC 차원, SRM\n일반화 오류 상한", ACCENT_GREEN),
    ("Hastie et al. (2009)", "Elements of Statistical\nLearning (ESL) Ch1-2",
     "편향-분산 분해\n최소제곱법 vs k-NN", ACCENT_PURPLE),
]
for i, (author, title, contrib, color) in enumerate(papers):
    x = Inches(0.2) + Inches(2.6) * i
    add_shape(s, x, Inches(2.2), Inches(2.4), Inches(4.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(2.4), Inches(0.5), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(2.4), Inches(0.4),
             author, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(2.9), Inches(2.1), Inches(1.2),
             title, font_size=12, color=LIGHT_GRAY)
    add_accent_line(s, x + Inches(0.2), Inches(4.2), Inches(2.0), color)
    add_text(s, x + Inches(0.15), Inches(4.4), Inches(2.1), Inches(2.0),
             contrib, font_size=12, color=DARK_GRAY)

# ============================================================
# 슬라이드 19: Domingos 12가지 교훈
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "Domingos (2012): 12가지 핵심 교훈")

lessons = [
    ("1", "Representation + Evaluation\n+ Optimization", ACCENT_BLUE),
    ("2", "It's Generalization\nThat Counts", ACCENT_BLUE),
    ("3", "Data Alone Is\nNot Enough", ACCENT_CYAN),
    ("4", "Overfitting Has\nMany Faces", ACCENT_CYAN),
    ("5", "Intuition Fails in\nHigh Dimensions", ACCENT_GREEN),
    ("6", "Theoretical Guarantees\nAre Not What They Seem", ACCENT_GREEN),
    ("7", "Feature Engineering\nIs the Key", ACCENT_ORANGE),
    ("8", "More Data Beats a\nCleverer Algorithm", ACCENT_ORANGE),
    ("9", "Learn Many Models,\nNot Just One", ACCENT_RED),
    ("10", "Simplicity ≠ Accuracy", ACCENT_RED),
    ("11", "Representable ≠\nLearnable", ACCENT_PURPLE),
    ("12", "Correlation ≠\nCausation", ACCENT_PURPLE),
]
for i, (num, text, color) in enumerate(lessons):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(3.15) * col
    y = Inches(2.2) + Inches(1.7) * row
    add_shape(s, x, y, Inches(2.9), Inches(1.4), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4),
             num, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.6), y + Inches(0.15), Inches(2.2), Inches(1.1),
             text, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 20: 평가 지표 개요
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "평가 지표 개요", "Evaluation Metrics")

# 회귀 지표
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "회귀 평가 지표", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(1.3), [
    "RMSE: √(평균(예측-실제)²) — 큰 오차에 민감",
    "MAE: 평균(|예측-실제|) — 이상치에 강건",
    "R²: 설명력 지표 (1에 가까울수록 좋음)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 분류 지표
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "분류 평가 지표", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(1.3), [
    "Accuracy: 전체 중 맞춘 비율 (불균형 주의)",
    "Precision: 양성 예측 중 실제 양성 비율",
    "Recall: 실제 양성 중 맞춘 비율 (FN 최소화)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 혼동 행렬
add_shape(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.8), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(4.6), Inches(5.2), Inches(0.4),
         "혼동 행렬 (Confusion Matrix)", font_size=15, color=ACCENT_CYAN, bold=True)
# 행렬 시각화
cells = [
    (Inches(2.8), Inches(5.2), "TP", ACCENT_GREEN),
    (Inches(4.3), Inches(5.2), "FN", ACCENT_RED),
    (Inches(2.8), Inches(5.9), "FP", ACCENT_RED),
    (Inches(4.3), Inches(5.9), "TN", ACCENT_GREEN),
]
add_text(s, Inches(2.8), Inches(4.9), Inches(1.3), Inches(0.3),
         "예측 Pos", font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.3), Inches(4.9), Inches(1.3), Inches(0.3),
         "예측 Neg", font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.3), Inches(5.3), Inches(1.3), Inches(0.3),
         "실제 Pos", font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.3), Inches(6.0), Inches(1.3), Inches(0.3),
         "실제 Neg", font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
for cx, cy, label, color in cells:
    add_shape(s, cx, cy, Inches(1.2), Inches(0.6), color)
    add_text(s, cx, cy, Inches(1.2), Inches(0.6),
             label, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# F1 & AUC
add_shape(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.8), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(7.1), Inches(4.6), Inches(5.2), Inches(0.4),
         "F1-Score & AUC-ROC", font_size=15, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(5.1), Inches(5.2), Inches(2.0), [
    "F1 = 2 × (Precision × Recall) / (Precision + Recall)",
    "  → 정밀도와 재현율의 조화 평균",
    "  → 클래스 불균형 시 Accuracy보다 유용",
    "",
    "AUC-ROC: ROC 곡선 아래 면적 (0~1)",
    "  → 1.0: 완벽한 분류기, 0.5: 랜덤 수준",
    "  → 임계값에 독립적인 종합 성능 지표",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(2))

# ============================================================
# 슬라이드 21: 실습 소개
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "실습 소개 (3개 구현 소스)")

labs = [
    ("실습 1", "01_bias_variance_tradeoff.py", "편향-분산 트레이드오프 시각화", [
        "다항 회귀로 모델 복잡도 변화 실험",
        "200개 데이터셋 반복으로 편향/분산 계산",
        "과소적합 ↔ 과적합 영역 직관적 확인",
        "세 가지 시각화: 적합, 트레이드오프 곡선, 변동성",
    ], ACCENT_BLUE),
    ("실습 2", "02_no_free_lunch_demo.py", "NFL 정리 데모", [
        "4가지 데이터 구조: 선형, 원형, XOR, 반달",
        "5가지 알고리즘: LR, k-NN, SVM, DT, RF",
        "Stratified 5-Fold로 공정 비교",
        "어떤 단일 알고리즘도 모든 곳에서 최상 아님",
    ], ACCENT_GREEN),
    ("실습 3", "03_cross_validation_demo.py", "교차검증 비교", [
        "Hold-Out, K-Fold, Stratified K-Fold, LOO 비교",
        "50회 반복으로 성능 추정 분산 측정",
        "K값에 따른 편향-분산 트레이드오프",
        "Iris 데이터셋 실전 예제",
    ], ACCENT_ORANGE),
]
for i, (title, file, desc, items, color) in enumerate(labs):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(4.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(3.9), Inches(0.5),
             title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.9), Inches(3.5), Inches(0.3),
             file, font_size=11, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(0.3),
             desc, font_size=13, color=WHITE)
    add_bullet_list(s, x + Inches(0.2), Inches(3.6), Inches(3.5), Inches(3.0),
                    items, font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# 슬라이드 22: 핵심 요약
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "핵심 요약 (Key Takeaways)")

summaries = [
    ("ML 정의", "데이터로부터 패턴을 자동 학습\nMitchell: T-P-E 프레임워크", ACCENT_BLUE),
    ("학습 유형", "지도 (분류/회귀) · 비지도 · 강화\nAI ⊃ ML ⊃ DL", ACCENT_CYAN),
    ("편향-분산", "EPE = Bias² + Var + σ²\n모델 복잡도의 최적 균형", ACCENT_GREEN),
    ("NFL 정리", "보편적 최상 알고리즘 없음\n문제에 맞는 선택이 핵심", ACCENT_ORANGE),
    ("교차검증", "Stratified K-Fold (K=5~10) 권장\n데이터 크기에 따라 선택", ACCENT_RED),
    ("ML 파이프라인", "문제정의→수집→전처리→EDA\n→모델링→평가→튜닝→배포", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(summaries):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.2) + Inches(2.5) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.1), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
             title, font_size=17, color=color, bold=True)
    add_accent_line(s, x + Inches(0.2), y + Inches(0.6), Inches(2.0), color)
    add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.2),
             desc, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 23: 핵심 수식 정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "핵심 수식 정리")

formulas = [
    ("편향-분산 분해", "EPE = Bias² + Variance + σ²", "총 오차의 세 가지 원천"),
    ("편향", "Bias(x) = f(x) - E[f̂(x)]", "모델 가정에 의한 체계적 오차"),
    ("분산", "Var(x) = E[(f̂(x) - E[f̂(x)])²]", "데이터 변동에 의한 예측 불안정"),
    ("NFL 정리", "Σf P(d|f,m,a₁) = Σf P(d|f,m,a₂)", "모든 알고리즘의 평균 성능 동일"),
    ("구조적 위험 최소화", "R(f) ≤ R_emp(f) + Φ(VCdim, n)", "일반화 오류 상한"),
    ("최적 예측기 (회귀)", "f*(x) = E[Y|X=x]", "MSE를 최소화하는 예측"),
]
for i, (name, formula, meaning) in enumerate(formulas):
    y = Inches(2.2) + Inches(0.85) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    add_shape(s, Inches(0.6), y, Inches(12.1), Inches(0.75), bg, radius=True)
    add_text(s, Inches(0.9), y + Inches(0.05), Inches(2.5), Inches(0.65),
             name, font_size=14, color=ACCENT_CYAN, bold=True)
    add_text(s, Inches(3.5), y + Inches(0.05), Inches(4.5), Inches(0.65),
             formula, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.2), y + Inches(0.05), Inches(4.3), Inches(0.65),
             meaning, font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 24: 복습 질문
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "복습 질문 (Review Questions)")

questions = [
    "Q1. Mitchell의 T-P-E 정의를 스팸 필터에 적용하여 설명하시오.",
    "Q2. 지도/비지도/강화학습을 데이터, 목표, 피드백 관점에서 비교하시오.",
    "Q3. 편향-분산 분해 공식의 각 항을 설명하고 모델 복잡도와의 관계를 서술하시오.",
    "Q4. NFL 정리의 핵심과 실무적 시사점 3가지를 서술하시오.",
    "Q5. Hold-Out, K-Fold, Stratified K-Fold, LOO의 장단점을 비교하시오.",
    "Q6. 과적합과 과소적합의 원인과 해결방법을 각 3가지 이상 서술하시오.",
    "Q7. VC 차원이란 무엇이며, 2차원 선형 분류기의 VC 차원이 3인 이유를 설명하시오.",
    "Q8. ML 파이프라인 8단계와 Data Leakage 방지 원칙을 설명하시오.",
    "Q9. ESL의 최소제곱법 vs k-NN 비교를 통한 편향-분산 트레이드오프를 설명하시오.",
    "Q10. Domingos의 \"More Data Beats a Cleverer Algorithm\"과\n      \"Feature Engineering Is the Key\"를 실무 예시와 함께 설명하시오.",
]
for i, q in enumerate(questions):
    y = Inches(2.0) + Inches(0.52) * i
    color = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED,
             ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE][i]
    add_shape(s, Inches(0.8), y + Inches(0.05), Inches(0.08), Inches(0.3), color)
    add_text(s, Inches(1.1), y, Inches(11.5), Inches(0.5),
             q, font_size=13, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 25: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.5), prs.slide_width, Inches(1.0),
         "Thank You", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), prs.slide_width, Inches(0.6),
         "1장: 머신러닝 개요", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(4.2), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(4.5), prs.slide_width, Inches(0.5),
         "다음 장: 2장 - 파이썬 기초", font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ── 저장 ──
output_path = os.path.join(os.path.dirname(__file__), "1장_머신러닝_개요_강의PPT.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
