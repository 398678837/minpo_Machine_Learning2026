"""1장 머신러닝 개요 - 확장 강의 PPT 생성 스크립트 (상세 버전)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = '맑은 고딕'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    """테이블을 카드 형태로 그리는 헬퍼"""
    total_w = sum(col_widths)
    # 헤더
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    # 행
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    """섹션 구분 슬라이드"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), accent)
    add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), accent)
    add_text(s, Inches(0), Inches(2.0), prs.slide_width, Inches(0.5),
             f"SECTION {section_num}", font_size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, Inches(5.5), Inches(2.7), Inches(2.3), accent)
    add_text(s, Inches(0), Inches(3.0), prs.slide_width, Inches(1.0),
             title, font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.5),
             subtitle, font_size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    return s


# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "기계학습 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "1장: 머신러닝 개요", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "Introduction to Machine Learning", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.3), Inches(11), Inches(0.8),
         "[ 확장 상세 버전 ]", font_size=20, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
         "핵심 키워드: ML 정의 · 학습 유형 · 편향-분산 트레이드오프 · NFL 정리 · 교차검증 · ML 파이프라인",
         font_size=14, color=DARK_GRAY)
add_shape(s, Inches(0), Inches(6.3), prs.slide_width, Inches(1.2), RGBColor(0x0A, 0x0A, 0x1A))
add_text(s, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
         "2026년 1학기", font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 2: 목차 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (Contents) - 1/2")
toc1 = [
    ("01", "머신러닝이란 - 정의, 전통 프로그래밍 vs ML, 사례", ACCENT_BLUE),
    ("02", "머신러닝의 역사 - 마일스톤, AI 겨울과 부활", ACCENT_CYAN),
    ("03", "학습 유형 - 지도학습, 비지도학습, 강화학습 상세", ACCENT_GREEN),
    ("04", "ML vs DL vs AI - 계층 구조, 비교, 선택 가이드", ACCENT_PURPLE),
    ("05", "편향-분산 트레이드오프 - 개념, 수식 유도, 모델별 특성", ACCENT_ORANGE),
    ("06", "No Free Lunch 정리 - 수학적 정의, 가정, 시사점", ACCENT_RED),
    ("07", "모델 평가 - Hold-Out, K-Fold, Stratified, LOO 비교", ACCENT_BLUE),
]
for i, (num, title, color) in enumerate(toc1):
    y = Inches(2.0) + Inches(0.7) * i
    add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.5), color, radius=True)
    add_text(s, Inches(1.2), y + Inches(0.05), Inches(0.55), Inches(0.4), num,
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), y + Inches(0.05), Inches(10), Inches(0.4), title,
             font_size=16, color=WHITE)

# ============================================================
# 슬라이드 3: 목차 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (Contents) - 2/2")
toc2 = [
    ("08", "ML 파이프라인과 데이터 품질 - 8단계, 피처 엔지니어링, 데이터 누출", ACCENT_CYAN),
    ("09", "핵심 논문 리뷰 5편 - Domingos, Mitchell, Wolpert, Vapnik, ESL", ACCENT_GREEN),
    ("10", "실습 1: 편향-분산 트레이드오프 구현 (코드 상세 해설)", ACCENT_BLUE),
    ("11", "실습 2: NFL 정리 데모 (코드 상세 해설)", ACCENT_ORANGE),
    ("12", "실습 3: 교차검증 비교 (코드 상세 해설)", ACCENT_RED),
    ("13", "응용사례 심화 - 의료, 금융, 제조, NLP", ACCENT_PURPLE),
    ("14", "핵심 요약, 수식 정리, 복습 질문 10개", ACCENT_CYAN),
]
for i, (num, title, color) in enumerate(toc2):
    y = Inches(2.0) + Inches(0.7) * i
    add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.5), color, radius=True)
    add_text(s, Inches(1.2), y + Inches(0.05), Inches(0.55), Inches(0.4), num,
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), y + Inches(0.05), Inches(10), Inches(0.4), title,
             font_size=16, color=WHITE)

# ============================================================
# 슬라이드 4: 학습 목표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "학습 목표 (Learning Objectives)")
objectives = [
    "머신러닝의 정의(Mitchell의 T-P-E)와 전통적 프로그래밍과의 차이를 설명할 수 있다",
    "지도학습, 비지도학습, 강화학습의 차이를 데이터/목표/피드백 관점에서 비교할 수 있다",
    "편향-분산 분해 공식을 유도하고, 각 항의 의미를 설명할 수 있다",
    "No Free Lunch 정리의 수학적 의미와 실무적 시사점을 서술할 수 있다",
    "4가지 교차검증 전략(Hold-Out, K-Fold, Stratified, LOO)을 비교하고 상황별 선택 기준을 제시할 수 있다",
    "ML 파이프라인 8단계를 이해하고, 데이터 누출(Data Leakage) 방지 원칙을 설명할 수 있다",
    "5편의 핵심 논문(Domingos, Mitchell, Wolpert, Vapnik, ESL)의 주요 기여를 요약할 수 있다",
    "3개 실습 코드(편향-분산, NFL, 교차검증)를 이해하고 실행할 수 있다",
]
for i, obj in enumerate(objectives):
    y = Inches(2.0) + Inches(0.62) * i
    colors = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE,
              ACCENT_RED, ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN]
    add_shape(s, Inches(0.8), y + Inches(0.08), Inches(0.35), Inches(0.35), colors[i], radius=True)
    add_text(s, Inches(0.8), y + Inches(0.08), Inches(0.35), Inches(0.35),
             str(i+1), font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.4), y, Inches(11), Inches(0.55),
             obj, font_size=15, color=LIGHT_GRAY)

# ============================================================
# SECTION 01: 머신러닝이란
# ============================================================
section_divider("머신러닝이란?", "Machine Learning Definition & Fundamentals", "01", ACCENT_BLUE)

# ============================================================
# 슬라이드 6: 머신러닝 정의 - Samuel & Mitchell
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "머신러닝의 정의", "Two Foundational Definitions")

add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Arthur Samuel (1959)", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.2),
         "\"Machine Learning is a field of study that gives\ncomputers the ability to learn without being\nexplicitly programmed.\"",
         font_size=14, color=LIGHT_GRAY)
add_text(s, Inches(0.9), Inches(3.7), Inches(5.2), Inches(0.4),
         "→ 명시적 프로그래밍 없이 컴퓨터가 학습하는 능력을 부여하는 학문", font_size=12, color=DARK_GRAY)

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "Tom Mitchell (1997) - 형식적 정의", font_size=16, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.2),
         "\"A computer program is said to learn from\nexperience E with respect to some class of tasks T\nand performance measure P, if its performance at\ntasks in T, as measured by P, improves with experience E.\"",
         font_size=13, color=LIGHT_GRAY)

# T-P-E 상세
labels = [
    ("T (Task)", "수행하고자 하는 작업", "이메일 스팸 여부 분류", "무엇을 할 것인가?", ACCENT_BLUE),
    ("P (Performance)", "성능 측정 지표", "정확하게 분류된 이메일 비율", "얼마나 잘 하는가?", ACCENT_GREEN),
    ("E (Experience)", "학습에 사용되는 데이터", "사용자가 레이블링한 이메일", "무엇으로 학습하는가?", ACCENT_ORANGE),
]
for i, (name, meaning, example, question, color) in enumerate(labels):
    x = Inches(0.6) + Inches(4.1) * i
    add_shape(s, x, Inches(4.6), Inches(3.8), Inches(2.7), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(4.7), Inches(3.4), Inches(0.4),
             name, font_size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, x + Inches(0.3), Inches(5.15), Inches(3.2), color)
    add_text(s, x + Inches(0.2), Inches(5.3), Inches(3.4), Inches(0.3),
             question, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(5.7), Inches(3.4), Inches(0.3),
             meaning, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(6.1), Inches(3.4), Inches(0.8),
             f"예: {example}", font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 7: 전통적 프로그래밍 vs 머신러닝
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "전통적 프로그래밍 vs 머신러닝", "패러다임의 전환")

add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.8), Inches(2.3), Inches(5.4), Inches(0.5),
         "전통적 프로그래밍", font_size=22, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(1.0), Inches(3.0), Inches(5.0), Inches(3.8), [
    "입력: 데이터 + 규칙 (Rules)",
    "출력: 결과 (Answers)",
    "사람이 규칙을 직접 코딩",
    "",
    "예시: if-else 기반 스팸 필터",
    '  if "당첨" in email: return "스팸"',
    '  if "무료" in email: return "스팸"',
    "",
    "한계: 새로운 패턴에 대응 어려움",
    "      규칙이 복잡해질수록 유지보수 어려움",
], font_size=14, color=LIGHT_GRAY)

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.0), Inches(2.3), Inches(5.4), Inches(0.5),
         "머신러닝", font_size=22, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(7.2), Inches(3.0), Inches(5.0), Inches(3.8), [
    "입력: 데이터 + 정답 (Labels)",
    "출력: 규칙/모델 (Model)",
    "데이터로부터 규칙을 자동 학습",
    "",
    "예시: 수만 건의 이메일로 학습",
    "  → 패턴을 자동으로 발견",
    "  → 새로운 스팸 유형도 감지",
    "",
    "장점: 환경 변화에 적응 가능",
    "      사람이 인지 못하는 패턴도 발견",
], font_size=14, color=LIGHT_GRAY)

add_shape(s, Inches(6.0), Inches(4.0), Inches(1.2), Inches(0.8), ACCENT_BLUE, radius=True)
add_text(s, Inches(6.0), Inches(4.05), Inches(1.2), Inches(0.7),
         "VS", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 8: 일상 속 머신러닝 사례 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "일상 속 머신러닝 사례 (1/2)")

# 추천 시스템
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "1. 추천 시스템 (Recommendation System)", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "넷플릭스, 유튜브, 스포티파이: 시청/청취 이력 분석 → 개인화 추천",
    "협업 필터링 (Collaborative Filtering) + 콘텐츠 기반 필터링",
    "넷플릭스: 추천으로 시청 시간의 약 80% 유도",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 자율주행
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "2. 자율주행 (Autonomous Driving)", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "테슬라, 웨이모: 카메라, 라이다, 레이더 센서 데이터 실시간 처리",
    "CNN 기반 객체 탐지 + 강화학습 기반 경로 계획",
    "수백만 km 주행 데이터로 학습, 보행자/차량/신호등 인식",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 의료 진단
add_shape(s, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(5.2), Inches(0.4),
         "3. 의료 진단 (Medical Diagnosis)", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(5.4), Inches(5.2), Inches(1.5), [
    "X-ray, CT, MRI에서 종양/폐렴/골절 자동 탐지",
    "AlphaFold: 단백질 구조 예측 문제 해결 → 생물학 혁신",
    "유방암 진단: 일부 ML 모델이 전문의보다 높은 정확도",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 플레이스홀더
add_shape(s, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(7.1), Inches(4.9), Inches(5.2), Inches(0.4),
         "핵심 포인트", font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(5.4), Inches(5.2), Inches(1.5), [
    "ML은 이미 일상 곳곳에 적용되어 있음",
    "다양한 데이터 유형: 텍스트, 이미지, 음성, 센서",
    "핵심: 데이터에서 패턴을 자동으로 학습하여 예측/결정",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# 슬라이드 9: 일상 속 머신러닝 사례 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "일상 속 머신러닝 사례 (2/2)")

add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "4. 음성 인식 (Speech Recognition)", font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "Siri, Google Assistant, Alexa: 음성→텍스트(STT) + 의도파악(NLU)",
    "딥러닝 기반 RNN, Transformer 모델 (수만 시간 음성 학습)",
    "OpenAI Whisper: 다국어 음성 인식에서 뛰어난 성능",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "5. 사기 탐지 (Fraud Detection)", font_size=16, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "신용카드 부정 사용, 보험 사기, 자금 세탁 실시간 탐지",
    "정상 패턴 학습 → 이상 거래(anomaly) 감지",
    "비지도학습(이상치 탐지) + 지도학습(분류) 병행, 초당 수천 건 분석",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# ML이 필요한 이유
add_shape(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.4),
         "머신러닝이 필요한 이유 (Mitchell, 1997)", font_size=18, color=ACCENT_ORANGE, bold=True)
reasons = [
    ("복잡한 규칙의 자동 발견", "사람이 규칙을 일일이 정의하기 어려운 문제 (이미지 인식, 음성 인식)", ACCENT_BLUE),
    ("환경 변화에 대한 적응", "새로운 데이터가 들어오면 모델을 재학습하여 변화에 대응", ACCENT_GREEN),
    ("대량 데이터에서 패턴 발견", "빅데이터 시대에 수작업 분석은 불가능 → 자동화된 분석", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(reasons):
    x = Inches(0.8) + Inches(4.0) * i
    add_shape(s, x, Inches(5.5), Inches(3.7), Inches(1.5), RGBColor(0x2D, 0x2D, 0x45), color, radius=True)
    add_text(s, x + Inches(0.15), Inches(5.55), Inches(3.4), Inches(0.4),
             f"{i+1}. {title}", font_size=13, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(5.95), Inches(3.4), Inches(0.8),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# SECTION 02: 머신러닝의 역사
# ============================================================
section_divider("머신러닝의 역사", "History of Machine Learning", "02", ACCENT_CYAN)

# ============================================================
# 슬라이드 11: 역사 타임라인 (1/2) - 1943~1997
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "머신러닝 주요 마일스톤 (1943~1997)")

milestones1 = [
    ("1943", "McCulloch & Pitts", "인공 뉴런 모델", "최초의 수학적 뉴런 모델\n신경망의 이론적 기초", ACCENT_BLUE),
    ("1950", "Alan Turing", "Computing Machinery\nand Intelligence", "튜링 테스트 제안\n\"기계가 생각할 수 있는가?\"", ACCENT_CYAN),
    ("1957", "Rosenblatt", "퍼셉트론", "최초의 학습 가능한\n신경망 모델", ACCENT_GREEN),
    ("1959", "Arthur Samuel", "체커 프로그램", "\"Machine Learning\"\n용어 최초 사용", ACCENT_GREEN),
    ("1969", "Minsky & Papert", "Perceptrons", "퍼셉트론 한계(XOR) 증명\n→ 1차 AI 겨울 촉발", ACCENT_RED),
    ("1986", "Rumelhart et al.", "역전파 알고리즘", "다층 신경망 학습 가능\n→ 신경망 부활", ACCENT_ORANGE),
    ("1995", "Vapnik", "SVM", "통계적 학습 이론 기반\n강력한 분류기", ACCENT_PURPLE),
    ("1997", "Mitchell / Wolpert", "ML 교재 / NFL 정리", "T-P-E 정의 정립\n보편적 최적 알고리즘 부정", ACCENT_BLUE),
]
for i, (year, author, event, desc, color) in enumerate(milestones1):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(3.15) * col
    y = Inches(2.2) + Inches(2.7) * row
    add_shape(s, x, y, Inches(2.9), Inches(2.4), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.7), Inches(0.35), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.7), Inches(0.35),
             year, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.9), y + Inches(0.1), Inches(1.9), Inches(0.35),
             author, font_size=11, color=color, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(0.35),
             event, font_size=13, color=WHITE, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(1.2),
             desc, font_size=11, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 12: 역사 타임라인 (2/2) - 2006~현재
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "머신러닝 주요 마일스톤 (2006~현재)")

milestones2 = [
    ("2006", "Hinton", "Deep Belief Networks", "딥러닝 부활 신호탄\n\"딥러닝\" 용어 대중화", ACCENT_BLUE),
    ("2012", "Krizhevsky et al.", "AlexNet (ImageNet)", "CNN 폭발적 성능 향상\n딥러닝 붐 시작", ACCENT_CYAN),
    ("2014", "Goodfellow", "GAN", "생성적 적대 신경망\n생성 모델의 혁신", ACCENT_GREEN),
    ("2016", "DeepMind", "AlphaGo vs 이세돌", "강화학습의 가능성을\n전 세계에 입증", ACCENT_ORANGE),
    ("2017", "Vaswani et al.", "Transformer", "Attention Is All You Need\nNLP 혁명의 시작", ACCENT_RED),
    ("2018~", "Google, OpenAI", "BERT, GPT 시리즈", "대규모 언어모델(LLM)\n시대 개막", ACCENT_PURPLE),
    ("2022~", "OpenAI", "ChatGPT, GPT-4", "생성형 AI의 대중화\nAGI 논의 활발", ACCENT_BLUE),
]
for i, (year, author, event, desc, color) in enumerate(milestones2):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(3.15) * col
    y = Inches(2.2) + Inches(2.7) * row
    add_shape(s, x, y, Inches(2.9), Inches(2.4), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.7), Inches(0.35), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.7), Inches(0.35),
             year, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.9), y + Inches(0.1), Inches(1.9), Inches(0.35),
             author, font_size=11, color=color, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(0.35),
             event, font_size=13, color=WHITE, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(1.2),
             desc, font_size=11, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 13: AI 겨울과 부활
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "AI 겨울과 부활", "AI Winters and Revivals")

phases = [
    ("1차 AI 겨울\n(1970년대)", "퍼셉트론의 한계\n(XOR 문제) 증명\n→ 신경망 연구 침체\nMinsky & Papert(1969)",
     ACCENT_RED, Inches(0.5)),
    ("부활\n(1980년대)", "역전파 알고리즘 발견\n→ 다층 신경망 학습 가능\nRumelhart, Hinton,\nWilliams (1986)",
     ACCENT_GREEN, Inches(3.3)),
    ("2차 AI 겨울\n(1990년대 초)", "신경망 학습 어려움\n(Vanishing Gradient)\n→ SVM, 앙상블 주류\nVapnik의 SVM(1995)",
     ACCENT_ORANGE, Inches(6.1)),
    ("딥러닝 부활\n(2006~2012)", "GPU 컴퓨팅 + 빅데이터\n+ ReLU 활성화 함수\n→ 딥러닝 폭발적 성장\nAlexNet(2012)",
     ACCENT_CYAN, Inches(8.9)),
]
# 타임라인 연결선
add_shape(s, Inches(0.8), Inches(3.7), Inches(11.7), Pt(4), ACCENT_BLUE)

for title, desc, color, x in phases:
    add_shape(s, x, Inches(2.2), Inches(2.5), Inches(1.3), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(2.3), Inches(2.2), Inches(1.1),
             title, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    # 도트
    add_shape(s, x + Inches(1.1), Inches(3.55), Inches(0.3), Inches(0.3), color, radius=True)
    # 설명
    add_shape(s, x, Inches(4.1), Inches(2.5), Inches(2.5), CARD_BG, radius=True)
    add_text(s, x + Inches(0.15), Inches(4.2), Inches(2.2), Inches(2.2),
             desc, font_size=12, color=LIGHT_GRAY)

add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
         "교훈: ML의 발전은 직선적이지 않다. 이론적 한계 발견 → 침체 → 새로운 돌파구 → 부활의 반복",
         font_size=14, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 03: 학습 유형
# ============================================================
section_divider("학습 유형", "Supervised / Unsupervised / Reinforcement Learning", "03", ACCENT_GREEN)

# ============================================================
# 슬라이드 15: 지도학습 상세
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "지도학습 (Supervised Learning)", "입력 X와 정답 y가 함께 주어진 학습")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.35),
         "정의: 입력 데이터 X와 정답 레이블 y가 함께 주어진 상태에서 매핑 함수 f: X → y를 학습하는 방법",
         font_size=16, color=WHITE, bold=True)
add_text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.35),
         "핵심: 예측값과 실제값의 차이(오차)를 줄이는 방향으로 반복 학습 → 회귀(연속값) + 분류(이산값)",
         font_size=14, color=LIGHT_GRAY)

# 회귀
add_card(s, Inches(0.6), Inches(3.5), Inches(5.8), Inches(1.8),
         "회귀 (Regression) - 연속값 예측", [
             "출력값 y가 연속적인 수치: 주택 가격, 기온, 매출 예측",
             "대표 알고리즘: 선형회귀, Ridge, Lasso, RF, GBM",
             "평가 지표: MSE, RMSE, MAE, R²",
         ], ACCENT_GREEN, ACCENT_GREEN)
# 분류
add_card(s, Inches(6.8), Inches(3.5), Inches(5.8), Inches(1.8),
         "분류 (Classification) - 범주 예측", [
             "출력값 y가 이산적 카테고리: 스팸/정상, 악성/양성",
             "이진분류 / 다중분류 / 다중레이블분류",
             "평가 지표: Accuracy, Precision, Recall, F1, AUC",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 산업 사례 테이블
add_text(s, Inches(0.6), Inches(5.5), Inches(5), Inches(0.4),
         "지도학습 산업 적용 사례", font_size=15, color=WHITE, bold=True)
add_table_slide(s,
    ["산업", "문제", "입력(X)", "출력(y)", "알고리즘"],
    [
        ["금융", "신용 점수 예측", "소득, 부채, 거래이력", "신용 등급(A~F)", "XGBoost"],
        ["제조", "제품 불량 예측", "센서, 온도, 압력", "양품/불량", "Random Forest"],
        ["마케팅", "고객 이탈 예측", "사용 패턴, 결제", "이탈/유지", "Gradient Boosting"],
    ],
    Inches(0.6), Inches(5.9), [2.0, 1.8, 2.5, 2.5, 2.4], font_size=12, header_font_size=13, row_height=0.45)


# ============================================================
# PLACEHOLDER for remaining slides - will be appended
# ============================================================

# ============================================================
# 슬라이드 16: 비지도학습 상세
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "비지도학습 (Unsupervised Learning)", "정답 레이블 없이 데이터 구조 발견")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.7), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5),
         "정의: 입력 데이터 X만 주어지고 정답 레이블 y가 없는 상태에서 데이터 자체의 구조나 패턴을 발견하는 학습 방법",
         font_size=16, color=WHITE, bold=True)

tasks = [
    ("군집화 (Clustering)", "유사한 데이터를 그룹화", "K-Means, DBSCAN\n계층적 군집화",
     "고객 세분화\n유전자 군집화", ACCENT_BLUE),
    ("차원 축소\n(Dimensionality Reduction)", "고차원 → 저차원 변환", "PCA, t-SNE\nAutoencoder",
     "시각화\n노이즈 제거", ACCENT_CYAN),
    ("이상치 탐지\n(Anomaly Detection)", "정상 패턴에서 벗어나는\n데이터 감지", "Isolation Forest\nOne-Class SVM",
     "사기 탐지\n네트워크 이상", ACCENT_ORANGE),
]
for i, (title, desc, algos, cases, color) in enumerate(tasks):
    x = Inches(0.5) + Inches(4.15) * i
    add_shape(s, x, Inches(3.2), Inches(3.9), Inches(4.0), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(3.3), Inches(3.5), Inches(0.6),
             title, font_size=15, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.9), Inches(3.5), Inches(0.5),
             desc, font_size=13, color=WHITE)
    add_text(s, x + Inches(0.2), Inches(4.5), Inches(3.5), Inches(0.3),
             "알고리즘:", font_size=12, color=ACCENT_CYAN, bold=True)
    add_text(s, x + Inches(0.2), Inches(4.8), Inches(3.5), Inches(0.8),
             algos, font_size=12, color=LIGHT_GRAY)
    add_text(s, x + Inches(0.2), Inches(5.6), Inches(3.5), Inches(0.3),
             "활용 사례:", font_size=12, color=ACCENT_GREEN, bold=True)
    add_text(s, x + Inches(0.2), Inches(5.9), Inches(3.5), Inches(0.8),
             cases, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 17: 강화학습 상세
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "강화학습 (Reinforcement Learning)", "에이전트가 환경과 상호작용하며 보상 최대화")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.7), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5),
         "정의: 에이전트(Agent)가 환경(Environment)과 상호작용하면서, 보상(Reward)을 최대화하는 행동 전략(Policy)을 학습",
         font_size=16, color=WHITE, bold=True)

# 구성 요소
components = [
    ("에이전트 (Agent)", "학습하고 행동하는 주체", ACCENT_BLUE),
    ("환경 (Environment)", "에이전트가 상호작용하는 세계", ACCENT_CYAN),
    ("상태 (State)", "현재 환경의 상황", ACCENT_GREEN),
    ("행동 (Action)", "에이전트가 취할 수 있는 동작", ACCENT_ORANGE),
    ("보상 (Reward)", "행동에 대한 피드백 신호", ACCENT_RED),
]
for i, (name, desc, color) in enumerate(components):
    x = Inches(0.4) + Inches(2.55) * i
    add_shape(s, x, Inches(3.2), Inches(2.35), Inches(1.2), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.1), Inches(3.3), Inches(2.15), Inches(0.4),
             name, font_size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(3.7), Inches(2.15), Inches(0.5),
             desc, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 상호작용 다이어그램
add_shape(s, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(4.8), Inches(5.2), Inches(0.4),
         "강화학습 상호작용 루프", font_size=15, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(0.9), Inches(5.3), Inches(5.2), Inches(1.5),
         "에이전트 --[행동(Action)]--> 환경\n환경 --[상태(State), 보상(Reward)]--> 에이전트\n\n이 과정을 반복하며 최적 정책(Policy)을 학습\n목표: 누적 보상의 기대값 최대화",
         font_size=14, color=LIGHT_GRAY)

# 산업 사례
add_shape(s, Inches(6.8), Inches(4.7), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(4.8), Inches(5.2), Inches(0.4),
         "강화학습 산업 사례", font_size=15, color=ACCENT_CYAN, bold=True)
add_table_slide(s,
    ["분야", "보상 설계", "알고리즘"],
    [
        ["게임 AI (AlphaGo)", "승리+1, 패배-1", "MCTS+정책신경망"],
        ["로봇 팔 제어", "목표물 잡기 성공/실패", "PPO, SAC"],
        ["데이터센터 냉각", "에너지 절감량", "DQN (DeepMind)"],
    ],
    Inches(6.8), Inches(5.3), [2.1, 1.9, 1.8], font_size=11, header_font_size=12, row_height=0.4)

# ============================================================
# 슬라이드 18: 세 가지 학습 유형 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "세 가지 학습 유형 비교 요약")

add_table_slide(s,
    ["구분", "지도학습", "비지도학습", "강화학습"],
    [
        ["데이터", "입력 + 정답", "입력만", "상태 + 보상"],
        ["목표", "예측 / 분류", "구조 발견", "보상 최대화"],
        ["피드백", "정답 레이블", "없음", "보상 신호"],
        ["대표 문제", "분류, 회귀", "군집화, 차원축소", "게임, 로봇 제어"],
        ["대표 알고리즘", "RF, XGB, SVM, LR", "K-Means, PCA", "Q-Learning, PPO"],
        ["난이도", "상대적 쉬움", "중간", "상대적 어려움"],
    ],
    Inches(0.8), Inches(2.2), [2.5, 2.8, 2.8, 2.8], row_height=0.55)

add_shape(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.3), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(1.1), Inches(5.9), Inches(11.1), Inches(0.4),
         "실무 선택 가이드", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(1.1), Inches(6.3), Inches(11.1), Inches(0.7), [
    "정답 레이블이 있으면 → 지도학습 (분류/회귀)  |  정답 없이 패턴 발견 → 비지도학습  |  순차적 의사결정 → 강화학습",
    "실무에서 가장 많이 사용: 지도학습 (데이터만 충분하면 높은 성능). 비지도학습은 탐색/전처리 단계에서 자주 활용",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# SECTION 04: ML vs DL vs AI
# ============================================================
section_divider("ML vs DL vs AI", "Hierarchy and Comparison", "04", ACCENT_PURPLE)

# ============================================================
# 슬라이드 20: ML vs DL vs AI 계층 구조
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "AI ⊃ ML ⊃ DL 계층 구조")

# 동심원 계층
circles = [
    (Inches(1.0), Inches(2.0), Inches(6.0), Inches(5.2), RGBColor(0x1E, 0x3A, 0x5F), ACCENT_BLUE),
    (Inches(1.8), Inches(2.6), Inches(4.4), Inches(4.0), RGBColor(0x1E, 0x4D, 0x3A), ACCENT_GREEN),
    (Inches(2.6), Inches(3.2), Inches(2.8), Inches(2.6), RGBColor(0x3A, 0x1E, 0x5F), ACCENT_PURPLE),
]
for x, y, w, h, fill, border in circles:
    add_shape(s, x, y, w, h, fill, border, radius=True)

add_text(s, Inches(1.3), Inches(2.2), Inches(3), Inches(0.4),
         "인공지능 (AI)", font_size=18, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(2.1), Inches(2.8), Inches(3), Inches(0.4),
         "머신러닝 (ML)", font_size=16, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(2.9), Inches(4.2), Inches(2.2), Inches(0.4),
         "딥러닝 (DL)", font_size=15, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)

# 설명 카드
descs = [
    ("AI (Artificial Intelligence)", "인간의 지능을 모방하는 모든 기술의 총칭\n규칙 기반, 전문가 시스템, 탐색 등 포함\n가장 넓은 개념", ACCENT_BLUE),
    ("ML (Machine Learning)", "AI의 하위 분야, 데이터 기반 학습에 초점\n명시적 규칙 대신 데이터에서 패턴 자동 학습\nSVM, RF, XGBoost 등", ACCENT_GREEN),
    ("DL (Deep Learning)", "ML의 하위 분야, 심층 신경망 활용\n여러 층의 비선형 변환으로 복잡한 표현 학습\nCNN, RNN, Transformer", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(descs):
    y = Inches(2.0) + Inches(1.8) * i
    add_shape(s, Inches(7.5), y, Inches(5.2), Inches(1.5), CARD_BG, color, radius=True)
    add_text(s, Inches(7.7), y + Inches(0.1), Inches(4.8), Inches(0.35),
             title, font_size=14, color=color, bold=True)
    add_text(s, Inches(7.7), y + Inches(0.45), Inches(4.8), Inches(0.9),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 21: ML vs DL 상세 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "전통적 ML vs 딥러닝 상세 비교")

add_table_slide(s,
    ["특성", "전통적 머신러닝", "딥러닝"],
    [
        ["피처 엔지니어링", "사람이 직접 설계 (수동)", "모델이 자동으로 학습"],
        ["데이터 요구량", "적은 데이터도 가능", "대량의 데이터 필요"],
        ["계산 자원", "CPU로 충분한 경우 많음", "GPU/TPU 필수적"],
        ["해석 가능성", "상대적으로 높음", "블랙박스 경향"],
        ["성능 상한", "데이터 증가시 성능 포화", "데이터 증가에 따라 지속 향상"],
        ["적합한 데이터", "정형 데이터 (테이블)", "비정형 (이미지, 텍스트, 음성)"],
        ["대표 알고리즘", "SVM, RF, XGBoost", "CNN, RNN, Transformer"],
    ],
    Inches(0.8), Inches(2.2), [3.0, 4.3, 4.3], row_height=0.5, font_size=14)

# 선택 가이드
add_shape(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.35),
         "언제 무엇을 사용할 것인가? (Domingos, 2012)", font_size=14, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7), [
    "정형 데이터 + 소~중규모 → XGBoost, LightGBM, RF  |  비정형 (이미지, 텍스트) → CNN, Transformer",
    "해석 가능성 중요 (의료, 금융) → 전통 ML or XAI  |  데이터 매우 적을 때 → 전통 ML + 도메인 지식 피처",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# SECTION 05: 편향-분산 트레이드오프
# ============================================================
section_divider("편향-분산 트레이드오프", "Bias-Variance Tradeoff", "05", ACCENT_ORANGE)

# ============================================================
# 슬라이드 23: 편향과 분산 개념
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "편향(Bias)과 분산(Variance) 개념")

# 편향
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.5),
         "편향 (Bias)", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.9), Inches(5.2), Inches(2.0), [
    "모델의 예측값과 실제 정답 사이의 체계적인 오차",
    "모델이 학습 데이터의 패턴을 얼마나 잘 포착하는지 측정",
    "",
    "높은 편향 → 모델이 너무 단순 → 과소적합 (Underfitting)",
    "낮은 편향 → 데이터 패턴을 잘 포착",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 분산
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.8), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.5),
         "분산 (Variance)", font_size=22, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.9), Inches(5.2), Inches(2.0), [
    "서로 다른 학습 데이터셋에서 예측값이 얼마나 변하는지",
    "모델의 예측 안정성을 측정",
    "",
    "높은 분산 → 학습 데이터에 민감 → 과적합 (Overfitting)",
    "낮은 분산 → 데이터가 바뀌어도 안정적 예측",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 과녁 비유
add_text(s, Inches(0.6), Inches(5.3), Inches(5), Inches(0.4),
         "직관적 비유: 과녁 맞추기 (과녁 중심 = 실제 정답)", font_size=15, color=WHITE, bold=True)
targets = [
    ("낮은 편향\n낮은 분산", "이상적 (Best)", ACCENT_GREEN),
    ("높은 편향\n낮은 분산", "과소적합", ACCENT_ORANGE),
    ("낮은 편향\n높은 분산", "과적합", ACCENT_RED),
    ("높은 편향\n높은 분산", "최악 (Worst)", DARK_GRAY),
]
for i, (label, status, color) in enumerate(targets):
    x = Inches(0.6) + Inches(3.1) * i
    add_shape(s, x, Inches(5.8), Inches(2.8), Inches(1.4), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(5.9), Inches(1.3), Inches(1.2),
             label, font_size=12, color=LIGHT_GRAY)
    add_text(s, x + Inches(1.4), Inches(6.1), Inches(1.3), Inches(0.8),
             status, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 24: 수식 유도 (1/2) - 문제 설정
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "편향-분산 분해: 수식 유도 (1/2)", "Problem Setup & Derivation")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "1단계: 데이터 생성 모델 가정", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.5),
         "y = f(x) + e,    e ~ N(0, s²)\n여기서 f(x)는 실제 함수 (ground truth), e은 줄일 수 없는 노이즈 (irreducible error)",
         font_size=15, color=LIGHT_GRAY)

add_shape(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(0.4),
         "2단계: 기대 예측 오류 (Expected Prediction Error) 정의", font_size=16, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.5),
         "학습 데이터셋 D로 학습한 모델의 예측: f_hat_D(x)\nEPE(x) = E_D[(y - f_hat_D(x))²]  =  E_D[(f(x) + e - f_hat_D(x))²]",
         font_size=15, color=LIGHT_GRAY)

add_shape(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4),
         "3단계: 평균 예측 도입", font_size=16, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.3),
         "f_bar(x) = E_D[f_hat_D(x)]  ← 여러 데이터셋에 대한 모델 예측의 기대값 (평균 예측)\n\n"
         "f_hat_D(x)를 다음과 같이 분해:\n"
         "f_hat_D(x) = f_bar(x) + (f_hat_D(x) - f_bar(x))\n"
         "              ↑ 평균 예측    ↑ 평균으로부터의 편차",
         font_size=14, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 25: 수식 유도 (2/2) - 최종 분해
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "편향-분산 분해: 수식 유도 (2/2)", "Final Decomposition")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "4단계: 전개 및 정리", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.0),
         "EPE(x) = E[(f(x) + e - f_hat(x))²]\n"
         "       = E[(f(x) - f_bar(x))² + (f_hat(x) - f_bar(x))² + e² + 교차항들]\n"
         "교차항들은 e과 f_hat이 독립이고, E[e]=0 이므로 모두 0이 됨",
         font_size=14, color=LIGHT_GRAY)

# 최종 공식 - 큰 박스
add_shape(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(1.5), RGBColor(0x15, 0x15, 0x30), ACCENT_CYAN, radius=True)
add_text(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.6),
         "EPE = Bias² + Variance + s²  (Irreducible Error)", font_size=28, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(0.5),
         "총 오차 = 모델 가정에 의한 체계적 오차 + 데이터 변동에 의한 불안정성 + 데이터 자체 노이즈",
         font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 각 항 설명 테이블
add_table_slide(s,
    ["항", "수식", "의미", "줄이는 방법"],
    [
        ["Bias²", "(f(x) - f_bar(x))²", "모델 가정에 의한 체계적 오차", "모델 복잡도 증가"],
        ["Variance", "E[(f_hat - f_bar)²]", "데이터 변동에 의한 예측 불안정", "모델 복잡도 감소, 앙상블"],
        ["s² (비가약)", "s²", "데이터 자체 노이즈", "줄일 수 없음"],
    ],
    Inches(0.6), Inches(5.8), [1.8, 2.8, 3.5, 3.0], row_height=0.5, font_size=13)

# ============================================================
# 슬라이드 26: 그래프 해석 + 모델별 특성
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "편향-분산 트레이드오프 그래프 & 모델별 특성")

# 그래프 설명 (왼쪽)
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "그래프 해석", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(4.0), [
    "X축: 모델 복잡도 (다항식 차수, 트리 깊이, 층 수)",
    "Y축: 오류 (Error)",
    "",
    "모델 복잡도 증가하면:",
    "  → 편향(Bias)은 감소 (복잡한 패턴 포착 가능)",
    "  → 분산(Variance)은 증가 (데이터에 민감해짐)",
    "",
    "최적의 모델 복잡도:",
    "  → Bias² + Variance의 합이 최소인 지점",
    "  → 이 지점에서 일반화 성능이 최대",
    "",
    "왼쪽(저복잡도) = 과소적합 영역",
    "오른쪽(고복잡도) = 과적합 영역",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 모델별 특성 (오른쪽)
add_text(s, Inches(6.8), Inches(2.2), Inches(5.5), Inches(0.4),
         "모델별 편향-분산 특성", font_size=16, color=ACCENT_CYAN, bold=True)
add_table_slide(s,
    ["모델", "편향", "분산", "특징"],
    [
        ["선형 회귀", "높음", "낮음", "단순, 안정적"],
        ["KNN (K 큼)", "높음", "낮음", "과도한 평활화"],
        ["KNN (K 작음)", "낮음", "높음", "노이즈 민감"],
        ["결정트리 (깊음)", "낮음", "높음", "과적합 경향"],
        ["랜덤 포레스트", "낮음", "적절", "앙상블→분산↓"],
        ["부스팅", "낮음", "적절~높음", "편향 감소 강점"],
    ],
    Inches(6.8), Inches(2.7), [1.7, 0.8, 1.0, 1.6], row_height=0.48, font_size=12, header_font_size=13)

add_shape(s, Inches(6.8), Inches(6.1), Inches(5.5), Inches(0.8), CARD_BG, radius=True)
add_text(s, Inches(7.0), Inches(6.15), Inches(5.1), Inches(0.65),
         "ESL (Hastie et al., 2009) Ch2:\n최소제곱법(높은 편향, 낮은 분산) vs k-NN(낮은 편향, 높은 분산) 비교",
         font_size=12, color=DARK_GRAY)

# ============================================================
# SECTION 06: NFL 정리
# ============================================================
section_divider("No Free Lunch 정리", "Wolpert & Macready (1997)", "06", ACCENT_RED)

# ============================================================
# 슬라이드 28: NFL 정리 - 내용과 수학적 정의
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "No Free Lunch 정리", "핵심 내용과 수학적 정의")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "핵심 주장 (Wolpert & Macready, 1997)", font_size=16, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.6),
         "모든 가능한 문제(목적 함수)에 대해 평균을 내면, 어떤 두 최적화(또는 학습) 알고리즘이든 동일한 성능을 보인다.\n"
         "→ \"최고의 알고리즘\"은 존재하지 않는다. 항상 문제에 맞는 알고리즘을 선택해야 한다.",
         font_size=15, color=LIGHT_GRAY)

# 수학적 정의
add_shape(s, Inches(0.6), Inches(3.8), Inches(12.1), Inches(2.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.4),
         "수학적 정의", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(4.3), Inches(5.5), Inches(1.8), [
    "X: 후보 해(candidate solution)의 집합",
    "Y: 비용 값(cost value)의 집합",
    "f: X → Y: 목적 함수",
    "a: 최적화 알고리즘",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(6.5), Inches(4.3), Inches(5.8), Inches(0.8), RGBColor(0x15, 0x15, 0x30))
add_text(s, Inches(6.5), Inches(4.35), Inches(5.8), Inches(0.7),
         "Sf P(d_m^y | f, m, a1) = Sf P(d_m^y | f, m, a2)",
         font_size=16, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.5), Inches(5.15), Inches(5.8), Inches(0.6),
         "모든 함수 f에 대해 합산하면, 모든 알고리즘 쌍\n(a1, a2)에 대해 위 등식이 성립",
         font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 직관적 의미
add_shape(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7), CARD_BG, ACCENT_ORANGE, radius=True)
add_bullet_list(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6), [
    "a1이 특정 문제에서 a2보다 우수하면, 반드시 다른 문제에서는 a2가 a1보다 우수  |  무작위 탐색도 평균적으로 정교한 알고리즘과 동일",
], font_size=13, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 29: NFL 가정과 시사점
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "NFL 정리: 핵심 가정과 실무적 시사점")

# 가정
add_text(s, Inches(0.6), Inches(2.1), Inches(5), Inches(0.4),
         "핵심 가정과 한계", font_size=16, color=WHITE, bold=True)
add_table_slide(s,
    ["가정", "의미", "현실적 해석"],
    [
        ["모든 함수에 균일 분포", "가능한 모든 f를 동등 고려", "현실 문제는 특정 구조를 가짐"],
        ["이전 평가 재방문 없음", "같은 점을 두 번 평가 안함", "대부분 알고리즘이 만족"],
        ["비용 함수가 결정적", "f가 deterministic", "확률적 문제는 별도 분석"],
    ],
    Inches(0.6), Inches(2.5), [3.0, 3.5, 4.5], row_height=0.5, font_size=13)

add_shape(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.6), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(4.25), Inches(11.5), Inches(0.5),
         "주의: NFL은 \"모든 알고리즘이 동등하다\"로 오해되지만, 현실 문제는 구조(structure)를 가지고 있어 특정 알고리즘이 유리!",
         font_size=14, color=ACCENT_ORANGE, bold=True)

# 시사점
add_text(s, Inches(0.6), Inches(5.0), Inches(5), Inches(0.4),
         "실무적 시사점 5가지", font_size=16, color=WHITE, bold=True)
implications = [
    ("1", "\"최고의 알고리즘\"은 없다", "항상 문제에 맞는 알고리즘 선택", ACCENT_BLUE),
    ("2", "사전 지식의 활용", "문제 구조(귀납적 편향)가 핵심", ACCENT_GREEN),
    ("3", "다양한 알고리즘 시도", "하나에 의존하지 말고 비교", ACCENT_ORANGE),
    ("4", "벤치마크의 한계", "특정 성능 ≠ 모든 문제 우수", ACCENT_RED),
    ("5", "AutoML의 근거", "자동 알고리즘 선택/튜닝 당위성", ACCENT_PURPLE),
]
for i, (num, title, desc, color) in enumerate(implications):
    x = Inches(0.4) + Inches(2.55) * i
    add_shape(s, x, Inches(5.5), Inches(2.35), Inches(1.8), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), Inches(5.6), Inches(0.35), Inches(0.35), color, radius=True)
    add_text(s, x + Inches(0.1), Inches(5.6), Inches(0.35), Inches(0.35),
             num, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.55), Inches(5.6), Inches(1.7), Inches(0.4),
             title, font_size=12, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(6.1), Inches(2.1), Inches(0.9),
             desc, font_size=11, color=LIGHT_GRAY)

# ============================================================
# SECTION 07: 모델 평가 - 교차검증
# ============================================================
section_divider("모델 평가: 교차검증", "Cross-Validation Strategies", "07", ACCENT_BLUE)

# ============================================================
# 슬라이드 31: 왜 데이터를 나누는가 + Hold-Out
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "왜 데이터를 나누는가? & Hold-Out")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "일반화 성능 (Generalization Performance) 평가의 핵심 원칙", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.5),
         "학습에 사용한 데이터와 평가에 사용하는 데이터를 반드시 분리해야 한다.\n학습 데이터로 평가하면 과적합 여부를 알 수 없다 → 낙관적 편향 (Optimistic Bias) 발생",
         font_size=14, color=LIGHT_GRAY)

# Hold-Out
add_shape(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(3.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(0.4),
         "Hold-Out (단순 분할)", font_size=18, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(0.9), Inches(4.3), Inches(5.5), Inches(0.4),
         "가장 단순한 방법: 데이터를 한 번만 학습/테스트로 분할 (예: 80:20)", font_size=14, color=WHITE)

add_bullet_list(s, Inches(0.9), Inches(4.8), Inches(5.5), Inches(2.0), [
    "장점:",
    "  • 계산 비용이 가장 낮음",
    "  • 구현이 간단",
    "단점:",
    "  • 분할에 따라 결과가 크게 달라짐 (높은 분산)",
    "  • 데이터를 비효율적으로 사용",
    "사용 시기: 데이터가 매우 많을 때, 빠른 프로토타이핑",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# Hold-Out 시각화
add_text(s, Inches(7.0), Inches(4.5), Inches(5.5), Inches(0.3),
         "Hold-Out 분할 예시:", font_size=13, color=WHITE, bold=True)
add_shape(s, Inches(7.0), Inches(4.9), Inches(4.0), Inches(0.4), ACCENT_BLUE)
add_text(s, Inches(7.0), Inches(4.9), Inches(4.0), Inches(0.4),
         "학습 데이터 (80%)", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
add_shape(s, Inches(11.0), Inches(4.9), Inches(1.0), Inches(0.4), ACCENT_RED)
add_text(s, Inches(11.0), Inches(4.9), Inches(1.0), Inches(0.4),
         "테스트", font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(7.0), Inches(5.6), Inches(5.5), Inches(0.3),
         "다른 분할 → 다른 결과:", font_size=13, color=ACCENT_RED, bold=True)
add_shape(s, Inches(7.0), Inches(6.0), Inches(3.0), Inches(0.3), ACCENT_BLUE)
add_shape(s, Inches(10.0), Inches(6.0), Inches(2.0), Inches(0.3), ACCENT_RED)
add_shape(s, Inches(7.0), Inches(6.4), Inches(4.5), Inches(0.3), ACCENT_BLUE)
add_shape(s, Inches(11.5), Inches(6.4), Inches(0.5), Inches(0.3), ACCENT_RED)

# ============================================================
# 슬라이드 32: K-Fold & Stratified K-Fold
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "K-Fold & Stratified K-Fold 교차검증")

# K-Fold
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "K-Fold 교차검증", font_size=18, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(0.5),
         "데이터를 K개 폴드로 나눠 K번 평가\n최종 성능 = (성능1+성능2+...+성능K) / K", font_size=13, color=WHITE)

# K-Fold 시각화
for fold_i in range(5):
    y = Inches(3.6) + Inches(0.28) * fold_i
    add_text(s, Inches(0.9), y, Inches(1.0), Inches(0.25),
             f"Fold {fold_i+1}:", font_size=10, color=DARK_GRAY)
    for block_j in range(5):
        x = Inches(2.0) + Inches(0.8) * block_j
        is_val = (block_j == fold_i)
        c = ACCENT_RED if is_val else RGBColor(0x2A, 0x4A, 0x6A)
        add_shape(s, x, y, Inches(0.7), Inches(0.22), c)

add_bullet_list(s, Inches(0.9), Inches(5.2), Inches(5.2), Inches(1.5), [
    "장점: 모든 데이터를 학습/평가에 활용",
    "단점: 클래스 불균형시 폴드별 비율 다를 수 있음",
    "K=5 또는 K=10이 표준",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# Stratified K-Fold
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "Stratified K-Fold 교차검증", font_size=18, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(0.5),
         "K-Fold + 각 폴드에서 클래스 비율을 원본과 동일하게 유지", font_size=13, color=WHITE)

add_bullet_list(s, Inches(7.1), Inches(3.5), Inches(5.2), Inches(3.0), [
    "장점:",
    "  • 클래스 불균형 데이터에서 안정적",
    "  • 가장 널리 권장되는 CV 전략",
    "  • sklearn의 기본 CV 전략",
    "",
    "단점:",
    "  • 회귀 문제에 직접 적용 불가 (연속값)",
    "",
    "사용 시기:",
    "  • 분류 문제, 특히 클래스 불균형",
    "  • 일반적으로 가장 추천",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 33: LOO + 전략 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "Leave-One-Out & 전략 비교표")

# LOO
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.2), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Leave-One-Out (LOO)", font_size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "한 번에 1개만 테스트 (K=N, 극단적 K-Fold)",
    "장점: 편향이 가장 낮음 (거의 전체 데이터로 학습)",
    "단점: 계산 비용 매우 높음 (N번 학습), 분산 높을 수 있음",
    "사용: 데이터가 매우 적을 때 (N < 50)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# K값 트레이드오프
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.2), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "K값에 따른 편향-분산 트레이드오프", font_size=15, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "K 작음 (K=2): 학습데이터 50% → 높은 편향, 낮은 분산",
    "K 큼 (K=N, LOO): 거의 전체 학습 → 낮은 편향, 높은 분산",
    "실무적 최적점: K = 5 ~ 10 가장 일반적 권장",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 비교 테이블
add_text(s, Inches(0.6), Inches(4.6), Inches(5), Inches(0.4),
         "네 가지 전략 비교", font_size=16, color=WHITE, bold=True)
add_table_slide(s,
    ["전략", "학습 데이터", "평가 횟수", "편향", "분산", "계산 비용"],
    [
        ["Hold-Out", "80%", "1회", "높음", "높음", "매우 낮음"],
        ["5-Fold", "80%(반복)", "5회", "중간", "중간", "중간"],
        ["10-Fold", "90%(반복)", "10회", "낮음", "중간~높", "중간~높"],
        ["LOO", "(N-1)/N", "N회", "최저", "높음", "매우 높음"],
    ],
    Inches(0.6), Inches(5.0), [1.8, 1.5, 1.3, 1.3, 1.3, 1.5], row_height=0.45, font_size=12, header_font_size=13)

# 데이터 크기별 권장
add_shape(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3), CARD_BG)
# 이 위치에 있어서 보이지 않으므로 생략

# ============================================================
# 슬라이드 34: 데이터 크기별 권장 + 평가지표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "데이터 크기별 권장 전략 & 평가 지표 개요")

# 권장 전략
recs = [
    ("대용량 (>10,000)", "Hold-Out 또는 3-Fold", "충분한 데이터로 계산 효율성 우선", ACCENT_BLUE),
    ("중간 (100~10,000)", "Stratified 5 or 10-Fold", "안정적 추정 + 적절한 비용", ACCENT_GREEN),
    ("소용량 (<100)", "LOO or Repeated K-Fold", "모든 데이터를 최대한 활용", ACCENT_ORANGE),
]
for i, (size, strat, reason, color) in enumerate(recs):
    x = Inches(0.5) + Inches(4.15) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(1.6), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.35),
             size, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.7), Inches(3.5), Inches(0.35),
             strat, font_size=14, color=WHITE, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.1), Inches(3.5), Inches(0.5),
             reason, font_size=12, color=DARK_GRAY)

# 평가 지표
add_text(s, Inches(0.6), Inches(4.1), Inches(5), Inches(0.4),
         "주요 평가 지표", font_size=16, color=WHITE, bold=True)

# 회귀
add_card(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.7),
         "회귀 평가 지표", [
             "RMSE = sqrt(mean((y_pred - y_true)²)) → 큰 오차에 민감",
             "MAE = mean(|y_pred - y_true|) → 이상치에 강건",
             "R² = 1 - SS_res/SS_tot → 설명력 (1에 가까울수록 좋음)",
             "",
             "RMSE vs MAE: 이상치 많으면 MAE 선호",
             "R² < 0 가능 (평균보다 못한 모델)",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 분류
add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.7),
         "분류 평가 지표", [
             "Accuracy = (TP+TN)/(TP+TN+FP+FN) → 불균형 주의!",
             "Precision = TP/(TP+FP) → 양성 예측의 정확도",
             "Recall = TP/(TP+FN) → 실제 양성 탐지율 (의료 핵심)",
             "F1 = 2*(Prec*Rec)/(Prec+Rec) → 조화 평균",
             "AUC-ROC = ROC 곡선 아래 면적 (0~1)",
             "  1.0=완벽, 0.5=랜덤 → 임계값 독립 종합 지표",
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# SECTION 08: ML 파이프라인
# ============================================================
section_divider("ML 파이프라인과 데이터 품질", "Machine Learning Pipeline & Data Quality", "08", ACCENT_CYAN)

# ============================================================
# 슬라이드 36: ML 파이프라인 8단계
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "머신러닝 파이프라인 8단계", "End-to-End ML Workflow")

steps = [
    ("1", "문제 정의", "회귀/분류 결정\n성공 기준 설정\n비즈니스→ML 변환", ACCENT_BLUE),
    ("2", "데이터 수집", "DB, API, 크롤링\n공개 데이터셋\n데이터 확보", ACCENT_CYAN),
    ("3", "데이터\n전처리", "결측치, 이상치\n인코딩, 스케일링\n원시→모델용 가공", ACCENT_GREEN),
    ("4", "EDA", "분포, 상관관계\n시각화\n데이터 특성 파악", ACCENT_GREEN),
    ("5", "모델링", "알고리즘 선택\n학습\n여러 모델 비교", ACCENT_ORANGE),
    ("6", "평가", "교차검증\n테스트셋 평가\n일반화 성능 측정", ACCENT_ORANGE),
    ("7", "HP 튜닝", "Grid/Random\nBayesian Search\n최적 파라미터", ACCENT_RED),
    ("8", "배포", "REST API\n모니터링\n재학습 파이프라인", ACCENT_PURPLE),
]
for i, (num, name, desc, color) in enumerate(steps):
    x = Inches(0.3) + Inches(1.6) * i
    if i < len(steps) - 1:
        add_shape(s, x + Inches(1.45), Inches(3.5), Inches(0.25), Inches(0.15), DARK_GRAY)
    add_shape(s, x, Inches(2.3), Inches(1.4), Inches(2.5), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.5), Inches(2.1), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.5), Inches(2.1), Inches(0.4), Inches(0.4),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), Inches(2.55), Inches(1.3), Inches(0.6),
             name, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), Inches(3.2), Inches(1.3), Inches(1.4),
             desc, font_size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 하단 핵심
add_shape(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(2.2), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.4),
         "\"Garbage In, Garbage Out\" - 데이터 품질의 중요성", font_size=16, color=ACCENT_ORANGE, bold=True)
add_table_slide(s,
    ["전처리 단계", "설명", "방법"],
    [
        ["결측치 처리", "누락된 데이터 처리", "삭제, 평균/중앙값 대체, 보간"],
        ["이상치 처리", "비정상적 극단값", "IQR, Z-score 기반 탐지"],
        ["인코딩", "범주형→수치 변환", "원핫 인코딩, 레이블 인코딩"],
        ["스케일링", "변수 크기/범위 통일", "StandardScaler, MinMaxScaler"],
    ],
    Inches(0.8), Inches(5.7), [2.5, 3.0, 5.5], row_height=0.4, font_size=12, header_font_size=13)

# ============================================================
# 슬라이드 37: 피처 엔지니어링 & 데이터 누출
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "피처 엔지니어링 & 데이터 누출 방지")

# 피처 엔지니어링
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Domingos 교훈 7: Feature Engineering Is the Key", font_size=14, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.5), [
    "알고리즘 선택보다 피처 설계가 더 큰 영향",
    "도메인 지식 기반 피처 구성이 핵심",
    "예: 주택 가격 예측",
    "  단순: 면적, 방 수, 위치",
    "  엔지니어링: 방당 면적, 지하철역 거리, 학군 등급",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 데이터 누출
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "Data Leakage (데이터 누출) 방지", font_size=14, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "학습 중 테스트 데이터 정보가 모델에 유입되는 현상",
    "→ 과도하게 낙관적인 성능 추정",
    "",
    "방지 원칙:",
    "  1. 스케일링: 학습셋 fit_transform, 테스트셋 transform만",
    "  2. 교차검증 루프 안에서 전처리 (Pipeline 활용)",
    "  3. 시계열: 미래 데이터를 학습에 사용하지 않음",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# 과적합 vs 과소적합
add_shape(s, Inches(0.6), Inches(4.8), Inches(3.9), Inches(2.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.8), Inches(4.9), Inches(3.5), Inches(0.4),
         "과소적합 (Underfitting)", font_size=15, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(0.8), Inches(5.3), Inches(3.5), Inches(1.8), [
    "학습/테스트 모두 낮은 성능",
    "모델이 너무 단순, 높은 편향",
    "해결: 복잡한 모델, 학습 시간 증가",
    "      피처 엔지니어링, 정규화 줄이기",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(4.8), Inches(4.8), Inches(3.5), Inches(2.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(5.0), Inches(4.9), Inches(3.1), Inches(0.4),
         "최적 모델 (Good Fit)", font_size=15, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(5.0), Inches(5.3), Inches(3.1), Inches(1.8), [
    "학습/테스트 모두 높은 성능",
    "Bias² + Variance 최소",
    "일반화 성능 최대",
    "교차검증으로 확인",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(8.6), Inches(4.8), Inches(3.9), Inches(2.5), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(8.8), Inches(4.9), Inches(3.5), Inches(0.4),
         "과적합 (Overfitting)", font_size=15, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(8.8), Inches(5.3), Inches(3.5), Inches(1.8), [
    "학습 높음, 테스트 낮음",
    "모델 너무 복잡, 높은 분산",
    "해결: 더 많은 데이터, 정규화(L1/L2)",
    "      드롭아웃, 조기 종료, 특성 선택",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# SECTION 09: 핵심 논문 리뷰
# ============================================================
section_divider("핵심 논문 리뷰 5편", "Key Paper Reviews", "09", ACCENT_GREEN)

# ============================================================
# 슬라이드 39: Domingos (2012) - 12가지 교훈 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "Domingos (2012): 12가지 핵심 교훈 (1/2)",
             "A Few Useful Things to Know About Machine Learning")

add_shape(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.6), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.5),
         "Pedro Domingos | Communications of the ACM, Vol. 55, No. 10, pp. 78-87 | ML 실무 12가지 교훈",
         font_size=13, color=LIGHT_GRAY)

lessons1 = [
    ("1", "Representation + Evaluation\n+ Optimization", "모든 학습 알고리즘의\n세 구성 요소", ACCENT_BLUE),
    ("2", "It's Generalization\nThat Counts", "일반화 성능이\n핵심 목표", ACCENT_BLUE),
    ("3", "Data Alone Is\nNot Enough", "사전 가정(귀납적 편향)이\n반드시 필요", ACCENT_CYAN),
    ("4", "Overfitting Has\nMany Faces", "과적합은 다양한\n형태로 나타남", ACCENT_CYAN),
    ("5", "Intuition Fails in\nHigh Dimensions", "고차원에서는 직관이\n작동하지 않음 (차원의 저주)", ACCENT_GREEN),
    ("6", "Theoretical Guarantees\nAre Not What They Seem", "이론적 보장은\n실제와 다를 수 있음", ACCENT_GREEN),
]
for i, (num, text, desc, color) in enumerate(lessons1):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.9) + Inches(2.3) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.0), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.6), y + Inches(0.15), Inches(3.2), Inches(0.8),
             text, font_size=13, color=WHITE)
    add_accent_line(s, x + Inches(0.2), y + Inches(1.0), Inches(3.0), color)
    add_text(s, x + Inches(0.2), y + Inches(1.15), Inches(3.5), Inches(0.7),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 40: Domingos (2012) - 12가지 교훈 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "Domingos (2012): 12가지 핵심 교훈 (2/2)")

lessons2 = [
    ("7", "Feature Engineering\nIs the Key", "피처 설계가\n알고리즘 선택보다 중요", ACCENT_ORANGE),
    ("8", "More Data Beats a\nCleverer Algorithm", "데이터 양이\n알고리즘 정교함보다 중요", ACCENT_ORANGE),
    ("9", "Learn Many Models,\nNot Just One", "앙상블이\n단일 모델보다 우수", ACCENT_RED),
    ("10", "Simplicity Does Not\nImply Accuracy", "단순한 모델이 반드시\n더 정확하지는 않음", ACCENT_RED),
    ("11", "Representable Does Not\nImply Learnable", "표현 가능 ≠ 학습 가능\n(학습 난이도는 별개)", ACCENT_PURPLE),
    ("12", "Correlation Does Not\nImply Causation", "상관관계 ≠ 인과관계\n(해석 주의)", ACCENT_PURPLE),
]
for i, (num, text, desc, color) in enumerate(lessons2):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.2) + Inches(2.3) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.0), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.6), y + Inches(0.15), Inches(3.2), Inches(0.8),
             text, font_size=13, color=WHITE)
    add_accent_line(s, x + Inches(0.2), y + Inches(1.0), Inches(3.0), color)
    add_text(s, x + Inches(0.2), y + Inches(1.15), Inches(3.5), Inches(0.7),
             desc, font_size=12, color=LIGHT_GRAY)

# 핵심 요약
add_shape(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(6.83), Inches(11.5), Inches(0.4),
         "실무 핵심: 일반화가 목표 | 피처 설계가 핵심 | 많은 데이터 > 정교한 알고리즘 | 앙상블 활용 | 상관≠인과",
         font_size=13, color=ACCENT_CYAN, bold=True)

# ============================================================
# 슬라이드 41: Mitchell, Wolpert, Vapnik, ESL 논문
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "핵심 논문 리뷰: Mitchell, Wolpert, Vapnik, ESL")

papers = [
    ("Mitchell (1997)", "Machine Learning (교재)", [
        "ML의 형식적 정의 (T-P-E) 정립",
        "개념 학습: 가설 공간 H, 버전 스페이스",
        "귀납적 편향(Inductive Bias):",
        "  모든 학습 알고리즘은 일반화를 위해",
        "  사전 가정이 필요 → NFL과 직결",
    ], ACCENT_BLUE),
    ("Wolpert & Macready (1997)", "NFL Theorems for Optimization", [
        "보편적 최적 알고리즘 불가능성 증명",
        "모든 함수에 평균하면 모든 알고리즘 동일",
        "\"이 특정 문제에 어떤 알고리즘이",
        "  가장 적합한가?\"가 올바른 질문",
    ], ACCENT_RED),
    ("Vapnik (1995)", "Statistical Learning Theory", [
        "경험적 위험: R_emp = (1/n)SUM L(f(xi),yi)",
        "실제 위험: R(f) = INT L(f(x),y) dP(x,y)",
        "VC 차원: shatter 가능한 최대 점 수",
        "SRM: R(f) <= R_emp + Phi(VCdim, n)",
        "핵심: 유한 데이터 → 일반화 조건?",
    ], ACCENT_GREEN),
    ("Hastie et al. (2009)", "ESL Ch1-2", [
        "EPE(f) = E[(Y-f(X))²]",
        "최적 예측기: f*(x) = E[Y|X=x]",
        "Bias-Variance Decomposition",
        "최소제곱법(높편향,낮분산) vs",
        "k-NN(낮편향,높분산) 대비 설명",
    ], ACCENT_PURPLE),
]
for i, (author, title, items, color) in enumerate(papers):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(6.3) * col
    y = Inches(2.2) + Inches(2.6) * row
    add_shape(s, x, y, Inches(6.0), Inches(2.3), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.35),
             author, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(5.6), Inches(0.3),
             title, font_size=12, color=WHITE)
    add_bullet_list(s, x + Inches(0.2), y + Inches(0.8), Inches(5.6), Inches(1.3),
                    items, font_size=11, color=LIGHT_GRAY, spacing=Pt(2))

# ============================================================
# SECTION 10: 실습 1 - 편향-분산 트레이드오프
# ============================================================
section_divider("실습 1: 편향-분산 트레이드오프", "01_bias_variance_tradeoff.py", "10", ACCENT_BLUE)

# ============================================================
# 슬라이드 43: 실습1 개요 + 핵심 코드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "실습 1: 편향-분산 트레이드오프 구현", "01_bias_variance_tradeoff.py")

add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(1.5), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "실습 목표", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(0.8), [
    "다양한 복잡도의 다항 회귀 모델로 편향², 분산, 총 오류 시각화",
    "과소적합/과적합 영역을 직관적으로 확인",
    "편향-분산 분해 공식을 실험적으로 검증",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(1.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "실험 설정", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(0.8), [
    "실제 함수: f(x) = sin(1.5*pi*x)",
    "노이즈: e ~ N(0, 0.3²), 200개 데이터셋 반복",
    "다항식 차수 1~15로 모델 복잡도 변화",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 핵심 코드 1: 데이터 생성
add_text(s, Inches(0.6), Inches(4.0), Inches(5), Inches(0.3),
         "핵심 코드 1: 실제 함수 & 데이터 생성", font_size=13, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8), [
    "def true_function(x):",
    "    return np.sin(1.5 * np.pi * x)",
    "",
    "def generate_data(n_samples=30, noise_std=0.3):",
    "    X = np.sort(np.random.uniform(0, 1, n_samples))",
    "    y = true_function(X) + \\",
    "        np.random.normal(0, noise_std, n_samples)",
    "    return X, y",
    "",
    "# y = f(x) + epsilon",
    "# epsilon ~ N(0, sigma^2)  <- 비가약 오류",
], font_size=11)

# 핵심 코드 2: 모델 학습
add_text(s, Inches(6.8), Inches(4.0), Inches(5), Inches(0.3),
         "핵심 코드 2: 다항 회귀 모델 학습", font_size=13, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.8), [
    "from sklearn.pipeline import make_pipeline",
    "from sklearn.preprocessing import PolynomialFeatures",
    "from sklearn.linear_model import LinearRegression",
    "",
    "model = make_pipeline(",
    "    PolynomialFeatures(degree, include_bias=True),",
    "    LinearRegression()",
    ")",
    "model.fit(X_train.reshape(-1, 1), y_train)",
    "predictions[i] = model.predict(X_test.reshape(-1,1))",
], font_size=11)

# ============================================================
# 슬라이드 44: 실습1 - 편향분산 계산 코드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "실습 1: 편향-분산 계산 핵심 로직")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(4.8), [
    "def compute_bias_variance(degrees, n_datasets=200, n_samples=30, noise_std=0.3):",
    "    X_test = np.linspace(0, 1, 100)   # 고정된 테스트 포인트",
    "    y_true = true_function(X_test)      # 실제 함수값",
    "    noise_var = noise_std ** 2          # 비가약 오류 = 0.09",
    "",
    "    for degree in degrees:",
    "        predictions = np.zeros((n_datasets, 100))  # 200개 데이터셋 x 100 테스트점",
    "",
    "        for i in range(n_datasets):            # 200번 반복",
    "            X_train, y_train = generate_data()  # 매번 새 데이터셋 생성",
    "            model = make_pipeline(PolynomialFeatures(degree), LinearRegression())",
    "            model.fit(X_train.reshape(-1,1), y_train)",
    "            predictions[i,:] = model.predict(X_test.reshape(-1,1))",
    "",
    "        # ============ 핵심: 편향-분산 분해 계산 ============",
    "        mean_pred = predictions.mean(axis=0)            # f_bar(x) = E[f_hat(x)]",
    "        bias_sq = np.mean((mean_pred - y_true) ** 2)    # Bias² = (f(x) - f_bar(x))²",
    "        var = np.mean(predictions.var(axis=0))           # Var = E[(f_hat - f_bar)²]",
    "        total = bias_sq + var + noise_var                # Total = Bias² + Var + sigma²",
], font_size=11)

add_shape(s, Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3), CARD_BG)
add_text(s, Inches(0.9), Inches(7.1), Inches(11.5), Inches(0.3),
         "핵심: 200개 데이터셋으로 반복 학습 → 예측의 평균(편향)과 변동(분산)을 실험적으로 계산",
         font_size=12, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 11: 실습 2 - NFL 정리 데모
# ============================================================
section_divider("실습 2: NFL 정리 데모", "02_no_free_lunch_demo.py", "11", ACCENT_ORANGE)

# ============================================================
# 슬라이드 46: 실습2 개요 + 핵심 코드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "11", "실습 2: NFL 정리 데모", "02_no_free_lunch_demo.py")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "실습 목표: 4가지 데이터 구조에 5가지 알고리즘 적용 → 단일 최상 알고리즘 없음을 실험적 확인",
         font_size=15, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(0.4),
         "Stratified 5-Fold 교차검증으로 공정한 비교 수행", font_size=14, color=LIGHT_GRAY)

# 데이터셋 설명
ds_items = [
    ("선형 분리 (Linear)", "선형 결정 경계\n→ 선형 모델에 유리", ACCENT_BLUE),
    ("원형 (Circles)", "동심원 구조\n→ RBF SVM에 유리", ACCENT_CYAN),
    ("XOR 패턴", "XOR 논리 연산\n→ 결정 트리에 유리", ACCENT_GREEN),
    ("반달 (Moons)", "반달 형태\n→ k-NN에 유리", ACCENT_ORANGE),
]
add_text(s, Inches(0.6), Inches(3.6), Inches(3), Inches(0.3),
         "4가지 데이터셋:", font_size=14, color=WHITE, bold=True)
for i, (name, desc, color) in enumerate(ds_items):
    x = Inches(0.5) + Inches(3.15) * i
    add_shape(s, x, Inches(3.95), Inches(2.9), Inches(1.3), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(4.0), Inches(2.6), Inches(0.35),
             name, font_size=13, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(4.4), Inches(2.6), Inches(0.7),
             desc, font_size=11, color=LIGHT_GRAY)

# 알고리즘 + 핵심 코드
add_text(s, Inches(0.6), Inches(5.5), Inches(5), Inches(0.3),
         "5가지 알고리즘 (각기 다른 귀납적 편향):", font_size=14, color=WHITE, bold=True)
add_code_block(s, Inches(0.6), Inches(5.85), Inches(5.8), Inches(1.4), [
    "LogisticRegression()     # 선형 결정 경계",
    "KNeighborsClassifier(5)  # 국소적 유사성",
    "SVC(kernel='rbf')        # 마진 최대화 + 비선형",
    "DecisionTreeClassifier() # 축 정렬 분할",
    "RandomForestClassifier() # 앙상블 + 분할",
], font_size=11)

add_code_block(s, Inches(6.8), Inches(5.85), Inches(5.8), Inches(1.4), [
    "# 공정한 비교: Stratified 5-Fold CV",
    "cv = StratifiedKFold(n_splits=5, shuffle=True)",
    "scores = cross_val_score(clf, X_scaled, y,",
    "                         cv=cv, scoring='accuracy')",
    "# 핵심 결론: 모든 데이터셋 1위인 알고리즘 없음!",
], font_size=11)

# ============================================================
# SECTION 12: 실습 3 - 교차검증 비교
# ============================================================
section_divider("실습 3: 교차검증 비교", "03_cross_validation_demo.py", "12", ACCENT_RED)

# ============================================================
# 슬라이드 48: 실습3 개요 + 핵심 코드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "12", "실습 3: 교차검증 전략 비교", "03_cross_validation_demo.py")

add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.0), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.35),
         "실습 목표: Hold-Out, K-Fold, Stratified K-Fold, LOO 비교 | 50회 반복으로 성능 추정 분산 측정 | K값 트레이드오프 확인",
         font_size=14, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.3),
         "데이터: make_classification(150샘플) + Iris 데이터셋", font_size=13, color=LIGHT_GRAY)

# 분산 비교 코드
add_text(s, Inches(0.6), Inches(3.4), Inches(6), Inches(0.3),
         "핵심 코드 1: 성능 추정 분산 비교 (50회 반복)", font_size=13, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(0.6), Inches(3.7), Inches(5.8), Inches(3.0), [
    "n_repeats = 50",
    "",
    "# Hold-Out: 분할마다 결과가 다름",
    "for i in range(n_repeats):",
    "    X_tr, X_te, y_tr, y_te = train_test_split(",
    "        X, y, test_size=0.2, random_state=i)",
    "    model.fit(X_tr, y_tr)",
    "    holdout_scores.append(model.score(X_te, y_te))",
    "",
    "# Stratified K-Fold: 가장 안정적",
    "for i in range(n_repeats):",
    "    cv = StratifiedKFold(n_splits=5, shuffle=True,",
    "                         random_state=i)",
    "    scores = cross_val_score(model, X, y, cv=cv)",
    "    skfold_scores.append(scores.mean())",
], font_size=10)

# K값 트레이드오프 코드
add_text(s, Inches(6.8), Inches(3.4), Inches(6), Inches(0.3),
         "핵심 코드 2: K값에 따른 트레이드오프", font_size=13, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(6.8), Inches(3.7), Inches(5.8), Inches(3.0), [
    "k_values = [2, 3, 5, 7, 10, 15, 20, 50, 100]",
    "",
    "for k in k_values:",
    "    repeat_means = []",
    "    for i in range(30):  # 30회 반복",
    "        cv = KFold(n_splits=k, shuffle=True,",
    "                   random_state=i)",
    "        scores = cross_val_score(model, X, y, cv=cv)",
    "        repeat_means.append(scores.mean())",
    "    means.append(np.mean(repeat_means))",
    "    stds.append(np.std(repeat_means))",
    "",
    "# 핵심 관찰:",
    "# K 작음 → 높은 편향, 낮은 분산",
    "# K 큼   → 낮은 편향, 높은 분산",
], font_size=10)

add_shape(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(6.93), Inches(11.5), Inches(0.3),
         "핵심 결론: Hold-Out 분산 가장 큼 | Stratified K-Fold 가장 안정적 | K=5~10 실무 권장",
         font_size=13, color=ACCENT_GREEN, bold=True)

# ============================================================
# SECTION 13: 응용사례 심화
# ============================================================
section_divider("응용사례 심화", "Healthcare, Finance, Manufacturing, NLP", "13", ACCENT_PURPLE)

# ============================================================
# 슬라이드 50: 응용사례 - 의료 & 금융
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "13", "응용사례 심화: 의료 & 금융")

# 의료
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "의료 (Healthcare) - 유방암 진단", font_size=16, color=ACCENT_GREEN, bold=True)
add_table_slide(s,
    ["항목", "내용"],
    [
        ["데이터", "Wisconsin Breast Cancer (569샘플, 30특성)"],
        ["클래스", "악성 37.3% / 양성 62.7%"],
        ["핵심 지표", "재현율(Recall) - 암 환자 놓치지 않기"],
        ["성능", "로지스틱 회귀만으로 정확도 97~98%"],
    ],
    Inches(0.8), Inches(2.8), [1.5, 3.8], row_height=0.4, font_size=12, header_font_size=13)

add_bullet_list(s, Inches(0.9), Inches(4.6), Inches(5.0), Inches(2.0), [
    "의료 AI 핵심 고려사항:",
    "  • FN(위음성) 최소화가 생명과 직결",
    "  • 정확도보다 재현율과 F1-Score 중시",
    "  • 모델 해석 가능성 필수 (의사 이해 필요)",
    "  • 규제 요건 충족 (FDA 승인 등)",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# 금융
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "금융 (Finance) - 사기 탐지", font_size=16, color=ACCENT_ORANGE, bold=True)
add_table_slide(s,
    ["항목", "내용"],
    [
        ["문제", "극심한 클래스 불균형 (사기 < 1%)"],
        ["핵심 과제", "FP 최소화 + 실제 사기 Recall"],
        ["주요 기법", "Isolation Forest, SMOTE, 앙상블"],
        ["특이점", "실시간 처리 + 패턴 지속 변화"],
    ],
    Inches(7.0), Inches(2.8), [1.5, 3.8], row_height=0.4, font_size=12, header_font_size=13)

add_bullet_list(s, Inches(7.1), Inches(4.6), Inches(5.0), Inches(2.0), [
    "금융 ML 핵심 고려사항:",
    "  • 시계열: Time Series Split, 미래 데이터 누출 방지",
    "  • 클래스 불균형: SMOTE, Cost-sensitive Learning",
    "  • 해석 가능성: 규제 당국 모델 설명 필요",
    "  • 실시간 추론 성능 요구",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 51: 응용사례 - 제조 & NLP
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "13", "응용사례 심화: 제조 & NLP")

# 제조
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "제조 (Manufacturing) - 예측 정비", font_size=16, color=ACCENT_BLUE, bold=True)
add_table_slide(s,
    ["항목", "내용"],
    [
        ["문제", "장비 고장 시점/여부 예측"],
        ["데이터", "IoT 센서 (진동, 온도, 압력, 소음)"],
        ["기법", "RF, LSTM, Survival Analysis"],
        ["효과", "비계획 정지 30~50% 감소"],
    ],
    Inches(0.8), Inches(2.8), [1.5, 3.8], row_height=0.4, font_size=12, header_font_size=13)

add_bullet_list(s, Inches(0.9), Inches(4.6), Inches(5.0), Inches(2.0), [
    "제조 ML 핵심:",
    "  • 센서 데이터의 시계열 특성",
    "  • 고장 데이터가 매우 희소 (불균형)",
    "  • Edge Computing으로 실시간 추론",
    "  • 도메인 전문가(공정 엔지니어) 협업 필수",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# NLP
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "NLP - 감성 분석", font_size=16, color=ACCENT_PURPLE, bold=True)
add_table_slide(s,
    ["항목", "내용"],
    [
        ["문제", "텍스트 분류 (긍정/부정/중립)"],
        ["데이터", "리뷰, 소셜 미디어, 뉴스"],
        ["전통 접근", "TF-IDF + SVM/로지스틱 회귀"],
        ["현대 접근", "BERT, GPT Fine-tuning"],
    ],
    Inches(7.0), Inches(2.8), [1.5, 3.8], row_height=0.4, font_size=12, header_font_size=13)

add_bullet_list(s, Inches(7.1), Inches(4.6), Inches(5.0), Inches(2.0), [
    "NLP ML 핵심:",
    "  • 전처리: 토큰화, 불용어 제거, 정규화",
    "  • 피처: BoW, TF-IDF, Word2Vec, BERT Embedding",
    "  • 전이 학습: 사전학습 모델 Fine-tuning",
    "  • 한국어 형태소 분석의 특수성",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# SECTION 14: 핵심 요약 + 복습 질문
# ============================================================
section_divider("핵심 요약 & 복습 질문", "Summary & Review Questions", "14", ACCENT_CYAN)

# ============================================================
# 슬라이드 53: 핵심 요약
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "14", "핵심 요약 (Key Takeaways)")

summaries = [
    ("ML 정의", "데이터에서 패턴 자동 학습\nMitchell T-P-E 프레임워크", ACCENT_BLUE),
    ("학습 유형", "지도(분류/회귀) · 비지도 · 강화\nAI ⊃ ML ⊃ DL", ACCENT_CYAN),
    ("편향-분산", "EPE = Bias² + Var + σ²\n모델 복잡도 최적 균형점", ACCENT_GREEN),
    ("NFL 정리", "보편적 최상 알고리즘 없음\n문제에 맞는 선택이 핵심", ACCENT_ORANGE),
    ("교차검증", "Stratified K-Fold (K=5~10)\n데이터 크기에 따라 선택", ACCENT_RED),
    ("ML 파이프라인", "문제정의→수집→전처리→EDA\n→모델링→평가→튜닝→배포", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(summaries):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.2) + Inches(2.5) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.1), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
             title, font_size=17, color=color, bold=True)
    add_accent_line(s, x + Inches(0.2), y + Inches(0.6), Inches(2.0), color)
    add_text(s, x + Inches(0.2), y + Inches(0.75), Inches(3.5), Inches(1.1),
             desc, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 54: 핵심 수식 정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "14", "핵심 수식 정리")

formulas = [
    ("편향-분산 분해", "EPE = Bias² + Variance + σ²", "총 오차의 세 가지 원천"),
    ("편향", "Bias(x) = f(x) - E[f_hat(x)]", "모델 가정에 의한 체계적 오차"),
    ("분산", "Var(x) = E[(f_hat(x) - E[f_hat(x)])²]", "데이터 변동에 의한 예측 불안정"),
    ("NFL 정리", "Sf P(d|f,m,a1) = Sf P(d|f,m,a2)", "모든 알고리즘의 평균 성능 동일"),
    ("구조적 위험 최소화", "R(f) <= R_emp(f) + Phi(VCdim, n)", "일반화 오류 상한 (Vapnik)"),
    ("최적 예측기 (회귀)", "f*(x) = E[Y|X=x]", "MSE를 최소화하는 예측 (ESL)"),
]
for i, (name, formula, meaning) in enumerate(formulas):
    y = Inches(2.2) + Inches(0.85) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    add_shape(s, Inches(0.6), y, Inches(12.1), Inches(0.75), bg, radius=True)
    add_text(s, Inches(0.9), y + Inches(0.05), Inches(2.5), Inches(0.65),
             name, font_size=14, color=ACCENT_CYAN, bold=True)
    add_text(s, Inches(3.5), y + Inches(0.05), Inches(4.5), Inches(0.65),
             formula, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.2), y + Inches(0.05), Inches(4.3), Inches(0.65),
             meaning, font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 55: 복습 질문 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "14", "복습 질문 (1/2)")

questions1 = [
    "Q1. Mitchell의 T-P-E 정의를 스팸 필터에 적용하여 설명하시오.",
    "Q2. 지도/비지도/강화학습을 데이터, 목표, 피드백 관점에서 비교하시오.",
    "Q3. 편향-분산 분해 공식의 각 항을 설명하고\n      모델 복잡도와의 관계를 서술하시오.",
    "Q4. NFL 정리의 핵심과 실무적 시사점 3가지를 서술하시오.",
    "Q5. Hold-Out, K-Fold, Stratified K-Fold, LOO의\n      장단점을 비교하시오.",
]
colors5 = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
for i, q in enumerate(questions1):
    y = Inches(2.2) + Inches(1.0) * i
    add_shape(s, Inches(0.8), y, Inches(11.5), Inches(0.85), CARD_BG, colors5[i], radius=True)
    add_shape(s, Inches(0.8), y + Inches(0.15), Inches(0.08), Inches(0.5), colors5[i])
    add_text(s, Inches(1.2), y + Inches(0.1), Inches(10.8), Inches(0.7),
             q, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 56: 복습 질문 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "14", "복습 질문 (2/2)")

questions2 = [
    "Q6. 과적합과 과소적합의 원인과 해결방법을 각 3가지 이상 서술하시오.\n      Domingos의 \"Overfitting Has Many Faces\"와 연결하여 설명하시오.",
    "Q7. VC 차원이란 무엇이며, 2차원 선형 분류기의\n      VC 차원이 3인 이유를 설명하시오.",
    "Q8. ML 파이프라인 8단계와 Data Leakage 방지 원칙을 설명하시오.",
    "Q9. ESL의 최소제곱법 vs k-NN 비교를 통한\n      편향-분산 트레이드오프를 설명하시오.",
    "Q10. Domingos의 \"More Data Beats a Cleverer Algorithm\"과\n       \"Feature Engineering Is the Key\"를 실무 예시와 함께 설명하시오.",
]
colors5b = [ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE]
for i, q in enumerate(questions2):
    y = Inches(2.2) + Inches(1.0) * i
    add_shape(s, Inches(0.8), y, Inches(11.5), Inches(0.85), CARD_BG, colors5b[i], radius=True)
    add_shape(s, Inches(0.8), y + Inches(0.15), Inches(0.08), Inches(0.5), colors5b[i])
    add_text(s, Inches(1.2), y + Inches(0.1), Inches(10.8), Inches(0.7),
             q, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 57: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.5), prs.slide_width, Inches(1.0),
         "Thank You", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), prs.slide_width, Inches(0.6),
         "1장: 머신러닝 개요 (확장 상세 버전)", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(4.2), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(4.5), prs.slide_width, Inches(0.5),
         "다음 장: 2장 - 파이썬 기초", font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ── 저장 ──
output_path = os.path.join(os.path.dirname(__file__), "1장_머신러닝_개요_강의PPT_확장.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
