# -*- coding: utf-8 -*-
"""
5장 로지스틱 회귀 강의 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)

def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    card = add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)


# ============================================================
# 슬라이드 1: 표지
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)

# 상단 장식선
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)

# 챕터 번호
add_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
         "CHAPTER 5", font_size=20, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

# 메인 타이틀
add_text(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(1.0),
         "로지스틱 회귀", font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.6), Inches(3.2), Inches(12), Inches(0.6),
         "Logistic Regression", font_size=28, color=ACCENT_CYAN, bold=False, align=PP_ALIGN.CENTER)

# 구분선
add_accent_line(slide, Inches(5.0), Inches(4.1), Inches(3.3), ACCENT_BLUE)

# 키워드
add_text(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.5),
         "분류 | 시그모이드 | 로그 손실 | 오즈비 | 소프트맥스 | 판별 vs 생성 모델 | 희귀 이벤트",
         font_size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 논문 참고
add_text(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.4),
         "Cox(1958) | Hosmer & Lemeshow(2000) | King & Zeng(2001) | Ng & Jordan(2002) | Menard(2002)",
         font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 하단 정보
add_text(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
         "기계학습 | 2026년 1학기", font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# 슬라이드 2: 목차 (TOC)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "0", "목차 (Table of Contents)")

toc_left = [
    "5.1  분류 문제의 본질 - 왜 선형회귀로 분류하면 안 되는가",
    "5.2  시그모이드와 로짓 - 오즈비, 로짓 변환",
    "5.3  MLE와 크로스엔트로피 - 비용 함수 유도",
    "5.4  경사하강법 - 배치/미니배치/SGD, 학습률",
    "5.5  다중 클래스 - 소프트맥스, OvR, OvO",
    "5.6  판별 vs 생성 모델 - Ng & Jordan(2002)",
    "5.7  희귀 이벤트 - King & Zeng(2001) 보정",
]
toc_right = [
    "5.8   모델 평가 - ROC/AUC, 혼동행렬, F1, MCC",
    "5.9   논문 리뷰 - Cox, Hosmer, Menard 등",
    "5.10  실습: 스크래치 구현",
    "5.11  실습: 판별 vs 생성 비교",
    "5.12  실습: 희귀 이벤트 보정",
    "5.13  응용사례 - 사기탐지, 질병진단, 고객이탈",
    "5.14  핵심 요약 + 복습 질문",
]

add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(5.0),
                toc_left, font_size=15, color=LIGHT_GRAY, spacing=Pt(10))
add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.8), Inches(5.0),
                toc_right, font_size=15, color=LIGHT_GRAY, spacing=Pt(10))


# ============================================================
# 슬라이드 3: 분류 문제의 본질
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.1", "분류 문제의 본질", "왜 선형회귀로 분류하면 안 되는가?")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.2),
         "선형 회귀 (Linear Regression)", [
             "문제 유형: 회귀 (연속값 예측)",
             "출력값: 연속적인 수치 (-inf ~ +inf)",
             "활성화 함수: 없음 (항등 함수)",
             "비용 함수: MSE (평균제곱오차)",
             "예시: 주택 가격, 보험료 예측",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.2),
         "로지스틱 회귀 (Logistic Regression)", [
             "문제 유형: 분류 (클래스 예측)",
             "출력값: 클래스 확률 (0~1 사이)",
             "활성화 함수: 시그모이드 (Sigmoid)",
             "비용 함수: Log Loss (로그 손실)",
             "예시: 생존/사망, 스팸/정상, 합격/불합격",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# 선형회귀 분류 문제점
add_card(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.6),
         "선형 회귀로 분류하면 발생하는 4가지 문제 (Menard, 2002)", [
             "1. 확률 범위 위반: 출력이 [0,1] 범위를 벗어남 (음수 or 1 초과 확률)",
             "2. 이분산성 (Heteroscedasticity): 오차 분산이 y_hat(1-y_hat)으로 예측값에 의존",
             "3. 비정규 오차: 잔차가 이산적 값만 취하여 정규성 가정 위배",
             "4. 비선형 관계: 확률과 설명 변수 간 실제 관계는 S-curve (포화 현상)",
             "=> 해결: 시그모이드 함수로 출력을 (0,1) 범위로 변환!",
         ], title_color=ACCENT_RED, border=ACCENT_RED)


# ============================================================
# 슬라이드 4: 시그모이드 함수와 핵심 아이디어
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.2", "시그모이드 함수와 로지스틱 회귀의 핵심 아이디어")

# 핵심 흐름도
add_shape(slide, Inches(0.6), Inches(2.1), Inches(12.1), Inches(0.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.6),
         "입력(X)  -->  선형 결합 z = w^T x + b  -->  시그모이드 sigma(z)  -->  확률 P  -->  클래스(0 or 1)",
         font_size=18, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

# 시그모이드 수식 카드
add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(2.0),
         "시그모이드 함수 수식", [
             "sigma(z) = 1 / (1 + e^(-z))",
             "",
             "z = w^T x + b  (선형 결합 결과)",
             "출력: 0과 1 사이의 확률값",
             "미분: sigma'(z) = sigma(z) * (1 - sigma(z))",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 시그모이드 성질 카드
add_card(slide, Inches(6.9), Inches(3.2), Inches(5.8), Inches(2.0),
         "시그모이드 함수의 핵심 성질", [
             "z -> +inf  :  sigma(z) -> 1 (양성 확률 높음)",
             "z = 0       :  sigma(z) = 0.5 (결정 경계)",
             "z -> -inf   :  sigma(z) -> 0 (음성 확률 높음)",
             "",
             "=> S자 곡선으로 포화 현상을 자연스럽게 표현",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# 확률 출력과 임계값
add_card(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.7),
         "확률 출력과 임계값 (Threshold)", [
             "기본 임계값 0.5:  P(y=1|X) >= 0.5 -> 클래스 1(양성),  P(y=1|X) < 0.5 -> 클래스 0(음성)",
             "질병 진단: 임계값 낮춤 (0.3) -> Recall 향상  |  스팸 필터: 임계값 높임 (0.7) -> Precision 향상",
             "sklearn:  model.predict() -> 클래스,  model.predict_proba() -> 확률  [Cox(1958)에 의해 체계화]",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# 슬라이드 5: 오즈비와 로짓 변환
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.2", "오즈비와 로짓 변환", "확률 -> 오즈 -> 로짓 -> 선형 결합")

add_card(slide, Inches(0.6), Inches(2.1), Inches(3.8), Inches(2.4),
         "오즈 (Odds)", [
             "Odds = p / (1 - p)",
             "",
             "p=0.75 -> Odds=3",
             "  (생존이 사망의 3배)",
             "p=0.50 -> Odds=1",
             "  (동일 확률)",
             "p=0.20 -> Odds=0.25",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(4.7), Inches(2.1), Inches(3.8), Inches(2.4),
         "로짓 변환 (Logit)", [
             "logit(p) = ln(p / (1-p)) = z",
             "= w^T x + b",
             "",
             "범위 변환:",
             "  확률 p: (0, 1)",
             "  오즈:  (0, +inf)",
             "  로짓: (-inf, +inf)",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(8.8), Inches(2.1), Inches(3.9), Inches(2.4),
         "오즈비 (Odds Ratio)", [
             "OR = e^(w_i)",
             "",
             "w_i > 0 -> OR > 1 -> 양성 확률 증가",
             "w_i = 0 -> OR = 1 -> 영향 없음",
             "w_i < 0 -> OR < 1 -> 양성 확률 감소",
             "",
             "피처 1단위 증가 시 오즈 변화량",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.4),
         "타이타닉 예시: 오즈비 해석", [
             "Sex_male 계수 = -2.569  ->  OR = e^(-2.569) = 0.077  ->  남성은 여성보다 생존 오즈가 92.3% 낮다",
             "Pclass 계수 = -1.182  ->  등급이 높아질수록(3등실 방향) 생존 오즈 감소",
             "",
             "핵심 의의: '확률을 로짓 변환하면 피처들의 선형 결합과 같다' (Cox, 1958)",
             "Menard(2002): 오즈비 외에 표준화 로짓 계수, 한계 효과(marginal effect) 등 다양한 해석법 비교",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)


# ============================================================
# 슬라이드 6: MLE와 크로스엔트로피
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.3", "MLE와 크로스엔트로피 (비용 함수)", "왜 MSE 대신 Log Loss를 사용하는가?")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.6),
         "최대우도추정 (MLE) 원리", [
             "단일 관측치 확률:",
             "  P(y_i|x_i; w) = y_hat^(y_i) * (1-y_hat)^(1-y_i)",
             "",
             "우도 함수 (전체):",
             "  L(w) = PROD y_hat^(y_i) * (1-y_hat)^(1-y_i)",
             "",
             "로그 우도 (곱 -> 합으로 변환):",
             "  l(w) = SUM [y*ln(y_hat) + (1-y)*ln(1-y_hat)]",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.6),
         "이진 교차 엔트로피 (Log Loss)", [
             "비용 함수 (부호 반전 + 평균):",
             "  J(w) = -(1/n) SUM [y*ln(y_hat)",
             "              + (1-y)*ln(1-y_hat)]",
             "",
             "MLE 최대화 = Log Loss 최소화 (동치)",
             "",
             "MSE를 쓰면 비볼록(non-convex)이 되어",
             "경사하강법으로 전역 최솟값 찾기 어려움",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2),
         "로그 손실의 직관적 이해", [
             "y=1, y_hat=0.9 -> Loss=-ln(0.9)=0.105 (정확 -> 낮은 손실)     y=1, y_hat=0.1 -> Loss=-ln(0.1)=2.303 (틀림 -> 높은 손실)",
             "y=0, y_hat=0.1 -> Loss=-ln(0.9)=0.105 (정확 -> 낮은 손실)     y=0, y_hat=0.9 -> Loss=-ln(0.1)=2.303 (틀림 -> 높은 손실)",
             "",
             "핵심: 예측이 틀릴수록 손실이 기하급수적으로 증가 -> 확신 있는 잘못된 예측에 큰 페널티",
         ], title_color=ACCENT_RED, border=ACCENT_RED)


# ============================================================
# 슬라이드 7: 그래디언트 유도와 경사하강법
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.4", "경사하강법과 그래디언트 유도", "배치 / 미니배치 / SGD 비교")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.0),
         "그래디언트 수식 (Cox의 스코어 방정식)", [
             "dJ/dw_j = (1/n) SUM (y_hat_i - y_i) * x_ij",
             "dJ/db   = (1/n) SUM (y_hat_i - y_i)",
             "",
             "벡터형: grad_w J = (1/n) X^T (y_hat - y)",
             "파라미터 업데이트: w <- w - alpha * grad_w J",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.0),
         "학습률 (Learning Rate) 선택", [
             "너무 작음 (0.001): 수렴이 매우 느림",
             "적절함 (0.01~0.1): 안정적으로 빠르게 수렴",
             "너무 큼 (10): 발산 (비용 함수 증가)",
             "",
             "탐색: {0.001, 0.003, 0.01, 0.03, 0.1, 0.3, 1.0}",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(0.6), Inches(4.4), Inches(3.8), Inches(2.8),
         "배치 GD (Batch)", [
             "사용 데이터: 전체 (n개)",
             "장점: 안정적 수렴",
             "       정확한 그래디언트",
             "단점: 대용량 데이터에서 느림",
             "       메모리 많이 사용",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(4.7), Inches(4.4), Inches(3.8), Inches(2.8),
         "미니배치 GD (Mini-batch)", [
             "사용 데이터: 일부 (B개, 32~256)",
             "장점: 배치와 SGD의 절충",
             "       GPU 병렬화 효율적",
             "단점: 배치 크기 B 하이퍼파라미터",
             "=> 실전에서 가장 많이 사용!",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(slide, Inches(8.8), Inches(4.4), Inches(3.9), Inches(2.8),
         "SGD (Stochastic)", [
             "사용 데이터: 단일 샘플 (1개)",
             "장점: 매우 빠른 업데이트",
             "       지역 최솟값 탈출 가능",
             "단점: 노이즈 큼",
             "       수렴 불안정",
         ], title_color=ACCENT_RED, border=ACCENT_RED)


# ============================================================
# 슬라이드 8: 다중 클래스 - 소프트맥스
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.5", "다중 클래스 분류", "소프트맥스, OvR, OvO 전략")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.2),
         "이진 분류 vs 다중 분류", [
             "이진 분류:  2개 클래스, 시그모이드, 1개 확률값",
             "  예: 생존/사망, 스팸/정상",
             "",
             "다중 분류:  K개 클래스(K>=3), 소프트맥스, K개 확률",
             "  예: 붓꽃 종 분류, 숫자 인식(0~9)",
             "  모든 확률의 합 = 1",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.2),
         "소프트맥스 함수 (Softmax)", [
             "P(y=k|X) = e^(z_k) / SUM_j e^(z_j)",
             "",
             "핵심 성질:",
             "  - 모든 출력이 (0, 1) 범위",
             "  - 모든 클래스 확률의 합 = 1",
             "  - K=2일 때 시그모이드와 동일",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(4.6), Inches(3.8), Inches(2.6),
         "OvR (One-vs-Rest)", [
             "K개 이진 분류기 학습",
             "각 클래스 vs 나머지 전부",
             "가장 높은 확률의 클래스 선택",
             "",
             "sklearn: multi_class='ovr'",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(4.7), Inches(4.6), Inches(3.8), Inches(2.6),
         "OvO (One-vs-One)", [
             "K(K-1)/2개 이진 분류기 학습",
             "모든 클래스 쌍에 대해 학습",
             "투표(voting)로 최종 결정",
             "",
             "개별 학습은 빠르나 분류기 수 많음",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(slide, Inches(8.8), Inches(4.6), Inches(3.9), Inches(2.6),
         "범주형 교차 엔트로피", [
             "J = -(1/n) SUM_i SUM_k",
             "    y_ik * ln(y_hat_ik)",
             "",
             "y_ik: 원-핫 인코딩",
             "(클래스 k에 속하면 1, 아니면 0)",
         ], title_color=ACCENT_RED, border=ACCENT_RED)


# ============================================================
# 슬라이드 9: 판별 vs 생성 모델
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.6", "판별 vs 생성 모델", "Ng & Jordan (2002) 핵심 비교")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.5),
         "판별 모델 (Discriminative Model)", [
             "학습 대상: P(y|x) 직접 학습",
             "대표 모델: 로지스틱 회귀, SVM, 신경망",
             "결정 경계를 직접 학습",
             "점근적으로 더 정확 (낮은 편향)",
             "수렴 속도: O(n) (느림)",
             "대표본에서 우수",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.5),
         "생성 모델 (Generative Model)", [
             "학습 대상: P(x|y)와 P(y) -> 베이즈 정리",
             "대표 모델: 나이브 베이즈, GMM, HMM",
             "데이터 분포 자체를 모델링",
             "적은 데이터에서 효과적 (낮은 분산)",
             "수렴 속도: O(log n) (빠름)",
             "소표본에서 우수",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.3),
         "Ng & Jordan (2002)의 핵심 발견: 전환점 (Crossover Point)", [
             "표본이 작을 때는 나이브 베이즈(생성)가 우수, 표본이 클 때는 로지스틱 회귀(판별)가 우수",
             "이 '전환점(crossover point)'이 항상 존재하며, 모델 가정의 정확도에 따라 달라짐",
             "나이브 베이즈 가정이 잘 성립하면 전환점이 늦게 / 위반되면 전환점이 일찍 나타남",
             "핵심 통찰: 두 모델은 동일한 모델 패밀리 공유, 파라미터 추정 방법만 다름 (결합 우도 vs 조건부 우도)",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# 슬라이드 10: 희귀 이벤트
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.7", "희귀 이벤트와 불균형 데이터", "King & Zeng (2001) 편향 보정")

add_card(slide, Inches(0.6), Inches(2.1), Inches(3.8), Inches(2.0),
         "희소 사건 실전 예시", [
             "국제 분쟁 발생: ~1-3%",
             "신용카드 사기: ~0.1%",
             "희귀 질환 진단: ~1%",
             "불량품 탐지: ~0.5%",
             "침입 탐지: ~0.01%",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(slide, Inches(4.7), Inches(2.1), Inches(3.8), Inches(2.0),
         "유한표본 편향 문제", [
             "MLE가 희소 사건의 확률을",
             "과소추정하는 방향으로 편향",
             "",
             "양성 비율이 낮을수록",
             "편향이 더 심해짐",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(8.8), Inches(2.1), Inches(3.9), Inches(2.0),
         "실용적 권고사항", [
             "> 20%: 표준 LR 사용 가능",
             "5~20%: 가중 LR 권고",
             "< 5%: 반드시 보정 적용",
             "",
             "평가: AUC, Recall, F1 사용",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(slide, Inches(0.6), Inches(4.4), Inches(5.8), Inches(2.8),
         "보정 방법 1: 사전 교정 (Prior Correction)", [
             "beta_0_corrected = beta_0 - ln[(1-tau)/tau * y_bar/(1-y_bar)]",
             "",
             "tau: 인구 수준의 사건 발생률",
             "y_bar: 표본에서의 사건 비율",
             "",
             "케이스-통제 설계에서 양성 사례를",
             "과대 표집했을 때의 편향 교정",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(4.4), Inches(5.8), Inches(2.8),
         "보정 방법 2: 가중 로지스틱 회귀", [
             "각 관측치에 역선택확률 가중치 부여:",
             "  양성(Y=1): w_1 = tau / y_bar",
             "  음성(Y=0): w_0 = (1-tau) / (1-y_bar)",
             "",
             "sklearn 구현:",
             "  LogisticRegression(class_weight='balanced')",
             "  -> 클래스 빈도에 반비례하는 가중치 자동 적용",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# 슬라이드 11: 모델 평가 - 혼동 행렬
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.8", "모델 평가 (1) - 혼동 행렬과 평가 지표")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.4),
         "혼동 행렬 (Confusion Matrix)", [
             "             예측:음성(0)    예측:양성(1)",
             "실제:음성(0)   TN              FP (제1종 오류)",
             "실제:양성(1)   FN (제2종 오류)  TP",
             "",
             "TP: 양성을 양성으로 올바르게 예측",
             "FN: 양성을 음성으로 잘못 예측 (가장 위험할 수 있음)",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.4),
         "주요 평가 지표 수식", [
             "Accuracy  = (TP + TN) / (TP + TN + FP + FN)",
             "Precision = TP / (TP + FP)   '예측의 품질'",
             "Recall    = TP / (TP + FN)   '검출력'",
             "F1 Score  = 2 * (Prec * Rec) / (Prec + Rec)",
             "",
             "MCC = (TP*TN - FP*FN) / sqrt((TP+FP)(TP+FN)(TN+FP)(TN+FN))",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.4),
         "임계값(Threshold) 조정 가이드", [
             "질병 진단: 임계값 낮춤 (0.3) -> Recall 강조",
             "스팸 필터: 임계값 높임 (0.7) -> Precision 강조",
             "사기 탐지: 임계값 낮춤 (0.1) -> Recall 극대화",
             "고객 이탈: 임계값 낮춤 (0.3) -> 이탈 방지 커버리지",
             "",
             "기본 0.5가 항상 최적은 아님!",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(6.9), Inches(4.8), Inches(5.8), Inches(2.4),
         "MCC의 장점 (불균형 데이터)", [
             "범위: [-1, +1]",
             "  +1: 완벽한 예측",
             "   0: 무작위 수준",
             "  -1: 완전히 반대 예측",
             "",
             "불균형 데이터에서 accuracy보다",
             "훨씬 신뢰할 수 있는 단일 지표",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)


# ============================================================
# 슬라이드 12: 모델 평가 - ROC/AUC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.8", "모델 평가 (2) - ROC 곡선, PR 곡선, 적합도 검정")

add_card(slide, Inches(0.6), Inches(2.1), Inches(3.8), Inches(2.6),
         "ROC 곡선과 AUC", [
             "x축: FPR = FP/(FP+TN)",
             "y축: TPR = TP/(TP+FN)",
             "",
             "AUC 해석:",
             "  0.5: 무작위 추측 수준",
             "  0.7~0.8: 수용 가능",
             "  0.8~0.9: 우수한 판별력",
             "  > 0.9: 뛰어난 판별력",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(4.7), Inches(2.1), Inches(3.8), Inches(2.6),
         "PR 곡선 (Precision-Recall)", [
             "x축: Recall",
             "y축: Precision",
             "AP: PR 곡선 아래 면적",
             "",
             "불균형 데이터에서는",
             "ROC보다 PR 곡선이",
             "더 유용할 수 있음",
             "(양성 클래스에 초점)",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(8.8), Inches(2.1), Inches(3.9), Inches(2.6),
         "Hosmer-Lemeshow 적합도 검정", [
             "1. 예측 확률 기준 g개 그룹 분할",
             "2. 관측 vs 기대 사건 수 비교",
             "3. 검정 통계량 C 산출",
             "",
             "H0: 모델이 데이터를 잘 설명",
             "p-value 높으면 적합도 좋음",
             "",
             "Hosmer & Lemeshow(2000)",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2),
         "Menard(2002): 유사 R-squared 비교", [
             "McFadden's R^2: [0, 1), 로그 우도 기반, 최대값이 1에 도달하지 않음",
             "Cox & Snell R^2: [0, <1), 최대값이 1 미만",
             "Nagelkerke R^2: [0, 1], Cox & Snell을 정규화, 최대값 1 가능",
             "Count R^2: [0, 1], 정확 분류 비율, 가장 직관적",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)


# ============================================================
# 슬라이드 13: 전처리 - 상관관계 분석
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.A", "데이터 전처리 (1) - 상관관계 분석", "corr(), heatmap으로 변수 간 관계 파악")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.4),
         "상관계수 (Pearson Correlation)", [
             "+1에 가까움: 강한 양의 상관관계",
             " 0에 가까움: 상관관계 없음",
             "-1에 가까움: 강한 음의 상관관계",
             "",
             "pandas: df.corr()",
             "주의: 수치형 변수만 대상",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.4),
         "Heatmap 시각화 (seaborn)", [
             "sns.heatmap(df.corr(),",
             "    cmap='coolwarm',  # 파랑-빨강",
             "    vmin=-1, vmax=1,  # 범위 고정",
             "    annot=True,       # 수치 표시",
             "    fmt='.2f')        # 소수점 2자리",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.4),
         "타이타닉 데이터: Survived와의 상관계수 해석", [
             "Pclass:  -0.336  ->  등급 숫자가 커질수록(3등실) 생존율 감소",
             "Age:     -0.070  ->  나이가 많을수록 생존율 약간 감소",
             "SibSp:   -0.034  ->  형제/배우자 수와 생존율은 거의 무관",
             "Parch:   +0.083  ->  부모/자녀 수가 많을수록 생존율 약간 증가",
             "",
             "=> 수치형 변수 외에 카테고리 변수(Sex, Embarked)도 함께 고려해야 함",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# 슬라이드 14: 전처리 - 카테고리 변수, 더미 변수
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.A", "데이터 전처리 (2) - 카테고리 변수 & 불필요 변수 제거")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.6),
         "원-핫 인코딩 & 더미 변수", [
             "원-핫 인코딩: 각 고유값을 별도 이진(0/1) 열로 변환",
             "  Embarked -> Embarked_C, Embarked_Q, Embarked_S",
             "",
             "더미 변수 (drop_first=True):",
             "  첫 번째 열 제거 -> 다중공선성 방지",
             "  C는 기준 범주 (Embarked_Q=0, Embarked_S=0)",
             "",
             "pd.get_dummies(df, columns=[...], drop_first=True)",
             "=> 로지스틱 회귀에 권장!",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.6),
         "불필요한 변수 제거 (nunique)", [
             "판단 기준: 고유값(unique) 수",
             "",
             "Name:      889개 -> 제거 (패턴 학습 불가)",
             "Sex:         2개 -> 유지 (이진 변수)",
             "Ticket:    680개 -> 제거 (차원 폭발 위험)",
             "Embarked:    3개 -> 유지 (적은 범주)",
             "",
             "data.drop(['Name', 'Ticket'], axis=1)",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2),
         "다중공선성 (Multicollinearity) 진단 - Menard(2002)", [
             "다중공선성: 독립 변수 간 높은 상관관계 -> 계수 추정 불안정, 해석 어려움",
             "진단 기준: VIF > 10, Tolerance < 0.1, 조건지수 > 30이면 심각한 다중공선성",
             "원-핫 인코딩의 완전 다중공선성: 한 열이 나머지 열로부터 완벽 예측 가능 -> drop_first=True로 해결",
             "Hosmer & Lemeshow(2000): 목적적 변수 선택(7단계)으로 교란변수 확인 (20% 이상 계수 변화 시 재포함)",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)


# ============================================================
# 슬라이드 15: 피처 엔지니어링
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.B", "피처 엔지니어링", "Feature Engineering - 모델 성능의 핵심")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.4),
         "피처 엔지니어링이란?", [
             "기존 데이터를 변환/결합하여",
             "새로운 피처를 생성하는 과정",
             "",
             "중요한 이유:",
             "  1. 좋은 피처 -> 단순한 모델로도 높은 성능",
             "  2. 도메인 지식을 데이터에 반영",
             "  3. 차원 축소 (변수 합치기)",
             "  4. 해석력 향상",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.4),
         "일반적인 피처 엔지니어링 기법", [
             "변수 결합: SibSp + Parch -> family",
             "구간화(Binning): Age -> 유아/청소년/성인/노인",
             "로그 변환: 왜도가 큰 변수에 log 적용",
             "다항 피처: x1 * x2, x1^2 등 상호작용",
             "텍스트 파생: Name -> 호칭(Mr/Mrs/Miss)",
             "",
             "=> '정답'은 없다. 반복 실험으로 최적 피처 탐색",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.4),
         "타이타닉 예시: family = SibSp + Parch", [
             "SibSp: 함께 탑승한 형제자매/배우자 수  |  Parch: 함께 탑승한 부모/자녀 수",
             "두 변수 모두 '함께 탑승한 가족 수'라는 공통된 의미 -> 하나로 합치기",
             "",
             "성능 비교: 피처 엔지니어링 전 78.09%  ->  후 79.21%  (+1.1%p 향상)",
             "data['family'] = data['SibSp'] + data['Parch'] / data.drop(['SibSp', 'Parch'], axis=1)",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# 슬라이드 16: 논문 리뷰 (1)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.9", "논문 리뷰 (1)", "Cox(1958) & Hosmer-Lemeshow(2000)")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.6),
         "Cox (1958) - Regression Analysis of Binary Sequences", [
             "로지스틱 회귀의 이론적 토대를 정립한 기념비적 논문",
             "",
             "핵심 기여:",
             "  - 로짓 변환의 체계적 도입",
             "    log(p/(1-p)) = beta_0 + beta_1*x_1 + ...",
             "  - MLE 적용 방법론과 점근적 성질 체계화",
             "  - 프로빗 모델 대비 오즈비 해석의 장점 강조",
             "  - 뉴턴-랩슨 최적화 방법 제안",
             "",
             "=> 왜 시그모이드를 쓰는지, 왜 로그 손실을 쓰는지,",
             "   왜 오즈비로 계수를 해석하는지의 근본적 이유",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(4.6),
         "Hosmer & Lemeshow (2000) - Applied Logistic Regression", [
             "가장 권위 있는 실용적 로지스틱 회귀 교과서",
             "",
             "핵심 기여:",
             "  - 체계적 모델 구축 프레임워크:",
             "    변수 선택 -> 모델 적합 -> 진단 -> 평가",
             "  - 목적적 변수 선택 (7단계 과정)",
             "    1) 단변량 분석 (p < 0.25)",
             "    2) 다변량 모델 적합",
             "    3) 유의하지 않은 변수 제거 (p > 0.05)",
             "    4) 교란변수 확인 (20% 계수 변화)",
             "  - Hosmer-Lemeshow 적합도 검정 표준화",
             "  - 진단 통계량: 잔차, 레버리지, Cook's distance",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# 슬라이드 17: 논문 리뷰 (2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.9", "논문 리뷰 (2)", "King & Zeng(2001), Ng & Jordan(2002), Menard(2002)")

add_card(slide, Inches(0.6), Inches(2.1), Inches(3.8), Inches(4.6),
         "King & Zeng (2001)", [
             "Logistic Regression in",
             "Rare Events Data",
             "",
             "- 희소 사건에서 MLE",
             "  편향 문제 분석",
             "- 확률 과소추정 이론적 증명",
             "- 사전 교정 방법 제시",
             "- 가중 로지스틱 회귀 제안",
             "- Firth(1993) 편향 축소",
             "  방법 적용",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(slide, Inches(4.7), Inches(2.1), Inches(3.8), Inches(4.6),
         "Ng & Jordan (2002)", [
             "Discriminative vs",
             "Generative Classifiers",
             "",
             "- 판별/생성 모델의 근본적",
             "  차이 규명 (NIPS)",
             "- 나이브 베이즈와 로지스틱",
             "  회귀: 동일 모델 패밀리",
             "- O(log n) vs O(n) 수렴",
             "- 전환점 존재 증명",
             "- 편향-분산 트레이드오프",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(8.8), Inches(2.1), Inches(3.9), Inches(4.6),
         "Menard (2002)", [
             "Applied Logistic",
             "Regression Analysis",
             "",
             "- OLS와의 체계적 비교",
             "- 다양한 계수 해석법:",
             "  로짓, 오즈비, 표준화 계수,",
             "  한계 효과 비교",
             "- 유사 R-squared 정리:",
             "  McFadden, Cox & Snell,",
             "  Nagelkerke, Count R^2",
             "- 다중공선성 진단 기준 정리",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)


# ============================================================
# 슬라이드 18: 실습 소개 - 스크래치 구현
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.10", "실습: 로지스틱 회귀 스크래치 구현", "Breast Cancer 데이터셋 / Cox(1958) 기반")

add_card(slide, Inches(0.6), Inches(2.1), Inches(3.8), Inches(2.4),
         "구현 내용", [
             "1. sigmoid() 함수 구현",
             "2. compute_cost() 로그 손실",
             "3. compute_gradients() 그래디언트",
             "4. logistic_regression_gd()",
             "   경사하강법 학습 루프",
             "5. predict() 예측 함수",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(4.7), Inches(2.1), Inches(3.8), Inches(2.4),
         "핵심 구현 포인트", [
             "수치 안정성:",
             "  np.clip(z, -500, 500)",
             "  epsilon = 1e-15 (log(0) 방지)",
             "",
             "벡터화 연산:",
             "  dw = (1/m) * X^T @ error",
             "  -> for 루프 없이 효율적 계산",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(8.8), Inches(2.1), Inches(3.9), Inches(2.4),
         "sklearn 비교 결과", [
             "직접 구현 (경사하강법):",
             "  -> ~97% 정확도 달성",
             "",
             "sklearn (L-BFGS):",
             "  -> ~97% 정확도 달성",
             "",
             "차이: 최적화 알고리즘 차이",
             "      + L2 정규화 유무",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.4),
         "시각화 내용", [
             "시그모이드 함수 곡선  |  로그 손실 함수 곡선  |  비용 함수 수렴 곡선  |  2D 결정 경계 (PCA)",
             "",
             "학습률(0.01, 0.05, 0.1, 0.5, 1.0)에 따른 수렴 속도 비교",
             "30차원 -> PCA 2차원 축소 후 결정 경계 비교 (직접 구현 vs sklearn)",
             "=> 특성 표준화(StandardScaler) 필수: 경사하강법 수렴을 위해",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)


# ============================================================
# 슬라이드 19: 실습 소개 - 판별 vs 생성 비교
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.11-12", "실습: 판별 vs 생성 비교 & 희귀 이벤트 보정")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.4),
         "실습 5.11: Ng & Jordan(2002) 재현", [
             "학습 데이터 크기 1%~100% 변화시키며",
             "로지스틱 회귀 vs 나이브 베이즈 비교",
             "",
             "3가지 실험:",
             "  1. 독립적 특성 (NB 가정 성립)",
             "  2. 상관된 특성 (NB 가정 위반)",
             "  3. 고차원 데이터 (50개 특성)",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.4),
         "실습 5.12: King & Zeng(2001) 재현", [
             "불균형 비율별 MLE 편향 분석:",
             "  50% -> 20% -> 10% -> 5% -> 2% -> 1%",
             "",
             "3가지 모델 비교:",
             "  1. 표준 로지스틱 회귀",
             "  2. 가중 LR (class_weight='balanced')",
             "  3. 사전 교정 (Prior Correction)",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(slide, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.4),
         "핵심 관찰 포인트 (판별 vs 생성)", [
             "전환점(crossover) 확인:",
             "  소표본 -> 나이브 베이즈 우수",
             "  대표본 -> 로지스틱 회귀 우수",
             "",
             "NB 가정 위반 시:",
             "  전환점이 더 일찍 나타남",
             "  로지스틱 회귀가 빠르게 역전",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(4.8), Inches(5.8), Inches(2.4),
         "핵심 관찰 포인트 (희귀 이벤트)", [
             "양성 비율 낮을수록 과소추정 심화",
             "  1% 양성 -> 수십% 과소추정",
             "",
             "ROC 곡선 vs PR 곡선 비교:",
             "  불균형 데이터 -> PR 곡선이 더 유용",
             "",
             "임계값 0.5에서 양성 거의 검출 못함",
             "  -> 임계값을 양성 비율에 맞춤",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# 슬라이드 20: 응용사례 - 타이타닉 실습
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.C", "실습: 타이타닉 생존자 예측", "전처리 -> 모델링 -> 평가 -> 피처 엔지니어링 전체 파이프라인")

add_card(slide, Inches(0.6), Inches(2.1), Inches(3.0), Inches(2.2),
         "Step 1: 데이터 탐색", [
             "data.head()",
             "data.info()",
             "data.describe()",
             "data.corr() + heatmap",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(3.8), Inches(2.1), Inches(3.0), Inches(2.2),
         "Step 2: 전처리", [
             "nunique()로 고유값 확인",
             "Name, Ticket 제거",
             "get_dummies(",
             "  drop_first=True)",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(7.0), Inches(2.1), Inches(3.0), Inches(2.2),
         "Step 3: 모델링", [
             "train_test_split(0.8/0.2)",
             "LogisticRegression()",
             "model.fit(X_train, y_train)",
             "pred = model.predict(X_test)",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(10.2), Inches(2.1), Inches(2.5), Inches(2.2),
         "Step 4: 평가", [
             "accuracy_score()",
             "predict_proba()",
             "coef_[0] 계수",
             "Odds Ratio",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(slide, Inches(0.6), Inches(4.6), Inches(5.8), Inches(2.6),
         "sklearn 핵심 코드", [
             "from sklearn.linear_model import LogisticRegression",
             "from sklearn.model_selection import train_test_split",
             "from sklearn.metrics import accuracy_score",
             "",
             "X_train, X_test, y_train, y_test = train_test_split(",
             "    X, y, test_size=0.2, random_state=100)",
             "model = LogisticRegression()",
             "model.fit(X_train, y_train)",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(4.6), Inches(5.8), Inches(2.6),
         "계수 해석 주의사항", [
             "model.coef_ -> 2차원 배열 shape=(1, n_features)",
             "  -> model.coef_[0]으로 1차원 추출 필요",
             "model.intercept_ -> 절편(bias)",
             "",
             "계수 양수: 피처 증가 시 생존 확률 증가",
             "계수 음수: 피처 증가 시 생존 확률 감소",
             "절대값 클수록 영향력 큼",
             "오즈비: np.exp(model.coef_[0])",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# 슬라이드 21: 응용사례 - 스팸 필터링
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.13", "응용사례 (1) - 이메일 스팸 필터링 & 사기 탐지")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.4),
         "스팸 필터링 시스템", [
             "이메일 수신 -> 피처 추출 -> LR 모델 -> P(스팸) -> 분류",
             "",
             "주요 피처: keyword_freq, capital_ratio,",
             "  special_char_count, link_count, sender_score",
             "",
             "핵심: 정상 메일을 스팸으로 분류(FP)가 치명적",
             "  -> 임계값 높이거나 Precision 중시",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.4),
         "사기 탐지 (Fraud Detection)", [
             "양성 비율 극히 낮음 (~0.1%)",
             "  -> King & Zeng(2001) 희소 사건",
             "실시간 처리 필요",
             "비대칭 비용: 사기 놓침 >> 정상 차단",
             "",
             "가중 LR: class_weight={0:1, 1:100}",
             "평가: PR-AUC, Recall, F1 필수",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(slide, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.4),
         "질병 진단 (Disease Diagnosis)", [
             "오진 비용 매우 높음 (특히 FN)",
             "해석 가능성 중요: 의사가 근거 이해",
             "",
             "로지스틱 회귀의 장점:",
             "  - 오즈비로 위험 인자 정량화",
             "  - '발생 확률 73%' 직관적 소통",
             "  - Cox(1958) 기반 의학 연구 표준",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(6.9), Inches(4.8), Inches(5.8), Inches(2.4),
         "고객 이탈 예측 (Customer Churn)", [
             "이탈률: 보통 15~30%",
             "FN 비용 > FP 비용 (이탈 놓침이 더 위험)",
             "신규 유치 비용 = 유지 비용의 5~7배",
             "",
             "전략: 예측 확률 상위 20% 고객에게",
             "  리텐션 마케팅 (쿠폰, 할인)",
             "임계값 0.3으로 낮춰 Recall 향상",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# 슬라이드 22: 고객 이탈 - Confusion Matrix 분석
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.13", "응용사례 - 고객 이탈 Confusion Matrix 분석")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.4),
         "혼동 행렬 비즈니스 해석", [
             "TN (유지->유지): 정상 (대응 불필요)",
             "FP (유지->이탈): 불필요한 마케팅 비용 발생",
             "FN (이탈->유지): 이탈 방지 실패 (가장 위험!)",
             "TP (이탈->이탈): 리텐션 마케팅 대상 (성공)",
             "",
             "FN의 비용: 미래 수익 손실 (수백만원 LTV)",
             "FP의 비용: 불필요한 쿠폰 (수만원)",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.4),
         "이탈 예측에서 Recall이 중요한 이유", [
             "Precision: 이탈 예측 중 실제 이탈 비율",
             "  -> 마케팅 비용 효율성 측정",
             "",
             "Recall: 실제 이탈 중 식별된 비율",
             "  -> 이탈 방지 커버리지 측정",
             "",
             "=> Recall을 높이는 방향으로 조정",
             "   (임계값을 낮추는 방향)",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.4),
         "실습 코드 하이라이트", [
             "confusion_matrix(y_test, pred)  ->  tn, fp, fn, tp = cm.ravel()",
             "classification_report(y_test, pred, target_names=['유지', '이탈'])",
             "ConfusionMatrixDisplay.from_predictions(y_test, pred, normalize='true')",
             "",
             "피처 중요도: coef_churn = pd.Series(model.coef_[0], index=X.columns)",
             "  양수 계수 -> 이탈 확률 증가 요인  /  음수 계수 -> 이탈 확률 감소 요인",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)


# ============================================================
# 슬라이드 23: 핵심 요약
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.14", "핵심 요약", "5장 전체 내용 정리")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.3),
         "로지스틱 회귀 알고리즘", [
             "유형: 지도학습 - 분류 (Classification)",
             "핵심 함수: sigma(z) = 1/(1+e^(-z))",
             "비용 함수: Log Loss (Binary Cross Entropy)",
             "최적화: 경사하강법 or L-BFGS (준뉴턴법)",
             "출력: 클래스 확률 (0~1)",
             "장점: 간단, 확률 해석, 계수 해석, 효율적",
             "단점: 선형 결정 경계만 가능",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.3),
         "5개 논문의 핵심 메시지", [
             "Cox(1958): 로짓+MLE = LR의 수학적 기초",
             "Hosmer(2000): 체계적 모델 구축/진단/평가",
             "King(2001): 희소 사건 MLE 편향, 보정 필수",
             "Ng(2002): 소표본=생성, 대표본=판별 모델",
             "Menard(2002): OLS 비교, 다양한 계수 해석법",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.5),
         "전처리 체크리스트", [
             "1. 상관관계 분석: corr(), heatmap",
             "2. 고유값 많은 불필요 변수 제거: nunique(), drop()",
             "3. 카테고리 -> 더미 변수: get_dummies(drop_first=True)",
             "4. 피처 엔지니어링: 도메인 지식 활용",
             "5. 특성 표준화: StandardScaler (경사하강법 시)",
             "6. 다중공선성 확인: VIF > 10 주의",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(6.9), Inches(4.7), Inches(5.8), Inches(2.5),
         "모델 평가 체크리스트", [
             "기본: accuracy_score, confusion_matrix",
             "분류 보고서: classification_report",
             "곡선: ROC/AUC, PR 곡선/AP",
             "불균형 데이터: MCC, F1, Recall 중시",
             "임계값: 문제 특성에 맞게 조정",
             "적합도: Hosmer-Lemeshow 검정",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)


# ============================================================
# 슬라이드 24: 핵심 수식 정리
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.14", "핵심 수식 정리")

formulas = [
    ("시그모이드 함수",       "sigma(z) = 1 / (1 + e^(-z))", ACCENT_CYAN),
    ("시그모이드 미분",       "sigma'(z) = sigma(z) * (1 - sigma(z))", ACCENT_CYAN),
    ("로짓 변환",            "logit(p) = ln(p / (1-p)) = w^T x + b", ACCENT_GREEN),
    ("오즈비",               "OR = e^(w_i)", ACCENT_GREEN),
    ("이진 교차 엔트로피",    "J = -(1/n) SUM [y*ln(y_hat) + (1-y)*ln(1-y_hat)]", ACCENT_ORANGE),
    ("그래디언트",            "grad_w J = (1/n) X^T (y_hat - y)", ACCENT_ORANGE),
    ("소프트맥스",            "P(y=k|X) = e^(z_k) / SUM_j e^(z_j)", ACCENT_PURPLE),
    ("사전 교정 (King)",     "beta_0_c = beta_0 - ln[(1-tau)/tau * y_bar/(1-y_bar)]", ACCENT_RED),
    ("F1 Score",             "F1 = 2 * Precision * Recall / (Precision + Recall)", ACCENT_BLUE),
    ("MCC",                  "MCC = (TP*TN - FP*FN) / sqrt((TP+FP)(TP+FN)(TN+FP)(TN+FN))", ACCENT_BLUE),
]

y_pos = Inches(2.1)
for name, formula, color in formulas:
    add_shape(slide, Inches(0.6), y_pos, Inches(3.0), Inches(0.42), CARD_BG, color, radius=True)
    add_text(slide, Inches(0.8), y_pos + Inches(0.03), Inches(2.6), Inches(0.36),
             name, font_size=12, color=color, bold=True)
    add_text(slide, Inches(3.8), y_pos + Inches(0.03), Inches(9.0), Inches(0.36),
             formula, font_size=14, color=WHITE, font_name='Consolas')
    y_pos += Inches(0.50)


# ============================================================
# 슬라이드 25: sklearn 핵심 코드 요약
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.14", "sklearn 핵심 코드 요약")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.8),
         "기본 학습 & 예측 코드", [
             "from sklearn.linear_model import LogisticRegression",
             "from sklearn.model_selection import train_test_split",
             "from sklearn.metrics import accuracy_score,",
             "    classification_report, roc_auc_score",
             "",
             "# 데이터 분리",
             "X_train, X_test, y_train, y_test = train_test_split(",
             "    X, y, test_size=0.2, random_state=42)",
             "",
             "# 모델 학습",
             "model = LogisticRegression(max_iter=1000)",
             "model.fit(X_train, y_train)",
             "",
             "# 예측",
             "pred = model.predict(X_test)        # 클래스",
             "prob = model.predict_proba(X_test)   # 확률",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(4.8),
         "평가 & 계수 해석 코드", [
             "# 평가",
             "accuracy_score(y_test, pred)",
             "roc_auc_score(y_test, prob[:, 1])",
             "classification_report(y_test, pred)",
             "",
             "# 계수 확인",
             "model.coef_[0]        # 피처별 계수",
             "model.intercept_       # 절편",
             "np.exp(model.coef_[0]) # 오즈비",
             "",
             "# 불균형 데이터",
             "LogisticRegression(class_weight='balanced')",
             "",
             "# 다중 분류",
             "LogisticRegression(",
             "    multi_class='multinomial',",
             "    solver='lbfgs')",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# 슬라이드 26: 주요 용어 정리
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.14", "주요 용어 정리")

terms_left = [
    "로지스틱 회귀 (Logistic Regression) - 분류를 위한 선형 모델",
    "시그모이드 함수 (Sigmoid) - 실수를 0~1 확률로 변환",
    "임계값 (Threshold) - 확률을 클래스로 변환하는 기준값",
    "로그 손실 (Log Loss) - 로지스틱 회귀의 비용 함수",
    "이진 교차 엔트로피 (Binary CE) - 로그 손실의 다른 이름",
    "오즈 (Odds) - 발생/미발생 확률의 비",
    "오즈비 (Odds Ratio) - 피처 변화에 따른 오즈 변화량",
    "로짓 (Logit) - 오즈의 자연로그, 시그모이드 역함수",
]

terms_right = [
    "소프트맥스 (Softmax) - 다중 분류용 활성화 함수",
    "더미 변수 (Dummy Variable) - 카테고리를 이진 열로 변환",
    "판별 모델 (Discriminative) - P(y|x) 직접 학습",
    "생성 모델 (Generative) - P(x|y)P(y) 학습, 베이즈 정리",
    "혼동 행렬 (Confusion Matrix) - 예측 결과 세부 분류표",
    "정밀도 (Precision) - 양성 예측 중 실제 양성 비율",
    "재현율 (Recall) - 실제 양성 중 양성 예측 비율",
    "MCC - 불균형에 강건한 단일 평가 지표 [-1, +1]",
]

add_bullet_list(slide, Inches(0.6), Inches(2.1), Inches(6.0), Inches(5.0),
                terms_left, font_size=14, color=LIGHT_GRAY, spacing=Pt(10))
add_bullet_list(slide, Inches(6.9), Inches(2.1), Inches(6.0), Inches(5.0),
                terms_right, font_size=14, color=LIGHT_GRAY, spacing=Pt(10))


# ============================================================
# 슬라이드 27: 복습 질문
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "5.14", "복습 질문", "10개 핵심 질문으로 이해도 점검")

questions_left = [
    "Q1. 선형 회귀로 이진 분류 시 발생하는 4가지 문제와 로지스틱 회귀의 해결법은?",
    "Q2. 시그모이드 미분 sigma'(z) = sigma(z)(1-sigma(z))을 증명하시오.",
    "Q3. 오즈비(OR) = 1.5의 의미를 타이타닉 예시로 설명하시오.",
    "Q4. MLE 우도 최대화 = 로그 손실 최소화의 동치성을 보이시오.",
    "Q5. Ng(2002): O(log n) vs O(n) 수렴이 모델 선택에 주는 시사점은?",
]

questions_right = [
    "Q6. King(2001): 양성 비율 1% 데이터에서 MLE 편향과 2가지 교정법?",
    "Q7. 정확도 98%인 사기 탐지 모델(양성 2%)이 좋은 모델인가?",
    "Q8. K=2일 때 소프트맥스가 시그모이드로 환원됨을 보이시오.",
    "Q9. Hosmer-Lemeshow 적합도 검정의 절차와 귀무가설?",
    "Q10. 직접 구현 LR과 sklearn LR의 성능 차이 원인 2가지?",
]

add_bullet_list(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(5.0),
                questions_left, font_size=14, color=LIGHT_GRAY, spacing=Pt(14))
add_bullet_list(slide, Inches(6.9), Inches(2.3), Inches(5.8), Inches(5.0),
                questions_right, font_size=14, color=LIGHT_GRAY, spacing=Pt(14))


# ============================================================
# 슬라이드 28: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)

# 상단 장식선
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)

add_text(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(0.8),
         "Thank You", font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.0), Inches(3.2), Inches(3.3), ACCENT_CYAN)

add_text(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.6),
         "5장: 로지스틱 회귀 (Logistic Regression)", font_size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
         "다음 장: 6장 - KNN (K-Nearest Neighbors)", font_size=20, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
         "핵심 논문: Cox(1958), Hosmer & Lemeshow(2000), King & Zeng(2001), Ng & Jordan(2002), Menard(2002)",
         font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# 저장
# ============================================================
output_path = r"D:\26년1학기\기계학습\5장\5장_로지스틱회귀_강의PPT.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
