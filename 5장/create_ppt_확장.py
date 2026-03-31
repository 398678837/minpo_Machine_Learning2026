"""5장 로지스틱 회귀 - 확장 강의 PPT 생성 스크립트 (상세 버전)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = '맑은 고딕'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    """테이블을 카드 형태로 그리는 헬퍼"""
    total_w = sum(col_widths)
    # 헤더
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    # 행
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    """섹션 구분 슬라이드"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), accent)
    add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), accent)
    add_text(s, Inches(0), Inches(2.0), prs.slide_width, Inches(0.5),
             f"SECTION {section_num}", font_size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, Inches(5.5), Inches(2.7), Inches(2.3), accent)
    add_text(s, Inches(0), Inches(3.0), prs.slide_width, Inches(1.0),
             title, font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.5),
             subtitle, font_size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    return s


# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "기계학습 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "5장: 로지스틱 회귀", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "Logistic Regression", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.3), Inches(11), Inches(0.8),
         "[ 확장 상세 버전 ]", font_size=20, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
         "핵심 키워드: 분류 · 시그모이드 · 로그 손실 · 오즈비 · 소프트맥스 · 판별 모델 · 희귀 이벤트",
         font_size=14, color=DARK_GRAY)
add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
         "참고 논문: Cox(1958), Hosmer & Lemeshow(2000), King & Zeng(2001), Ng & Jordan(2002), Menard(2002)",
         font_size=13, color=DARK_GRAY)
add_shape(s, Inches(0), Inches(6.3), prs.slide_width, Inches(1.2), RGBColor(0x0A, 0x0A, 0x1A))
add_text(s, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
         "2026년 1학기", font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 2: 목차 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (1/2)", "Table of Contents")

toc_left = [
    "5.1  분류 문제의 본질 - 왜 선형회귀로 분류하면 안 되는지",
    "5.2  시그모이드와 로짓 - 오즈비, 로짓 변환",
    "5.3  MLE와 크로스엔트로피 - 우도함수, 로그우도, 경사하강법",
    "5.4  경사하강법 - 배치/미니배치/SGD 비교, 학습률 선택",
    "5.5  다중 클래스 - 소프트맥스, OvR, OvO 전략",
    "5.6  판별 vs 생성 모델 - Ng&Jordan(2002) 기반 비교",
    "5.7  희귀 이벤트 - King&Zeng(2001) 불균형 데이터 편향 보정",
]
add_bullet_list(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5.0),
                toc_left, font_size=18, color=LIGHT_GRAY, spacing=Pt(14))

# ============================================================
# 슬라이드 3: 목차 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (2/2)", "Table of Contents")

toc_right = [
    "5.8   모델 평가 - ROC/AUC, 혼동행렬, Precision/Recall/F1/MCC",
    "5.9   논문 리뷰 - Cox(1958), Hosmer&Lemeshow, Menard 등",
    "5.10  실습: 로지스틱 회귀 스크래치 구현",
    "5.11  실습: 판별 vs 생성 비교",
    "5.12  실습: 희귀 이벤트 보정",
    "5.13  응용사례 - 사기탐지, 질병진단, 고객이탈",
    "5.14  핵심 요약 + 복습 질문 10개",
]
add_bullet_list(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5.0),
                toc_right, font_size=18, color=LIGHT_GRAY, spacing=Pt(14))

# ============================================================
# 슬라이드 4: 학습 목표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "학습 목표", "Learning Objectives")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "이론적 목표", [
             "로지스틱 회귀의 수학적 원리 (시그모이드, 로짓, MLE) 이해",
             "로그 손실(크로스엔트로피)과 그래디언트 유도 과정 이해",
             "판별 모델 vs 생성 모델의 이론적 비교 (Ng & Jordan)",
             "희소 사건에서의 MLE 편향과 보정 방법 (King & Zeng)",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "실습적 목표", [
             "로지스틱 회귀를 밑바닥부터(scratch) 구현",
             "sklearn LogisticRegression 활용 및 성능 비교",
             "ROC/AUC, 혼동행렬 등 분류 모델 평가 수행",
             "사기탐지/질병진단/고객이탈 응용사례 분석",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.5), Inches(4.8), Inches(12.1), Inches(2.0),
         "참고 논문 5편", [
             "Cox(1958) - 로지스틱 회귀의 수학적 기초 정립",
             "Hosmer & Lemeshow(2000) - 실용적 모델 구축 프레임워크",
             "King & Zeng(2001) - 희소 사건 편향 분석 및 보정",
             "Ng & Jordan(2002) - 판별 vs 생성 모델 비교",
             "Menard(2002) - OLS 비교를 통한 로지스틱 회귀 입문",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 5: 섹션 1 - 분류 문제의 본질
# ============================================================
section_divider("분류 문제의 본질", "왜 선형회귀로 분류하면 안 되는지", "5.1", ACCENT_BLUE)

# ============================================================
# 슬라이드 6: 회귀 vs 분류
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.1", "회귀와 분류의 근본적 차이")

add_table_slide(s,
    ["구분", "선형 회귀", "로지스틱 회귀"],
    [
        ["문제 유형", "회귀 (Regression)", "분류 (Classification)"],
        ["출력값", "연속 수치 (-inf ~ +inf)", "클래스 확률 (0~1)"],
        ["예시", "주택 가격, 보험료 예측", "생존/사망, 스팸/정상"],
        ["활성화 함수", "없음 (항등 함수)", "시그모이드 (Sigmoid)"],
        ["비용 함수", "MSE (평균제곱오차)", "Log Loss (로그 손실)"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 4.5, 4.5])

add_text(s, Inches(0.6), Inches(5.5), Inches(11), Inches(0.5),
         "* 로지스틱 회귀는 이름에 '회귀'가 있지만 실제로는 분류(Classification) 알고리즘이다.",
         font_size=14, color=ACCENT_ORANGE)

# ============================================================
# 슬라이드 7: 선형 회귀로 분류하면 발생하는 문제
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.1", "선형 회귀로 분류 시 문제점 (Menard, 2002)")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3),
         "문제 1: 확률 범위 위반", [
             "선형 회귀 출력: y_hat = w^T x + b",
             "범위: (-inf, +inf) -> 확률 [0,1] 위반",
             "예측값이 0 미만 또는 1 초과 가능",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.3),
         "문제 2: 이분산성 (Heteroscedasticity)", [
             "이진 종속 변수의 오차 분산: y_hat(1-y_hat)",
             "예측값에 따라 분산이 달라짐",
             "OLS의 등분산성 가정 위배",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.3),
         "문제 3: 비정규 오차", [
             "실제값이 0 또는 1뿐 -> 잔차가 이산적",
             "정규성 가정 완전히 위배",
             "t-검정, F-검정 등 통계 추론 부적절",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.3),
         "문제 4: 비선형 관계", [
             "실제 확률과 설명 변수: S자 곡선(S-curve)",
             "0이나 1에 가까울수록 포화(saturation) 현상",
             "직선으로는 이 관계를 표현할 수 없음",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 8: 로지스틱 회귀의 핵심 아이디어
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.1", "로지스틱 회귀의 핵심 아이디어 - Cox(1958)")

add_text(s, Inches(0.6), Inches(2.0), Inches(11), Inches(0.5),
         "시그모이드 함수로 모든 문제를 해결!", font_size=22, color=ACCENT_CYAN, bold=True)

add_shape(s, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.2), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(1.0), Inches(3.0), Inches(11), Inches(0.8),
         "입력(X)  -->  선형 결합(z = w^T x + b)  -->  시그모이드(sigma)  -->  확률(P)  -->  클래스(0 또는 1)",
         font_size=18, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_card(s, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.5),
         "시그모이드의 역할", [
             "(-inf, +inf) 범위를 (0, 1)로 압축",
             "확률 범위 위반 문제 해결",
             "S자 곡선으로 비선형 관계 표현",
             "미분이 깔끔하여 경사하강법에 적합",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.5),
         "Cox(1958)의 기여", [
             "'The Regression Analysis of Binary Sequences'",
             "로짓 변환의 체계적 도입",
             "MLE 기반 파라미터 추정 정립",
             "이 논문이 로지스틱 회귀의 이론적 토대",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 9: 섹션 2 - 시그모이드와 로짓
# ============================================================
section_divider("시그모이드와 로짓", "오즈비, 로짓 변환", "5.2", ACCENT_CYAN)

# ============================================================
# 슬라이드 10: 시그모이드 함수
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.2", "시그모이드 함수 (Sigmoid Function)")

add_shape(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.0), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(1.0), Inches(2.1), Inches(5), Inches(0.8),
         "sigma(z) = 1 / (1 + e^(-z)),   z = w^T x + b",
         font_size=20, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_table_slide(s,
    ["입력값 z", "sigma(z)", "해석"],
    [
        ["z -> +inf", "1에 가까움", "양성 클래스(1) 확률 높음"],
        ["z = 0", "0.5", "결정 경계 (Decision Boundary)"],
        ["z -> -inf", "0에 가까움", "음성 클래스(0) 확률 높음"],
    ],
    left=Inches(0.6), top=Inches(3.3), col_widths=[3.0, 3.0, 5.5])

add_shape(s, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.2), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(1.0), Inches(5.6), Inches(5), Inches(0.4),
         "미분 공식 (매우 중요!):", font_size=14, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(1.0), Inches(6.0), Inches(5), Inches(0.5),
         "sigma'(z) = sigma(z) * (1 - sigma(z))",
         font_size=18, color=ACCENT_GREEN, bold=True, font_name='Consolas')

# ============================================================
# 슬라이드 11: 오즈와 오즈비
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.2", "오즈(Odds)와 오즈비(Odds Ratio)")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.0),
         "오즈(Odds) 정의", [
             "Odds = P(사건 발생) / P(사건 미발생) = p / (1-p)",
             "생존 p=0.75 -> Odds = 3 ('생존 확률이 3배')",
             "합격 p=0.50 -> Odds = 1 ('동일')",
             "발병 p=0.20 -> Odds = 0.25",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "오즈비(Odds Ratio) = e^(w_i)", [
             "피처가 1단위 증가 시 오즈 변화 배수",
             "w_i > 0 -> OR > 1 -> 양성 확률 증가",
             "w_i = 0 -> OR = 1 -> 영향 없음",
             "w_i < 0 -> OR < 1 -> 양성 확률 감소",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.5), Inches(4.3), Inches(12.1), Inches(2.5),
         "의료 사례 - 흡연과 폐암 / 타이타닉", [
             "[흡연] 흡연자 폐암 Odds = 0.136, 비흡연자 = 0.0101 -> OR = 13.5배",
             "[타이타닉] Sex_male 계수 = -2.569 -> OR = e^(-2.569) = 0.077",
             "  => 남성은 여성보다 생존 오즈가 약 92.3% 낮다",
             "[타이타닉] Pclass 계수 = -1.182 -> 등급 높아질수록 생존 오즈 감소",
             "Menard(2002): 오즈비 외에 표준화 로짓 계수, 한계 효과 등 다양한 해석법 비교",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 12: 로짓 변환
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.2", "로짓 변환 (Logit Transformation)")

add_shape(s, Inches(0.8), Inches(2.0), Inches(7.0), Inches(0.9), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(1.0), Inches(2.1), Inches(6.5), Inches(0.7),
         "logit(p) = ln(p / (1-p)) = z = w^T x + b",
         font_size=20, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_text(s, Inches(0.6), Inches(3.2), Inches(11), Inches(0.5),
         "범위의 일치 - 로짓 변환의 핵심적 의의", font_size=20, color=ACCENT_CYAN, bold=True)

add_table_slide(s,
    ["변환 단계", "범위", "의미"],
    [
        ["확률 p", "(0, 1)", "원래 확률 공간"],
        ["오즈 p/(1-p)", "(0, +inf)", "비율 공간"],
        ["로짓 ln[p/(1-p)]", "(-inf, +inf)", "선형 결합 가능!"],
        ["선형 결합 z", "(-inf, +inf)", "w^T x + b와 동일 범위"],
    ],
    left=Inches(0.6), top=Inches(3.8), col_widths=[3.5, 3.0, 5.0])

add_text(s, Inches(0.6), Inches(6.3), Inches(11), Inches(0.5),
         "핵심: '확률을 로짓 변환하면 피처들의 선형 결합과 같다' - Cox(1958)",
         font_size=16, color=ACCENT_ORANGE, bold=True)

# ============================================================
# 슬라이드 13: 섹션 3 - MLE와 크로스엔트로피
# ============================================================
section_divider("MLE와 크로스엔트로피", "우도함수, 로그우도, 경사하강법", "5.3", ACCENT_GREEN)

# ============================================================
# 슬라이드 14: 왜 MSE를 쓰면 안 되는가
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.3", "왜 MSE를 비용 함수로 쓰면 안 되는가?")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5),
         "MSE + 시그모이드 = 비볼록(Non-convex)", [
             "시그모이드의 비선형성 때문에",
             "MSE 비용 함수가 비볼록 함수가 됨",
             "",
             "비볼록 함수의 문제:",
             "  - 지역 최솟값(local minimum)이 다수 존재",
             "  - 경사하강법으로 전역 최솟값 보장 불가",
             "  - 초기값에 따라 다른 해에 수렴",
             "",
             "해결: 로그 손실(Log Loss) 사용",
             "  -> 볼록(convex) 함수가 되어",
             "  -> 전역 최솟값을 보장",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5),
         "로그 손실의 직관적 이해", [
             "실제=1, 예측=0.9 -> 손실 = -ln(0.9) = 0.105 (낮음)",
             "실제=1, 예측=0.1 -> 손실 = -ln(0.1) = 2.303 (높음!)",
             "실제=0, 예측=0.1 -> 손실 = -ln(0.9) = 0.105 (낮음)",
             "실제=0, 예측=0.9 -> 손실 = -ln(0.1) = 2.303 (높음!)",
             "",
             "핵심 특성:",
             "  - 예측이 틀릴수록 손실이 기하급수적 증가",
             "  - 확신에 찬 잘못된 예측에 큰 페널티",
             "  - 올바른 예측에는 낮은 손실",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 15: MLE (최대우도추정)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.3", "최대우도추정 (MLE)의 원리")

add_text(s, Inches(0.6), Inches(2.0), Inches(11), Inches(0.4),
         "단일 관측치의 확률:", font_size=16, color=ACCENT_CYAN, bold=True)
add_shape(s, Inches(0.8), Inches(2.5), Inches(7.0), Inches(0.7), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(1.0), Inches(2.6), Inches(6.5), Inches(0.5),
         "P(y_i | x_i; w) = y_hat_i^(y_i) * (1 - y_hat_i)^(1-y_i)",
         font_size=17, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_bullet_list(s, Inches(0.8), Inches(3.4), Inches(5.0), Inches(0.8), [
    "y_i=1이면: P = y_hat_i (양성 확률 자체)",
    "y_i=0이면: P = 1 - y_hat_i (음성 확률)",
], font_size=14, color=LIGHT_GRAY)

add_text(s, Inches(0.6), Inches(4.3), Inches(11), Inches(0.4),
         "우도 함수 (Likelihood Function):", font_size=16, color=ACCENT_CYAN, bold=True)
add_shape(s, Inches(0.8), Inches(4.8), Inches(8.0), Inches(0.7), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(1.0), Inches(4.9), Inches(7.5), Inches(0.5),
         "L(w) = PROD_i [y_hat_i^(y_i) * (1 - y_hat_i)^(1-y_i)]",
         font_size=17, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_text(s, Inches(0.6), Inches(5.8), Inches(11), Inches(0.4),
         "로그 우도 (Log-Likelihood):", font_size=16, color=ACCENT_CYAN, bold=True)
add_shape(s, Inches(0.8), Inches(6.2), Inches(10.0), Inches(0.7), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(1.0), Inches(6.3), Inches(9.5), Inches(0.5),
         "l(w) = SUM_i [y_i * ln(y_hat_i) + (1-y_i) * ln(1-y_hat_i)]",
         font_size=17, color=ACCENT_GREEN, bold=True, font_name='Consolas')

# ============================================================
# 슬라이드 16: 로그 우도에서 크로스엔트로피로
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.3", "로그 우도 -> 이진 교차 엔트로피 (Log Loss)")

add_text(s, Inches(0.6), Inches(2.0), Inches(11), Inches(0.4),
         "MLE: 로그 우도를 최대화 -> 최적화 관행상 최소화 문제로 변환 (부호 반전 + 평균)",
         font_size=15, color=LIGHT_GRAY)

add_shape(s, Inches(0.8), Inches(2.6), Inches(11.0), Inches(1.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(1.0), Inches(2.7), Inches(10.5), Inches(0.8),
         "J(w) = -(1/n) * SUM_i [y_i * ln(y_hat_i) + (1-y_i) * ln(1-y_hat_i)]",
         font_size=20, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_shape(s, Inches(0.8), Inches(4.0), Inches(11.0), Inches(1.0), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(1.0), Inches(4.1), Inches(10.5), Inches(0.8),
         "핵심 통찰: MLE 우도 최대화 = 로그 손실(크로스엔트로피) 최소화  [Cox, 1958]",
         font_size=18, color=ACCENT_ORANGE, bold=True)

add_text(s, Inches(0.6), Inches(5.3), Inches(11), Inches(0.4),
         "그래디언트 유도:", font_size=18, color=ACCENT_CYAN, bold=True)

add_code_block(s, Inches(0.8), Inches(5.8), Inches(11.0), Inches(1.2), [
    "dJ/dw_j = (1/n) * SUM_i (y_hat_i - y_i) * x_ij",
    "dJ/db   = (1/n) * SUM_i (y_hat_i - y_i)",
    "",
    "벡터 형태: nabla_w J = (1/n) * X^T * (y_hat - y)",
], font_size=14)

# ============================================================
# 슬라이드 17: 섹션 4 - 경사하강법
# ============================================================
section_divider("경사하강법", "배치/미니배치/SGD 비교, 학습률 선택", "5.4", ACCENT_PURPLE)

# ============================================================
# 슬라이드 18: 경사하강법의 원리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.4", "경사하강법 (Gradient Descent)의 원리")

add_shape(s, Inches(0.8), Inches(2.0), Inches(7.0), Inches(1.2), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(1.0), Inches(2.1), Inches(6.5), Inches(0.4),
         "파라미터 업데이트 규칙:", font_size=14, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(1.0), Inches(2.5), Inches(6.5), Inches(0.6),
         "w <- w - alpha * nabla_w J(w)\nb <- b - alpha * dJ/db",
         font_size=17, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_text(s, Inches(0.6), Inches(3.5), Inches(11), Inches(0.5),
         "배치 / 미니배치 / SGD 비교", font_size=20, color=ACCENT_CYAN, bold=True)

add_table_slide(s,
    ["방법", "데이터 사용량", "장점", "단점"],
    [
        ["배치 GD", "전체 (n개)", "안정적 수렴", "대용량에서 느림"],
        ["미니배치 GD", "일부 (32~256)", "GPU 효율적, 절충", "B 하이퍼파라미터"],
        ["SGD", "1개", "빠른 업데이트", "불안정한 수렴"],
    ],
    left=Inches(0.6), top=Inches(4.0), col_widths=[2.5, 2.5, 3.0, 3.5])

add_text(s, Inches(0.6), Inches(6.3), Inches(11), Inches(0.5),
         "* 실전에서는 미니배치 GD가 가장 많이 사용됨. sklearn은 기본적으로 L-BFGS(준뉴턴법) 사용.",
         font_size=14, color=ACCENT_ORANGE)

# ============================================================
# 슬라이드 19: 학습률 선택
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.4", "학습률(Learning Rate) 선택의 중요성")

add_table_slide(s,
    ["학습률", "현상", "결과"],
    [
        ["너무 작음 (0.001)", "수렴이 매우 느림", "반복 횟수 과도하게 증가"],
        ["적절함 (0.01~0.1)", "안정적으로 빠르게 수렴", "최적 비용에 도달"],
        ["너무 큼 (10)", "발산 (divergence)", "비용 함수가 증가"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.0, 3.5, 5.0])

add_card(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(2.8),
         "학습률 선택 가이드라인", [
             "1. 로그 스케일로 탐색:",
             "   {0.001, 0.003, 0.01, 0.03, 0.1, 0.3, 1.0}",
             "2. 비용 함수 수렴 곡선 확인",
             "3. 학습률 스케줄링: 점진적으로 감소",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.8),
         "특성 표준화의 필요성", [
             "x_j^scaled = (x_j - mu_j) / sigma_j",
             "",
             "표준화 없으면:",
             "  - 등고선이 타원형 -> 지그재그 경로",
             "  - 수렴 속도 크게 저하",
             "표준화 하면:",
             "  - 등고선이 원형 -> 직선 경로",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 20: 경사하강법 코드 구현
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.4", "경사하강법 코드 구현 (핵심부)")

add_code_block(s, Inches(0.5), Inches(2.0), Inches(12.0), Inches(5.0), [
    "def sigmoid(z):",
    "    z = np.clip(z, -500, 500)  # 수치 안정성",
    "    return 1.0 / (1.0 + np.exp(-z))",
    "",
    "def compute_cost(X, y, weights, bias):",
    "    m = X.shape[0]",
    "    z = np.dot(X, weights) + bias",
    "    y_hat = sigmoid(z)",
    "    epsilon = 1e-15",
    "    y_hat = np.clip(y_hat, epsilon, 1 - epsilon)  # log(0) 방지",
    "    cost = -(1.0/m) * np.sum(y*np.log(y_hat) + (1-y)*np.log(1-y_hat))",
    "    return cost",
    "",
    "def compute_gradients(X, y, weights, bias):",
    "    m = X.shape[0]",
    "    y_hat = sigmoid(np.dot(X, weights) + bias)",
    "    error = y_hat - y",
    "    dw = (1.0/m) * np.dot(X.T, error)  # nabla_w J",
    "    db = (1.0/m) * np.sum(error)        # dJ/db",
    "    return dw, db",
], font_size=12)

# ============================================================
# 슬라이드 21: 섹션 5 - 다중 클래스
# ============================================================
section_divider("다중 클래스 분류", "소프트맥스, OvR, OvO 전략", "5.5", ACCENT_ORANGE)

# ============================================================
# 슬라이드 22: 이진 분류에서 다중 분류로
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.5", "이진 분류에서 다중 분류로")

add_table_slide(s,
    ["구분", "이진 분류", "다중 분류"],
    [
        ["클래스 수", "2개", "K개 (K >= 3)"],
        ["활성화 함수", "시그모이드 (Sigmoid)", "소프트맥스 (Softmax)"],
        ["출력", "1개의 확률값", "K개의 확률값 (합=1)"],
        ["예시", "생존/사망, 스팸/정상", "붓꽃 종 분류, 숫자 인식(0~9)"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 4.5, 4.5])

add_shape(s, Inches(0.8), Inches(4.5), Inches(8.0), Inches(1.0), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(1.0), Inches(4.55), Inches(7.5), Inches(0.4),
         "소프트맥스 함수:", font_size=14, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(1.0), Inches(5.0), Inches(7.5), Inches(0.4),
         "P(y=k|X) = e^(z_k) / SUM_j e^(z_j),   z_k = w_k^T x + b_k",
         font_size=17, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_bullet_list(s, Inches(0.8), Inches(5.8), Inches(11), Inches(1.2), [
    "모든 출력이 (0, 1) 범위",
    "모든 클래스에 대한 확률의 합 = 1",
    "K=2일 때 시그모이드 함수와 수학적으로 동일",
], font_size=15, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 23: 소프트맥스 비용 함수와 OvR/OvO
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.5", "범주형 교차 엔트로피 & OvR/OvO 전략")

add_shape(s, Inches(0.8), Inches(2.0), Inches(10.0), Inches(0.9), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(1.0), Inches(2.05), Inches(9.5), Inches(0.4),
         "범주형 교차 엔트로피 (Categorical Cross Entropy):", font_size=14, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(1.0), Inches(2.45), Inches(9.5), Inches(0.4),
         "J = -(1/n) * SUM_i SUM_k y_ik * ln(y_hat_ik)",
         font_size=17, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_card(s, Inches(0.5), Inches(3.2), Inches(5.8), Inches(2.5),
         "OvR (One-vs-Rest) 전략", [
             "각 클래스: '해당 클래스 vs 나머지 전부'",
             "K개 클래스 -> K개의 이진 분류기 학습",
             "예측 시: 가장 높은 확률 분류기 선택",
             "sklearn: multi_class='ovr'",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(3.2), Inches(5.8), Inches(2.5),
         "OvO (One-vs-One) 전략", [
             "모든 클래스 쌍에 대해 이진 분류기 학습",
             "K개 클래스 -> K(K-1)/2개 분류기",
             "예측 시: 투표(voting)로 최종 결정",
             "개별 학습은 빠르지만 분류기 수가 많음",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_code_block(s, Inches(0.5), Inches(6.0), Inches(12.1), Inches(1.2), [
    "# sklearn 다중 분류",
    "model_softmax = LogisticRegression(multi_class='multinomial', solver='lbfgs', max_iter=200)",
    "model_ovr = LogisticRegression(multi_class='ovr', solver='lbfgs', max_iter=200)",
], font_size=12)

# ============================================================
# 슬라이드 24: 섹션 6 - 판별 vs 생성 모델
# ============================================================
section_divider("판별 vs 생성 모델", "Ng & Jordan(2002) 기반 비교", "5.6", ACCENT_RED)

# ============================================================
# 슬라이드 25: 두 가지 근본적 접근법
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.6", "두 가지 근본적 접근법 - Ng & Jordan(2002)")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
         "판별 모델 (Discriminative Model)", [
             "학습 대상: P(y|x) 직접 학습",
             "대표 모델: 로지스틱 회귀, SVM, 신경망",
             "결정 경계를 직접 학습",
             "점근적으로 더 정확 (낮은 편향)",
             "",
             "추정 방법: 조건부 우도 P(y|x) 최대화",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(3.0),
         "생성 모델 (Generative Model)", [
             "학습 대상: P(x|y)와 P(y) 학습 후",
             "  -> 베이즈 정리로 P(y|x) 유도",
             "대표 모델: 나이브 베이즈, GMM, HMM",
             "적은 데이터에서도 효과적",
             "",
             "추정 방법: 결합 우도 P(x, y) 최대화",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.5), Inches(5.3), Inches(12.1), Inches(1.8),
         "핵심 발견: 나이브 베이즈와 로지스틱 회귀의 이론적 관계", [
             "나이브 베이즈 가정(특성 독립성, 지수족 분포) 하에서:",
             "  ln[P(y=1|x)/P(y=0|x)] = theta_0 + SUM_j theta_j * x_j  ->  이것이 바로 로지스틱 회귀의 형태!",
             "두 모델은 동일한 모델 패밀리를 공유하지만, 파라미터 추정 방법이 다르다.",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 26: 수렴 속도와 전환점
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.6", "수렴 속도와 전환점 (Crossover Point)")

add_table_slide(s,
    ["특성", "나이브 베이즈 (생성)", "로지스틱 회귀 (판별)"],
    [
        ["점근적 오차", "높거나 같음", "낮거나 같음 (우수)"],
        ["수렴 속도", "O(log n) - 빠름!", "O(n) - 느림"],
        ["소표본 성능", "우수 (낮은 분산)", "열등 (높은 분산)"],
        ["대표본 성능", "열등 (높은 편향)", "우수 (낮은 편향)"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.0, 4.3, 4.3])

add_shape(s, Inches(0.8), Inches(4.5), Inches(11.0), Inches(1.0), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(1.0), Inches(4.6), Inches(10.5), Inches(0.8),
         "전환점(Crossover Point): 소표본에서는 NB가, 대표본에서는 LR이 우수.\n이 전환점은 데이터 특성과 모델 가정의 정확도에 따라 달라진다.",
         font_size=16, color=ACCENT_ORANGE, bold=True)

add_card(s, Inches(0.5), Inches(5.8), Inches(5.8), Inches(1.3),
         "NB 가정이 잘 성립하면", [
             "전환점이 늦게 나타남",
             "더 많은 데이터가 필요해야 LR이 역전",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(6.8), Inches(5.8), Inches(5.8), Inches(1.3),
         "NB 가정이 위반되면", [
             "전환점이 일찍 나타남",
             "LR이 빠르게 우세해짐",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 27: 편향-분산 트레이드오프 관점
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.6", "편향-분산 트레이드오프 관점")

add_table_slide(s,
    ["관점", "나이브 베이즈", "로지스틱 회귀"],
    [
        ["편향 (Bias)", "높음 (강한 모델 가정)", "낮음 (유연한 추정)"],
        ["분산 (Variance)", "낮음 (적은 파라미터)", "높음 (많은 파라미터)"],
        ["소표본", "낮은 분산이 유리", "높은 분산이 불리"],
        ["대표본", "높은 편향이 불리", "낮은 편향이 유리"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.0, 4.3, 4.3])

add_card(s, Inches(0.5), Inches(4.5), Inches(12.1), Inches(2.5),
         "Ng & Jordan(2002)의 실용적 가이드라인", [
             "데이터가 적을 때 (수백 개 이하): 나이브 베이즈(생성 모델) 먼저 시도",
             "데이터가 충분할 때 (수천 개 이상): 로지스틱 회귀(판별 모델) 사용",
             "두 모델을 모두 학습시키고 교차검증으로 비교하는 것이 가장 안전",
             "전환점은 데이터마다 다르므로, 실험적으로 확인 필요",
             "나이브 베이즈 가정이 잘 맞는 경우 (텍스트 분류 등)에는 NB가 대표본에서도 우수할 수 있음",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 28: 섹션 7 - 희귀 이벤트
# ============================================================
section_divider("희귀 이벤트", "King & Zeng(2001) 불균형 데이터 편향 보정", "5.7", ACCENT_RED)

# ============================================================
# 슬라이드 29: 희귀 이벤트 문제 정의
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.7", "희귀 이벤트 문제 - King & Zeng(2001)")

add_table_slide(s,
    ["분야", "사례", "양성 비율"],
    [
        ["정치학", "국제 분쟁 발생", "~1-3%"],
        ["금융", "신용카드 사기", "~0.1%"],
        ["의학", "희귀 질환 진단", "~1%"],
        ["제조", "불량품 탐지", "~0.5%"],
        ["정보보안", "침입 탐지", "~0.01%"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.0, 4.0, 4.5])

add_card(s, Inches(0.5), Inches(5.0), Inches(12.1), Inches(2.0),
         "유한표본 편향 (Finite-Sample Bias)", [
             "핵심: 양성 비율이 낮을수록 MLE 편향이 커짐",
             "편향 방향: 예측 확률의 과소추정 (실제보다 더 낮게 추정)",
             "bias(beta_hat) = (X^T W_hat X)^(-1) X^T W_hat xi",
             "결과: 희소 사건의 발생 확률을 과소추정하는 경향",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 30: 사전 교정과 가중 로지스틱 회귀
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.7", "편향 보정 방법")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.8),
         "방법 1: 사전 교정 (Prior Correction)", [
             "인구 수준 사건 발생률 tau와",
             "표본 사건 비율 y_bar를 이용:",
             "",
             "beta_0_corrected = beta_0",
             "  - ln[(1-tau)/tau * y_bar/(1-y_bar)]",
             "",
             "케이스-통제 설계에서 과대 표집 편향 교정",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.8),
         "방법 2: 가중 로지스틱 회귀", [
             "역선택확률 가중치(IPW) 부여:",
             "양성(Y=1) 가중치: w1 = tau / y_bar",
             "음성(Y=0) 가중치: w0 = (1-tau) / (1-y_bar)",
             "",
             "sklearn: class_weight='balanced'",
             "  -> 클래스 빈도에 반비례하는 가중치 자동 계산",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_text(s, Inches(0.6), Inches(5.1), Inches(11), Inches(0.4),
         "실용적 권고사항 (King & Zeng, 2001):", font_size=18, color=ACCENT_ORANGE, bold=True)

add_table_slide(s,
    ["양성 비율", "권고사항"],
    [
        ["> 20%", "표준 로지스틱 회귀 사용 가능"],
        ["5% ~ 20%", "가중 로지스틱 회귀 권고"],
        ["< 5%", "반드시 보정 방법 적용 (가중치 또는 사전 교정)"],
    ],
    left=Inches(0.6), top=Inches(5.6), col_widths=[3.5, 8.0])

# ============================================================
# 슬라이드 31: 섹션 8 - 모델 평가
# ============================================================
section_divider("모델 평가", "ROC/AUC, 혼동행렬, Precision/Recall/F1/MCC", "5.8", ACCENT_CYAN)

# ============================================================
# 슬라이드 32: 혼동 행렬
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.8", "혼동 행렬 (Confusion Matrix)")

add_table_slide(s,
    ["", "예측: 음성(0)", "예측: 양성(1)"],
    [
        ["실제: 음성(0)", "TN (True Negative)", "FP (False Positive) - 제1종 오류"],
        ["실제: 양성(1)", "FN (False Negative) - 제2종 오류", "TP (True Positive)"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.0, 4.3, 4.3], row_height=0.7)

add_card(s, Inches(0.5), Inches(3.8), Inches(5.8), Inches(3.0),
         "4가지 지표", [
             "TN: 실제 음성을 음성으로 올바르게 예측",
             "FP: 실제 음성을 양성으로 잘못 예측 (제1종 오류)",
             "FN: 실제 양성을 음성으로 잘못 예측 (제2종 오류)",
             "TP: 실제 양성을 양성으로 올바르게 예측",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.0),
         "비용 비대칭의 예", [
             "[질병 진단] FN(환자 놓침) >>> FP(정상인 재검사)",
             "  -> Recall 극대화 (임계값 낮춤)",
             "[스팸 필터] FP(정상 메일 차단) >> FN(스팸 통과)",
             "  -> Precision 극대화 (임계값 높임)",
             "[사기 탐지] FN(사기 놓침) >>> FP(정상 거래 재확인)",
             "  -> Recall 극대화",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 33: 주요 평가 지표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.8", "주요 평가 지표")

add_card(s, Inches(0.3), Inches(2.0), Inches(4.0), Inches(2.0),
         "Accuracy (정확도)", [
             "(TP+TN) / (TP+TN+FP+FN)",
             "불균형 데이터에서 무의미 가능",
             "양성 1%면 모두 음성 예측해도 99%",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(4.5), Inches(2.0), Inches(4.0), Inches(2.0),
         "Precision (정밀도)", [
             "TP / (TP + FP)",
             "양성 예측 중 실제 양성 비율",
             "'예측의 품질'을 측정",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.7), Inches(2.0), Inches(4.3), Inches(2.0),
         "Recall (재현율/민감도)", [
             "TP / (TP + FN)",
             "실제 양성 중 양성으로 예측한 비율",
             "'검출력'을 측정",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.3), Inches(4.3), Inches(4.0), Inches(2.5),
         "F1 Score", [
             "2 * (Precision * Recall) / (P + R)",
             "정밀도와 재현율의 조화평균",
             "두 지표의 균형을 측정",
             "0~1 범위, 1이 최선",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(4.5), Inches(4.3), Inches(8.5), Inches(2.5),
         "MCC (Matthews Correlation Coefficient)", [
             "MCC = (TP*TN - FP*FN) / sqrt((TP+FP)(TP+FN)(TN+FP)(TN+FN))",
             "범위: [-1, +1]. +1=완벽, 0=무작위, -1=완전 반대",
             "불균형 데이터에서 accuracy보다 신뢰할 수 있는 단일 지표",
             "양성/음성 비율에 영향을 덜 받음",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 34: ROC 곡선과 AUC
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.8", "ROC 곡선과 AUC")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "ROC 곡선 (Receiver Operating Characteristic)", [
             "x축: FPR (위양성률) = FP / (FP + TN)",
             "y축: TPR (진양성률=Recall) = TP / (TP + FN)",
             "다양한 임계값에서의 FPR-TPR 쌍을 연결",
             "대각선(y=x)이 무작위 기준선",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "PR 곡선 (Precision-Recall)", [
             "불균형 데이터에서 ROC보다 유용",
             "x축: Recall, y축: Precision",
             "AP (Average Precision): PR 곡선 아래 면적",
             "양성 클래스에 초점을 맞춤",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_text(s, Inches(0.6), Inches(4.8), Inches(11), Inches(0.4),
         "AUC (Area Under the ROC Curve) 해석:", font_size=18, color=ACCENT_ORANGE, bold=True)

add_table_slide(s,
    ["AUC 범위", "판별력", "해석"],
    [
        ["0.5", "무작위 추측 수준", "모델이 분류 능력 없음"],
        ["0.7 ~ 0.8", "수용 가능", "기본적 분류 능력"],
        ["0.8 ~ 0.9", "우수", "좋은 분류 성능"],
        ["> 0.9", "뛰어남", "매우 높은 분류 성능"],
    ],
    left=Inches(0.6), top=Inches(5.3), col_widths=[2.5, 3.5, 5.5])

# ============================================================
# 슬라이드 35: 임계값 조정
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.8", "임계값(Threshold) 조정 전략")

add_table_slide(s,
    ["상황", "임계값 방향", "강조 지표", "이유"],
    [
        ["질병 진단", "낮춤 (0.3)", "Recall", "놓치면 생명 위험"],
        ["스팸 필터링", "높임 (0.7)", "Precision", "정상 메일 손실 방지"],
        ["사기 탐지", "낮춤 (0.1)", "Recall", "사기 놓치면 큰 손실"],
        ["고객 이탈 예측", "낮춤 (0.3)", "Recall", "이탈 방지 실패 비용 큼"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.8, 2.5, 2.5, 3.7])

add_card(s, Inches(0.5), Inches(4.5), Inches(12.1), Inches(2.5),
         "Hosmer-Lemeshow 적합도 검정 (Hosmer & Lemeshow, 2000)", [
             "1. 예측 확률 기준으로 데이터를 g개 그룹(보통 10개)으로 분할",
             "2. 각 그룹에서 관측된 사건 수와 기대 사건 수를 비교",
             "3. 검정 통계량: C_hat = SUM_k (O_k - E_k)^2 / [E_k(1 - E_k/n_k)]",
             "4. 귀무가설: '모델이 데이터를 잘 설명한다'",
             "5. p-value가 높으면 모델 적합도가 좋다고 판단",
             "* Hosmer & Lemeshow(2000)의 c-통계량(concordance statistic) = AUC",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 36: 섹션 9 - 논문 리뷰
# ============================================================
section_divider("논문 리뷰", "Cox(1958), Hosmer&Lemeshow, King&Zeng, Ng&Jordan, Menard", "5.9", ACCENT_PURPLE)

# ============================================================
# 슬라이드 37: Cox(1958) + Hosmer & Lemeshow(2000)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.9", "논문 리뷰 (1/2)")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "Cox (1958)", [
             "'The Regression Analysis of Binary Sequences'",
             "로짓 변환의 체계적 도입",
             "MLE 기반 파라미터 추정 정립",
             "프로빗 모델 대비 로짓 모델의 장점 강조",
             "의학, 생물학, 사회과학 응용",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "Hosmer & Lemeshow (2000)", [
             "'Applied Logistic Regression' (2nd ed.)",
             "체계적 모델 구축 프레임워크",
             "7단계 목적적 변수 선택(Purposeful Selection)",
             "Hosmer-Lemeshow 적합도 검정 제안",
             "진단 통계량: 피어슨 잔차, dfbeta, Cook's D",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(0.5), Inches(4.8), Inches(12.1), Inches(2.2),
         "Hosmer & Lemeshow 7단계 변수 선택 과정", [
             "1. 단변량 분석 (p < 0.25 기준)  ->  2. 다변량 모델 적합",
             "3. 유의하지 않은 변수 제거 (p > 0.05)  ->  4. 교란변수 확인 (20% 이상 계수 변화 시 재포함)",
             "5. 연속형 변수의 로짓 선형성 확인  ->  6. 교호작용 추가",
             "7. 최종 모델 적합도 평가 (Hosmer-Lemeshow 검정)",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 38: King&Zeng(2001), Ng&Jordan(2002), Menard(2002)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.9", "논문 리뷰 (2/2)")

add_card(s, Inches(0.3), Inches(2.0), Inches(4.0), Inches(2.5),
         "King & Zeng (2001)", [
             "'Logistic Regression in",
             " Rare Events Data'",
             "MLE 편향 이론적 증명",
             "사전 교정(Prior Correction)",
             "가중 MLE 제안",
             "Firth(1993) 편향 축소 적용",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(4.5), Inches(2.0), Inches(4.2), Inches(2.5),
         "Ng & Jordan (2002)", [
             "'On Discriminative vs.",
             " Generative classifiers'",
             "동일 모델 패밀리 공유 증명",
             "O(log n) vs O(n) 수렴",
             "전환점(crossover) 확인",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.9), Inches(2.0), Inches(4.1), Inches(2.5),
         "Menard (2002)", [
             "'Applied Logistic Regression",
             " Analysis' (2nd ed.)",
             "OLS와의 체계적 비교",
             "다양한 계수 해석법 비교",
             "유사 R-squared 정리",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_text(s, Inches(0.6), Inches(4.8), Inches(11), Inches(0.4),
         "Menard(2002)가 정리한 유사 R-squared 지표:", font_size=16, color=ACCENT_CYAN, bold=True)

add_table_slide(s,
    ["지표", "범위", "특징"],
    [
        ["McFadden's R^2", "[0, 1)", "로그 우도 기반, 최대값 1 미도달"],
        ["Cox & Snell R^2", "[0, <1)", "최대값이 1 미만"],
        ["Nagelkerke R^2", "[0, 1]", "Cox & Snell 정규화, 최대 1 가능"],
        ["Count R^2", "[0, 1]", "정확 분류 비율, 가장 직관적"],
    ],
    left=Inches(0.6), top=Inches(5.3), col_widths=[3.5, 2.5, 5.5])

# ============================================================
# 슬라이드 39: 섹션 10 - 실습: 스크래치 구현
# ============================================================
section_divider("실습: 로지스틱 회귀 스크래치 구현", "Cox(1958) 기반 - Breast Cancer 데이터셋", "5.10", ACCENT_GREEN)

# ============================================================
# 슬라이드 40: 스크래치 구현 - 핵심 함수
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.10", "스크래치 구현 - 핵심 함수 3개")

add_code_block(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0), [
    "def sigmoid(z):",
    "    z = np.clip(z, -500, 500)",
    "    return 1.0 / (1.0 + np.exp(-z))",
    "",
    "def compute_cost(X, y, w, b):",
    "    m = X.shape[0]",
    "    y_hat = sigmoid(np.dot(X, w) + b)",
    "    eps = 1e-15",
    "    y_hat = np.clip(y_hat, eps, 1-eps)",
    "    return -(1/m) * np.sum(",
    "        y*np.log(y_hat) +",
    "        (1-y)*np.log(1-y_hat))",
    "",
    "def compute_gradients(X, y, w, b):",
    "    m = X.shape[0]",
    "    y_hat = sigmoid(np.dot(X, w) + b)",
    "    err = y_hat - y",
    "    dw = (1/m) * np.dot(X.T, err)",
    "    db = (1/m) * np.sum(err)",
    "    return dw, db",
], font_size=11)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.0),
         "코드 해설", [
             "[sigmoid] np.clip으로 수치 안정성 확보",
             "  z가 극단적이면 exp(-z) 오버플로 발생",
             "",
             "[compute_cost] epsilon=1e-15로 log(0) 방지",
             "  MLE의 음의 로그우도와 수학적으로 동일",
             "",
             "[compute_gradients] 벡터화 연산으로 효율적",
             "  error = y_hat - y로 모든 샘플 오차 한번에 계산",
             "  dw = (1/m) * X^T @ error",
             "  db = (1/m) * sum(error)",
             "",
             "이 그래디언트는 Cox(1958)의",
             "스코어 방정식에서 유도됨",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 41: 스크래치 구현 - 학습 루프
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.10", "스크래치 구현 - 경사하강법 학습 루프")

add_code_block(s, Inches(0.5), Inches(2.0), Inches(7.5), Inches(4.8), [
    "def logistic_regression_gd(X, y, lr=0.1, n_iter=1000):",
    "    m, n = X.shape",
    "    weights = np.zeros(n)  # 파라미터 초기화",
    "    bias = 0.0",
    "    cost_history = []",
    "",
    "    for i in range(n_iter):",
    "        cost = compute_cost(X, y, weights, bias)",
    "        cost_history.append(cost)",
    "",
    "        dw, db = compute_gradients(X, y, weights, bias)",
    "        weights -= lr * dw  # w = w - alpha * dw",
    "        bias -= lr * db     # b = b - alpha * db",
    "",
    "    return weights, bias, cost_history",
    "",
    "def predict(X, weights, bias, threshold=0.5):",
    "    z = np.dot(X, weights) + bias",
    "    probs = sigmoid(z)",
    "    return (probs >= threshold).astype(int), probs",
], font_size=12)

add_card(s, Inches(8.3), Inches(2.0), Inches(4.5), Inches(4.8),
         "실행 결과 (Breast Cancer)", [
             "데이터: sklearn Breast Cancer",
             "  - 569 samples, 30 features",
             "  - 악성(0): 212, 양성(1): 357",
             "",
             "직접 구현 (GD):",
             "  정확도: ~97%",
             "",
             "sklearn (L-BFGS):",
             "  정확도: ~97%",
             "",
             "미세 차이 원인:",
             "  1. 최적화 알고리즘 차이",
             "  2. sklearn 기본 L2 정규화",
             "  3. 수렴 기준 차이",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

# ============================================================
# 슬라이드 42: 스크래치 구현 - 시각화 코드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.10", "스크래치 구현 - 시각화: 시그모이드 & 비용 수렴")

add_code_block(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5), [
    "# 시그모이드 함수 시각화",
    "z_vals = np.linspace(-10, 10, 300)",
    "sig_vals = sigmoid(z_vals)",
    "plt.plot(z_vals, sig_vals, 'b-', lw=2.5)",
    "plt.axhline(y=0.5, color='r', ls='--')",
    "plt.xlabel('z'); plt.ylabel('sigma(z)')",
    "plt.title('Sigmoid Function')",
], font_size=11)

add_code_block(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5), [
    "# 비용 함수 수렴 곡선",
    "plt.plot(range(len(cost_history)),",
    "         cost_history, 'b-', lw=1.5)",
    "plt.xlabel('Iteration')",
    "plt.ylabel('Cost (Log Loss)')",
    "plt.title('Cost Convergence')",
    "plt.grid(True, alpha=0.3)",
], font_size=11)

add_code_block(s, Inches(0.5), Inches(4.8), Inches(12.1), Inches(2.3), [
    "# 결정 경계 시각화 (PCA로 2D 축소)",
    "pca = PCA(n_components=2)",
    "X_train_2d = pca.fit_transform(X_train_scaled)",
    "weights_2d, bias_2d, _ = logistic_regression_gd(X_train_2d, y_train, lr=0.1, n_iter=1000)",
    "",
    "xx, yy = np.meshgrid(np.linspace(x_min, x_max, 300), np.linspace(y_min, y_max, 300))",
    "grid = np.c_[xx.ravel(), yy.ravel()]",
    "prob_grid = sigmoid(np.dot(grid, weights_2d) + bias_2d).reshape(xx.shape)",
    "plt.contourf(xx, yy, prob_grid, levels=50, cmap='RdYlBu')  # 확률 등고선",
    "plt.contour(xx, yy, prob_grid, levels=[0.5], colors='black')  # 결정 경계",
], font_size=11)

# ============================================================
# 슬라이드 43: 섹션 11 - 판별 vs 생성 비교 실습
# ============================================================
section_divider("실습: 판별 vs 생성 비교", "Ng & Jordan(2002) 핵심 결과 재현", "5.11", ACCENT_BLUE)

# ============================================================
# 슬라이드 44: 판별 vs 생성 실험 코드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.11", "판별 vs 생성 비교 - 실험 코드")

add_code_block(s, Inches(0.5), Inches(2.0), Inches(12.1), Inches(4.8), [
    "def compare_models_by_sample_size(X, y, train_fractions, n_repeats=50):",
    "    \"\"\"다양한 학습 데이터 크기에서 로지스틱 회귀 vs 나이브 베이즈 비교\"\"\"",
    "    for frac in train_fractions:",
    "        n_train = max(int(n_total * (1 - test_size) * frac), 4)",
    "        for seed in range(n_repeats):",
    "            # 로지스틱 회귀 (판별 모델)",
    "            lr_model = LogisticRegression(max_iter=2000, solver='lbfgs')",
    "            lr_model.fit(X_train_s, y_train)",
    "            lr_scores.append(accuracy_score(y_test, lr_model.predict(X_test_s)))",
    "            ",
    "            # 나이브 베이즈 (생성 모델)",
    "            nb_model = GaussianNB()",
    "            nb_model.fit(X_train_s, y_train)",
    "            nb_scores.append(accuracy_score(y_test, nb_model.predict(X_test_s)))",
    "",
    "# 3가지 실험 실행",
    "# 실험 1: 독립적 특성 (NB 가정 성립) -> NB 소표본 우세, 전환점 관찰",
    "# 실험 2: 상관된 특성 (NB 가정 위반) -> LR이 더 빨리 우세",
    "# 실험 3: 고차원 (50개 특성) -> NB의 O(log n) 수렴 이점 부각",
], font_size=12)

add_bullet_list(s, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4), [
    "* 각 크기에서 50회 반복 실험하여 평균과 표준편차 산출 -> Ng & Jordan(2002) 실험 방법론 재현",
], font_size=13, color=ACCENT_ORANGE)

# ============================================================
# 슬라이드 45: 판별 vs 생성 실험 결과 해석
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.11", "판별 vs 생성 비교 - 실험 결과 해석")

add_card(s, Inches(0.3), Inches(2.0), Inches(4.0), Inches(2.5),
         "실험 1: 독립적 특성", [
             "NB 가정이 성립하는 데이터",
             "소표본에서 NB가 우세",
             "전환점(crossover) 관찰",
             "대표본에서 LR이 역전",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(s, Inches(4.5), Inches(2.0), Inches(4.2), Inches(2.5),
         "실험 2: 상관된 특성", [
             "NB 가정이 위반되는 데이터",
             "상관된 특성(n_redundant=10)",
             "전환점이 더 일찍 나타남",
             "LR이 빠르게 우세해짐",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.9), Inches(2.0), Inches(4.1), Inches(2.5),
         "실험 3: 고차원 데이터", [
             "50개 특성 (30개 유의미)",
             "소표본에서 LR 불안정",
             "NB의 O(log n) 수렴 이점",
             "차원의 저주 효과 확인",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.5), Inches(4.8), Inches(12.1), Inches(2.2),
         "종합 결론 (Ng & Jordan, 2002 확인)", [
             "1. 생성 모델(NB)은 O(log n)으로 수렴: 소표본에서 데이터 효율적",
             "2. 판별 모델(LR)은 점근적으로 더 낮은 오차: 대표본에서 우수",
             "3. 전환점은 데이터 특성(NB 가정 성립 여부, 차원 수)에 따라 달라짐",
             "4. 실전: 두 모델을 모두 학습시키고 교차검증으로 비교하는 것이 가장 안전",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 46: 섹션 12 - 희귀 이벤트 보정 실습
# ============================================================
section_divider("실습: 희귀 이벤트 보정", "King & Zeng(2001) 핵심 결과 재현", "5.12", ACCENT_RED)

# ============================================================
# 슬라이드 47: 희귀 이벤트 보정 코드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.12", "희귀 이벤트 보정 - 핵심 코드")

add_code_block(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.8), [
    "# 사전 교정 (Prior Correction) 구현",
    "def prior_correction(model, pop_rate, sample_rate):",
    "    correction = np.log(",
    "        ((1-pop_rate)/pop_rate) *",
    "        (sample_rate/(1-sample_rate))",
    "    )",
    "    return correction",
    "",
    "# 보정된 예측 확률",
    "logits = X @ model.coef_[0] + model.intercept_[0]",
    "corrected = logits - correction",
    "prob = 1 / (1 + np.exp(-corrected))",
], font_size=11)

add_code_block(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.8), [
    "# 3가지 모델 비교",
    "# 모델 1: 표준 로지스틱 회귀",
    "model_std = LogisticRegression(max_iter=2000)",
    "",
    "# 모델 2: 가중 로지스틱 회귀",
    "model_wt = LogisticRegression(",
    "    class_weight='balanced', max_iter=2000)",
    "",
    "# 모델 3: 사전 교정 적용",
    "prob_corrected = predict_with_prior_correction(",
    "    model_std, X_test, pop_rate, sample_rate)",
], font_size=11)

add_card(s, Inches(0.5), Inches(5.1), Inches(12.1), Inches(2.0),
         "실험 결과: 양성 비율별 MLE 편향", [
             "50%: 편향 거의 없음 | 20%: 미세한 과소추정",
             "10%: 눈에 띄는 과소추정 | 5%: 상당한 과소추정",
             "2%: 심각한 과소추정 | 1%: 매우 심각한 과소추정",
             "-> King & Zeng(2001)의 이론적 예측과 일치: 양성 비율이 낮을수록 과소추정 편향 증가",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 48: 희귀 이벤트 - ROC/PR 곡선과 임계값 분석
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.12", "희귀 이벤트 - ROC/PR 곡선 & 임계값 분석")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3),
         "ROC 곡선 비교 (양성 비율 2%)", [
             "표준 LR: 높은 AUC (판별력 자체는 우수)",
             "가중 LR (balanced): AUC 유사 또는 소폭 변동",
             "사전 교정: AUC 동일 (기울기만 이동)",
             "=> AUC만으로는 보정 효과를 구분하기 어려움",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.3),
         "PR 곡선 비교 (더 유용!)", [
             "불균형 데이터에서는 PR 곡선이 더 유용",
             "가중 LR: AP(Average Precision) 향상",
             "양성 클래스 검출에 초점",
             "기준선: 양성 비율 = 0.02",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_text(s, Inches(0.6), Inches(4.6), Inches(11), Inches(0.4),
         "임계값 분석 결과 (양성 비율 2%):", font_size=16, color=ACCENT_ORANGE, bold=True)

add_table_slide(s,
    ["임계값", "정확도", "정밀도", "재현율", "F1"],
    [
        ["0.01", "낮음", "낮음", "높음", "중간"],
        ["0.02", "중간", "중간", "높음", "양호"],
        ["0.05", "높음", "높음", "중간", "양호"],
        ["0.50", "98%", "높음/0", "매우 낮음", "낮음"],
    ],
    left=Inches(0.6), top=Inches(5.1), col_widths=[2.0, 2.2, 2.2, 2.2, 2.9])

add_text(s, Inches(0.6), Inches(7.0), Inches(11), Inches(0.4),
         "* 임계값 0.5에서는 양성을 거의 검출 못함! 양성 비율에 맞게 임계값 조정 필수.",
         font_size=13, color=ACCENT_RED)

# ============================================================
# 슬라이드 49: 섹션 13 - 응용사례
# ============================================================
section_divider("응용사례", "사기탐지, 질병진단, 고객이탈", "5.13", ACCENT_ORANGE)

# ============================================================
# 슬라이드 50: 사기 탐지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.13", "응용사례 1: 사기 탐지 (Fraud Detection)")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3),
         "문제 특성", [
             "양성 비율 극히 낮음 (~0.1%)",
             "  -> King & Zeng(2001) 희소 사건",
             "실시간 처리 요구 (거래 즉시 판단)",
             "비대칭 비용: 사기 놓침 >>> 정상 거래 차단",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.3),
         "로지스틱 회귀 적용", [
             "가중 LR 사용 (King & Zeng 권고)",
             "class_weight={0:1, 1:100}  # 사기 가중치 100배",
             "주요 피처: 거래금액, 빈도, 위치, 시간대",
             "평가: PR-AUC, Recall, F1 (accuracy 부적절)",
         ], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_code_block(s, Inches(0.5), Inches(4.6), Inches(12.1), Inches(2.5), [
    "# 사기 탐지 모델 예시",
    "fraud_model = LogisticRegression(",
    "    class_weight={0: 1, 1: 100},  # 사기(1) 가중치 100배",
    "    max_iter=1000, random_state=42",
    ")",
    "fraud_model.fit(X_train, y_train)",
    "",
    "# 임계값을 낮게 설정 (0.1): Recall 극대화",
    "prob = fraud_model.predict_proba(X_test)[:, 1]",
    "pred = (prob >= 0.1).astype(int)  # 임계값 0.5 대신 0.1",
], font_size=12)

# ============================================================
# 슬라이드 51: 질병 진단
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.13", "응용사례 2: 질병 진단 (Disease Diagnosis)")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.0),
         "문제 특성", [
             "유병률에 따라 불균형 정도 다양",
             "FN(환자 놓침) 비용이 매우 높음",
             "해석 가능성 중요 (의사의 판단 근거)",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "로지스틱 회귀의 장점", [
             "오즈비 해석: 위험 인자 영향 정량화",
             "확률 출력: '당뇨 확률 73%' 직관적 소통",
             "계수 해석: Cox(1958) 오즈비가 의학 표준",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_code_block(s, Inches(0.5), Inches(4.3), Inches(12.1), Inches(2.8), [
    "# 질병 진단 모델 + 오즈비 계산",
    "model = LogisticRegression()",
    "model.fit(X_train, y_train)",
    "",
    "# 오즈비 계산 및 해석",
    "odds_ratios = np.exp(model.coef_[0])",
    "feature_importance = pd.DataFrame({",
    "    'Feature': feature_names,",
    "    'Coefficient': model.coef_[0],",
    "    'Odds Ratio': odds_ratios",
    "}).sort_values('Odds Ratio', ascending=False)",
    "",
    "# 예: BMI의 OR=1.85 -> BMI 1단위 증가 시 당뇨 오즈 85% 증가",
    "# 임계값: 0.3으로 낮춤 -> Recall 극대화 (놓치면 생명 위험)",
], font_size=12)

# ============================================================
# 슬라이드 52: 고객 이탈 예측
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.13", "응용사례 3: 고객 이탈 예측 (Customer Churn)")

add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.0),
         "문제 특성", [
             "이탈률 15~30% (중간 수준 불균형)",
             "FN(이탈 고객 놓침) > FP(불필요한 마케팅)",
             "신규 고객 유치 비용 = 유지 비용의 5~7배",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "비즈니스 전략", [
             "예측 확률 상위 20% -> 리텐션 마케팅",
             "높은 계수 피처 분석 -> 서비스 개선",
             "임계값 0.3으로 낮춤 -> Recall 확보",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_code_block(s, Inches(0.5), Inches(4.3), Inches(12.1), Inches(2.8), [
    "# 고객 이탈 예측 모델",
    "model_churn = LogisticRegression(max_iter=1000, random_state=42)",
    "model_churn.fit(X_train, y_train)",
    "",
    "pred = model_churn.predict(X_test)",
    "prob = model_churn.predict_proba(X_test)[:, 1]",
    "",
    "# 혼동 행렬 분석",
    "cm = confusion_matrix(y_test, pred)",
    "tn, fp, fn, tp = cm.ravel()",
    "print(f'TN (유지->유지): {tn}명')",
    "print(f'FP (유지->이탈): {fp}명  <- 불필요한 마케팅 비용')",
    "print(f'FN (이탈->유지): {fn}명  <- 이탈 방지 실패 (가장 위험!)')",
    "print(f'TP (이탈->이탈): {tp}명  <- 리텐션 마케팅 대상')",
], font_size=12)

# ============================================================
# 슬라이드 53: 섹션 14 - 핵심 요약
# ============================================================
section_divider("핵심 요약", "수식 정리, 용어, 복습 질문", "5.14", ACCENT_BLUE)

# ============================================================
# 슬라이드 54: 알고리즘 요약
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.14", "로지스틱 회귀 알고리즘 요약")

add_table_slide(s,
    ["항목", "내용"],
    [
        ["유형", "지도학습 - 분류 (Classification)"],
        ["핵심 함수", "시그모이드: sigma(z) = 1/(1+e^(-z))"],
        ["비용 함수", "로그 손실 (Log Loss / Binary Cross Entropy)"],
        ["최적화", "경사하강법 또는 L-BFGS (준뉴턴법)"],
        ["출력", "클래스 확률 (0~1 사이)"],
        ["장점", "구현 간단, 확률 해석, 계수로 피처 영향도, 연산 효율적"],
        ["단점", "선형 결정 경계만 가능, 복잡한 비선형 패턴 한계"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.0, 8.5])

add_text(s, Inches(0.6), Inches(5.7), Inches(11), Inches(0.4),
         "5개 논문의 핵심 메시지:", font_size=18, color=ACCENT_ORANGE, bold=True)

add_table_slide(s,
    ["논문", "핵심 메시지"],
    [
        ["Cox (1958)", "로짓 변환 + MLE = 로지스틱 회귀의 수학적 기초"],
        ["Hosmer & Lemeshow (2000)", "체계적 모델 구축, 진단, 적합도 평가"],
        ["King & Zeng (2001)", "희소 사건에서 MLE 편향 존재, 보정 필수"],
        ["Ng & Jordan (2002)", "소표본 -> 생성 모델, 대표본 -> 판별 모델"],
    ],
    left=Inches(0.6), top=Inches(6.1), col_widths=[4.0, 7.5], row_height=0.3, font_size=12)

# ============================================================
# 슬라이드 55: 주요 수식 정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.14", "주요 수식 정리")

add_table_slide(s,
    ["수식", "의미"],
    [
        ["sigma(z) = 1/(1+e^(-z))", "시그모이드 함수"],
        ["logit(p) = ln(p/(1-p))", "로짓 (시그모이드 역함수)"],
        ["J = -(1/n)SUM[y*ln(y_hat)+(1-y)*ln(1-y_hat)]", "이진 교차 엔트로피"],
        ["nabla_w J = (1/n)*X^T*(y_hat-y)", "그래디언트"],
        ["OR = e^(w_i)", "오즈비"],
        ["P(y=k|X) = e^(z_k) / SUM_j e^(z_j)", "소프트맥스"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[6.5, 5.0])

add_code_block(s, Inches(0.5), Inches(5.2), Inches(12.1), Inches(2.0), [
    "# sklearn 핵심 코드 요약",
    "from sklearn.linear_model import LogisticRegression",
    "model = LogisticRegression(max_iter=1000)",
    "model.fit(X_train, y_train)",
    "pred = model.predict(X_test)            # 클래스 예측",
    "prob = model.predict_proba(X_test)       # 확률 예측",
    "accuracy_score(y_test, pred)             # 정확도",
    "roc_auc_score(y_test, prob[:, 1])        # AUC",
    "np.exp(model.coef_[0])                   # 오즈비",
], font_size=12)

# ============================================================
# 슬라이드 56: 복습 질문 10개
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "5.14", "복습 질문 10개")

add_bullet_list(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.2), [
    "Q1. 선형 회귀로 분류 시 4가지 문제점과 로지스틱 회귀의 해결 방법?",
    "Q2. 시그모이드 미분이 sigma(z)(1-sigma(z))임을 증명하시오.",
    "Q3. 오즈비(OR)=1.5의 의미를 타이타닉 예시로 설명하시오.",
    "Q4. MLE 우도 최대화 = 로그 손실 최소화의 동치성을 보이시오.",
    "Q5. Ng&Jordan: NB는 O(log n), LR은 O(n) 수렴의 시사점?",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(10))

add_bullet_list(s, Inches(6.5), Inches(2.0), Inches(6.0), Inches(5.2), [
    "Q6. King&Zeng: 양성 1% 데이터에서 MLE 편향과 2가지 교정법?",
    "Q7. 양성 2% 사기 탐지에서 정확도 98%가 좋은 모델인가?",
    "Q8. K=2일 때 소프트맥스가 시그모이드로 환원됨을 보이시오.",
    "Q9. Hosmer-Lemeshow 검정의 절차와 귀무/대립 가설?",
    "Q10. 직접 구현(GD)과 sklearn(L-BFGS) 성능 차이의 원인 2가지?",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(10))

# ============================================================
# 슬라이드 57: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.0), prs.slide_width, Inches(1.0),
         "Thank You", font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(3.2), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(3.5), prs.slide_width, Inches(0.8),
         "5장: 로지스틱 회귀 (Logistic Regression)", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.3), prs.slide_width, Inches(0.5),
         "확장 상세 버전 | 57 Slides", font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.0), prs.slide_width, Inches(0.8),
         "핵심: 시그모이드 + 로그 손실 + MLE = 로지스틱 회귀\n"
         "Cox(1958) | Hosmer&Lemeshow(2000) | King&Zeng(2001) | Ng&Jordan(2002) | Menard(2002)",
         font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 저장
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), "5장_로지스틱회귀_강의PPT_확장.pptx")
prs.save(output_path)
print(f"[완료] {output_path} 저장됨  (총 {len(prs.slides)}장)")
