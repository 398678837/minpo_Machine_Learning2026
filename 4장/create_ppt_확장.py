"""4장 선형회귀 - 확장 강의 PPT 생성 스크립트 (상세 버전)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = '맑은 고딕'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    total_w = sum(col_widths)
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(1), Inches(2.5), Inches(11.333), Inches(3), SECTION_BG, accent)
    add_accent_line(s, Inches(5.5), Inches(2.6), Inches(2.333), accent)
    add_text(s, Inches(1), Inches(2.8), Inches(11.333), Inches(0.5),
             f"SECTION {section_num}", font_size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.4), Inches(11.333), Inches(1),
             title, font_size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.4), Inches(11.333), Inches(0.6),
             subtitle, font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG, ACCENT_BLUE)
add_accent_line(s, Inches(4), Inches(1.5), Inches(5.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.8), Inches(11.333), Inches(0.5),
         "CHAPTER 4", font_size=22, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.5), Inches(11.333), Inches(1.2),
         "선형회귀", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.333), Inches(0.5),
         "Linear Regression: OLS, Ridge, Lasso, Elastic Net", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.8),
         "OLS부터 정규화 기법까지 체계적 이해 + 직접 구현으로 알고리즘 내부 동작 파악",
         font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
         "기계학습 | Machine Learning", font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 2: 목차
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, None, "목차 (Table of Contents)", "Chapter 4 전체 구성")
toc_left = [
    "4.1  역사와 의의",
    "4.2  OLS 이론 (정규방정식, 기하학적 해석)",
    "4.3  다중공선성 (VIF, 조건수)",
    "4.4  Ridge 회귀 (L2 정규화)",
    "4.5  Lasso 회귀 (L1 정규화)",
    "4.6  Elastic Net (L1+L2)",
    "4.7  좌표 하강법 (glmnet)",
]
toc_right = [
    "4.8   SCAD와 Oracle 특성",
    "4.9   회귀 진단",
    "4.10  실습: Ridge 회귀 구현",
    "4.11  실습: Lasso 회귀 구현",
    "4.12  실습: Elastic Net 구현",
    "4.13  응용사례",
    "4.14  핵심 요약 + 복습 질문",
]
add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
         "Part I: 이론", toc_left, ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
         "Part II: 실습 & 응용", toc_right, ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# 슬라이드 3: 학습 목표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, None, "학습 목표")
objectives = [
    "1. OLS의 정규방정식을 행렬 미분으로 유도하고 기하학적으로 해석할 수 있다",
    "2. 가우스-마르코프 정리의 의미와 한계를 설명할 수 있다",
    "3. Ridge, Lasso, Elastic Net의 수학적 정의와 차이를 비교할 수 있다",
    "4. L1 vs L2 정규화의 기하학적 차이를 설명할 수 있다",
    "5. 좌표 하강법과 연성 임계값 연산자를 직접 구현할 수 있다",
    "6. SCAD의 오라클 성질과 Lasso 대비 장점을 설명할 수 있다",
    "7. 회귀 진단(잔차 분석, Cook's Distance)을 수행할 수 있다",
]
add_bullet_list(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(5),
                objectives, font_size=18, color=LIGHT_GRAY, spacing=Pt(12))

# ============================================================
# SECTION 1: 역사와 의의
# ============================================================
section_divider("역사와 의의", "200년 역사의 가장 기본적인 ML 알고리즘", 1, ACCENT_BLUE)

# 슬라이드 5: 최소제곱법의 탄생
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 1, "최소제곱법의 탄생", "Legendre (1805) vs Gauss (1809)")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "르장드르 (Legendre, 1805)", [
             "최초로 최소제곱법을 공식 발표",
             "혜성 궤도를 관측 데이터로부터 결정",
             "잔차 제곱합을 최소화하는 원리 제안",
             "Nouvelles méthodes... (1805)",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "가우스 (Gauss, 1809)", [
             "1795년부터 사용해왔다고 주장",
             "Theoria Motus Corporum Coelestium (1809)",
             "정규분포 가정 하에서 MLE와 동일함을 증명",
             "→ 확률론적 정당성 부여",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0),
         "선형회귀의 현대적 위치", [
             "기초 알고리즘: 지도학습에서 연속값을 예측하는 가장 기본 모델",
             "이론적 토대: 비용함수 최적화, 편향-분산 트레이드오프, 정규화의 출발점",
             "확장의 기반: 로지스틱 회귀, 신경망, 커널 방법의 구성 요소",
             "벤치마크: 새로운 모델의 성능 비교 기준선 (baseline)",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 슬라이드 6: 수학적 정의
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 1, "선형회귀의 수학적 정의", "독립변수와 종속변수 사이의 선형 관계")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(2.5),
         "스칼라 표기", [
             "y = β₀ + β₁x₁ + β₂x₂ + ··· + βₚxₚ + ε",
             "",
             "y: 종속변수 (타겟)",
             "x₁,...,xₚ: 독립변수 (특성)",
             "β₀: 절편, β₁,...,βₚ: 회귀 계수",
             "ε: 오차항",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(2.5),
         "행렬 표기", [
             "y = Xβ + ε",
             "",
             "y ∈ R^n: 반응 벡터 (n개 관측치)",
             "X ∈ R^(n×(p+1)): 설계 행렬",
             "β ∈ R^(p+1): 회귀 계수 벡터",
             "ε ∈ R^n: 오차, εᵢ ~iid~ N(0, σ²)",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_text(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.8),
         "핵심: '선형'은 파라미터 β에 대해 선형이라는 의미. y = β₁x + β₂x²도 선형회귀 (x²를 새 특성으로 취급)",
         font_size=16, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 2: OLS 이론
# ============================================================
section_divider("OLS 이론", "정규방정식 유도, 기하학적 해석, 가우스-마르코프 정리", 2, ACCENT_GREEN)

# 슬라이드 8: 정규방정식 유도
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "정규방정식 유도 (행렬 미분)", "RSS를 최소화하는 β를 찾는다")
add_card(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(5.0),
         "단계별 유도", [
             "Step 1. 잔차 제곱합 (RSS) 정의",
             "  RSS(β) = ||y - Xβ||² = (y - Xβ)ᵀ(y - Xβ)",
             "",
             "Step 2. 전개",
             "  RSS(β) = yᵀy - 2βᵀXᵀy + βᵀXᵀXβ",
             "",
             "Step 3. β에 대해 편미분 (행렬 미분 규칙 적용)",
             "  ∂RSS/∂β = -2Xᵀy + 2XᵀXβ",
             "  사용된 공식: ∂(aᵀx)/∂x = a,  ∂(xᵀAx)/∂x = 2Ax (A 대칭)",
             "",
             "Step 4. 0으로 놓기 → 정규방정식 (Normal Equation)",
             "  XᵀXβ = Xᵀy",
             "",
             "Step 5. 해 (XᵀX가 가역일 때)",
             "  β̂_OLS = (XᵀX)⁻¹Xᵀy      ← 핵심 수식",
         ], ACCENT_GREEN, ACCENT_GREEN)

# 슬라이드 9: 헤시안과 볼록성
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "2차 미분(헤시안)과 볼록성 확인", "정규방정식의 해가 최솟값임을 보장")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.5),
         "헤시안 (2차 미분)", [
             "∂²RSS / ∂β∂βᵀ = 2XᵀX",
             "",
             "XᵀX는 양의 반정치(positive semi-definite)",
             "  → RSS는 볼록(convex) 함수",
             "  → 정규방정식의 해는 전역 최솟값",
             "",
             "XᵀX가 양정치(positive definite)이면",
             "  → 유일한 최솟값",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.5),
         "볼록 최적화의 이점", [
             "1. 지역 최솟값 = 전역 최솟값",
             "  → 어디서 시작해도 동일한 해",
             "",
             "2. 경사하강법이 수렴 보장",
             "  → 학습률만 적절하면 됨",
             "",
             "3. Ridge도 볼록 (λI 추가해도 볼록 유지)",
             "  → Lasso도 볼록 (|β|가 볼록)",
             "",
             "4. SCAD는 비볼록 → 지역 최솟값 문제",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 슬라이드 10: 기하학적 해석
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "OLS의 기하학적 해석", "직교 사영 (Orthogonal Projection)")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5),
         "직교 사영 관점", [
             "y는 R^n 공간의 벡터",
             "C(X) = X의 열공간 (column space)",
             "",
             "ŷ = X(XᵀX)⁻¹Xᵀy = Hy",
             "",
             "H = X(XᵀX)⁻¹Xᵀ",
             "  = 햇 행렬 (hat matrix)",
             "  = 사영 행렬 (projection matrix)",
             "",
             "ŷ는 y를 C(X)에 직교 사영한 것",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.0),
         "핵심 성질", [
             "잔차 Xᵀe = 0 (잔차와 열공간 직교)",
             "피타고라스: ||y||² = ||ŷ||² + ||e||²",
             "  → TSS = ESS + RSS",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(6.7), Inches(4.3), Inches(6.1), Inches(2.2),
         "R²의 기하학적 의미", [
             "R² = cos²θ  (y와 ŷ 사이 각도)",
             "",
             "R² = 1: y가 C(X) 안에 있음 (완벽 적합)",
             "R² = 0: y가 C(X)에 직교 (설명력 없음)",
             "0 < R² < 1: 부분적 설명",
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# 슬라이드 11: 가우스-마르코프 정리
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "가우스-마르코프 정리", "OLS 추정량의 최적성 보장")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(3.0),
         "정리의 조건 (4가지)", [
             "1. E[εᵢ] = 0 (오차 기대값 0)",
             "2. Var(εᵢ) = σ² (등분산성)",
             "3. Cov(εᵢ, εⱼ) = 0, i≠j (무상관)",
             "4. X는 고정된(non-random) 행렬",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(3.0),
         "BLUE: 최적성의 의미", [
             "Best: 모든 선형 비편향 추정량 중 분산 최소",
             "Linear: β̂이 y의 선형 함수",
             "Unbiased: E[β̂] = β (편향 없음)",
             "Estimator: β의 추정량",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
         "핵심 한계와 Ridge의 동기", [
             "가우스-마르코프는 OLS가 비편향 추정량 중 최적임을 보장할 뿐!",
             "편향을 허용하면 MSE = Var + Bias² 관점에서 OLS보다 나은 추정량이 존재",
             "→ Hoerl & Kennard (1970): 약간의 편향으로 분산을 크게 줄여 전체 MSE 감소",
         ], ACCENT_RED, ACCENT_RED)

# ============================================================
# SECTION 3: 다중공선성
# ============================================================
section_divider("다중공선성", "VIF, 조건수, 해결 방법", 3, ACCENT_ORANGE)

# 슬라이드 13: 다중공선성
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 3, "다중공선성 (Multicollinearity)", "독립변수 간 높은 선형 상관관계 문제")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(2.5),
         "문제점", [
             "XᵀX가 특이(singular)에 가까워지면:",
             "(XᵀX)⁻¹의 원소가 매우 커짐",
             "Var(β̂) = σ²(XᵀX)⁻¹이 폭증",
             "회귀 계수 추정이 불안정 (높은 분산)",
         ], ACCENT_RED, ACCENT_RED)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(2.5),
         "VIF (Variance Inflation Factor)", [
             "VIFⱼ = 1 / (1 - Rⱼ²)",
             "",
             "VIF = 1: 다중공선성 없음",
             "VIF 1~5: 약한 (일반적으로 허용)",
             "VIF 5~10: 중간 (주의 필요)",
             "VIF > 10: 심각 (조치 필요)",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.2),
         "조건수 (Condition Number)", [
             "κ(XᵀX) = λ_max / λ_min",
             "",
             "조건수 > 30 → 다중공선성 의심",
             "조건수가 크면 역행렬 계산 불안정",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_table_slide(s,
    ["해결 방법", "설명"],
    [
        ["변수 제거", "VIF 높은 변수를 수동 제거"],
        ["주성분 회귀", "PCA로 차원 축소 후 회귀"],
        ["Ridge 회귀", "L2 정규화로 조건수 개선"],
        ["Lasso 회귀", "L1 정규화로 변수 자동 제거"],
    ],
    Inches(6.8), Inches(4.8), [2.5, 3.3],
    row_height=0.45, font_size=12)

# ============================================================
# SECTION 4: Ridge 회귀
# ============================================================
section_divider("Ridge 회귀", "L2 정규화: Hoerl & Kennard (1970)", 4, ACCENT_PURPLE)

# 슬라이드 15: Ridge 핵심 아이디어
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 4, "Ridge 회귀: 핵심 아이디어", "편향-분산 트레이드오프의 고전적 예시")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
         "편향-분산 트레이드오프", [
             "MSE(β̂) = Var(β̂) + [Bias(β̂)]²",
             "",
             "OLS: Bias = 0, Var = 클 수 있음",
             "Ridge: Bias > 0 (약간), Var = 크게 감소",
             "→ 전체 MSE 감소!",
             "",
             "핵심 통찰 (Hoerl & Kennard, 1970):",
             "  가우스-마르코프의 '허점' 이용",
             "  비편향 중 최적이지만...",
             "  약간의 편향을 허용하면",
             "  MSE가 더 좋아질 수 있다!",
         ], ACCENT_PURPLE, ACCENT_PURPLE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(5.0),
         "L2 정규화의 수학적 정의", [
             "최적화 문제:",
             "  min ||y - Xβ||² + λ||β||²",
             "",
             "닫힌 형태 해:",
             "  β̂_ridge = (XᵀX + λI)⁻¹Xᵀy",
             "",
             "유도:",
             "  ∂/∂β [...] = -2Xᵀy + 2XᵀXβ + 2λβ = 0",
             "  (XᵀX + λI)β = Xᵀy",
             "",
             "핵심 효과:",
             "  λI를 더해 대각 원소 증가",
             "  → 조건수 감소",
             "  → 수치적 안정성 확보",
             "  → 계수를 0 방향으로 축소",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 슬라이드 16: 고유값 분해와 베이지안 해석
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 4, "Ridge: 고유값 분해 & 베이지안 해석", "두 가지 관점에서의 이해")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "고유값 분해를 통한 이해", [
             "XᵀX = VDVᵀ에서",
             "Ridge는 각 방향의 계수를 dⱼ/(dⱼ+λ) 축소",
             "",
             "dⱼ 큼: 축소 비율 ≈ 1 (거의 변화 없음)",
             "dⱼ 작음: 축소 비율 ≈ 0 (강하게 축소)",
             "→ 불안정한 방향을 선택적으로 축소",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "베이지안 해석", [
             "사전분포: β ~ N(0, τ²I)",
             "우도: y|X,β ~ N(Xβ, σ²I)",
             "",
             "MAP 추정 = Ridge (λ = σ²/τ²)",
             "→ Ridge = 계수에 정규분포 사전분포 부여",
             "τ² 작을수록 → λ 커짐 → 더 많이 축소",
         ], ACCENT_PURPLE, ACCENT_PURPLE)
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0),
         "릿지 트레이스 (Ridge Trace)", [
             "x축: λ, y축: 각 회귀 계수 값을 그린 그래프",
             "λ = 0 (OLS)에서 시작 → λ 증가 → 모든 계수가 0으로 수렴",
             "계수가 안정화되는 λ 값을 시각적으로 선택 (교차 검증으로 정밀화)",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# SECTION 5: Lasso 회귀
# ============================================================
section_divider("Lasso 회귀", "L1 정규화: Tibshirani (1996)", 5, ACCENT_RED)

# 슬라이드 18: Lasso 핵심
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "Lasso: L1 정규화와 변수 선택", "축소(shrinkage) + 변수 선택을 동시에!")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
         "L1 정규화의 수학적 정의", [
             "min (1/2n)||y - Xβ||² + λ||β||₁",
             "",
             "||β||₁ = Σ|βⱼ| (L1 노름)",
             "",
             "닫힌 형태 해 없음 (원점에서 미분 불가능)",
             "→ 좌표 하강법(Coordinate Descent) 필요",
         ], ACCENT_RED, ACCENT_RED)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.0),
         "Lasso의 혁신 (Tibshirani, 1996)", [
             "Ridge의 L2를 L1으로 교체한 간단한 변경이",
             "근본적으로 다른 성질을 만들어냄:",
             "",
             "1. 회귀 계수 축소 (shrinkage)",
             "2. 변수 선택 (variable selection) ← 핵심!",
             "3. 일부 계수를 정확히 0으로 만듦",
             "",
             "Google Scholar 인용 50,000+ (2024 기준)",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 슬라이드 19: 연성 임계값 연산자
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "연성 임계값 연산자 (Soft Thresholding)", "Lasso 해의 핵심 메커니즘")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(4.8),
         "수식 및 동작", [
             "S(z, γ) = sign(z) · max(|z| - γ, 0)",
             "",
             "세 구간별 동작:",
             "",
             "|z| ≤ γ:   결과 = 0",
             "  → 작은 계수를 정확히 0으로 (변수 제거!)",
             "",
             "z > γ:    결과 = z - γ",
             "  → 양수 계수를 γ만큼 축소",
             "",
             "z < -γ:   결과 = z + γ",
             "  → 음수 계수를 γ만큼 축소",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(2.3),
         "Ridge vs Lasso 축소 방식 비교", [
             "Ridge: 비례 축소  β̂ = dⱼ/(dⱼ+λ) · β̂_OLS",
             "  → 0에 가깝게 만들지만 절대 0이 아님",
             "Lasso: 이동 축소  β̂ = sign(β̂_OLS)·max(|β̂|-λ,0)",
             "  → γ보다 작은 계수는 정확히 0!",
         ], ACCENT_CYAN, ACCENT_CYAN)
add_code_block(s, Inches(6.8), Inches(4.6), Inches(6), Inches(2.2), [
    "# Python 구현 (1줄!)",
    "def soft_threshold(z, gamma):",
    "    return np.sign(z) * np.maximum(",
    "        np.abs(z) - gamma, 0.0)",
], font_size=13)

# 슬라이드 20: L1 vs L2 기하학
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "L1 vs L2 기하학적 비교", "왜 Lasso는 변수 선택이 가능한가?")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
         "L2 제약 (Ridge) = 원(Circle)", [
             "제약: β₁² + β₂² ≤ t",
             "",
             "RSS 등고선(타원)과 원의 접점:",
             "  → 일반적으로 좌표축 위에 있지 않음",
             "  → 계수가 정확히 0이 되기 어려움",
             "",
             "결과: 모든 계수가 비영(非零)",
             "  → 변수 선택 불가능",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(5.0),
         "L1 제약 (Lasso) = 마름모(Diamond)", [
             "제약: |β₁| + |β₂| ≤ t",
             "",
             "마름모의 꼭짓점이 좌표축 위에 위치!",
             "  → RSS 등고선과의 접점이",
             "    꼭짓점에서 만날 확률이 높음",
             "  → 꼭짓점: 하나 이상의 좌표 = 0",
             "",
             "결과: 일부 계수가 정확히 0",
             "  → 변수가 자동으로 선택됨!",
         ], ACCENT_RED, ACCENT_RED)

# 슬라이드 21: 베이지안 해석 비교
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "베이지안 해석: Ridge vs Lasso", "사전분포의 차이가 변수 선택을 결정")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
         "Ridge = 정규분포 사전분포", [
             "p(β) = (1/√2πτ²) exp(-β²/2τ²)",
             "",
             "종 모양: 0 근처에서 부드럽게 감소",
             "→ 계수를 0 근처로 축소하지만",
             "   정확히 0으로는 만들지 못함",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.0),
         "Lasso = 라플라스 사전분포", [
             "p(β) = (λ/2) exp(-λ|β|)",
             "",
             "원점에서 뾰족 + 꼬리 두꺼움",
             "→ 작은 계수에 높은 사전 확률 (희소성)",
             "→ 큰 계수에도 적당한 확률 (편향 제한)",
         ], ACCENT_RED, ACCENT_RED)
add_text(s, Inches(0.8), Inches(5.5), Inches(11), Inches(1.0),
         "핵심: 라플라스 분포가 원점에서 더 뾰족하기 때문에 작은 계수를 0으로 만드는 사전 믿음이 더 강함",
         font_size=16, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 6: Elastic Net
# ============================================================
section_divider("Elastic Net", "L1 + L2 혼합 정규화: Zou & Hastie (2005)", 6, ACCENT_CYAN)

# 슬라이드 23: Elastic Net
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 6, "Elastic Net: Lasso의 두 가지 한계 해결", "Zou & Hastie (2005)")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "Lasso의 한계", [
             "1. p > n 문제: 변수 수 > 표본 수이면",
             "   Lasso는 최대 n개 변수만 선택 가능",
             "2. 그룹화 효과 부재: 상관된 변수 중",
             "   하나만 임의로 선택, 나머지 제거",
         ], ACCENT_RED, ACCENT_RED)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "Elastic Net의 해결", [
             "min (1/2n)||y-Xβ||² + λ[α||β||₁ + (1-α)/2 ||β||₂²]",
             "",
             "α = 1: 순수 Lasso",
             "α = 0: 순수 Ridge",
             "0 < α < 1: Elastic Net (최적의 절충)",
         ], ACCENT_CYAN, ACCENT_CYAN)
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         "그룹화 효과 정리 (Zou & Hastie, 2005)", [
             "|β̂ᵢ - β̂ⱼ| ≤ [1 / λ(1-α)] · √(2(1-r)) · ||y||",
             "→ 상관관계 높을수록 (r→1) 두 계수의 차이가 작아짐",
             "→ Elastic Net은 '상관된 변수들은 함께 선택하거나 함께 제거'하는 원칙을 따름",
             "→ 순수 Lasso (α=1)에서는 분모=0이 되어 이 성질 성립하지 않음",
         ], ACCENT_GREEN, ACCENT_GREEN)

# 슬라이드 24: 세 방법 종합 비교
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 6, "Ridge vs Lasso vs Elastic Net 종합 비교")
add_table_slide(s,
    ["특성", "Ridge (L2)", "Lasso (L1)", "Elastic Net (L1+L2)"],
    [
        ["패널티", "λ||β||₂²", "λ||β||₁", "λ[α||β||₁+(1-α)||β||₂²/2]"],
        ["변수 선택", "X (모든 계수 비영)", "O (일부=0)", "O (일부=0)"],
        ["그룹화 효과", "O (동시 축소)", "X (하나만 선택)", "O (동시 선택)"],
        ["닫힌 형태 해", "O", "X", "X"],
        ["p>n 지원", "O", "최대 n개", "O"],
        ["베이지안 사전", "정규분포", "라플라스", "정규+라플라스"],
        ["제약 영역", "원", "마름모", "둥근 마름모"],
    ],
    Inches(0.5), Inches(2.0), [2.5, 2.5, 2.5, 3.5],
    row_height=0.5, font_size=12, header_font_size=13)

# ============================================================
# SECTION 7: 좌표 하강법
# ============================================================
section_divider("좌표 하강법", "glmnet 알고리즘: Friedman, Hastie, Tibshirani (2010)", 7, ACCENT_GREEN)

# 슬라이드 26: 좌표 하강법
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 7, "좌표 하강법 (Coordinate Descent)", "한 번에 하나의 변수만 최적화")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(5.0),
         "업데이트 규칙", [
             "j번째 좌표 업데이트:",
             "",
             "β̃ⱼ = S(ρⱼ, αλ) / (Xⱼᵀ Xⱼ/n + λ(1-α))",
             "",
             "ρⱼ = Xⱼᵀ rⱼ / n",
             "rⱼ = y - Σ_{k≠j} Xₖβₖ  (부분 잔차)",
             "S(z,γ) = sign(z)·max(|z|-γ,0)",
             "",
             "알고리즘:",
             "  for iteration in range(max_iter):",
             "    for j in range(n_features):",
             "      부분 잔차 계산 → 연성 임계값 적용",
             "    수렴 확인 (max|β_new - β_old| < tol)",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(5.0),
         "핵심 최적화 기법 3가지", [
             "1. 따뜻한 시작 (Warm Start)",
             "  λ₁ > λ₂ > ... 순차적으로 풀 때",
             "  이전 해를 초기값으로 사용",
             "  인접 λ의 해가 유사 → 빠른 수렴",
             "",
             "2. 활성 집합 (Active Set)",
             "  0이 아닌 계수만 우선 업데이트",
             "  Lasso에서 대부분 계수=0이므로 효과적",
             "",
             "3. λ_max 계산",
             "  λ_max = (1/αn) max_j |Xⱼᵀ(y-ȳ)|",
             "  모든 계수=0이 되는 최소 λ",
             "  여기서 시작해 λ를 줄여감",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# SECTION 8: SCAD와 Oracle 특성
# ============================================================
section_divider("SCAD와 Oracle 특성", "Fan & Li (2001): 비볼록 패널티와 이상적 성질", 8, ACCENT_PURPLE)

# 슬라이드 28: SCAD
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 8, "SCAD: 오라클 성질을 가진 패널티", "Fan & Li (2001)")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "Lasso의 근본적 한계", [
             "L1 패널티는 계수 크기에 관계없이",
             "동일한 양(λ)만큼 축소:",
             "→ 작은 계수: 적절히 0으로 제거 (OK)",
             "→ 큰 계수: 불필요하게 과도한 축소 (문제!)",
         ], ACCENT_RED, ACCENT_RED)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "오라클 성질 (Oracle Properties)", [
             "1. 일관성 있는 변수 선택:",
             "   P(선택된 변수 = 진정한 변수) → 1",
             "2. 점근적 정규성:",
             "   비영 계수의 추정이 OLS와 동일 분포",
             "Lasso는 오라클 성질을 만족하지 못함!",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_table_slide(s,
    ["구간", "패널티 동작", "효과"],
    [
        ["|θ| ≤ λ", "L1과 동일", "작은 계수를 0으로 제거"],
        ["λ < |θ| ≤ aλ", "패널티 점차 감소", "중간 계수 점진적 축소"],
        ["|θ| > aλ", "패널티 미분 = 0", "큰 계수에 추가 축소 없음!"],
    ],
    Inches(0.5), Inches(4.8), [2.5, 3.5, 5.5],
    row_height=0.55, font_size=13)

# 슬라이드 29: 정규화 방법의 계보
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 8, "정규화 방법의 계보", "좋은 패널티의 3가지 조건")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "좋은 패널티의 3가지 조건 (Fan & Li)", [
             "1. 비편향성: 큰 계수에 편향 없을 것",
             "2. 희소성: 작은 계수를 0으로 만들 것",
             "3. 연속성: 추정량이 데이터의 연속 함수일 것",
             "세 조건 동시 만족 → 반드시 비볼록!",
         ], ACCENT_PURPLE, ACCENT_PURPLE)
add_table_slide(s,
    ["방법", "연도", "패널티", "변수선택", "오라클", "볼록"],
    [
        ["Ridge", "1970", "L2", "X", "X", "O"],
        ["Lasso", "1996", "L1", "O", "X", "O"],
        ["SCAD", "2001", "비볼록", "O", "O", "X"],
        ["Elastic Net", "2005", "L1+L2", "O", "X", "O"],
        ["Adaptive Lasso", "2006", "가중 L1", "O", "O", "O"],
        ["MCP", "2010", "비볼록", "O", "O", "X"],
    ],
    Inches(0.5), Inches(4.8), [2.0, 1.2, 1.8, 1.5, 1.5, 1.5],
    row_height=0.45, font_size=12)

# ============================================================
# SECTION 9: 회귀 진단
# ============================================================
section_divider("회귀 진단", "잔차 분석, Q-Q Plot, Cook's Distance", 9, ACCENT_ORANGE)

# 슬라이드 31: 잔차 분석
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 9, "잔차 분석 (Residual Analysis)", "모델 적합성 평가")
add_table_slide(s,
    ["잔차 유형", "공식", "용도"],
    [
        ["일반 잔차", "eᵢ = yᵢ - ŷᵢ", "기본 분석"],
        ["표준화 잔차", "rᵢ = eᵢ / (σ̂√(1-hᵢᵢ))", "이상치 탐지"],
        ["스튜던트화 잔차", "tᵢ = eᵢ / (σ̂₍ᵢ₎√(1-hᵢᵢ))", "정밀 이상치 탐지"],
    ],
    Inches(0.5), Inches(2.0), [2.5, 4.5, 4.5],
    row_height=0.55, font_size=13)
add_card(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(3.0),
         "잔차 플롯 확인 사항", [
             "1. 예측값 vs 잔차: 패턴이 없어야 함 (등분산성)",
             "2. 잔차의 정규성: Q-Q plot으로 확인",
             "3. 자기상관: 시계열 데이터에서 잔차 독립성",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_card(s, Inches(6.7), Inches(4.0), Inches(6.1), Inches(3.0),
         "Q-Q Plot", [
             "x축: 이론적 정규분포의 분위수",
             "y축: 잔차의 분위수",
             "점들이 45도 직선 위 → 정규분포 따름",
             "직선 벗어남 → 분포 왜곡 판단",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 슬라이드 32: Cook's Distance
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 9, "Cook's Distance", "영향력 있는 관측치 탐지")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(3.0),
         "Cook's Distance 공식", [
             "Dᵢ = (rᵢ² / p) · (hᵢᵢ / (1-hᵢᵢ)²)",
             "",
             "rᵢ: 표준화 잔차",
             "hᵢᵢ: 레버리지 (햇 행렬 대각 원소)",
             "p: 모수의 수",
             "",
             "판정: Dᵢ > 1 → 영향력 큰 관측치",
             "보수적: Dᵢ > 4/n",
         ], ACCENT_RED, ACCENT_RED)
add_table_slide(s,
    ["진단 항목", "확인 방법", "위반 시 조치"],
    [
        ["선형성", "예측값 vs 잔차 플롯", "변수 변환, 다항 회귀"],
        ["등분산성", "예측값 vs 잔차 플롯", "WLS, 변수 변환"],
        ["정규성", "Q-Q plot, Shapiro-Wilk", "변수 변환, 로버스트 회귀"],
        ["독립성", "Durbin-Watson 검정", "GLS, 시계열 모델"],
        ["다중공선성", "VIF, 조건수", "변수 제거, Ridge"],
        ["이상치", "Cook's distance", "제거 또는 로버스트 회귀"],
    ],
    Inches(6.8), Inches(2.0), [2.0, 2.0, 2.0],
    row_height=0.5, font_size=11, header_font_size=12)

# ============================================================
# SECTION 10: 실습 - Ridge 구현
# ============================================================
section_divider("실습: Ridge 회귀 구현", "닫힌 형태 해를 직접 코딩", 10, ACCENT_BLUE)

# 슬라이드 34: Ridge 직접 구현
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 10, "Ridge 회귀 직접 구현", "01_ridge_regression.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(4.8), [
    "class RidgeRegressionFromScratch:",
    '    """닫힌 형태 해: w = (X\'X + λI)⁻¹X\'y"""',
    "",
    "    def __init__(self, alpha=1.0):",
    "        self.alpha = alpha",
    "",
    "    def fit(self, X, y):",
    "        n, p = X.shape",
    "        self.X_mean_ = np.mean(X, axis=0)",
    "        self.y_mean_ = np.mean(y)",
    "        Xc = X - self.X_mean_",
    "        yc = y - self.y_mean_",
    "",
    "        # 핵심: (X'X + αI)β = X'y",
    "        XtX = Xc.T @ Xc",
    "        reg = XtX + self.alpha * np.eye(p)",
    "        Xty = Xc.T @ yc",
    "        self.coef_ = np.linalg.solve(reg, Xty)",
    "        self.intercept_ = (",
    "            self.y_mean_ - self.X_mean_ @ self.coef_)",
    "        return self",
], font_size=11)
add_card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.8),
         "구현 핵심 포인트", [
             "1. 데이터 중심화 (centering)",
             "   Xc = X - mean(X)",
             "   yc = y - mean(y)",
             "   → 절편 별도 처리",
             "",
             "2. np.linalg.solve 사용",
             "   역행렬(inv) 대신 선형시스템 풀기",
             "   → 수치적으로 더 안정적",
             "",
             "3. sklearn과 결과 일치 확인",
             "   계수 최대 차이: ~10⁻¹² 수준",
             "",
             "4. 다중공선성 효과 확인",
             "   조건수: Ridge로 크게 감소",
             "   κ(XᵀX+λI) << κ(XᵀX)",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 슬라이드 35: Ridge 정규화 경로
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 10, "정규화 경로 (Ridge Trace) & 교차 검증", "최적 λ 선택")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(2.5), [
    "# 정규화 경로 시각화",
    "alphas = np.logspace(-2, 4, 200)",
    "coefs_path = []",
    "for alpha in alphas:",
    "    ridge = RidgeFromScratch(alpha=alpha)",
    "    ridge.fit(X_scaled, y)",
    "    coefs_path.append(ridge.coef_.copy())",
], font_size=12)
add_code_block(s, Inches(0.6), Inches(5.0), Inches(6), Inches(2.0), [
    "# 교차 검증으로 최적 λ",
    "ridge_cv = RidgeCV(alphas=alphas, cv=5,",
    "    scoring='neg_mean_squared_error')",
    "ridge_cv.fit(X_scaled, y)",
    "print(f'최적 λ: {ridge_cv.alpha_:.4f}')",
], font_size=12)
add_card(s, Inches(6.5), Inches(2.2), Inches(6.3), Inches(4.8),
         "Ridge Trace 해석", [
             "x축: log₁₀(λ), y축: 각 회귀 계수",
             "",
             "λ = 0 (왼쪽): OLS 추정치",
             "  → 다중공선성으로 불안정 (큰 진동)",
             "",
             "λ 증가: 모든 계수가 0으로 수렴",
             "  → 편향 증가, 분산 감소",
             "",
             "최적 λ: CV MSE가 최소인 점",
             "  → 편향-분산 트레이드오프의 최적점",
             "",
             "OLS vs Ridge 비교:",
             "  OLS 계수의 L2 노름 >> Ridge",
             "  Ridge로 조건수 크게 감소",
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# SECTION 11: 실습 - Lasso 구현
# ============================================================
section_divider("실습: Lasso 회귀 구현", "좌표 하강법 직접 코딩", 11, ACCENT_RED)

# 슬라이드 37: Lasso 직접 구현
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 11, "Lasso 좌표 하강법 직접 구현", "02_lasso_regression.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(4.8), [
    "class LassoFromScratch:",
    "    def fit(self, X, y):",
    "        n, p = X.shape",
    "        Xc = X - np.mean(X, axis=0)",
    "        yc = y - np.mean(y)",
    "        col_norms = np.sum(Xc**2, axis=0)/n",
    "        w = np.zeros(p)",
    "        residual = yc.copy()",
    "",
    "        for _ in range(self.max_iter):",
    "            w_old = w.copy()",
    "            for j in range(p):",
    "                # 부분 잔차에 j번째 기여 복원",
    "                residual += Xc[:,j] * w[j]",
    "                # X_j와 부분 잔차의 내적/n",
    "                rho = Xc[:,j] @ residual / n",
    "                # 연성 임계값 적용!",
    "                w[j] = soft_threshold(rho,",
    "                    self.alpha) / col_norms[j]",
    "                # 잔차 업데이트",
    "                residual -= Xc[:,j] * w[j]",
    "            if np.max(np.abs(w-w_old))<self.tol:",
    "                break",
], font_size=10)
add_card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.8),
         "좌표 하강법 핵심 이해", [
             "1. 부분 잔차 (partial residual)",
             "   rⱼ = y - Σ_{k≠j} Xₖwₖ",
             "   '다른 변수들의 기여를 뺀 잔차'",
             "",
             "2. 효율적 잔차 업데이트",
             "   residual += Xⱼ·w_old[j]  (복원)",
             "   residual -= Xⱼ·w_new[j]  (갱신)",
             "   → 매번 전체 계산 대신 차이만 반영",
             "",
             "3. 변수 선택 확인",
             "   0인 계수 수 → 실제와 비교",
             "   직접 구현 vs sklearn 결과 일치",
             "",
             "4. 정규화 경로",
             "   λ 감소 → 선택 변수 수 증가",
             "   빨간 실선: 실제 비영 계수 수",
         ], ACCENT_RED, ACCENT_RED)

# ============================================================
# SECTION 12: 실습 - Elastic Net 구현
# ============================================================
section_divider("실습: Elastic Net 구현", "L1+L2 혼합 정규화 + 그룹화 효과", 12, ACCENT_CYAN)

# 슬라이드 39: Elastic Net 직접 구현
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 12, "Elastic Net 좌표 하강법 구현", "03_elastic_net.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(3.0), [
    "# Elastic Net 업데이트 규칙",
    "# Lasso와의 차이: 분모에 L2 항 추가!",
    "l1_penalty = alpha * l1_ratio",
    "l2_penalty = alpha * (1 - l1_ratio)",
    "",
    "for j in range(p):",
    "    residual += Xc[:,j] * w[j]",
    "    rho = Xc[:,j] @ residual / n",
    "    # 핵심: 분모에 l2_penalty 추가",
    "    denom = col_norms[j] + l2_penalty",
    "    w[j] = soft_threshold(rho, l1_penalty) / denom",
    "    residual -= Xc[:,j] * w[j]",
], font_size=12)
add_card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.8),
         "그룹화 효과 실험 결과", [
             "상관된 특성 그룹 데이터:",
             "  그룹1 (x0,x1,x2): 높은 상관, β=2.0",
             "  그룹2 (x3,x4): 높은 상관, β=-1.5",
             "  노이즈 (x5~x14): β=0",
             "",
             "Ridge: 모든 계수 비영, 그룹 내 유사",
             "Lasso: 변수 선택 O, 그룹 중 하나만!",
             "Elastic Net: 변수 선택 + 그룹 동시 선택",
             "",
             "그룹 내 계수 분산:",
             "  Ridge < Elastic Net << Lasso",
             "  → EN이 그룹화 효과 가장 우수",
         ], ACCENT_CYAN, ACCENT_CYAN)

# ============================================================
# SECTION 13: 응용사례
# ============================================================
section_divider("응용사례", "부동산 가격 예측 & 광고비-매출 분석", 13, ACCENT_GREEN)

# 슬라이드 41: 부동산 가격 예측
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 13, "응용 1: California Housing 가격 예측", "sklearn 내장 데이터셋")
add_table_slide(s,
    ["모델", "RMSE", "R²", "비영 계수 수"],
    [
        ["OLS", "~0.73", "~0.58", "8"],
        ["Ridge (α=1)", "~0.73", "~0.58", "8"],
        ["Lasso (α=0.01)", "~0.74", "~0.57", "6~7"],
    ],
    Inches(0.5), Inches(2.0), [3.0, 2.0, 2.0, 2.5],
    row_height=0.55, font_size=13)
add_card(s, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.8),
         "주요 계수 해석", [
             "MedInc (중위 소득): 강한 양(+)",
             "HouseAge (주택 연령): 양(+) - 도심 효과",
             "AveRooms (평균 방 수): 양(+) - 넓은 집",
             "Latitude (위도): 음(-) - 북쪽 가격 하락",
             "Longitude (경도): 음(-) - 해안 쪽 높음",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(6.7), Inches(4.2), Inches(6.1), Inches(2.8),
         "한계점 & 교훈", [
             "R² ≈ 0.58: 선형 모델로 58%만 설명",
             "위도/경도의 비선형 지리 패턴 포착 어려움",
             "→ RF, XGBoost 등 비선형 모델로 개선",
             "",
             "교훈: 선형회귀는 해석 가능한 baseline",
             "성능 vs 해석가능성 트레이드오프",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 슬라이드 42: 광고비-매출 분석
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 13, "응용 2: 광고비-매출 회귀 분석", "마케팅 분석 사례")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
         "데이터 구조", [
             "TV 광고비: 10~300 (주효과)",
             "라디오 광고비: 0~50 (주효과)",
             "신문 광고비: 0~100 (약한 효과)",
             "매출 = 0.05·TV + 0.1·Radio + 0.005·News",
             "        + 0.001·TV·Radio + 5 + 노이즈",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.0),
         "Elastic Net 결과 (변수 선택)", [
             "TV 광고: 선택 (가장 큰 양의 효과)",
             "라디오: 선택 (단위당 효율 가장 높음)",
             "신문: 거의 0으로 축소 → 실질적 제거!",
             "",
             "마케팅 시사점:",
             "  신문 예산 → TV/라디오로 재배분이 합리적",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_text(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
         "핵심: 선형회귀의 해석 가능성이 마케팅 의사결정에 직접 활용됨. 각 계수가 '광고비 1단위 증가 시 매출 변화'를 의미",
         font_size=15, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 14: 핵심 요약
# ============================================================
section_divider("핵심 요약 + 복습 질문", "4장 전체 정리", 14, ACCENT_BLUE)

# 슬라이드 44: 이론 요약표
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 14, "이론 요약표", "핵심 수식과 특성 한눈에 보기")
add_table_slide(s,
    ["개념", "핵심 수식", "핵심 특성"],
    [
        ["OLS", "β̂ = (XᵀX)⁻¹Xᵀy", "BLUE, 비편향"],
        ["Ridge", "β̂ = (XᵀX+λI)⁻¹Xᵀy", "L2, 축소, 닫힌 해"],
        ["Lasso", "min ||y-Xβ||²/2n + λ||β||₁", "L1, 변수 선택"],
        ["Elastic Net", "L1+L2 혼합", "그룹화 효과"],
        ["SCAD", "비볼록 패널티", "오라클 성질"],
    ],
    Inches(0.5), Inches(2.0), [2.0, 4.5, 4.5],
    row_height=0.55, font_size=13)

# 슬라이드 45: 핵심 논문 연대기
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 14, "핵심 논문 연대기 & 방법 선택 가이드")
add_table_slide(s,
    ["연도", "논문", "핵심 기여"],
    [
        ["1970", "Hoerl & Kennard", "Ridge, 편향-분산 트레이드오프"],
        ["1996", "Tibshirani", "Lasso, L1 변수 선택"],
        ["2001", "Fan & Li", "SCAD, 오라클 성질"],
        ["2005", "Zou & Hastie", "Elastic Net, 그룹화 효과"],
        ["2010", "Friedman et al.", "glmnet, 좌표 하강법"],
    ],
    Inches(0.5), Inches(2.0), [1.2, 3.0, 7.3],
    row_height=0.5, font_size=13)
add_card(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
         "방법 선택 가이드", [
             "다중공선성 심함? → Ridge 또는 Elastic Net",
             "변수 선택 필요? → Lasso 또는 Elastic Net",
             "상관된 변수 그룹? → Elastic Net (그룹화 효과)",
             "p > n? → Elastic Net",
             "변수 수 적고 다중공선성 없음? → OLS",
         ], ACCENT_GREEN, ACCENT_GREEN)

# 슬라이드 46: 복습 질문 (1/2)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 14, "복습 질문 (1/2)", "수식 유도 + 이론")
questions_1 = [
    "Q1. RSS를 β에 대해 행렬 미분하여 정규방정식을 유도하라",
    "Q2. 가우스-마르코프 정리의 한계: 왜 Ridge가 OLS보다 나을 수 있는가?",
    "Q3. Ridge의 닫힌 형태 해를 유도하라. λI가 조건수를 개선하는 이유?",
    "Q4. L1(마름모) vs L2(원) 제약 영역에서 변수 선택 차이를 기하학적으로 설명",
    "Q5. 연성 임계값 연산자 S(z,γ)의 세 구간 동작과 좌표 하강법에서의 역할",
]
add_bullet_list(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                questions_1, font_size=17, color=LIGHT_GRAY, spacing=Pt(14))

# 슬라이드 47: 복습 질문 (2/2)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 14, "복습 질문 (2/2)", "심화 + 진단")
questions_2 = [
    "Q6. Elastic Net의 그룹화 효과 정리를 설명하라",
    "Q7. 좌표 하강법의 효율성: warm start와 active set 전략",
    "Q8. SCAD의 오라클 성질 두 가지 조건과 Lasso 대비 장점",
    "Q9. Ridge(정규분포)와 Lasso(라플라스)의 베이지안 해석",
    "Q10. 잔차 분석: (a) 깔때기 패턴 (b) Q-Q 꼬리 벗어남 (c) Cook's D > 1",
]
add_bullet_list(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                questions_2, font_size=17, color=LIGHT_GRAY, spacing=Pt(14))

# ============================================================
# 슬라이드 48: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG, ACCENT_BLUE)
add_accent_line(s, Inches(4), Inches(2.5), Inches(5.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.8), Inches(11.333), Inches(1),
         "Thank You", font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.0), Inches(11.333), Inches(0.5),
         "Chapter 4: 선형회귀 (Linear Regression)", font_size=22, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
         "다음 장: 로지스틱 회귀 (Logistic Regression)", font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 저장
# ============================================================
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "4장_선형회귀_강의PPT_확장.pptx")
prs.save(out_path)
print(f"[완료] {out_path}")
print(f"[슬라이드 수] {len(prs.slides)}")
