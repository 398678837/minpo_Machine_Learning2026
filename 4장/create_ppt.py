"""4장 선형회귀 - 강의 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    card = add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))


def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)


# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "기계학습 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "4장: 선형회귀 (Linear Regression)", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "OLS, Ridge, Lasso, Elastic Net, SCAD", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
         "핵심 키워드: 정규방정식 · 가우스-마르코프 · 다중공선성 · L1/L2 정규화 · 좌표 하강법 · 오라클 성질 · 회귀 진단",
         font_size=14, color=DARK_GRAY)
add_shape(s, Inches(0), Inches(6.3), prs.slide_width, Inches(1.2), RGBColor(0x0A, 0x0A, 0x1A))
add_text(s, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
         "2026년 1학기", font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 2: 목차
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (Contents)")

toc = [
    ("01", "역사와 의의 - 최소제곱법의 탄생", ACCENT_BLUE),
    ("02", "수학적 정의 - 선형회귀 모델", ACCENT_CYAN),
    ("03", "OLS 이론 - 정규방정식과 기하학적 해석", ACCENT_GREEN),
    ("04", "가우스-마르코프 정리와 BLUE", ACCENT_PURPLE),
    ("05", "다중공선성과 VIF", ACCENT_ORANGE),
    ("06", "Ridge 회귀 - L2 정규화", ACCENT_RED),
    ("07", "Lasso 회귀 - L1 정규화와 변수 선택", ACCENT_BLUE),
    ("08", "Elastic Net - L1+L2 혼합", ACCENT_CYAN),
    ("09", "좌표 하강법과 SCAD", ACCENT_GREEN),
    ("10", "회귀 진단과 실습", ACCENT_PURPLE),
    ("11", "응용사례와 핵심 요약", ACCENT_ORANGE),
]
for i, (num, title, color) in enumerate(toc):
    y = Inches(1.95) + Inches(0.46) * i
    add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.38), color, radius=True)
    add_text(s, Inches(1.2), y, Inches(0.55), Inches(0.38), num,
             font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), y, Inches(8), Inches(0.38), title,
             font_size=17, color=WHITE)

# ============================================================
# 슬라이드 3: 역사와 의의 - 최소제곱법의 탄생
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "역사와 의의 - 최소제곱법의 탄생", "The Birth of Least Squares")

# Legendre
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "르장드르 (Adrien-Marie Legendre, 1805)", font_size=15, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.2), [
    "최소제곱법을 최초로 공식 발표",
    "혜성 궤도 결정 문제에서 잔차 제곱합 최소화 원리 제안",
    "저서: Nouvelles methodes pour la determination des orbites des cometes",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# Gauss
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "가우스 (Carl Friedrich Gauss, 1809)", font_size=15, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.2), [
    "1795년부터 사용했다고 주장 (우선권 논쟁)",
    "정규 분포 가정 하에서 MLE와 동일함을 증명",
    "확률론적 정당성 부여 (Theoria Motus, 1809)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 통계학에서 기계학습으로
add_text(s, Inches(0.6), Inches(4.5), Inches(5), Inches(0.4),
         "통계학에서 기계학습으로", font_size=18, color=WHITE, bold=True)

roles = [
    ("기초 알고리즘", "지도학습에서 연속값 예측의\n가장 기본적인 모델", ACCENT_BLUE),
    ("이론적 토대", "비용 함수 최적화, 편향-분산\n트레이드오프의 출발점", ACCENT_GREEN),
    ("확장의 기반", "로지스틱 회귀, 신경망,\n커널 방법의 구성 요소", ACCENT_ORANGE),
    ("해석 가능성", "회귀 계수를 통해 특성의\n영향력 직접 해석 가능", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(roles):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(5.0), Inches(2.9), Inches(2.1), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(5.1), Inches(2.5), Inches(0.4),
             title, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(5.5), Inches(2.5), Inches(1.3),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 4: 수학적 정의
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "선형회귀의 수학적 정의", "Mathematical Formulation")

# 수식 박스
add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.4), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.6), Inches(2.3), Inches(12.1), Inches(0.5),
         "y = b0 + b1*x1 + b2*x2 + ... + bp*xp + e", font_size=22, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.85), Inches(12.1), Inches(0.5),
         "행렬 표기:   y = Xb + e", font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 구성 요소 설명
components = [
    ("y", "R^n", "반응 벡터\n(n개 관측치)", ACCENT_BLUE),
    ("X", "R^(n x (p+1))", "설계 행렬\n(첫 열은 절편용 1 벡터)", ACCENT_CYAN),
    ("b", "R^(p+1)", "회귀 계수 벡터\n(추정 대상)", ACCENT_GREEN),
    ("e", "R^n", "오차 벡터\nei ~ iid N(0, s^2)", ACCENT_ORANGE),
]
for i, (sym, space, desc, color) in enumerate(components):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(4.0), Inches(2.9), Inches(2.8), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.9), Inches(3.75), Inches(1.1), Inches(0.5), color, radius=True)
    add_text(s, x + Inches(0.9), Inches(3.75), Inches(1.1), Inches(0.5),
             sym, font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(4.4), Inches(2.5), Inches(0.4),
             space, font_size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(4.9), Inches(2.5), Inches(1.5),
             desc, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 5: OLS 이론 (1) - 정규방정식 유도
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "OLS 이론 (1) - 정규방정식 유도", "Normal Equation Derivation")

# RSS 정의
add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.8), RGBColor(0x15, 0x15, 0x30))
add_text(s, Inches(0.6), Inches(2.25), Inches(12.1), Inches(0.7),
         "RSS(b) = ||y - Xb||^2 = (y - Xb)^T (y - Xb) = y^Ty - 2b^T X^Ty + b^T X^TX b",
         font_size=16, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# 유도 과정
steps = [
    ("Step 1: 행렬 미분", "dRSS/db = -2X^Ty + 2X^TXb", ACCENT_BLUE),
    ("Step 2: 0으로 놓기", "X^TXb = X^Ty  (정규방정식)", ACCENT_GREEN),
    ("Step 3: 해 구하기", "b_OLS = (X^TX)^(-1) X^Ty", ACCENT_ORANGE),
]
for i, (step, formula, color) in enumerate(steps):
    y_pos = Inches(3.3) + Inches(0.9) * i
    add_shape(s, Inches(0.6), y_pos, Inches(3.5), Inches(0.75), CARD_BG, color, radius=True)
    add_text(s, Inches(0.8), y_pos + Inches(0.05), Inches(3.1), Inches(0.35),
             step, font_size=13, color=color, bold=True)
    add_text(s, Inches(0.8), y_pos + Inches(0.38), Inches(3.1), Inches(0.35),
             formula, font_size=13, color=WHITE)

# 핵심 결과 박스
add_shape(s, Inches(4.5), Inches(3.3), Inches(8.2), Inches(1.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(4.5), Inches(3.35), Inches(8.2), Inches(0.5),
         "b_OLS = (X^TX)^(-1) X^Ty", font_size=28, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.7), Inches(3.95), Inches(7.8), Inches(0.7),
         "X^TX가 가역(invertible)일 때 유일한 해가 존재\n헤시안: d^2RSS/db^2 = 2X^TX (양의 반정치) => RSS는 볼록 함수 => 전역 최솟값",
         font_size=13, color=LIGHT_GRAY)

# 행렬 미분 규칙
add_shape(s, Inches(4.5), Inches(5.1), Inches(8.2), Inches(1.7), CARD_BG, radius=True)
add_text(s, Inches(4.7), Inches(5.2), Inches(7.8), Inches(0.4),
         "사용된 행렬 미분 공식", font_size=14, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(4.7), Inches(5.6), Inches(7.8), Inches(1.0), [
    "d/dx (a^T x) = a",
    "d/dx (x^T A x) = 2Ax  (A가 대칭일 때)",
    "=> RSS의 2차 형식 구조를 활용한 직접 미분",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# 슬라이드 6: OLS 이론 (2) - 기하학적 해석
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "OLS 이론 (2) - 기하학적 해석", "Geometric Interpretation: Orthogonal Projection")

# 직교 사영
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "직교 사영 (Orthogonal Projection)", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.8), [
    "y: R^n 공간의 벡터",
    "C(X): X의 열들이 생성하는 부분공간 (열공간)",
    "y_hat = Xb_OLS: y를 C(X)에 직교 사영한 것",
    "잔차 e = y - y_hat 는 C(X)에 직교: X^Te = 0",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(5))

# 햇 행렬
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "햇 행렬 (Hat Matrix)", font_size=16, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(0.5),
         "H = X(X^TX)^(-1)X^T", font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(7.1), Inches(3.4), Inches(5.2), Inches(1.2), [
    "y_hat = Hy  (y에 '모자'를 씌워 예측값을 만듦)",
    "사영 행렬: H^2 = H, H^T = H",
    "h_ii (대각 원소) = 레버리지 (영향력 측정)",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 핵심 성질
add_shape(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.6), Inches(5.05), Inches(12.1), Inches(0.4),
         "피타고라스 정리:  ||y||^2 = ||y_hat||^2 + ||e||^2   =>   TSS = ESS + RSS",
         font_size=18, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.4),
         "R^2 = cos^2(theta) : y와 y_hat 사이 각도의 코사인 제곱 = 모델의 설명력",
         font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 잔차 직교성 요약
add_shape(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.9), CARD_BG, radius=True)
add_text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.7),
         "잔차 직교성: X^T(y - Xb_hat) = 0  <=>  정규방정식 X^TXb = X^Ty 와 동일한 조건",
         font_size=15, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 7: 가우스-마르코프 정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "가우스-마르코프 정리", "Gauss-Markov Theorem: OLS is BLUE")

# BLUE 설명
blue_items = [
    ("B", "Best", "모든 선형 비편향 추정량 중\n분산이 가장 작다", ACCENT_BLUE),
    ("L", "Linear", "b_hat가 y의 선형 함수\nb_hat = Cy (C는 상수 행렬)", ACCENT_CYAN),
    ("U", "Unbiased", "E[b_hat] = b\n(편향이 없다)", ACCENT_GREEN),
    ("E", "Estimator", "b의 추정량\n(모수를 추정하는 통계량)", ACCENT_ORANGE),
]
for i, (letter, word, desc, color) in enumerate(blue_items):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(2.2), Inches(2.9), Inches(2.5), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(1.0), Inches(1.95), Inches(0.9), Inches(0.5), color, radius=True)
    add_text(s, x + Inches(1.0), Inches(1.95), Inches(0.9), Inches(0.5),
             letter, font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.6), Inches(2.5), Inches(0.4),
             word, font_size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(1.3),
             desc, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 4가지 조건
add_text(s, Inches(0.6), Inches(4.9), Inches(3), Inches(0.4),
         "가우스-마르코프 조건 4가지", font_size=15, color=WHITE, bold=True)

conditions = [
    "1. E[ei] = 0  (오차의 기대값이 0)",
    "2. Var(ei) = s^2  (등분산성)",
    "3. Cov(ei, ej) = 0, i != j  (오차 간 무상관)",
    "4. X는 고정된(non-random) 행렬",
]
add_bullet_list(s, Inches(0.6), Inches(5.3), Inches(5.5), Inches(2.0),
                conditions, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 분산-공분산
add_shape(s, Inches(6.5), Inches(5.0), Inches(6.2), Inches(1.2), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(6.5), Inches(5.05), Inches(6.2), Inches(0.5),
         "Var(b_OLS) = s^2 (X^TX)^(-1)", font_size=22, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.7), Inches(5.6), Inches(5.8), Inches(0.5),
         "OLS 추정량의 분산-공분산 행렬", font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 중요한 제한
add_shape(s, Inches(6.5), Inches(6.4), Inches(6.2), Inches(0.8), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(6.7), Inches(6.45), Inches(5.8), Inches(0.7),
         "주의: 가우스-마르코프는 '비편향' 중 최적. 편향을 허용하면\nMSE 관점에서 OLS보다 나은 추정량 존재 => Ridge 회귀의 핵심 동기!",
         font_size=12, color=ACCENT_RED)

# ============================================================
# 슬라이드 8: 다중공선성
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "다중공선성 (Multicollinearity)", "정의, VIF, 조건수, 해결 방법")

# 정의
add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "다중공선성이란?", font_size=16, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.6),
         "독립 변수들 사이에 높은 선형 상관관계 존재 => X^TX가 특이에 가까움 => (X^TX)^(-1) 원소 폭증 => Var(b_hat) 폭증 => 계수 추정 불안정",
         font_size=14, color=LIGHT_GRAY)

# VIF 테이블
add_text(s, Inches(0.6), Inches(3.7), Inches(5), Inches(0.4),
         "VIF (Variance Inflation Factor)", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.6), Inches(4.1), Inches(5), Inches(0.4),
         "VIF_j = 1 / (1 - R_j^2)", font_size=16, color=WHITE, bold=True)

vif_headers = ["VIF 값", "해석"]
vif_rows = [
    ["1", "다중공선성 없음"],
    ["1~5", "약한 (일반적으로 허용)"],
    ["5~10", "중간 수준 (주의 필요)"],
    ["10 이상", "심각 (조치 필요)"],
]
for j, h in enumerate(vif_headers):
    x = Inches(0.6) + Inches(2.5) * j
    add_shape(s, x, Inches(4.6), Inches(2.3), Inches(0.45), ACCENT_BLUE)
    add_text(s, x, Inches(4.6), Inches(2.3), Inches(0.45),
             h, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, row in enumerate(vif_rows):
    y_pos = Inches(5.1) + Inches(0.45) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.6) + Inches(2.5) * j
        add_shape(s, x, y_pos, Inches(2.3), Inches(0.4), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y_pos, Inches(2.3), Inches(0.4),
                 cell, font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 조건수 & 해결 방법
add_shape(s, Inches(6.2), Inches(3.7), Inches(6.5), Inches(1.2), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(6.4), Inches(3.8), Inches(6.1), Inches(0.4),
         "조건수 (Condition Number)", font_size=14, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(6.4), Inches(4.2), Inches(6.1), Inches(0.6),
         "k(X^TX) = lambda_max / lambda_min\n조건수 > 30 이면 다중공선성 의심, 역행렬 계산 수치적 불안정",
         font_size=13, color=LIGHT_GRAY)

# 해결 방법
add_shape(s, Inches(6.2), Inches(5.1), Inches(6.5), Inches(1.9), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(6.4), Inches(5.2), Inches(6.1), Inches(0.4),
         "해결 방법 4가지", font_size=14, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(6.4), Inches(5.6), Inches(6.1), Inches(1.3), [
    "1. 변수 제거: VIF가 높은 변수를 수동 제거",
    "2. 주성분 회귀(PCR): PCA로 차원 축소 후 회귀",
    "3. Ridge 회귀: L2 정규화로 조건수 개선",
    "4. Lasso 회귀: L1 정규화로 불필요 변수 자동 제거",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 9: Ridge 회귀 (1) - 핵심 아이디어
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "Ridge 회귀 (1) - 핵심 아이디어", "Hoerl & Kennard (1970): Biased but Better")

# 핵심 통찰
add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.5), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "MSE = Variance + Bias^2", font_size=26, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.7),
         "OLS: 편향=0이지만 분산이 클 수 있음. Ridge: 약간의 편향 도입 => 분산을 더 크게 감소 => 전체 MSE 감소!",
         font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# L2 정규화 수식
add_shape(s, Inches(0.6), Inches(4.0), Inches(5.8), Inches(1.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(4.1), Inches(5.2), Inches(0.4),
         "L2 정규화 최적화 문제", font_size=14, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(0.9), Inches(4.5), Inches(5.2), Inches(0.8),
         "min { ||y - Xb||^2 + lambda * ||b||^2 }\n\n||b||^2 = sum(bj^2) : L2 패널티 (계수 제곱합)",
         font_size=14, color=WHITE)

# 닫힌 형태 해
add_shape(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(1.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.1), Inches(4.1), Inches(5.2), Inches(0.4),
         "닫힌 형태 해 (Closed-form Solution)", font_size=14, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(7.1), Inches(4.5), Inches(5.2), Inches(0.8),
         "b_ridge = (X^TX + lambda*I)^(-1) X^Ty\n\nlambda >= 0 : 정규화 강도 (하이퍼파라미터)",
         font_size=14, color=WHITE)

# 핵심 효과
add_text(s, Inches(0.6), Inches(5.7), Inches(5), Inches(0.4),
         "lambda*I 를 더하는 핵심 효과 3가지", font_size=15, color=WHITE, bold=True)

effects = [
    ("1", "조건수 감소", "k(X^TX+lambda*I) < k(X^TX)", ACCENT_BLUE),
    ("2", "수치적 안정성", "역행렬 계산이 안정화", ACCENT_GREEN),
    ("3", "계수 축소", "0 방향으로 shrinkage", ACCENT_ORANGE),
]
for i, (num, title, desc, color) in enumerate(effects):
    x = Inches(0.6) + Inches(4.15) * i
    add_shape(s, x, Inches(6.15), Inches(3.8), Inches(1.1), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(6.2), Inches(3.4), Inches(0.4),
             f"{num}. {title}", font_size=13, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(6.6), Inches(3.4), Inches(0.5),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 10: Ridge 회귀 (2) - 고유값 분해 & 베이지안
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "Ridge 회귀 (2) - 고유값 분해와 베이지안 해석")

# 고유값 분해
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "고유값 분해를 통한 이해", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(0.4),
         "축소 비율: d_j / (d_j + lambda)", font_size=16, color=WHITE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(3.3), Inches(5.2), Inches(1.0), [
    "d_j 큰 방향: 비율 -> 1 (거의 변하지 않음)",
    "d_j 작은 방향: 비율 -> 0 (강하게 축소)",
    "=> 데이터 분산이 작은(불안정한) 방향을 선택적 축소",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 베이지안 해석
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "베이지안 해석", font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5), [
    "사전분포: b ~ N(0, tau^2 * I)  (정규분포)",
    "우도: y|X,b ~ N(Xb, sigma^2 * I)",
    "MAP 추정 = Ridge (lambda = sigma^2/tau^2)",
    "",
    "=> Ridge = 회귀 계수에 정규 분포 사전 분포를",
    "   부여한 베이지안 추정",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 릿지 트레이스
add_shape(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.4), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.4),
         "릿지 트레이스 (Ridge Trace) - Hoerl & Kennard", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(5.4), Inches(5.2), Inches(1.5), [
    "x축: lambda (정규화 강도)",
    "y축: 각 회귀 계수 b_j의 값",
    "lambda=0 (OLS)에서 시작, lambda 증가에 따라 모든 계수가 0으로 수렴",
    "계수가 안정화되는 lambda를 시각적으로 선택 가능",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 축소 비율 시각화 (텍스트 기반)
add_text(s, Inches(6.5), Inches(5.4), Inches(6.0), Inches(0.4),
         "축소 비율 d_j/(d_j+lambda) 이해:", font_size=13, color=WHITE, bold=True)
shrink_data = [
    ("d_j = 100", "lambda=1: 0.99", "lambda=10: 0.91", "lambda=100: 0.50"),
    ("d_j = 10", "lambda=1: 0.91", "lambda=10: 0.50", "lambda=100: 0.09"),
    ("d_j = 1", "lambda=1: 0.50", "lambda=10: 0.09", "lambda=100: 0.01"),
]
for i, (dj, v1, v2, v3) in enumerate(shrink_data):
    y_pos = Inches(5.85) + Inches(0.4) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    add_shape(s, Inches(6.5), y_pos, Inches(1.5), Inches(0.35), bg)
    add_text(s, Inches(6.5), y_pos, Inches(1.5), Inches(0.35),
             dj, font_size=11, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_shape(s, Inches(8.0), y_pos, Inches(1.5), Inches(0.35), bg)
    add_text(s, Inches(8.0), y_pos, Inches(1.5), Inches(0.35),
             v1, font_size=11, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_shape(s, Inches(9.5), y_pos, Inches(1.5), Inches(0.35), bg)
    add_text(s, Inches(9.5), y_pos, Inches(1.5), Inches(0.35),
             v2, font_size=11, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
    add_shape(s, Inches(11.0), y_pos, Inches(1.5), Inches(0.35), bg)
    add_text(s, Inches(11.0), y_pos, Inches(1.5), Inches(0.35),
             v3, font_size=11, color=ACCENT_RED, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 11: Lasso 회귀 (1) - L1 정규화
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "Lasso 회귀 (1) - L1 정규화", "Tibshirani (1996): Shrinkage + Selection")

# 핵심 혁신
add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "Lasso = Least Absolute Shrinkage and Selection Operator", font_size=16, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.6),
         "Ridge의 L2 패널티(||b||^2)를 L1 패널티(||b||_1)로 교체 => 축소(shrinkage) + 변수 선택(selection) 동시 수행\nGoogle Scholar 인용 50,000회+ (통계학/ML 역사상 가장 영향력 있는 논문 중 하나)",
         font_size=14, color=LIGHT_GRAY)

# L1 수식
add_shape(s, Inches(0.6), Inches(3.8), Inches(5.8), Inches(1.5), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(3.9), Inches(5.2), Inches(0.4),
         "L1 정규화 최적화 문제", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(4.3), Inches(5.2), Inches(0.8),
         "min { (1/2n)||y - Xb||^2 + lambda * ||b||_1 }\n\n||b||_1 = sum(|bj|) : L1 노름 (절댓값 합)",
         font_size=14, color=WHITE)

# 핵심 차이
add_shape(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(1.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.1), Inches(3.9), Inches(5.2), Inches(0.4),
         "Ridge와의 핵심 차이", font_size=14, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(4.3), Inches(5.2), Inches(0.9), [
    "L1 노름은 원점에서 미분 불가능",
    "닫힌 형태 해가 존재하지 않음 (직교 설계 제외)",
    "좌표 하강법 등 반복 알고리즘 필요",
    "일부 계수를 정확히 0으로 만듦 => 변수 선택!",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 연성 임계값 연산자
add_shape(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.6), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.4),
         "연성 임계값 연산자 (Soft Thresholding): S(z, gamma) = sign(z) * max(|z| - gamma, 0)",
         font_size=15, color=ACCENT_CYAN, bold=True)

thresh_items = [
    ("|z| <= lambda", "=> 계수를 정확히 0", "변수 제거", ACCENT_RED),
    ("z > lambda", "=> z - lambda", "양수 방향 축소", ACCENT_GREEN),
    ("z < -lambda", "=> z + lambda", "음수 방향 축소", ACCENT_BLUE),
]
for i, (cond, result, effect, color) in enumerate(thresh_items):
    x = Inches(0.8) + Inches(4.0) * i
    add_shape(s, x, Inches(6.2), Inches(3.6), Inches(0.85), RGBColor(0x2D, 0x2D, 0x45), color, radius=True)
    add_text(s, x + Inches(0.15), Inches(6.25), Inches(3.3), Inches(0.35),
             f"{cond}  {result}", font_size=12, color=WHITE)
    add_text(s, x + Inches(0.15), Inches(6.6), Inches(3.3), Inches(0.35),
             effect, font_size=12, color=color, bold=True)

# ============================================================
# 슬라이드 12: Lasso 회귀 (2) - 기하학적 설명 & 베이지안
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "Lasso 회귀 (2) - 기하학적 설명", "L1: 마름모 vs L2: 원 - 왜 변수 선택이 되는가?")

# L2 (원)
add_shape(s, Inches(0.6), Inches(2.2), Inches(3.8), Inches(3.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.6), Inches(2.25), Inches(3.8), Inches(0.45),
         "L2 제약 (Ridge): 원", font_size=16, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(0.8), Inches(2.8), Inches(3.4), Inches(2.2), [
    "b1^2 + b2^2 <= t",
    "제약 영역이 원(circle/sphere)",
    "RSS 등고선과 원의 접점은",
    "일반적으로 좌표축 위에 없음",
    "=> 계수가 정확히 0이 되기 어려움",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# L1 (마름모)
add_shape(s, Inches(4.8), Inches(2.2), Inches(3.8), Inches(3.0), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(4.8), Inches(2.25), Inches(3.8), Inches(0.45),
         "L1 제약 (Lasso): 마름모", font_size=16, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(5.0), Inches(2.8), Inches(3.4), Inches(2.2), [
    "|b1| + |b2| <= t",
    "제약 영역이 마름모(diamond)",
    "마름모의 꼭짓점이 좌표축 위",
    "RSS 등고선이 꼭짓점에서 접할 확률 높음",
    "=> 꼭짓점에서 b1=0 또는 b2=0",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# 비교 화살표
add_shape(s, Inches(9.0), Inches(2.2), Inches(3.8), Inches(3.0), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(9.0), Inches(2.25), Inches(3.8), Inches(0.45),
         "베이지안 해석", font_size=16, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(9.2), Inches(2.8), Inches(3.4), Inches(2.2), [
    "Lasso = 라플라스 사전 분포",
    "p(bj) = (lambda/2) exp(-lambda|bj|)",
    "",
    "정규분포보다 원점에서 더 뾰족",
    "꼬리가 더 두꺼움",
    "=> 0 근처의 작은 계수: 높은 사전 확률",
    "=> 큰 계수: 적당한 사전 확률 유지",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(2))

# Ridge vs Lasso 축소 방식 비교
add_shape(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.7), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.4),
         "축소 방식 비교", font_size=15, color=ACCENT_CYAN, bold=True)

compare_h = ["", "Ridge (L2)", "Lasso (L1)"]
compare_rows = [
    ["축소 방식", "비례 축소: bj * d_j/(d_j+lambda)", "절대값 축소: sign(bj)*max(|bj|-lambda, 0)"],
    ["0 생성", "X (0에 가까워질 뿐 정확히 0 아님)", "O (|bj| <= lambda 이면 정확히 0)"],
    ["변수 선택", "불가능", "가능 (희소 해)"],
]
for j, h in enumerate(compare_h):
    x = Inches(0.9) + Inches(3.7) * j
    w = Inches(1.5) if j == 0 else Inches(4.5)
    if j == 0:
        x = Inches(0.9)
        w = Inches(1.5)
    elif j == 1:
        x = Inches(2.4)
        w = Inches(4.8)
    else:
        x = Inches(7.2)
        w = Inches(5.3)
    add_shape(s, x, Inches(6.0), w, Inches(0.35), ACCENT_BLUE if j > 0 else DARK_GRAY)
    add_text(s, x, Inches(6.0), w, Inches(0.35),
             h, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(compare_rows):
    y_pos = Inches(6.4) + Inches(0.28) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    positions = [(Inches(0.9), Inches(1.5)), (Inches(2.4), Inches(4.8)), (Inches(7.2), Inches(5.3))]
    for j, (cell, (x, w)) in enumerate(zip(row, positions)):
        add_shape(s, x, y_pos, w, Inches(0.25), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y_pos, w, Inches(0.25),
                 cell, font_size=10, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 13: Elastic Net
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "Elastic Net - L1+L2 혼합 정규화", "Zou & Hastie (2005)")

# Lasso의 한계
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(1.6), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Lasso의 두 가지 한계", font_size=15, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(1.0), [
    "1. p > n 문제: 최대 n개 변수만 선택 가능",
    "2. 그룹화 효과 부재: 상관 변수 중 하나만 선택",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# Elastic Net 수식
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(1.6), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "Elastic Net 최적화 문제", font_size=15, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(7.1), Inches(2.75), Inches(5.2), Inches(0.9),
         "min { (1/2n)||y-Xb||^2\n   + lambda * [alpha*||b||_1 + (1-alpha)/2 * ||b||^2] }",
         font_size=14, color=WHITE)

# alpha 파라미터
add_shape(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.9), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.4),
         "alpha (l1_ratio) 파라미터:  alpha=1 => 순수 Lasso  |  alpha=0 => 순수 Ridge  |  0<alpha<1 => Elastic Net",
         font_size=15, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.3),
         "lambda > 0: 전체 정규화 강도   |   alpha: L1과 L2의 혼합 비율 (sklearn에서 l1_ratio)",
         font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 그룹화 효과 정리
add_shape(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4),
         "그룹화 효과 정리 (Grouping Effect Theorem)", font_size=15, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.5),
         "|b_i - b_j| <= [1 / (lambda*(1-alpha))] * sqrt(2*(1-r)) * ||y||", font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8), [
    "상관관계 높을수록 (r->1): 두 계수 차이 작아짐 => 상관 변수들을 함께 선택하거나 함께 제거",
    "순수 Lasso (alpha=1): 분모가 0 => 이 성질 성립하지 않음 => Elastic Net만의 고유한 장점",
    "제약 영역: L1 마름모 + L2 원 = 모서리가 둥근 마름모 (rounded diamond) => 희소성 + 안정성",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 14: Ridge vs Lasso vs Elastic Net 종합 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "Ridge vs Lasso vs Elastic Net 종합 비교", "Comprehensive Comparison")

headers = ["특성", "Ridge (L2)", "Lasso (L1)", "Elastic Net (L1+L2)"]
rows = [
    ["패널티", "lambda*||b||_2^2", "lambda*||b||_1", "lambda[alpha*||b||_1+(1-a)/2*||b||_2^2]"],
    ["변수 선택", "X (모든 계수 비영)", "O (일부 계수=0)", "O (일부 계수=0)"],
    ["그룹화 효과", "O (상관 변수 동시 축소)", "X (하나만 선택)", "O (상관 변수 동시 선택)"],
    ["닫힌 형태 해", "O", "X", "X"],
    ["p > n 지원", "O", "최대 n개 변수", "O"],
    ["베이지안 사전분포", "정규 분포", "라플라스 분포", "정규+라플라스 혼합"],
    ["제약 영역", "원 (circle)", "마름모 (diamond)", "둥근 마름모"],
]
header_colors = [DARK_GRAY, ACCENT_BLUE, ACCENT_RED, ACCENT_GREEN]

for j, (h, hc) in enumerate(zip(headers, header_colors)):
    x = Inches(0.5) + Inches(3.1) * j
    w = Inches(1.5) if j == 0 else Inches(3.5)
    if j == 0:
        x = Inches(0.5)
        w = Inches(1.8)
    elif j == 1:
        x = Inches(2.3)
        w = Inches(3.3)
    elif j == 2:
        x = Inches(5.6)
        w = Inches(3.3)
    else:
        x = Inches(8.9)
        w = Inches(3.8)
    add_shape(s, x, Inches(2.2), w, Inches(0.55), hc)
    add_text(s, x, Inches(2.2), w, Inches(0.55),
             h, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    y_pos = Inches(2.85) + Inches(0.58) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    positions = [(Inches(0.5), Inches(1.8)), (Inches(2.3), Inches(3.3)),
                 (Inches(5.6), Inches(3.3)), (Inches(8.9), Inches(3.8))]
    for j, (cell, (x, w)) in enumerate(zip(row, positions)):
        add_shape(s, x, y_pos, w, Inches(0.5), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        if "O" == cell.strip():
            fc = ACCENT_GREEN
        elif "X" == cell.strip():
            fc = ACCENT_RED
        add_text(s, x, y_pos, w, Inches(0.5),
                 cell, font_size=11, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 15: 좌표 하강법
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "좌표 하강법 (Coordinate Descent)", "Friedman, Hastie, Tibshirani (2010) - glmnet")

# 알고리즘 원리
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "알고리즘 원리", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.8), [
    "한 번에 하나의 변수(좌표)만 최적화",
    "모든 변수에 대해 반복 순환",
    "",
    "j번째 좌표 업데이트:",
    "  b_j <- S(rho_j, alpha*lambda) / (X_j'X_j/n + lambda*(1-alpha))",
    "",
    "rho_j = X_j' * r^(j) / n  (부분 잔차와의 내적)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(2))

# 핵심 최적화 기법
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "핵심 최적화 기법 3가지", font_size=16, color=ACCENT_GREEN, bold=True)

techniques = [
    ("따뜻한 시작\n(Warm Start)", "이전 lambda의 해를\n다음 lambda 초기값으로", ACCENT_BLUE),
    ("활성 집합\n(Active Set)", "0이 아닌 계수만\n우선 업데이트", ACCENT_GREEN),
    ("lambda_max\n계산", "모든 계수=0인 최소\nlambda부터 시작", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(techniques):
    y_pos = Inches(2.85) + Inches(0.6) * i
    add_shape(s, Inches(7.1), y_pos, Inches(2.2), Inches(0.5), color, radius=True)
    add_text(s, Inches(7.1), y_pos, Inches(2.2), Inches(0.5),
             title, font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.4), y_pos, Inches(3.0), Inches(0.5),
             desc, font_size=11, color=LIGHT_GRAY)

# 효율성
add_shape(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.4),
         "좌표 하강법이 효율적인 이유", font_size=15, color=ACCENT_CYAN, bold=True)

reasons = [
    ("닫힌 형태", "각 좌표별 최적화가 연성 임계값\n연산자로 한 번에 계산 가능", ACCENT_BLUE),
    ("전역 수렴", "목적 함수가 볼록 => \n전역 수렴 보장", ACCENT_GREEN),
    ("효율적 연산", "한 좌표: O(n), 한 순환: O(np)\nLARS보다 빠르고 확장성 우수", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(reasons):
    x = Inches(0.8) + Inches(4.0) * i
    add_shape(s, x, Inches(5.6), Inches(3.6), Inches(1.4), RGBColor(0x2D, 0x2D, 0x45), color, radius=True)
    add_text(s, x + Inches(0.2), Inches(5.65), Inches(3.2), Inches(0.4),
             title, font_size=13, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(6.05), Inches(3.2), Inches(0.8),
             desc, font_size=12, color=LIGHT_GRAY)

# lambda_max 공식
add_text(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3),
         "lambda_max = (1 / alpha*n) * max_j |X_j^T (y - y_bar)|   (모든 계수가 0이 되는 최소 lambda)",
         font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 16: SCAD와 Oracle 특성
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "SCAD와 Oracle 특성", "Fan & Li (2001): 오라클 성질을 가진 패널티")

# Lasso 한계
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(1.4), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "Lasso의 근본적 한계", font_size=15, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(0.8), [
    "L1 패널티: 계수 크기와 무관하게 동일한 양(lambda) 축소",
    "작은 계수: 적절히 0으로 제거 (OK)",
    "큰 계수: 불필요하게 과도한 축소 => 편향 발생!",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 오라클 성질
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(1.4), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "오라클 성질 (Oracle Properties) 2가지", font_size=15, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(0.8), [
    "1. 일관성 있는 변수 선택: 진정한 변수를 올바르게 식별",
    "2. 점근적 정규성: 0 아닌 계수 추정이 진정한 모델의",
    "   OLS와 동일한 점근 분포 (편향 없음)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# Lasso는 오라클 X
add_shape(s, Inches(0.6), Inches(3.8), Inches(12.1), Inches(0.6), RGBColor(0x3A, 0x1A, 0x1A), ACCENT_RED, radius=True)
add_text(s, Inches(0.6), Inches(3.85), Inches(12.1), Inches(0.5),
         "Lasso는 오라클 성질을 만족하지 못한다 - 큰 계수에 대한 과도한 축소(편향) 때문",
         font_size=14, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)

# SCAD 3구간
add_text(s, Inches(0.6), Inches(4.6), Inches(5), Inches(0.4),
         "SCAD 패널티 3구간", font_size=15, color=ACCENT_PURPLE, bold=True)

scad_headers = ["구간", "패널티 동작", "효과"]
scad_rows = [
    ["|theta| <= lambda", "L1과 동일", "작은 계수를 0으로 제거"],
    ["lambda < |theta| <= a*lambda", "패널티 점차 감소", "중간 계수 점진적 축소"],
    ["|theta| > a*lambda", "패널티 미분 = 0", "큰 계수에 추가 축소 없음 (편향 없음)"],
]
for j, h in enumerate(scad_headers):
    x = Inches(0.6) + Inches(3.8) * j
    w = Inches(3.6)
    if j == 0:
        w = Inches(2.8)
    elif j == 1:
        x = Inches(3.4)
        w = Inches(3.0)
    else:
        x = Inches(6.4)
        w = Inches(3.5)
    add_shape(s, x, Inches(5.0), w, Inches(0.45), ACCENT_PURPLE)
    add_text(s, x, Inches(5.0), w, Inches(0.45),
             h, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(scad_rows):
    y_pos = Inches(5.5) + Inches(0.5) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    positions = [(Inches(0.6), Inches(2.8)), (Inches(3.4), Inches(3.0)), (Inches(6.4), Inches(3.5))]
    for j, (cell, (x, w)) in enumerate(zip(row, positions)):
        add_shape(s, x, y_pos, w, Inches(0.45), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y_pos, w, Inches(0.45),
                 cell, font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# a = 3.7 참고
add_text(s, Inches(10.2), Inches(5.0), Inches(2.5), Inches(0.45),
         "a = 3.7 (통상)", font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 17: 좋은 패널티의 3 조건 + 계보 테이블
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "좋은 패널티의 3 조건 + 정규화 방법의 계보")

# 3 조건
conds = [
    ("비편향성\n(Unbiasedness)", "큰 계수에 대해\n편향이 없거나 거의 없을 것", ACCENT_BLUE),
    ("희소성\n(Sparsity)", "작은 계수를\n0으로 만들 것", ACCENT_GREEN),
    ("연속성\n(Continuity)", "추정량이 데이터의\n연속 함수일 것", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(conds):
    x = Inches(0.5) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.8), Inches(1.5), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(2.3), Inches(3.4), Inches(0.5),
             title, font_size=15, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.9), Inches(3.4), Inches(0.6),
             desc, font_size=13, color=LIGHT_GRAY)

add_text(s, Inches(0.6), Inches(3.85), Inches(12), Inches(0.3),
         "3 조건을 동시 만족하는 패널티는 반드시 비볼록(nonconvex) => L1은 볼록이므로 비편향성 불만족",
         font_size=13, color=ACCENT_RED, align=PP_ALIGN.CENTER)

# 계보 테이블
lineage_h = ["방법", "연도", "패널티 유형", "변수 선택", "오라클", "볼록성"]
lineage_rows = [
    ["Ridge", "1970", "L2", "X", "X", "O"],
    ["Lasso", "1996", "L1", "O", "X", "O"],
    ["SCAD", "2001", "비볼록", "O", "O", "X"],
    ["Elastic Net", "2005", "L1+L2", "O", "X", "O"],
    ["Adaptive Lasso", "2006", "가중 L1", "O", "O", "O"],
    ["MCP", "2010", "비볼록", "O", "O", "X"],
]
for j, h in enumerate(lineage_h):
    x = Inches(0.5) + Inches(2.1) * j
    w = Inches(2.0)
    add_shape(s, x, Inches(4.3), w, Inches(0.5), ACCENT_BLUE)
    add_text(s, x, Inches(4.3), w, Inches(0.5),
             h, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(lineage_rows):
    y_pos = Inches(4.85) + Inches(0.42) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.5) + Inches(2.1) * j
        w = Inches(2.0)
        add_shape(s, x, y_pos, w, Inches(0.38), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        if cell == "O":
            fc = ACCENT_GREEN
        elif cell == "X":
            fc = ACCENT_RED
        add_text(s, x, y_pos, w, Inches(0.38),
                 cell, font_size=11, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 18: 회귀 진단 - 잔차 분석
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "회귀 진단 - 잔차 분석", "Residual Analysis, Q-Q Plot, Cook's Distance")

# 잔차의 종류
add_text(s, Inches(0.6), Inches(2.0), Inches(5), Inches(0.4),
         "잔차의 종류", font_size=16, color=WHITE, bold=True)

res_h = ["잔차 유형", "공식", "용도"]
res_rows = [
    ["일반 잔차", "ei = yi - yi_hat", "기본 분석"],
    ["표준화 잔차", "ri = ei / (s_hat * sqrt(1-hii))", "이상치 탐지"],
    ["스튜던트화 잔차", "ti = ei / (s_hat(i) * sqrt(1-hii))", "더 정밀한 이상치 탐지"],
]
for j, h in enumerate(res_h):
    x = Inches(0.6) + Inches(4.0) * j
    w = Inches(3.8)
    add_shape(s, x, Inches(2.5), w, Inches(0.45), ACCENT_BLUE)
    add_text(s, x, Inches(2.5), w, Inches(0.45),
             h, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, row in enumerate(res_rows):
    y_pos = Inches(3.0) + Inches(0.48) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.6) + Inches(4.0) * j
        w = Inches(3.8)
        add_shape(s, x, y_pos, w, Inches(0.43), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y_pos, w, Inches(0.43),
                 cell, font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# Q-Q Plot
add_shape(s, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(4.8), Inches(5.2), Inches(0.4),
         "Q-Q Plot (Quantile-Quantile Plot)", font_size=15, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(5.2), Inches(5.2), Inches(1.8), [
    "x축: 이론적 정규 분포의 분위수",
    "y축: 잔차의 분위수",
    "점들이 45도 직선 위에 놓이면 정규 분포",
    "직선에서 벗어나면 분포 왜곡 판단",
    "양쪽 꼬리 벗어남: 정규성 위반 (두꺼운 꼬리)",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# Cook's Distance
add_shape(s, Inches(6.8), Inches(4.7), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.1), Inches(4.8), Inches(5.2), Inches(0.4),
         "Cook's Distance (쿡의 거리)", font_size=15, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(7.1), Inches(5.25), Inches(5.2), Inches(0.4),
         "Di = (ri^2 / p) * hii / (1-hii)^2", font_size=15, color=WHITE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(5.7), Inches(5.2), Inches(1.3), [
    "각 관측치의 영향력(influence) 측정",
    "ri: 표준화 잔차, hii: 레버리지",
    "판정: Di > 1 (일반 기준), Di > 4/n (보수적)",
    "높은 Di: 영향력 큰 관측치 => 제거 검토",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 19: 회귀 진단 체크리스트
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "회귀 진단 종합 체크리스트", "Regression Diagnostics Checklist")

diag_h = ["항목", "확인 방법", "위반 시 조치"]
diag_rows = [
    ["선형성", "예측값 vs 잔차 플롯", "변수 변환, 다항 회귀"],
    ["등분산성", "예측값 vs 잔차 플롯 (깔때기 패턴)", "WLS, 변수 변환"],
    ["정규성", "Q-Q plot, Shapiro-Wilk 검정", "변수 변환, 로버스트 회귀"],
    ["독립성", "Durbin-Watson 검정", "GLS, 시계열 모델"],
    ["다중공선성", "VIF, 조건수", "변수 제거, Ridge 회귀"],
    ["이상치/영향점", "Cook's distance, 레버리지", "제거 또는 로버스트 회귀"],
]

for j, h in enumerate(diag_h):
    x = Inches(0.6) + Inches(4.0) * j
    w = Inches(3.8)
    add_shape(s, x, Inches(2.2), w, Inches(0.55), ACCENT_BLUE)
    add_text(s, x, Inches(2.2), w, Inches(0.55),
             h, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

colors_diag = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED, ACCENT_PURPLE]
for i, row in enumerate(diag_rows):
    y_pos = Inches(2.85) + Inches(0.65) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.6) + Inches(4.0) * j
        w = Inches(3.8)
        add_shape(s, x, y_pos, w, Inches(0.55), bg)
        if j == 0:
            # 색깔 인디케이터
            add_shape(s, x, y_pos, Inches(0.08), Inches(0.55), colors_diag[i])
            fc = ACCENT_CYAN
        else:
            fc = LIGHT_GRAY
        add_text(s, x + Inches(0.15), y_pos, w - Inches(0.15), Inches(0.55),
                 cell, font_size=13, color=fc, bold=(j == 0))

# 하단 요약
add_shape(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5), CARD_BG, radius=True)
add_text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
         "잔차 플롯 핵심: 패턴이 없어야 한다! (무작위 분포) | 레버리지: hii (X 공간에서의 극단성) | 영향력: Cook's D (레버리지 x 잔차)",
         font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 20: 실습 소개
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "실습 소개 (3개 구현 소스)")

labs = [
    ("실습 1", "01_ridge_regression.py", "Ridge 회귀 직접 구현", [
        "다중공선성 데이터 생성 (5 기본 + 5 상관)",
        "닫힌 형태 해 직접 구현 vs sklearn 비교",
        "릿지 트레이스 시각화 (정규화 경로)",
        "교차검증 최적 lambda 선택 (RidgeCV)",
        "OLS vs Ridge: 조건수 개선 효과 확인",
    ], ACCENT_BLUE),
    ("실습 2", "02_lasso_regression.py", "Lasso 회귀 좌표 하강법 구현", [
        "연성 임계값 연산자 직접 구현",
        "좌표 하강법 알고리즘 구현 (Friedman et al.)",
        "희소 데이터에서 변수 선택 효과 확인",
        "정규화 경로: 선택 변수 수 변화 시각화",
        "LassoCV로 최적 lambda 및 변수 선택 검증",
    ], ACCENT_GREEN),
    ("실습 3", "03_elastic_net.py", "Elastic Net 구현 & 비교", [
        "Elastic Net 좌표 하강법 직접 구현",
        "상관 그룹 데이터로 그룹화 효과 시연",
        "Ridge vs Lasso vs Elastic Net 비교",
        "그룹 내 계수 분산으로 그룹화 효과 정량화",
        "L1/L2/Elastic Net 제약 영역 시각화",
    ], ACCENT_ORANGE),
]
for i, (title, file, desc, items, color) in enumerate(labs):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.0), Inches(3.9), Inches(5.2), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.0), Inches(3.9), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(2.05), Inches(3.9), Inches(0.5),
             title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.7), Inches(3.5), Inches(0.3),
             file, font_size=11, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.0), Inches(3.5), Inches(0.3),
             desc, font_size=13, color=WHITE)
    add_bullet_list(s, x + Inches(0.2), Inches(3.4), Inches(3.5), Inches(3.5),
                    items, font_size=11, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# 슬라이드 21: 응용사례
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "11", "응용사례", "Real-World Applications of Linear Regression")

apps = [
    ("의료", "질병 진행 예측\n약물 반응 모델링\n바이오마커 선택 (Lasso)\n유전체 데이터 분석 (Elastic Net)", ACCENT_BLUE),
    ("금융", "자산 가격 모델링 (CAPM)\n리스크 팩터 분석\n포트폴리오 최적화\n신용 스코어링", ACCENT_CYAN),
    ("제조", "품질 관리 (공정 변수 영향)\n에너지 소비 예측\n설비 수명 예측\n불량률 분석", ACCENT_GREEN),
    ("마케팅", "광고비-매출 분석\n고객 생애 가치(CLV) 예측\n가격 탄력성 추정\n채널 기여도 분석", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(apps):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(2.2), Inches(2.9), Inches(4.5), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(2.9), Inches(0.6), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(2.9), Inches(0.5),
             title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(3.3),
             desc, font_size=13, color=LIGHT_GRAY)

# 하단 참고
add_shape(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4), CARD_BG, radius=True)
add_text(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4),
         "핵심: 해석 가능성이 중요한 도메인에서 선형회귀 + 정규화는 여전히 강력한 첫 번째 선택지",
         font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 22: 핵심 수식 정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "핵심 수식 정리", "Key Formulas Summary")

formulas = [
    ("OLS", "b_OLS = (X^TX)^(-1) X^Ty", "BLUE, 비편향, 닫힌 형태 해"),
    ("Ridge", "b_ridge = (X^TX + lambda*I)^(-1) X^Ty", "L2 정규화, 축소, 조건수 개선"),
    ("Lasso", "min (1/2n)||y-Xb||^2 + lambda*||b||_1", "L1 정규화, 변수 선택, 좌표 하강법"),
    ("Elastic Net", "min (1/2n)||y-Xb||^2 + lambda[a||b||_1+(1-a)/2*||b||_2^2]", "L1+L2 혼합, 그룹화 효과"),
    ("VIF", "VIF_j = 1 / (1 - R_j^2)", "다중공선성 진단 (10 이상 심각)"),
    ("Cook's D", "Di = (ri^2/p) * hii/(1-hii)^2", "영향력 관측치 탐지 (Di>1 주의)"),
]
for i, (name, formula, meaning) in enumerate(formulas):
    y_pos = Inches(2.2) + Inches(0.85) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    add_shape(s, Inches(0.6), y_pos, Inches(12.1), Inches(0.75), bg, radius=True)
    add_text(s, Inches(0.9), y_pos + Inches(0.05), Inches(2.0), Inches(0.65),
             name, font_size=15, color=ACCENT_CYAN, bold=True)
    add_text(s, Inches(3.0), y_pos + Inches(0.05), Inches(5.5), Inches(0.65),
             formula, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.6), y_pos + Inches(0.05), Inches(4.0), Inches(0.65),
             meaning, font_size=12, color=DARK_GRAY)

# ============================================================
# 슬라이드 23: 핵심 요약
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "핵심 요약 (Key Takeaways)")

summaries = [
    ("OLS 이론", "정규방정식: b=(X^TX)^(-1)X^Ty\n기하학적: 직교 사영 + 햇 행렬", ACCENT_BLUE),
    ("가우스-마르코프", "OLS = BLUE (비편향 중 최소분산)\n편향 허용 => MSE 더 작아질 수 있음", ACCENT_CYAN),
    ("Ridge 회귀", "L2 정규화, 축소, 닫힌 형태 해\n베이지안: 정규 분포 사전 분포", ACCENT_GREEN),
    ("Lasso 회귀", "L1 정규화, 변수 선택 + 축소 동시\n연성 임계값, 마름모 제약 영역", ACCENT_ORANGE),
    ("Elastic Net", "L1+L2 혼합, 그룹화 효과\np>n 지원, 둥근 마름모 제약", ACCENT_RED),
    ("SCAD & Oracle", "비볼록 패널티, 오라클 성질 만족\n큰 계수에 편향 없음", ACCENT_PURPLE),
    ("좌표 하강법", "glmnet (Friedman 2010)\n따뜻한 시작 + 활성 집합", ACCENT_BLUE),
    ("회귀 진단", "잔차 분석, Q-Q Plot, Cook's D\n선형성/등분산/정규성/독립성 점검", ACCENT_CYAN),
]
for i, (title, desc, color) in enumerate(summaries):
    col = i % 4
    row = i // 4
    x = Inches(0.4) + Inches(3.15) * col
    y = Inches(2.0) + Inches(2.6) * row
    add_shape(s, x, y, Inches(2.95), Inches(2.2), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), y + Inches(0.1), Inches(2.65), Inches(0.4),
             title, font_size=15, color=color, bold=True)
    add_accent_line(s, x + Inches(0.15), y + Inches(0.55), Inches(1.8), color)
    add_text(s, x + Inches(0.15), y + Inches(0.65), Inches(2.65), Inches(1.3),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 24: 방법 선택 가이드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "방법 선택 가이드 (Decision Tree)")

# 의사결정 트리
add_shape(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(5.2), CARD_BG, ACCENT_BLUE, radius=True)

# Root
add_shape(s, Inches(4.5), Inches(2.3), Inches(4.3), Inches(0.6), ACCENT_BLUE, radius=True)
add_text(s, Inches(4.5), Inches(2.3), Inches(4.3), Inches(0.6),
         "다중공선성이 심한가?", font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# YES branch
add_shape(s, Inches(1.2), Inches(3.3), Inches(0.8), Inches(0.4), ACCENT_GREEN, radius=True)
add_text(s, Inches(1.2), Inches(3.3), Inches(0.8), Inches(0.4),
         "YES", font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(s, Inches(1.0), Inches(3.9), Inches(4.8), Inches(0.5), RGBColor(0x2D, 0x2D, 0x45), ACCENT_CYAN, radius=True)
add_text(s, Inches(1.0), Inches(3.9), Inches(4.8), Inches(0.5),
         "변수 선택이 필요한가?", font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_shape(s, Inches(1.0), Inches(4.7), Inches(2.2), Inches(0.8), ACCENT_GREEN, radius=True)
add_text(s, Inches(1.0), Inches(4.7), Inches(2.2), Inches(0.4),
         "YES => Elastic Net", font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(5.05), Inches(2.2), Inches(0.4),
         "상관 그룹 있으면 특히 유리", font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_shape(s, Inches(3.5), Inches(4.7), Inches(2.2), Inches(0.8), ACCENT_BLUE, radius=True)
add_text(s, Inches(3.5), Inches(4.7), Inches(2.2), Inches(0.4),
         "NO => Ridge", font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(3.5), Inches(5.05), Inches(2.2), Inches(0.4),
         "조건수 개선, 모든 변수 유지", font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# NO branch
add_shape(s, Inches(7.3), Inches(3.3), Inches(0.8), Inches(0.4), ACCENT_RED, radius=True)
add_text(s, Inches(7.3), Inches(3.3), Inches(0.8), Inches(0.4),
         "NO", font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(s, Inches(6.8), Inches(3.9), Inches(5.5), Inches(0.5), RGBColor(0x2D, 0x2D, 0x45), ACCENT_ORANGE, radius=True)
add_text(s, Inches(6.8), Inches(3.9), Inches(5.5), Inches(0.5),
         "변수 수가 많은가? (p > n?)", font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_shape(s, Inches(6.8), Inches(4.7), Inches(2.5), Inches(0.8), ACCENT_ORANGE, radius=True)
add_text(s, Inches(6.8), Inches(4.7), Inches(2.5), Inches(0.4),
         "YES => Lasso / Elastic Net", font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.8), Inches(5.05), Inches(2.5), Inches(0.4),
         "변수 선택으로 차원 축소", font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_shape(s, Inches(9.6), Inches(4.7), Inches(2.5), Inches(0.8), ACCENT_PURPLE, radius=True)
add_text(s, Inches(9.6), Inches(4.7), Inches(2.5), Inches(0.4),
         "NO => OLS", font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(9.6), Inches(5.05), Inches(2.5), Inches(0.4),
         "비편향, BLUE, 해석 용이", font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 핵심 논문 연대기
add_text(s, Inches(1.0), Inches(5.8), Inches(11), Inches(0.4),
         "핵심 논문 연대기", font_size=14, color=WHITE, bold=True)

timeline = [
    ("1970", "Hoerl & Kennard\nRidge 회귀", ACCENT_BLUE),
    ("1996", "Tibshirani\nLasso", ACCENT_RED),
    ("2001", "Fan & Li\nSCAD", ACCENT_PURPLE),
    ("2005", "Zou & Hastie\nElastic Net", ACCENT_GREEN),
    ("2010", "Friedman et al.\nglmnet", ACCENT_ORANGE),
]
add_shape(s, Inches(1.2), Inches(6.5), Inches(10.8), Pt(3), ACCENT_CYAN)
for i, (year, desc, color) in enumerate(timeline):
    x = Inches(1.0) + Inches(2.35) * i
    add_shape(s, x + Inches(0.7), Inches(6.35), Inches(0.3), Inches(0.3), color, radius=True)
    add_text(s, x, Inches(6.0), Inches(1.7), Inches(0.35),
             year, font_size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(6.7), Inches(1.7), Inches(0.7),
             desc, font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 25: 복습 질문
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "복습 질문 (Review Questions)")

questions = [
    "Q1. OLS의 비용함수 RSS(b)를 행렬 미분하여 정규방정식을 유도하고, 헤시안으로 최솟값임을 확인하라.",
    "Q2. 가우스-마르코프 정리의 한계를 설명하고, Ridge가 OLS보다 MSE 관점에서 나을 수 있는 이유를 서술하라.",
    "Q3. Ridge의 닫힌 형태 해를 유도하고, lambda*I가 조건수를 개선하는 이유를 설명하라.",
    "Q4. L1(마름모) vs L2(원) 제약 영역의 기하학으로 Lasso의 변수 선택 원리를 설명하라.",
    "Q5. 연성 임계값 연산자 S(z,gamma)의 세 구간 동작을 설명하고, 좌표 하강법에서의 역할을 서술하라.",
    "Q6. Elastic Net의 그룹화 효과 정리를 기술하고, Lasso의 그룹화 부재 문제를 어떻게 해결하는지 설명하라.",
    "Q7. 좌표 하강법의 따뜻한 시작(warm start)과 활성 집합(active set)이 계산 속도를 향상시키는 원리를 설명하라.",
    "Q8. SCAD의 오라클 성질 2가지를 기술하고, Lasso가 이를 만족하지 못하는 이유를 설명하라.",
    "Q9. Ridge와 Lasso를 각각 베이지안 MAP 추정 관점에서 해석하라 (정규 vs 라플라스 사전 분포).",
    "Q10. 회귀 진단: 깔때기 잔차 패턴, Q-Q 꼬리 이탈, Cook's D > 1 관측치에 대한 조치를 서술하라.",
]
for i, q in enumerate(questions):
    y_pos = Inches(2.0) + Inches(0.52) * i
    color = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED,
             ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE][i]
    add_shape(s, Inches(0.8), y_pos + Inches(0.05), Inches(0.08), Inches(0.3), color)
    add_text(s, Inches(1.1), y_pos, Inches(11.5), Inches(0.5),
             q, font_size=13, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 26: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.5), prs.slide_width, Inches(1.0),
         "Thank You", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), prs.slide_width, Inches(0.6),
         "4장: 선형회귀 (Linear Regression)", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(4.2), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(4.5), prs.slide_width, Inches(0.5),
         "다음 장: 5장 - 로지스틱 회귀", font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ── 저장 ──
output_path = os.path.join(os.path.dirname(__file__), "4장_선형회귀_강의PPT.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
