"""4장 선형회귀 - 교재 실습 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)
ACCENT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = '맑은 고딕'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)


# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG, ACCENT_BLUE)
add_accent_line(s, Inches(4), Inches(1.3), Inches(5.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.6), Inches(11.333), Inches(0.5),
         "CHAPTER 04", font_size=22, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.3), Inches(11.333), Inches(1.2),
         "선형 회귀 - 보험료 예측하기", font_size=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.6), Inches(11.333), Inches(0.5),
         "Linear Regression: Predicting Insurance Charges", font_size=22, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5), Inches(4.5), Inches(3.333), ACCENT_CYAN)
add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.4),
         "교수자: Jung, Minpo", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.4), Inches(11.333), Inches(0.4),
         "교과목: Machine Learning  |  학기: 2026년도 1학기", font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 2: 학습 목표 & 학습 순서
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "학습 목표 & 학습 순서")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "학습 목표",
         ["선형 회귀 모델로 보험 데이터셋을 학습해",
          "보험료를 예측하고,",
          "선형 회귀의 작동 원리를 이해한다."],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5),
         "학습 순서",
         ["4.1  문제 정의",
          "4.2  라이브러리 및 데이터 불러오기",
          "4.3  데이터 확인하기",
          "4.4  전처리: 학습셋/시험셋 나누기",
          "4.5  데이터 모델링",
          "4.6  모델을 활용해 예측하기",
          "4.7  모델 평가하기",
          "4.8  이해하기: 선형 회귀 원리"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 파이프라인 화살표
arrow_y = Inches(5.2)
steps = ["4.1", "4.2", "4.3", "4.4", "4.5", "4.6", "4.7", "4.8"]
for i, st in enumerate(steps):
    x = Inches(0.6 + i * 0.75)
    add_shape(s, x, arrow_y, Inches(0.6), Inches(0.45), ACCENT_BLUE, radius=True)
    add_text(s, x, arrow_y, Inches(0.6), Inches(0.45), st, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s, x + Inches(0.55), arrow_y, Inches(0.25), Inches(0.45), "\u2192",
                 font_size=16, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 3: 선형 회귀 소개
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "선형 회귀 소개", "Linear Regression Overview")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.2),
         "Linear Regression 이란?",
         ["가장 기초적인 머신러닝 모델",
          "연속형 변수를 예측 (예: 165.5cm, 172.3cm 등)",
          "입력 변수와 출력 변수 사이의 선형 관계를 학습"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.2),
         "TOP 10 알고리즘 선정 이유",
         ["가장 기초적인 알고리즘으로 반드시 알아야 함",
          "다른 알고리즘의 베이스라인(Baseline)으로 활용",
          "결과 해석이 용이하여 실무에서 자주 사용"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# 장단점 테이블
add_table_slide(s,
    headers=["구분", "내용"],
    rows=[["장점", "간단하고 직관적, 해석이 쉬움"],
          ["장점", "학습 속도가 빠르고 구현이 용이"],
          ["단점", "예측력이 낮을 수 있음"],
          ["단점", "독립변수와 종속변수 간 선형 관계 전제"]],
    left=Inches(0.6), top=Inches(4.6),
    col_widths=[2.0, 10.0], row_height=0.45, font_size=13)

# ============================================================
# 슬라이드 4: 4.1 문제 정의
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.1", "문제 정의", "Problem Definition")

add_card(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.8),
         "미션 (Mission)",
         ["보험 데이터셋을 이용하여",
          "보험사에서 청구할 보험료(charges)를 예측하라!"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(1.8),
         "프로젝트 정보",
         ["난이도: \u2605\u2606\u2606",
          "알고리즘: Linear Regression",
          "데이터셋: insurance.csv"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.1), Inches(3.8), Inches(2.5),
         "데이터 정보",
         ["데이터: insurance.csv",
          "종속변수(y): charges",
          "독립변수(X): age, sex, bmi,",
          "  children, smoker"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(4.1), Inches(3.8), Inches(2.5),
         "평가지표",
         ["RMSE",
          "(Root Mean Squared Error)",
          "예측값과 실제값의 차이를",
          "제곱 평균의 제곱근으로 측정"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(8.8), Inches(4.1), Inches(3.8), Inches(2.5),
         "사용 라이브러리",
         ["numpy", "pandas", "seaborn",
          "matplotlib", "sklearn"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 슬라이드 5: 4.2 라이브러리 및 데이터 불러오기
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.2", "라이브러리 및 데이터 불러오기", "Import Libraries & Load Data")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(11.5), Inches(2.5), [
    "import pandas as pd",
    "",
    "file_url = 'https://media.githubusercontent.com/media/",
    "            musthave-ML10/data_source/main/insurance.csv'",
    "",
    "data = pd.read_csv(file_url)",
], font_size=14)

add_card(s, Inches(0.6), Inches(5.0), Inches(5.5), Inches(1.8),
         "pandas",
         ["데이터 분석을 위한 핵심 라이브러리",
          "DataFrame 구조로 데이터를 다룸",
          "CSV, Excel 등 다양한 형식 지원"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.5), Inches(5.0), Inches(5.6), Inches(1.8),
         "pd.read_csv()",
         ["URL 또는 로컬 경로에서 CSV 파일을 읽어",
          "DataFrame으로 변환하는 함수",
          "data 변수에 저장하여 사용"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 6: 4.3 데이터 확인하기 (1)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.3", "데이터 확인하기 (1)", "data 출력 및 data.head()")

add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "data \u2192 1338 rows \u00d7 6 columns", font_size=18, color=ACCENT_CYAN, bold=True)

add_card(s, Inches(0.6), Inches(2.5), Inches(5.5), Inches(2.2),
         "변수 설명",
         ["age     : 나이 (연속형)",
          "sex     : 성별 (범주형, 1=남, 0=여)",
          "bmi     : 체질량지수 (연속형)",
          "children: 자녀 수 (연속형)",
          "smoker  : 흡연 여부 (범주형, 1=예, 0=아니오)",
          "charges : 보험료 (종속변수, 연속형)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# data.head() 테이블
add_text(s, Inches(0.6), Inches(5.0), Inches(4), Inches(0.4),
         "data.head() 결과", font_size=16, color=ACCENT_GREEN, bold=True)

add_table_slide(s,
    headers=["age", "sex", "bmi", "children", "smoker", "charges"],
    rows=[["19", "0", "27.90", "0", "1", "16884.92"],
          ["18", "1", "33.77", "1", "0", "1725.55"],
          ["28", "1", "33.00", "3", "0", "4449.46"],
          ["33", "1", "22.70", "0", "0", "21984.47"],
          ["32", "1", "28.88", "0", "0", "3866.86"]],
    left=Inches(6.5), top=Inches(2.5),
    col_widths=[1.0, 0.8, 1.0, 1.0, 1.0, 1.6],
    row_height=0.42, font_size=12, header_font_size=12)

# ============================================================
# 슬라이드 7: 4.3 데이터 확인하기 (2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.3", "데이터 확인하기 (2)", "연속형 vs 범주형 변수 & data.info()")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.0),
         "연속형 변수 vs 범주형 변수",
         ["\u25b6 연속형(Continuous): 연속적 숫자값 (age, bmi, children, charges)",
          "\u25b6 범주형(Categorical): 범주/그룹을 나타내는 값 (sex, smoker)",
          "",
          "  sex   : 1 = 남자(male),  0 = 여자(female)",
          "  smoker: 1 = 흡연자(yes), 0 = 비흡연자(no)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
         "data.info() 결과",
         ["<class 'pandas.core.frame.DataFrame'>",
          "RangeIndex: 1338 entries, 0 to 1337",
          "Data columns (total 6 columns):",
          "",
          "  #  Column    Non-Null Count  Dtype",
          "  0  age       1338 non-null   int64",
          "  1  sex       1338 non-null   int64",
          "  2  bmi       1338 non-null   float64",
          "  3  children  1338 non-null   int64",
          "  4  smoker    1338 non-null   int64",
          "  5  charges   1338 non-null   float64",
          "",
          "dtypes: float64(2), int64(4)",
          "결측치(null) 없음 \u2192 전처리 부담 적음"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.5),
         "핵심 포인트",
         ["\u2022 6개 컬럼, 1338개 행",
          "\u2022 결측치(Missing Value)가 없음",
          "\u2022 sex, smoker는 이미 숫자로 인코딩 완료",
          "\u2022 별도의 인코딩 전처리 불필요"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 8: 4.3 데이터 확인하기 (3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.3", "데이터 확인하기 (3)", "data.describe() & 사분위수")

add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "round(data.describe(), 2)", font_size=16, color=ACCENT_GREEN, bold=True)

add_table_slide(s,
    headers=["", "age", "sex", "bmi", "children", "smoker", "charges"],
    rows=[["count", "1338", "1338", "1338", "1338", "1338", "1338"],
          ["mean", "39.21", "0.51", "30.66", "1.09", "0.20", "13270.42"],
          ["std", "14.05", "0.50", "6.10", "1.21", "0.40", "12110.01"],
          ["min", "18.00", "0.00", "15.96", "0.00", "0.00", "1121.87"],
          ["25%", "27.00", "0.00", "26.30", "0.00", "0.00", "4740.29"],
          ["50%", "39.00", "1.00", "30.40", "1.00", "0.00", "9382.03"],
          ["75%", "51.00", "1.00", "34.69", "2.00", "0.00", "16639.91"],
          ["max", "64.00", "1.00", "53.13", "5.00", "1.00", "63770.43"]],
    left=Inches(0.6), top=Inches(2.5),
    col_widths=[1.2, 1.2, 1.0, 1.2, 1.2, 1.2, 1.8],
    row_height=0.4, font_size=11, header_font_size=12)

add_card(s, Inches(0.6), Inches(6.0), Inches(12.0), Inches(1.2),
         "사분위수 (Quantile) 개념",
         ["25% (Q1): 하위 25% 지점  |  50% (Q2/중앙값): 정중앙  |  75% (Q3): 상위 25% 지점",
          "charges 평균(13270) > 중앙값(9382) \u2192 오른쪽으로 치우친 분포 (고액 보험료 존재)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 9: 4.4 전처리 - 학습셋과 시험셋 나누기 (개념)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.4", "전처리: 학습셋과 시험셋 나누기", "Train/Test Split Concept")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.0),
         "1단계: 독립변수 / 종속변수 분리",
         ["독립변수(X): age, sex, bmi, children, smoker",
          "종속변수(y): charges (예측 대상)",
          "",
          "\u2192 X는 원인, y는 결과"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "2단계: 학습셋 / 시험셋 분리",
         ["학습셋(Train): 모델이 학습하는 데이터",
          "시험셋(Test): 모델 성능을 평가하는 데이터",
          "",
          "\u2192 보통 8:2 또는 7:3 비율로 분리"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 4개 데이터셋 카드
labels = [("X_train", "학습용 독립변수", ACCENT_BLUE),
          ("y_train", "학습용 종속변수", ACCENT_CYAN),
          ("X_test", "시험용 독립변수", ACCENT_ORANGE),
          ("y_test", "시험용 종속변수", ACCENT_RED)]
for i, (name, desc, clr) in enumerate(labels):
    x = Inches(0.6 + i * 3.15)
    add_shape(s, x, Inches(4.5), Inches(2.8), Inches(1.5), CARD_BG, clr, radius=True)
    add_text(s, x, Inches(4.6), Inches(2.8), Inches(0.5), name,
             font_size=22, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.2), Inches(2.8), Inches(0.5), desc,
             font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 10: 4.4 데이터셋 나누기 (코드)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.4", "데이터셋 나누기 (코드)", "Train/Test Split Code")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(11.5), Inches(3.0), [
    "# 독립변수(X)와 종속변수(y) 분리",
    "X = data[['age', 'sex', 'bmi', 'children', 'smoker']]",
    "y = data['charges']",
    "",
    "# 학습셋과 시험셋 분리",
    "from sklearn.model_selection import train_test_split",
    "",
    "X_train, X_test, y_train, y_test = train_test_split(",
    "    X, y, test_size=0.2, random_state=100)",
], font_size=14)

add_card(s, Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.5),
         "test_size = 0.2",
         ["전체 데이터의 20%를 시험셋으로 분리",
          "나머지 80%가 학습셋이 됨"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.5), Inches(5.5), Inches(5.6), Inches(1.5),
         "random_state = 100",
         ["랜덤 샘플링을 고정하여 재현성 보장",
          "같은 숫자를 넣으면 항상 같은 결과"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 11: 4.5 데이터 모델링
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.5", "데이터 모델링", "Model Training")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(1.0),
         "모델링 3단계",
         ["\u2460 알고리즘 선택 (LinearRegression)  \u2192  \u2461 모델 생성 (model 객체)  \u2192  \u2462 모델 학습 (fit)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_code_block(s, Inches(0.6), Inches(3.3), Inches(11.5), Inches(2.5), [
    "# 1. 알고리즘 선택 & 모델 생성",
    "from sklearn.linear_model import LinearRegression",
    "",
    "model = LinearRegression()",
    "",
    "# 2. 모델 학습",
    "model.fit(X_train, y_train)",
], font_size=14)

add_card(s, Inches(0.6), Inches(6.1), Inches(5.5), Inches(1.0),
         "LinearRegression()",
         ["sklearn이 제공하는 선형 회귀 모델 클래스"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(6.5), Inches(6.1), Inches(5.6), Inches(1.0),
         "model.fit(X_train, y_train)",
         ["학습셋(X, y)을 넣어 모델을 학습시킴"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 12: 4.6 모델을 활용해 예측하기
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.6", "모델을 활용해 예측하기", "Prediction")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(11.5), Inches(1.5), [
    "pred = model.predict(X_test)",
], font_size=16)

add_card(s, Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.5),
         "predict() 함수",
         ["\u2022 학습된 모델에 시험셋(X_test)을 입력",
          "\u2022 독립변수(X)만 입력해야 함",
          "\u2022 종속변수(y)는 넣지 않음!",
          "",
          "\u2192 pred에는 예측된 보험료가 저장됨"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.5),
         "주의사항",
         ["\u2022 X_test만 넣어야 함 (y_test 넣으면 안 됨)",
          "\u2022 학습에 사용하지 않은 데이터로 예측해야",
          "  모델의 일반화 성능을 평가할 수 있음",
          "",
          "\u2192 pred 결과를 y_test와 비교하여 평가"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 13: 4.7.1 테이블로 평가하기
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.7.1", "테이블로 평가하기", "Evaluation with Table")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(7.0), Inches(1.8), [
    "comparison = pd.DataFrame({",
    "    'actual': y_test,",
    "    'pred': pred",
    "})",
    "comparison  # 268 rows x 2 columns",
], font_size=13)

add_card(s, Inches(8.0), Inches(2.2), Inches(4.6), Inches(1.8),
         "결과 해석",
         ["268개 시험 데이터에 대해",
          "실제값(actual)과 예측값(pred)을",
          "나란히 비교하는 DataFrame 생성"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# 예시 테이블
add_text(s, Inches(0.6), Inches(4.3), Inches(4), Inches(0.4),
         "comparison 결과 (일부)", font_size=15, color=ACCENT_GREEN, bold=True)

add_table_slide(s,
    headers=["index", "actual", "pred"],
    rows=[["764", "9095.07", "11834.68"],
          ["887", "5272.18", "8292.14"],
          ["890", "29330.98", "34069.01"],
          ["1293", "9301.89", "12077.35"],
          ["259", "44202.65", "36553.34"]],
    left=Inches(0.6), top=Inches(4.8),
    col_widths=[1.5, 2.5, 2.5],
    row_height=0.4, font_size=13)

add_card(s, Inches(7.5), Inches(4.5), Inches(5.2), Inches(2.5),
         "actual vs pred 비교",
         ["\u2022 actual: 실제 보험료 (정답)",
          "\u2022 pred: 모델이 예측한 보험료",
          "",
          "\u2022 값이 비슷하면 좋은 예측",
          "\u2022 차이가 크면 예측 실패",
          "\u2022 테이블만으로는 전체 패턴 파악 어려움"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 14: 4.7.2 그래프로 평가하기
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.7.2", "그래프로 평가하기", "Evaluation with Scatter Plot")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(2.0), [
    "import matplotlib.pyplot as plt",
    "import seaborn as sns",
    "",
    "plt.figure(figsize=(10, 10))",
    "sns.scatterplot(x='actual', y='pred',",
    "                data=comparison)",
], font_size=13)

add_card(s, Inches(7.5), Inches(2.2), Inches(5.2), Inches(4.5),
         "산점도(Scatter Plot) 해석",
         ["\u25b6 x축: 실제값(actual), y축: 예측값(pred)",
          "",
          "\u2022 녹색 점선(대각선) 근처:",
          "   actual \u2248 pred \u2192 정확한 예측",
          "",
          "\u2022 대각선 위쪽:",
          "   pred > actual \u2192 과대예측",
          "",
          "\u2022 대각선 아래쪽:",
          "   pred < actual \u2192 과소예측",
          "",
          "\u2022 점들이 대각선에 가까울수록 좋은 모델"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 시각적 영역 표현
add_shape(s, Inches(0.6), Inches(4.5), Inches(6.5), Inches(2.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.8), Inches(4.6), Inches(6.1), Inches(0.4),
         "3개 영역 구분", font_size=14, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(0.8), Inches(5.0), Inches(6.1), Inches(1.8), [
    "\u2705 대각선 근처 = 예측 정확 (ideal zone)",
    "\u26a0\ufe0f 대각선 위쪽 = 과대예측 (over-prediction)",
    "\u26a0\ufe0f 대각선 아래쪽 = 과소예측 (under-prediction)",
    "\u2192 대부분의 점이 대각선 부근에 모여 있으나 이탈값도 존재",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

# ============================================================
# 슬라이드 15: 4.7.3 RMSE 평가
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.7.3", "RMSE 평가", "Root Mean Squared Error")

# MAE, MSE, RMSE 테이블
add_table_slide(s,
    headers=["지표", "수식", "특징"],
    rows=[["MAE", "|실제-예측| 평균", "직관적, 이상치에 덜 민감"],
          ["MSE", "(실제-예측)² 평균", "큰 오차에 더 큰 패널티"],
          ["RMSE", "√MSE", "단위가 원래와 같아 해석 용이"]],
    left=Inches(0.6), top=Inches(2.2),
    col_widths=[1.5, 3.5, 4.5],
    row_height=0.5, font_size=14)

add_card(s, Inches(0.6), Inches(4.2), Inches(5.5), Inches(1.5),
         "Table A vs Table B 비교",
         ["\u2022 절댓값 합(MAE)이 같아도 제곱 합(MSE)은 다를 수 있음",
          "\u2022 큰 오차 하나가 MSE를 크게 증가시킴",
          "\u2022 RMSE는 큰 오차에 민감한 지표"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_code_block(s, Inches(0.6), Inches(5.9), Inches(11.5), Inches(1.3), [
    "from sklearn.metrics import mean_squared_error",
    "",
    "mean_squared_error(y_test, pred) ** 0.5   # RMSE: 5684.93",
], font_size=14)

add_card(s, Inches(6.5), Inches(4.2), Inches(5.6), Inches(1.5),
         "결과: RMSE = 5684.93",
         ["\u2022 예측값과 실제값의 평균 오차가 약 $5,685",
          "\u2022 보험료 범위(1121~63770) 대비 판단",
          "\u2022 단독으로는 좋고 나쁨 판단 어려움"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 16: R² (결정 계수)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.7", "R\u00b2 (결정 계수)", "Coefficient of Determination")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.8),
         "R\u00b2 = SSR / SST",
         ["\u2022 SST (총변동): 실제값과 평균의 차이 제곱합",
          "\u2022 SSR (설명된 변동): 예측값과 평균의 차이 제곱합",
          "\u2022 SSE (설명 안 된 변동): 실제값과 예측값의 차이 제곱합",
          "",
          "SST = SSR + SSE",
          "R\u00b2 = SSR / SST = 1 - (SSE / SST)",
          "",
          "\u2022 0 ~ 1 사이 값, 1에 가까울수록 좋음"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_code_block(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(1.5), [
    "model.score(X_train, y_train)",
    "# 결과: 0.7368",
], font_size=14)

add_card(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(2.0),
         "R\u00b2 해석 기준",
         ["\u2022 R\u00b2 = 0.7368 \u2192 약 73.7% 설명력",
          "\u2022 독립변수가 종속변수 변동의 73.7%를 설명",
          "",
          "\u2022 0.7 ~ 0.8 이면 괜찮은(decent) 수준",
          "\u2022 0.9 이상이면 매우 좋은 모델"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(5.2), Inches(12.0), Inches(1.5),
         "R\u00b2 vs RMSE",
         ["\u2022 R\u00b2는 상대적 지표 (0~1), RMSE는 절대적 지표 (단위: 달러)",
          "\u2022 R\u00b2가 높으면 RMSE는 상대적으로 낮음",
          "\u2022 두 지표를 함께 보면 모델 성능을 종합적으로 판단 가능"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 17: 4.8 이해하기 - 선형 회귀 원리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.8", "이해하기: 선형 회귀 원리", "How Linear Regression Works")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.0),
         "y = ax + b",
         ["\u2022 a: 기울기 (coefficient, 계수)",
          "\u2022 b: y절편 (intercept)",
          "\u2022 독립변수(x)와 종속변수(y)의 선형 관계",
          "\u2022 여러 독립변수: y = a\u2081x\u2081 + a\u2082x\u2082 + ... + b"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "손실 함수 (Loss Function)",
         ["\u2022 예측값과 실제값의 차이를 최소화",
          "\u2022 OLS: (실제-예측)\u00b2 합을 최소화하는 선 찾기",
          "\u2022 최적의 기울기(a)와 절편(b)을 찾는 것이 학습"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# 보험료 예측 수식
add_shape(s, Inches(0.6), Inches(4.3), Inches(12.0), Inches(1.5), CODE_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.8), Inches(4.4), Inches(11.6), Inches(0.4),
         "보험료 예측 수식 (학습된 모델)", font_size=15, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(0.8), Inches(4.9), Inches(11.6), Inches(0.8),
         "charges = 264.8\u00d7age + 17.3\u00d7sex + 297.5\u00d7bmi + 469.3\u00d7children + 23469.3\u00d7smoker - 11577.0",
         font_size=16, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_card(s, Inches(0.6), Inches(6.0), Inches(12.0), Inches(1.2),
         "수식 해석",
         ["age 1 증가 \u2192 charges 약 265 증가  |  smoker=1(흡연) \u2192 charges 약 23,469 증가  |  절편(b) = -11,577"],
         title_color=ACCENT_YELLOW, border=ACCENT_YELLOW)

# ============================================================
# 슬라이드 18: 4.8 계수와 절편 확인
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.8", "계수와 절편 확인", "Coefficients & Intercept")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(7.0), Inches(3.0), [
    "# 계수(기울기) 확인",
    "model.coef_",
    "# array([ 264.8, 17.3, 297.5, 469.3, 23469.3])",
    "",
    "# 계수를 보기 좋게 출력",
    "pd.Series(model.coef_, index=X.columns)",
    "# age        264.8",
    "# sex         17.3",
    "# bmi        297.5",
    "# children   469.3",
    "# smoker   23469.3",
    "",
    "# 절편(y-intercept) 확인",
    "model.intercept_",
    "# -11576.999976112367",
], font_size=12)

add_card(s, Inches(8.0), Inches(2.2), Inches(4.6), Inches(3.0),
         "계수 해석",
         ["\u2022 age: 나이 1세 증가 \u2192 보험료 약 $265 증가",
          "\u2022 sex: 성별에 따른 차이 약 $17 (매우 작음)",
          "\u2022 bmi: BMI 1 증가 \u2192 보험료 약 $298 증가",
          "\u2022 children: 자녀 1명 증가 \u2192 약 $469 증가",
          "\u2022 smoker: 흡연자 \u2192 약 $23,469 증가!",
          "",
          "\u2192 흡연 여부가 보험료에 가장 큰 영향"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(0.6), Inches(5.5), Inches(12.0), Inches(1.5),
         "절편 (Intercept) = -11,577",
         ["\u2022 모든 독립변수가 0일 때의 기본값 (실제로는 발생하지 않는 이론적 값)",
          "\u2022 계수(coef_)와 절편(intercept_)을 합치면 전체 예측 수식이 완성됨"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 19: 학습 마무리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
slide_header(s, "", "학습 마무리", "Chapter 04 Summary")

# 파이프라인 다이어그램
steps_full = [
    ("문제 정의", ACCENT_BLUE),
    ("데이터 로드", ACCENT_CYAN),
    ("데이터 확인", ACCENT_GREEN),
    ("전처리", ACCENT_ORANGE),
    ("모델링", ACCENT_PURPLE),
    ("예측", ACCENT_RED),
    ("평가", ACCENT_YELLOW),
]
for i, (label, clr) in enumerate(steps_full):
    x = Inches(0.6 + i * 1.75)
    add_shape(s, x, Inches(2.3), Inches(1.5), Inches(0.8), clr, radius=True)
    add_text(s, x, Inches(2.3), Inches(1.5), Inches(0.8), label,
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps_full) - 1:
        add_text(s, x + Inches(1.45), Inches(2.3), Inches(0.35), Inches(0.8), "\u2192",
                 font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_card(s, Inches(0.6), Inches(3.5), Inches(5.8), Inches(3.5),
         "핵심 요약",
         ["\u2022 선형 회귀: y = a\u2081x\u2081 + a\u2082x\u2082 + ... + b",
          "\u2022 연속형 변수 예측에 사용되는 기초 알고리즘",
          "\u2022 sklearn으로 간단히 구현 가능",
          "   LinearRegression().fit() \u2192 predict()",
          "",
          "\u2022 평가: RMSE = 5,685 / R\u00b2 = 0.7368",
          "\u2022 흡연 여부(smoker)가 보험료에 가장 큰 영향",
          "",
          "\u2022 장점: 간단, 해석 용이, 빠른 학습",
          "\u2022 한계: 비선형 관계 포착 어려움"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(3.5), Inches(5.8), Inches(3.5),
         "전체 코드 파이프라인",
         ["1. import pandas as pd",
          "2. data = pd.read_csv(file_url)",
          "3. data.head() / data.info() / data.describe()",
          "4. X, y 분리 \u2192 train_test_split()",
          "5. model = LinearRegression()",
          "6. model.fit(X_train, y_train)",
          "7. pred = model.predict(X_test)",
          "8. pd.DataFrame({'actual':y_test, 'pred':pred})",
          "9. sns.scatterplot(x='actual', y='pred', ...)",
          "10. mean_squared_error(y_test, pred)**0.5",
          "11. model.score() / model.coef_ / model.intercept_"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ── 저장 ──
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "4장_선형회귀_교재실습.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
