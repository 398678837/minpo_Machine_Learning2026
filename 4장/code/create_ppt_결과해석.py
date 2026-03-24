"""04 Linear Regression 결과해석 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)
ACCENT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=14, color=LIGHT_GRAY, bold=False, spacing=Pt(4), font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = font_name; p.space_after = spacing
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = '맑은 고딕'; p.space_after = spacing; p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))


def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)


def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)


def add_table_custom(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                     row_height=0.45, font_size=12, header_font_size=13):
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.45), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.45),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    for i, row in enumerate(rows):
        y = top + Inches(0.45) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)
            cx += Inches(w)


# ============================================================
# Slide 1 - Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG)
add_accent_line(s, Inches(4.5), Inches(1.8), Inches(4.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.0), Inches(11.333), Inches(1.2),
         "04 Linear Regression\n실행 결과 상세 해석", font_size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(4.5), Inches(3.5), Inches(4.333), ACCENT_CYAN)

add_text(s, Inches(1), Inches(4.0), Inches(11.333), Inches(0.5),
         "교수자: Jung, Minpo", font_size=20, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.6), Inches(11.333), Inches(0.5),
         "교과목: Machine Learning", font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.1), Inches(11.333), Inches(0.5),
         "학기: 2026년도 1학기", font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.7), Inches(11.333), Inches(0.5),
         "Chapter 04 - 보험료 예측 선형 회귀 모델 결과 분석", font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2 - 1. 데이터 로드 및 전체 구조 확인
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "데이터 로드 및 전체 구조 확인", "1338명 보험 가입자 데이터 (6개 변수)")

# Code block
add_code_block(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(2.2), [
    "      age  sex     bmi  children  smoker      charges",
    "0      19    0  27.900         0       1  16884.92400",
    "...                                                  ",
    "1337   61    0  29.070         0       1  29141.36030",
    "",
    "[1338 rows x 6 columns]"
], font_size=11)

# Interpretation card
add_card(s, Inches(6.5), Inches(2.0), Inches(6.2), Inches(2.2),
         "해석 포인트",
         ["총 1338명의 보험 가입자 데이터 정상 로드",
          "6개 변수: age, sex, bmi, children, smoker, charges",
          "인덱스: 0부터 1337까지 (파이썬 기본 인덱스)"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# Row interpretation cards
add_card(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.5),
         "첫 번째 데이터 (인덱스 0)",
         ["19세, 여성(0), BMI 27.9, 자녀 없음, 흡연자(1)",
          "보험료: $16,884",
          "젊은 나이의 흡연 여성 → 비교적 높은 보험료"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.5),
         "마지막 데이터 (인덱스 1337)",
         ["61세, 여성(0), BMI 29.07, 자녀 없음, 흡연자(1)",
          "보험료: $29,141",
          "고령 흡연 여성 → 나이 영향으로 더 높은 보험료"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# Slide 3 - 2. 상위 5개 데이터 확인 (head)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "상위 5개 데이터 확인 (head)", "head() 함수로 데이터 미리보기")

# Table
headers = ["idx", "age", "sex", "bmi", "children", "smoker", "charges"]
rows = [
    ["0", "19", "0", "27.900", "0", "1", "$16,884.92"],
    ["1", "18", "1", "33.770", "1", "0", "$1,725.55"],
    ["2", "28", "1", "33.000", "3", "0", "$4,449.46"],
    ["3", "33", "1", "22.705", "0", "0", "$21,984.47"],
    ["4", "32", "1", "28.880", "0", "0", "$3,866.86"],
]
add_table_custom(s, headers, rows, Inches(0.6), Inches(2.0),
                 [0.7, 0.8, 0.8, 1.1, 1.1, 1.1, 1.8], row_height=0.42, font_size=12)

# Encoding cards
add_card(s, Inches(0.6), Inches(4.6), Inches(3.5), Inches(1.5),
         "sex 인코딩",
         ["0 = 여성 (Female)", "1 = 남성 (Male)", "범주형 → 숫자 변환 완료"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(4.4), Inches(4.6), Inches(3.5), Inches(1.5),
         "smoker 인코딩",
         ["0 = 비흡연자", "1 = 흡연자", "범주형 → 숫자 변환 완료"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# Insight
add_card(s, Inches(8.2), Inches(4.6), Inches(4.5), Inches(2.5),
         "인덱스 0 vs 3 비교 인사이트",
         ["0번: 19세, 흡연자 → $16,884",
          "3번: 33세, 비흡연자 → $21,984",
          "나이가 14살 더 많지만 비흡연자인 3번이 더 높은 보험료",
          "→ 흡연 여부 외에도 다른 요인이 복합적으로 작용"],
         title_color=ACCENT_YELLOW, border=ACCENT_YELLOW)


# ============================================================
# Slide 4 - 3. 데이터 구조 정보 (info)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "데이터 구조 정보 (info)", "데이터 타입, 결측치, 메모리 사용량 확인")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(3.5), [
    "RangeIndex: 1338 entries, 0 to 1337",
    "Data columns (total 6 columns):",
    " #   Column    Non-Null Count  Dtype  ",
    "---  ------    --------------  -----  ",
    " 0   age       1338 non-null   int64  ",
    " 1   sex       1338 non-null   int64  ",
    " 2   bmi       1338 non-null   float64",
    " 3   children  1338 non-null   int64  ",
    " 4   smoker    1338 non-null   int64  ",
    " 5   charges   1338 non-null   float64",
    "dtypes: float64(2), int64(4)",
    "memory usage: 62.8 KB"
], font_size=11)

add_card(s, Inches(6.5), Inches(2.0), Inches(6.2), Inches(1.5),
         "결측치(Missing Value) 없음",
         ["모든 변수가 1338 non-null → 결측치 전혀 없음",
          "실무에서는 매우 드문 경우 (교재용 사전 정제 데이터)",
          "별도 데이터 클리닝 없이 바로 모델링 가능"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(6.5), Inches(3.8), Inches(6.2), Inches(1.8),
         "자료형(Dtype) 해석",
         ["int64 (정수): age, sex, children, smoker",
          "  → sex, smoker는 범주형이지만 숫자로 인코딩됨",
          "float64 (실수): bmi, charges",
          "  → 소수점이 있는 연속형 변수"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(6.5), Inches(5.9), Inches(6.2), Inches(1.2),
         "메모리 사용량: 62.8 KB",
         ["매우 작은 데이터셋 (실무에서는 수 GB 이상도 흔함)",
          "학습용으로 적합한 크기"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# Slide 5 - 4. 기술 통계량 (describe) - 변수별 해석 1
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04-1", "기술 통계량 (describe) - age, sex, bmi", "주요 변수별 상세 해석")

# age card
add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(4.5),
         "age (연령)",
         ["mean = 39.21  → 평균 연령 약 39세",
          "std = 14.05  → 표준편차 14세 (분포 넓음)",
          "min = 18.00  → 최소 18세 (성인만 포함)",
          "50% = 39.00  → 중앙값 = 평균 (대칭 분포)",
          "max = 64.00  → 최대 64세",
          "",
          "평균과 중앙값이 거의 일치 →",
          "연령은 비교적 균등한 분포를 가짐"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

# sex card
add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.5),
         "sex (성별)",
         ["mean = 0.51",
          "→ 남성(1) 비율 약 51%",
          "→ 남녀 비율이 거의 균등",
          "",
          "이진 변수(0/1)이므로",
          "mean = 남성 비율로 해석",
          "",
          "성별 균형이 잡힌 데이터"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# bmi card
add_card(s, Inches(8.8), Inches(2.0), Inches(3.9), Inches(4.5),
         "bmi (체질량지수)",
         ["mean = 30.66",
          "→ WHO 기준 BMI 30+ = '비만'",
          "→ 평균이 비만 경계에 위치",
          "",
          "min = 15.96  → 저체중 수준",
          "max = 53.13  → 고도 비만",
          "",
          "미국 보험 가입자 데이터로",
          "비만율이 높은 특성이 반영됨"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# Slide 6 - 4. 기술 통계량 (describe) - 변수별 해석 2
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04-2", "기술 통계량 (describe) - children, smoker, charges", "종속변수 charges의 우편향 분포가 핵심")

# children card
add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(2.5),
         "children (자녀 수)",
         ["mean = 1.09  → 평균 자녀 약 1명",
          "max = 5.00  → 최대 5명",
          "75% = 2.00  → 75%가 자녀 2명 이하"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

# smoker card
add_card(s, Inches(0.6), Inches(4.8), Inches(3.8), Inches(2.3),
         "smoker (흡연 여부)",
         ["mean = 0.20",
          "→ 흡연자 비율 약 20%",
          "→ 비흡연자가 80%로 압도적"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# charges card
add_card(s, Inches(4.7), Inches(2.0), Inches(4.0), Inches(5.1),
         "charges (보험 청구비용) - 종속변수",
         ["mean = $13,270.42  (평균 보험료)",
          "std = $12,110.01  (표준편차 ≒ 평균! 매우 넓은 분포)",
          "min = $1,121.87  (최저)",
          "50% = $9,382.03  (중앙값)",
          "max = $63,770.43  (최고 = 최소의 약 57배)",
          "",
          "핵심: 평균($13,270) > 중앙값($9,382)",
          "→ 우측으로 치우친 분포 (right-skewed)"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# Insight box
add_card(s, Inches(9.0), Inches(2.0), Inches(3.7), Inches(5.1),
         "핵심 인사이트",
         ["charges의 평균 > 중앙값",
          "→ 우편향(right-skewed) 분포",
          "",
          "소수의 고액 청구자가",
          "평균을 끌어올리는 구조",
          "",
          "이는 흡연자 그룹의",
          "높은 보험료 때문일",
          "가능성이 높다",
          "",
          "선형 회귀에서 이러한",
          "비대칭 분포는 예측",
          "정확도에 영향을 줌"],
         title_color=ACCENT_YELLOW, border=ACCENT_YELLOW)


# ============================================================
# Slide 7 - 5. 예측값 vs 실제값 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "예측값 vs 실제값 비교", "시험셋 268개 (전체 1338의 20%) 예측 결과")

# Table
headers = ["인덱스", "실제값(actual)", "예측값(pred)", "오차", "판정"]
rows = [
    ["12", "$1,827", "$4,765", "+$2,938", "과대예측(2.6배)"],
    ["306", "$20,178", "$4,958", "-$15,220", "심각한 과소예측"],
    ["318", "$7,421", "$8,299", "+$878", "비교적 정확"],
    ["815", "$1,878", "$3,079", "+$1,201", "과대예측"],
    ["157", "$15,518", "$24,166", "+$8,648", "과대예측"],
    ["1015", "$12,125", "$11,638", "-$487", "매우 정확"],
]
add_table_custom(s, headers, rows, Inches(0.6), Inches(2.1),
                 [1.0, 1.8, 1.8, 1.5, 2.0], row_height=0.42, font_size=12)

# Case analysis cards
add_card(s, Inches(7.2), Inches(2.1), Inches(5.5), Inches(2.0),
         "인덱스 306 - 심각한 과소예측",
         ["실제 $20,178 → 예측 $4,958 (오차 -$15,220)",
          "아마 흡연자인데 모델이 제대로 포착 못한 케이스",
          "선형 회귀의 한계를 보여주는 사례"],
         title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(7.2), Inches(4.4), Inches(5.5), Inches(1.8),
         "인덱스 1015 - 매우 정확한 예측",
         ["실제 $12,125 → 예측 $11,638 (오차 -$487)",
          "오차율 약 4%로 매우 우수한 예측",
          "비흡연자 중가 보험료 구간에서 모델이 잘 작동"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# Slide 8 - 6. 산점도 시각화 해석
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "산점도 시각화 해석", "x축 = 실제값, y축 = 예측값 — 완벽한 모델이면 45도 대각선 위에 위치")

# Three region cards
add_card(s, Inches(0.6), Inches(2.1), Inches(3.9), Inches(2.0),
         "저가 보험료 구간 ($0~$15,000)",
         ["점들이 대각선 근처에 비교적 잘 분포",
          "→ 예측이 양호한 구간",
          "비흡연자 대부분이 이 구간에 위치"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(4.8), Inches(2.1), Inches(3.9), Inches(2.0),
         "중가 보험료 구간 ($15,000~$35,000)",
         ["점들이 대각선에서 벗어나기 시작",
          "→ 과대/과소예측 혼재",
          "흡연자와 비흡연자 경계 구간"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(9.0), Inches(2.1), Inches(3.7), Inches(2.0),
         "고가 보험료 구간 ($35,000+)",
         ["점들이 크게 분산됨",
          "→ 예측이 불안정한 구간",
          "고액 흡연자 보험료 예측 난항"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# 3 clusters explanation
add_card(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.7),
         "3개 클러스터가 보이는 이유",
         ["산점도에서 데이터가 3개의 뚜렷한 군집(cluster)으로 나뉘는 것이 관찰됨",
          "",
          "Cluster 1: 비흡연자 (저가 보험료) — 대각선에 비교적 가까이 위치",
          "Cluster 2: 흡연자 + 낮은 BMI — 중간 보험료 구간",
          "Cluster 3: 흡연자 + 높은 BMI — 고가 보험료 구간",
          "",
          "원인: 흡연자/비흡연자 그룹 간 보험료 차이가 크기 때문",
          "선형 회귀는 이런 그룹 간 비선형적 차이를 완벽히 포착하기 어렵다"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)


# ============================================================
# Slide 9 - 7. RMSE 해석
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "RMSE (Root Mean Squared Error)", "평균 예측 오차의 크기를 나타내는 지표")

# RMSE value
add_shape(s, Inches(0.6), Inches(2.1), Inches(4.5), Inches(1.2), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.6), Inches(2.2), Inches(4.5), Inches(0.5),
         "RMSE = $5,684.93", font_size=28, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.8), Inches(4.5), Inches(0.4),
         "RMSE / 평균 = 5,685 / 13,270 = 약 42.8%", font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Interpretation
add_card(s, Inches(5.5), Inches(2.1), Inches(7.2), Inches(1.6),
         "해석",
         ["\"평균적으로 예측이 실제값에서 약 $5,685 정도 벗어난다\"",
          "평균 보험료($13,270) 대비 약 42.8% 오차 → 상당히 큰 오차",
          "RMSE의 절대적 기준은 없으며, 다른 모델과 상대 비교에 사용"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

# Formula
add_code_block(s, Inches(0.6), Inches(3.6), Inches(7.0), Inches(2.0), [
    "RMSE 계산 과정:",
    "",
    "RMSE = sqrt( (1/n) * sum( (yi - y_hat_i)^2 ) )",
    "",
    "1. 각 데이터에 대해 (실제값 - 예측값) 차이 계산",
    "2. 차이를 제곱 (부호 제거 + 큰 오차에 패널티)",
    "3. 제곱 값들의 평균 구함 (MSE)",
    "4. 평균에 루트 (원래 스케일로 복원)"
], font_size=12)

# Comparison purpose
add_card(s, Inches(8.0), Inches(3.9), Inches(4.7), Inches(2.8),
         "RMSE 비교 목적",
         ["향후 다른 모델 적용 시 RMSE 비교:",
          "",
          "Ridge Regression → RMSE ?",
          "Lasso Regression → RMSE ?",
          "Random Forest → RMSE ?",
          "XGBoost → RMSE ?",
          "",
          "이 RMSE($5,685)를 낮추는 것이 목표"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# Slide 10 - 8. R² 결정 계수
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "R² 결정 계수 (Coefficient of Determination)", "모델의 설명력을 나타내는 핵심 지표")

# R² value
add_shape(s, Inches(0.6), Inches(2.1), Inches(4.5), Inches(1.2), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.6), Inches(2.2), Inches(4.5), Inches(0.5),
         "R² = 0.7368 (73.7%)", font_size=28, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.8), Inches(4.5), Inches(0.4),
         "독립변수 5개가 charges 변동의 73.7% 설명", font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# SST/SSR/SSE
add_card(s, Inches(5.5), Inches(2.1), Inches(7.2), Inches(1.5),
         "R² = SSR / SST = 1 - SSE / SST",
         ["SST (총 변동량): 종속변수의 총 변동량 (평균값 대비)",
          "SSR (회귀 변동량): 모델이 설명하는 변동량",
          "SSE (오차 변동량): 모델이 설명하지 못하는 변동량"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

# R² range table
headers = ["R² 범위", "해석"]
rows_table = [
    ["0.9 이상", "매우 우수한 모델"],
    ["0.7 ~ 0.9", "양호한 모델  ← 우리 모델 (0.74)"],
    ["0.5 ~ 0.7", "보통"],
    ["0.5 미만", "설명력 부족"],
]
add_table_custom(s, headers, rows_table, Inches(0.6), Inches(3.6),
                 [2.5, 5.0], row_height=0.5, font_size=14, header_font_size=15)

# Remaining 26.3%
add_card(s, Inches(8.5), Inches(3.8), Inches(4.2), Inches(3.2),
         "남은 26.3%는?",
         ["현재 모델에 포함되지 않은 요인:",
          "",
          "- 기저 질환 유무",
          "- 거주 지역",
          "- 직업",
          "- 운동 습관",
          "",
          "또는 변수 간 비선형 관계",
          "(예: 흡연자의 BMI 효과가",
          " 비흡연자와 다를 수 있음)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# Slide 11 - 9. 회귀 계수 해석
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "회귀 계수 (Coefficients) 해석", "각 독립변수가 보험료에 미치는 영향력")

# Coefficients table
headers = ["변수", "계수", "해석"]
rows_coef = [
    ["smoker", "+23,469.28", "가장 큰 영향력. 흡연자는 비흡연자 대비 ~$23,469 높음"],
    ["children", "+469.34", "자녀 1명 증가 시 보험료 약 $469 증가"],
    ["bmi", "+297.51", "BMI 1 증가 시 보험료 약 $298 증가"],
    ["age", "+264.80", "나이 1살 증가 시 보험료 약 $265 증가"],
    ["sex", "+17.34", "남성이 여성 대비 약 $17 높음 → 사실상 영향 미미"],
]
add_table_custom(s, headers, rows_coef, Inches(0.6), Inches(2.1),
                 [1.5, 1.8, 6.0], row_height=0.48, font_size=13, header_font_size=14)

# 4 key insights
add_card(s, Inches(0.6), Inches(4.8), Inches(3.0), Inches(2.3),
         "1. smoker 압도적 1위",
         ["계수 23,469로 다른 변수보다",
          "수십~수백 배 큼",
          "보험료의 가장 큰 결정 요인"],
         title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(3.9), Inches(4.8), Inches(3.0), Inches(2.3),
         "2. sex는 거의 무의미",
         ["계수 17로 실질적 영향 없음",
          "성별보다 건강 관련 변수가",
          "보험료에 훨씬 중요"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(7.2), Inches(4.8), Inches(3.0), Inches(2.3),
         "3. BMI와 나이 비슷",
         ["각각 298, 265로",
          "유사한 크기의 영향",
          "건강/연령 요인이 비슷한 가중치"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(10.5), Inches(4.8), Inches(2.2), Inches(2.3),
         "4. 자녀 수 유의미",
         ["1명당 약 $469 증가",
          "부양가족 반영"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ============================================================
# Slide 12 - 10. 절편 해석
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "절편 (Intercept) 해석", "회귀 수식의 상수항")

# Intercept value
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.5), Inches(1.4), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.6), Inches(2.3), Inches(5.5), Inches(0.6),
         "Intercept = -$11,577.00", font_size=30, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.95), Inches(5.5), Inches(0.5),
         "(-11576.999976112374)", font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Explanation
add_card(s, Inches(6.5), Inches(2.2), Inches(6.2), Inches(2.5),
         "절편의 의미",
         ["수식에서 모든 독립변수가 0일 때의 charges 값",
          "",
          "\"0세, 여성, BMI 0, 자녀 0, 비흡연\" 인 사람의 보험료",
          "→ 실제로 이런 사람은 존재하지 않음",
          "",
          "따라서 절편 자체에 현실적 의미를 부여하기보다",
          "수학적 보정값으로 이해하는 것이 적절"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

add_card(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(2.0),
         "음수인 이유",
         ["다른 계수들의 합이 양수로 크기 때문에 절편이 음수가 되어야",
          "전체 예측값이 합리적인 범위에 들어온다",
          "",
          "예시: 18세 비흡연자의 경우",
          "  264.80×18 + 17.34×0 + 297.51×25 + 469.34×0 + 23469.28×0 + (-11577.00) = $4,230",
          "  절편이 0이었다면 $15,807로 과대 예측됨"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)


# ============================================================
# Slide 13 - 11. 최종 예측 수식
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "11", "최종 예측 수식", "5개 독립변수와 절편으로 구성된 선형 회귀 모델")

# Formula
add_code_block(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.0), [
    "charges = 264.80 × age + 17.34 × sex + 297.51 × bmi + 469.34 × children + 23469.28 × smoker - 11577.00"
], font_size=14)

# Example 1: Non-smoker
add_card(s, Inches(0.6), Inches(3.4), Inches(5.8), Inches(3.7),
         "예측 예시 1: 30세 남성 BMI 25, 자녀 1명, 비흡연자",
         ["charges = 264.80×30 + 17.34×1 + 297.51×25",
          "        + 469.34×1 + 23469.28×0 - 11577.00",
          "",
          "= 7,944.00 + 17.34 + 7,437.75",
          "  + 469.34 + 0 - 11,577.00",
          "",
          "= $4,291.43",
          "",
          "→ 비흡연자의 합리적인 보험료 수준"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# Example 2: Smoker
add_card(s, Inches(6.8), Inches(3.4), Inches(5.8), Inches(3.7),
         "예측 예시 2: 같은 조건의 흡연자",
         ["charges = $4,291.43 + 23,469.28",
          "",
          "= $27,760.71",
          "",
          "→ 흡연 여부 하나만 바뀌어도",
          "  보험료가 약 6.5배 차이!",
          "",
          "$4,291 vs $27,761",
          "",
          "흡연이 보험료에 미치는",
          "압도적 영향력을 확인"],
         title_color=ACCENT_RED, border=ACCENT_RED)


# ============================================================
# Slide 14 - 12. 종합 평가
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "12", "종합 평가", "Linear Regression 모델의 성능, 강점, 한계, 개선 방향")

# Performance summary
headers_perf = ["지표", "값", "평가"]
rows_perf = [
    ["RMSE", "$5,684.93", "평균 보험료 대비 약 43% 오차"],
    ["R²", "0.7368", "양호 (73.7% 설명력)"],
]
add_table_custom(s, headers_perf, rows_perf, Inches(0.6), Inches(2.1),
                 [1.5, 2.0, 4.0], row_height=0.48, font_size=14, header_font_size=15)

# Strengths
add_card(s, Inches(0.6), Inches(3.3), Inches(3.8), Inches(2.0),
         "강점",
         ["해석이 매우 직관적: 각 변수의",
          "  영향을 정량적으로 파악 가능",
          "훈련 속도가 매우 빠름",
          "흡연 = 핵심 요인 명확히 포착"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# Limitations
add_card(s, Inches(4.7), Inches(3.3), Inches(3.8), Inches(2.0),
         "한계",
         ["흡연자 내 세밀한 예측 부족",
          "  (비선형 관계 미포착)",
          "우편향 분포 → 고액 구간 오차 큼",
          "교호작용(interaction) 미고려"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# Improvement directions
add_card(s, Inches(8.8), Inches(3.3), Inches(3.9), Inches(2.0),
         "개선 방향",
         ["1. 로그 변환: charges에 log 취해 분포 정규화",
          "2. 교호작용: smoker×bmi 등 상호작용 항",
          "3. 비선형 모델: RF, XGBoost 등",
          "4. 정규화: Ridge, Lasso, Elastic Net"],
         title_color=ACCENT_CYAN, border=ACCENT_BLUE)

# Final note
add_shape(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.5), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
         "결론: Linear Regression은 기본 모델로서 데이터의 핵심 패턴을 잘 포착하였으나,",
         font_size=16, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
         "비선형성과 교호작용을 반영하는 고급 모델로의 확장이 필요하다. (다음 챕터에서 계속)",
         font_size=16, color=ACCENT_PURPLE)

# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "04_Linear Regression_결과해석.pptx")
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
