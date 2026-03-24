"""04 Linear Regression 运行结果详细解析 - PPT 生成脚本 (简体中文)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 色상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)

FONT_NAME = 'Microsoft YaHei'


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = FONT_NAME; p.space_after = spacing; p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))


def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)


def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)


def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)
            cx += Inches(w)


# ============================================================
# Slide 1 - Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Inches(7.5), SECTION_BG)
add_accent_line(s, Inches(4.5), Inches(1.8), Inches(4.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.0), Inches(11.333), Inches(1.2),
         "04 Linear Regression\n运行结果详细解析", font_size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.6), Inches(11.333), Inches(0.5),
         "授课教师: Jung, Minpo", font_size=22, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.2), Inches(11.333), Inches(0.5),
         "课程: Machine Learning  |  学期: 2026年第1学期", font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.8), Inches(11.333), Inches(0.5),
         "对象: 博士课程（中国博士课程学生）", font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 2 - 1. 数据加载与整体结构确认
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "数据加载与整体结构确认", "1338 rows x 6 columns — 保险投保人数据")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.6), [
    "      age  sex     bmi  children  smoker      charges",
    "0      19    0  27.900         0       1  16884.92400",
    "...",
    "1337   61    0  29.070         0       1  29141.36030",
    "",
    "[1338 rows x 6 columns]"
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(2.0),
         "数据概要", [
             "共1338名保险投保人数据已正常加载",
             "6个变量: age, sex, bmi, children, smoker, charges",
             "索引: 从0到1337 (Python默认索引)"
         ], ACCENT_CYAN, ACCENT_BLUE)

add_card(s, Inches(0.6), Inches(4.0), Inches(5.8), Inches(3.0),
         "6个变量说明", [
             "age — 年龄 (18~64岁)",
             "sex — 性别 (0=女性, 1=男性)",
             "bmi — 体质量指数 (Body Mass Index)",
             "children — 子女数 (0~5)",
             "smoker — 吸烟与否 (0=非吸烟, 1=吸烟)",
             "charges — 保险理赔费用 (因变量, 美元)"
         ], ACCENT_GREEN, ACCENT_GREEN)

add_card(s, Inches(7.0), Inches(4.3), Inches(5.8), Inches(2.5),
         "首行/末行数据解读", [
             "索引0: 19岁, 女性(0), BMI 27.9, 无子女, 吸烟者(1)",
             "  -> 保险费 $16,884",
             "索引1337: 61岁, 女性(0), BMI 29.07, 无子女, 吸烟者(1)",
             "  -> 保险费 $29,141",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# Slide 3 - 2. 前5行数据确认 (head)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "前5行数据确认 (head)", "head() 输出 — 数据样本检查")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(2.2), [
    "   age  sex     bmi  children  smoker      charges",
    "0   19    0  27.900         0       1  16884.92400",
    "1   18    1  33.770         1       0   1725.55230",
    "2   28    1  33.000         3       0   4449.46200",
    "3   33    1  22.705         0       0  21984.47061",
    "4   32    1  28.880         0       0   3866.85520"
], font_size=12)

add_card(s, Inches(8.5), Inches(2.0), Inches(4.3), Inches(1.6),
         "编码说明", [
             "sex: 0 = 女性, 1 = 男性",
             "smoker: 0 = 非吸烟者, 1 = 吸烟者",
             "(分类型 -> 数值转换已完成)"
         ], ACCENT_CYAN, ACCENT_BLUE)

add_card(s, Inches(0.6), Inches(4.6), Inches(5.8), Inches(2.5),
         "索引0号 vs 3号 对比分析", [
             "0号: 19岁, 吸烟者 -> $16,884",
             "3号: 33岁, 非吸烟者 -> $21,984",
             "",
             "虽然年龄大14岁，但非吸烟者3号的保险费更高",
             "-> 表明除吸烟与否外，其他因素也在综合发挥作用"
         ], ACCENT_ORANGE, ACCENT_ORANGE)

add_card(s, Inches(7.0), Inches(4.6), Inches(5.8), Inches(2.5),
         "核心发现", [
             "吸烟与否(smoker)并非唯一决定因素",
             "年龄、BMI、子女数等变量共同影响保险费",
             "-> 需要多元回归(Multiple Regression)分析"
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# Slide 4 - 3. 数据结构信息 (info)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "数据结构信息 (info)", "数据类型、缺失值、内存使用量确认")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(3.0), [
    "RangeIndex: 1338 entries, 0 to 1337",
    "Data columns (total 6 columns):",
    " #   Column    Non-Null Count  Dtype",
    "---  ------    --------------  -----",
    " 0   age       1338 non-null   int64",
    " 1   sex       1338 non-null   int64",
    " 2   bmi       1338 non-null   float64",
    " 3   children  1338 non-null   int64",
    " 4   smoker    1338 non-null   int64",
    " 5   charges   1338 non-null   float64",
    "dtypes: float64(2), int64(4)",
    "memory usage: 62.8 KB"
], font_size=11)

add_card(s, Inches(7.5), Inches(2.0), Inches(5.3), Inches(1.5),
         "无缺失值", [
             "所有变量均为 1338 non-null",
             "实际工作中非常罕见 -> 教材预先清洗过的数据",
         ], ACCENT_GREEN, ACCENT_GREEN)

add_card(s, Inches(7.5), Inches(3.8), Inches(5.3), Inches(1.8),
         "数据类型 (Dtype)", [
             "int64 (整数): age, sex, children, smoker",
             "  -> sex和smoker原为分类变量, 已编码为数值",
             "float64 (浮点数): bmi, charges",
             "  -> 带小数点的连续型变量"
         ], ACCENT_CYAN, ACCENT_BLUE)

add_card(s, Inches(7.5), Inches(5.9), Inches(5.3), Inches(1.2),
         "内存使用量: 62.8 KB", [
             "非常小的数据集 (实际工作中数GB以上也常见)",
             "无需额外清洗, 可直接进入建模阶段"
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# Slide 5 - 4. 描述性统计 (describe) - 变量解析 1
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "描述性统计 (describe) — 变量解析 1", "age, sex, bmi 详细解析")

# age card
add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(5.0),
         "age (年龄)", [
             "mean = 39.21 -> 平均年龄约39岁",
             "std = 14.05 -> 标准差14岁, 分布较广",
             "min = 18.00 -> 最小18岁 (仅成年人)",
             "50% = 39.00 -> 中位数=平均值",
             "  -> 暗示对称分布",
             "max = 64.00 -> 最大64岁",
             "",
             "年龄范围: 18~64岁",
             "分布接近均匀/对称"
         ], ACCENT_CYAN, ACCENT_BLUE)

# sex card
add_card(s, Inches(4.8), Inches(2.0), Inches(3.8), Inches(5.0),
         "sex (性别)", [
             "mean = 0.51",
             "  -> 男性(1)比例约51%",
             "  -> 男女比例几乎均等",
             "",
             "编码方式:",
             "  0 = 女性, 1 = 男性",
             "",
             "性别分布均衡,",
             "不存在严重的类别不平衡问题"
         ], ACCENT_GREEN, ACCENT_GREEN)

# bmi card
add_card(s, Inches(9.0), Inches(2.0), Inches(3.8), Inches(5.0),
         "bmi (体质量指数)", [
             "mean = 30.66",
             "  -> WHO标准: BMI>=30 为'肥胖'",
             "  -> 平均值位于肥胖临界值",
             "",
             "min = 15.96 -> 体重偏低水平",
             "max = 53.13 -> 重度肥胖",
             "",
             "BMI分布范围很广,",
             "涵盖偏低到重度肥胖"
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# Slide 6 - 4. 描述性统计 (describe) - 变量解析 2
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "描述性统计 (describe) — 变量解析 2", "children, smoker, charges 详细解析")

# children card
add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(3.0),
         "children (子女数)", [
             "mean = 1.09 -> 平均约1个子女",
             "max = 5.00 -> 最多5个",
             "75% = 2.00 -> 75%的人子女在2个以下",
         ], ACCENT_CYAN, ACCENT_BLUE)

# smoker card
add_card(s, Inches(4.8), Inches(2.0), Inches(3.8), Inches(3.0),
         "smoker (吸烟与否)", [
             "mean = 0.20",
             "  -> 吸烟者比例约20%",
             "  -> 非吸烟者占80%, 占压倒性多数",
         ], ACCENT_GREEN, ACCENT_GREEN)

# charges card
add_card(s, Inches(9.0), Inches(2.0), Inches(3.8), Inches(3.0),
         "charges (保险理赔费用) — 因变量", [
             "mean = $13,270  |  std = $12,110",
             "  -> 标准差与平均值几乎相同, 分布非常广",
             "min = $1,122  |  max = $63,770",
             "  -> 最高约为最低的57倍",
             "50% = $9,382 -> 远低于平均值",
         ], ACCENT_RED, ACCENT_RED)

# 核心发现
add_card(s, Inches(0.6), Inches(5.3), Inches(12.2), Inches(1.8),
         "核心发现: charges 呈右偏(right-skewed)分布", [
             "平均值($13,270) > 中位数($9,382) -> 少数高额理赔者拉高了平均值",
             "这很可能由吸烟者群体的高保险费造成",
             "对线性回归来说, 右偏分布可能导致高额区间预测不佳 -> 后续可考虑对数变换(log transform)"
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# Slide 7 - 5. 预测值 vs 实际值对比
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "预测值 vs 实际值对比", "268个测试集样本 (占全部1338的20%)")

headers = ["索引", "实际值 (actual)", "预测值 (pred)", "误差", "判定"]
rows = [
    ["12", "$1,827", "$4,765", "+$2,938", "过高预测(2.6倍)"],
    ["306", "$20,178", "$4,958", "-$15,220", "严重过低预测"],
    ["318", "$7,421", "$8,299", "+$878", "相对准确"],
    ["815", "$1,878", "$3,079", "+$1,201", "过高预测"],
    ["157", "$15,518", "$24,166", "+$8,648", "过高预测"],
    ["1015", "$12,125", "$11,638", "-$487", "非常准确"],
]
add_table_slide(s, headers, rows, Inches(0.6), Inches(2.1), [1.2, 2.2, 2.2, 2.0, 2.8],
                row_height=0.45, font_size=13)

add_card(s, Inches(0.6), Inches(5.2), Inches(5.8), Inches(2.0),
         "索引306号 — 严重过低预测", [
             "实际 $20,178 却预测 $4,958",
             "此人可能是吸烟者, 但模型未能正确捕捉",
             "误差高达 -$15,220"
         ], ACCENT_RED, ACCENT_RED)

add_card(s, Inches(7.0), Inches(5.2), Inches(5.8), Inches(2.0),
         "索引1015号 — 非常准确", [
             "实际 $12,125 vs 预测 $11,638",
             "误差仅 $487 (约4%)",
             "展示了模型在中间区间的良好性能"
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# Slide 8 - 6. 散点图可视化解析
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "散点图可视化解析", "x轴 = 实际值, y轴 = 预测值 | 完美模型: 所有点在45度对角线上")

# 3个区域
add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(4.5),
         "3个区域分析", [
             "低额区间 ($0~$15,000):",
             "  -> 点聚集在对角线附近, 预测良好",
             "",
             "中额区间 ($15,000~$35,000):",
             "  -> 点开始偏离对角线",
             "  -> 过高/过低预测混合出现",
             "",
             "高额区间 ($35,000以上):",
             "  -> 点大幅分散, 预测不稳定",
         ], ACCENT_CYAN, ACCENT_BLUE)

# 3个聚类
add_card(s, Inches(4.8), Inches(2.0), Inches(3.8), Inches(4.5),
         "3个聚类现象", [
             "聚类1: 非吸烟者 (低额)",
             "  -> 集中在 $0~$15,000",
             "",
             "聚类2: 非吸烟者 (中额)",
             "  -> 年龄/BMI较高的非吸烟者",
             "",
             "聚类3: 吸烟者 (高额)",
             "  -> $20,000以上, 分散度大",
         ], ACCENT_GREEN, ACCENT_GREEN)

# 结论
add_card(s, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.5),
         "关键结论", [
             "吸烟者/非吸烟者群体之间",
             "保险费差异极大",
             "",
             "线性回归难以完美捕捉",
             "这种群体间的非线性差异",
             "",
             "-> 需要更复杂的模型",
             "   (决策树、随机森林等)",
             "   或添加交互项"
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# Slide 9 - 7. RMSE 解析
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "RMSE (Root Mean Squared Error) 解析", "模型预测误差的衡量指标")

# RMSE 值
add_shape(s, Inches(0.6), Inches(2.1), Inches(4.0), Inches(1.2), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.6), Inches(2.2), Inches(4.0), Inches(0.5),
         "RMSE = $5,684.93", font_size=28, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.8), Inches(4.0), Inches(0.4),
         "预测值与实际值平均偏差约 $5,685", font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# RMSE/均值
add_shape(s, Inches(5.0), Inches(2.1), Inches(3.5), Inches(1.2), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(5.0), Inches(2.2), Inches(3.5), Inches(0.5),
         "RMSE/均值 = 42.8%", font_size=24, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(5.0), Inches(2.8), Inches(3.5), Inches(0.4),
         "5,685 / 13,270 = 42.8%", font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 公式
add_card(s, Inches(9.0), Inches(2.1), Inches(3.8), Inches(1.2),
         "RMSE 公式", [
             "RMSE = sqrt(mean((yi - yi_hat)^2))",
         ], ACCENT_CYAN, ACCENT_BLUE)

# 计算过程
add_card(s, Inches(0.6), Inches(3.8), Inches(5.5), Inches(3.2),
         "RMSE 计算步骤", [
             "1. 对每个数据计算 (实际值 - 预测值) 的差",
             "2. 对差值进行平方 (消除符号 + 惩罚大误差)",
             "3. 计算平方值的平均 (MSE)",
             "4. 对平均值取平方根 (恢复到原始尺度)",
             "",
             "RMSE对大误差更敏感 (因为平方运算)",
             "相比MAE, 更严厉地惩罚离群值"
         ], ACCENT_CYAN, ACCENT_BLUE)

# 相对比较
add_card(s, Inches(6.5), Inches(3.8), Inches(6.3), Inches(3.2),
         "相对比较说明", [
             "RMSE没有绝对的好/坏标准",
             "主要用于与其他模型的相对比较",
             "",
             "后续目标: 通过以下方法降低RMSE",
             "  -> Ridge / Lasso 回归",
             "  -> Random Forest",
             "  -> XGBoost",
             "  -> 特征工程 (交互项、对数变换等)"
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# Slide 10 - 8. R² 决定系数
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "R² 决定系数 (Coefficient of Determination)", "模型解释力的衡量指标")

# R² 值
add_shape(s, Inches(0.6), Inches(2.1), Inches(5.0), Inches(1.2), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.6), Inches(2.15), Inches(5.0), Inches(0.6),
         "R² = 0.7368 (约73.7%)", font_size=28, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.8), Inches(5.0), Inches(0.4),
         "5个自变量解释了charges变动的约73.7%", font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# SST/SSR/SSE
add_card(s, Inches(6.0), Inches(2.1), Inches(6.8), Inches(1.2),
         "R² = SSR/SST = 1 - SSE/SST", [
             "SST(总变动) | SSR(模型解释的变动) | SSE(误差变动)"
         ], ACCENT_CYAN, ACCENT_BLUE)

# R²范围表
headers_r2 = ["R² 范围", "评价"]
rows_r2 = [
    ["0.9 以上", "非常优秀的模型"],
    ["0.7 ~ 0.9", "良好的模型  <-- 我们的模型"],
    ["0.5 ~ 0.7", "一般"],
    ["0.5 以下", "解释力不足"],
]
add_table_slide(s, headers_r2, rows_r2, Inches(0.6), Inches(3.8), [3.0, 5.0],
                row_height=0.5, font_size=14)

# 剩余解释
add_card(s, Inches(9.0), Inches(3.8), Inches(3.8), Inches(3.2),
         "剩余26.3%的来源", [
             "当前模型未包含的因素:",
             "  -> 基础疾病有无",
             "  -> 居住地区",
             "  -> 职业、运动习惯等",
             "",
             "或变量间的非线性关系:",
             "  -> 吸烟者的BMI效应",
             "     可能与非吸烟者不同"
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# Slide 11 - 9. 回归系数解析
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "回归系数 (Coefficients) 解析", "各自变量对保险费的影响大小")

# 系数表
headers_coef = ["变量", "系数", "解析"]
rows_coef = [
    ["smoker", "+23,469.28", "影响最大: 吸烟者保险费高约$23,469"],
    ["children", "+469.34", "每增加1个子女, 增加约$469"],
    ["bmi", "+297.51", "BMI每增加1, 增加约$298"],
    ["age", "+264.80", "年龄每增加1岁, 增加约$265"],
    ["sex", "+17.34", "男性比女性高约$17 (影响微乎其微)"],
]
add_table_slide(s, headers_coef, rows_coef, Inches(0.6), Inches(2.1), [2.0, 2.2, 6.2],
                row_height=0.5, font_size=13)

# 4个核心发现
add_card(s, Inches(0.6), Inches(5.0), Inches(6.0), Inches(2.2),
         "核心发现", [
             "1. 吸烟(smoker)是压倒性首要因素",
             "   系数23,469, 比其他变量大数十至数百倍",
             "2. 性别(sex)几乎无意义",
             "   系数仅17, 影响可忽略",
         ], ACCENT_CYAN, ACCENT_BLUE)

add_card(s, Inches(7.0), Inches(5.0), Inches(5.8), Inches(2.2),
         "补充发现", [
             "3. BMI和年龄影响相近",
             "   分别为298和265, 程度相似",
             "4. 子女数也有显著影响",
             "   每增加1个子女约增加$469",
             "   反映根据受抚养人数确定保险费的机制"
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# Slide 12 - 10. 截距解析
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "截距 (Intercept) 解析", "回归方程的常数项")

add_shape(s, Inches(3.0), Inches(2.2), Inches(7.0), Inches(1.2), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(3.0), Inches(2.3), Inches(7.0), Inches(0.6),
         "Intercept = -$11,577", font_size=34, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(3.0), Inches(2.9), Inches(7.0), Inches(0.4),
         "所有自变量为0时的charges预测值", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_card(s, Inches(0.6), Inches(4.0), Inches(5.8), Inches(3.0),
         "截距的含义", [
             "数学定义: 所有自变量 = 0 时的 charges 值",
             "",
             "实际意义有限:",
             "  -> 不存在'0岁, 女性, BMI=0, 无子女, 非吸烟者'",
             "  -> 应理解为数学上的校正值",
             "  -> 而非现实中的实际情况"
         ], ACCENT_CYAN, ACCENT_BLUE)

add_card(s, Inches(7.0), Inches(4.0), Inches(5.8), Inches(3.0),
         "截距为负的原因", [
             "其他系数之和为较大的正数",
             "  -> age(264.8) + bmi(297.5) + ...",
             "",
             "截距需要为负数才能使整体预测值",
             "落在合理范围内",
             "",
             "这是回归方程自动调整的结果,",
             "属于正常现象"
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# Slide 13 - 11. 最终预测公式
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "11", "最终预测公式", "回归方程与预测示例")

# 公式
add_shape(s, Inches(0.6), Inches(2.1), Inches(12.2), Inches(0.9), CODE_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.8), Inches(2.2), Inches(11.8), Inches(0.7),
         "charges = 264.80 x age + 17.34 x sex + 297.51 x bmi + 469.34 x children + 23469.28 x smoker - 11577.00",
         font_size=16, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER, font_name='Consolas')

# 非吸烟者示例
add_card(s, Inches(0.6), Inches(3.3), Inches(5.8), Inches(3.8),
         "示例: 30岁, 男性, BMI 25, 1子女, 非吸烟者", [
             "= 264.80 x 30 + 17.34 x 1 + 297.51 x 25",
             "  + 469.34 x 1 + 23469.28 x 0 - 11577.00",
             "",
             "= 7944.00 + 17.34 + 7437.75",
             "  + 469.34 + 0 - 11577.00",
             "",
             "= $4,291.43",
         ], ACCENT_CYAN, ACCENT_BLUE)

# 吸烟者示例
add_card(s, Inches(7.0), Inches(3.3), Inches(5.8), Inches(3.8),
         "同等条件, 吸烟者", [
             "= $4,291.43 + $23,469.28",
             "",
             "= $27,760.71",
             "",
             "仅凭吸烟与否一项,",
             "保险费就相差约 6.5 倍!",
             "",
             "非吸烟: $4,291  vs  吸烟: $27,761"
         ], ACCENT_RED, ACCENT_RED)

# ============================================================
# Slide 14 - 12. 综合评价
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "12", "综合评价", "模型性能总结、优势与局限、改进方向")

# 性能指标
add_shape(s, Inches(0.6), Inches(2.0), Inches(3.0), Inches(0.8), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.6), Inches(2.05), Inches(3.0), Inches(0.35),
         "RMSE = $5,684.93", font_size=18, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.4), Inches(3.0), Inches(0.35),
         "约43%误差", font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_shape(s, Inches(4.0), Inches(2.0), Inches(3.0), Inches(0.8), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(4.0), Inches(2.05), Inches(3.0), Inches(0.35),
         "R² = 0.7368", font_size=18, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.0), Inches(2.4), Inches(3.0), Inches(0.35),
         "73.7%解释力 (良好)", font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 优势
add_card(s, Inches(0.6), Inches(3.2), Inches(3.8), Inches(4.0),
         "优势 (3点)", [
             "1. 解释性非常直观",
             "   可定量把握各变量对保险费的影响",
             "",
             "2. 训练速度非常快",
             "   适合作为基线模型(Baseline)",
             "",
             "3. 明确捕捉到核心因素",
             "   吸烟是保险费的首要决定因素"
         ], ACCENT_GREEN, ACCENT_GREEN)

# 局限
add_card(s, Inches(4.7), Inches(3.2), Inches(3.8), Inches(4.0),
         "局限 (3点)", [
             "1. 非线性关系捕捉不足",
             "   吸烟者内部的精细预测不佳",
             "",
             "2. 高额区间预测误差大",
             "   charges右偏分布的影响",
             "",
             "3. 未考虑交互效应",
             "   例: 吸烟者的BMI效应",
             "   可能与非吸烟者不同"
         ], ACCENT_RED, ACCENT_RED)

# 改进方向
add_card(s, Inches(8.8), Inches(3.2), Inches(4.0), Inches(4.0),
         "改进方向 (4点)", [
             "1. 对数变换",
             "   对charges取log使分布正态化",
             "",
             "2. 添加交互项",
             "   smoker x bmi 等交互作用项",
             "",
             "3. 非线性模型",
             "   决策树、随机森林、XGBoost",
             "",
             "4. 正则化技术",
             "   Ridge, Lasso, Elastic Net"
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# ============================================================
# Save
# ============================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "04_Linear Regression_结果解析.pptx")
prs.save(out_path)
print("PPT saved successfully!")
print(f"Total slides: {len(prs.slides)}")
