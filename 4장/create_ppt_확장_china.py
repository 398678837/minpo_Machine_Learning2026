"""第4章 线性回归 - 扩展讲义PPT生成脚本（详细版）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色调色板 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    total_w = sum(col_widths)
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(1), Inches(2.5), Inches(11.333), Inches(3), SECTION_BG, accent)
    add_accent_line(s, Inches(5.5), Inches(2.6), Inches(2.333), accent)
    add_text(s, Inches(1), Inches(2.8), Inches(11.333), Inches(0.5),
             f"SECTION {section_num}", font_size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.4), Inches(11.333), Inches(1),
             title, font_size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.4), Inches(11.333), Inches(0.6),
             subtitle, font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片 1: 封面
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG, ACCENT_BLUE)
add_accent_line(s, Inches(4), Inches(1.5), Inches(5.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.8), Inches(11.333), Inches(0.5),
         "CHAPTER 4", font_size=22, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.5), Inches(11.333), Inches(1.2),
         "线性回归", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.333), Inches(0.5),
         "Linear Regression: OLS, Ridge, Lasso, Elastic Net", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.8),
         "从OLS到正则化技术的系统理解 + 亲手实现算法内部运作机制",
         font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
         "机器学习 | Machine Learning", font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片 2: 目录
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, None, "目录 (Table of Contents)", "第4章 整体结构")
toc_left = [
    "4.1  历史与意义",
    "4.2  OLS理论（正规方程、几何解释）",
    "4.3  多重共线性（VIF、条件数）",
    "4.4  Ridge回归（L2正则化）",
    "4.5  Lasso回归（L1正则化）",
    "4.6  Elastic Net（L1+L2）",
    "4.7  坐标下降法（glmnet）",
]
toc_right = [
    "4.8   SCAD与Oracle特性",
    "4.9   回归诊断",
    "4.10  实践：Ridge回归实现",
    "4.11  实践：Lasso回归实现",
    "4.12  实践：Elastic Net实现",
    "4.13  应用案例",
    "4.14  核心总结 + 复习题",
]
add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
         "Part I: 理论", toc_left, ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
         "Part II: 实践 & 应用", toc_right, ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# 幻灯片 3: 学习目标
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, None, "学习目标")
objectives = [
    "1. 能够通过矩阵微分推导OLS的正规方程并进行几何解释",
    "2. 能够说明高斯-马尔可夫定理的意义与局限性",
    "3. 能够比较Ridge、Lasso、Elastic Net的数学定义与差异",
    "4. 能够解释L1与L2正则化的几何差异",
    "5. 能够亲手实现坐标下降法和软阈值算子",
    "6. 能够说明SCAD的Oracle性质及相对于Lasso的优势",
    "7. 能够执行回归诊断（残差分析、Cook's Distance）",
]
add_bullet_list(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(5),
                objectives, font_size=18, color=LIGHT_GRAY, spacing=Pt(12))

# ============================================================
# SECTION 1: 历史与意义
# ============================================================
section_divider("历史与意义", "拥有200年历史的最基本ML算法", 1, ACCENT_BLUE)

# 幻灯片 5: 最小二乘法的诞生
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 1, "最小二乘法的诞生", "Legendre (1805) vs Gauss (1809)")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "勒让德 (Legendre, 1805)", [
             "首次正式发表最小二乘法",
             "根据观测数据确定彗星轨道",
             "提出最小化残差平方和的原理",
             "Nouvelles méthodes... (1805)",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "高斯 (Gauss, 1809)", [
             "声称自1795年起就开始使用",
             "Theoria Motus Corporum Coelestium (1809)",
             "证明在正态分布假设下与MLE等价",
             "→ 赋予概率论上的合理性",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0),
         "线性回归的现代地位", [
             "基础算法：监督学习中预测连续值的最基本模型",
             "理论基础：代价函数优化、偏差-方差权衡、正则化的出发点",
             "扩展基础：逻辑回归、神经网络、核方法的构成要素",
             "基准线：新模型性能比较的基准（baseline）",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 幻灯片 6: 数学定义
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 1, "线性回归的数学定义", "自变量与因变量之间的线性关系")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(2.5),
         "标量表示", [
             "y = β₀ + β₁x₁ + β₂x₂ + ··· + βₚxₚ + ε",
             "",
             "y: 因变量（目标）",
             "x₁,...,xₚ: 自变量（特征）",
             "β₀: 截距, β₁,...,βₚ: 回归系数",
             "ε: 误差项",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(2.5),
         "矩阵表示", [
             "y = Xβ + ε",
             "",
             "y ∈ R^n: 响应向量（n个观测值）",
             "X ∈ R^(n×(p+1)): 设计矩阵",
             "β ∈ R^(p+1): 回归系数向量",
             "ε ∈ R^n: 误差, εᵢ ~iid~ N(0, σ²)",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_text(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.8),
         "核心：'线性'是指对参数β为线性。y = β₁x + β₂x²也是线性回归（将x²视为新特征）",
         font_size=16, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 2: OLS理论
# ============================================================
section_divider("OLS理论", "正规方程推导、几何解释、高斯-马尔可夫定理", 2, ACCENT_GREEN)

# 幻灯片 8: 正规方程推导
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "正规方程推导（矩阵微分）", "求使RSS最小化的β")
add_card(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(5.0),
         "逐步推导", [
             "Step 1. 残差平方和（RSS）定义",
             "  RSS(β) = ||y - Xβ||² = (y - Xβ)ᵀ(y - Xβ)",
             "",
             "Step 2. 展开",
             "  RSS(β) = yᵀy - 2βᵀXᵀy + βᵀXᵀXβ",
             "",
             "Step 3. 对β求偏导（应用矩阵微分法则）",
             "  ∂RSS/∂β = -2Xᵀy + 2XᵀXβ",
             "  使用的公式: ∂(aᵀx)/∂x = a,  ∂(xᵀAx)/∂x = 2Ax (A对称)",
             "",
             "Step 4. 令其为0 → 正规方程 (Normal Equation)",
             "  XᵀXβ = Xᵀy",
             "",
             "Step 5. 解（当XᵀX可逆时）",
             "  β̂_OLS = (XᵀX)⁻¹Xᵀy      ← 核心公式",
         ], ACCENT_GREEN, ACCENT_GREEN)

# 幻灯片 9: 海森矩阵与凸性
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "二阶导数（海森矩阵）与凸性验证", "保证正规方程的解为最小值")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.5),
         "海森矩阵（二阶导数）", [
             "∂²RSS / ∂β∂βᵀ = 2XᵀX",
             "",
             "XᵀX为半正定矩阵(positive semi-definite)",
             "  → RSS是凸(convex)函数",
             "  → 正规方程的解为全局最小值",
             "",
             "若XᵀX为正定矩阵(positive definite)",
             "  → 唯一最小值",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.5),
         "凸优化的优势", [
             "1. 局部最小值 = 全局最小值",
             "  → 从任何起点出发都能得到相同的解",
             "",
             "2. 梯度下降法保证收敛",
             "  → 只需适当的学习率",
             "",
             "3. Ridge也是凸的（加λI仍保持凸性）",
             "  → Lasso也是凸的（|β|是凸的）",
             "",
             "4. SCAD是非凸的 → 存在局部最小值问题",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 幻灯片 10: 几何解释
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "OLS的几何解释", "正交投影 (Orthogonal Projection)")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5),
         "正交投影视角", [
             "y是R^n空间中的向量",
             "C(X) = X的列空间 (column space)",
             "",
             "ŷ = X(XᵀX)⁻¹Xᵀy = Hy",
             "",
             "H = X(XᵀX)⁻¹Xᵀ",
             "  = 帽子矩阵 (hat matrix)",
             "  = 投影矩阵 (projection matrix)",
             "",
             "ŷ是y在C(X)上的正交投影",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.0),
         "核心性质", [
             "残差 Xᵀe = 0（残差与列空间正交）",
             "勾股定理: ||y||² = ||ŷ||² + ||e||²",
             "  → TSS = ESS + RSS",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(6.7), Inches(4.3), Inches(6.1), Inches(2.2),
         "R²的几何意义", [
             "R² = cos²θ  (y与ŷ之间的夹角)",
             "",
             "R² = 1: y在C(X)内（完美拟合）",
             "R² = 0: y与C(X)正交（无解释力）",
             "0 < R² < 1: 部分解释",
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# 幻灯片 11: 高斯-马尔可夫定理
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "高斯-马尔可夫定理", "OLS估计量的最优性保证")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(3.0),
         "定理的条件（4个）", [
             "1. E[εᵢ] = 0（误差期望为0）",
             "2. Var(εᵢ) = σ²（同方差性）",
             "3. Cov(εᵢ, εⱼ) = 0, i≠j（不相关）",
             "4. X是固定的（non-random）矩阵",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(3.0),
         "BLUE: 最优性的含义", [
             "Best: 在所有线性无偏估计量中方差最小",
             "Linear: β̂是y的线性函数",
             "Unbiased: E[β̂] = β（无偏）",
             "Estimator: β的估计量",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
         "核心局限与Ridge的动机", [
             "高斯-马尔可夫仅保证OLS在无偏估计量中是最优的！",
             "若允许偏差，从MSE = Var + Bias²的角度，存在优于OLS的估计量",
             "→ Hoerl & Kennard (1970): 以少量偏差大幅降低方差，从而减小总MSE",
         ], ACCENT_RED, ACCENT_RED)

# ============================================================
# SECTION 3: 多重共线性
# ============================================================
section_divider("多重共线性", "VIF、条件数、解决方法", 3, ACCENT_ORANGE)

# 幻灯片 13: 多重共线性
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 3, "多重共线性 (Multicollinearity)", "自变量之间高度线性相关的问题")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(2.5),
         "问题", [
             "当XᵀX接近奇异(singular)时：",
             "(XᵀX)⁻¹的元素变得非常大",
             "Var(β̂) = σ²(XᵀX)⁻¹急剧增大",
             "回归系数估计不稳定（高方差）",
         ], ACCENT_RED, ACCENT_RED)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(2.5),
         "VIF（方差膨胀因子）", [
             "VIFⱼ = 1 / (1 - Rⱼ²)",
             "",
             "VIF = 1: 无多重共线性",
             "VIF 1~5: 轻微（一般可接受）",
             "VIF 5~10: 中等（需注意）",
             "VIF > 10: 严重（需处理）",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.2),
         "条件数 (Condition Number)", [
             "κ(XᵀX) = λ_max / λ_min",
             "",
             "条件数 > 30 → 怀疑存在多重共线性",
             "条件数越大，逆矩阵计算越不稳定",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_table_slide(s,
    ["解决方法", "说明"],
    [
        ["删除变量", "手动删除VIF较高的变量"],
        ["主成分回归", "PCA降维后进行回归"],
        ["Ridge回归", "L2正则化改善条件数"],
        ["Lasso回归", "L1正则化自动删除变量"],
    ],
    Inches(6.8), Inches(4.8), [2.5, 3.3],
    row_height=0.45, font_size=12)

# ============================================================
# SECTION 4: Ridge回归
# ============================================================
section_divider("Ridge回归", "L2正则化: Hoerl & Kennard (1970)", 4, ACCENT_PURPLE)

# 幻灯片 15: Ridge核心思想
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 4, "Ridge回归：核心思想", "偏差-方差权衡的经典案例")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
         "偏差-方差权衡", [
             "MSE(β̂) = Var(β̂) + [Bias(β̂)]²",
             "",
             "OLS: Bias = 0, Var = 可能很大",
             "Ridge: Bias > 0（略微）, Var = 大幅减小",
             "→ 总MSE减小！",
             "",
             "核心洞察 (Hoerl & Kennard, 1970):",
             "  利用高斯-马尔可夫的'漏洞'",
             "  在无偏估计量中是最优的，但是...",
             "  如果允许少量偏差",
             "  MSE可以更优！",
         ], ACCENT_PURPLE, ACCENT_PURPLE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(5.0),
         "L2正则化的数学定义", [
             "优化问题:",
             "  min ||y - Xβ||² + λ||β||²",
             "",
             "闭式解:",
             "  β̂_ridge = (XᵀX + λI)⁻¹Xᵀy",
             "",
             "推导:",
             "  ∂/∂β [...] = -2Xᵀy + 2XᵀXβ + 2λβ = 0",
             "  (XᵀX + λI)β = Xᵀy",
             "",
             "核心效果:",
             "  加λI使对角元素增大",
             "  → 条件数减小",
             "  → 数值稳定性提高",
             "  → 系数向0方向收缩",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 幻灯片 16: 特征值分解与贝叶斯解释
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 4, "Ridge: 特征值分解 & 贝叶斯解释", "从两个角度理解")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "通过特征值分解理解", [
             "在XᵀX = VDVᵀ中",
             "Ridge将每个方向的系数按 dⱼ/(dⱼ+λ) 收缩",
             "",
             "dⱼ大: 收缩比例 ≈ 1（几乎不变）",
             "dⱼ小: 收缩比例 ≈ 0（强烈收缩）",
             "→ 选择性地收缩不稳定方向",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "贝叶斯解释", [
             "先验分布: β ~ N(0, τ²I)",
             "似然: y|X,β ~ N(Xβ, σ²I)",
             "",
             "MAP估计 = Ridge (λ = σ²/τ²)",
             "→ Ridge = 对系数赋予正态先验分布",
             "τ²越小 → λ越大 → 收缩越多",
         ], ACCENT_PURPLE, ACCENT_PURPLE)
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0),
         "岭迹 (Ridge Trace)", [
             "x轴: λ, y轴: 各回归系数值的图形",
             "从λ = 0 (OLS)开始 → λ增大 → 所有系数收敛于0",
             "选择系数趋于稳定的λ值（通过交叉验证精确化）",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# SECTION 5: Lasso回归
# ============================================================
section_divider("Lasso回归", "L1正则化: Tibshirani (1996)", 5, ACCENT_RED)

# 幻灯片 18: Lasso核心
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "Lasso: L1正则化与变量选择", "收缩(shrinkage) + 变量选择同时实现！")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
         "L1正则化的数学定义", [
             "min (1/2n)||y - Xβ||² + λ||β||₁",
             "",
             "||β||₁ = Σ|βⱼ| (L1范数)",
             "",
             "无闭式解（在原点处不可微）",
             "→ 需要坐标下降法(Coordinate Descent)",
         ], ACCENT_RED, ACCENT_RED)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.0),
         "Lasso的创新 (Tibshirani, 1996)", [
             "将Ridge的L2替换为L1这一简单变更",
             "产生了根本性不同的性质：",
             "",
             "1. 回归系数收缩 (shrinkage)",
             "2. 变量选择 (variable selection) ← 核心！",
             "3. 将部分系数精确置为0",
             "",
             "Google Scholar引用 50,000+ (截至2024年)",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 幻灯片 19: 软阈值算子
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "软阈值算子 (Soft Thresholding)", "Lasso解的核心机制")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(4.8),
         "公式及运作", [
             "S(z, γ) = sign(z) · max(|z| - γ, 0)",
             "",
             "三个区间的运作：",
             "",
             "|z| ≤ γ:   结果 = 0",
             "  → 将小系数精确置为0（变量删除！）",
             "",
             "z > γ:    结果 = z - γ",
             "  → 将正系数收缩γ",
             "",
             "z < -γ:   结果 = z + γ",
             "  → 将负系数收缩γ",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(2.3),
         "Ridge vs Lasso 收缩方式比较", [
             "Ridge: 比例收缩  β̂ = dⱼ/(dⱼ+λ) · β̂_OLS",
             "  → 接近0但永远不为0",
             "Lasso: 平移收缩  β̂ = sign(β̂_OLS)·max(|β̂|-λ,0)",
             "  → 小于γ的系数精确为0！",
         ], ACCENT_CYAN, ACCENT_CYAN)
add_code_block(s, Inches(6.8), Inches(4.6), Inches(6), Inches(2.2), [
    "# Python实现（1行！）",
    "def soft_threshold(z, gamma):",
    "    return np.sign(z) * np.maximum(",
    "        np.abs(z) - gamma, 0.0)",
], font_size=13)

# 幻灯片 20: L1 vs L2 几何
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "L1 vs L2 几何比较", "为什么Lasso能进行变量选择？")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
         "L2约束 (Ridge) = 圆(Circle)", [
             "约束: β₁² + β₂² ≤ t",
             "",
             "RSS等高线（椭圆）与圆的切点：",
             "  → 通常不在坐标轴上",
             "  → 系数难以精确为0",
             "",
             "结果: 所有系数非零",
             "  → 无法进行变量选择",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(5.0),
         "L1约束 (Lasso) = 菱形(Diamond)", [
             "约束: |β₁| + |β₂| ≤ t",
             "",
             "菱形的顶点位于坐标轴上！",
             "  → RSS等高线的切点",
             "    很可能在顶点处",
             "  → 顶点: 一个或多个坐标 = 0",
             "",
             "结果: 部分系数精确为0",
             "  → 变量被自动选择！",
         ], ACCENT_RED, ACCENT_RED)

# 幻灯片 21: 贝叶斯解释比较
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "贝叶斯解释: Ridge vs Lasso", "先验分布的差异决定变量选择")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
         "Ridge = 正态先验分布", [
             "p(β) = (1/√2πτ²) exp(-β²/2τ²)",
             "",
             "钟形曲线: 在0附近平滑衰减",
             "→ 将系数向0收缩",
             "   但无法精确置为0",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.0),
         "Lasso = 拉普拉斯先验分布", [
             "p(β) = (λ/2) exp(-λ|β|)",
             "",
             "在原点处尖锐 + 厚尾",
             "→ 对小系数赋予高先验概率（稀疏性）",
             "→ 对大系数也给予适当概率（限制偏差）",
         ], ACCENT_RED, ACCENT_RED)
add_text(s, Inches(0.8), Inches(5.5), Inches(11), Inches(1.0),
         "核心: 拉普拉斯分布在原点处更尖锐，因此将小系数置为0的先验信念更强",
         font_size=16, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 6: Elastic Net
# ============================================================
section_divider("Elastic Net", "L1 + L2混合正则化: Zou & Hastie (2005)", 6, ACCENT_CYAN)

# 幻灯片 23: Elastic Net
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 6, "Elastic Net: 解决Lasso的两个局限", "Zou & Hastie (2005)")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "Lasso的局限", [
             "1. p > n问题: 当变量数 > 样本数时",
             "   Lasso最多只能选择n个变量",
             "2. 缺乏分组效应: 在相关变量中",
             "   只任意选择一个，删除其余",
         ], ACCENT_RED, ACCENT_RED)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "Elastic Net的解决方案", [
             "min (1/2n)||y-Xβ||² + λ[α||β||₁ + (1-α)/2 ||β||₂²]",
             "",
             "α = 1: 纯Lasso",
             "α = 0: 纯Ridge",
             "0 < α < 1: Elastic Net（最优折中）",
         ], ACCENT_CYAN, ACCENT_CYAN)
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         "分组效应定理 (Zou & Hastie, 2005)", [
             "|β̂ᵢ - β̂ⱼ| ≤ [1 / λ(1-α)] · √(2(1-r)) · ||y||",
             "→ 相关性越高（r→1），两个系数的差异越小",
             "→ Elastic Net遵循'相关变量一起选择或一起删除'的原则",
             "→ 纯Lasso (α=1)时分母=0，此性质不成立",
         ], ACCENT_GREEN, ACCENT_GREEN)

# 幻灯片 24: 三种方法综合比较
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 6, "Ridge vs Lasso vs Elastic Net 综合比较")
add_table_slide(s,
    ["特性", "Ridge (L2)", "Lasso (L1)", "Elastic Net (L1+L2)"],
    [
        ["惩罚项", "λ||β||₂²", "λ||β||₁", "λ[α||β||₁+(1-α)||β||₂²/2]"],
        ["变量选择", "X（所有系数非零）", "O（部分=0）", "O（部分=0）"],
        ["分组效应", "O（同时收缩）", "X（只选一个）", "O（同时选择）"],
        ["闭式解", "O", "X", "X"],
        ["p>n支持", "O", "最多n个", "O"],
        ["贝叶斯先验", "正态分布", "拉普拉斯", "正态+拉普拉斯"],
        ["约束区域", "圆", "菱形", "圆角菱形"],
    ],
    Inches(0.5), Inches(2.0), [2.5, 2.5, 2.5, 3.5],
    row_height=0.5, font_size=12, header_font_size=13)

# ============================================================
# SECTION 7: 坐标下降法
# ============================================================
section_divider("坐标下降法", "glmnet算法: Friedman, Hastie, Tibshirani (2010)", 7, ACCENT_GREEN)

# 幻灯片 26: 坐标下降法
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 7, "坐标下降法 (Coordinate Descent)", "每次只优化一个变量")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(5.0),
         "更新规则", [
             "第j个坐标更新：",
             "",
             "β̃ⱼ = S(ρⱼ, αλ) / (Xⱼᵀ Xⱼ/n + λ(1-α))",
             "",
             "ρⱼ = Xⱼᵀ rⱼ / n",
             "rⱼ = y - Σ_{k≠j} Xₖβₖ  (部分残差)",
             "S(z,γ) = sign(z)·max(|z|-γ,0)",
             "",
             "算法：",
             "  for iteration in range(max_iter):",
             "    for j in range(n_features):",
             "      计算部分残差 → 应用软阈值",
             "    检查收敛 (max|β_new - β_old| < tol)",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(5.0),
         "3个核心优化技巧", [
             "1. 热启动 (Warm Start)",
             "  按λ₁ > λ₂ > ...顺序求解时",
             "  使用上一个解作为初始值",
             "  相邻λ的解相似 → 快速收敛",
             "",
             "2. 活跃集 (Active Set)",
             "  优先更新非零系数",
             "  Lasso中大部分系数=0，因此非常高效",
             "",
             "3. λ_max计算",
             "  λ_max = (1/αn) max_j |Xⱼᵀ(y-ȳ)|",
             "  使所有系数=0的最小λ",
             "  从此处开始逐渐减小λ",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# SECTION 8: SCAD与Oracle特性
# ============================================================
section_divider("SCAD与Oracle特性", "Fan & Li (2001): 非凸惩罚与理想性质", 8, ACCENT_PURPLE)

# 幻灯片 28: SCAD
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 8, "SCAD: 具有Oracle性质的惩罚", "Fan & Li (2001)")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "Lasso的根本局限", [
             "L1惩罚无论系数大小",
             "都收缩相同的量(λ)：",
             "→ 小系数: 适当地置为0（OK）",
             "→ 大系数: 不必要的过度收缩（问题！）",
         ], ACCENT_RED, ACCENT_RED)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "Oracle性质 (Oracle Properties)", [
             "1. 一致的变量选择：",
             "   P(选中变量 = 真实变量) → 1",
             "2. 渐近正态性：",
             "   非零系数的估计与OLS同分布",
             "Lasso不满足Oracle性质！",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_table_slide(s,
    ["区间", "惩罚行为", "效果"],
    [
        ["|θ| ≤ λ", "与L1相同", "将小系数置为0"],
        ["λ < |θ| ≤ aλ", "惩罚逐渐减小", "中等系数逐步收缩"],
        ["|θ| > aλ", "惩罚导数 = 0", "对大系数不再收缩！"],
    ],
    Inches(0.5), Inches(4.8), [2.5, 3.5, 5.5],
    row_height=0.55, font_size=13)

# 幻灯片 29: 正则化方法的谱系
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 8, "正则化方法的谱系", "好的惩罚的3个条件")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "好的惩罚的3个条件 (Fan & Li)", [
             "1. 无偏性: 对大系数无偏差",
             "2. 稀疏性: 将小系数置为0",
             "3. 连续性: 估计量是数据的连续函数",
             "同时满足三个条件 → 必定非凸！",
         ], ACCENT_PURPLE, ACCENT_PURPLE)
add_table_slide(s,
    ["方法", "年份", "惩罚", "变量选择", "Oracle", "凸"],
    [
        ["Ridge", "1970", "L2", "X", "X", "O"],
        ["Lasso", "1996", "L1", "O", "X", "O"],
        ["SCAD", "2001", "非凸", "O", "O", "X"],
        ["Elastic Net", "2005", "L1+L2", "O", "X", "O"],
        ["Adaptive Lasso", "2006", "加权L1", "O", "O", "O"],
        ["MCP", "2010", "非凸", "O", "O", "X"],
    ],
    Inches(0.5), Inches(4.8), [2.0, 1.2, 1.8, 1.5, 1.5, 1.5],
    row_height=0.45, font_size=12)

# ============================================================
# SECTION 9: 回归诊断
# ============================================================
section_divider("回归诊断", "残差分析、Q-Q Plot、Cook's Distance", 9, ACCENT_ORANGE)

# 幻灯片 31: 残差分析
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 9, "残差分析 (Residual Analysis)", "模型拟合度评估")
add_table_slide(s,
    ["残差类型", "公式", "用途"],
    [
        ["普通残差", "eᵢ = yᵢ - ŷᵢ", "基本分析"],
        ["标准化残差", "rᵢ = eᵢ / (σ̂√(1-hᵢᵢ))", "异常值检测"],
        ["学生化残差", "tᵢ = eᵢ / (σ̂₍ᵢ₎√(1-hᵢᵢ))", "精确异常值检测"],
    ],
    Inches(0.5), Inches(2.0), [2.5, 4.5, 4.5],
    row_height=0.55, font_size=13)
add_card(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(3.0),
         "残差图检查项目", [
             "1. 预测值 vs 残差: 应无模式（同方差性）",
             "2. 残差的正态性: 通过Q-Q plot确认",
             "3. 自相关: 时间序列数据中残差的独立性",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_card(s, Inches(6.7), Inches(4.0), Inches(6.1), Inches(3.0),
         "Q-Q Plot", [
             "x轴: 理论正态分布的分位数",
             "y轴: 残差的分位数",
             "点在45度直线上 → 服从正态分布",
             "偏离直线 → 判断分布偏态",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 幻灯片 32: Cook's Distance
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 9, "Cook's Distance", "检测有影响力的观测值")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(3.0),
         "Cook's Distance公式", [
             "Dᵢ = (rᵢ² / p) · (hᵢᵢ / (1-hᵢᵢ)²)",
             "",
             "rᵢ: 标准化残差",
             "hᵢᵢ: 杠杆值（帽子矩阵对角元素）",
             "p: 参数个数",
             "",
             "判定: Dᵢ > 1 → 高影响力观测值",
             "保守标准: Dᵢ > 4/n",
         ], ACCENT_RED, ACCENT_RED)
add_table_slide(s,
    ["诊断项目", "检查方法", "违反时的措施"],
    [
        ["线性", "预测值 vs 残差图", "变量变换、多项式回归"],
        ["同方差性", "预测值 vs 残差图", "WLS、变量变换"],
        ["正态性", "Q-Q plot, Shapiro-Wilk", "变量变换、稳健回归"],
        ["独立性", "Durbin-Watson检验", "GLS、时间序列模型"],
        ["多重共线性", "VIF、条件数", "删除变量、Ridge"],
        ["异常值", "Cook's distance", "删除或稳健回归"],
    ],
    Inches(6.8), Inches(2.0), [2.0, 2.0, 2.0],
    row_height=0.5, font_size=11, header_font_size=12)

# ============================================================
# SECTION 10: 实践 - Ridge实现
# ============================================================
section_divider("实践: Ridge回归实现", "直接编写闭式解代码", 10, ACCENT_BLUE)

# 幻灯片 34: Ridge直接实现
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 10, "Ridge回归直接实现", "01_ridge_regression.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(4.8), [
    "class RidgeRegressionFromScratch:",
    '    """闭式解: w = (X\'X + λI)⁻¹X\'y"""',
    "",
    "    def __init__(self, alpha=1.0):",
    "        self.alpha = alpha",
    "",
    "    def fit(self, X, y):",
    "        n, p = X.shape",
    "        self.X_mean_ = np.mean(X, axis=0)",
    "        self.y_mean_ = np.mean(y)",
    "        Xc = X - self.X_mean_",
    "        yc = y - self.y_mean_",
    "",
    "        # 核心: (X'X + αI)β = X'y",
    "        XtX = Xc.T @ Xc",
    "        reg = XtX + self.alpha * np.eye(p)",
    "        Xty = Xc.T @ yc",
    "        self.coef_ = np.linalg.solve(reg, Xty)",
    "        self.intercept_ = (",
    "            self.y_mean_ - self.X_mean_ @ self.coef_)",
    "        return self",
], font_size=11)
add_card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.8),
         "实现核心要点", [
             "1. 数据中心化 (centering)",
             "   Xc = X - mean(X)",
             "   yc = y - mean(y)",
             "   → 截距单独处理",
             "",
             "2. 使用np.linalg.solve",
             "   用求解线性系统代替逆矩阵(inv)",
             "   → 数值上更稳定",
             "",
             "3. 与sklearn结果一致性验证",
             "   系数最大差异: ~10⁻¹²级别",
             "",
             "4. 多重共线性效果验证",
             "   条件数: Ridge大幅降低",
             "   κ(XᵀX+λI) << κ(XᵀX)",
         ], ACCENT_BLUE, ACCENT_BLUE)

# 幻灯片 35: Ridge正则化路径
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 10, "正则化路径 (Ridge Trace) & 交叉验证", "选择最优λ")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(2.5), [
    "# 正则化路径可视化",
    "alphas = np.logspace(-2, 4, 200)",
    "coefs_path = []",
    "for alpha in alphas:",
    "    ridge = RidgeFromScratch(alpha=alpha)",
    "    ridge.fit(X_scaled, y)",
    "    coefs_path.append(ridge.coef_.copy())",
], font_size=12)
add_code_block(s, Inches(0.6), Inches(5.0), Inches(6), Inches(2.0), [
    "# 交叉验证选择最优λ",
    "ridge_cv = RidgeCV(alphas=alphas, cv=5,",
    "    scoring='neg_mean_squared_error')",
    "ridge_cv.fit(X_scaled, y)",
    "print(f'最优λ: {ridge_cv.alpha_:.4f}')",
], font_size=12)
add_card(s, Inches(6.5), Inches(2.2), Inches(6.3), Inches(4.8),
         "Ridge Trace解读", [
             "x轴: log₁₀(λ), y轴: 各回归系数",
             "",
             "λ = 0（左侧）: OLS估计值",
             "  → 因多重共线性而不稳定（大幅振荡）",
             "",
             "λ增大: 所有系数收敛于0",
             "  → 偏差增大，方差减小",
             "",
             "最优λ: CV MSE最小的点",
             "  → 偏差-方差权衡的最优点",
             "",
             "OLS vs Ridge比较:",
             "  OLS系数的L2范数 >> Ridge",
             "  Ridge大幅降低条件数",
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# SECTION 11: 实践 - Lasso实现
# ============================================================
section_divider("实践: Lasso回归实现", "直接编写坐标下降法代码", 11, ACCENT_RED)

# 幻灯片 37: Lasso直接实现
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 11, "Lasso坐标下降法直接实现", "02_lasso_regression.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(4.8), [
    "class LassoFromScratch:",
    "    def fit(self, X, y):",
    "        n, p = X.shape",
    "        Xc = X - np.mean(X, axis=0)",
    "        yc = y - np.mean(y)",
    "        col_norms = np.sum(Xc**2, axis=0)/n",
    "        w = np.zeros(p)",
    "        residual = yc.copy()",
    "",
    "        for _ in range(self.max_iter):",
    "            w_old = w.copy()",
    "            for j in range(p):",
    "                # 恢复第j个特征对部分残差的贡献",
    "                residual += Xc[:,j] * w[j]",
    "                # X_j与部分残差的内积/n",
    "                rho = Xc[:,j] @ residual / n",
    "                # 应用软阈值！",
    "                w[j] = soft_threshold(rho,",
    "                    self.alpha) / col_norms[j]",
    "                # 更新残差",
    "                residual -= Xc[:,j] * w[j]",
    "            if np.max(np.abs(w-w_old))<self.tol:",
    "                break",
], font_size=10)
add_card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.8),
         "坐标下降法核心理解", [
             "1. 部分残差 (partial residual)",
             "   rⱼ = y - Σ_{k≠j} Xₖwₖ",
             "   '减去其他变量贡献后的残差'",
             "",
             "2. 高效残差更新",
             "   residual += Xⱼ·w_old[j]  (恢复)",
             "   residual -= Xⱼ·w_new[j]  (更新)",
             "   → 每次只反映差异而非全部重新计算",
             "",
             "3. 变量选择验证",
             "   零系数个数 → 与实际对比",
             "   自行实现 vs sklearn结果一致",
             "",
             "4. 正则化路径",
             "   λ减小 → 选中变量数增加",
             "   红色实线: 实际非零系数个数",
         ], ACCENT_RED, ACCENT_RED)

# ============================================================
# SECTION 12: 实践 - Elastic Net实现
# ============================================================
section_divider("实践: Elastic Net实现", "L1+L2混合正则化 + 分组效应", 12, ACCENT_CYAN)

# 幻灯片 39: Elastic Net直接实现
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 12, "Elastic Net坐标下降法实现", "03_elastic_net.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(3.0), [
    "# Elastic Net更新规则",
    "# 与Lasso的区别: 分母加L2项！",
    "l1_penalty = alpha * l1_ratio",
    "l2_penalty = alpha * (1 - l1_ratio)",
    "",
    "for j in range(p):",
    "    residual += Xc[:,j] * w[j]",
    "    rho = Xc[:,j] @ residual / n",
    "    # 核心: 分母加l2_penalty",
    "    denom = col_norms[j] + l2_penalty",
    "    w[j] = soft_threshold(rho, l1_penalty) / denom",
    "    residual -= Xc[:,j] * w[j]",
], font_size=12)
add_card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.8),
         "分组效应实验结果", [
             "相关特征组数据：",
             "  组1 (x0,x1,x2): 高相关, β=2.0",
             "  组2 (x3,x4): 高相关, β=-1.5",
             "  噪声 (x5~x14): β=0",
             "",
             "Ridge: 所有系数非零，组内相似",
             "Lasso: 变量选择O，组内只选一个！",
             "Elastic Net: 变量选择 + 组内同时选择",
             "",
             "组内系数方差：",
             "  Ridge < Elastic Net << Lasso",
             "  → EN的分组效应最优",
         ], ACCENT_CYAN, ACCENT_CYAN)

# ============================================================
# SECTION 13: 应用案例
# ============================================================
section_divider("应用案例", "房价预测 & 广告费-销售额分析", 13, ACCENT_GREEN)

# 幻灯片 41: 房价预测
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 13, "应用1: California Housing房价预测", "sklearn内置数据集")
add_table_slide(s,
    ["模型", "RMSE", "R²", "非零系数数"],
    [
        ["OLS", "~0.73", "~0.58", "8"],
        ["Ridge (α=1)", "~0.73", "~0.58", "8"],
        ["Lasso (α=0.01)", "~0.74", "~0.57", "6~7"],
    ],
    Inches(0.5), Inches(2.0), [3.0, 2.0, 2.0, 2.5],
    row_height=0.55, font_size=13)
add_card(s, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.8),
         "主要系数解读", [
             "MedInc（中位收入）: 强正(+)",
             "HouseAge（房龄）: 正(+) - 市中心效应",
             "AveRooms（平均房间数）: 正(+) - 大房子",
             "Latitude（纬度）: 负(-) - 北部价格下降",
             "Longitude（经度）: 负(-) - 靠海侧较高",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(6.7), Inches(4.2), Inches(6.1), Inches(2.8),
         "局限性 & 经验教训", [
             "R² ≈ 0.58: 线性模型仅能解释58%",
             "难以捕捉纬度/经度的非线性地理模式",
             "→ 可用RF、XGBoost等非线性模型改善",
             "",
             "教训: 线性回归是可解释的baseline",
             "性能 vs 可解释性的权衡",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 幻灯片 42: 广告费-销售额分析
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 13, "应用2: 广告费-销售额回归分析", "市场营销分析案例")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
         "数据结构", [
             "电视广告费: 10~300（主效应）",
             "广播广告费: 0~50（主效应）",
             "报纸广告费: 0~100（弱效应）",
             "销售额 = 0.05·TV + 0.1·Radio + 0.005·News",
             "        + 0.001·TV·Radio + 5 + 噪声",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.0),
         "Elastic Net结果（变量选择）", [
             "电视广告: 选中（最大正效应）",
             "广播: 选中（单位效率最高）",
             "报纸: 几乎收缩为0 → 实质上被删除！",
             "",
             "营销启示：",
             "  将报纸预算重新分配给电视/广播更合理",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_text(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
         "核心: 线性回归的可解释性直接用于营销决策。每个系数代表'广告费增加1单位时销售额的变化'",
         font_size=15, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 14: 核心总结
# ============================================================
section_divider("核心总结 + 复习题", "第4章 整体回顾", 14, ACCENT_BLUE)

# 幻灯片 44: 理论总结表
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 14, "理论总结表", "核心公式与特性一览")
add_table_slide(s,
    ["概念", "核心公式", "核心特性"],
    [
        ["OLS", "β̂ = (XᵀX)⁻¹Xᵀy", "BLUE, 无偏"],
        ["Ridge", "β̂ = (XᵀX+λI)⁻¹Xᵀy", "L2, 收缩, 闭式解"],
        ["Lasso", "min ||y-Xβ||²/2n + λ||β||₁", "L1, 变量选择"],
        ["Elastic Net", "L1+L2混合", "分组效应"],
        ["SCAD", "非凸惩罚", "Oracle性质"],
    ],
    Inches(0.5), Inches(2.0), [2.0, 4.5, 4.5],
    row_height=0.55, font_size=13)

# 幻灯片 45: 核心论文年表
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 14, "核心论文年表 & 方法选择指南")
add_table_slide(s,
    ["年份", "论文", "核心贡献"],
    [
        ["1970", "Hoerl & Kennard", "Ridge, 偏差-方差权衡"],
        ["1996", "Tibshirani", "Lasso, L1变量选择"],
        ["2001", "Fan & Li", "SCAD, Oracle性质"],
        ["2005", "Zou & Hastie", "Elastic Net, 分组效应"],
        ["2010", "Friedman et al.", "glmnet, 坐标下降法"],
    ],
    Inches(0.5), Inches(2.0), [1.2, 3.0, 7.3],
    row_height=0.5, font_size=13)
add_card(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
         "方法选择指南", [
             "多重共线性严重？ → Ridge或Elastic Net",
             "需要变量选择？ → Lasso或Elastic Net",
             "存在相关变量组？ → Elastic Net（分组效应）",
             "p > n？ → Elastic Net",
             "变量少且无多重共线性？ → OLS",
         ], ACCENT_GREEN, ACCENT_GREEN)

# 幻灯片 46: 复习题 (1/2)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 14, "复习题 (1/2)", "公式推导 + 理论")
questions_1 = [
    "Q1. 对β进行矩阵微分RSS，推导正规方程",
    "Q2. 高斯-马尔可夫定理的局限：为什么Ridge可能优于OLS？",
    "Q3. 推导Ridge的闭式解。λI改善条件数的原因？",
    "Q4. 从L1（菱形）vs L2（圆）约束区域的几何角度解释变量选择的差异",
    "Q5. 软阈值算子S(z,γ)的三个区间运作及其在坐标下降法中的作用",
]
add_bullet_list(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                questions_1, font_size=17, color=LIGHT_GRAY, spacing=Pt(14))

# 幻灯片 47: 复习题 (2/2)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 14, "复习题 (2/2)", "深入 + 诊断")
questions_2 = [
    "Q6. 解释Elastic Net的分组效应定理",
    "Q7. 坐标下降法的效率：warm start与active set策略",
    "Q8. SCAD的Oracle性质两个条件及相对于Lasso的优势",
    "Q9. Ridge（正态分布）与Lasso（拉普拉斯）的贝叶斯解释",
    "Q10. 残差分析: (a) 漏斗形模式 (b) Q-Q尾部偏离 (c) Cook's D > 1",
]
add_bullet_list(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                questions_2, font_size=17, color=LIGHT_GRAY, spacing=Pt(14))

# ============================================================
# 幻灯片 48: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG, ACCENT_BLUE)
add_accent_line(s, Inches(4), Inches(2.5), Inches(5.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.8), Inches(11.333), Inches(1),
         "Thank You", font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.0), Inches(11.333), Inches(0.5),
         "第4章: 线性回归 (Linear Regression)", font_size=22, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
         "下一章: 逻辑回归 (Logistic Regression)", font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 保存
# ============================================================
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "4장_선형회귀_강의PPT_확장_china.pptx")
prs.save(out_path)
print(f"[Done] {out_path}")
print(f"[Slides] {len(prs.slides)}")
