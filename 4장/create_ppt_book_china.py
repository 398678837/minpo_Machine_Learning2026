"""第4章 线性回归 - 教材实习 PPT 生成脚本"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色调色板 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)
ACCENT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)


# ============================================================
# 幻灯片 1: 封面
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG, ACCENT_BLUE)
add_accent_line(s, Inches(4), Inches(1.3), Inches(5.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.6), Inches(11.333), Inches(0.5),
         "CHAPTER 04", font_size=22, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.3), Inches(11.333), Inches(1.2),
         "线性回归 - 预测保险费", font_size=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.6), Inches(11.333), Inches(0.5),
         "Linear Regression: Predicting Insurance Charges", font_size=22, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5), Inches(4.5), Inches(3.333), ACCENT_CYAN)
add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.4),
         "授课教师: Jung, Minpo", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.4), Inches(11.333), Inches(0.4),
         "课程: Machine Learning  |  学期: 2026年度 第1学期", font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片 2: 学习目标 & 学习顺序
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "学习目标 & 学习顺序")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "学习目标",
         ["使用线性回归模型学习保险数据集，",
          "预测保险费，",
          "并理解线性回归的工作原理。"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5),
         "学习顺序",
         ["4.1  问题定义",
          "4.2  导入库及加载数据",
          "4.3  数据确认",
          "4.4  预处理: 划分训练集/测试集",
          "4.5  数据建模",
          "4.6  使用模型进行预测",
          "4.7  模型评估",
          "4.8  理解: 线性回归原理"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 流水线箭头
arrow_y = Inches(5.2)
steps = ["4.1", "4.2", "4.3", "4.4", "4.5", "4.6", "4.7", "4.8"]
for i, st in enumerate(steps):
    x = Inches(0.6 + i * 0.75)
    add_shape(s, x, arrow_y, Inches(0.6), Inches(0.45), ACCENT_BLUE, radius=True)
    add_text(s, x, arrow_y, Inches(0.6), Inches(0.45), st, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s, x + Inches(0.55), arrow_y, Inches(0.25), Inches(0.45), "\u2192",
                 font_size=16, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片 3: 线性回归简介
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "线性回归简介", "Linear Regression Overview")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.2),
         "什么是 Linear Regression？",
         ["最基础的机器学习模型",
          "预测连续型变量（例：165.5cm、172.3cm 等）",
          "学习输入变量与输出变量之间的线性关系"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.2),
         "入选 TOP 10 算法的原因",
         ["作为最基础的算法，必须掌握",
          "可作为其他算法的基线（Baseline）使用",
          "结果易于解释，在实际工作中经常使用"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# 优缺点表格
add_table_slide(s,
    headers=["类别", "内容"],
    rows=[["优点", "简单直观，易于解释"],
          ["优点", "学习速度快，实现简便"],
          ["缺点", "预测能力可能较低"],
          ["缺点", "前提条件：自变量与因变量之间存在线性关系"]],
    left=Inches(0.6), top=Inches(4.6),
    col_widths=[2.0, 10.0], row_height=0.45, font_size=13)

# ============================================================
# 幻灯片 4: 4.1 问题定义
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.1", "问题定义", "Problem Definition")

add_card(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(1.8),
         "任务 (Mission)",
         ["利用保险数据集，",
          "预测保险公司收取的保险费(charges)！"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(1.8),
         "项目信息",
         ["难度: \u2605\u2606\u2606",
          "算法: Linear Regression",
          "数据集: insurance.csv"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.1), Inches(3.8), Inches(2.5),
         "数据信息",
         ["数据: insurance.csv",
          "因变量(y): charges",
          "自变量(X): age, sex, bmi,",
          "  children, smoker"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(4.1), Inches(3.8), Inches(2.5),
         "评估指标",
         ["RMSE",
          "(Root Mean Squared Error)",
          "通过预测值与实际值的差异",
          "的均方根来衡量"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(8.8), Inches(4.1), Inches(3.8), Inches(2.5),
         "使用的库",
         ["numpy", "pandas", "seaborn",
          "matplotlib", "sklearn"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 幻灯片 5: 4.2 导入库及加载数据
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.2", "导入库及加载数据", "Import Libraries & Load Data")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(11.5), Inches(2.5), [
    "import pandas as pd",
    "",
    "file_url = 'https://media.githubusercontent.com/media/",
    "            musthave-ML10/data_source/main/insurance.csv'",
    "",
    "data = pd.read_csv(file_url)",
], font_size=14)

add_card(s, Inches(0.6), Inches(5.0), Inches(5.5), Inches(1.8),
         "pandas",
         ["数据分析的核心库",
          "以 DataFrame 结构处理数据",
          "支持 CSV、Excel 等多种格式"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.5), Inches(5.0), Inches(5.6), Inches(1.8),
         "pd.read_csv()",
         ["从 URL 或本地路径读取 CSV 文件",
          "并转换为 DataFrame 的函数",
          "存储在 data 变量中使用"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 6: 4.3 数据确认 (1)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.3", "数据确认 (1)", "data 输出及 data.head()")

add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "data \u2192 1338 rows \u00d7 6 columns", font_size=18, color=ACCENT_CYAN, bold=True)

add_card(s, Inches(0.6), Inches(2.5), Inches(5.5), Inches(2.2),
         "变量说明",
         ["age     : 年龄（连续型）",
          "sex     : 性别（分类型，1=男，0=女）",
          "bmi     : 体质指数（连续型）",
          "children: 子女数（连续型）",
          "smoker  : 吸烟与否（分类型，1=是，0=否）",
          "charges : 保险费（因变量，连续型）"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# data.head() 表格
add_text(s, Inches(0.6), Inches(5.0), Inches(4), Inches(0.4),
         "data.head() 结果", font_size=16, color=ACCENT_GREEN, bold=True)

add_table_slide(s,
    headers=["age", "sex", "bmi", "children", "smoker", "charges"],
    rows=[["19", "0", "27.90", "0", "1", "16884.92"],
          ["18", "1", "33.77", "1", "0", "1725.55"],
          ["28", "1", "33.00", "3", "0", "4449.46"],
          ["33", "1", "22.70", "0", "0", "21984.47"],
          ["32", "1", "28.88", "0", "0", "3866.86"]],
    left=Inches(6.5), top=Inches(2.5),
    col_widths=[1.0, 0.8, 1.0, 1.0, 1.0, 1.6],
    row_height=0.42, font_size=12, header_font_size=12)

# ============================================================
# 幻灯片 7: 4.3 数据确认 (2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.3", "数据确认 (2)", "连续型 vs 分类型变量 & data.info()")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.0),
         "连续型变量 vs 分类型变量",
         ["\u25b6 连续型(Continuous): 连续数值（age, bmi, children, charges）",
          "\u25b6 分类型(Categorical): 表示类别/分组的值（sex, smoker）",
          "",
          "  sex   : 1 = 男性(male),  0 = 女性(female)",
          "  smoker: 1 = 吸烟者(yes), 0 = 非吸烟者(no)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
         "data.info() 结果",
         ["<class 'pandas.core.frame.DataFrame'>",
          "RangeIndex: 1338 entries, 0 to 1337",
          "Data columns (total 6 columns):",
          "",
          "  #  Column    Non-Null Count  Dtype",
          "  0  age       1338 non-null   int64",
          "  1  sex       1338 non-null   int64",
          "  2  bmi       1338 non-null   float64",
          "  3  children  1338 non-null   int64",
          "  4  smoker    1338 non-null   int64",
          "  5  charges   1338 non-null   float64",
          "",
          "dtypes: float64(2), int64(4)",
          "无缺失值(null) \u2192 预处理负担较小"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.5),
         "关键要点",
         ["\u2022 6个列，1338行",
          "\u2022 无缺失值(Missing Value)",
          "\u2022 sex、smoker 已完成数值编码",
          "\u2022 无需额外的编码预处理"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 8: 4.3 数据确认 (3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.3", "数据确认 (3)", "data.describe() & 四分位数")

add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
         "round(data.describe(), 2)", font_size=16, color=ACCENT_GREEN, bold=True)

add_table_slide(s,
    headers=["", "age", "sex", "bmi", "children", "smoker", "charges"],
    rows=[["count", "1338", "1338", "1338", "1338", "1338", "1338"],
          ["mean", "39.21", "0.51", "30.66", "1.09", "0.20", "13270.42"],
          ["std", "14.05", "0.50", "6.10", "1.21", "0.40", "12110.01"],
          ["min", "18.00", "0.00", "15.96", "0.00", "0.00", "1121.87"],
          ["25%", "27.00", "0.00", "26.30", "0.00", "0.00", "4740.29"],
          ["50%", "39.00", "1.00", "30.40", "1.00", "0.00", "9382.03"],
          ["75%", "51.00", "1.00", "34.69", "2.00", "0.00", "16639.91"],
          ["max", "64.00", "1.00", "53.13", "5.00", "1.00", "63770.43"]],
    left=Inches(0.6), top=Inches(2.5),
    col_widths=[1.2, 1.2, 1.0, 1.2, 1.2, 1.2, 1.8],
    row_height=0.4, font_size=11, header_font_size=12)

add_card(s, Inches(0.6), Inches(6.0), Inches(12.0), Inches(1.2),
         "四分位数 (Quantile) 概念",
         ["25% (Q1): 下位25%位置  |  50% (Q2/中位数): 正中间  |  75% (Q3): 上位25%位置",
          "charges 均值(13270) > 中位数(9382) \u2192 右偏分布（存在高额保险费）"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 9: 4.4 预处理 - 划分训练集和测试集（概念）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.4", "预处理: 划分训练集和测试集", "Train/Test Split Concept")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.0),
         "第1步: 分离自变量 / 因变量",
         ["自变量(X): age, sex, bmi, children, smoker",
          "因变量(y): charges（预测对象）",
          "",
          "\u2192 X 是原因，y 是结果"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "第2步: 分离训练集 / 测试集",
         ["训练集(Train): 模型用于学习的数据",
          "测试集(Test): 用于评估模型性能的数据",
          "",
          "\u2192 通常按 8:2 或 7:3 比例划分"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 4个数据集卡片
labels = [("X_train", "训练用自变量", ACCENT_BLUE),
          ("y_train", "训练用因变量", ACCENT_CYAN),
          ("X_test", "测试用自变量", ACCENT_ORANGE),
          ("y_test", "测试用因变量", ACCENT_RED)]
for i, (name, desc, clr) in enumerate(labels):
    x = Inches(0.6 + i * 3.15)
    add_shape(s, x, Inches(4.5), Inches(2.8), Inches(1.5), CARD_BG, clr, radius=True)
    add_text(s, x, Inches(4.6), Inches(2.8), Inches(0.5), name,
             font_size=22, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.2), Inches(2.8), Inches(0.5), desc,
             font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片 10: 4.4 数据集划分（代码）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.4", "数据集划分（代码）", "Train/Test Split Code")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(11.5), Inches(3.0), [
    "# 分离自变量(X)和因变量(y)",
    "X = data[['age', 'sex', 'bmi', 'children', 'smoker']]",
    "y = data['charges']",
    "",
    "# 划分训练集和测试集",
    "from sklearn.model_selection import train_test_split",
    "",
    "X_train, X_test, y_train, y_test = train_test_split(",
    "    X, y, test_size=0.2, random_state=100)",
], font_size=14)

add_card(s, Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.5),
         "test_size = 0.2",
         ["将全部数据的20%划分为测试集",
          "剩余80%作为训练集"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.5), Inches(5.5), Inches(5.6), Inches(1.5),
         "random_state = 100",
         ["固定随机采样以保证可重复性",
          "输入相同数字则始终得到相同结果"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 11: 4.5 数据建模
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.5", "数据建模", "Model Training")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(1.0),
         "建模三步骤",
         ["\u2460 选择算法 (LinearRegression)  \u2192  \u2461 创建模型 (model 对象)  \u2192  \u2462 训练模型 (fit)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_code_block(s, Inches(0.6), Inches(3.3), Inches(11.5), Inches(2.5), [
    "# 1. 选择算法 & 创建模型",
    "from sklearn.linear_model import LinearRegression",
    "",
    "model = LinearRegression()",
    "",
    "# 2. 训练模型",
    "model.fit(X_train, y_train)",
], font_size=14)

add_card(s, Inches(0.6), Inches(6.1), Inches(5.5), Inches(1.0),
         "LinearRegression()",
         ["sklearn 提供的线性回归模型类"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(6.5), Inches(6.1), Inches(5.6), Inches(1.0),
         "model.fit(X_train, y_train)",
         ["输入训练集(X, y)来训练模型"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 12: 4.6 使用模型进行预测
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.6", "使用模型进行预测", "Prediction")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(11.5), Inches(1.5), [
    "pred = model.predict(X_test)",
], font_size=16)

add_card(s, Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.5),
         "predict() 函数",
         ["\u2022 将测试集(X_test)输入已训练的模型",
          "\u2022 只需输入自变量(X)",
          "\u2022 不要输入因变量(y)！",
          "",
          "\u2192 pred 中存储了预测的保险费"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.5),
         "注意事项",
         ["\u2022 只能输入 X_test（不能输入 y_test）",
          "\u2022 必须使用未参与训练的数据进行预测",
          "  才能评估模型的泛化性能",
          "",
          "\u2192 将 pred 结果与 y_test 比较进行评估"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 13: 4.7.1 通过表格评估
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.7.1", "通过表格评估", "Evaluation with Table")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(7.0), Inches(1.8), [
    "comparison = pd.DataFrame({",
    "    'actual': y_test,",
    "    'pred': pred",
    "})",
    "comparison  # 268 rows x 2 columns",
], font_size=13)

add_card(s, Inches(8.0), Inches(2.2), Inches(4.6), Inches(1.8),
         "结果解读",
         ["针对268条测试数据，",
          "创建实际值(actual)与预测值(pred)",
          "并排对比的 DataFrame"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# 示例表格
add_text(s, Inches(0.6), Inches(4.3), Inches(4), Inches(0.4),
         "comparison 结果（部分）", font_size=15, color=ACCENT_GREEN, bold=True)

add_table_slide(s,
    headers=["index", "actual", "pred"],
    rows=[["764", "9095.07", "11834.68"],
          ["887", "5272.18", "8292.14"],
          ["890", "29330.98", "34069.01"],
          ["1293", "9301.89", "12077.35"],
          ["259", "44202.65", "36553.34"]],
    left=Inches(0.6), top=Inches(4.8),
    col_widths=[1.5, 2.5, 2.5],
    row_height=0.4, font_size=13)

add_card(s, Inches(7.5), Inches(4.5), Inches(5.2), Inches(2.5),
         "actual vs pred 对比",
         ["\u2022 actual: 实际保险费（正确答案）",
          "\u2022 pred: 模型预测的保险费",
          "",
          "\u2022 值越接近则预测越好",
          "\u2022 差异大则预测失败",
          "\u2022 仅凭表格难以把握整体模式"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 14: 4.7.2 通过图表评估
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.7.2", "通过图表评估", "Evaluation with Scatter Plot")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(2.0), [
    "import matplotlib.pyplot as plt",
    "import seaborn as sns",
    "",
    "plt.figure(figsize=(10, 10))",
    "sns.scatterplot(x='actual', y='pred',",
    "                data=comparison)",
], font_size=13)

add_card(s, Inches(7.5), Inches(2.2), Inches(5.2), Inches(4.5),
         "散点图(Scatter Plot) 解读",
         ["\u25b6 x轴: 实际值(actual)，y轴: 预测值(pred)",
          "",
          "\u2022 绿色虚线（对角线）附近：",
          "   actual \u2248 pred \u2192 预测准确",
          "",
          "\u2022 对角线上方：",
          "   pred > actual \u2192 过高预测",
          "",
          "\u2022 对角线下方：",
          "   pred < actual \u2192 过低预测",
          "",
          "\u2022 点越靠近对角线，模型越好"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 视觉区域表现
add_shape(s, Inches(0.6), Inches(4.5), Inches(6.5), Inches(2.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.8), Inches(4.6), Inches(6.1), Inches(0.4),
         "三个区域划分", font_size=14, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(0.8), Inches(5.0), Inches(6.1), Inches(1.8), [
    "\u2705 对角线附近 = 预测准确 (ideal zone)",
    "\u26a0\ufe0f 对角线上方 = 过高预测 (over-prediction)",
    "\u26a0\ufe0f 对角线下方 = 过低预测 (under-prediction)",
    "\u2192 大部分点集中在对角线附近，但也存在离群值",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

# ============================================================
# 幻灯片 15: 4.7.3 RMSE 评估
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.7.3", "RMSE 评估", "Root Mean Squared Error")

# MAE, MSE, RMSE 表格
add_table_slide(s,
    headers=["指标", "公式", "特点"],
    rows=[["MAE", "|实际-预测| 均值", "直观，对异常值不太敏感"],
          ["MSE", "(实际-预测)\u00b2 均值", "对大误差施加更大惩罚"],
          ["RMSE", "\u221aMSE", "单位与原始数据相同，易于解读"]],
    left=Inches(0.6), top=Inches(2.2),
    col_widths=[1.5, 3.5, 4.5],
    row_height=0.5, font_size=14)

add_card(s, Inches(0.6), Inches(4.2), Inches(5.5), Inches(1.5),
         "Table A vs Table B 对比",
         ["\u2022 即使绝对值之和(MAE)相同，平方和(MSE)也可能不同",
          "\u2022 一个大误差会大幅增加 MSE",
          "\u2022 RMSE 是对大误差敏感的指标"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_code_block(s, Inches(0.6), Inches(5.9), Inches(11.5), Inches(1.3), [
    "from sklearn.metrics import mean_squared_error",
    "",
    "mean_squared_error(y_test, pred) ** 0.5   # RMSE: 5684.93",
], font_size=14)

add_card(s, Inches(6.5), Inches(4.2), Inches(5.6), Inches(1.5),
         "结果: RMSE = 5684.93",
         ["\u2022 预测值与实际值的平均误差约为 $5,685",
          "\u2022 需结合保险费范围(1121~63770)来判断",
          "\u2022 单独来看难以判断好坏"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 幻灯片 16: R\u00b2 (决定系数)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.7", "R\u00b2 (决定系数)", "Coefficient of Determination")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.8),
         "R\u00b2 = SSR / SST",
         ["\u2022 SST (总变异): 实际值与均值之差的平方和",
          "\u2022 SSR (已解释变异): 预测值与均值之差的平方和",
          "\u2022 SSE (未解释变异): 实际值与预测值之差的平方和",
          "",
          "SST = SSR + SSE",
          "R\u00b2 = SSR / SST = 1 - (SSE / SST)",
          "",
          "\u2022 值在 0~1 之间，越接近1越好"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_code_block(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(1.5), [
    "model.score(X_train, y_train)",
    "# 结果: 0.7368",
], font_size=14)

add_card(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(2.0),
         "R\u00b2 解读标准",
         ["\u2022 R\u00b2 = 0.7368 \u2192 约 73.7% 的解释力",
          "\u2022 自变量解释了因变量变异的 73.7%",
          "",
          "\u2022 0.7~0.8 属于还不错(decent)的水平",
          "\u2022 0.9 以上属于非常好的模型"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(5.2), Inches(12.0), Inches(1.5),
         "R\u00b2 vs RMSE",
         ["\u2022 R\u00b2 是相对指标 (0~1)，RMSE 是绝对指标（单位：美元）",
          "\u2022 R\u00b2 越高则 RMSE 相对越低",
          "\u2022 两个指标结合来看，可以综合判断模型性能"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 17: 4.8 理解: 线性回归原理
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.8", "理解: 线性回归原理", "How Linear Regression Works")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.0),
         "y = ax + b",
         ["\u2022 a: 斜率 (coefficient, 系数)",
          "\u2022 b: y截距 (intercept)",
          "\u2022 自变量(x)与因变量(y)的线性关系",
          "\u2022 多个自变量: y = a\u2081x\u2081 + a\u2082x\u2082 + ... + b"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "损失函数 (Loss Function)",
         ["\u2022 最小化预测值与实际值的差异",
          "\u2022 OLS: 找到使(实际-预测)\u00b2之和最小的线",
          "\u2022 找到最优斜率(a)和截距(b)就是学习过程"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# 保险费预测公式
add_shape(s, Inches(0.6), Inches(4.3), Inches(12.0), Inches(1.5), CODE_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.8), Inches(4.4), Inches(11.6), Inches(0.4),
         "保险费预测公式（已训练模型）", font_size=15, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(0.8), Inches(4.9), Inches(11.6), Inches(0.8),
         "charges = 264.8\u00d7age + 17.3\u00d7sex + 297.5\u00d7bmi + 469.3\u00d7children + 23469.3\u00d7smoker - 11577.0",
         font_size=16, color=ACCENT_GREEN, bold=True, font_name='Consolas')

add_card(s, Inches(0.6), Inches(6.0), Inches(12.0), Inches(1.2),
         "公式解读",
         ["age 增加1 \u2192 charges 约增加265  |  smoker=1(吸烟) \u2192 charges 约增加23,469  |  截距(b) = -11,577"],
         title_color=ACCENT_YELLOW, border=ACCENT_YELLOW)

# ============================================================
# 幻灯片 18: 4.8 系数与截距确认
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4.8", "系数与截距确认", "Coefficients & Intercept")

add_code_block(s, Inches(0.6), Inches(2.2), Inches(7.0), Inches(3.0), [
    "# 确认系数（斜率）",
    "model.coef_",
    "# array([ 264.8, 17.3, 297.5, 469.3, 23469.3])",
    "",
    "# 以易读格式输出系数",
    "pd.Series(model.coef_, index=X.columns)",
    "# age        264.8",
    "# sex         17.3",
    "# bmi        297.5",
    "# children   469.3",
    "# smoker   23469.3",
    "",
    "# 确认截距(y-intercept)",
    "model.intercept_",
    "# -11576.999976112367",
], font_size=12)

add_card(s, Inches(8.0), Inches(2.2), Inches(4.6), Inches(3.0),
         "系数解读",
         ["\u2022 age: 年龄增加1岁 \u2192 保险费约增加 $265",
          "\u2022 sex: 性别差异约 $17（非常小）",
          "\u2022 bmi: BMI 增加1 \u2192 保险费约增加 $298",
          "\u2022 children: 子女增加1人 \u2192 约增加 $469",
          "\u2022 smoker: 吸烟者 \u2192 约增加 $23,469！",
          "",
          "\u2192 吸烟与否对保险费影响最大"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(0.6), Inches(5.5), Inches(12.0), Inches(1.5),
         "截距 (Intercept) = -11,577",
         ["\u2022 所有自变量为0时的基础值（实际上不会发生的理论值）",
          "\u2022 将系数(coef_)与截距(intercept_)结合即可得到完整的预测公式"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 19: 学习总结
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
slide_header(s, "", "学习总结", "Chapter 04 Summary")

# 流水线图
steps_full = [
    ("问题定义", ACCENT_BLUE),
    ("数据加载", ACCENT_CYAN),
    ("数据确认", ACCENT_GREEN),
    ("预处理", ACCENT_ORANGE),
    ("建模", ACCENT_PURPLE),
    ("预测", ACCENT_RED),
    ("评估", ACCENT_YELLOW),
]
for i, (label, clr) in enumerate(steps_full):
    x = Inches(0.6 + i * 1.75)
    add_shape(s, x, Inches(2.3), Inches(1.5), Inches(0.8), clr, radius=True)
    add_text(s, x, Inches(2.3), Inches(1.5), Inches(0.8), label,
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps_full) - 1:
        add_text(s, x + Inches(1.45), Inches(2.3), Inches(0.35), Inches(0.8), "\u2192",
                 font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_card(s, Inches(0.6), Inches(3.5), Inches(5.8), Inches(3.5),
         "核心摘要",
         ["\u2022 线性回归: y = a\u2081x\u2081 + a\u2082x\u2082 + ... + b",
          "\u2022 用于连续型变量预测的基础算法",
          "\u2022 可用 sklearn 简单实现",
          "   LinearRegression().fit() \u2192 predict()",
          "",
          "\u2022 评估: RMSE = 5,685 / R\u00b2 = 0.7368",
          "\u2022 吸烟与否(smoker)对保险费影响最大",
          "",
          "\u2022 优点: 简单、易于解释、学习速度快",
          "\u2022 局限: 难以捕捉非线性关系"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(3.5), Inches(5.8), Inches(3.5),
         "完整代码流水线",
         ["1. import pandas as pd",
          "2. data = pd.read_csv(file_url)",
          "3. data.head() / data.info() / data.describe()",
          "4. X, y 分离 \u2192 train_test_split()",
          "5. model = LinearRegression()",
          "6. model.fit(X_train, y_train)",
          "7. pred = model.predict(X_test)",
          "8. pd.DataFrame({'actual':y_test, 'pred':pred})",
          "9. sns.scatterplot(x='actual', y='pred', ...)",
          "10. mean_squared_error(y_test, pred)**0.5",
          "11. model.score() / model.coef_ / model.intercept_"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)


# ── 保存 ──
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "4장_선형회귀_교재실습_china.pptx")
prs.save(output_path)
print(f"PPT created: {output_path}")
print(f"Total slides: {len(prs.slides)}")
