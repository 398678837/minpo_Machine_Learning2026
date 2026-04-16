"""7장 나이브 베이즈 - 확장 강의 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)

def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = '맑은 고딕'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), accent)
    add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), accent)
    add_text(s, Inches(0), Inches(2.0), prs.slide_width, Inches(0.5),
             f"SECTION {section_num}", font_size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, Inches(5.5), Inches(2.7), Inches(2.3), accent)
    add_text(s, Inches(0), Inches(3.0), prs.slide_width, Inches(1.0),
             title, font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.5),
             subtitle, font_size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    return s

# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(1.2), prs.slide_width, Inches(0.5),
         "MACHINE LEARNING", font_size=18, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.0), Inches(1.9), Inches(3.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.2), prs.slide_width, Inches(1.2),
         "7장: 나이브 베이즈\n(Naive Bayes)", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.8), prs.slide_width, Inches(0.5),
         "베이즈 정리 | 조건부 독립 | 텍스트 분류 | 스팸 필터링", font_size=18, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.0), prs.slide_width, Inches(0.4),
         "참고: Rish(2001), McCallum & Nigam(1998), Rennie et al.(2003), Zhang(2004), Metsis et al.(2006)",
         font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.8), prs.slide_width, Inches(0.4),
         "기계학습 | 2026년 1학기", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 2: 목차
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "", "목차 (Table of Contents)")
sections = [
    ("7.1", "베이즈 정리 (Bayes' Theorem)", ACCENT_BLUE),
    ("7.2", "빈도주의 vs 베이지안", ACCENT_CYAN),
    ("7.3", "나이브 가정 (Naive Assumption)", ACCENT_GREEN),
    ("7.4", "나이브 베이즈 변형들", ACCENT_ORANGE),
    ("7.5", "텍스트 분류 -- BoW와 TF-IDF", ACCENT_PURPLE),
    ("7.6", "스팸 필터링", ACCENT_RED),
    ("7.7", "라플라스 스무딩과 로그 확률 트릭", ACCENT_BLUE),
    ("7.8", "생성 모델 vs 판별 모델", ACCENT_CYAN),
    ("7.9", "논문 리뷰 (Rish, McCallum, Rennie)", ACCENT_GREEN),
    ("7.10~12", "실습: GaussianNB 구현 / 텍스트 분류 / 스팸 필터", ACCENT_ORANGE),
    ("7.13", "응용 사례 (감정분석, 문서분류, 의료진단)", ACCENT_PURPLE),
    ("7.14", "핵심 요약 + 복습 질문", ACCENT_RED),
]
for i, (num, title, clr) in enumerate(sections):
    y = Inches(1.9) + Inches(0.43) * i
    add_shape(s, Inches(1.5), y, Inches(0.08), Inches(0.32), clr)
    add_text(s, Inches(1.8), y, Inches(1.2), Inches(0.35), num, font_size=14, color=clr, bold=True)
    add_text(s, Inches(3.0), y, Inches(8), Inches(0.35), title, font_size=15, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 3: 학습 목표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "", "학습 목표 (Learning Objectives)")
objectives = [
    "1. 베이즈 정리를 유도하고, 사전확률/우도/사후확률의 의미를 설명할 수 있다",
    "2. 빈도주의 vs 베이지안 관점의 차이를 이해한다",
    "3. 나이브(조건부 독립) 가정의 의미와 왜 깨져도 잘 작동하는지 설명할 수 있다",
    "4. GaussianNB, MultinomialNB, BernoulliNB, ComplementNB를 구분하고 적절히 선택할 수 있다",
    "5. CountVectorizer와 TF-IDF의 원리를 수식으로 설명할 수 있다",
    "6. 라플라스 스무딩과 로그 확률 트릭의 필요성과 수식을 이해한다",
    "7. 생성 모델(나이브 베이즈)과 판별 모델(로지스틱 회귀)의 관계를 파악한다",
    "8. 스팸 필터링 파이프라인을 구축하고 특성 중요도를 분석할 수 있다",
]
add_bullet_list(s, Inches(1.0), Inches(1.9), Inches(11), Inches(5), objectives,
                font_size=17, color=LIGHT_GRAY, spacing=Pt(10))

# ============================================================
# SECTION 1: 베이즈 정리
# ============================================================
section_divider("베이즈 정리", "Bayes' Theorem -- 사전지식을 증거로 갱신", "1")

# 슬라이드 5: 조건부 확률에서 베이즈 정리까지
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "1", "조건부 확률에서 베이즈 정리까지")
add_text(s, Inches(0.8), Inches(1.9), Inches(11), Inches(0.4),
         "18세기 영국 수학자 토마스 베이즈(Thomas Bayes)가 제안", font_size=16, color=LIGHT_GRAY)
add_text(s, Inches(0.8), Inches(2.4), Inches(11), Inches(0.4),
         "사전 지식(prior)을 새로운 증거(evidence)로 갱신하여 사후 확률(posterior)을 얻는 공식",
         font_size=17, color=ACCENT_CYAN, bold=True)
add_card(s, Inches(0.8), Inches(3.0), Inches(5.5), Inches(2.0), "유도 과정", [
    "조건부 확률의 정의:",
    "  P(A|B) = P(A and B) / P(B)",
    "  P(B|A) = P(A and B) / P(A)",
    "P(A and B) 소거:",
    "  P(A|B) * P(B) = P(B|A) * P(A)",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_card(s, Inches(6.8), Inches(3.0), Inches(5.5), Inches(2.0), "베이즈 정리 (최종)", [
    "",
    "  P(A|B) = P(B|A) * P(A) / P(B)",
    "",
    "사후확률 = 우도 x 사전확률 / 증거",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_text(s, Inches(0.8), Inches(5.3), Inches(11), Inches(0.8),
         "핵심: 데이터를 관측한 후 믿음을 합리적으로 갱신하는 수학적 원리",
         font_size=18, color=ACCENT_ORANGE, bold=True)

# 슬라이드 6: 각 항의 의미
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "1", "베이즈 정리 -- 각 항의 의미")
headers = ["용어", "영문", "수식", "설명"]
rows = [
    ["사후확률", "Posterior", "P(A|B)", "데이터 B 관측 후 A가 참일 확률"],
    ["우도", "Likelihood", "P(B|A)", "A가 참일 때 데이터 B가 관측될 확률"],
    ["사전확률", "Prior", "P(A)", "데이터 관측 전 A가 참일 확률"],
    ["증거", "Evidence", "P(B)", "데이터 B의 전체 관측 확률 (정규화 상수)"],
]
add_table_slide(s, headers, rows, Inches(0.8), Inches(2.0), [2.0, 2.0, 2.5, 5.5],
                row_height=0.55, font_size=14)
add_text(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.6),
         "직관: 사후확률 = (우도 x 사전확률) / 증거",
         font_size=20, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

# 슬라이드 7: 의료진단 사례
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "1", "의료진단 사례 -- 민감도 / 특이도 / 유병률")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.2), "상황 설정", [
    "유병률(Prevalence): P(D) = 0.001 (0.1%)",
    "민감도(Sensitivity): P(+|D) = 0.99 (99%)",
    "특이도(Specificity): P(-|~D) = 0.95 (95%)",
    "",
    "질문: 양성 판정 -> 실제 질병 확률?",
], title_color=ACCENT_RED, border=ACCENT_RED)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.2), "베이즈 정리 적용", [
    "P(D|+) = P(+|D) * P(D) / P(+)",
    "",
    "P(+) = P(+|D)*P(D) + P(+|~D)*P(~D)",
    "     = 0.99*0.001 + 0.05*0.999",
    "     = 0.05094",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(0.8), Inches(4.4), Inches(5.5), Inches(1.6), "계산 결과", [
    "P(D|+) = 0.99 * 0.001 / 0.05094",
    "       = 0.0194 (약 1.94%)",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_card(s, Inches(6.8), Inches(4.4), Inches(5.5), Inches(1.6), "핵심 교훈", [
    "99% 정확도 검사로 양성이어도 실제 질병 확률 ~2%",
    "사전확률(유병률)이 사후확률에 결정적 영향",
    "베이즈 정리가 직관과 다른 결과의 대표적 사례",
], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# 슬라이드 8: 스팸 메일 사례
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "1", "스팸 메일 사례 -- 베이즈 정리 적용")
headers = ["구분", "전체 메일 수", '"free" 포함 수']
rows = [
    ["스팸", "400 (40%)", "200 (스팸 중 50%)"],
    ["정상", "600 (60%)", "60 (정상 중 10%)"],
    ["합계", "1000 (100%)", "260 (전체 26%)"],
]
add_table_slide(s, headers, rows, Inches(0.8), Inches(2.0), [2.5, 3.5, 3.5],
                row_height=0.55, font_size=14)
add_card(s, Inches(0.8), Inches(4.2), Inches(11.2), Inches(2.0), "베이즈 정리 계산", [
    'P(spam | "free") = P("free" | spam) * P(spam) / P("free")',
    "                  = 0.5 * 0.4 / 0.26 = 0.769 (약 76.9%)",
    "",
    '"free"가 포함된 이메일이 스팸일 확률은 약 76.9%',
], title_color=ACCENT_CYAN, border=ACCENT_BLUE)

# ============================================================
# SECTION 2: 빈도주의 vs 베이지안
# ============================================================
section_divider("빈도주의 vs 베이지안", "Frequentist vs Bayesian -- 확률의 두 가지 철학", "2")

# 슬라이드 10: 빈도주의
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "2", "빈도주의 (Frequentist)")
add_card(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(3.5), "빈도주의 핵심 개념", [
    "확률의 정의: 동일한 실험을 무한히 반복했을 때 사건이 발생하는 상대적 빈도",
    "",
    "파라미터 관점: 파라미터는 고정된 미지의 상수이다",
    "",
    "추론 방법: 최대우도추정(MLE), 신뢰구간, p-값",
    "",
    '예시: "이 동전을 10,000번 던지면 앞면이 약 5,000번 나올 것이다"',
    "",
    "수식: theta_MLE = argmax_theta P(D|theta)",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)
add_text(s, Inches(0.8), Inches(5.7), Inches(11), Inches(0.5),
         "데이터만으로 결론을 내리며, 사전 정보를 사용하지 않음",
         font_size=17, color=ACCENT_ORANGE, bold=True)

# 슬라이드 11: 베이지안
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "2", "베이지안 (Bayesian)")
add_card(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(3.5), "베이지안 핵심 개념", [
    "확률의 정의: 불확실성에 대한 주관적 믿음의 정도 (degree of belief)",
    "",
    "파라미터 관점: 파라미터 자체가 확률분포를 가지는 확률변수이다",
    "",
    "추론 방법: 베이즈 정리로 사전분포를 사후분포로 갱신",
    "",
    '예시: "이 동전이 공정할 확률이 80%라고 믿는다"',
    "",
    "수식: P(theta|D) = P(D|theta) * P(theta) / P(D)",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_text(s, Inches(0.8), Inches(5.7), Inches(11), Inches(0.5),
         "사전 지식을 활용하여 데이터와 결합 -- 소량 데이터에서 강점",
         font_size=17, color=ACCENT_ORANGE, bold=True)

# 슬라이드 12: 비교표
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "2", "빈도주의 vs 베이지안 비교")
headers = ["관점", "빈도주의", "베이지안"]
rows = [
    ["확률 정의", "장기적 빈도", "믿음의 정도"],
    ["파라미터", "고정된 상수", "확률변수"],
    ["사전 정보", "사용하지 않음", "사전분포로 반영"],
    ["대표 추정법", "MLE", "MAP, 사후분포"],
    ["소량 데이터", "추정 불안정", "사전분포가 보완"],
    ["대량 데이터", "수렴 보장", "사전분포 영향 감소"],
]
add_table_slide(s, headers, rows, Inches(1.0), Inches(2.0), [2.5, 4.5, 4.5],
                row_height=0.55, font_size=14)
add_text(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
         "나이브 베이즈: 베이지안 접근 기반 + 빈도 기반 파라미터 추정 = 혼합적 접근",
         font_size=17, color=ACCENT_CYAN, bold=True)

# ============================================================
# SECTION 3: 나이브 가정
# ============================================================
section_divider("나이브 가정", "Naive Assumption -- 조건부 독립과 그 정당화", "3")

# 슬라이드 14: 조건부 독립이란
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "3", "조건부 독립이란")
add_text(s, Inches(0.8), Inches(1.9), Inches(11), Inches(0.5),
         '"나이브(Naive)" = "순진한" -- 모든 특성이 클래스가 주어졌을 때 조건부 독립이라고 가정',
         font_size=17, color=ACCENT_CYAN, bold=True)
add_card(s, Inches(0.8), Inches(2.6), Inches(5.5), Inches(2.0), "조건부 독립의 수학적 정의", [
    "클래스 C가 주어졌을 때:",
    "",
    "P(x1, x2, ..., xn | C) = PROD_i P(xi | C)",
    "",
    "결합확률이 주변확률의 곱으로 분해",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_card(s, Inches(6.8), Inches(2.6), Inches(5.5), Inches(2.0), "나이브 베이즈 분류 규칙", [
    "",
    "C_hat = argmax_C  P(C) * PROD_i P(xi | C)",
    "",
    "분모 P(x1,...,xn)은 모든 클래스에 동일",
    "-> argmax 계산 시 무시 가능",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_text(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.8),
         "파라미터 수: 결합분포 O(k^n) -> 독립가정 O(n*k)  (n: 특성 수, k: 클래스 수)",
         font_size=17, color=ACCENT_ORANGE, bold=True)

# 슬라이드 15: 왜 독립 가정은 현실에서 깨지는가
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "3", "왜 현실에서 독립 가정은 깨지는가")
add_text(s, Inches(0.8), Inches(1.9), Inches(11), Inches(0.4),
         "텍스트 데이터에서 단어들은 명백히 서로 독립이 아니다:", font_size=17, color=LIGHT_GRAY)
add_card(s, Inches(0.8), Inches(2.6), Inches(3.5), Inches(1.8), "양의 상관관계 예시", [
    '"machine" -> "learning" 동시 출현',
    '"New" -> "York" 동시 출현',
    '"free" -> "win", "prize" 동시 출현',
], title_color=ACCENT_RED, border=ACCENT_RED)
add_card(s, Inches(4.6), Inches(2.6), Inches(3.8), Inches(1.8), "문법적 의존성", [
    "형용사와 명사의 수식 관계",
    "주어-동사 일치",
    "전치사구의 관계",
], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)
add_card(s, Inches(8.7), Inches(2.6), Inches(3.8), Inches(1.8), "의미적 의존성", [
    "동의어 관계: fast = quick",
    "반의어 관계: good vs bad",
    "상위어/하위어: animal -> dog",
], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)
add_text(s, Inches(0.8), Inches(4.8), Inches(11), Inches(1.0),
         "그럼에도 불구하고 나이브 베이즈는 놀라울 만큼 잘 작동한다! -- Why?",
         font_size=22, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

# 슬라이드 16: Zhang(2004)의 최적성 정리
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "3", "Zhang(2004)의 최적성 정리")
add_card(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(2.0), "핵심 정리: 이진 분류에서 최적 조건", [
    "나이브 베이즈가 베이즈 최적 분류기와 동일한 결과를 내는 조건:",
    "",
    "SUM_{i!=j} dep_ij(c1) = SUM_{i!=j} dep_ij(c2)",
    "",
    "의존성 측도: dep_ij(c) = P(xi,xj|c) - P(xi|c)*P(xj|c)",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)
add_text(s, Inches(0.8), Inches(4.2), Inches(11), Inches(0.4),
         "해석: 의존성의 총량이 클래스 간에 동일하게 분포하면, 의존성이 아무리 강해도 NB는 최적",
         font_size=17, color=ACCENT_ORANGE, bold=True)

# 슬라이드 17: 왜 깨져도 잘 작동하는가 -- 3가지 이유
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "3", "깨져도 잘 작동하는 3가지 이유")
add_card(s, Inches(0.8), Inches(1.9), Inches(3.6), Inches(3.5),
         "1. 분류 vs 확률 추정 (Rish 2001)", [
    "정확한 확률값이 아닌 클래스 간",
    "상대적 순위만 정확하면 OK",
    "",
    "P(C1|x)=0.8 이든 0.95 이든",
    "C1으로 올바르게 분류하면 됨",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(4.7), Inches(1.9), Inches(3.6), Inches(3.5),
         "2. 의존성의 상쇄 효과", [
    "양의 상관관계와 음의 상관관계가",
    "혼재할 때 서로 효과가 상쇄",
    "",
    "결과적으로 전체 성능이 유지됨",
    "Zhang(2004)의 이론적 근거",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_card(s, Inches(8.6), Inches(1.9), Inches(3.6), Inches(3.5),
         "3. 계산 효율성 (과적합 방지)", [
    "독립 가정 -> 파라미터 수 O(n*k)",
    "결합 분포 -> 파라미터 수 O(k^n)",
    "",
    "적은 파라미터 = 적은 과적합",
    "소량 데이터에서 강건",
], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# SECTION 4: 나이브 베이즈 변형들
# ============================================================
section_divider("나이브 베이즈 변형들", "GaussianNB | MultinomialNB | BernoulliNB | ComplementNB", "4")

# 슬라이드 19: GaussianNB
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "GaussianNB -- 가우시안 나이브 베이즈")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.8), "핵심 개념", [
    "적용 대상: 연속형(continuous) 실수 데이터",
    "",
    "분포 가정: 각 특성이 클래스별 정규분포",
    "",
    "P(xi|C) = 1/sqrt(2*pi*sigma^2)",
    "          * exp(-(xi - mu)^2 / (2*sigma^2))",
    "",
    "사용 예: Iris 분류, 키/몸무게 기반 성별 예측",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(1.6), "파라미터 추정", [
    "mu_Ci = (1/Nc) * SUM x_ji  (클래스 C 특성 i의 평균)",
    "sigma^2_Ci = (1/Nc) * SUM (x_ji - mu)^2  (표본 분산)",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_code_block(s, Inches(6.8), Inches(3.8), Inches(5.5), Inches(1.2), [
    "from sklearn.naive_bayes import GaussianNB",
    "model = GaussianNB()",
    "model.fit(X_train, y_train)",
])

# 슬라이드 20: MultinomialNB
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "MultinomialNB -- 다항 나이브 베이즈")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.8), "핵심 개념", [
    "적용 대상: 카운트(count) 또는 빈도 데이터",
    "",
    "분포 가정: 각 특성이 다항분포",
    "",
    "P(xi|C) = (count(xi,C) + alpha)",
    "        / SUM_j (count(xj,C) + alpha)",
    "",
    "핵심: 단어의 출현 빈도 기반 확률 계산",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(1.6), "McCallum & Nigam(1998)", [
    "텍스트 분류에서 BernoulliNB보다 우수",
    "어휘가 클수록 MultinomialNB 이점 증가",
    "사용 예: 스팸 필터링, 감성 분석, 문서 분류",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_code_block(s, Inches(6.8), Inches(3.8), Inches(5.5), Inches(1.2), [
    "from sklearn.naive_bayes import MultinomialNB",
    "model = MultinomialNB(alpha=1.0)",
    "model.fit(X_train_tfidf, y_train)",
])

# 슬라이드 21: BernoulliNB
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "BernoulliNB -- 베르누이 나이브 베이즈")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(3.0), "핵심 개념", [
    "적용 대상: 이진(binary) 데이터 (0 또는 1)",
    "",
    "분포 가정: 각 특성이 베르누이 분포",
    "",
    "P(x|C) = PROD_i P(xi|C)^xi",
    "         * (1 - P(xi|C))^(1-xi)",
    "",
    "핵심: 단어의 존재 여부(있다/없다)만 고려",
    "단어 부재 정보도 명시적 활용!",
], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(1.6), "McCallum & Nigam(1998) 발견", [
    "어휘 사전이 작을 때: BernoulliNB 경쟁력 O",
    "어휘가 클수록: MultinomialNB가 우수",
    "MultinomialNB와의 핵심 차이: 부재 정보 활용",
], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)
add_code_block(s, Inches(6.8), Inches(3.8), Inches(5.5), Inches(1.2), [
    "from sklearn.naive_bayes import BernoulliNB",
    "model = BernoulliNB(alpha=1.0)",
    "model.fit(X_train_binary, y_train)",
])

# 슬라이드 22: ComplementNB
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "ComplementNB -- 보완 나이브 베이즈 (Rennie 2003)")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.8), "핵심 아이디어", [
    "각 클래스 c의 파라미터를 직접 추정 대신,",
    "보완 클래스(complement class ~c)의 데이터 사용",
    "",
    "theta_ci = (alpha_i + SUM_{j:yj!=c} f_ji)",
    "         / (alpha + SUM_i' SUM_{j:yj!=c} f_ji')",
    "",
    "분류 규칙:",
    "c_hat = argmin_c SUM_i fi * log(theta_ci)",
], title_color=ACCENT_RED, border=ACCENT_RED)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(1.6), "장점", [
    "클래스 불균형 문제를 완화",
    "표준 MultinomialNB 대비 약 5~8% 정확도 향상",
    "SVM에 근접하는 성능 달성",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_code_block(s, Inches(6.8), Inches(3.8), Inches(5.5), Inches(1.2), [
    "from sklearn.naive_bayes import ComplementNB",
    "model = ComplementNB(alpha=1.0)",
    "model.fit(X_train_tfidf, y_train)",
])

# 슬라이드 23: 변형 비교 요약
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "4", "나이브 베이즈 변형 비교 요약")
headers = ["종류", "데이터 유형", "분포 가정", "대표 사용 사례"]
rows = [
    ["GaussianNB", "연속형 (실수)", "정규분포", "수치형 특성 분류"],
    ["MultinomialNB", "카운트 (정수)", "다항분포", "텍스트 분류 (빈도 기반)"],
    ["BernoulliNB", "이진 (0/1)", "베르누이분포", "텍스트 분류 (존재 유무)"],
    ["ComplementNB", "카운트/빈도", "보완 다항분포", "불균형 텍스트 분류"],
]
add_table_slide(s, headers, rows, Inches(0.8), Inches(2.2), [2.5, 2.5, 3.0, 4.0],
                row_height=0.6, font_size=14)
add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.8),
         "실무 팁: 텍스트 분류에서는 ComplementNB + TfidfVectorizer 조합이 가장 안정적",
         font_size=18, color=ACCENT_ORANGE, bold=True)

# ============================================================
# SECTION 5: 텍스트 분류 -- BoW와 TF-IDF
# ============================================================
section_divider("텍스트 분류", "Bag of Words와 TF-IDF -- 텍스트를 수치로 변환", "5")

# 슬라이드 25: Bag of Words 모델
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "5", "Bag of Words (BoW) 모델")
add_text(s, Inches(0.8), Inches(1.9), Inches(11), Inches(0.4),
         "텍스트를 단어의 출현 빈도로 표현하는 가장 기본적인 벡터화 방법",
         font_size=17, color=ACCENT_CYAN, bold=True)
add_card(s, Inches(0.8), Inches(2.6), Inches(5.5), Inches(2.5), "동작 원리", [
    "1. 전체 문서에서 모든 고유 단어 수집",
    "   -> 어휘 사전(Vocabulary) 구축",
    "",
    "2. 각 문서에서 단어 등장 횟수 계산",
    "   -> 문서-단어 행렬(DTM) 생성",
    "",
    "특징: 단어의 순서를 무시, 빈도만 고려",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_card(s, Inches(6.8), Inches(2.6), Inches(5.5), Inches(2.5), "예시: 문서-단어 행렬", [
    "문서1: 'free money free'",
    "  -> [card:0, free:2, gift:0, money:1, transfer:0]",
    "",
    "문서2: 'money transfer'",
    "  -> [card:0, free:0, gift:0, money:1, transfer:1]",
    "",
    "문서3: 'free gift card'",
    "  -> [card:1, free:1, gift:1, money:0, transfer:0]",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_text(s, Inches(0.8), Inches(5.4), Inches(11), Inches(0.6),
         "출력은 희소 행렬(Sparse Matrix) -- 대부분 값이 0이므로 메모리 효율적 저장",
         font_size=16, color=LIGHT_GRAY)

# 슬라이드 26: CountVectorizer
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "5", "CountVectorizer -- sklearn 구현")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(2.5), [
    "from sklearn.feature_extraction.text import CountVectorizer",
    "",
    "corpus = ['free money free', 'money transfer', 'free gift card']",
    "",
    "cv = CountVectorizer()",
    "X = cv.fit_transform(corpus)",
    "",
    "print(cv.vocabulary_)   # 단어-인덱스 매핑",
    "# {'card': 0, 'free': 1, 'gift': 2, 'money': 3, 'transfer': 4}",
    "",
    "print(X.toarray())      # 문서-단어 행렬",
    "# [[0, 2, 0, 1, 0],  <- 'free money free'",
    "#  [0, 0, 0, 1, 1],  <- 'money transfer'",
    "#  [1, 1, 1, 0, 0]]  <- 'free gift card'",
], font_size=13)
add_text(s, Inches(0.8), Inches(4.7), Inches(11), Inches(0.5),
         "한계: 모든 문서에 공통적으로 등장하는 단어('the', 'is' 등)에 높은 가중치 부여",
         font_size=17, color=ACCENT_RED, bold=True)

# 슬라이드 27: TF-IDF 수식
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "5", "TF-IDF -- Term Frequency x Inverse Document Frequency")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5), "TF-IDF 수식", [
    "TF-IDF(t, d) = TF(t, d) x IDF(t)",
    "",
    "TF(t,d) = f(t,d) / SUM_t' f(t',d)",
    "  -> 문서 d에서 단어 t의 빈도 비율",
    "",
    "IDF(t) = log(N / |{d : t in d}|)",
    "  -> 단어 t가 전체 문서에서 얼마나 드문지",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.5), "핵심 아이디어", [
    "TF 높음 + IDF 높음 = 중요한 단어",
    "  -> 특정 문서에서 자주 등장 & 전체적으로 드문",
    "",
    "TF 높음 + IDF 낮음 = 덜 중요한 단어",
    "  -> 모든 문서에서 공통적으로 등장",
    "",
    "IDF의 역할: 흔한 단어의 가중치를 낮춤",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_code_block(s, Inches(0.8), Inches(4.7), Inches(11.2), Inches(1.2), [
    "from sklearn.feature_extraction.text import TfidfVectorizer",
    "",
    "tfidf = TfidfVectorizer(max_features=5000, stop_words='english')",
    "X_tfidf = tfidf.fit_transform(corpus)",
])

# 슬라이드 28: CountVectorizer vs TfidfVectorizer 비교
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "5", "CountVectorizer vs TfidfVectorizer")
headers = ["특성", "CountVectorizer", "TfidfVectorizer"]
rows = [
    ["가중치", "단순 빈도 (정수)", "TF-IDF 가중치 (실수)"],
    ["흔한 단어 처리", "높은 값 그대로", "IDF로 가중치 감소"],
    ["희귀 단어 처리", "낮은 값 그대로", "IDF로 가중치 증가"],
    ["일반적 성능", "기본적", "대부분의 경우 우수"],
]
add_table_slide(s, headers, rows, Inches(1.0), Inches(2.2), [3.0, 4.5, 4.5],
                row_height=0.6, font_size=14)
add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.6),
         "실무: 텍스트 분류에서는 TfidfVectorizer가 대부분 CountVectorizer보다 우수",
         font_size=18, color=ACCENT_ORANGE, bold=True)

# ============================================================
# SECTION 6: 스팸 필터링
# ============================================================
section_divider("스팸 필터링", "Bayesian Spam Filtering -- 역사와 벤치마크", "6")

# 슬라이드 30: 역사 -- SpamBayes
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "6", "역사 -- SpamBayes와 초기 베이지안 필터")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(3.0), "Paul Graham (2002)", [
    '"A Plan for Spam" 에세이',
    "나이브 베이즈로 스팸의 99.5% 이상 분류",
    "",
    "베이지안 스팸 필터링의 본격적 시작",
    "",
    "단어의 조건부 확률로 스팸 여부 판별",
    "P(spam|word1, word2, ...) 계산",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(3.0), "SpamBayes 프로젝트", [
    "Python으로 작성된 오픈소스 스팸 필터",
    "",
    "주요 특징:",
    "  - 사용자별 학습 가능 (개인화)",
    '  - "스팸", "정상", "불확실" 3범주 분류',
    "  - 단어 출현 빈도 기반 베이지안 확률",
    "",
    "현대 이메일 클라이언트의 기초",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 슬라이드 31: Metsis et al. (2006) 벤치마크
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "6", "Metsis et al. (2006) -- 어떤 나이브 베이즈?")
add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.4),
         '"Spam Filtering with Naive Bayes -- Which Naive Bayes?"  6개 데이터셋 체계적 비교',
         font_size=15, color=LIGHT_GRAY)
headers = ["변형", "설명", "평균 AUC"]
rows = [
    ["MNB", "단어 빈도 기반", "0.961"],
    ["MBNB", "단어 존재/부재 기반", "0.943"],
    ["MNBB (최우수)", "빈도를 이진화 + MultinomialNB", "0.972"],
    ["MNB + TF-IDF", "TF-IDF 가중치 적용", "0.965"],
    ["Flexible Bayes", "가우시안 커널 밀도 추정", "0.938"],
]
add_table_slide(s, headers, rows, Inches(0.8), Inches(2.4), [3.0, 5.0, 2.5],
                row_height=0.5, font_size=13)
add_text(s, Inches(0.8), Inches(5.3), Inches(11), Inches(0.8),
         "핵심: MNBB(MultinomialNB + Boolean)가 최우수\n스팸 판별에서 빈도보다 존재 여부가 중요하지만, MultinomialNB의 확률 모델이 더 적합",
         font_size=16, color=ACCENT_CYAN, bold=True)

# 슬라이드 32: 비용 민감 평가
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "6", "비용 민감 평가 -- FP vs FN의 비대칭성")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5), "FP (정상 -> 스팸)", [
    "중요한 메일을 잃을 수 있음",
    "매우 심각한 오류",
    "",
    "예: 면접 통보, 계약서, 긴급 연락",
    "-> 한 번의 FP가 큰 손실",
], title_color=ACCENT_RED, border=ACCENT_RED)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.5), "FN (스팸 -> 정상)", [
    "스팸이 받은편지함에 들어옴",
    "불편하지만 덜 심각",
    "",
    "예: 광고 메일이 보임",
    "-> 사용자가 수동으로 삭제 가능",
], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)
add_card(s, Inches(0.8), Inches(4.7), Inches(11.2), Inches(1.5), "가중 정확도 (Metsis 2006)", [
    "WAcc = (lambda * TN + TP) / (lambda * (TN + FP) + (TP + FN))",
    "",
    "lambda > 1: FP의 상대적 비용 -- FP에 더 높은 패널티 부여",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# SECTION 7: 라플라스 스무딩과 로그 확률 트릭
# ============================================================
section_divider("라플라스 스무딩과 로그 확률", "Zero Probability & Underflow 문제의 해결", "7")

# 슬라이드 34: 제로 확률 문제
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "7", "제로 확률 문제 (Zero Probability Problem)")
add_card(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(2.0), "문제 상황", [
    '학습 데이터에서 "scholarship"이 스팸 클래스에서 한 번도 등장하지 않으면:',
    '',
    '  P("scholarship" | spam) = 0',
    "",
    "나이브 베이즈는 확률을 곱셈으로 결합하므로:",
    '  P(spam | x1, x2, ..., "scholarship", ...) = 0',
], title_color=ACCENT_RED, border=ACCENT_RED)
add_text(s, Inches(0.8), Inches(4.2), Inches(11), Inches(0.8),
         "하나라도 0이면 전체 결과가 0! -> 다른 모든 단어의 정보가 무시됨",
         font_size=20, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)

# 슬라이드 35: 라플라스 스무딩
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "7", "라플라스 스무딩 (Laplace Smoothing)")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5), "스무딩 수식", [
    "P(xi|C) = (count(xi,C) + alpha)",
    "        / (count(C) + alpha * |V|)",
    "",
    "count(xi,C): 클래스 C에서 단어 xi 출현 횟수",
    "count(C): 클래스 C의 전체 단어 수",
    "alpha: 스무딩 파라미터",
    "|V|: 전체 어휘(Vocabulary)의 크기",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
headers2 = ["alpha 값", "명칭", "효과"]
rows2 = [
    ["alpha = 1", "라플라스 스무딩", "표준 (기본값)"],
    ["0 < alpha < 1", "리드스톤 스무딩", "약한 스무딩"],
    ["alpha > 1", "강한 스무딩", "분포를 더 균일하게"],
]
add_table_slide(s, headers2, rows2, Inches(6.8), Inches(2.0), [2.0, 2.0, 2.0],
                row_height=0.5, font_size=12)
add_code_block(s, Inches(6.8), Inches(3.8), Inches(5.5), Inches(1.5), [
    "model = MultinomialNB(alpha=1.0)   # 라플라스",
    "model = MultinomialNB(alpha=0.5)   # 리드스톤",
    "model = MultinomialNB(alpha=0.01)  # 약한 스무딩",
])

# 슬라이드 36: 로그 확률 트릭
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "7", "로그 확률 트릭 (Log Probability Trick)")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5), "문제: 언더플로", [
    "많은 확률값(0~1)을 곱하면 -> 0에 수렴",
    "",
    "P(C) * PROD P(xi|C)",
    "= 0.001 * 0.01 * 0.005 * ... -> 0",
    "",
    "컴퓨터의 부동소수점 정밀도 한계",
], title_color=ACCENT_RED, border=ACCENT_RED)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.5), "해결: 로그 공간으로", [
    "곱셈 -> 로그 공간에서 덧셈으로 변환",
    "",
    "log P(C|x) ~ log P(C) + SUM log P(xi|C)",
    "",
    "C_hat = argmax_C [log P(C)",
    "                 + SUM_i log P(xi|C)]",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_card(s, Inches(0.8), Inches(4.7), Inches(11.2), Inches(1.5), "장점 3가지", [
    "1. 수치적 안정성 확보 (언더플로 방지)  |  2. 곱셈 -> 덧셈으로 계산 효율 향상  |  3. argmax 보존 (로그는 단조증가함수)",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# SECTION 8: 생성 모델 vs 판별 모델
# ============================================================
section_divider("생성 모델 vs 판별 모델", "Naive Bayes vs Logistic Regression -- Ng & Jordan (2001)", "8")

# 슬라이드 38: 생성 vs 판별 비교표
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "8", "생성 모델 vs 판별 모델")
headers = ["관점", "생성 모델 (Generative)", "판별 모델 (Discriminative)"]
rows = [
    ["대표 알고리즘", "나이브 베이즈", "로지스틱 회귀"],
    ["모델링 대상", "P(X,Y) = P(X|Y)*P(Y)", "P(Y|X) 직접 모델링"],
    ["학습 방식", "클래스별 데이터 분포 학습", "결정 경계를 직접 학습"],
    ["가정", "특성의 분포를 가정", "분포 가정 불필요"],
    ["소량 데이터", "상대적으로 강함", "과적합 위험"],
    ["대량 데이터", "점근적으로 열등", "점근적으로 우수"],
]
add_table_slide(s, headers, rows, Inches(0.5), Inches(2.0), [2.5, 5.0, 5.0],
                row_height=0.55, font_size=13)

# 슬라이드 39: 수학적 연결
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "8", "수학적 연결 -- Ng & Jordan (2001)")
add_text(s, Inches(0.8), Inches(1.9), Inches(11), Inches(0.4),
         "나이브 베이즈와 로지스틱 회귀는 동일한 모델 패밀리에 속한다!", font_size=18, color=ACCENT_CYAN, bold=True)
add_card(s, Inches(0.8), Inches(2.6), Inches(11.2), Inches(2.5), "로그 오즈(Log-Odds) 전개", [
    "이진 분류에서 나이브 베이즈의 로그 오즈:",
    "",
    "log P(Y=1|x)/P(Y=0|x) = log P(Y=1)/P(Y=0)",
    "                       + SUM_i log P(xi|Y=1)/P(xi|Y=0)",
    "",
    "이것은 로지스틱 회귀의 선형 모델 w^T*x + b 와 동일한 형태!",
    "",
    "핵심 차이: 파라미터를 추정하는 방식이 다르다",
    "  NB: 생성적 -- P(X|Y)와 P(Y)를 각각 추정 후 베이즈 정리 적용",
    "  LR: 판별적 -- P(Y|X)를 직접 최적화 (조건부 우도 최대화)",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 슬라이드 40: 실용적 선택 기준
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "8", "실용적 선택 기준")
headers = ["상황", "추천"]
rows = [
    ["학습 데이터가 적을 때", "나이브 베이즈 (수렴이 빠름)"],
    ["학습 데이터가 충분할 때", "로지스틱 회귀 (점근 성능 우수)"],
    ["특성 간 독립성이 높을 때", "나이브 베이즈"],
    ["특성 간 상관관계가 강할 때", "로지스틱 회귀"],
    ["해석 가능성이 중요할 때", "둘 다 해석 용이"],
    ["속도가 매우 중요할 때", "나이브 베이즈"],
]
add_table_slide(s, headers, rows, Inches(1.5), Inches(2.0), [5.0, 5.0],
                row_height=0.55, font_size=14)
add_text(s, Inches(0.8), Inches(5.6), Inches(11), Inches(0.5),
         "Ng & Jordan(2001): NB는 O(log n) 샘플에서 수렴, LR은 O(n) 샘플 필요",
         font_size=16, color=ACCENT_ORANGE, bold=True)

# ============================================================
# SECTION 9: 논문 리뷰
# ============================================================
section_divider("논문 리뷰", "Rish(2001), McCallum(1998), Rennie(2003)", "9")

# 슬라이드 42: Rish (2001)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "9", "Rish (2001) -- 나이브 베이즈의 경험적 연구")
add_text(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.3),
         "IJCAI 2001 Workshop | An Empirical Study of the Naive Bayes Classifier", font_size=14, color=DARK_GRAY)
add_card(s, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.0), "주요 발견", [
    "1. 의존성의 방향이 중요하다:",
    "   양의 상관 + 음의 상관 혼재 -> 상쇄",
    "   -> 성능 유지",
    "",
    "2. 균일한 의존성은 해롭지 않다:",
    "   모든 특성 쌍이 비슷한 수준의 의존성",
    "   -> 순위 보존",
    "",
    "3. 분류 vs 확률 추정의 구분:",
    "   확률 추정 부정확해도 순위 정확 -> OK",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.2), Inches(5.5), Inches(1.5), "핵심 기여", [
    "조건부 독립 가정 위반 시에도 왜 잘 작동하는지",
    "경험적으로 분석한 최초의 체계적 연구",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_card(s, Inches(6.8), Inches(4.0), Inches(5.5), Inches(1.2), "한계", [
    "소규모 합성 데이터 위주의 실험",
    "연속형 특성 분석 부족",
], title_color=ACCENT_RED, border=ACCENT_RED)

# 슬라이드 43: McCallum & Nigam (1998)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "9", "McCallum & Nigam (1998) -- 이벤트 모델 비교")
add_text(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.3),
         "AAAI-98 Workshop | A Comparison of Event Models for Naive Bayes Text Classification",
         font_size=14, color=DARK_GRAY)
headers = ["어휘 크기", "베르누이 모델", "다항 모델", "차이"]
rows = [
    ["100", "0.68", "0.65", "-0.03"],
    ["1,000", "0.78", "0.82", "+0.04"],
    ["10,000", "0.81", "0.87", "+0.06"],
    ["전체", "0.82", "0.89", "+0.07"],
]
add_table_slide(s, headers, rows, Inches(0.8), Inches(2.3), [2.5, 2.5, 2.5, 2.5],
                row_height=0.5, font_size=13)
add_card(s, Inches(0.8), Inches(4.7), Inches(11.2), Inches(1.6), "핵심 결론", [
    "어휘 크기가 클수록 다항 모델(MultinomialNB)이 베르누이 모델보다 일관되게 우수",
    "이 논문이 sklearn에서 MultinomialNB와 BernoulliNB를 별도 클래스로 제공하는 직접적 근거",
    "20 Newsgroups 데이터셋에서 실험",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# 슬라이드 44: Rennie et al. (2003) -- 5가지 문제점
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "9", "Rennie et al. (2003) -- 5가지 잘못된 가정")
add_text(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.3),
         "ICML 2003 | Tackling the Poor Assumptions of Naive Bayes Text Classifiers",
         font_size=14, color=DARK_GRAY)
headers = ["문제", "설명", "해결"]
rows = [
    ["조건부 독립", "단어 간 의존성 무시", "CNB로 완화"],
    ["균일한 특성 중요도", "모든 단어가 동등", "TF-IDF 가중치"],
    ["클래스 불균형", "다수 클래스 편향", "보완 클래스 사용"],
    ["문서 길이 영향", "긴 문서의 우도 왜곡", "L2 정규화"],
    ["버스티니스", "단어 반복 출현 무시", "로그 TF 변환"],
]
add_table_slide(s, headers, rows, Inches(0.8), Inches(2.3), [3.0, 4.0, 3.5],
                row_height=0.5, font_size=13)

# 슬라이드 45: Rennie et al. (2003) -- 실험 결과
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "9", "Rennie et al. (2003) -- 실험 결과")
add_text(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.3),
         "20 Newsgroups, 20 클래스 분류", font_size=14, color=DARK_GRAY)
headers = ["방법", "정확도"]
rows = [
    ["MultinomialNB (기본)", "0.774"],
    ["CNB (ComplementNB)", "0.819"],
    ["WCNB (전체 개선 적용)", "0.847"],
    ["SVM (선형)", "0.864"],
]
add_table_slide(s, headers, rows, Inches(2.5), Inches(2.2), [5.0, 3.0],
                row_height=0.55, font_size=15)
add_text(s, Inches(0.8), Inches(4.8), Inches(11), Inches(0.8),
         "ComplementNB는 MultinomialNB 대비 ~4.5%p 향상, SVM에 근접\n이 논문이 sklearn에 ComplementNB 클래스가 추가되는 직접적 근거",
         font_size=17, color=ACCENT_GREEN, bold=True)

# ============================================================
# SECTION 10: 실습 코드
# ============================================================
section_divider("실습 코드", "GaussianNB 스크래치 | 텍스트 분류 | 스팸 필터", "10")

# 슬라이드 47: GaussianNB 스크래치 -- 클래스 구조
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "GaussianNB 스크래치 -- 클래스 구조 (01_naive_bayes_scratch.py)")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(4.5), [
    "class GaussianNBScratch:",
    '    """가우시안 나이브 베이즈 분류기 (밑바닥 구현)"""',
    "",
    "    def __init__(self):",
    "        self.classes = None      # 클래스 레이블 배열",
    "        self.priors = None       # 각 클래스의 사전확률 P(c)",
    "        self.means = None        # 각 클래스, 각 특성의 평균 (mu)",
    "        self.variances = None    # 각 클래스, 각 특성의 분산 (sigma^2)",
    "",
    "    def fit(self, X, y):",
    '        """학습 데이터로부터 사전확률, 평균, 분산을 추정"""',
    "        self.classes = np.unique(y)",
    "        for idx, c in enumerate(self.classes):",
    "            X_c = X[y == c]",
    "            self.priors[idx] = X_c.shape[0] / X.shape[0]  # P(c) = Nc/N",
    "            self.means[idx] = X_c.mean(axis=0)            # 클래스별 평균",
    "            self.variances[idx] = X_c.var(axis=0) + 1e-9  # 분산 + epsilon",
], font_size=12)

# 슬라이드 48: GaussianNB -- 로그 우도 계산
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "GaussianNB -- 로그 우도 계산")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(2.5), [
    "def _gaussian_log_likelihood(self, x, mean, var):",
    '    """가우시안 분포의 로그 확률밀도를 계산"""',
    "    # log P(x|mu,sigma^2) = -0.5*log(2*pi*sigma^2)",
    "    #                      - (x-mu)^2 / (2*sigma^2)",
    "    log_likelihood = -0.5 * np.log(2*np.pi*var) \\",
    "                     - ((x - mean)**2) / (2*var)",
    "    return np.sum(log_likelihood)  # 독립가정 -> 합",
], font_size=13)
add_code_block(s, Inches(0.8), Inches(4.6), Inches(11.2), Inches(2.0), [
    "def _compute_log_posterior(self, x):",
    '    """각 클래스의 로그 사후확률 계산"""',
    "    for idx in range(self.n_classes):",
    "        log_prior = np.log(self.priors[idx])        # log P(c)",
    "        log_lik = self._gaussian_log_likelihood(     # sum log P(xi|c)",
    "            x, self.means[idx], self.variances[idx])",
    "        log_posteriors[idx] = log_prior + log_lik    # 합산",
], font_size=13)

# 슬라이드 49: GaussianNB -- predict_proba (log-sum-exp)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "GaussianNB -- predict_proba (Log-Sum-Exp 트릭)")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(3.0), [
    "def predict_proba(self, X):",
    '    """로그 사후확률을 softmax로 정규화하여 확률 변환"""',
    "    for x in X:",
    "        log_posteriors = self._compute_log_posterior(x)",
    "",
    "        # 수치 안정성을 위한 log-sum-exp trick",
    "        max_log = np.max(log_posteriors)",
    "        log_posteriors_shifted = log_posteriors - max_log",
    "        posteriors = np.exp(log_posteriors_shifted)",
    "        posteriors /= np.sum(posteriors)  # 정규화 (합=1)",
], font_size=13)
add_card(s, Inches(0.8), Inches(5.2), Inches(11.2), Inches(1.2), "핵심 포인트", [
    "fit(): 클래스별 P(C), mu, sigma^2 계산  |  predict(): argmax(log posterior)  |  predict_proba(): log-sum-exp -> softmax 정규화",
], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# 슬라이드 50: sklearn 비교 결과
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "GaussianNB -- sklearn과 비교 결과")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.0), "Iris 데이터셋", [
    "직접 구현 정확도: ~0.9778",
    "sklearn 정확도:   ~0.9778",
    "예측 일치율:      ~1.0000",
    "",
    "두 구현이 완전히 동일한 결과!",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.0), "Wine 데이터셋", [
    "직접 구현 정확도: ~0.9630",
    "sklearn 정확도:   ~0.9630",
    "",
    "13개 특성의 다차원 데이터에서도",
    "동일한 결과를 재현",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_text(s, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
         "밑바닥 구현을 통해 GaussianNB의 동작 원리를 완전히 이해",
         font_size=18, color=ACCENT_CYAN, bold=True)
add_code_block(s, Inches(0.8), Inches(5.0), Inches(11.2), Inches(1.5), [
    "# 메인 실행 예시",
    "my_gnb = GaussianNBScratch()",
    "my_gnb.fit(X_train, y_train)            # 학습",
    "my_pred = my_gnb.predict(X_test)          # 예측",
    "my_proba = my_gnb.predict_proba(X_test)   # 확률 추정",
])

# 슬라이드 51: 텍스트 분류 비교 -- 실험 설계
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "텍스트 분류 비교 (02_text_classification_comparison.py)")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.8), "실험 설계", [
    "데이터: 20 Newsgroups (4개 카테고리)",
    "  - rec.sport.baseball (야구)",
    "  - sci.med (의학)",
    "  - comp.graphics (그래픽)",
    "  - talk.politics.misc (정치)",
    "",
    "3가지 NB변형 x 2가지 벡터화기 = 6조합",
    "Alpha 튜닝 + 어휘 크기별 비교",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_code_block(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.8), [
    "# 실험 코드 구조",
    "vectorizers = {",
    "  'CountVec': CountVectorizer(...),",
    "  'TfidfVec': TfidfVectorizer(...),",
    "}",
    "nb_models = {",
    "  'MultinomialNB': MultinomialNB(alpha=1.0),",
    "  'BernoulliNB': BernoulliNB(alpha=1.0),",
    "  'ComplementNB': ComplementNB(alpha=1.0),",
    "}",
], font_size=12)
add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
         "실험 결과: ComplementNB + TfidfVectorizer가 일관되게 최우수",
         font_size=17, color=ACCENT_GREEN, bold=True)

# 슬라이드 52: Alpha 튜닝 결과
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "Alpha 하이퍼파라미터 튜닝 결과")
add_card(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(2.5), "Alpha 값에 따른 성능 변화", [
    "alpha = 0.001 ~ 10.0 범위에서 탐색 (TfidfVectorizer 사용)",
    "",
    "MultinomialNB: alpha=0.1~0.5 범위에서 최적",
    "BernoulliNB:   alpha=0.05~0.3 범위에서 최적",
    "ComplementNB:  alpha에 상대적으로 덜 민감, 전반적으로 안정적",
    "",
    "실무 팁: alpha 값은 교차 검증으로 최적화 (기본값 1.0이 항상 최적은 아님)",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(0.8), Inches(4.7), Inches(11.2), Inches(1.5), "어휘 크기별 결론 (McCallum & Nigam 1998 재현)", [
    "어휘 100: BernoulliNB = MultinomialNB  |  어휘 1K+: MultinomialNB > BernoulliNB  |  어휘 10K+: ComplementNB > MultinomialNB >> BernoulliNB",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# 슬라이드 53: 스팸 필터 -- 전처리 파이프라인
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "스팸 필터 (03_spam_filter.py) -- 전처리")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(2.8), [
    "def preprocess_text(text):",
    '    """텍스트 전처리: 소문자, 특수문자 제거, 숫자 처리"""',
    "    text = text.lower()",
    "    text = re.sub(r'\\$[\\d,]+', 'MONEY', text)  # $500 -> MONEY",
    "    text = re.sub(r'\\d+%', 'PERCENT', text)     # 50% -> PERCENT",
    "    text = re.sub(r'\\d+', 'NUM', text)           # 1000 -> NUM",
    "    text = re.sub(r'[^\\w\\s]', ' ', text)         # 특수문자 제거",
    "    text = re.sub(r'\\s+', ' ', text).strip()     # 다중 공백 정리",
    "    return text",
], font_size=13)
add_card(s, Inches(0.8), Inches(5.0), Inches(11.2), Inches(1.3), "전처리 핵심 포인트", [
    "금액($500) -> MONEY 토큰, 퍼센트(50%) -> PERCENT 토큰으로 대체하여 스팸 패턴 보존",
    "실제 스팸에서 금액/퍼센트 언급은 매우 강력한 스팸 지표",
], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# 슬라이드 54: 스팸 필터 -- TF-IDF 설정
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "스팸 필터 -- TF-IDF 설정과 모델 학습")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(3.0), [
    "# TF-IDF 벡터화 (Rennie 2003 권장 설정)",
    "tfidf = TfidfVectorizer(",
    "    max_features=5000,         # 상위 5000개 단어",
    "    stop_words='english',      # 불용어 제거",
    "    ngram_range=(1, 2),        # 유니그램 + 바이그램",
    "    min_df=2,                  # 최소 2개 문서에 출현",
    "    sublinear_tf=True          # 로그 TF 변환 (Rennie 2003)",
    ")",
    "",
    "# MultinomialNB 학습",
    "mnb = MultinomialNB(alpha=0.1)",
    "mnb.fit(X_train_tfidf, y_train)",
], font_size=13)
add_card(s, Inches(0.8), Inches(5.2), Inches(11.2), Inches(1.2), "핵심 설정 해설", [
    'sublinear_tf=True: Rennie(2003) 권장 로그 TF 변환  |  ngram_range=(1,2): "free gift" 같은 복합 패턴 포착  |  alpha=0.1: 약한 스무딩이 종종 최적',
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 슬라이드 55: 스팸 필터 -- 특성 중요도
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "스팸 필터 -- 특성 중요도 분석")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5), [
    "def analyze_feature_importance(model, vec):",
    "    log_probs = model.feature_log_prob_",
    "    log_prob_ham = log_probs[0]",
    "    log_prob_spam = log_probs[1]",
    "",
    "    # 스팸 판별력",
    "    spam_importance = (",
    "        log_prob_spam - log_prob_ham)",
], font_size=12)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(1.3), "스팸 지표 단어", [
    '"free", "win", "prize", "money", "claim"',
    '"congratulations", "guaranteed", "offer"',
], title_color=ACCENT_RED, border=ACCENT_RED)
add_card(s, Inches(6.8), Inches(3.5), Inches(5.5), Inches(1.3), "정상 지표 단어", [
    '"meeting", "please", "team", "review"',
    '"project", "schedule", "report"',
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(0.8), Inches(5.0), Inches(11.2), Inches(1.3), "판별 원리", [
    "log P(w|spam) - log P(w|ham) > 0 이면 스팸 지표 단어",
    "log P(w|spam) - log P(w|ham) < 0 이면 정상 지표 단어",
    "이 차이값이 클수록 해당 단어의 판별력이 강함",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# 슬라이드 56: 스팸 필터 -- NB 변형 비교 (Metsis 2006 재현)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "스팸 필터 -- NB 변형 비교 (Metsis 2006 재현)")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5), [
    "models = {",
    "  'MultinomialNB': MultinomialNB(alpha=0.1),",
    "  'BernoulliNB': BernoulliNB(alpha=0.1),",
    "  'ComplementNB': ComplementNB(alpha=0.1),",
    "  'MNB+Boolean': MultinomialNB(alpha=0.1),",
    "}",
    "# MNB+Boolean: binary=True로 이진화",
    "# -> Metsis(2006) MNBB 최우수 변형",
], font_size=12)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.5), "Metsis(2006) MNBB 변형", [
    "TfidfVectorizer에서 binary=True 옵션 사용",
    "단어 빈도를 이진(0/1)으로 변환한 후",
    "MultinomialNB의 확률 모델 적용",
    "",
    "스팸 판별: 빈도보다 존재 여부가 중요",
    "BUT MultinomialNB의 모델이 더 적합",
    "-> BernoulliNB의 모델보다 우수",
], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# 슬라이드 57: 스팸 필터 -- 예측 데모
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "10", "스팸 필터 -- 새로운 이메일 예측 데모")
add_code_block(s, Inches(0.8), Inches(1.9), Inches(11.2), Inches(4.5), [
    "new_emails = [",
    '    "Congratulations! You won a free trip to Hawaii! Click here",',
    '    "Hi, can we schedule a meeting for next Tuesday?",',
    '    "URGENT: Your account will be suspended unless you verify",',
    '    "Thanks for the update. I will review the document",',
    "]",
    "",
    "new_clean = [preprocess_text(e) for e in new_emails]",
    "new_tfidf = tfidf.transform(new_clean)",
    "new_pred = mnb.predict(new_tfidf)",
    "new_proba = mnb.predict_proba(new_tfidf)",
    "",
    '# 결과: "Congratulations..." -> 스팸 (P=0.95+)',
    '#       "Hi, can we..."       -> 정상 (P=0.02)',
    '#       "URGENT: Your..."     -> 스팸 (P=0.88+)',
    '#       "Thanks for..."       -> 정상 (P=0.05)',
], font_size=12)

# ============================================================
# SECTION 11: 응용 사례
# ============================================================
section_divider("응용 사례", "감정분석 | 문서분류 | 의료진단", "11")

# 슬라이드 59: 감정분석
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "11", "감정분석 (Sentiment Analysis)")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.0), "감정분석 종류", [
    "이진 분류: 긍정 / 부정 (영화 리뷰, 상품 후기)",
    "다중 분류: 긍정 / 중립 / 부정 (SNS 여론)",
    "감성 점수: 연속적 수치 (-1.0 ~ +1.0)",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.0), "나이브 베이즈의 장점", [
    "감성 키워드(great, terrible 등) 빈도로 분류",
    "빠른 학습/예측 -> 대규모 텍스트 실시간 처리",
    "적은 학습 데이터로도 상당한 성능 달성",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)
add_card(s, Inches(0.8), Inches(4.2), Inches(11.2), Inches(1.5), "활용 분야", [
    "마케팅: 소비자 반응 분석  |  금융: 뉴스 기반 주가 예측  |  고객 서비스: 불만 고객 자동 감지  |  정치: 선거 여론 분석",
], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# 슬라이드 60: 문서 분류
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "11", "문서 분류 (Document Classification)")
headers = ["분야", "설명"]
rows = [
    ["뉴스 분류", "정치, 경제, 스포츠, 문화 등"],
    ["이메일 분류", "업무, 개인, 프로모션, 소셜"],
    ["고객 문의 분류", "배송, 반품, 결제, 기술 지원"],
    ["법률 문서 분류", "계약서, 소송 자료, 특허"],
]
add_table_slide(s, headers, rows, Inches(1.5), Inches(2.0), [4.0, 6.0],
                row_height=0.55, font_size=14)
add_card(s, Inches(0.8), Inches(4.5), Inches(11.2), Inches(1.8), "나이브 베이즈가 문서 분류에 적합한 이유", [
    "1. 수만~수십만 개의 고차원 특성 공간을 효율적으로 처리",
    "2. 대량의 문서를 실시간으로 분류 가능",
    "3. 다중 클래스 지원이 자연스러움",
    "4. 카테고리당 소수의 학습 데이터로도 합리적 결과",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# 슬라이드 61: 의료진단
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "11", "의료진단 (Medical Diagnosis)")
add_card(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5), "증상 기반 질병 진단", [
    "환자의 증상 벡터 x = (x1, x2, ..., xn)",
    "",
    "P(D | x1,...,xn) ~ P(D) * PROD P(xi|D)",
    "",
    "각 증상이 질병 확률에 독립적으로 기여",
    "여러 질병 후보를 확률 순으로 제시",
], title_color=ACCENT_BLUE, border=ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.5), Inches(2.5), "장점 및 사례", [
    "소규모 의료 데이터에서의 강건성",
    "각 증상의 기여도를 확률로 해석 (설명 가능)",
    "",
    "실제 연구 사례:",
    "  유방암 진단, 심장병 예측, 당뇨병 분류",
    "  SVM/신경망과 비슷한 성능 + 해석 용이",
], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# SECTION 12: 핵심 요약
# ============================================================
section_divider("핵심 요약", "Key Concepts Summary + 복습 질문", "12")

# 슬라이드 63: 핵심 개념 정리 1
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "12", "핵심 개념 정리 (1/2)")
headers = ["개념", "설명"]
rows = [
    ["베이즈 정리", "P(A|B) = P(B|A)*P(A) / P(B)"],
    ["나이브 가정", "모든 특성이 클래스에 대해 조건부 독립"],
    ["GaussianNB", "연속형 데이터, 정규분포 가정"],
    ["MultinomialNB", "텍스트 분류 최적, 단어 빈도 기반"],
    ["BernoulliNB", "단어 존재/부재, 소규모 어휘에서 경쟁력"],
    ["ComplementNB", "Rennie(2003), 보완 클래스, SVM에 근접"],
]
add_table_slide(s, headers, rows, Inches(0.8), Inches(2.0), [3.0, 9.0],
                row_height=0.55, font_size=14)

# 슬라이드 64: 핵심 개념 정리 2
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "12", "핵심 개념 정리 (2/2)")
headers = ["개념", "설명"]
rows = [
    ["라플라스 스무딩", "제로 확률 문제 해결, alpha 파라미터"],
    ["로그 확률 트릭", "언더플로 방지, 곱셈 -> 덧셈 변환"],
    ["CountVectorizer", "텍스트 -> 단어 출현 빈도 벡터"],
    ["TF-IDF", "단어 중요도를 반영한 가중 벡터화"],
    ["생성 모델", "NB: P(X|Y)*P(Y)를 모델링"],
    ["판별 모델", "LR: P(Y|X)를 직접 모델링"],
]
add_table_slide(s, headers, rows, Inches(0.8), Inches(2.0), [3.0, 9.0],
                row_height=0.55, font_size=14)

# 슬라이드 65: 전체 워크플로우
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "12", "전체 워크플로우")
steps = [
    ("텍스트 데이터", ACCENT_BLUE),
    ("전처리 (특수기호 제거, 소문자, 불용어)", ACCENT_CYAN),
    ("벡터화 (CountVectorizer / TfidfVectorizer)", ACCENT_GREEN),
    ("모델링 (MultinomialNB / ComplementNB + alpha 튜닝)", ACCENT_ORANGE),
    ("평가 (정확도, 정밀도, 재현율, F1, 혼동행렬)", ACCENT_PURPLE),
]
for i, (step, clr) in enumerate(steps):
    y = Inches(2.0) + Inches(1.0) * i
    add_shape(s, Inches(2.0), y, Inches(9.0), Inches(0.7), CARD_BG, clr, radius=True)
    add_text(s, Inches(2.3), y + Inches(0.1), Inches(8.5), Inches(0.5),
             step, font_size=18, color=WHITE, bold=True)
    if i < len(steps) - 1:
        add_text(s, Inches(6.2), y + Inches(0.65), Inches(1), Inches(0.35),
                 "v", font_size=20, color=clr, bold=True, align=PP_ALIGN.CENTER)

# 슬라이드 66: 알고리즘 선택 가이드
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "12", "알고리즘 선택 가이드")
headers = ["상황", "추천"]
rows = [
    ["텍스트 분류 (빈도 기반)", "MultinomialNB + CountVectorizer"],
    ["텍스트 분류 (가중 빈도)", "MultinomialNB + TfidfVectorizer"],
    ["텍스트 분류 (존재 유무)", "BernoulliNB"],
    ["불균형 텍스트 분류", "ComplementNB"],
    ["수치형 특성 분류", "GaussianNB"],
    ["대규모 데이터 빠른 분류", "나이브 베이즈 전 종류"],
    ["특성 간 상관관계 강한 경우", "로지스틱 회귀 또는 SVM 권장"],
]
add_table_slide(s, headers, rows, Inches(1.0), Inches(2.0), [5.5, 5.5],
                row_height=0.53, font_size=14)

# 슬라이드 67: 복습 질문 1-5
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "12", "복습 질문 (1~5)")
questions_1 = [
    "Q1. 베이즈 정리의 각 항(사전확률, 우도, 사후확률, 증거)을",
    "    의료진단 사례를 들어 설명하시오.",
    "",
    "Q2. 나이브 베이즈에서 '나이브'의 의미는 무엇이며, 이 가정이",
    "    위반되어도 잘 작동하는 이유를 Zhang(2004)를 인용하여 설명하시오.",
    "",
    "Q3. 라플라스 스무딩이 필요한 이유와 수식을 설명하고,",
    "    alpha=1과 alpha=0.1의 차이를 논하시오.",
    "",
    "Q4. GaussianNB, MultinomialNB, BernoulliNB, ComplementNB의",
    "    차이점과 각각의 적합한 데이터 유형을 비교하시오.",
    "",
    "Q5. CountVectorizer와 TfidfVectorizer의 차이점을 수식과 함께",
    "    설명하고, TF-IDF가 더 우수한 상황을 논하시오.",
]
add_bullet_list(s, Inches(1.0), Inches(1.9), Inches(11), Inches(5),
                questions_1, font_size=15, color=LIGHT_GRAY, spacing=Pt(3))

# 슬라이드 68: 복습 질문 6-10
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "12", "복습 질문 (6~10)")
questions_2 = [
    "Q6. Rennie et al.(2003)이 식별한 나이브 베이즈의 5가지 잘못된 가정을",
    "    나열하고, 각각의 해결 방안을 설명하시오.",
    "",
    "Q7. 스팸 필터링에서 FP와 FN 중 어느 것이 더 심각한 문제인지 논하고,",
    "    혼동행렬의 TP, TN, FP, FN을 스팸 필터링 맥락에서 설명하시오.",
    "",
    "Q8. 나이브 베이즈(생성 모델)와 로지스틱 회귀(판별 모델)의 차이를 설명하고,",
    "    소량 데이터 vs 대량 데이터에서의 성능 특성을 비교하시오.",
    "",
    "Q9. 로그 확률 트릭이 필요한 이유를 수치 예제를 들어 설명하고,",
    "    argmax 연산이 로그 변환 후에도 보존되는 이유를 서술하시오.",
    "",
    "Q10. McCallum & Nigam(1998)에서 어휘가 커질수록 MultinomialNB가",
    "     BernoulliNB보다 우수한 이유와, Metsis(2006)의 MNBB 변형이",
    "     최우수인 이유를 논하시오.",
]
add_bullet_list(s, Inches(1.0), Inches(1.9), Inches(11), Inches(5),
                questions_2, font_size=15, color=LIGHT_GRAY, spacing=Pt(3))

# 슬라이드 69: 참고 문헌
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, "", "참고 문헌 (References)")
refs = [
    "1. Rish, I. (2001). An Empirical Study of the Naive Bayes Classifier.",
    "   IJCAI Workshop.",
    "",
    "2. McCallum, A. & Nigam, K. (1998). A Comparison of Event Models",
    "   for Naive Bayes Text Classification. AAAI Workshop.",
    "",
    "3. Rennie, J.D.M. et al. (2003). Tackling the Poor Assumptions of",
    "   Naive Bayes Text Classifiers. ICML, pp. 616-623.",
    "",
    "4. Zhang, H. (2004). The Optimality of Naive Bayes. FLAIRS Conference.",
    "",
    "5. Metsis, V. et al. (2006). Spam Filtering with Naive Bayes --",
    "   Which Naive Bayes? CEAS.",
    "",
    "6. Ng, A.Y. & Jordan, M.I. (2001). On Discriminative vs. Generative",
    "   Classifiers. NIPS.",
]
add_bullet_list(s, Inches(1.0), Inches(1.9), Inches(11), Inches(5),
                refs, font_size=14, color=LIGHT_GRAY, spacing=Pt(2))

# 슬라이드 70: 마무리
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.0), prs.slide_width, Inches(0.5),
         "CHAPTER 7 COMPLETE", font_size=20, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.0), Inches(2.7), Inches(3.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(3.0), prs.slide_width, Inches(1.0),
         "Thank You", font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.5),
         "나이브 베이즈: 단순하지만 강력한 확률적 분류기", font_size=20, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.2), prs.slide_width, Inches(0.8),
         '"나이브 가정이 깨져도 작동하는 것은, 분류에는 정확한 확률값이 아니라\n올바른 순위만 필요하기 때문이다." -- Rish (2001)',
         font_size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 저장
# ============================================================
out = os.path.join(os.path.dirname(__file__), "7장_나이브베이즈_강의PPT_확장.pptx")
prs.save(out)
print(f"Created: {out}")
print(f"Total slides: {len(prs.slides)}")
