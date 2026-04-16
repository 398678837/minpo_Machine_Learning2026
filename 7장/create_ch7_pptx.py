# -*- coding: utf-8 -*-
"""
Chapter 7: 나이브 베이즈 - 스팸 여부 판단하기
Dark Academia Style PPTX Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Dark Academia Color Palette ──
BG_COLOR = RGBColor(0x1A, 0x12, 0x08)
TITLE_GOLD = RGBColor(0xC9, 0xA8, 0x4C)
BODY_PARCHMENT = RGBColor(0xD4, 0xBF, 0x9A)
BORDER_GOLD = RGBColor(0x3D, 0x2E, 0x10)
ACCENT_GOLD = RGBColor(0x8A, 0x73, 0x40)
MUTED_TEXT = RGBColor(0xA0, 0x90, 0x70)
CODE_BG = RGBColor(0x12, 0x0E, 0x06)
CODE_TEXT = RGBColor(0xE8, 0xD5, 0xA0)
HIGHLIGHT = RGBColor(0xE8, 0xC4, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG2 = RGBColor(0x0E, 0x0A, 0x05)
ACCENT_RED = RGBColor(0xC0, 0x5C, 0x3C)
ACCENT_GREEN = RGBColor(0x6B, 0x9E, 0x5C)
ACCENT_BLUE = RGBColor(0x5C, 0x7E, 0xA0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_border_frame(slide, outer_margin=Inches(0.3), inner_margin=Inches(0.45)):
    """Add double inset border - Dark Academia signature"""
    # Outer border
    shapes = slide.shapes
    outer = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        outer_margin, outer_margin,
        W - 2 * outer_margin, H - 2 * outer_margin
    )
    outer.fill.background()
    outer.line.color.rgb = BORDER_GOLD
    outer.line.width = Pt(1.5)

    # Inner border
    inner = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        inner_margin, inner_margin,
        W - 2 * inner_margin, H - 2 * inner_margin
    )
    inner.fill.background()
    inner.line.color.rgb = ACCENT_GOLD
    inner.line.width = Pt(0.75)


def add_decorative_line(slide, top, left=Inches(1.5), width=Inches(10.333)):
    """Thin gold horizontal rule"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_GOLD
    line.line.fill.background()


def add_textbox(slide, left, top, width, height, text, font_name="Georgia",
                font_size=Pt(16), color=BODY_PARCHMENT, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(4)
    if line_spacing:
        p.line_spacing = Pt(int(font_size.pt * line_spacing))
    return txBox


def add_multi_text(slide, left, top, width, height, texts, default_size=Pt(15)):
    """texts: list of (text, font_size, color, bold, italic, font_name)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(texts):
        text = item[0]
        size = item[1] if len(item) > 1 else default_size
        color = item[2] if len(item) > 2 else BODY_PARCHMENT
        bold = item[3] if len(item) > 3 else False
        italic = item[4] if len(item) > 4 else False
        fname = item[5] if len(item) > 5 else "Georgia"
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = fname
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.space_after = Pt(4)
        p.line_spacing = Pt(int(size.pt * 1.5))
    return txBox


def add_code_block(slide, left, top, width, height, code_text):
    """Dark code block with monospace font"""
    # Background rectangle
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = CODE_BG
    rect.line.color.rgb = ACCENT_GOLD
    rect.line.width = Pt(0.5)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(13)
        p.font.color.rgb = CODE_TEXT
        p.space_after = Pt(1)
        p.line_spacing = Pt(18)
    return rect


def add_accent_box(slide, left, top, width, height, text, box_color=BORDER_GOLD,
                   text_color=TITLE_GOLD, font_size=Pt(14)):
    """Accent box with colored border"""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x20, 0x18, 0x0C)
    rect.line.color.rgb = box_color
    rect.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = "Georgia"
        p.font.size = font_size
        p.font.color.rgb = text_color
        p.space_after = Pt(3)
        p.line_spacing = Pt(int(font_size.pt * 1.4))
    return rect


def add_footnote(slide, text, top=Inches(7.0)):
    add_textbox(slide, Inches(1.0), top, Inches(11), Inches(0.4),
                text, font_name="Consolas", font_size=Pt(9),
                color=MUTED_TEXT, alignment=PP_ALIGN.RIGHT)


def add_slide_number(slide, num, total=None):
    txt = f"{num}" if total is None else f"{num} / {total}"
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1), Inches(0.4),
                txt, font_name="Consolas", font_size=Pt(9),
                color=MUTED_TEXT, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG2)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.8), Inches(11.333), Inches(0.5),
            "CHAPTER 7", font_name="Consolas", font_size=Pt(14),
            color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11.333), Inches(1.2),
            "나이브 베이즈", font_name="Georgia", font_size=Pt(54),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(0.8),
            "Naive Bayes Classification", font_name="Georgia", font_size=Pt(28),
            color=ACCENT_GOLD, italic=True, alignment=PP_ALIGN.CENTER)

add_decorative_line(slide, Inches(4.2), Inches(4), Inches(5.333))

add_textbox(slide, Inches(1), Inches(4.6), Inches(11.333), Inches(0.8),
            "스팸 여부 판단하기", font_name="Georgia", font_size=Pt(24),
            color=BODY_PARCHMENT, alignment=PP_ALIGN.CENTER)

add_multi_text(slide, Inches(1), Inches(5.8), Inches(11.333), Inches(1.2), [
    ("조건부 확률 기반의 분류 모델로 문자 데이터셋을 분석하여", Pt(14), MUTED_TEXT, False, False),
    ("스팸 문자를 필터링하는 나이브 베이즈를 깊게 이해합니다.", Pt(14), MUTED_TEXT, False, False),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 2: 학습 목표 & 순서
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "학습 목표 & 학습 순서", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_accent_box(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2),
               "학습 목표: 나이브 베이즈 모델로 문자 데이터셋을 분석해\n"
               "스팸 문자를 필터링하고, 나이브 베이즈를 더 깊게 이해합니다.",
               text_color=HIGHLIGHT, font_size=Pt(16))

# Learning steps
steps = [
    ("7.1", "문제 정의", "스팸 문자 여부 판별 미션 이해"),
    ("7.2", "라이브러리 및 데이터 불러오기", "pandas, numpy, matplotlib, seaborn, sklearn, nltk"),
    ("7.3", "특수 기호 제거하기", "string.punctuation을 활용한 노이즈 제거"),
    ("7.4", "불용어 제거하기", "NLTK stopwords를 사용한 불용어 필터링"),
    ("7.5", "목표 컬럼 형태 변경하기", "spam/ham → 1/0 숫자 변환"),
    ("7.6", "카운트 기반 벡터화하기", "CountVectorizer로 단어를 벡터로 변환"),
    ("7.7", "모델링 및 예측/평가하기", "MultinomialNB, 정확도, 혼동 행렬"),
    ("7.8", "이해하기: 나이브 베이즈", "베이즈 정리, 사전/사후 확률"),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(3.2) + Inches(i * 0.5)
    # Number circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.2), y, Inches(0.35), Inches(0.35)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT_GOLD
    circ.line.fill.background()
    # Number text
    add_textbox(slide, Inches(1.2), y - Pt(1), Inches(0.35), Inches(0.35),
                num, font_name="Consolas", font_size=Pt(9),
                color=BG_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Inches(1.7), y, Inches(3.5), Inches(0.4),
                title, font_size=Pt(14), color=TITLE_GOLD, bold=True)
    # Description
    add_textbox(slide, Inches(5.5), y, Inches(6.5), Inches(0.4),
                desc, font_size=Pt(12), color=MUTED_TEXT, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: 나이브 베이즈 소개
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "나이브 베이즈 소개", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.8), Inches(11.6), Inches(2.5), [
    ("나이브 베이즈(Naive Bayes)란?", Pt(20), HIGHLIGHT, True, False),
    ("", Pt(8), BODY_PARCHMENT),
    ("베이즈 정리를 적용한 조건부 확률 기반의 분류 모델입니다.", Pt(16), BODY_PARCHMENT),
    ("", Pt(8), BODY_PARCHMENT),
    ("조건부 확률: A가 일어났을 때 B가 일어날 확률", Pt(15), BODY_PARCHMENT, False, True),
    ("예) '무료'라는 단어가 들어있을 때 해당 메일이 스팸일 확률", Pt(14), MUTED_TEXT, False, True),
    ("", Pt(8), BODY_PARCHMENT),
    ("스팸 필터링을 위한 대표적인 모델로, 자연어 처리(NLP)가", Pt(15), BODY_PARCHMENT),
    ("목적일 때 여전히 좋은 선택이 될 수 있습니다.", Pt(15), BODY_PARCHMENT),
    ("(딥러닝보다 간단한 방법으로 자연어 처리를 원할 때 적합)", Pt(13), MUTED_TEXT, False, True),
])

# Pros/Cons boxes
add_accent_box(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(2.0),
               "장점\n"
               "• 비교적 간단한 알고리즘, 속도가 빠름\n"
               "• 작은 훈련셋으로도 잘 예측\n"
               "• 독립변수가 모두 독립적이면 우수한 성능",
               box_color=ACCENT_GREEN, text_color=BODY_PARCHMENT, font_size=Pt(13))

add_accent_box(slide, Inches(6.8), Inches(5.0), Inches(5.5), Inches(2.0),
               "단점\n"
               "• 독립변수가 각각 독립적임을 전제\n"
               "  (실제 데이터에서 그런 경우가 많지 않음)\n"
               "• 숫자형 변수가 많을 때는 적합하지 않음",
               box_color=ACCENT_RED, text_color=BODY_PARCHMENT, font_size=Pt(13))

# ════════════════════════════════════════════════════════════════
# SLIDE 4: 유용한 곳 & TOP10 선정 이유
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "나이브 베이즈가 유용한 곳", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.8), Inches(11.6), Inches(3.0), [
    ("TOP 10 선정 이유", Pt(20), HIGHLIGHT, True, False),
    ("", Pt(8), BODY_PARCHMENT),
    ("범용성이 높지는 않지만, 독립변수들이 모두 독립적이라면", Pt(16), BODY_PARCHMENT),
    ("충분히 경쟁력이 있는 알고리즘입니다.", Pt(16), BODY_PARCHMENT),
    ("특히 딥러닝을 제외하고 자연어 처리에 가장 적합한 알고리즘입니다.", Pt(16), BODY_PARCHMENT),
])

# Useful situations
useful_items = [
    "각 독립변수들이 모두 독립적이고 그 중요도가 비슷할 때",
    "자연어 처리(NLP)에서 간단하지만 좋은 성능을 보여줌",
    "범주 형태의 변수가 많을 때 적합",
    "숫자형 변수가 많은 때는 적합하지 않음",
]

for i, item in enumerate(useful_items):
    y_pos = Inches(4.2) + Inches(i * 0.55)
    # bullet dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.5), y_pos + Inches(0.08), Inches(0.12), Inches(0.12)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = TITLE_GOLD
    dot.line.fill.background()
    add_textbox(slide, Inches(1.8), y_pos, Inches(10), Inches(0.45),
                item, font_size=Pt(15), color=BODY_PARCHMENT)

add_accent_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8),
               "핵심: 조건부 확률로 데이터를 분류하는 알고리즘",
               text_color=HIGHLIGHT, font_size=Pt(16))

# ════════════════════════════════════════════════════════════════
# SLIDE 5: 문제 정의
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.1 문제 정의 : 한눈에 보는 예측 목표", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_accent_box(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.0),
               "미션: 문자 데이터셋을 이용해 스팸 여부를 판단하라!",
               text_color=HIGHLIGHT, font_size=Pt(20))

# Mission details
details = [
    ("알고리즘", "나이브 베이즈 (Naive Bayes)"),
    ("데이터셋", "spam.csv"),
    ("종속변수", "target (스팸 여부)"),
    ("독립변수", "text (문자 내용)"),
    ("문제 유형", "분류 (Classification)"),
    ("평가지표", "정확도(Accuracy), 혼동행렬(Confusion Matrix)"),
    ("모델", "MultinomialNB"),
]

for i, (label, value) in enumerate(details):
    y_pos = Inches(3.1) + Inches(i * 0.48)
    add_textbox(slide, Inches(1.5), y_pos, Inches(2.5), Inches(0.4),
                label, font_size=Pt(14), color=ACCENT_GOLD, bold=True)
    add_textbox(slide, Inches(4.2), y_pos, Inches(8), Inches(0.4),
                value, font_size=Pt(14), color=BODY_PARCHMENT)

add_multi_text(slide, Inches(0.8), Inches(6.4), Inches(11.6), Inches(0.8), [
    ("데이터 소개: 스팸 문자 데이터로, 독립변수 text 하나에 긴 문장 형태의 데이터가 들어있어", Pt(13), MUTED_TEXT, False, True),
    ("많은 전처리 작업이 필요합니다. 각 문장에 들어간 단어를 활용하여 스팸 여부를 예측합니다.", Pt(13), MUTED_TEXT, False, True),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 6: 라이브러리 및 데이터 불러오기
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.2 라이브러리 및 데이터 불러오기", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_code_block(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.8),
               "import pandas as pd\n"
               "import numpy as np\n"
               "import matplotlib.pyplot as plt\n"
               "import seaborn as sns\n"
               "\n"
               "file_url = 'https://media...spam.csv'\n"
               "data = pd.read_csv(file_url)")

# Library descriptions
libs = [
    ("pandas", "데이터프레임 조작"),
    ("numpy", "수치 연산"),
    ("matplotlib", "시각화"),
    ("seaborn", "통계 시각화"),
    ("sklearn", "머신러닝 모델 & 평가"),
    ("nltk", "자연어 처리 (불용어 등)"),
]

for i, (lib, desc) in enumerate(libs):
    y_pos = Inches(1.8) + Inches(i * 0.45)
    add_textbox(slide, Inches(6.8), y_pos, Inches(2.0), Inches(0.4),
                lib, font_name="Consolas", font_size=Pt(14), color=TITLE_GOLD, bold=True)
    add_textbox(slide, Inches(9.0), y_pos, Inches(3.5), Inches(0.4),
                desc, font_size=Pt(13), color=BODY_PARCHMENT)

# Data preview
add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.6), Inches(0.5),
            "데이터 확인: data.head()", font_size=Pt(16), color=HIGHLIGHT, bold=True)

add_code_block(slide, Inches(0.8), Inches(5.4), Inches(11.6), Inches(1.8),
               "# target: 종속변수 (spam/ham)\n"
               "# text: 독립변수 (문자 내용)\n"
               "\n"
               "# 0  ham   Go until jurong point, crazy...\n"
               "# 1  ham   Ok lar... Joking wif u oni...\n"
               "# 2  spam  Free entry in 2 a wkly comp...\n"
               "\n"
               "data['target'].unique()  # array(['ham', 'spam'])")

# ════════════════════════════════════════════════════════════════
# SLIDE 7: 전처리 개요
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "전처리 파이프라인 개요", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

# Pipeline steps
pipeline_steps = [
    ("STEP 1", "특수 기호 제거", "쉼표, 마침표 등 노이즈 제거\nstring.punctuation 활용", Inches(1.8)),
    ("STEP 2", "불용어 제거", "분석에 도움 안 되는 단어 제거\nNLTK stopwords 활용", Inches(3.2)),
    ("STEP 3", "목표 컬럼 변환", "spam → 1, ham → 0\nmap() 함수 활용", Inches(4.6)),
    ("STEP 4", "카운트 벡터화", "단어를 숫자 벡터로 변환\nCountVectorizer 활용", Inches(6.0)),
]

for i, (step, title, desc, y) in enumerate(pipeline_steps):
    # Step box
    add_accent_box(slide, Inches(0.8), y, Inches(1.8), Inches(1.0),
                   step, text_color=HIGHLIGHT, font_size=Pt(16))
    add_accent_box(slide, Inches(3.0), y, Inches(3.0), Inches(1.0),
                   title, text_color=TITLE_GOLD, font_size=Pt(18))
    add_textbox(slide, Inches(6.5), y + Inches(0.1), Inches(6.0), Inches(0.9),
                desc, font_size=Pt(14), color=BODY_PARCHMENT)
    # Arrow
    if i < len(pipeline_steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(1.6), y + Inches(1.05),
            Inches(0.3), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_GOLD
        arrow.line.fill.background()

# ════════════════════════════════════════════════════════════════
# SLIDE 8: 특수 기호 제거 (1)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.3 특수 기호 제거하기 (1/2)", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(1.2), [
    ("왜 특수 기호를 제거하는가?", Pt(18), HIGHLIGHT, True),
    ("자연어를 다룰 때 데이터의 기준은 단어입니다.", Pt(15), BODY_PARCHMENT),
    ("단어를 처리할 때 쉼표, 마침표 등의 특수 기호는 노이즈가 되므로 제거해야 합니다.", Pt(15), BODY_PARCHMENT),
])

add_code_block(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(2.0),
               "import string\n"
               "\n"
               "# 특수 기호 목록 확인\n"
               "string.punctuation\n"
               "# !\"#$%&'()*+,-./:;<=>?@[\\]^_`{|}~")

add_code_block(slide, Inches(6.8), Inches(3.3), Inches(5.5), Inches(2.0),
               "# 문자열에서 특수 기호 제거\n"
               "sample_string = data['text'].loc[0]\n"
               "\n"
               "for i in sample_string:\n"
               "    if i not in string.punctuation:\n"
               "        print(i)  # 특수기호 아닌 것만")

add_multi_text(slide, Inches(0.8), Inches(5.5), Inches(11.6), Inches(1.8), [
    ("핵심 로직", Pt(16), TITLE_GOLD, True),
    ("① 문자열에서 문자를 하나씩 꺼냄", Pt(14), BODY_PARCHMENT),
    ("② 특수 기호인지 판단 (in / not in string.punctuation)", Pt(14), BODY_PARCHMENT),
    ("③ 특수 기호가 아닌 문자들만 리스트에 저장", Pt(14), BODY_PARCHMENT),
    ("④ join()으로 리스트를 다시 문자열로 합침", Pt(14), BODY_PARCHMENT),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 9: 특수 기호 제거 (2)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.3 특수 기호 제거하기 (2/2)", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(0.5),
            "remove_punc() 함수 만들기", font_size=Pt(18), color=HIGHLIGHT, bold=True)

add_code_block(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.6),
               "def remove_punc(x):\n"
               "    new_string = []  # 빈 리스트\n"
               "    for i in x:       # 문자열 순회\n"
               "        if i not in string.punctuation:\n"
               "            new_string.append(i)\n"
               "    new_string = ''.join(new_string)\n"
               "    return new_string")

add_code_block(slide, Inches(6.8), Inches(2.3), Inches(5.5), Inches(2.6),
               "# 테스트\n"
               "remove_punc(sample_string)\n"
               "# 'Go until jurong point crazy\n"
               "#  Available only in bugis n great\n"
               "#  world la e buffet Cine there\n"
               "#  got amore wat'")

add_textbox(slide, Inches(0.8), Inches(5.1), Inches(11.6), Inches(0.5),
            "데이터프레임 전체에 적용: apply() 함수", font_size=Pt(18), color=HIGHLIGHT, bold=True)

add_code_block(slide, Inches(0.8), Inches(5.7), Inches(11.6), Inches(1.4),
               "# 주의! remove_punc(data['text'])는 한 줄에 모든 문자가 합쳐짐\n"
               "# → apply()를 사용하여 각 행마다 별도로 적용\n"
               "\n"
               "data['text'] = data['text'].apply(remove_punc)  # 각 행에 함수 적용")

# ════════════════════════════════════════════════════════════════
# SLIDE 10: in / join 함수 설명
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "핵심 함수: in, join(), apply()", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

# in operator
add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.8), Inches(0.5),
            "in 연산자", font_size=Pt(20), color=HIGHLIGHT, bold=True)
add_code_block(slide, Inches(0.8), Inches(2.3), Inches(3.8), Inches(1.8),
               "'a' in 'apple'   # True\n"
               "'b' in 'apple'   # False\n"
               "\n"
               "# 문자열에 속하는지\n"
               "# 아닌지 판단")

# join function
add_textbox(slide, Inches(5.0), Inches(1.7), Inches(3.5), Inches(0.5),
            "join() 함수", font_size=Pt(20), color=HIGHLIGHT, bold=True)
add_code_block(slide, Inches(5.0), Inches(2.3), Inches(3.5), Inches(1.8),
               "sample = ['a','p','p','l','e']\n"
               "'_'.join(sample)\n"
               "# 'a_p_p_l_e'\n"
               "\n"
               "''.join(sample)\n"
               "# 'apple'")

# apply function
add_textbox(slide, Inches(8.8), Inches(1.7), Inches(3.8), Inches(0.5),
            "apply() 함수", font_size=Pt(20), color=HIGHLIGHT, bold=True)
add_code_block(slide, Inches(8.8), Inches(2.3), Inches(3.8), Inches(1.8),
               "# 데이터의 각 행에\n"
               "# 별도로 함수를 적용\n"
               "\n"
               "data['text'].apply(func)\n"
               "# 각 행 → func → 결과")

add_accent_box(slide, Inches(0.8), Inches(4.5), Inches(11.6), Inches(2.8),
               "apply() 사용 시 주의사항\n\n"
               "• remove_punc(data['text']): 전체를 하나의 문자열로 처리 → 한 행에 모든 문자가 합쳐짐 ✗\n"
               "• data['text'].apply(remove_punc): 각 행마다 별도로 함수 적용 → 올바른 결과 ✓\n\n"
               "apply()는 데이터프레임의 한 줄 한 줄을 따로 함수에 적용시킬 수 있습니다.\n"
               "함수가 한 줄의 문자열에만 작동하도록 설계되었기 때문에 반드시 apply()를 사용해야 합니다.",
               text_color=BODY_PARCHMENT, font_size=Pt(14))

# ════════════════════════════════════════════════════════════════
# SLIDE 11: 불용어 제거 (1)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.4 불용어 제거하기 (1/2)", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(1.8), [
    ("불용어(Stopwords)란?", Pt(20), HIGHLIGHT, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("자연어 분석에 큰 도움이 안 되는 단어 (I, you, he, she, the, a 등)", Pt(15), BODY_PARCHMENT),
    ("이러한 단어를 제거하면 데이터를 가볍게 만들 수 있습니다.", Pt(15), BODY_PARCHMENT),
    ("자연어 처리에서는 각 단어가 하나의 독립변수처럼 작용하므로", Pt(15), BODY_PARCHMENT),
    ("불용어를 제거해 분석의 부담을 줄입니다.", Pt(15), BODY_PARCHMENT),
])

add_code_block(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(2.5),
               "import nltk\n"
               "nltk.download('stopwords')\n"
               "\n"
               "from nltk.corpus import stopwords\n"
               "\n"
               "# 영어 불용어 목록 확인\n"
               "stopwords.words('english')\n"
               "# ['i','me','my','myself','we',...]")

add_multi_text(slide, Inches(6.8), Inches(3.8), Inches(5.5), Inches(2.5), [
    ("NLTK stopwords 지원 언어 (24개)", Pt(14), TITLE_GOLD, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("arabic, azerbaijani, bengali, danish,", Pt(12), MUTED_TEXT),
    ("dutch, english, finnish, french,", Pt(12), MUTED_TEXT),
    ("german, greek, hungarian, indonesian,", Pt(12), MUTED_TEXT),
    ("italian, kazakh, nepali, norwegian,", Pt(12), MUTED_TEXT),
    ("portuguese, romanian, russian, slovene,", Pt(12), MUTED_TEXT),
    ("spanish, swedish, tajik, turkish", Pt(12), MUTED_TEXT),
    ("", Pt(6), BODY_PARCHMENT),
    ("※ 한국어는 미지원 → 별도 리스트 필요", Pt(13), ACCENT_RED),
])

add_accent_box(slide, Inches(0.8), Inches(6.5), Inches(11.6), Inches(0.6),
               "한국어 불용어: www.ranks.nl/stopwords/korean 에서 확인 가능",
               text_color=MUTED_TEXT, font_size=Pt(13))

# ════════════════════════════════════════════════════════════════
# SLIDE 12: 불용어 제거 (2)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.4 불용어 제거하기 (2/2)", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5),
            "불용어 제거 과정", font_size=Pt(18), color=HIGHLIGHT, bold=True)

add_multi_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(1.5), [
    ("① split()으로 문장을 단어 단위 리스트로 변환", Pt(14), BODY_PARCHMENT),
    ("② 각 단어가 불용어인지 판단", Pt(14), BODY_PARCHMENT),
    ("③ 불용어가 아니면 소문자로 변환하여 저장", Pt(14), BODY_PARCHMENT),
    ("④ ' '.join()으로 문자열로 합침", Pt(14), BODY_PARCHMENT),
])

add_code_block(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.0),
               "# split() : 공백 기준 단어 분리\n"
               "sample_string.split()\n"
               "# ['Go','until','jurong','point',...]\n"
               "\n"
               "# 특정 문자 기준 분리도 가능\n"
               "'This is not - SPAM'.split('-')\n"
               "# ['This is not ', ' SPAM']")

add_textbox(slide, Inches(0.8), Inches(3.9), Inches(11.6), Inches(0.5),
            "stop_words() 함수 구현", font_size=Pt(18), color=HIGHLIGHT, bold=True)

add_code_block(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.6),
               "def stop_words(x):\n"
               "    new_string = []\n"
               "    for i in x.split():\n"
               "        if i.lower() not in \\\n"
               "           stopwords.words('english'):\n"
               "            new_string.append(i.lower())\n"
               "    new_string = ' '.join(new_string)\n"
               "    return new_string")

add_code_block(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.6),
               "# 데이터에 적용\n"
               "data['text'] = data['text'].apply(stop_words)\n"
               "\n"
               "# 결과 확인\n"
               "# 0  go jurong point crazy available...\n"
               "# 1  ok lar joking wif u oni\n"
               "# 2  free entry 2 wkly comp win fa...")

# ════════════════════════════════════════════════════════════════
# SLIDE 13: 대소문자 변환
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "대소문자 변환 함수", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(1.2), [
    ("왜 소문자 변환이 필요한가?", Pt(18), HIGHLIGHT, True),
    ("파이썬은 대소문자를 구분합니다. stopwords의 단어들은 모두 소문자로 되어 있으므로", Pt(15), BODY_PARCHMENT),
    ("비교할 단어도 소문자여야 정확한 판단이 가능합니다. (Go ≠ go)", Pt(15), BODY_PARCHMENT),
])

# Three function boxes
funcs = [
    ("lower()", "소문자로 변환", "sample_word = 'NaiveBayes'\nsample_word.lower()\n# 'naivebayes'"),
    ("upper()", "대문자로 변환", "sample_word.upper()\n# 'NAIVEBAYES'"),
    ("capitalize()", "첫 글자만 대문자", "sample_word.capitalize()\n# 'Naivebayes'"),
]

for i, (func_name, desc, code) in enumerate(funcs):
    x_pos = Inches(0.8) + Inches(i * 4.2)
    add_accent_box(slide, x_pos, Inches(3.3), Inches(3.8), Inches(0.8),
                   f"{func_name}  —  {desc}",
                   text_color=TITLE_GOLD, font_size=Pt(15))
    add_code_block(slide, x_pos, Inches(4.3), Inches(3.8), Inches(1.5), code)

add_accent_box(slide, Inches(0.8), Inches(6.2), Inches(11.6), Inches(0.8),
               "불용어 비교 시 i.lower()를 사용하여 대소문자 구분 없이 정확하게 불용어를 제거합니다.",
               text_color=BODY_PARCHMENT, font_size=Pt(15))

# ════════════════════════════════════════════════════════════════
# SLIDE 14: 목표 컬럼 형태 변경
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.5 목표 컬럼 형태 변경하기", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(1.0), [
    ("목표 컬럼(target)을 문자에서 숫자로 변환", Pt(18), HIGHLIGHT, True),
    ("문자 형식도 에러를 유발하지는 않지만, 해석에 문제가 생길 수 있으므로 숫자로 변환합니다.", Pt(15), BODY_PARCHMENT),
])

add_textbox(slide, Inches(0.8), Inches(2.9), Inches(11.6), Inches(0.5),
            "map() 함수 사용법", font_size=Pt(18), color=TITLE_GOLD, bold=True)

add_code_block(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(2.5),
               "# map() - 딕셔너리로 매핑\n"
               "sample1 = pd.Series(['a','b','c'])\n"
               "sample1.map({'a':'apple',\n"
               "             'b':'banana',\n"
               "             'c':'cherry'})\n"
               "# 0  apple\n"
               "# 1  banana\n"
               "# 2  cherry")

add_code_block(slide, Inches(6.8), Inches(3.5), Inches(5.5), Inches(2.5),
               "# map() - 함수 적용도 가능\n"
               "sample2 = pd.Series(['a','b','c'])\n"
               "def add_i(x):\n"
               "    return x + 'i'\n"
               "\n"
               "sample2.map(add_i)\n"
               "# 0  ai\n"
               "# 1  bi\n"
               "# 2  ci")

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.6), Inches(0.5),
            "스팸/햄 변환 적용", font_size=Pt(18), color=HIGHLIGHT, bold=True)

add_code_block(slide, Inches(0.8), Inches(6.7), Inches(11.6), Inches(0.6),
               "data['target'] = data['target'].map({'spam': 1, 'ham': 0})  # spam→1, ham→0")

# ════════════════════════════════════════════════════════════════
# SLIDE 15: 카운트 기반 벡터화 개념
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.6 카운트 기반 벡터화 — 개념", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(1.5), [
    ("카운트 기반 벡터화(Count Vectorization)란?", Pt(20), HIGHLIGHT, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("문자를 개수 기반으로 벡터화하는 방식입니다.", Pt(16), BODY_PARCHMENT),
    ("데이터 전체에 존재하는 모든 단어를 사전처럼 모은 뒤 인덱스를 부여하고,", Pt(16), BODY_PARCHMENT),
    ("문장마다 속한 단어가 있는 인덱스를 카운트합니다.", Pt(16), BODY_PARCHMENT),
])

# Example visualization
add_accent_box(slide, Inches(0.8), Inches(3.5), Inches(11.6), Inches(3.5),
               "예시: 카운트 기반 벡터화 과정\n\n"
               "원본 데이터:\n"
               "  data[0] = 'brown dog white cat brown bear'\n"
               "  data[1] = 'white dog black dog'\n\n"
               "① 모든 단어를 확인해 컬럼으로 삼기:\n"
               "   brown  black  white  cat  bear  dog\n\n"
               "② 각 문장마다 출현한 단어 수 확인:\n"
               "   data[0]:  2      0      1     1    1     1   → 총 단어 6종\n"
               "   data[1]:  0      1      1     0    0     2   → 총 단어 3종\n\n"
               "③ 인덱스 부여: 0=brown, 1=black, 2=white, 3=cat, 4=bear, 5=dog",
               text_color=CODE_TEXT, font_size=Pt(13))

# ════════════════════════════════════════════════════════════════
# SLIDE 16: 카운트 기반 벡터화 코드
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.6 카운트 기반 벡터화 — 코드", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_code_block(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5),
               "from sklearn.feature_extraction.text \\\n"
               "    import CountVectorizer\n"
               "\n"
               "x = data['text']    # 독립변수\n"
               "y = data['target']  # 종속변수\n"
               "\n"
               "cv = CountVectorizer()  # 객체 생성\n"
               "cv.fit(x)  # 학습하기")

add_code_block(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.5),
               "# 단어와 인덱스 확인\n"
               "cv.vocabulary_\n"
               "# {'go': 3791, 'jurong': 4687,\n"
               "#  'point': 6433, 'crazy': 2497,\n"
               "#  'available': 1414, 'bugis': 1881,\n"
               "#  'great': 3888, 'world': 9184, ...}")

add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.6), Inches(0.5),
            "transform()으로 데이터 변환", font_size=Pt(18), color=HIGHLIGHT, bold=True)

add_code_block(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(2.2),
               "x = cv.transform(x)  # 변환\n"
               "print(x)\n"
               "# (0, 1181) 1  ← 0행, 인덱스1181 1회\n"
               "# (0, 1414) 1\n"
               "# (0, 3791) 1  ← go\n"
               "# (0, 4687) 1  ← jurong\n"
               "# (0, 6433) 1  ← point")

add_multi_text(slide, Inches(6.8), Inches(5.0), Inches(5.5), Inches(2.2), [
    ("출력 형태 해석", Pt(16), TITLE_GOLD, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("(행번호, 인덱스) 출현횟수", Pt(15), HIGHLIGHT, True, False, "Consolas"),
    ("", Pt(6), BODY_PARCHMENT),
    ("• 0번째 행에 인덱스 1181인 단어가 1번 등장", Pt(13), BODY_PARCHMENT),
    ("• 인덱스는 cv.vocabulary_에서 확인", Pt(13), BODY_PARCHMENT),
    ("• 출현하지 않은 단어는 포함되지 않음", Pt(13), BODY_PARCHMENT),
    ("  (희소 행렬/Sparse Matrix 형태)", Pt(13), MUTED_TEXT),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 17: 모델링 및 예측
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.7 모델링 및 예측하기", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_code_block(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.0),
               "from sklearn.model_selection \\\n"
               "    import train_test_split\n"
               "\n"
               "# 훈련셋/시험셋 분할\n"
               "x_train, x_test, y_train, y_test = \\\n"
               "    train_test_split(x, y,\n"
               "        test_size=0.2,\n"
               "        random_state=100)")

add_code_block(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(3.0),
               "from sklearn.naive_bayes \\\n"
               "    import MultinomialNB\n"
               "\n"
               "# 모델 생성 및 학습\n"
               "model = MultinomialNB()\n"
               "model.fit(x_train, y_train)\n"
               "\n"
               "# 예측\n"
               "pred = model.predict(x_test)")

add_accent_box(slide, Inches(0.8), Inches(5.0), Inches(11.6), Inches(1.0),
               "MultinomialNB: 다항 분포에 대한 Naive Bayes 알고리즘\n"
               "다른 NB 모델: GaussianNB (정규분포), BernoulliNB (베르누이 분포)\n"
               "데이터 특성을 명확히 알기 어렵다면 세 가지 모두 사용하여 가장 좋은 모델 선택",
               text_color=BODY_PARCHMENT, font_size=Pt(14))

add_accent_box(slide, Inches(0.8), Inches(6.2), Inches(11.6), Inches(0.8),
               "사용법: fit()으로 학습 → predict()로 예측 (다른 모델과 동일한 패턴)",
               text_color=HIGHLIGHT, font_size=Pt(16))

# ════════════════════════════════════════════════════════════════
# SLIDE 18: 평가 - 정확도
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.7 평가하기 — 정확도 & 혼동 행렬", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_code_block(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.2),
               "from sklearn.metrics import \\\n"
               "    accuracy_score, confusion_matrix\n"
               "\n"
               "# 정확도 계산\n"
               "accuracy_score(y_test, pred)\n"
               "# 0.9856502242152466\n"
               "# → 약 98.6% 정확도!")

add_code_block(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.2),
               "# 혼동 행렬 출력\n"
               "print(confusion_matrix(y_test, pred))\n"
               "# [[965  12]\n"
               "#  [  4 134]]\n"
               "\n"
               "# 히트맵으로 시각화\n"
               "sns.heatmap(confusion_matrix(\n"
               "    y_test, pred), annot=True, fmt='d')")

# Accuracy calculation
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11.6), Inches(0.5),
            "정확도 계산", font_size=Pt(18), color=HIGHLIGHT, bold=True)

add_accent_box(slide, Inches(0.8), Inches(4.8), Inches(11.6), Inches(2.2),
               "정확도 = 정확한 예측 건수 / 전체 경우 수\n\n"
               "       = (965 + 134) / (965 + 134 + 12 + 4)\n"
               "       = 1099 / 1115\n"
               "       ≈ 98.6%\n\n"
               "accuracy_score() 계산 결과와 동일한 수치",
               text_color=BODY_PARCHMENT, font_size=Pt(15))

# ════════════════════════════════════════════════════════════════
# SLIDE 19: 혼동 행렬 상세
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "혼동 행렬 (Confusion Matrix) 상세", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

# Build confusion matrix visual
# Header
add_textbox(slide, Inches(5.0), Inches(1.7), Inches(4), Inches(0.4),
            "예측값", font_size=Pt(16), color=TITLE_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.5), Inches(2.1), Inches(2), Inches(0.4),
            "0", font_size=Pt(16), color=TITLE_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(6.5), Inches(2.1), Inches(2), Inches(0.4),
            "1", font_size=Pt(16), color=TITLE_GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2.0), Inches(2.5), Inches(2), Inches(2.5),
            "실\n젯\n값", font_size=Pt(16), color=TITLE_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.5), Inches(2.6), Inches(1), Inches(0.4),
            "0", font_size=Pt(16), color=TITLE_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.5), Inches(3.8), Inches(1), Inches(0.4),
            "1", font_size=Pt(16), color=TITLE_GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Cells
cells = [
    (Inches(4.5), Inches(2.5), "965", "TN\n(True Negative)", ACCENT_GREEN),
    (Inches(6.5), Inches(2.5), "12", "FP\n(False Positive)", ACCENT_RED),
    (Inches(4.5), Inches(3.7), "4", "FN\n(False Negative)", ACCENT_RED),
    (Inches(6.5), Inches(3.7), "134", "TP\n(True Positive)", ACCENT_GREEN),
]

for left, top, num, label, color in cells:
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2), Inches(1.1)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x20, 0x18, 0x0C)
    rect.line.color.rgb = color
    rect.line.width = Pt(2)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.05), Inches(1.8), Inches(0.5),
                num, font_size=Pt(28), color=HIGHLIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.5), Inches(1.8), Inches(0.6),
                label, font_size=Pt(11), color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

# Legend
add_multi_text(slide, Inches(0.8), Inches(5.2), Inches(11.6), Inches(2.0), [
    ("혼동 행렬 해석", Pt(16), TITLE_GOLD, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("• TN (True Negative) = 965: 실젯값 0, 예측값 0 → 정확한 예측 (음성을 음성으로)", Pt(13), BODY_PARCHMENT),
    ("• TP (True Positive) = 134: 실젯값 1, 예측값 1 → 정확한 예측 (양성을 양성으로)", Pt(13), BODY_PARCHMENT),
    ("• FP (False Positive) = 12: 실젯값 0, 예측값 1 → 1종 오류 (음성을 양성으로)", Pt(13), ACCENT_RED),
    ("• FN (False Negative) = 4: 실젯값 1, 예측값 0 → 2종 오류 (양성을 음성으로)", Pt(13), ACCENT_RED),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 20: 1종/2종 오류
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "1종 오류 vs 2종 오류", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

# Type 1
add_accent_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.2),
               "1종 오류 (Type 1 Error)\nFalse Positive\n\n"
               "실제 음성인 것을 양성으로 예측하는 오류\n\n"
               "암 진단 예: 실제 암이 아닌데 암이라고 진단\n"
               "→ 추가 진단으로 곧 아님을 알게 됨",
               box_color=ACCENT_BLUE, text_color=BODY_PARCHMENT, font_size=Pt(14))

# Type 2
add_accent_box(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.2),
               "2종 오류 (Type 2 Error)\nFalse Negative\n\n"
               "실제 양성인 것을 음성으로 예측하는 오류\n\n"
               "암 진단 예: 실제 암인데 아니라고 진단\n"
               "→ 치료 시기를 놓칠 수 있음 (더 위험!)",
               box_color=ACCENT_RED, text_color=BODY_PARCHMENT, font_size=Pt(14))

add_accent_box(slide, Inches(0.8), Inches(4.3), Inches(11.6), Inches(1.2),
               "암기 Tip: '양성 1종, 음성 2종'\n"
               "• 1종 오류 = False Positive (거짓 양성) → 양성으로 잘못 예측\n"
               "• 2종 오류 = False Negative (거짓 음성) → 음성으로 잘못 예측",
               text_color=HIGHLIGHT, font_size=Pt(15))

# Spam context
add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.6), Inches(0.5),
            "스팸 문자에서의 1종/2종 오류", font_size=Pt(18), color=TITLE_GOLD, bold=True)

add_multi_text(slide, Inches(0.8), Inches(6.2), Inches(5.5), Inches(1.0), [
    ("1종 오류: 스팸이 아닌데 스팸으로 분류", Pt(14), BODY_PARCHMENT),
    ("→ 중요 문자가 필터링될 수 있음", Pt(13), MUTED_TEXT, False, True),
])

add_multi_text(slide, Inches(6.8), Inches(6.2), Inches(5.5), Inches(1.0), [
    ("2종 오류: 스팸인데 스팸이 아니라고 분류", Pt(14), BODY_PARCHMENT),
    ("→ 스팸을 받게 되지만 중요 문자 차단 적음", Pt(13), MUTED_TEXT, False, True),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 21: 베이즈 정리 수식
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG2)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "7.8 베이즈 정리 (Bayes' Theorem)", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

# Main formula
add_accent_box(slide, Inches(2.5), Inches(1.8), Inches(8.3), Inches(1.2),
               "P(A|B) = P(B|A) × P(A) / P(B)",
               text_color=HIGHLIGHT, font_size=Pt(28))

# Explanation of each term
terms = [
    ("P(A|B)", "사후확률", "B가 발생했을 때 A가 발생할 확률\n특정 단어가 등장했을 때 스팸일 확률"),
    ("P(A)", "사전확률", "B의 발생 유무와 관련없이 A가 발생할 확률\n전체 문자 중 스팸 문자의 비율"),
    ("P(B|A)", "우도(Likelihood)", "A가 발생했을 때 B가 발생할 확률\n스팸 메일인 경우 특정 단어가 들어있을 확률"),
    ("P(B)", "증거(Evidence)", "전체에서 B가 발생할 확률\n전체 문자에서 특정 단어가 들어있을 확률"),
]

for i, (symbol, name, desc) in enumerate(terms):
    y_pos = Inches(3.3) + Inches(i * 1.0)
    add_textbox(slide, Inches(0.8), y_pos, Inches(1.5), Inches(0.4),
                symbol, font_name="Consolas", font_size=Pt(18),
                color=HIGHLIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), y_pos, Inches(2.0), Inches(0.4),
                name, font_size=Pt(16), color=TITLE_GOLD, bold=True)
    add_textbox(slide, Inches(4.8), y_pos, Inches(7.5), Inches(0.9),
                desc, font_size=Pt(13), color=BODY_PARCHMENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 22: 베이즈 정리 예시
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "베이즈 정리 — 그림으로 이해하기", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

# Setup conditions
add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(1.5), [
    ("예시 조건", Pt(18), HIGHLIGHT, True),
    ("• 전체 문자 중 스팸 비율: 30% (햄 70%)", Pt(15), BODY_PARCHMENT),
    ("• 스팸에서 특정 단어 X가 포함된 경우: 50%", Pt(15), BODY_PARCHMENT),
    ("• 햄에서 특정 단어 X가 포함된 경우: 10%", Pt(15), BODY_PARCHMENT),
])

# Visual boxes for spam/ham
# Spam box
spam_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.5), Inches(3.5), Inches(1.5)
)
spam_box.fill.solid()
spam_box.fill.fore_color.rgb = RGBColor(0x80, 0x30, 0x20)
spam_box.line.color.rgb = ACCENT_RED
spam_box.line.width = Pt(2)
add_textbox(slide, Inches(1.1), Inches(3.55), Inches(3.3), Inches(1.4),
            "스팸 (30%)\n특정 단어 X 포함: 50%",
            font_size=Pt(15), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Ham box
ham_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.5), Inches(3.5), Inches(1.5)
)
ham_box.fill.solid()
ham_box.fill.fore_color.rgb = RGBColor(0x30, 0x60, 0x30)
ham_box.line.color.rgb = ACCENT_GREEN
ham_box.line.width = Pt(2)
add_textbox(slide, Inches(5.1), Inches(3.55), Inches(3.3), Inches(1.4),
            "햄 (70%)\n특정 단어 X 포함: 10%",
            font_size=Pt(15), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Calculation
add_accent_box(slide, Inches(0.8), Inches(5.3), Inches(11.6), Inches(1.8),
               "계산: 특정 단어 X가 있을 때 스팸일 확률 = ?\n\n"
               "• P(A) = 스팸 확률 = 0.3\n"
               "• P(B|A) = 스팸 중 X 포함 = 0.5\n"
               "• P(B) = 전체 중 X 포함 = (0.3 × 0.5) + (0.7 × 0.1) = 0.22\n\n"
               "P(A|B) = P(B|A) × P(A) / P(B) = 0.5 × 0.3 / 0.22 ≈ 0.6818 (약 68%)",
               text_color=BODY_PARCHMENT, font_size=Pt(14))

# ════════════════════════════════════════════════════════════════
# SLIDE 23: 나이브 베이즈의 '나이브'
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "나이브 베이즈의 핵심 가정", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(2.5), [
    ("왜 'Naive(나이브)'인가?", Pt(22), HIGHLIGHT, True),
    ("", Pt(8), BODY_PARCHMENT),
    ("나이브 베이즈는 모든 독립변수(feature)가 서로 독립이라고 가정합니다.", Pt(16), BODY_PARCHMENT),
    ("이 가정은 실제 데이터에서 거의 성립하지 않지만,", Pt(16), BODY_PARCHMENT),
    ("이 '순진한(naive)' 가정 덕분에 계산이 매우 단순해집니다.", Pt(16), BODY_PARCHMENT),
    ("", Pt(8), BODY_PARCHMENT),
    ("독립 가정이 성립하면:", Pt(16), TITLE_GOLD, True),
    ("P(x₁, x₂, ..., xₙ | C) = P(x₁|C) × P(x₂|C) × ... × P(xₙ|C)", Pt(15), HIGHLIGHT, False, False, "Consolas"),
])

add_accent_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5),
               "나이브 베이즈의 3가지 변형\n\n"
               "1. MultinomialNB\n"
               "   다항 분포 기반, 텍스트 분류에 적합\n"
               "   단어 빈도(count)를 특징으로 사용\n\n"
               "2. GaussianNB\n"
               "   정규(가우시안) 분포 기반\n"
               "   연속형 변수에 적합",
               text_color=BODY_PARCHMENT, font_size=Pt(13))

add_accent_box(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.5),
               "3. BernoulliNB\n"
               "   베르누이 분포 기반\n"
               "   이진(0/1) 특징에 적합\n\n"
               "모델 선택 가이드:\n"
               "• 텍스트 데이터 → MultinomialNB\n"
               "• 연속 수치 데이터 → GaussianNB\n"
               "• 이진 데이터 → BernoulliNB\n"
               "• 잘 모르면 → 셋 다 비교!",
               text_color=BODY_PARCHMENT, font_size=Pt(13))

# ════════════════════════════════════════════════════════════════
# SLIDE 24: 전체 코드 정리
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "전체 코드 흐름 정리", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_code_block(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.3),
               "# ① 라이브러리 & 데이터\n"
               "import pandas as pd, numpy as np\n"
               "import matplotlib.pyplot as plt\n"
               "import seaborn as sns\n"
               "import string, nltk\n"
               "from nltk.corpus import stopwords\n"
               "\n"
               "data = pd.read_csv(file_url)\n"
               "\n"
               "# ② 특수 기호 제거\n"
               "def remove_punc(x):\n"
               "    new_string = []\n"
               "    for i in x:\n"
               "        if i not in string.punctuation:\n"
               "            new_string.append(i)\n"
               "    return ''.join(new_string)\n"
               "\n"
               "data['text'] = data['text'].apply(remove_punc)")

add_code_block(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(5.3),
               "# ③ 불용어 제거\n"
               "def stop_words(x):\n"
               "    new_string = []\n"
               "    for i in x.split():\n"
               "        if i.lower() not in \\\n"
               "           stopwords.words('english'):\n"
               "            new_string.append(i.lower())\n"
               "    return ' '.join(new_string)\n"
               "data['text'] = data['text'].apply(stop_words)\n"
               "\n"
               "# ④ 타겟 변환 & 벡터화\n"
               "data['target'] = data['target'].map(\n"
               "    {'spam':1, 'ham':0})\n"
               "x = data['text']; y = data['target']\n"
               "cv = CountVectorizer()\n"
               "cv.fit(x); x = cv.transform(x)\n"
               "\n"
               "# ⑤ 모델링 & 평가\n"
               "model = MultinomialNB()\n"
               "model.fit(x_train, y_train)\n"
               "pred = model.predict(x_test)")

# ════════════════════════════════════════════════════════════════
# SLIDE 25: 학습 마무리
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "학습 마무리", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

# Summary flow
summary_steps = [
    ("문제 정의", "스팸 문자를 구분하는 모델을 만듭니다."),
    ("라이브러리", "pandas, numpy, matplotlib, seaborn, sklearn, nltk 임포트"),
    ("특수 기호 제거", "문자 메시지의 마침표, 느낌표 등 특수 기호 제거"),
    ("불용어 제거", "예측에 변별력이 없는 불용어를 제거"),
    ("목표 변환", "spam/ham → 1/0 숫자로 변경"),
    ("벡터화", "문장을 각 단어의 카운트 기반으로 벡터화"),
    ("모델링", "MultinomialNB 분류기 사용"),
    ("결과", "98.6%의 높은 정확도 달성!"),
]

for i, (title, desc) in enumerate(summary_steps):
    y_pos = Inches(1.8) + Inches(i * 0.65)
    # Step number
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.0), y_pos + Inches(0.05),
        Inches(0.35), Inches(0.35)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT_GOLD
    circ.line.fill.background()
    add_textbox(slide, Inches(1.0), y_pos + Inches(0.02), Inches(0.35), Inches(0.35),
                str(i+1), font_name="Consolas", font_size=Pt(12),
                color=BG_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.6), y_pos, Inches(2.5), Inches(0.5),
                title, font_size=Pt(15), color=TITLE_GOLD, bold=True)
    add_textbox(slide, Inches(4.3), y_pos, Inches(8.0), Inches(0.5),
                desc, font_size=Pt(14), color=BODY_PARCHMENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 26: 핵심 용어 정리
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "핵심 용어 & API 정리", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

terms_list = [
    ("나이브 베이즈 분류기", "조건부 확률 기반 모델, 자연어와 같이 변수 개수가 많을 때 유용"),
    ("1종 오류 (Type 1, FP)", "실제 음성인 것을 양성으로 예측하는 오류"),
    ("2종 오류 (Type 2, FN)", "실제 양성인 것을 음성으로 예측하는 오류"),
    ("사후확률", "사건 A가 발생한 상황에서 사건 B가 발생할 확률"),
    ("사전확률", "사건 A와 상관없이 사건 B가 발생할 확률"),
    ("베이즈 정리", "사전확률과 사후확률 사이의 관계를 나타내는 정리"),
]

for i, (term, desc) in enumerate(terms_list):
    y_pos = Inches(1.8) + Inches(i * 0.6)
    add_textbox(slide, Inches(0.8), y_pos, Inches(4.0), Inches(0.5),
                term, font_size=Pt(15), color=HIGHLIGHT, bold=True)
    add_textbox(slide, Inches(5.0), y_pos, Inches(7.5), Inches(0.5),
                desc, font_size=Pt(13), color=BODY_PARCHMENT)

add_decorative_line(slide, Inches(5.3), Inches(0.8), Inches(11.6))

apis = [
    ("string.punctuation", "특수 기호 목록 출력"),
    ("MultinomialNB()", "다항분포 나이브 베이즈 알고리즘"),
    ("nltk stopwords", "불용어 목록"),
    ("CountVectorizer()", "카운트 기반 벡터화 알고리즘"),
    ("confusion_matrix()", "혼동 행렬"),
    ("cv.vocabulary_", "벡터화 객체에서 단어와 인덱스 확인"),
]

for i, (api, desc) in enumerate(apis):
    y_pos = Inches(5.5) + Inches(i * 0.33)
    add_textbox(slide, Inches(0.8), y_pos, Inches(4.0), Inches(0.3),
                api, font_name="Consolas", font_size=Pt(12), color=TITLE_GOLD, bold=True)
    add_textbox(slide, Inches(5.0), y_pos, Inches(7.5), Inches(0.3),
                desc, font_size=Pt(12), color=MUTED_TEXT)

# ════════════════════════════════════════════════════════════════
# SLIDE 27: 연습 문제
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "연습 문제", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

questions = [
    ("Q1", "자연어 처리에서 의미 없이 빈번하게 발생하는 단어를 의미하는 것은?\n"
           "① 불용어  ② 특수기호  ③ 벡터  ④ 인덱스", "정답: ① 불용어"),
    ("Q2", "자연어를 머신러닝 알고리즘이 받아들일 수 있도록 각 단어와\n"
           "출현 빈도로 정리하는 함수는?\n"
           "① get_dummies()  ② CountVectorizer()  ③ StandardScaler()  ④ value_counts()",
     "정답: ② CountVectorizer()"),
    ("Q3", "혼동 행렬 해석 — False Negative는 총 64건? (0: 64,3 / 1: 16,17)\n"
           "→ FN은 실젯값 1, 예측값 0 → 총 16건",
     "정답: ④ (FN=64건은 오답, 실제 FN=16건)"),
    ("Q4", "나이브 베이즈를 가장 잘 설명한 것은?\n"
           "① 조건부 확률 기반, 사전/사후확률 활용  ② 선형 관계 전제\n"
           "③ 독립변수 적을 때  ④ 상관관계 강할 때도 잘 작동",
     "정답: ① 조건부 확률 기반"),
]

for i, (qnum, question, answer) in enumerate(questions):
    y_pos = Inches(1.7) + Inches(i * 1.4)
    add_textbox(slide, Inches(0.8), y_pos, Inches(0.5), Inches(0.4),
                qnum, font_size=Pt(16), color=HIGHLIGHT, bold=True)
    add_textbox(slide, Inches(1.4), y_pos, Inches(8.0), Inches(0.9),
                question, font_size=Pt(13), color=BODY_PARCHMENT)
    add_textbox(slide, Inches(9.5), y_pos + Inches(0.3), Inches(3.5), Inches(0.4),
                answer, font_size=Pt(12), color=ACCENT_GREEN, bold=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 28: 심화 - CountVectorizer 상세
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG2)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "심화: CountVectorizer 내부 동작", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(1.5), [
    ("CountVectorizer의 작동 원리 (scikit-learn)", Pt(18), HIGHLIGHT, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("1. 토큰화(Tokenization): 문장을 단어 단위로 분리", Pt(15), BODY_PARCHMENT),
    ("2. 어휘 구축(Vocabulary Building): 모든 고유 단어에 인덱스 부여", Pt(15), BODY_PARCHMENT),
    ("3. 벡터 인코딩(Encoding): 각 문서에서 단어 출현 횟수 카운트", Pt(15), BODY_PARCHMENT),
    ("4. 결과: 희소 행렬(Sparse Matrix) 형태로 저장", Pt(15), BODY_PARCHMENT),
])

add_code_block(slide, Inches(0.8), Inches(3.6), Inches(5.5), Inches(3.2),
               "from sklearn.feature_extraction.text \\\n"
               "    import CountVectorizer\n"
               "\n"
               "# 주요 매개변수\n"
               "cv = CountVectorizer(\n"
               "    max_features=5000,  # 최대 단어 수\n"
               "    ngram_range=(1,2),  # n-gram 범위\n"
               "    stop_words='english', # 내장 불용어\n"
               "    max_df=0.95,  # 문서빈도 상한\n"
               "    min_df=2      # 문서빈도 하한\n"
               ")")

add_multi_text(slide, Inches(6.8), Inches(3.6), Inches(5.5), Inches(3.2), [
    ("주요 매개변수 설명", Pt(16), TITLE_GOLD, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("• max_features: 빈도 상위 N개 단어만 사용", Pt(13), BODY_PARCHMENT),
    ("• ngram_range: 단어 조합 범위", Pt(13), BODY_PARCHMENT),
    ("  (1,1)=단어, (1,2)=단어+2단어조합", Pt(12), MUTED_TEXT),
    ("• stop_words: 불용어 제거 옵션", Pt(13), BODY_PARCHMENT),
    ("• max_df: 너무 많은 문서에 나오는 단어 제외", Pt(13), BODY_PARCHMENT),
    ("• min_df: 너무 적은 문서에 나오는 단어 제외", Pt(13), BODY_PARCHMENT),
    ("", Pt(6), BODY_PARCHMENT),
    ("TF-IDF: 단순 카운트 대신 가중치 부여", Pt(14), ACCENT_BLUE, True),
    ("→ TfidfVectorizer로 성능 개선 가능", Pt(13), MUTED_TEXT),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 29: 심화 - TF-IDF와 비교
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG2)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "심화: TF-IDF vs Count Vectorization", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_accent_box(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5),
               "Count Vectorization\n\n"
               "• 단순히 단어 출현 횟수를 카운트\n"
               "• 모든 단어를 동등하게 취급\n"
               "• 자주 나오는 일반적 단어도 높은 값\n"
               "• 간단하지만 정보 손실 가능\n\n"
               "사용: MultinomialNB와 주로 사용",
               text_color=BODY_PARCHMENT, font_size=Pt(14))

add_accent_box(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.5),
               "TF-IDF Vectorization\n\n"
               "• TF: 문서 내 단어 빈도\n"
               "• IDF: 전체 문서에서의 희귀도\n"
               "• TF × IDF = 중요도 가중치\n"
               "• 흔한 단어는 낮은 가중치\n\n"
               "사용: 더 정교한 텍스트 분류에 적합",
               text_color=BODY_PARCHMENT, font_size=Pt(14))

add_code_block(slide, Inches(0.8), Inches(4.5), Inches(11.6), Inches(2.5),
               "# TF-IDF 사용 예시\n"
               "from sklearn.feature_extraction.text import TfidfVectorizer\n"
               "\n"
               "tfidf = TfidfVectorizer(max_features=5000)\n"
               "x_tfidf = tfidf.fit_transform(data['text'])\n"
               "\n"
               "# TF-IDF 수식\n"
               "# TF(t,d) = (단어 t가 문서 d에 나타난 횟수) / (문서 d의 총 단어 수)\n"
               "# IDF(t) = log(전체 문서 수 / 단어 t를 포함하는 문서 수)\n"
               "# TF-IDF(t,d) = TF(t,d) × IDF(t)")

# ════════════════════════════════════════════════════════════════
# SLIDE 30: 심화 - 실전 구현 팁
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG2)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "심화: 실전 구현 & 성능 개선 팁", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_code_block(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.8),
               "# sklearn Pipeline으로 깔끔하게\n"
               "from sklearn.pipeline import Pipeline\n"
               "\n"
               "pipe = Pipeline([\n"
               "    ('vectorizer', CountVectorizer()),\n"
               "    ('classifier', MultinomialNB())\n"
               "])\n"
               "\n"
               "pipe.fit(x_train_text, y_train)\n"
               "pred = pipe.predict(x_test_text)")

add_code_block(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.8),
               "# 3가지 NB 모델 비교\n"
               "from sklearn.naive_bayes import (\n"
               "    MultinomialNB,\n"
               "    GaussianNB,\n"
               "    BernoulliNB\n"
               ")\n"
               "\n"
               "models = {\n"
               "    'Multinomial': MultinomialNB(),\n"
               "    'Bernoulli': BernoulliNB(binarize=True),\n"
               "}")

add_multi_text(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.3), [
    ("성능 개선 전략", Pt(18), HIGHLIGHT, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("1. TF-IDF 사용: 단순 카운트보다 효과적", Pt(14), BODY_PARCHMENT),
    ("2. n-gram 활용: 단어 조합으로 문맥 반영", Pt(14), BODY_PARCHMENT),
    ("3. 불용어 커스터마이징: 도메인별 불용어", Pt(14), BODY_PARCHMENT),
    ("4. 스테밍/레마타이제이션: 단어 원형 추출", Pt(14), BODY_PARCHMENT),
    ("5. 하이퍼파라미터 튜닝: alpha(스무딩)", Pt(14), BODY_PARCHMENT),
])

add_multi_text(slide, Inches(6.8), Inches(4.8), Inches(5.5), Inches(2.3), [
    ("관련 논문 & 참고자료", Pt(18), HIGHLIGHT, True),
    ("", Pt(6), BODY_PARCHMENT),
    ("• McCallum & Nigam (1998)", Pt(13), BODY_PARCHMENT),
    ("  'Comparison of Event Models for NB'", Pt(12), MUTED_TEXT, False, True),
    ("• Rennie et al. (2003)", Pt(13), BODY_PARCHMENT),
    ("  'Tackling the Poor Assumptions of NB'", Pt(12), MUTED_TEXT, False, True),
    ("• scikit-learn Documentation:", Pt(13), BODY_PARCHMENT),
    ("  sklearn.naive_bayes module", Pt(12), MUTED_TEXT, False, True),
    ("• NLTK Documentation:", Pt(13), BODY_PARCHMENT),
    ("  nltk.org for stopwords & tokenizers", Pt(12), MUTED_TEXT, False, True),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 31: 심화 - 나이브 베이즈 수학적 전개
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG2)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "심화: 나이브 베이즈의 수학적 전개", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(5.5), [
    ("분류 문제에서의 베이즈 정리", Pt(18), HIGHLIGHT, True),
    ("", Pt(8), BODY_PARCHMENT),
    ("목표: 주어진 특징 벡터 x = (x₁, x₂, ..., xₙ)에 대해 클래스 C를 예측", Pt(15), BODY_PARCHMENT),
    ("", Pt(6), BODY_PARCHMENT),
    ("P(C|x₁, x₂, ..., xₙ) = P(x₁, x₂, ..., xₙ|C) × P(C) / P(x₁, x₂, ..., xₙ)", Pt(14), HIGHLIGHT, False, False, "Consolas"),
    ("", Pt(8), BODY_PARCHMENT),
    ("나이브 가정 (독립 가정) 적용:", Pt(16), TITLE_GOLD, True),
    ("P(x₁, x₂, ..., xₙ|C) = P(x₁|C) × P(x₂|C) × ... × P(xₙ|C)", Pt(14), HIGHLIGHT, False, False, "Consolas"),
    ("", Pt(8), BODY_PARCHMENT),
    ("따라서:", Pt(16), TITLE_GOLD, True),
    ("P(C|x) ∝ P(C) × ∏ᵢ P(xᵢ|C)", Pt(14), HIGHLIGHT, False, False, "Consolas"),
    ("", Pt(8), BODY_PARCHMENT),
    ("최종 분류:", Pt(16), TITLE_GOLD, True),
    ("ŷ = argmax_C  P(C) × ∏ᵢ P(xᵢ|C)", Pt(14), HIGHLIGHT, False, False, "Consolas"),
    ("", Pt(8), BODY_PARCHMENT),
    ("Multinomial NB에서의 P(xᵢ|C) 계산:", Pt(16), TITLE_GOLD, True),
    ("P(xᵢ|C) = (클래스 C에서 단어 xᵢ의 출현 횟수 + α) / (클래스 C의 총 단어 수 + α × |V|)", Pt(13), BODY_PARCHMENT),
    ("여기서 α는 라플라스 스무딩 파라미터 (기본값 1.0), |V|는 어휘 크기", Pt(13), MUTED_TEXT, False, True),
])

# ════════════════════════════════════════════════════════════════
# SLIDE 32: 심화 코드 - 전체 파이프라인
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG2)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.333), Inches(0.8),
            "심화: 완전한 스팸 분류 파이프라인", font_name="Georgia", font_size=Pt(36),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(1.5))

add_code_block(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.2),
               "import re\n"
               "from sklearn.pipeline import Pipeline\n"
               "from sklearn.feature_extraction.text \\\n"
               "    import TfidfVectorizer\n"
               "from sklearn.naive_bayes \\\n"
               "    import MultinomialNB\n"
               "from sklearn.model_selection \\\n"
               "    import cross_val_score\n"
               "from sklearn.metrics \\\n"
               "    import classification_report\n"
               "\n"
               "# 전처리 함수 (정규표현식 버전)\n"
               "def preprocess(text):\n"
               "    text = re.sub(r'[^a-zA-Z\\s]','',text)\n"
               "    text = text.lower().strip()\n"
               "    return text\n"
               "\n"
               "data['text'] = data['text'].apply(preprocess)")

add_code_block(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(5.2),
               "# Pipeline 구성\n"
               "pipe = Pipeline([\n"
               "    ('tfidf', TfidfVectorizer(\n"
               "        max_features=5000,\n"
               "        ngram_range=(1, 2),\n"
               "        stop_words='english'\n"
               "    )),\n"
               "    ('clf', MultinomialNB(alpha=0.1))\n"
               "])\n"
               "\n"
               "# 교차 검증\n"
               "scores = cross_val_score(\n"
               "    pipe, data['text'], data['target'],\n"
               "    cv=5, scoring='accuracy'\n"
               ")\n"
               "print(f'CV Accuracy: {scores.mean():.4f}')\n"
               "\n"
               "# 상세 평가\n"
               "pipe.fit(x_train, y_train)\n"
               "print(classification_report(\n"
               "    y_test, pipe.predict(x_test)))")

# ════════════════════════════════════════════════════════════════
# SLIDE 33: 감사합니다 / Thank You
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG2)
add_border_frame(slide)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.0),
            "Chapter 7", font_name="Consolas", font_size=Pt(16),
            color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(2.8), Inches(11.333), Inches(1.0),
            "나이브 베이즈", font_name="Georgia", font_size=Pt(48),
            color=TITLE_GOLD, bold=True, italic=True, alignment=PP_ALIGN.CENTER)

add_decorative_line(slide, Inches(4.0), Inches(4), Inches(5.333))

add_textbox(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.8),
            "수고하셨습니다", font_name="Georgia", font_size=Pt(28),
            color=BODY_PARCHMENT, alignment=PP_ALIGN.CENTER)

add_multi_text(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(1.2), [
    ("다음 단계: 재현율(Recall), 정밀도(Precision), F1 Score → 10장", Pt(14), MUTED_TEXT, False, False),
    ("AUC (Area Under Curve) → 11장", Pt(14), MUTED_TEXT, False, False),
])

# ════════════════════════════════════════════════════════════════
# Add slide numbers
# ════════════════════════════════════════════════════════════════
total = len(prs.slides)
for i, slide in enumerate(prs.slides):
    add_slide_number(slide, i + 1, total)

# Save
output_path = r"F:\minpodata\기계학습\기계학습\7장\7장_책버전.pptx"
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Total slides: {total}")
