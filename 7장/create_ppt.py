# -*- coding: utf-8 -*-
"""
7장 나이브 베이즈 강의 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)

def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    card = add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

# ============================================================
# Slide 1: 표지
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG)
add_accent_line(slide, Inches(1.5), Inches(2.0), Inches(2.0), ACCENT_BLUE)
add_text(slide, Inches(1.5), Inches(2.1), Inches(10), Inches(0.5),
         "CHAPTER 7", font_size=20, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(1.0),
         "나이브 베이즈", font_size=48, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.6),
         "Naive Bayes", font_size=28, color=ACCENT_CYAN, bold=False)
add_accent_line(slide, Inches(1.5), Inches(4.5), Inches(8.0), ACCENT_BLUE)
add_text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.4),
         "베이즈 정리 | 조건부 독립 | 텍스트 분류 | 스팸 필터링 | 라플라스 스무딩",
         font_size=16, color=DARK_GRAY)
add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
         "참고 논문: Rish(2001), McCallum & Nigam(1998), Rennie et al.(2003), Zhang(2004), Metsis et al.(2006)",
         font_size=13, color=DARK_GRAY)
add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
         "기계학습 | 2026년 1학기", font_size=14, color=DARK_GRAY)

# ============================================================
# Slide 2: 목차
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "", "목차", "Chapter 7 Overview")

toc_left = [
    ("01", "베이즈 정리 (Bayes' Theorem)", ACCENT_BLUE),
    ("02", "빈도주의 vs 베이지안", ACCENT_CYAN),
    ("03", "나이브 가정과 최적성", ACCENT_GREEN),
    ("04", "나이브 베이즈 변형들", ACCENT_ORANGE),
    ("05", "텍스트 분류: BoW와 TF-IDF", ACCENT_PURPLE),
]
toc_right = [
    ("06", "스팸 필터링", ACCENT_RED),
    ("07", "라플라스 스무딩 & 로그 확률", ACCENT_BLUE),
    ("08", "생성 모델 vs 판별 모델", ACCENT_CYAN),
    ("09", "논문 리뷰 & 실습 소개", ACCENT_GREEN),
    ("10", "응용 사례 & 핵심 요약", ACCENT_ORANGE),
]

for idx, (num, title, clr) in enumerate(toc_left):
    y = Inches(2.2) + Inches(idx * 0.9)
    add_shape(slide, Inches(0.8), y, Inches(0.6), Inches(0.6), clr)
    add_text(slide, Inches(0.8), y + Inches(0.1), Inches(0.6), Inches(0.4),
             num, font_size=18, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.6), y + Inches(0.1), Inches(4.5), Inches(0.4),
             title, font_size=17, color=WHITE, bold=False)

for idx, (num, title, clr) in enumerate(toc_right):
    y = Inches(2.2) + Inches(idx * 0.9)
    add_shape(slide, Inches(7.0), y, Inches(0.6), Inches(0.6), clr)
    add_text(slide, Inches(7.0), y + Inches(0.1), Inches(0.6), Inches(0.4),
             num, font_size=18, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.8), y + Inches(0.1), Inches(4.5), Inches(0.4),
             title, font_size=17, color=WHITE, bold=False)

# ============================================================
# Slide 3: 7.1 베이즈 정리 - 공식과 유도
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "01", "베이즈 정리 (Bayes' Theorem)", "사전 지식을 새로운 증거로 갱신하여 사후 확률을 계산")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.5),
         "유도 과정", [
             "조건부 확률의 정의:",
             "  P(A|B) = P(A n B) / P(B)",
             "  P(B|A) = P(A n B) / P(A)",
             "두 식에서 P(A n B)를 소거하면:",
             "  P(A n B) = P(A|B)*P(B) = P(B|A)*P(A)",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_shape(slide, Inches(7.0), Inches(2.1), Inches(5.7), Inches(2.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(slide, Inches(7.2), Inches(2.2), Inches(5.3), Inches(0.4),
         "베이즈 정리 (핵심 공식)", font_size=15, color=ACCENT_CYAN, bold=True)
add_text(slide, Inches(7.2), Inches(2.8), Inches(5.3), Inches(0.8),
         "P(A|B) = P(B|A) * P(A) / P(B)", font_size=26, color=ACCENT_GREEN, bold=True,
         align=PP_ALIGN.CENTER, font_name='Consolas')
add_text(slide, Inches(7.2), Inches(3.7), Inches(5.3), Inches(0.6),
         "사후확률 = (우도 x 사전확률) / 증거", font_size=16, color=LIGHT_GRAY,
         align=PP_ALIGN.CENTER)

# 각 항의 의미 표
add_text(slide, Inches(0.6), Inches(4.9), Inches(12), Inches(0.4),
         "각 항의 의미", font_size=18, color=ACCENT_CYAN, bold=True)

terms = [
    ("사후확률 (Posterior)", "P(A|B)", "데이터 B를 관측한 후 A가 참일 확률"),
    ("우도 (Likelihood)", "P(B|A)", "A가 참일 때 데이터 B가 관측될 확률"),
    ("사전확률 (Prior)", "P(A)", "데이터 관측 전 A가 참일 확률"),
    ("증거 (Evidence)", "P(B)", "데이터 B가 관측될 전체 확률 (정규화 상수)"),
]
for idx, (term, formula, desc) in enumerate(terms):
    x_start = Inches(0.6) + Inches(idx * 3.1)
    add_shape(slide, x_start, Inches(5.4), Inches(2.9), Inches(1.7), CARD_BG, border_color=None, radius=True)
    add_text(slide, x_start + Inches(0.15), Inches(5.5), Inches(2.6), Inches(0.35),
             term, font_size=12, color=ACCENT_ORANGE, bold=True)
    add_text(slide, x_start + Inches(0.15), Inches(5.85), Inches(2.6), Inches(0.35),
             formula, font_size=16, color=WHITE, bold=True, font_name='Consolas')
    add_text(slide, x_start + Inches(0.15), Inches(6.25), Inches(2.6), Inches(0.7),
             desc, font_size=11, color=LIGHT_GRAY)

# ============================================================
# Slide 4: 의료진단 사례
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "01", "베이즈 정리 - 의료진단 사례", "직관과 다른 결과: 99% 정확도 검사의 함정")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.5), Inches(2.3),
         "상황 설정", [
             "유병률 (Prevalence): P(D) = 0.001 (0.1%)",
             "민감도 (Sensitivity): P(+|D) = 0.99 (99%)",
             "특이도 (Specificity): P(-|~D) = 0.95 (95%)",
             "",
             "Q: 양성 판정 시 실제 질병 확률은?",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(slide, Inches(6.5), Inches(2.1), Inches(6.2), Inches(2.3),
         "베이즈 정리 적용", [
             "P(D|+) = P(+|D)*P(D) / P(+)",
             "P(+) = P(+|D)*P(D) + P(+|~D)*P(~D)",
             "     = 0.99*0.001 + 0.05*0.999 = 0.05094",
             "P(D|+) = 0.00099 / 0.05094",
             "       = 0.0194  (약 1.94%)",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_shape(slide, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.4), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.4),
         "핵심 교훈", font_size=17, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.6), [
    "99% 정확도 검사로 양성이 나와도, 유병률이 낮으면 실제 질병 확률은 약 2%에 불과",
    "사전확률(유병률)이 사후확률에 결정적 영향을 미친다",
    "이것이 베이즈 정리가 직관과 다른 결과를 내는 대표적 사례",
    "스크리닝 검사에서 '양성'이라고 해서 바로 확진이 아닌 이유",
], font_size=14, color=LIGHT_GRAY)

# ============================================================
# Slide 5: 스팸 메일 사례
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "01", "베이즈 정리 - 스팸 메일 사례", '"free"가 포함된 이메일이 스팸일 확률은?')

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.5), Inches(2.5),
         "학습 데이터", [
             "전체 메일: 1,000통",
             "  스팸: 400통 (40%), 'free' 포함 200통 (50%)",
             "  정상: 600통 (60%), 'free' 포함 60통 (10%)",
             "  'free' 포함 메일 합계: 260통 (26%)",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_shape(slide, Inches(6.5), Inches(2.1), Inches(6.2), Inches(2.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(slide, Inches(6.7), Inches(2.2), Inches(5.8), Inches(0.4),
         "베이즈 정리 적용", font_size=15, color=ACCENT_GREEN, bold=True)
add_text(slide, Inches(6.7), Inches(2.7), Inches(5.8), Inches(0.5),
         'P(spam|"free") = P("free"|spam)*P(spam) / P("free")',
         font_size=15, color=WHITE, bold=True, font_name='Consolas')
add_text(slide, Inches(6.7), Inches(3.3), Inches(5.8), Inches(0.5),
         '= 0.5 * 0.4 / 0.26 = 0.769', font_size=18, color=ACCENT_GREEN, bold=True, font_name='Consolas')
add_text(slide, Inches(6.7), Inches(3.9), Inches(5.8), Inches(0.5),
         '"free"가 포함된 이메일이 스팸일 확률: 약 76.9%',
         font_size=14, color=LIGHT_GRAY)

add_card(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.1),
         "나이브 베이즈 분류기의 핵심 아이디어", [
             "단어들의 출현 패턴을 기반으로 각 클래스의 사후 확률을 계산하고, 가장 높은 클래스로 분류",
             "P(spam) = 사전확률 (전체 메일 중 스팸 비율)",
             'P("free"|spam) = 우도 (스팸에서 "free"가 등장할 확률)',
             "분모 P(B)는 모든 클래스에 동일하므로 argmax 계산 시 무시 가능",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# Slide 6: 빈도주의 vs 베이지안
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "02", "빈도주의 vs 베이지안", "확률을 해석하는 두 가지 철학적 관점")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(3.2),
         "빈도주의 (Frequentist)", [
             "확률 = 동일 실험의 무한 반복 시 상대적 빈도",
             "파라미터: 고정된 미지의 상수",
             "추론: MLE, 신뢰구간, p-값",
             "예: '이 동전을 10,000번 던지면 앞면 약 5,000번'",
             "",
             "추정식: theta_MLE = argmax P(D|theta)",
             "",
             "장점: 객관적, 재현 가능",
             "단점: 사전 정보 활용 불가, 소량 데이터 불안정",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(3.2),
         "베이지안 (Bayesian)", [
             "확률 = 불확실성에 대한 주관적 믿음의 정도",
             "파라미터: 확률분포를 가지는 확률변수",
             "추론: 베이즈 정리로 사전분포 -> 사후분포 갱신",
             "예: '이 동전이 공정할 확률이 80%라고 믿는다'",
             "",
             "추정식: P(theta|D) = P(D|theta)*P(theta) / P(D)",
             "",
             "장점: 사전 정보 반영, 소량 데이터에서 강건",
             "단점: 사전분포 선택의 주관성",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_shape(slide, Inches(0.6), Inches(5.6), Inches(12.0), Inches(1.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(slide, Inches(0.8), Inches(5.7), Inches(11.6), Inches(0.35),
         "나이브 베이즈 분류기의 위치", font_size=15, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(6.1), Inches(11.6), Inches(0.9), [
    "이름 그대로 베이지안 접근에 기반하되, 실제 구현에서는 빈도 기반으로 파라미터를 추정하는 혼합적 접근",
    "대량 데이터: 빈도주의와 베이지안 모두 수렴 | 소량 데이터: 사전분포가 보완 역할",
], font_size=13, color=LIGHT_GRAY)

# ============================================================
# Slide 7: 나이브 가정 - 조건부 독립
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "03", "나이브 가정 (Naive Assumption)", "모든 특성이 클래스가 주어졌을 때 조건부 독립")

add_shape(slide, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.8), CARD_BG, ACCENT_BLUE, radius=True)
add_text(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.35),
         "조건부 독립 가정의 수학적 정의", font_size=15, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.5),
         "P(x1, x2, ..., xn | C) = P(x1|C) * P(x2|C) * ... * P(xn|C)", font_size=18, color=WHITE, bold=True, font_name='Consolas')
add_text(slide, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.4),
         "분류 규칙:  C_hat = argmax_C [ P(C) * PROD_i P(xi|C) ]      (분모 P(x)는 모든 클래스에 동일하므로 무시)",
         font_size=14, color=LIGHT_GRAY, font_name='Consolas')

add_card(slide, Inches(0.6), Inches(4.2), Inches(5.8), Inches(1.8),
         "왜 현실에서 독립 가정은 깨지는가?", [
             '"machine" + "learning" -> 함께 등장할 확률 높음',
             '"New" + "York" -> 함께 등장할 확률 높음',
             '"free" + "win" + "prize" -> 스팸에서 동시 등장',
             "텍스트 데이터에서 단어들은 명백히 서로 독립이 아님",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(slide, Inches(6.8), Inches(4.2), Inches(5.9), Inches(1.8),
         "그럼에도 잘 작동하는 3가지 이유", [
             "1. 분류 vs 확률 추정 구분 (Rish 2001): 상대적 순위만 정확하면 OK",
             "2. 의존성의 상쇄 효과: 양/음의 상관관계가 서로 상쇄",
             "3. 계산 효율성: 파라미터 수 O(n*k)로 과적합 위험 감소",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_shape(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.9), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.35),
         "Zhang(2004)의 최적성 정리", font_size=14, color=ACCENT_PURPLE, bold=True)
add_text(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.35),
         "의존성의 총량이 클래스 간에 동일하게 분포하면, 의존성이 아무리 강해도 나이브 베이즈는 베이즈 최적 분류기와 동일한 결과를 낸다",
         font_size=12, color=LIGHT_GRAY)

# ============================================================
# Slide 8: 나이브 베이즈 변형들 (1) - Gaussian & Multinomial
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "04", "나이브 베이즈 변형들 (1/2)", "GaussianNB와 MultinomialNB")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.8),
         "GaussianNB - 가우시안 나이브 베이즈", [
             "적용 대상: 연속형(continuous) 실수 데이터",
             "",
             "분포 가정: 각 특성이 클래스가 주어졌을 때 정규분포",
             "P(xi|C) = (1/sqrt(2*pi*sigma^2))",
             "         * exp(-(xi-mu)^2 / (2*sigma^2))",
             "",
             "파라미터 추정:",
             "  mu = 클래스 C에서 특성 i의 표본 평균",
             "  sigma^2 = 클래스 C에서 특성 i의 표본 분산",
             "",
             "사용 예: Iris 분류, 키/몸무게 기반 성별 예측",
             "코드: GaussianNB().fit(X_train, y_train)",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.8),
         "MultinomialNB - 다항 나이브 베이즈", [
             "적용 대상: 카운트(count) 또는 빈도 데이터",
             "",
             "분포 가정: 각 특성이 다항분포를 따름",
             "P(xi|C) = (count(xi,C) + alpha)",
             "        / (SUM count(xj,C) + alpha*|V|)",
             "",
             "핵심: 단어의 출현 빈도를 기반으로 확률 계산",
             "McCallum & Nigam(1998): 텍스트 분류에서",
             "BernoulliNB보다 우수함을 증명",
             "",
             "사용 예: 스팸 필터링, 감성 분석, 문서 분류",
             "코드: MultinomialNB(alpha=1.0)",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# Slide 9: 나이브 베이즈 변형들 (2) - Bernoulli & Complement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "04", "나이브 베이즈 변형들 (2/2)", "BernoulliNB와 ComplementNB")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(3.0),
         "BernoulliNB - 베르누이 나이브 베이즈", [
             "적용 대상: 이진(binary) 데이터 (0 또는 1)",
             "분포 가정: 각 특성이 베르누이 분포",
             "P(x|C) = PROD P(xi|C)^xi * (1-P(xi|C))^(1-xi)",
             "",
             "핵심: 단어의 존재 여부(있다/없다)만 고려",
             "단어의 부재 정보도 명시적으로 활용 (vs MultinomialNB)",
             "McCallum(1998): 어휘 작을 때 경쟁력, 클수록 열세",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(3.0),
         "ComplementNB - 보완 나이브 베이즈", [
             "제안: Rennie et al. (2003), ICML",
             "핵심: 보완 클래스(complement class)의 데이터로 추정",
             "c_hat = argmin_c SUM fi * log(theta_ci)",
             "",
             "장점:",
             "  클래스 불균형 문제 완화",
             "  MultinomialNB 대비 약 5~8% 정확도 향상",
             "  SVM에 근접하는 성능 달성",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# 비교 요약 표
add_text(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.4),
         "변형 비교 요약", font_size=17, color=ACCENT_CYAN, bold=True)

variants = [
    ("GaussianNB", "연속형 (실수)", "정규분포", "수치형 특성 분류", ACCENT_BLUE),
    ("MultinomialNB", "카운트 (정수)", "다항분포", "텍스트 분류 (빈도)", ACCENT_ORANGE),
    ("BernoulliNB", "이진 (0/1)", "베르누이분포", "텍스트 분류 (존재)", ACCENT_GREEN),
    ("ComplementNB", "카운트/빈도", "보완 다항분포", "불균형 텍스트 분류", ACCENT_PURPLE),
]
headers = ["종류", "데이터 유형", "분포 가정", "대표 사용 사례"]
col_widths = [Inches(2.2), Inches(2.5), Inches(2.8), Inches(4.3)]
x_positions = [Inches(0.8)]
for w in col_widths[:-1]:
    x_positions.append(x_positions[-1] + w)

y_header = Inches(5.85)
for i, h in enumerate(headers):
    add_text(slide, x_positions[i], y_header, col_widths[i], Inches(0.3),
             h, font_size=11, color=ACCENT_CYAN, bold=True)

for row_idx, (name, dtype, dist, usecase, clr) in enumerate(variants):
    y_row = Inches(6.15) + Inches(row_idx * 0.3)
    vals = [name, dtype, dist, usecase]
    for col_idx, val in enumerate(vals):
        c = clr if col_idx == 0 else LIGHT_GRAY
        add_text(slide, x_positions[col_idx], y_row, col_widths[col_idx], Inches(0.3),
                 val, font_size=10, color=c, bold=(col_idx == 0))

# ============================================================
# Slide 10: 텍스트 분류 - BoW와 CountVectorizer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "05", "텍스트 분류 - Bag of Words", "텍스트를 단어의 출현 빈도로 표현하는 벡터화 방법")

add_card(slide, Inches(0.6), Inches(2.1), Inches(6.0), Inches(2.5),
         "Bag of Words (BoW) 모델", [
             "텍스트를 단어의 출현 빈도로 표현",
             "단어의 순서를 무시하고 빈도만 고려",
             "",
             "CountVectorizer 동작:",
             "  1. 전체 문서에서 어휘 사전(Vocabulary) 구축",
             "  2. 각 문서를 문서-단어 행렬(DTM)로 변환",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(7.0), Inches(2.1), Inches(5.7), Inches(2.5),
         "CountVectorizer 예시", [
             '문서1: "free money free" -> [0, 2, 0, 1, 0]',
             '문서2: "money transfer"  -> [0, 0, 0, 1, 1]',
             '문서3: "free gift card"  -> [1, 1, 1, 0, 0]',
             "",
             "어휘: [card, free, gift, money, transfer]",
             "출력: 희소 행렬 (Sparse Matrix) 형태",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.2),
         "코드 예시 & 희소 행렬", [
             "from sklearn.feature_extraction.text import CountVectorizer",
             "cv = CountVectorizer()",
             "X = cv.fit_transform(corpus)    # 어휘 구축 + 변환 동시 수행",
             "print(cv.vocabulary_)            # 단어-인덱스 매핑 딕셔너리",
             "print(X.toarray())              # 희소 -> 밀집 행렬 변환",
             "",
             "희소 행렬: 대부분의 값이 0인 행렬에서 0이 아닌 값의 위치와 값만 저장 -> 메모리 효율적",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# Slide 11: TF-IDF
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "05", "TF-IDF (Term Frequency - Inverse Document Frequency)", "단어의 중요도를 반영한 가중 벡터화")

add_shape(slide, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.3), CARD_BG, ACCENT_BLUE, radius=True)
add_text(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.35),
         "TF-IDF 핵심 수식", font_size=15, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.5),
         "TF-IDF(t,d) = TF(t,d) x IDF(t)     |     TF = f(t,d) / SUM f(t\',d)     |     IDF = log(N / |{d: t in d}|)",
         font_size=16, color=WHITE, bold=True, font_name='Consolas')

add_card(slide, Inches(0.6), Inches(3.7), Inches(3.8), Inches(1.8),
         "TF (Term Frequency)", [
             "특정 문서 d에서 단어 t가",
             "등장한 빈도 (정규화)",
             "높을수록 해당 문서에서 중요",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(4.7), Inches(3.7), Inches(3.8), Inches(1.8),
         "IDF (Inverse Document Freq)", [
             "단어 t가 전체 문서에서",
             "얼마나 드물게 등장하는지",
             "드물수록(IDF 높음) 중요한 단어",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(8.8), Inches(3.7), Inches(3.9), Inches(1.8),
         "핵심 아이디어", [
             "자주 등장(TF높) + 드물게(IDF높)",
             "  => 중요한 단어",
             "모든 문서에 공통(IDF낮)",
             "  => 덜 중요한 단어 ('the','is')",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# CountVec vs TfidfVec 비교
add_text(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.4),
         "CountVectorizer vs TfidfVectorizer", font_size=16, color=ACCENT_CYAN, bold=True)

comparisons = [
    ("가중치", "단순 빈도 (정수)", "TF-IDF 가중치 (실수)"),
    ("흔한 단어", "높은 값 그대로", "IDF로 가중치 감소"),
    ("희귀 단어", "낮은 값 그대로", "IDF로 가중치 증가"),
    ("일반 성능", "기본적", "대부분의 경우 우수"),
]
for idx, (aspect, cv_val, tfidf_val) in enumerate(comparisons):
    y = Inches(6.2) + Inches(idx * 0.28)
    add_text(slide, Inches(0.8), y, Inches(2.0), Inches(0.28),
             aspect, font_size=11, color=ACCENT_ORANGE, bold=True)
    add_text(slide, Inches(3.0), y, Inches(4.0), Inches(0.28),
             cv_val, font_size=11, color=LIGHT_GRAY)
    add_text(slide, Inches(7.5), y, Inches(5.0), Inches(0.28),
             tfidf_val, font_size=11, color=ACCENT_GREEN)

# ============================================================
# Slide 12: 스팸 필터링 역사 및 Metsis(2006)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "06", "스팸 필터링", "SpamBayes에서 Metsis(2006) 벤치마크까지")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.0),
         "역사: SpamBayes와 초기 베이지안 필터", [
             "2002년 Paul Graham 'A Plan for Spam' 에세이",
             "나이브 베이즈로 스팸 99.5% 이상 정확 분류",
             "SpamBayes: Python 기반 오픈소스 프로젝트",
             "사용자별 학습 가능, 3범주 분류 (스팸/정상/불확실)",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.0),
         "비용 민감 평가", [
             "FP (정상->스팸): 중요 메일 손실 - 매우 심각!",
             "FN (스팸->정상): 스팸이 받은편지함에 - 불편",
             "",
             "WAcc = (lambda*TN + TP) / (lambda*(TN+FP) + (TP+FN))",
             "lambda > 1: FP의 상대적 비용",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_shape(slide, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.8), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.35),
         "Metsis et al. (2006) 벤치마크 결과", font_size=15, color=ACCENT_ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.3),
         "5가지 나이브 베이즈 변형을 6개 스팸 데이터셋(Ling-Spam, SpamAssassin, Enron 등)에서 비교",
         font_size=12, color=LIGHT_GRAY)

results_data = [
    ("MNB (단어 빈도 기반)", "0.961", LIGHT_GRAY),
    ("MBNB (단어 존재/부재)", "0.943", LIGHT_GRAY),
    ("MNBB (빈도 이진화 + MNB) -- 최우수", "0.972", ACCENT_GREEN),
    ("MNB + TF-IDF", "0.965", LIGHT_GRAY),
    ("Flexible Bayes", "0.938", LIGHT_GRAY),
]
for idx, (name, auc_val, clr) in enumerate(results_data):
    y = Inches(5.3) + Inches(idx * 0.32)
    add_text(slide, Inches(1.0), y, Inches(7.0), Inches(0.3),
             name, font_size=13, color=clr, bold=(clr == ACCENT_GREEN))
    add_text(slide, Inches(8.5), y, Inches(2.0), Inches(0.3),
             f"AUC: {auc_val}", font_size=13, color=clr, bold=(clr == ACCENT_GREEN),
             font_name='Consolas')

add_text(slide, Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.3),
         "핵심: 단어의 빈도보다 존재 여부가 중요하지만, BernoulliNB보다 MultinomialNB의 확률 모델이 더 적합",
         font_size=12, color=ACCENT_CYAN, bold=True)

# ============================================================
# Slide 13: 라플라스 스무딩
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "07", "라플라스 스무딩 (Laplace Smoothing)", "제로 확률 문제의 해결")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.2),
         "제로 확률 문제 (Zero Probability)", [
             '학습 데이터에서 "scholarship"이 한번도 등장 안함:',
             '  P("scholarship"|spam) = 0',
             "",
             "나이브 베이즈는 확률을 곱셈으로 결합하므로",
             "하나라도 0이면 전체 결과가 0!",
             "다른 모든 단어가 스팸을 가리켜도 결과 = 0",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_shape(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.2), CARD_BG, ACCENT_GREEN, radius=True)
add_text(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(0.35),
         "라플라스 스무딩 수식", font_size=15, color=ACCENT_GREEN, bold=True)
add_text(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(0.5),
         "P(xi|C) = (count(xi,C) + alpha)\n         / (count(C) + alpha * |V|)",
         font_size=15, color=WHITE, bold=True, font_name='Consolas')
add_bullet_list(slide, Inches(7.0), Inches(3.4), Inches(5.5), Inches(0.8), [
    "count(xi,C): 클래스 C에서 단어 xi의 출현 횟수",
    "alpha: 스무딩 파라미터, |V|: 어휘 크기",
], font_size=11, color=LIGHT_GRAY)

# alpha 값 비교
add_text(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(0.35),
         "alpha 값에 따른 스무딩 종류", font_size=16, color=ACCENT_CYAN, bold=True)

alpha_items = [
    ("alpha = 1", "라플라스 스무딩", "표준 (기본값)", ACCENT_BLUE),
    ("0 < alpha < 1", "리드스톤 스무딩", "약한 스무딩", ACCENT_GREEN),
    ("alpha > 1", "강한 스무딩", "분포를 더 균일하게", ACCENT_ORANGE),
]
for idx, (alpha, name, effect, clr) in enumerate(alpha_items):
    x = Inches(0.6) + Inches(idx * 4.2)
    add_shape(slide, x, Inches(5.05), Inches(3.8), Inches(0.9), CARD_BG, clr, radius=True)
    add_text(slide, x + Inches(0.15), Inches(5.1), Inches(3.5), Inches(0.3),
             f"{alpha}  ({name})", font_size=13, color=clr, bold=True)
    add_text(slide, x + Inches(0.15), Inches(5.4), Inches(3.5), Inches(0.4),
             effect, font_size=12, color=LIGHT_GRAY)

add_card(slide, Inches(0.6), Inches(6.2), Inches(12.1), Inches(1.0),
         "sklearn 코드", [
             "MultinomialNB(alpha=1.0)  # 라플라스    |    MultinomialNB(alpha=0.5)  # 리드스톤    |    MultinomialNB(alpha=0.01)  # 약한 스무딩",
         ], title_color=ACCENT_PURPLE)

# ============================================================
# Slide 14: 로그 확률 트릭
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "07", "로그 확률 트릭 (Log Probability Trick)", "수치적 안정성을 위한 핵심 기법")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.2),
         "언더플로(Underflow) 문제", [
             "많은 확률값(0~1 사이)을 곱하면:",
             "P(C) * P(x1|C) * P(x2|C) * ...",
             "= 0.001 * 0.01 * 0.005 * ... -> 0",
             "",
             "컴퓨터가 표현할 수 있는 최소값보다 작아짐",
             "결과가 0으로 수렴하여 비교 불가",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_shape(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.2), CARD_BG, ACCENT_GREEN, radius=True)
add_text(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(0.35),
         "해결: 곱셈 -> 로그 공간에서 덧셈", font_size=15, color=ACCENT_GREEN, bold=True)
add_text(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(0.5),
         "log P(C|x) ~ log P(C) + SUM log P(xi|C)", font_size=16, color=WHITE, bold=True, font_name='Consolas')
add_text(slide, Inches(7.0), Inches(3.3), Inches(5.5), Inches(0.5),
         "C_hat = argmax [log P(C) + SUM log P(xi|C)]", font_size=15, color=ACCENT_CYAN, bold=True, font_name='Consolas')
add_text(slide, Inches(7.0), Inches(3.8), Inches(5.5), Inches(0.3),
         "argmax는 로그 변환에 의해 보존 (단조증가함수)", font_size=12, color=LIGHT_GRAY)

add_card(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.5),
         "로그 확률 트릭의 3가지 장점", [
             "1. 수치적 안정성 확보: 언더플로 방지. log(0.001) = -6.9, log(0.01) = -4.6 -> 합산으로 처리 가능",
             "2. 계산 효율 향상: 곱셈(느림) -> 덧셈(빠름)으로 변환",
             "3. argmax 보존: log는 단조증가함수이므로 최대값의 위치가 동일하게 유지됨",
             "",
             "구현 시 추가 팁: log-sum-exp trick으로 정규화된 확률도 안정적으로 계산 가능",
             "  max_log = max(log_posteriors);  posteriors = exp(log_posteriors - max_log);  posteriors /= sum(posteriors)",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# Slide 15: 생성 모델 vs 판별 모델
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "08", "생성 모델 vs 판별 모델", "나이브 베이즈와 로지스틱 회귀의 관계")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.5),
         "생성 모델 (Generative) - 나이브 베이즈", [
             "모델링: P(X,Y) = P(X|Y) * P(Y)",
             "클래스별 데이터 분포를 학습",
             "특성의 분포를 가정 (가우시안, 다항 등)",
             "",
             "장점: 소량 데이터에서 강건, 수렴이 빠름",
             "단점: 분포 가정이 맞지 않으면 성능 저하",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.5),
         "판별 모델 (Discriminative) - 로지스틱 회귀", [
             "모델링: P(Y|X)를 직접 모델링",
             "결정 경계를 직접 학습",
             "분포 가정 불필요",
             "",
             "장점: 대량 데이터에서 점근적으로 우수",
             "단점: 소량 데이터에서 과적합 위험",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_shape(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.2), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.35),
         "수학적 연결 - Ng & Jordan(2001)", font_size=14, color=ACCENT_PURPLE, bold=True)
add_text(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
         "NB 로그 오즈: log P(Y=1|x)/P(Y=0|x) = log P(Y=1)/P(Y=0) + SUM log P(xi|Y=1)/P(xi|Y=0)  =  w^T x + b  (로지스틱 회귀와 동일 형태)",
         font_size=12, color=LIGHT_GRAY, font_name='Consolas')

# 실용적 선택 기준
add_text(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.35),
         "실용적 선택 기준", font_size=15, color=ACCENT_CYAN, bold=True)
criteria = [
    ("학습 데이터 적음", "나이브 베이즈"),
    ("학습 데이터 충분", "로지스틱 회귀"),
    ("특성 독립성 높음", "나이브 베이즈"),
    ("특성 상관관계 강함", "로지스틱 회귀"),
    ("속도 중요", "나이브 베이즈"),
]
for idx, (situation, rec) in enumerate(criteria):
    x = Inches(0.8) + Inches(idx * 2.5)
    add_text(slide, x, Inches(6.65), Inches(2.3), Inches(0.25),
             situation, font_size=10, color=LIGHT_GRAY)
    add_text(slide, x, Inches(6.9), Inches(2.3), Inches(0.25),
             f"-> {rec}", font_size=10, color=ACCENT_GREEN, bold=True)

# ============================================================
# Slide 16: 논문 리뷰 (1) - Rish(2001) & McCallum(1998)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "09", "논문 리뷰 (1/2)", "Rish(2001)와 McCallum & Nigam(1998)")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.8),
         "Rish (2001) - IJCAI Workshop", [
             '"An Empirical Study of the Naive Bayes Classifier"',
             "",
             "핵심 기여: 조건부 독립 위반 시에도 NB가",
             "잘 작동하는 이유를 경험적으로 분석",
             "",
             "주요 발견:",
             "1. 의존성의 방향이 중요 - 양/음 상관관계 상쇄",
             "2. 균일한 의존성은 해롭지 않음 - 순위 보존",
             "3. 분류 vs 확률 추정 구분",
             "   -> 확률이 부정확해도 순위가 맞으면 OK",
             "",
             "한계: 소규모 합성 데이터 위주, 연속형 분석 부족",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.8),
         "McCallum & Nigam (1998) - AAAI Workshop", [
             '"A Comparison of Event Models for NB Text Classification"',
             "",
             "핵심 기여: 두 가지 이벤트 모델을 비교",
             "  MultinomialNB vs BernoulliNB",
             "",
             "20 Newsgroups 실험 결과:",
             "  어휘 100:    Bernoulli 0.68 vs Multi 0.65",
             "  어휘 1,000:  Bernoulli 0.78 vs Multi 0.82",
             "  어휘 10,000: Bernoulli 0.81 vs Multi 0.87",
             "  전체:        Bernoulli 0.82 vs Multi 0.89",
             "",
             "결론: 어휘 클수록 MultinomialNB가 일관되게 우수",
             "sklearn의 별도 클래스 제공 근거",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# Slide 17: 논문 리뷰 (2) - Rennie(2003) & Zhang(2004)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "09", "논문 리뷰 (2/2)", "Rennie et al.(2003)와 Zhang(2004)")

add_card(slide, Inches(0.6), Inches(2.1), Inches(7.5), Inches(4.8),
         "Rennie et al. (2003) - ICML", [
             '"Tackling the Poor Assumptions of Naive Bayes Text Classifiers"',
             "",
             "5가지 잘못된 가정 식별 및 해결:",
             "  1. 조건부 독립       -> CNB(Complement NB)로 완화",
             "  2. 균일한 특성 중요도 -> TF-IDF 가중치",
             "  3. 클래스 불균형     -> 보완 클래스 사용",
             "  4. 문서 길이 영향    -> L2 정규화",
             "  5. 버스티니스       -> 로그 TF 변환",
             "",
             "20 Newsgroups 실험 결과:",
             "  MultinomialNB(기본): 0.774",
             "  CNB:                0.819",
             "  WCNB(전체 개선):     0.847",
             "  SVM(선형):          0.864  -- SVM에 근접!",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(8.5), Inches(2.1), Inches(4.2), Inches(2.2),
         "Zhang (2004) - FLAIRS", [
             '"The Optimality of Naive Bayes"',
             "",
             "핵심 정리:",
             "이진 분류에서 의존성 총량이",
             "클래스 간에 동일하게 분포하면",
             "NB는 베이즈 최적 분류기와 동일",
         ], title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(slide, Inches(8.5), Inches(4.6), Inches(4.2), Inches(2.3),
         "Metsis et al. (2006) - CEAS", [
             '"Spam Filtering with NB"',
             '"Which Naive Bayes?"',
             "",
             "6개 스팸 데이터셋에서",
             "5가지 NB 변형 비교",
             "MNBB(MNB+Boolean)가 최우수",
             "AUC: 0.972",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# Slide 18: 실습 소개 (1) - GaussianNB 스크래치
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "09", "실습: GaussianNB 스크래치 구현", "밑바닥부터 구현하고 sklearn과 비교")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.5),
         "구현 핵심: GaussianNBScratch 클래스", [
             "fit(): 클래스별 사전확률, 평균, 분산 추정",
             "  P(C) = N_c / N",
             "  mu = 표본 평균, sigma^2 = 표본 분산 + 1e-9",
             "",
             "_gaussian_log_likelihood():",
             "  log P(x|mu,sigma^2) = -0.5*log(2pi*var)",
             "                       - (x-mu)^2/(2*var)",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.5),
         "예측 및 확률 계산", [
             "_compute_log_posterior():",
             "  log P(c|x) ~ log P(c) + SUM log P(xi|c)",
             "",
             "predict(): argmax로 클래스 선택",
             "",
             "predict_proba(): log-sum-exp trick",
             "  max_log = max(log_posteriors)",
             "  posteriors = exp(log_post - max_log) / sum",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.2),
         "실험 결과: Iris & Wine 데이터셋", [
             "Iris (3클래스, 4특성): 직접 구현과 sklearn의 정확도 및 예측 결과가 동일하게 일치",
             "Wine (3클래스, 13특성): 대규모 특성에서도 동일한 결과 확인",
             "",
             "핵심 포인트:",
             "  1. 분산에 epsilon(1e-9)을 더해 수치 안정성 확보 (0으로 나누기 방지)",
             "  2. 로그 공간에서 계산하여 언더플로 방지",
             "  3. log-sum-exp trick으로 정규화된 확률을 안정적으로 계산",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# Slide 19: 실습 소개 (2) - 텍스트 분류 비교
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "09", "실습: 텍스트 분류 & 스팸 필터", "MultinomialNB vs BernoulliNB vs ComplementNB")

add_card(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.3),
         "실습 2: 텍스트 분류 비교 (20 Newsgroups)", [
             "3가지 NB 변형 x 2가지 벡터화기 = 6가지 조합",
             "Alpha 하이퍼파라미터 튜닝 (0.001 ~ 10.0)",
             "어휘 크기별 성능 비교 (100 ~ 전체)",
             "",
             "결론: ComplementNB + TfidfVectorizer가 최우수",
             "어휘 클수록 MultinomialNB 이점 증가 (McCallum 재현)",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.3),
         "실습 3: 스팸 필터 구현", [
             "TF-IDF + 나이브 베이즈 완전한 스팸 필터",
             "전처리: 금액->MONEY, 퍼센트->PERCENT 토큰화",
             "sublinear_tf=True (Rennie 2003 권장)",
             "ngram_range=(1,2) 바이그램 포함",
             "",
             "특성 중요도: log P(w|spam) - log P(w|ham)",
             "  스팸 지표: free, win, prize, claim, now",
         ], title_color=ACCENT_RED, border=ACCENT_RED)

add_card(slide, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.5),
         "전체 텍스트 분류 파이프라인", [
             "[텍스트 데이터] -> [전처리: 특수기호 제거 -> 소문자 변환 -> 불용어 제거]",
             "    -> [벡터화: CountVectorizer 또는 TfidfVectorizer]",
             "    -> [모델링: MultinomialNB / ComplementNB (alpha 튜닝)]",
             "    -> [평가: 정확도, 정밀도, 재현율, F1, 혼동행렬]",
             "",
             "전처리 함수: remove_punc() -> string.punctuation으로 특수기호 제거",
             "            stop_words() -> NLTK stopwords로 불용어 제거 & 소문자 변환",
             "벡터화 출력: 희소 행렬 (Sparse Matrix) - 0이 아닌 값만 저장하여 메모리 효율적",
         ], title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# Slide 20: 응용 사례
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "10", "나이브 베이즈 응용 사례", "감정분석, 문서분류, 의료진단")

add_card(slide, Inches(0.6), Inches(2.1), Inches(3.8), Inches(4.6),
         "감정분석 (Sentiment Analysis)", [
             "텍스트에서 감정/태도/의견 파악",
             "",
             "종류:",
             "  이진: 긍정/부정",
             "  다중: 긍정/중립/부정",
             "  점수: -1.0 ~ +1.0",
             "",
             "활용: 마케팅, 금융(주가예측),",
             "고객서비스, 정치(여론분석)",
             "",
             "NB 장점: 감성 키워드 빈도 기반,",
             "빠른 처리, 적은 데이터 OK",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(4.7), Inches(2.1), Inches(3.8), Inches(4.6),
         "문서 분류 (Document Classification)", [
             "문서를 카테고리에 자동 할당",
             "",
             "분야:",
             "  뉴스: 정치/경제/스포츠/문화",
             "  이메일: 업무/개인/프로모션",
             "  고객문의: 배송/반품/결제",
             "  법률: 계약서/소송/특허",
             "",
             "NB 적합 이유:",
             "  고차원 효율적 처리",
             "  실시간 대량 분류 가능",
             "  다중 클래스 자연 지원",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(slide, Inches(8.8), Inches(2.1), Inches(3.9), Inches(4.6),
         "의료진단 (Medical Diagnosis)", [
             "증상 기반 질병 진단",
             "P(D|x1,...,xn) ~ P(D) PROD P(xi|D)",
             "",
             "장점:",
             "  소규모 의료 데이터에 강건",
             "  각 증상의 기여도를 확률 해석",
             "  -> 의료진에게 설명 가능",
             "  여러 질병 후보를 확률 순 제시",
             "",
             "사례: 유방암 진단, 심장병 예측,",
             "당뇨병 분류 등에서 SVM/신경망과",
             "비슷한 성능 + 해석 용이",
         ], title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# Slide 21: 핵심 요약
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "10", "핵심 요약", "Chapter 7 Key Takeaways")

summary_items = [
    ("베이즈 정리", "P(A|B) = P(B|A)*P(A)/P(B)", ACCENT_BLUE),
    ("나이브 가정", "모든 특성이 클래스에 대해 조건부 독립", ACCENT_CYAN),
    ("GaussianNB", "연속형 데이터, 정규분포 가정", ACCENT_GREEN),
    ("MultinomialNB", "텍스트 분류에 최적, 단어 빈도 기반", ACCENT_ORANGE),
    ("BernoulliNB", "단어 존재/부재 기반, 소규모 어휘 경쟁력", ACCENT_PURPLE),
    ("ComplementNB", "Rennie(2003), 보완 클래스, SVM에 근접", ACCENT_RED),
    ("라플라스 스무딩", "제로 확률 문제 해결, alpha 파라미터", ACCENT_BLUE),
    ("로그 확률 트릭", "언더플로 방지, 곱셈을 덧셈으로 변환", ACCENT_CYAN),
    ("CountVectorizer", "텍스트를 단어 출현 빈도 벡터로 변환", ACCENT_GREEN),
    ("TF-IDF", "단어 중요도를 반영한 가중 벡터화", ACCENT_ORANGE),
    ("생성 모델 (NB)", "P(X|Y)*P(Y)를 모델링", ACCENT_PURPLE),
    ("판별 모델 (LR)", "P(Y|X)를 직접 모델링", ACCENT_RED),
]

for idx, (concept, desc, clr) in enumerate(summary_items):
    col = idx % 2
    row = idx // 2
    x = Inches(0.6) + Inches(col * 6.3)
    y = Inches(2.1) + Inches(row * 0.8)
    add_shape(slide, x, y, Inches(5.9), Inches(0.7), CARD_BG, clr, radius=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.2), Inches(0.3),
             concept, font_size=13, color=clr, bold=True)
    add_text(slide, x + Inches(2.4), y + Inches(0.05), Inches(3.3), Inches(0.55),
             desc, font_size=11, color=LIGHT_GRAY)

# ============================================================
# Slide 22: 핵심 수식
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "10", "핵심 수식 정리", "Chapter 7 Key Formulas")

formulas = [
    ("베이즈 정리", "P(A|B) = P(B|A) * P(A) / P(B)", ACCENT_BLUE),
    ("나이브 베이즈 분류", "C_hat = argmax_C [ P(C) * PROD P(xi|C) ]", ACCENT_CYAN),
    ("가우시안 우도", "P(xi|C) = (1/sqrt(2pi*var)) * exp(-(xi-mu)^2 / (2*var))", ACCENT_GREEN),
    ("다항 우도", "P(xi|C) = (count(xi,C) + alpha) / (count(C) + alpha*|V|)", ACCENT_ORANGE),
    ("TF-IDF", "TF-IDF(t,d) = TF(t,d) * IDF(t) = TF * log(N/df)", ACCENT_PURPLE),
    ("로그 확률", "C_hat = argmax [ log P(C) + SUM log P(xi|C) ]", ACCENT_RED),
    ("로그 오즈 (NB=LR)", "log P(Y=1|x)/P(Y=0|x) = w^T x + b", ACCENT_BLUE),
    ("정밀도", "Precision = TP / (TP + FP)", ACCENT_GREEN),
    ("재현율", "Recall = TP / (TP + FN)", ACCENT_ORANGE),
    ("F1 Score", "F1 = 2 * Precision * Recall / (Precision + Recall)", ACCENT_CYAN),
]

for idx, (name, formula, clr) in enumerate(formulas):
    col = idx % 2
    row = idx // 2
    x = Inches(0.6) + Inches(col * 6.3)
    y = Inches(2.1) + Inches(row * 1.0)
    add_shape(slide, x, y, Inches(5.9), Inches(0.85), CARD_BG, clr, radius=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(5.6), Inches(0.3),
             name, font_size=12, color=clr, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.35), Inches(5.6), Inches(0.4),
             formula, font_size=13, color=WHITE, font_name='Consolas')

# ============================================================
# Slide 23: 알고리즘 선택 가이드
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "10", "알고리즘 선택 가이드", "상황별 최적 알고리즘 추천")

guide_items = [
    ("텍스트 분류 (빈도 기반)", "MultinomialNB + CountVectorizer", ACCENT_BLUE),
    ("텍스트 분류 (가중 빈도)", "MultinomialNB + TfidfVectorizer", ACCENT_CYAN),
    ("텍스트 분류 (존재 유무)", "BernoulliNB", ACCENT_GREEN),
    ("불균형 텍스트 분류", "ComplementNB", ACCENT_ORANGE),
    ("수치형 특성 분류", "GaussianNB", ACCENT_PURPLE),
    ("대규모 데이터 빠른 분류", "나이브 베이즈 전 종류", ACCENT_RED),
    ("특성 간 상관관계 강함", "로지스틱 회귀 또는 SVM 권장", ACCENT_ORANGE),
]

for idx, (situation, recommendation, clr) in enumerate(guide_items):
    y = Inches(2.2) + Inches(idx * 0.72)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.6), CARD_BG, clr, radius=True)
    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(5.0), Inches(0.4),
             situation, font_size=15, color=clr, bold=True)
    add_text(slide, Inches(6.5), y + Inches(0.1), Inches(5.5), Inches(0.4),
             f"-> {recommendation}", font_size=15, color=WHITE)

add_shape(slide, Inches(0.6), Inches(7.0 - 0.7), Inches(12.1), Inches(0.6), CARD_BG, radius=True)
add_text(slide, Inches(0.8), Inches(7.0 - 0.65), Inches(11.7), Inches(0.4),
         "Tip: Alpha 스무딩 파라미터는 0.1~1.0 범위에서 최적인 경우가 많음. 교차 검증으로 튜닝 권장",
         font_size=13, color=ACCENT_CYAN)

# ============================================================
# Slide 24: 복습 질문 (1/2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "10", "복습 질문 (1/2)", "Review Questions")

questions_1 = [
    "Q1. 베이즈 정리의 각 항(사전확률, 우도, 사후확률, 증거)을\n     의료진단 사례를 들어 설명하시오.",
    "Q2. 나이브 베이즈에서 '나이브'의 의미와 이 가정이 위반되어도\n     잘 작동하는 이유를 Zhang(2004) 결과를 인용하여 설명하시오.",
    "Q3. 라플라스 스무딩이 필요한 이유와 수식을 설명하고,\n     alpha=1과 alpha=0.1의 차이를 논하시오.",
    "Q4. GaussianNB, MultinomialNB, BernoulliNB, ComplementNB의\n     차이점과 각각의 적합한 데이터 유형을 비교하시오.",
    "Q5. CountVectorizer와 TfidfVectorizer의 차이점을 수식과 함께\n     설명하고, TF-IDF가 더 우수한 상황을 논하시오.",
]

colors_q = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_PURPLE]
for idx, q in enumerate(questions_1):
    y = Inches(2.1) + Inches(idx * 1.0)
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(0.85), CARD_BG, colors_q[idx], radius=True)
    add_text(slide, Inches(0.85), y + Inches(0.1), Inches(11.6), Inches(0.65),
             q, font_size=14, color=LIGHT_GRAY)

# ============================================================
# Slide 25: 복습 질문 (2/2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "10", "복습 질문 (2/2)", "Review Questions")

questions_2 = [
    "Q6. Rennie et al.(2003)이 식별한 나이브 베이즈의 5가지\n     잘못된 가정을 나열하고, 각각의 해결 방안을 설명하시오.",
    "Q7. 스팸 필터링에서 FP와 FN 중 더 심각한 문제는 무엇인지\n     논하고, 혼동행렬의 TP/TN/FP/FN을 스팸 맥락에서 설명하시오.",
    "Q8. 나이브 베이즈(생성 모델)와 로지스틱 회귀(판별 모델)의\n     차이, 소량 vs 대량 데이터에서의 성능 특성을 비교하시오.",
    "Q9. 로그 확률 트릭이 필요한 이유를 수치 예제로 설명하고,\n     argmax가 로그 변환 후에도 보존되는 이유를 서술하시오.",
    "Q10. McCallum(1998)에서 어휘 커질수록 MultinomialNB가 우수한\n      이유와 Metsis(2006)의 MNBB 최우수 발견 이유를 논하시오.",
]

colors_q2 = [ACCENT_RED, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE]
for idx, q in enumerate(questions_2):
    y = Inches(2.1) + Inches(idx * 1.0)
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(0.85), CARD_BG, colors_q2[idx], radius=True)
    add_text(slide, Inches(0.85), y + Inches(0.1), Inches(11.6), Inches(0.65),
             q, font_size=14, color=LIGHT_GRAY)

# ============================================================
# Slide 26: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG)

add_accent_line(slide, Inches(4.0), Inches(2.2), Inches(5.3), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.0),
         "Thank You", font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
         "7장: 나이브 베이즈 (Naive Bayes)", font_size=22, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(4.0), Inches(4.2), Inches(5.3), ACCENT_BLUE)

add_text(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5),
         "다음 장: 8장 - 의사결정 나무 (Decision Tree)",
         font_size=20, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
         "핵심: 베이즈 정리 + 조건부 독립 가정 = 단순하지만 강력한 분류기",
         font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
         "텍스트 분류의 기본 도구: CountVectorizer -> TF-IDF -> MultinomialNB/ComplementNB",
         font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 저장
# ============================================================
output_path = r"D:\26년1학기\기계학습\7장\7장_나이브베이즈_강의PPT.pptx"
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
