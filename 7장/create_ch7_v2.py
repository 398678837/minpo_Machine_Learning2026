# -*- coding: utf-8 -*-
"""
Chapter 7: 나이브 베이즈 - 스팸 여부 판단하기
Dark Academia Style PPTX Generator v2
- 겹침 방지를 위한 엄격한 레이아웃 시스템
- 더 상세한 내용 포함
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ══════════════════════════════════════════════════════
# 레이아웃 시스템 - 슬라이드 크기 기반 안전 영역 정의
# ══════════════════════════════════════════════════════
SLIDE_W = 13.333  # inches
SLIDE_H = 7.5     # inches

# 안전 영역 (border frame 안쪽)
SAFE_L = 0.7      # 왼쪽 안전 마진
SAFE_R = 0.7      # 오른쪽 안전 마진
SAFE_T = 0.5      # 상단 안전 마진
SAFE_B = 0.4      # 하단 안전 마진
CONTENT_W = SLIDE_W - SAFE_L - SAFE_R   # 11.933
CONTENT_H = SLIDE_H - SAFE_T - SAFE_B   # 6.6

# 제목 영역
TITLE_TOP = 0.55
TITLE_H = 0.55
DECO_LINE_TOP = TITLE_TOP + TITLE_H + 0.05  # 1.15
CONTENT_START = DECO_LINE_TOP + 0.2  # 1.35 - 컨텐츠 시작점

# 2열 레이아웃
COL_GAP = 0.3
HALF_W = (CONTENT_W - COL_GAP) / 2  # 5.8165
COL1_L = SAFE_L
COL2_L = SAFE_L + HALF_W + COL_GAP

# ══════════════════════════════════════════════════════
# Dark Academia Color Palette
# ══════════════════════════════════════════════════════
BG = RGBColor(0x1A, 0x12, 0x08)
BG_DEEP = RGBColor(0x0E, 0x0A, 0x05)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
PARCHMENT = RGBColor(0xD4, 0xBF, 0x9A)
BORDER = RGBColor(0x3D, 0x2E, 0x10)
ACCENT = RGBColor(0x8A, 0x73, 0x40)
MUTED = RGBColor(0xA0, 0x90, 0x70)
CODE_BG_C = RGBColor(0x12, 0x0E, 0x06)
CODE_FG = RGBColor(0xE8, 0xD5, 0xA0)
HI = RGBColor(0xE8, 0xC4, 0x5C)      # Highlight
RED = RGBColor(0xC0, 0x5C, 0x3C)
GREEN = RGBColor(0x6B, 0x9E, 0x5C)
BLUE = RGBColor(0x5C, 0x7E, 0xA0)
CARD_BG = RGBColor(0x20, 0x18, 0x0C)

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
W = prs.slide_width
H = prs.slide_height


# ══════════════════════════════════════════════════════
# Helper Functions
# ══════════════════════════════════════════════════════

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def draw_border(slide):
    """Double inset border"""
    for m, w in [(0.25, Pt(1.5)), (0.4, Pt(0.75))]:
        r = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(m), Inches(m),
            Inches(SLIDE_W - 2*m), Inches(SLIDE_H - 2*m)
        )
        r.fill.background()
        r.line.color.rgb = BORDER if w == Pt(1.5) else ACCENT
        r.line.width = w


def hline(slide, top, left=None, width=None):
    """Thin gold decorative line"""
    left = Inches(left) if left else Inches(SAFE_L + 0.5)
    width = Inches(width) if width else Inches(CONTENT_W - 1.0)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(top), width, Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()


def tb(slide, l, t, w, h, text, sz=15, color=PARCHMENT, bold=False, italic=False,
       align=PP_ALIGN.LEFT, font="Georgia", spacing=1.4):
    """Single textbox - all dimensions in inches (converted internally)"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    p.space_after = Pt(2)
    if spacing:
        p.line_spacing = Pt(int(sz * spacing))
    return box


def mtb(slide, l, t, w, h, items):
    """Multi-paragraph textbox. items: list of (text, sz, color, bold, italic, font)
    All dimensions in inches."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        txt = item[0]
        sz = item[1] if len(item) > 1 else 15
        clr = item[2] if len(item) > 2 else PARCHMENT
        bld = item[3] if len(item) > 3 else False
        ita = item[4] if len(item) > 4 else False
        fnt = item[5] if len(item) > 5 else "Georgia"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.name = fnt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.italic = ita
        p.space_after = Pt(2)
        p.line_spacing = Pt(int(sz * 1.35))
    return box


def code(slide, l, t, w, h, text):
    """Code block with dark background. Dims in inches."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = CODE_BG_C
    rect.line.color.rgb = ACCENT
    rect.line.width = Pt(0.5)
    pad = 0.15
    box = slide.shapes.add_textbox(
        Inches(l + pad), Inches(t + 0.1), Inches(w - 2*pad), Inches(h - 0.2)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(12)
        p.font.color.rgb = CODE_FG
        p.space_after = Pt(0)
        p.line_spacing = Pt(17)
    return rect


def card(slide, l, t, w, h, text, border=BORDER, tc=PARCHMENT, sz=14):
    """Accent card box. Dims in inches."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.color.rgb = border
    rect.line.width = Pt(1)
    pad = 0.15
    box = slide.shapes.add_textbox(
        Inches(l + pad), Inches(t + 0.1), Inches(w - 2*pad), Inches(h - 0.2)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Georgia"
        p.font.size = Pt(sz)
        p.font.color.rgb = tc
        p.space_after = Pt(2)
        p.line_spacing = Pt(int(sz * 1.35))
    return rect


def title_slide_setup(slide, title_text, bg_color=BG):
    """Standard slide: bg + border + title + deco line. Returns content_start Y."""
    set_bg(slide, bg_color)
    draw_border(slide)
    tb(slide, SAFE_L, TITLE_TOP, CONTENT_W, TITLE_H, title_text,
       sz=32, color=GOLD, bold=True, italic=True, align=PP_ALIGN.CENTER)
    hline(slide, DECO_LINE_TOP)
    return CONTENT_START


def slide_num(slide, n, total):
    tb(slide, SLIDE_W - 1.5, SLIDE_H - 0.45, 1.2, 0.35,
       f"{n} / {total}", sz=9, color=MUTED, font="Consolas", align=PP_ALIGN.RIGHT)


def bullet_list(slide, l, t, w, items, sz=14, color=PARCHMENT, gap=0.38):
    """Simple bullet list. Returns Y after last item."""
    y = t
    for item in items:
        # dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(l), Inches(y + 0.07), Inches(0.1), Inches(0.1)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.fill.background()
        tb(slide, l + 0.2, y, w - 0.2, gap, item, sz=sz, color=color)
        y += gap
    return y


def numbered_list(slide, l, t, w, items, sz=14, color=PARCHMENT, gap=0.4):
    """Numbered list with gold circles. Returns Y after last."""
    y = t
    for i, item in enumerate(items):
        # circle
        c = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(l), Inches(y + 0.02), Inches(0.28), Inches(0.28)
        )
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
        c.line.fill.background()
        tb(slide, l, y, 0.28, 0.28, str(i+1),
           sz=10, color=BG, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
        tb(slide, l + 0.38, y, w - 0.38, gap, item, sz=sz, color=color)
        y += gap
    return y


# ══════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════
slides_data = []  # will count at end


# ━━━━━━━━━━━ S1: TITLE ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DEEP)
draw_border(s)
tb(s, SAFE_L, 0.9, CONTENT_W, 0.4, "C H A P T E R  7",
   sz=14, color=MUTED, font="Consolas", align=PP_ALIGN.CENTER)
tb(s, SAFE_L, 1.6, CONTENT_W, 0.9, "나이브 베이즈",
   sz=52, color=GOLD, bold=True, italic=True, align=PP_ALIGN.CENTER)
tb(s, SAFE_L, 2.7, CONTENT_W, 0.6, "Naive Bayes Classification",
   sz=26, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)
hline(s, 3.5, SLIDE_W/2 - 2.5, 5.0)
tb(s, SAFE_L, 3.8, CONTENT_W, 0.6, "스팸 여부 판단하기",
   sz=22, color=PARCHMENT, align=PP_ALIGN.CENTER)
mtb(s, SAFE_L, 5.0, CONTENT_W, 1.0, [
    ("나이브 베이즈 모델로 문자 데이터셋을 분석해", 14, MUTED),
    ("스팸 문자를 필터링하고, 나이브 베이즈를 더 깊게 이해합니다.", 14, MUTED),
])
tb(s, SAFE_L, 6.2, CONTENT_W, 0.4, "머신러닝  |  지도학습  |  분류",
   sz=12, color=ACCENT, font="Consolas", align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━ S2: 학습 목표 & 순서 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "학습 목표 & 학습 순서")

card(s, SAFE_L, y, CONTENT_W, 0.7,
     "학습 목표: 나이브 베이즈 모델로 문자 데이터셋을 분석해 스팸 문자를 필터링하고, 나이브 베이즈를 더 깊게 이해합니다.",
     tc=HI, sz=15)
y += 0.85

steps = [
    ("7.1", "문제 정의", "스팸 문자 여부 판별 미션"),
    ("7.2", "라이브러리 및 데이터 불러오기", "pandas, numpy, seaborn, sklearn, nltk"),
    ("7.3", "특수 기호 제거하기", "string.punctuation 활용"),
    ("7.4", "불용어 제거하기", "NLTK stopwords 사용"),
    ("7.5", "목표 컬럼 형태 변경하기", "spam/ham → 1/0 변환"),
    ("7.6", "카운트 기반 벡터화하기", "CountVectorizer"),
    ("7.7", "모델링 및 예측/평가하기", "MultinomialNB, 혼동행렬"),
    ("7.8", "이해하기: 나이브 베이즈", "베이즈 정리 수식과 예시"),
]
for num, title, desc in steps:
    # number badge
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SAFE_L + 0.1), Inches(y + 0.03), Inches(0.3), Inches(0.3))
    c.fill.solid()
    c.fill.fore_color.rgb = ACCENT
    c.line.fill.background()
    tb(s, SAFE_L + 0.1, y + 0.01, 0.3, 0.3, num, sz=9, color=BG, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    tb(s, SAFE_L + 0.55, y, 3.5, 0.35, title, sz=14, color=GOLD, bold=True)
    tb(s, SAFE_L + 4.2, y, 7.0, 0.35, desc, sz=12, color=MUTED, italic=True)
    y += 0.42


# ━━━━━━━━━━━ S3: 나이브 베이즈 소개 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "나이브 베이즈(Naive Bayes) 소개")

mtb(s, SAFE_L, y, CONTENT_W, 1.6, [
    ("나이브 베이즈란?", 20, HI, True),
    ("", 4, PARCHMENT),
    ("베이즈 정리를 적용한 조건부 확률 기반의 분류 모델입니다.", 15, PARCHMENT),
    ("", 4, PARCHMENT),
    ("조건부 확률: A가 일어났을 때 B가 일어날 확률", 15, PARCHMENT, False, True),
    ("예) '무료'라는 단어가 들어있을 때 해당 메일이 스팸일 확률", 13, MUTED, False, True),
    ("", 4, PARCHMENT),
    ("스팸 필터링을 위한 대표적 모델. 딥러닝보다 간단한 방법으로", 15, PARCHMENT),
    ("자연어 처리를 원할 때 여전히 좋은 선택이 됩니다.", 15, PARCHMENT),
])
y += 1.75

# 장단점 - 2열
card(s, COL1_L, y, HALF_W, 2.0,
     "장점\n\n"
     "• 비교적 간단한 알고리즘, 속도 빠름\n"
     "• 작은 훈련셋으로도 잘 예측\n"
     "• 독립변수들이 독립적이면\n"
     "  다른 알고리즘보다 우수",
     border=GREEN, tc=PARCHMENT, sz=13)
card(s, COL2_L, y, HALF_W, 2.0,
     "단점\n\n"
     "• 독립변수가 각각 독립적임을 전제\n"
     "  (실제 데이터에서 그런 경우 드묾)\n"
     "• 숫자형 변수가 많을 때는 비적합\n"
     "• 범용성이 높지 않음",
     border=RED, tc=PARCHMENT, sz=13)
y += 2.1

card(s, SAFE_L, y, CONTENT_W, 0.55,
     "핵심: 조건부 확률(P(A|B))로 데이터를 분류하는 알고리즘",
     tc=HI, sz=15)


# ━━━━━━━━━━━ S4: 유용한 곳 & TOP10 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "나이브 베이즈가 유용한 곳")

card(s, SAFE_L, y, CONTENT_W, 1.0,
     "TOP 10 선정 이유\n"
     "범용성이 높지는 않지만, 독립변수들이 모두 독립적이라면 충분히 경쟁력 있는 알고리즘.\n"
     "특히 딥러닝을 제외하고 자연어 처리에 가장 적합한 알고리즘입니다.",
     tc=HI, sz=14)
y += 1.2

tb(s, SAFE_L, y, CONTENT_W, 0.35, "유용한 상황", sz=18, color=GOLD, bold=True)
y += 0.4

items = [
    "각 독립변수들이 모두 독립적이고 그 중요도가 비슷할 때 유용",
    "자연어 처리(NLP)에서 간단하지만 좋은 성능",
    "범주 형태의 변수가 많을 때 적합",
    "숫자형 변수가 많을 때는 적합하지 않음",
]
y = bullet_list(s, SAFE_L + 0.3, y, CONTENT_W - 0.3, items, sz=15)
y += 0.2

tb(s, SAFE_L, y, CONTENT_W, 0.35, "다른 알고리즘과의 비교", sz=18, color=GOLD, bold=True)
y += 0.4

card(s, COL1_L, y, HALF_W, 1.4,
     "나이브 베이즈\n"
     "• 조건부 확률 기반\n"
     "• 독립 가정 필요\n"
     "• 텍스트 분류에 강력\n"
     "• 훈련 속도 매우 빠름",
     border=GOLD, tc=PARCHMENT, sz=13)
card(s, COL2_L, y, HALF_W, 1.4,
     "로지스틱 회귀\n"
     "• 선형 관계 전제\n"
     "• 독립 가정 불필요\n"
     "• 범용적 사용 가능\n"
     "• 훈련 속도 빠름",
     border=BLUE, tc=PARCHMENT, sz=13)


# ━━━━━━━━━━━ S5: 문제 정의 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.1 문제 정의 : 한눈에 보는 예측 목표")

card(s, SAFE_L, y, CONTENT_W, 0.6,
     "미션: 문자 데이터셋을 이용해 스팸 여부를 판단하라!",
     tc=HI, sz=20)
y += 0.75

# 미션 상세 - 왼쪽
details = [
    ("알고리즘", "나이브 베이즈 (Naive Bayes)"),
    ("데이터셋", "spam.csv"),
    ("종속변수", "target (스팸 여부: spam / ham)"),
    ("독립변수", "text (문자 내용 - 자연어)"),
    ("문제 유형", "분류 (Classification)"),
    ("평가지표", "정확도, 혼동 행렬 (Confusion Matrix)"),
    ("모델", "MultinomialNB"),
]
for label, value in details:
    tb(s, SAFE_L + 0.2, y, 2.2, 0.32, label, sz=13, color=ACCENT, bold=True)
    tb(s, SAFE_L + 2.6, y, 8.5, 0.32, value, sz=13, color=PARCHMENT)
    y += 0.35
y += 0.1

card(s, SAFE_L, y, CONTENT_W, 0.9,
     "데이터 소개\n"
     "스팸 문자 데이터로, 독립변수는 text 하나밖에 없습니다.\n"
     "그러나 이 하나의 변수에 긴 문장 형태의 데이터가 들어 있어 많은 전처리 작업이 필요합니다.\n"
     "각 문장에 들어간 단어들을 활용하여 문자가 스팸인지 아닌지를 예측합니다.",
     tc=MUTED, sz=12)
y += 1.0

tb(s, SAFE_L, y, CONTENT_W, 0.35,
   "사용 라이브러리: numpy, pandas, seaborn, matplotlib, scikit-learn, nltk",
   sz=12, color=MUTED, italic=True)


# ━━━━━━━━━━━ S6: 라이브러리 & 데이터 불러오기 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.2 라이브러리 및 데이터 불러오기")

code(s, COL1_L, y, HALF_W, 2.2,
     "import pandas as pd\n"
     "import numpy as np\n"
     "import matplotlib.pyplot as plt\n"
     "import seaborn as sns\n"
     "\n"
     "file_url = 'https://media.github\n"
     "  usercontent.com/.../spam.csv'\n"
     "\n"
     "data = pd.read_csv(file_url)")

# 라이브러리 설명 - 오른쪽
libs = [
    ("pandas", "데이터프레임 조작 & 분석"),
    ("numpy", "수치 연산 라이브러리"),
    ("matplotlib", "기본 시각화 라이브러리"),
    ("seaborn", "통계 기반 시각화"),
    ("sklearn", "머신러닝 모델, 평가, 전처리"),
    ("nltk", "자연어 처리 (불용어 등)"),
]
ly = y
for lib, desc in libs:
    tb(s, COL2_L, ly, 2.0, 0.3, lib, sz=13, color=GOLD, bold=True, font="Consolas")
    tb(s, COL2_L + 2.1, ly, HALF_W - 2.1, 0.3, desc, sz=12, color=PARCHMENT)
    ly += 0.33
y += 2.4

tb(s, SAFE_L, y, CONTENT_W, 0.35, "데이터 확인", sz=18, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 2.3,
     "data.head()\n"
     "\n"
     "#    target  text\n"
     "# 0  ham     Go until jurong point...\n"
     "# 1  ham     Ok lar... Joking wif u...\n"
     "# 2  spam    Free entry in 2 a wkly...\n"
     "# 3  ham     U dun say so early hor...\n"
     "# 4  ham     Nah I don't think he...")

mtb(s, COL2_L, y, HALF_W, 2.3, [
    ("데이터 구조", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("• target: 종속변수 (spam / ham)", 14, PARCHMENT),
    ("• text: 독립변수 (문자 내용)", 14, PARCHMENT),
    ("", 4, PARCHMENT),
    ("data['target'].unique()", 13, CODE_FG, False, False, "Consolas"),
    ("# array(['ham', 'spam'])", 12, MUTED, False, False, "Consolas"),
    ("", 4, PARCHMENT),
    ("→ spam: 스팸 문자", 13, PARCHMENT),
    ("→ ham: 스팸이 아닌 정상 문자", 13, PARCHMENT),
    ("→ text 컬럼: 자연어 형태", 13, PARCHMENT),
])


# ━━━━━━━━━━━ S7: 전처리 파이프라인 개요 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "전처리 파이프라인 전체 흐름")

steps_pipe = [
    ("STEP 1", "특수 기호 제거", "쉼표, 마침표 등 노이즈 제거\nstring.punctuation 활용", GOLD),
    ("STEP 2", "불용어 제거", "분석에 도움 안 되는 단어 제거\nNLTK stopwords 활용", BLUE),
    ("STEP 3", "목표 컬럼 변환", "spam → 1, ham → 0 변환\nmap() 함수 활용", GREEN),
    ("STEP 4", "카운트 벡터화", "단어를 숫자 벡터로 변환\nCountVectorizer 활용", HI),
]

for i, (step_name, title, desc, clr) in enumerate(steps_pipe):
    card(s, SAFE_L, y, 2.0, 1.0, step_name, border=clr, tc=clr, sz=16)
    card(s, SAFE_L + 2.3, y, 3.0, 1.0, title, border=clr, tc=GOLD, sz=16)
    tb(s, SAFE_L + 5.6, y + 0.15, 6.0, 0.7, desc, sz=13, color=PARCHMENT)

    if i < len(steps_pipe) - 1:
        arrow = s.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(SAFE_L + 0.8), Inches(y + 1.05),
            Inches(0.3), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

    y += 1.3


# ━━━━━━━━━━━ S8: 특수 기호 제거 - 개념 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.3 특수 기호 제거하기 — 개념")

mtb(s, SAFE_L, y, CONTENT_W, 1.0, [
    ("왜 특수 기호를 제거하는가?", 18, HI, True),
    ("", 4, PARCHMENT),
    ("자연어를 다룰 때 데이터의 기준은 단어입니다.", 15, PARCHMENT),
    ("쉼표, 마침표 등의 특수 기호는 단어 처리 시 노이즈가 되므로 반드시 제거해야 합니다.", 15, PARCHMENT),
])
y += 1.1

# 특수 기호 목록
tb(s, SAFE_L, y, CONTENT_W, 0.35, "특수 기호 목록 확인", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, SAFE_L, y, CONTENT_W, 0.8,
     "import string\n"
     "string.punctuation\n"
     "# 출력: !\"#$%&'()*+,-./:;<=>?@[\\]^_`{|}~")
y += 0.95

# 제거 과정 설명
tb(s, SAFE_L, y, CONTENT_W, 0.35, "특수 기호 제거 과정 (4단계)", sz=16, color=GOLD, bold=True)
y += 0.4

process_items = [
    "문자열에서 문자를 하나씩 꺼냄 (for i in sample_string)",
    "특수 기호인지 판단 (if i not in string.punctuation)",
    "특수 기호가 아닌 문자들만 리스트에 저장 (new_string.append(i))",
    "join()으로 리스트를 다시 문자열로 합침 (''.join(new_string))",
]
y = numbered_list(s, SAFE_L + 0.2, y, CONTENT_W - 0.3, process_items, sz=14)
y += 0.15

card(s, SAFE_L, y, CONTENT_W, 0.55,
     "핵심: 문자열 → 문자 단위 순회 → 특수기호 필터링 → 리스트 → join()으로 재합성",
     tc=HI, sz=14)


# ━━━━━━━━━━━ S9: in 연산자 & join() 설명 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "핵심 함수 ① : in 연산자 & join()")

# in 연산자
tb(s, COL1_L, y, HALF_W, 0.35, "in 연산자 — 문자열 포함 여부 판단", sz=16, color=HI, bold=True)
y1 = y + 0.4

code(s, COL1_L, y1, HALF_W, 1.6,
     "# 문자열에 특정 문자가 있는지 확인\n"
     "'a' in 'apple'    # True\n"
     "'b' in 'apple'    # False\n"
     "\n"
     "# if절과 함께 사용\n"
     "if i not in string.punctuation:\n"
     "    # 특수 기호가 아닌 경우만 처리")

# join 함수
tb(s, COL2_L, y, HALF_W, 0.35, "join() — 리스트를 문자열로 합침", sz=16, color=HI, bold=True)

code(s, COL2_L, y1, HALF_W, 1.6,
     "# 리스트의 문자들을 합쳐서 문자열로\n"
     "sample = ['a','p','p','l','e']\n"
     "\n"
     "'_'.join(sample)  # 'a_p_p_l_e'\n"
     "''.join(sample)   # 'apple'\n"
     "' '.join(sample)  # 'a p p l e'\n"
     "# 앞의 따옴표 사이 문자가 구분자")
y2 = y1 + 1.75

tb(s, SAFE_L, y2, CONTENT_W, 0.35, "join() 사용 시 주의사항", sz=16, color=GOLD, bold=True)
y2 += 0.4

mtb(s, SAFE_L, y2, CONTENT_W, 1.4, [
    ("• 특수 기호 제거 시: ''.join(new_string) — 빈칸 없이 합침 (문자 하나씩 모았으므로)", 14, PARCHMENT),
    ("• 불용어 제거 시: ' '.join(new_string) — 공백으로 합침 (단어 단위로 모았으므로)", 14, PARCHMENT),
    ("", 4, PARCHMENT),
    ("join() 앞 따옴표 사이의 빈칸 유무에 따라 결과가 완전히 달라집니다!", 14, HI, True),
    ("", 4, PARCHMENT),
    ("join() 코드는 for문이 모두 끝난 뒤에 실행되어야 합니다.", 14, PARCHMENT),
    ("→ 들여쓰기를 하지 않도록 주의! (for문 바깥에 위치해야 함)", 13, RED),
])


# ━━━━━━━━━━━ S10: 특수 기호 제거 - 코드 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.3 특수 기호 제거하기 — 코드 구현")

tb(s, SAFE_L, y, CONTENT_W, 0.35, "Step 1: 샘플 문자열로 테스트", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 2.3,
     "# 문자열 하나 선택\n"
     "sample_string = data['text'].loc[0]\n"
     "\n"
     "# 특수 기호 제외 & 리스트 저장\n"
     "new_string = []\n"
     "for i in sample_string:\n"
     "    if i not in string.punctuation:\n"
     "        new_string.append(i)\n"
     "new_string = ''.join(new_string)")

code(s, COL2_L, y, HALF_W, 2.3,
     "# 함수로 만들기\n"
     "def remove_punc(x):\n"
     "    new_string = []    # 빈 리스트\n"
     "    for i in x:        # 문자열 순회\n"
     "        if i not in string.punctuation:\n"
     "            new_string.append(i)\n"
     "    new_string = ''.join(new_string)\n"
     "    return new_string  # 반환")
y += 2.45

tb(s, SAFE_L, y, CONTENT_W, 0.35, "Step 2: 데이터 전체에 적용", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 1.3,
     "# 잘못된 방법 (한 행에 모두 합쳐짐!)\n"
     "remove_punc(data['text'])\n"
     "# → 모든 행의 문자가 하나로 합쳐짐\n"
     "# → 이대로 업데이트하면 큰일!")

code(s, COL2_L, y, HALF_W, 1.3,
     "# 올바른 방법: apply() 사용\n"
     "data['text'] = \\\n"
     "    data['text'].apply(remove_punc)\n"
     "# → 각 행마다 별도로 함수 적용!")
y += 1.45

card(s, SAFE_L, y, CONTENT_W, 0.55,
     "remove_punc()은 한 줄의 문자열에만 작동 → 반드시 apply()로 각 행에 개별 적용해야 합니다.",
     tc=RED, sz=14)


# ━━━━━━━━━━━ S11: apply() 함수 상세 설명 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "핵심 함수 ② : apply() 함수")

mtb(s, SAFE_L, y, CONTENT_W, 1.0, [
    ("apply() — 데이터의 각 행에 함수를 별도로 적용", 18, HI, True),
    ("", 4, PARCHMENT),
    ("함수가 데이터의 한 행마다 별도로 적용되어야 할 때, apply()를 사용하면", 15, PARCHMENT),
    ("각 행마다 함수를 적용할 수 있습니다.", 15, PARCHMENT),
])
y += 1.1

code(s, COL1_L, y, HALF_W, 2.0,
     "# apply() 기본 사용법\n"
     "data = pd.Series([[1,2], [3,4,5]])\n"
     "\n"
     "def check_len(x):\n"
     "    return len(x)\n"
     "\n"
     "data.apply(check_len)\n"
     "# 0    2\n"
     "# 1    3")

mtb(s, COL2_L, y, HALF_W, 2.0, [
    ("apply() 동작 원리", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("1. data 시리즈의 각 행을 하나씩 꺼냄", 14, PARCHMENT),
    ("2. 꺼낸 값을 함수의 인수로 전달", 14, PARCHMENT),
    ("3. 함수의 반환값을 새 시리즈로 모음", 14, PARCHMENT),
    ("", 4, PARCHMENT),
    ("→ [1,2] → check_len → 2", 13, CODE_FG, False, False, "Consolas"),
    ("→ [3,4,5] → check_len → 3", 13, CODE_FG, False, False, "Consolas"),
])
y += 2.15

tb(s, SAFE_L, y, CONTENT_W, 0.35, "특수 기호 제거에 적용", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, SAFE_L, y, CONTENT_W, 1.5,
     "# apply()로 각 행마다 remove_punc() 적용\n"
     "data['text'].apply(remove_punc)\n"
     "\n"
     "# 결과:\n"
     "# 0    Go until jurong point crazy Available only in ...\n"
     "# 1    Ok lar Joking wif u oni\n"
     "# 2    Free entry in 2 a wkly comp to win FA Cup fina...\n"
     "\n"
     "data['text'] = data['text'].apply(remove_punc)  # 결과 업데이트")


# ━━━━━━━━━━━ S12: 불용어 제거 - 개념 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.4 불용어 제거하기 — 개념")

mtb(s, SAFE_L, y, CONTENT_W, 1.3, [
    ("불용어(Stopwords)란?", 20, HI, True),
    ("", 4, PARCHMENT),
    ("자연어 분석에 큰 도움이 안 되는 단어를 의미합니다. (I, you, he, she, the, a 등)", 15, PARCHMENT),
    ("이러한 단어를 제거하면 데이터를 가볍게 만들 수 있습니다.", 15, PARCHMENT),
    ("", 4, PARCHMENT),
    ("자연어 처리에서는 각 단어가 하나의 독립변수처럼 작용하기 때문에", 15, PARCHMENT),
    ("불용어를 제거해 분석의 부담을 줄입니다. (컬럼은 2개지만 분석 시 방대하게 펼쳐짐)", 15, PARCHMENT),
])
y += 1.5

tb(s, SAFE_L, y, CONTENT_W, 0.35, "불용어 제거의 이유", sz=16, color=GOLD, bold=True)
y += 0.4

items_sw = [
    "예측에 변별력이 없는 단어 제거 → 모델 성능 유지하면서 데이터 경량화",
    "각 단어가 독립변수가 되므로 불필요한 차원을 줄이는 효과",
    "분석 목적에 따라 불용어가 달라질 수 있음 (스팸 vs 감정 분석 등)",
]
y = bullet_list(s, SAFE_L + 0.2, y, CONTENT_W - 0.3, items_sw, sz=14)
y += 0.2

tb(s, SAFE_L, y, CONTENT_W, 0.35, "NLTK stopwords 라이브러리", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 1.5,
     "import nltk\n"
     "nltk.download('stopwords')\n"
     "\n"
     "from nltk.corpus import stopwords\n"
     "stopwords.words('english')\n"
     "# ['i','me','my','myself','we',...]")

mtb(s, COL2_L, y, HALF_W, 1.5, [
    ("지원 언어: 24개", 14, GOLD, True),
    ("arabic, english, french, german,", 12, MUTED),
    ("spanish, turkish 등", 12, MUTED),
    ("", 4, PARCHMENT),
    ("※ 한국어는 미지원", 13, RED, True),
    ("→ 별도 리스트 필요", 12, MUTED),
    ("→ www.ranks.nl/stopwords/korean", 12, MUTED),
])


# ━━━━━━━━━━━ S13: 불용어 제거 - split & lower ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.4 불용어 제거 — split()과 대소문자 변환")

# split()
tb(s, SAFE_L, y, CONTENT_W, 0.35, "split() — 문장을 단어 단위로 분리", sz=16, color=HI, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 1.5,
     "# 공백 기준 분리 (기본)\n"
     "sample_string.split()\n"
     "# ['Go','until','jurong','point',\n"
     "#  'crazy','Available','only',...]\n"
     "\n"
     "# 특정 문자 기준 분리\n"
     "'This is not - SPAM'.split('-')\n"
     "# ['This is not ', ' SPAM']")

mtb(s, COL2_L, y, HALF_W, 1.5, [
    ("split() vs for문 차이", 14, GOLD, True),
    ("", 4, PARCHMENT),
    ("for i in '문장':", 13, CODE_FG, False, False, "Consolas"),
    ("  → 문자 하나씩 순회 (문, 장)", 13, PARCHMENT),
    ("", 4, PARCHMENT),
    ("for i in '문장'.split():", 13, CODE_FG, False, False, "Consolas"),
    ("  → 단어 단위로 순회", 13, PARCHMENT),
    ("", 4, PARCHMENT),
    ("불용어 제거에는 단어 단위 필요!", 13, RED, True),
])
y += 1.65

# lower/upper/capitalize
tb(s, SAFE_L, y, CONTENT_W, 0.35, "대소문자 변환 — 왜 필요한가?", sz=16, color=HI, bold=True)
y += 0.4

mtb(s, SAFE_L, y, CONTENT_W, 0.6, [
    ("파이썬은 대소문자를 구분합니다. stopwords의 단어들은 모두 소문자로 되어 있으므로", 14, PARCHMENT),
    ("비교할 단어도 소문자여야 정확한 판단이 가능합니다. (Go ≠ go)", 14, PARCHMENT),
])
y += 0.7

# 3 functions side by side - 3열 레이아웃
third_w = (CONTENT_W - 0.4) / 3  # 약 3.84
for i, (func, desc, ex) in enumerate([
    ("lower()", "소문자로 변환", "'NaiveBayes'.lower()\n# 'naivebayes'"),
    ("upper()", "대문자로 변환", "'NaiveBayes'.upper()\n# 'NAIVEBAYES'"),
    ("capitalize()", "첫 글자만 대문자", "'NaiveBayes'.capitalize()\n# 'Naivebayes'"),
]):
    x_pos = SAFE_L + i * (third_w + 0.2)
    card(s, x_pos, y, third_w, 0.45, f"{func} — {desc}", tc=GOLD, sz=12)
    code(s, x_pos, y + 0.5, third_w, 0.7, ex)


# ━━━━━━━━━━━ S14: 불용어 제거 - 코드 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.4 불용어 제거하기 — 코드 구현")

tb(s, SAFE_L, y, CONTENT_W, 0.35, "불용어 제거 과정", sz=16, color=GOLD, bold=True)
y += 0.4

process_sw = [
    "split()으로 문장을 단어 단위 리스트로 변환",
    "각 단어를 소문자로 변환(lower()) 후 불용어인지 판단",
    "불용어가 아니면 소문자로 변환하여 리스트에 저장",
    "' '.join()으로 공백 기준 문자열로 합침 (단어 사이 빈칸!)",
]
y = numbered_list(s, SAFE_L + 0.2, y, CONTENT_W - 0.3, process_sw, sz=13, gap=0.35)
y += 0.15

code(s, COL1_L, y, HALF_W, 2.4,
     "def stop_words(x):\n"
     "    new_string = []  # 새 리스트\n"
     "    for i in x.split():  # 단어 순회\n"
     "        if i.lower() not in \\\n"
     "           stopwords.words('english'):\n"
     "            # 소문자 변환 후 불용어 아니면\n"
     "            new_string.append(i.lower())\n"
     "    new_string = ' '.join(new_string)\n"
     "    return new_string")

code(s, COL2_L, y, HALF_W, 2.4,
     "# 데이터에 적용\n"
     "data['text'] = \\\n"
     "    data['text'].apply(stop_words)\n"
     "\n"
     "# 결과 확인\n"
     "# 0  go jurong point crazy available...\n"
     "# 1  ok lar joking wif u oni\n"
     "# 2  free entry 2 wkly comp win fa...\n"
     "# 3  u dun say early hor u c already...")
y += 2.55

card(s, SAFE_L, y, CONTENT_W, 0.55,
     "주의: join() 앞 따옴표 사이에 공백(' ')을 넣어야 단어 사이에 빈칸이 생깁니다.",
     tc=RED, sz=14)


# ━━━━━━━━━━━ S15: 목표 컬럼 변경 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.5 목표 컬럼 형태 변경하기")

mtb(s, SAFE_L, y, CONTENT_W, 0.8, [
    ("target 컬럼을 문자 → 숫자로 변환", 18, HI, True),
    ("", 4, PARCHMENT),
    ("문자 형식도 에러를 유발하지 않지만, 해석에 문제가 생길 수 있으므로 숫자로 변환합니다.", 15, PARCHMENT),
    ("변환: spam → 1, ham → 0", 15, PARCHMENT),
])
y += 0.95

tb(s, SAFE_L, y, CONTENT_W, 0.35, "map() 함수", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 2.0,
     "# map() - 딕셔너리 매핑\n"
     "sample1 = pd.Series(['a','b','c'])\n"
     "sample1.map({\n"
     "    'a': 'apple',\n"
     "    'b': 'banana',\n"
     "    'c': 'cherry'\n"
     "})\n"
     "# 0  apple\n"
     "# 1  banana\n"
     "# 2  cherry")

code(s, COL2_L, y, HALF_W, 2.0,
     "# map() - 함수 적용도 가능\n"
     "sample2 = pd.Series(['a','b','c'])\n"
     "\n"
     "def add_i(x):\n"
     "    return x + 'i'\n"
     "\n"
     "sample2.map(add_i)\n"
     "# 0  ai\n"
     "# 1  bi\n"
     "# 2  ci")
y += 2.15

tb(s, SAFE_L, y, CONTENT_W, 0.35, "스팸/햄 변환 적용", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, SAFE_L, y, CONTENT_W, 1.1,
     "# spam → 1, ham → 0 변환\n"
     "data['target'] = data['target'].map({'spam': 1, 'ham': 0})\n"
     "\n"
     "data['target']\n"
     "# 0    0      (ham)\n"
     "# 1    0      (ham)\n"
     "# 2    1      (spam)\n"
     "# Name: target, Length: 5574, dtype: int64")


# ━━━━━━━━━━━ S16: 카운트 기반 벡터화 - 개념 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.6 카운트 기반 벡터화 — 개념")

mtb(s, SAFE_L, y, CONTENT_W, 1.0, [
    ("카운트 기반 벡터화(Count Vectorization)란?", 18, HI, True),
    ("", 4, PARCHMENT),
    ("문자를 개수 기반으로 벡터화하는 방식입니다.", 15, PARCHMENT),
    ("데이터 전체에 존재하는 모든 단어를 사전처럼 모은 뒤 인덱스를 부여하고,", 15, PARCHMENT),
    ("문장마다 속한 단어가 있는 인덱스를 카운트합니다.", 15, PARCHMENT),
])
y += 1.15

card(s, SAFE_L, y, CONTENT_W, 4.3,
     "예시: 카운트 기반 벡터화 과정\n"
     "\n"
     "원본 데이터:\n"
     "  data[0] = 'brown dog white cat brown bear'\n"
     "  data[1] = 'white dog black dog'\n"
     "\n"
     "① 모든 단어를 확인해 컬럼으로 삼기:\n"
     "          brown  black  white  cat  bear  dog\n"
     "\n"
     "② 각 문장마다 출현한 단어 수 확인:\n"
     "   data[0]:   2      0      1     1    1     1\n"
     "   data[1]:   0      1      1     0    0     2\n"
     "\n"
     "③ 인덱스 부여:\n"
     "   0=brown, 1=black, 2=white, 3=cat, 4=bear, 5=dog\n"
     "\n"
     "④ 변환된 데이터 (희소 행렬 형태):\n"
     "   data[0]: (0,0) 2  (0,2) 1  (0,3) 1  (0,4) 1  (0,5) 1\n"
     "   data[1]: (1,1) 1  (1,2) 1  (1,5) 2\n"
     "   → 출현하지 않은 단어의 인덱스는 포함하지 않음",
     tc=CODE_FG, sz=12)


# ━━━━━━━━━━━ S17: 카운트 기반 벡터화 - 코드 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.6 카운트 기반 벡터화 — 코드")

code(s, COL1_L, y, HALF_W, 2.2,
     "from sklearn.feature_extraction.text \\\n"
     "    import CountVectorizer\n"
     "\n"
     "x = data['text']     # 독립변수\n"
     "y = data['target']   # 종속변수\n"
     "\n"
     "cv = CountVectorizer()  # 객체 생성\n"
     "cv.fit(x)  # 학습 (어휘 구축)")

code(s, COL2_L, y, HALF_W, 2.2,
     "# 단어와 인덱스 확인\n"
     "cv.vocabulary_\n"
     "# {'go': 3791,\n"
     "#  'jurong': 4687,\n"
     "#  'point': 6433,\n"
     "#  'crazy': 2497,\n"
     "#  'available': 1414, ...}\n"
     "# 인덱스 = 일종의 ID (순서 무관)")
y += 2.35

tb(s, SAFE_L, y, CONTENT_W, 0.35, "transform()으로 데이터 변환", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 2.0,
     "x = cv.transform(x)  # 변환\n"
     "print(x)\n"
     "# (0, 1181) 1\n"
     "# (0, 1414) 1\n"
     "# (0, 3791) 1  ← 'go'\n"
     "# (0, 4687) 1  ← 'jurong'\n"
     "# (0, 6433) 1  ← 'point'\n"
     "# (1, 4655) 1\n"
     "# ...")

mtb(s, COL2_L, y, HALF_W, 2.0, [
    ("출력 형태 해석", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("(행번호, 인덱스) 출현횟수", 14, HI, True, False, "Consolas"),
    ("", 4, PARCHMENT),
    ("(0, 1181) 1 의 의미:", 13, PARCHMENT),
    ("→ 0번째 행에 인덱스 1181인", 13, PARCHMENT),
    ("  단어가 1번 등장", 13, PARCHMENT),
    ("", 4, PARCHMENT),
    ("• 출현하지 않은 단어는 미포함", 13, MUTED),
    ("• 희소 행렬(Sparse Matrix) 형태", 13, MUTED),
    ("  → 메모리 효율적 저장", 13, MUTED),
])
y += 2.15

card(s, SAFE_L, y, CONTENT_W, 0.5,
     "CountVectorizer 패턴: 객체 생성 → fit() 학습 → transform() 변환 (스케일링과 동일한 패턴!)",
     tc=HI, sz=14)


# ━━━━━━━━━━━ S18: 벡터화 검증 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.6 카운트 기반 벡터화 — 검증")

mtb(s, SAFE_L, y, CONTENT_W, 0.8, [
    ("변환 결과를 실제로 검증해봅시다", 16, HI, True),
    ("0번째 행의 데이터를 확인하고, 각 단어의 인덱스가 일치하는지 확인합니다.", 14, PARCHMENT),
])
y += 0.9

code(s, COL1_L, y, HALF_W, 2.0,
     "# 0번째 행 원본 텍스트\n"
     "data.loc[0]['text']\n"
     "# 'go jurong point crazy available\n"
     "#  bugis n great world la e\n"
     "#  buffet cine got amore wat'\n"
     "\n"
     "# 단어별 인덱스 확인\n"
     "print(cv.vocabulary_['go'])    # 3791\n"
     "print(cv.vocabulary_['jurong'])# 4687\n"
     "print(cv.vocabulary_['point']) # 6433")

mtb(s, COL2_L, y, HALF_W, 2.0, [
    ("변환 결과에서 확인", 14, GOLD, True),
    ("", 4, PARCHMENT),
    ("(0, 3791) 1  ← go 1회", 13, CODE_FG, False, False, "Consolas"),
    ("(0, 4687) 1  ← jurong 1회", 13, CODE_FG, False, False, "Consolas"),
    ("(0, 6433) 1  ← point 1회", 13, CODE_FG, False, False, "Consolas"),
    ("", 4, PARCHMENT),
    ("→ 인덱스와 출현 횟수가 정확히 일치!", 14, GREEN, True),
    ("", 4, PARCHMENT),
    ("Tip: print()를 써야 위 형태로 보임", 12, MUTED),
    ("주피터에서 여러 줄 출력 시 print() 필수", 12, MUTED),
])
y += 2.15

card(s, SAFE_L, y, CONTENT_W, 1.5,
     "vocabulary_ 속성 정리\n\n"
     "• cv.vocabulary_ : 딕셔너리 형태 → 단어를 키로, 인덱스를 값으로 저장\n"
     "• 인덱스는 큰 의미 없는 일종의 ID (알파벳 순서 등으로 부여됨)\n"
     "• 전체 데이터에서 고유한 모든 단어에 대해 인덱스가 존재\n\n"
     "이 방식을 카운트 기반 벡터화(CountVectorize)라고 합니다.\n"
     "문자를 이런 식의 형태(희소 행렬)로 변환시켜 머신러닝 모델에 입력합니다.",
     tc=PARCHMENT, sz=13)


# ━━━━━━━━━━━ S19: 모델링 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.7 모델링 및 예측하기")

tb(s, SAFE_L, y, CONTENT_W, 0.35, "훈련셋 / 시험셋 분할", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, SAFE_L, y, CONTENT_W, 1.2,
     "from sklearn.model_selection import train_test_split\n"
     "\n"
     "x_train, x_test, y_train, y_test = train_test_split(\n"
     "    x, y, test_size=0.2, random_state=100\n"
     ")  # 훈련 80%, 시험 20%")
y += 1.35

tb(s, SAFE_L, y, CONTENT_W, 0.35, "MultinomialNB 모델링", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 1.8,
     "from sklearn.naive_bayes \\\n"
     "    import MultinomialNB\n"
     "\n"
     "model = MultinomialNB()   # 객체 생성\n"
     "model.fit(x_train, y_train)  # 학습\n"
     "pred = model.predict(x_test) # 예측")

mtb(s, COL2_L, y, HALF_W, 1.8, [
    ("MultinomialNB란?", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("다항 분포에 대한 Naive Bayes 알고리즘", 14, PARCHMENT),
    ("텍스트 분류에 가장 적합한 NB 모델", 14, PARCHMENT),
    ("", 4, PARCHMENT),
    ("다른 NB 모델들:", 14, GOLD, True),
    ("• GaussianNB — 정규분포 (연속형)", 13, PARCHMENT),
    ("• BernoulliNB — 베르누이 (이진형)", 13, PARCHMENT),
    ("", 4, PARCHMENT),
    ("잘 모르면 → 세 가지 모두 비교!", 13, HI, True),
])
y += 1.95

card(s, SAFE_L, y, CONTENT_W, 0.5,
     "사용 패턴: fit(x_train, y_train) → predict(x_test) — 다른 sklearn 모델과 동일!",
     tc=HI, sz=14)


# ━━━━━━━━━━━ S20: 평가 - 정확도 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.7 평가하기 — 정확도(Accuracy)")

code(s, COL1_L, y, HALF_W, 1.5,
     "from sklearn.metrics import \\\n"
     "    accuracy_score, confusion_matrix\n"
     "\n"
     "# 정확도 계산\n"
     "accuracy_score(y_test, pred)\n"
     "# 0.9856502242152466\n"
     "# → 약 98.6% 정확도!")

mtb(s, COL2_L, y, HALF_W, 1.5, [
    ("정확도(Accuracy) 수식", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("정확도 = 정확한 예측 건수 / 전체 경우 수", 15, HI),
    ("", 4, PARCHMENT),
    ("= (965 + 134) / (965 + 12 + 4 + 134)", 14, PARCHMENT, False, False, "Consolas"),
    ("= 1099 / 1115", 14, PARCHMENT, False, False, "Consolas"),
    ("≈ 98.6%", 18, GREEN, True, False, "Consolas"),
])
y += 1.65

tb(s, SAFE_L, y, CONTENT_W, 0.35, "혼동 행렬 (Confusion Matrix)", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 1.3,
     "print(confusion_matrix(y_test, pred))\n"
     "# [[965  12]\n"
     "#  [  4 134]]\n"
     "\n"
     "# 히트맵 시각화\n"
     "sns.heatmap(confusion_matrix(\n"
     "    y_test, pred), annot=True, fmt='d')")

mtb(s, COL2_L, y, HALF_W, 1.3, [
    ("혼동 행렬 해석", 14, GOLD, True),
    ("", 4, PARCHMENT),
    ("         예측0  예측1", 13, GOLD, False, False, "Consolas"),
    ("실젯값0 [ 965    12 ]", 13, CODE_FG, False, False, "Consolas"),
    ("실젯값1 [   4   134 ]", 13, CODE_FG, False, False, "Consolas"),
    ("", 4, PARCHMENT),
    ("배경 음영 영역(965, 134) = 정확한 예측", 12, GREEN),
    ("흰색 영역(12, 4) = 잘못된 예측", 12, RED),
])
y += 1.45

card(s, SAFE_L, y, CONTENT_W, 1.2,
     "fmt 매개변수\n\n"
     "• fmt='d' : 정수 형태로 표시 (소수점 없음)\n"
     "• fmt='.2f' : 소수점 둘째 자리까지\n"
     "• 설정하지 않으면 과학적 표기법으로 나타날 수 있으므로 필요할 때 설정",
     tc=PARCHMENT, sz=13)


# ━━━━━━━━━━━ S21: 혼동 행렬 상세 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "혼동 행렬(Confusion Matrix) 상세 해석")

# 2x2 matrix visual
mx = SAFE_L + 2.5
my = y
cell_w = 2.2
cell_h = 1.0
gap_c = 0.1

# headers
tb(s, mx + cell_w/2, my, cell_w * 2 + gap_c, 0.35, "예측값",
   sz=15, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
tb(s, mx, my + 0.35, cell_w, 0.3, "0",
   sz=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
tb(s, mx + cell_w + gap_c, my + 0.35, cell_w, 0.3, "1",
   sz=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
tb(s, mx - 1.2, my + 0.7 + cell_h/2, 1.0, 0.3, "실젯값",
   sz=15, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
tb(s, mx - 0.4, my + 0.7, 0.3, cell_h, "0",
   sz=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
tb(s, mx - 0.4, my + 0.7 + cell_h + gap_c, 0.3, cell_h, "1",
   sz=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

cells_data = [
    (0, 0, "965", "TN", GREEN),
    (1, 0, "12", "FP", RED),
    (0, 1, "4", "FN", RED),
    (1, 1, "134", "TP", GREEN),
]
for ci, ri, num, label, clr in cells_data:
    cx = mx + ci * (cell_w + gap_c)
    cy = my + 0.7 + ri * (cell_h + gap_c)
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(cx), Inches(cy), Inches(cell_w), Inches(cell_h))
    r.fill.solid()
    r.fill.fore_color.rgb = CARD_BG
    r.line.color.rgb = clr
    r.line.width = Pt(2)
    tb(s, cx + 0.1, cy + 0.05, cell_w - 0.2, 0.45, num,
       sz=26, color=HI, bold=True, align=PP_ALIGN.CENTER)
    tb(s, cx + 0.1, cy + 0.55, cell_w - 0.2, 0.4, label,
       sz=12, color=MUTED, align=PP_ALIGN.CENTER, font="Consolas")

# 오른쪽 설명
rx = mx + 2 * cell_w + gap_c + 0.5
mtb(s, rx, my + 0.3, CONTENT_W - (rx - SAFE_L), 2.5, [
    ("각 셀의 의미", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("TN = 965", 14, GREEN, True),
    ("  음성을 음성으로 (정확)", 13, PARCHMENT),
    ("TP = 134", 14, GREEN, True),
    ("  양성을 양성으로 (정확)", 13, PARCHMENT),
    ("", 4, PARCHMENT),
    ("FP = 12 (1종 오류)", 14, RED, True),
    ("  음성을 양성으로 (오류)", 13, PARCHMENT),
    ("FN = 4 (2종 오류)", 14, RED, True),
    ("  양성을 음성으로 (오류)", 13, PARCHMENT),
])

y2 = my + 3.1

card(s, SAFE_L, y2, CONTENT_W, 1.5,
     "명칭 규칙\n\n"
     "• True/False: 예측이 맞았는지 (참 = 정확, 거짓 = 오류)\n"
     "• Positive/Negative: 예측값 기준 (1이면 양성, 0이면 음성)\n\n"
     "• 양성(Positive) = 1 (스팸)\n"
     "• 음성(Negative) = 0 (햄, 정상)",
     tc=PARCHMENT, sz=13)


# ━━━━━━━━━━━ S22: 1종/2종 오류 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "1종 오류 vs 2종 오류")

# Type 1
card(s, COL1_L, y, HALF_W, 2.0,
     "1종 오류 (Type 1 Error)\nFalse Positive — 거짓 양성\n\n"
     "실제 음성(0)인 것을\n"
     "양성(1)으로 예측하는 오류\n\n"
     "암기: '양성 1종'",
     border=BLUE, tc=PARCHMENT, sz=14)
card(s, COL2_L, y, HALF_W, 2.0,
     "2종 오류 (Type 2 Error)\nFalse Negative — 거짓 음성\n\n"
     "실제 양성(1)인 것을\n"
     "음성(0)으로 예측하는 오류\n\n"
     "암기: '음성 2종'",
     border=RED, tc=PARCHMENT, sz=14)
y += 2.15

tb(s, SAFE_L, y, CONTENT_W, 0.35, "실제 사례: 암 진단", sz=16, color=GOLD, bold=True)
y += 0.4

card(s, COL1_L, y, HALF_W, 1.3,
     "1종 오류: 암이 아닌데 암이라고 진단\n\n"
     "→ 추가 진단으로 곧 아님을 알게 됨\n"
     "→ 상대적으로 덜 위험",
     border=BLUE, tc=PARCHMENT, sz=13)
card(s, COL2_L, y, HALF_W, 1.3,
     "2종 오류: 암인데 아니라고 진단\n\n"
     "→ 환자가 암을 모르고 지냄\n"
     "→ 치료 시기를 놓칠 수 있음 (더 위험!)",
     border=RED, tc=PARCHMENT, sz=13)
y += 1.45

tb(s, SAFE_L, y, CONTENT_W, 0.35, "스팸 문자에서의 1종/2종 오류", sz=16, color=GOLD, bold=True)
y += 0.4

card(s, COL1_L, y, HALF_W, 0.9,
     "1종 오류: 스팸 아닌데 스팸으로 분류\n"
     "→ 중요 문자가 필터링될 수 있음",
     border=BLUE, tc=PARCHMENT, sz=13)
card(s, COL2_L, y, HALF_W, 0.9,
     "2종 오류: 스팸인데 스팸 아니라고 분류\n"
     "→ 스팸을 받지만 중요 문자 차단 적음",
     border=RED, tc=PARCHMENT, sz=13)


# ━━━━━━━━━━━ S23: 베이즈 정리 수식 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "7.8 베이즈 정리 (Bayes' Theorem)", BG_DEEP)

# Main formula
card(s, SAFE_L + 1.5, y, CONTENT_W - 3.0, 0.8,
     "P(A|B) = P(B|A) × P(A) / P(B)",
     tc=HI, sz=26)
y += 1.0

# Terms explanation
terms = [
    ("P(A|B)", "사후확률\n(Posterior)", "B가 발생했을 때 A가 발생할 확률\n→ 특정 단어가 등장했을 때 스팸일 확률"),
    ("P(A)", "사전확률\n(Prior)", "B의 발생 유무와 관련 없이 A가 발생할 확률\n→ 전체 문자 중 스팸 문자의 비율"),
    ("P(B|A)", "우도\n(Likelihood)", "A가 발생했을 때 B가 발생할 확률\n→ 스팸 메일에서 특정 단어가 포함될 확률"),
    ("P(B)", "증거\n(Evidence)", "전체에서 B가 발생할 확률\n→ 전체 문자에서 특정 단어가 포함될 확률"),
]

for symbol, name, desc in terms:
    tb(s, SAFE_L + 0.2, y, 1.5, 0.7, symbol, sz=16, color=HI, bold=True,
       align=PP_ALIGN.CENTER, font="Consolas")
    tb(s, SAFE_L + 1.8, y, 1.8, 0.7, name, sz=13, color=GOLD, bold=True)
    tb(s, SAFE_L + 3.8, y, CONTENT_W - 4.0, 0.7, desc, sz=12, color=PARCHMENT)
    y += 0.7

y += 0.1
card(s, SAFE_L, y, CONTENT_W, 0.8,
     "베이즈 정리: 두 확률 변수의 사전확률과 사후확률 사이의 관계를 나타내는 정리.\n"
     "사후확률을 구할 때 사용됩니다. 나이브 베이즈는 이를 기반으로 분류를 수행합니다.",
     tc=PARCHMENT, sz=14)


# ━━━━━━━━━━━ S24: 베이즈 정리 예시 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "베이즈 정리 — 그림으로 이해하기", BG_DEEP)

tb(s, SAFE_L, y, CONTENT_W, 0.35, "예시 조건 설정", sz=16, color=GOLD, bold=True)
y += 0.4

items_cond = [
    "전체 문자 중 스팸 비율: 30% (햄 70%)",
    "스팸에서 특정 단어 X가 포함된 경우: 50%",
    "햄에서 특정 단어 X가 포함된 경우: 10%",
]
y = bullet_list(s, SAFE_L + 0.2, y, CONTENT_W - 0.3, items_cond, sz=14)
y += 0.15

# Visual boxes
card(s, COL1_L, y, HALF_W, 1.1,
     "스팸 (30%)\n특정 단어 X 포함: 50%",
     border=RED, tc=PARCHMENT, sz=16)
card(s, COL2_L, y, HALF_W, 1.1,
     "햄 (70%)\n특정 단어 X 포함: 10%",
     border=GREEN, tc=PARCHMENT, sz=16)
y += 1.25

# Calculation
tb(s, SAFE_L, y, CONTENT_W, 0.35, "계산: 특정 단어 X가 있을 때 스팸일 확률 = ?", sz=16, color=HI, bold=True)
y += 0.4

card(s, SAFE_L, y, CONTENT_W, 2.3,
     "P(A) = 스팸 확률 = 0.3\n"
     "P(B|A) = 스팸 중 X 포함 확률 = 0.5\n"
     "P(B) = 전체 중 X 포함 확률\n"
     "     = (스팸 × 스팸에서 X) + (햄 × 햄에서 X)\n"
     "     = (0.3 × 0.5) + (0.7 × 0.1) = 0.15 + 0.07 = 0.22\n"
     "\n"
     "P(A|B) = P(B|A) × P(A) / P(B)\n"
     "       = 0.5 × 0.3 / 0.22\n"
     "       = 0.15 / 0.22\n"
     "       ≈ 0.6818 (약 68.2%)\n"
     "\n"
     "→ 특정 단어 X가 포함된 문자가 스팸일 확률: 약 68%",
     tc=CODE_FG, sz=13)


# ━━━━━━━━━━━ S25: 나이브의 의미 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "나이브 베이즈의 '나이브(Naive)' 가정")

mtb(s, SAFE_L, y, CONTENT_W, 1.5, [
    ("왜 'Naive(나이브, 순진한)'인가?", 20, HI, True),
    ("", 4, PARCHMENT),
    ("나이브 베이즈는 모든 독립변수(feature)가 서로 독립이라고 가정합니다.", 15, PARCHMENT),
    ("이 가정은 실제 데이터에서 거의 성립하지 않지만,", 15, PARCHMENT),
    ("이 '순진한' 가정 덕분에 계산이 매우 단순해집니다.", 15, PARCHMENT),
    ("", 4, PARCHMENT),
    ("독립 가정 적용 시:", 15, GOLD, True),
    ("P(x1,x2,...,xn|C) = P(x1|C) × P(x2|C) × ... × P(xn|C)", 14, HI, False, False, "Consolas"),
])
y += 1.65

tb(s, SAFE_L, y, CONTENT_W, 0.35, "나이브 베이즈의 3가지 변형", sz=16, color=GOLD, bold=True)
y += 0.4

third_w2 = (CONTENT_W - 0.4) / 3
models_nb = [
    ("MultinomialNB", "다항 분포 기반\n\n단어 빈도(count) 사용\n텍스트 분류에 최적\n\n이번 예제에서 사용"),
    ("GaussianNB", "정규(가우시안) 분포\n\n연속형 변수에 적합\n수치 데이터에 사용\n\n.toarray() 필요"),
    ("BernoulliNB", "베르누이 분포\n\n이진(0/1) 특징 사용\nbinarize 파라미터\n\nBernoulliNB(binarize=True)"),
]
for i, (name, desc) in enumerate(models_nb):
    x_pos = SAFE_L + i * (third_w2 + 0.2)
    card(s, x_pos, y, third_w2, 0.5, name, tc=GOLD, sz=14)
    card(s, x_pos, y + 0.55, third_w2, 2.0, desc, tc=PARCHMENT, sz=12)

y += 2.7
card(s, SAFE_L, y, CONTENT_W, 0.5,
     "데이터 특성을 명확히 알기 어렵다면 세 가지 모두 사용하여 가장 결과가 좋은 모델을 선택!",
     tc=HI, sz=14)


# ━━━━━━━━━━━ S26: 수학적 전개 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "심화: 나이브 베이즈의 수학적 전개", BG_DEEP)

mtb(s, SAFE_L, y, CONTENT_W, 5.5, [
    ("분류 문제에서의 베이즈 정리 적용", 16, HI, True),
    ("", 6, PARCHMENT),
    ("목표: 주어진 특징 벡터 x에 대해 클래스 C를 예측", 14, PARCHMENT),
    ("", 4, PARCHMENT),
    ("Step 1. 베이즈 정리 적용", 14, GOLD, True),
    ("P(C|x1,...,xn) = P(x1,...,xn|C) × P(C) / P(x1,...,xn)", 13, CODE_FG, False, False, "Consolas"),
    ("", 6, PARCHMENT),
    ("Step 2. 나이브(독립) 가정 적용", 14, GOLD, True),
    ("P(x1,...,xn|C) = P(x1|C) × P(x2|C) × ... × P(xn|C)", 13, CODE_FG, False, False, "Consolas"),
    ("", 6, PARCHMENT),
    ("Step 3. 정리", 14, GOLD, True),
    ("P(C|x) ∝ P(C) × ∏ P(xi|C)", 13, CODE_FG, False, False, "Consolas"),
    ("(P(x)는 모든 클래스에 대해 동일하므로 비교 시 무시 가능)", 12, MUTED),
    ("", 6, PARCHMENT),
    ("Step 4. 최종 분류 규칙", 14, GOLD, True),
    ("y_hat = argmax_C  P(C) × ∏ P(xi|C)", 13, CODE_FG, False, False, "Consolas"),
    ("→ 각 클래스에 대해 사후확률을 계산, 가장 큰 클래스 선택", 12, MUTED),
    ("", 6, PARCHMENT),
    ("MultinomialNB에서 P(xi|C) 계산:", 14, GOLD, True),
    ("P(xi|C) = (클래스 C에서 단어 xi 출현 횟수 + α)", 13, CODE_FG, False, False, "Consolas"),
    ("        / (클래스 C의 총 단어 수 + α × |V|)", 13, CODE_FG, False, False, "Consolas"),
    ("", 4, PARCHMENT),
    ("여기서 α = 라플라스 스무딩 (기본값 1.0)", 12, MUTED),
    ("|V| = 전체 어휘(vocabulary) 크기", 12, MUTED),
])


# ━━━━━━━━━━━ S27: 전체 코드 정리 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "전체 코드 흐름 정리")

code(s, COL1_L, y, HALF_W, 5.5,
     "# ── ① 라이브러리 & 데이터 ──\n"
     "import pandas as pd, numpy as np\n"
     "import matplotlib.pyplot as plt\n"
     "import seaborn as sns\n"
     "import string, nltk\n"
     "from nltk.corpus import stopwords\n"
     "nltk.download('stopwords')\n"
     "data = pd.read_csv(file_url)\n"
     "\n"
     "# ── ② 특수 기호 제거 ──\n"
     "def remove_punc(x):\n"
     "    new_string = []\n"
     "    for i in x:\n"
     "        if i not in string.punctuation:\n"
     "            new_string.append(i)\n"
     "    return ''.join(new_string)\n"
     "data['text'] = \\\n"
     "    data['text'].apply(remove_punc)\n"
     "\n"
     "# ── ③ 불용어 제거 ──\n"
     "def stop_words(x):\n"
     "    new_string = []\n"
     "    for i in x.split():\n"
     "        if i.lower() not in \\\n"
     "           stopwords.words('english'):\n"
     "            new_string.append(i.lower())\n"
     "    return ' '.join(new_string)\n"
     "data['text'] = \\\n"
     "    data['text'].apply(stop_words)")

code(s, COL2_L, y, HALF_W, 5.5,
     "# ── ④ 타겟 변환 ──\n"
     "data['target'] = data['target'].map(\n"
     "    {'spam': 1, 'ham': 0})\n"
     "\n"
     "# ── ⑤ 카운트 벡터화 ──\n"
     "from sklearn.feature_extraction.text \\\n"
     "    import CountVectorizer\n"
     "x = data['text']\n"
     "y = data['target']\n"
     "cv = CountVectorizer()\n"
     "cv.fit(x)\n"
     "x = cv.transform(x)\n"
     "\n"
     "# ── ⑥ 훈련/시험 분할 ──\n"
     "from sklearn.model_selection \\\n"
     "    import train_test_split\n"
     "x_train, x_test, y_train, y_test = \\\n"
     "    train_test_split(x, y,\n"
     "    test_size=0.2, random_state=100)\n"
     "\n"
     "# ── ⑦ 모델링 & 평가 ──\n"
     "from sklearn.naive_bayes \\\n"
     "    import MultinomialNB\n"
     "model = MultinomialNB()\n"
     "model.fit(x_train, y_train)\n"
     "pred = model.predict(x_test)\n"
     "\n"
     "from sklearn.metrics import \\\n"
     "    accuracy_score, confusion_matrix\n"
     "print(accuracy_score(y_test, pred))")


# ━━━━━━━━━━━ S28: 학습 마무리 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "학습 마무리 — 8단계 요약")

summary = [
    ("문제 정의", "스팸 문자를 구분하는 모델을 만듭니다"),
    ("라이브러리 임포트", "pandas, numpy, matplotlib, seaborn, sklearn, nltk"),
    ("특수 기호 제거", "string.punctuation + remove_punc() + apply()"),
    ("불용어 제거", "NLTK stopwords + stop_words() + apply()"),
    ("목표 컬럼 변환", "map({'spam': 1, 'ham': 0})"),
    ("카운트 벡터화", "CountVectorizer: fit() → transform()"),
    ("모델링", "MultinomialNB: fit() → predict()"),
    ("평가 결과", "약 98.6%의 높은 정확도 달성!"),
]

for i, (title, desc) in enumerate(summary):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(SAFE_L + 0.2), Inches(y + 0.02), Inches(0.3), Inches(0.3))
    c.fill.solid()
    c.fill.fore_color.rgb = ACCENT
    c.line.fill.background()
    tb(s, SAFE_L + 0.2, y, 0.3, 0.3, str(i+1),
       sz=11, color=BG, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    tb(s, SAFE_L + 0.65, y, 2.8, 0.35, title, sz=14, color=GOLD, bold=True)
    tb(s, SAFE_L + 3.6, y, CONTENT_W - 3.6, 0.35, desc, sz=13, color=PARCHMENT)
    y += 0.42

y += 0.15
card(s, SAFE_L, y, CONTENT_W, 1.0,
     "과제 안내\n"
     "• 베이즈 정리를 직접 손으로 풀어보기\n"
     "• MultinomialNB 대신 GaussianNB, BernoulliNB(binarize=True)를 사용하여 비교",
     tc=PARCHMENT, sz=14)


# ━━━━━━━━━━━ S29: 핵심 용어 & API ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "핵심 용어 & API 정리")

terms_list = [
    ("나이브 베이즈 분류기", "조건부 확률 기반 모델. 변수 개수 많을 때 유용"),
    ("1종 오류 (Type 1, FP)", "실제 음성을 양성으로 예측하는 오류"),
    ("2종 오류 (Type 2, FN)", "실제 양성을 음성으로 예측하는 오류"),
    ("사후확률 (Posterior)", "A 발생 상황에서 B 발생 확률"),
    ("사전확률 (Prior)", "A와 상관없이 B 발생 확률"),
    ("베이즈 정리", "사전/사후확률 관계를 나타내는 정리"),
]
for term, desc in terms_list:
    tb(s, SAFE_L, y, 4.0, 0.35, term, sz=14, color=HI, bold=True)
    tb(s, SAFE_L + 4.2, y, CONTENT_W - 4.2, 0.35, desc, sz=13, color=PARCHMENT)
    y += 0.38

hline(s, y + 0.05)
y += 0.2

apis = [
    ("string.punctuation", "특수 기호 목록 출력"),
    ("MultinomialNB()", "다항분포 나이브 베이즈"),
    ("nltk stopwords", "불용어 목록"),
    ("CountVectorizer()", "카운트 기반 벡터화"),
    ("confusion_matrix()", "혼동 행렬"),
    ("cv.vocabulary_", "단어 → 인덱스 딕셔너리"),
    ("accuracy_score()", "정확도 계산"),
    ("apply()", "각 행에 함수 개별 적용"),
    ("map()", "딕셔너리/함수로 값 매핑"),
    ("split()", "문자열 → 단어 리스트 분리"),
]
for api, desc in apis:
    tb(s, SAFE_L, y, 3.8, 0.3, api, sz=11, color=GOLD, bold=True, font="Consolas")
    tb(s, SAFE_L + 4.0, y, CONTENT_W - 4.0, 0.3, desc, sz=11, color=MUTED)
    y += 0.3


# ━━━━━━━━━━━ S30: 연습 문제 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "연습 문제")

questions = [
    ("Q1", "자연어 처리에서 의미 없이 빈번하게 발생하는 단어를 의미하는 것은?\n"
           "① 불용어  ② 특수기호  ③ 벡터  ④ 인덱스",
     "정답: ①"),
    ("Q2", "자연어를 머신러닝이 받아들일 수 있도록 단어와 출현 빈도로 정리하는 함수는?\n"
           "① get_dummies()  ② CountVectorizer()  ③ StandardScaler()  ④ value_counts()",
     "정답: ②"),
    ("Q3", "혼동 행렬 (0:64,3 / 1:16,17) — False Negative는 총 64건이다?\n"
           "→ FN은 실젯값 1, 예측값 0 영역이므로 16건 (64건은 TN)",
     "정답: ④ 틀림 (FN=16)"),
    ("Q4", "나이브 베이즈를 가장 잘 설명한 것은?\n"
           "① 조건부 확률 기반, 사전/사후확률 활용  ② 선형 관계 전제\n"
           "③ 독립변수 적을 때  ④ 상관관계 강할 때",
     "정답: ①"),
]

for qnum, question, answer in questions:
    tb(s, SAFE_L, y, 0.5, 0.35, qnum, sz=15, color=HI, bold=True)
    tb(s, SAFE_L + 0.6, y, CONTENT_W - 3.5, 0.8, question, sz=12, color=PARCHMENT)
    tb(s, SAFE_L + CONTENT_W - 2.8, y, 2.8, 0.35, answer, sz=12, color=GREEN, bold=True)
    y += 1.25


# ━━━━━━━━━━━ S31: 심화 - CountVectorizer 상세 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "심화: CountVectorizer 매개변수", BG_DEEP)

code(s, COL1_L, y, HALF_W, 2.5,
     "from sklearn.feature_extraction.text \\\n"
     "    import CountVectorizer\n"
     "\n"
     "# 주요 매개변수\n"
     "cv = CountVectorizer(\n"
     "    max_features=5000,   # 상위 N개만\n"
     "    ngram_range=(1, 2),  # 1~2단어 조합\n"
     "    stop_words='english',# 내장 불용어\n"
     "    max_df=0.95, # 빈도 너무 높은 단어 제외\n"
     "    min_df=2     # 빈도 너무 낮은 단어 제외\n"
     ")")

mtb(s, COL2_L, y, HALF_W, 2.5, [
    ("매개변수 설명", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("max_features: 빈도 상위 N개 단어만", 13, PARCHMENT),
    ("ngram_range: 단어 조합 범위", 13, PARCHMENT),
    ("  (1,1)=단어, (1,2)=단어+2단어조합", 12, MUTED),
    ("  예: 'free', 'free entry'", 12, MUTED),
    ("stop_words: 불용어 제거 옵션", 13, PARCHMENT),
    ("max_df: 너무 흔한 단어 제외", 13, PARCHMENT),
    ("min_df: 너무 드문 단어 제외", 13, PARCHMENT),
    ("", 4, PARCHMENT),
    ("→ 적절한 설정으로 성능 향상 가능", 13, HI, True),
])
y += 2.65

tb(s, SAFE_L, y, CONTENT_W, 0.35, "TF-IDF: 더 정교한 벡터화", sz=16, color=GOLD, bold=True)
y += 0.4

code(s, COL1_L, y, HALF_W, 1.8,
     "from sklearn.feature_extraction.text \\\n"
     "    import TfidfVectorizer\n"
     "\n"
     "tfidf = TfidfVectorizer(max_features=5000)\n"
     "x_tfidf = tfidf.fit_transform(data['text'])\n"
     "\n"
     "# TF-IDF = TF(빈도) × IDF(희귀도)\n"
     "# 흔한 단어는 낮은 가중치 부여")

mtb(s, COL2_L, y, HALF_W, 1.8, [
    ("Count vs TF-IDF", 14, GOLD, True),
    ("", 4, PARCHMENT),
    ("Count: 단순 출현 횟수", 13, PARCHMENT),
    ("  → 모든 단어 동등 취급", 12, MUTED),
    ("", 4, PARCHMENT),
    ("TF-IDF: 중요도 가중치 부여", 13, PARCHMENT),
    ("  → 흔한 단어는 낮은 가중치", 12, MUTED),
    ("  → 더 정교한 분류 가능", 12, MUTED),
    ("", 4, PARCHMENT),
    ("MultinomialNB와 함께 사용 가능", 13, HI),
])


# ━━━━━━━━━━━ S32: 심화 - Pipeline & 성능 개선 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "심화: Pipeline & 성능 개선 전략", BG_DEEP)

code(s, COL1_L, y, HALF_W, 2.8,
     "from sklearn.pipeline import Pipeline\n"
     "from sklearn.feature_extraction.text \\\n"
     "    import TfidfVectorizer\n"
     "from sklearn.naive_bayes \\\n"
     "    import MultinomialNB\n"
     "from sklearn.model_selection \\\n"
     "    import cross_val_score\n"
     "\n"
     "# Pipeline 구성\n"
     "pipe = Pipeline([\n"
     "    ('tfidf', TfidfVectorizer(\n"
     "        max_features=5000,\n"
     "        ngram_range=(1, 2))),\n"
     "    ('clf', MultinomialNB(alpha=0.1))\n"
     "])")

code(s, COL2_L, y, HALF_W, 2.8,
     "# 교차 검증\n"
     "scores = cross_val_score(\n"
     "    pipe, data['text'], data['target'],\n"
     "    cv=5, scoring='accuracy'\n"
     ")\n"
     "print(f'Mean: {scores.mean():.4f}')\n"
     "\n"
     "# 상세 평가\n"
     "from sklearn.metrics \\\n"
     "    import classification_report\n"
     "pipe.fit(x_train_text, y_train)\n"
     "print(classification_report(\n"
     "    y_test, pipe.predict(x_test_text)))")
y += 2.95

tb(s, SAFE_L, y, CONTENT_W, 0.35, "성능 개선 전략 5가지", sz=16, color=GOLD, bold=True)
y += 0.4

improve_items = [
    "TF-IDF 사용: 단순 카운트보다 효과적인 가중치 벡터화",
    "n-gram 활용: 단어 조합으로 문맥 반영 (ngram_range=(1,2))",
    "불용어 커스터마이징: 도메인별 불용어 리스트 구성",
    "스테밍/레마타이제이션: 단어 원형 추출로 어휘 축소",
    "alpha 튜닝: 라플라스 스무딩 파라미터 조정 (default=1.0)",
]
y = numbered_list(s, SAFE_L + 0.2, y, CONTENT_W - 0.3, improve_items, sz=13, gap=0.35)


# ━━━━━━━━━━━ S33: 심화 - 완전한 파이프라인 코드 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "심화: 완전한 스팸 분류 파이프라인", BG_DEEP)

code(s, COL1_L, y, HALF_W, 5.5,
     "import re\n"
     "import pandas as pd\n"
     "from sklearn.pipeline import Pipeline\n"
     "from sklearn.feature_extraction.text \\\n"
     "    import TfidfVectorizer\n"
     "from sklearn.naive_bayes import (\n"
     "    MultinomialNB, BernoulliNB)\n"
     "from sklearn.model_selection import (\n"
     "    train_test_split, cross_val_score)\n"
     "from sklearn.metrics import (\n"
     "    classification_report,\n"
     "    confusion_matrix)\n"
     "\n"
     "# 전처리 (정규표현식 버전)\n"
     "def preprocess(text):\n"
     "    text = re.sub(r'[^a-zA-Z\\s]',\n"
     "                  '', text)\n"
     "    return text.lower().strip()\n"
     "\n"
     "data = pd.read_csv('spam.csv')\n"
     "data['text'] = \\\n"
     "    data['text'].apply(preprocess)\n"
     "data['target'] = data['target'].map(\n"
     "    {'spam': 1, 'ham': 0})")

code(s, COL2_L, y, HALF_W, 5.5,
     "X = data['text']\n"
     "y = data['target']\n"
     "X_tr, X_te, y_tr, y_te = \\\n"
     "    train_test_split(X, y,\n"
     "    test_size=0.2, random_state=42)\n"
     "\n"
     "# 여러 모델 비교\n"
     "models = {\n"
     "    'MNB': Pipeline([\n"
     "        ('tf', TfidfVectorizer(\n"
     "            max_features=5000,\n"
     "            ngram_range=(1,2))),\n"
     "        ('clf', MultinomialNB(\n"
     "            alpha=0.1))]),\n"
     "    'BNB': Pipeline([\n"
     "        ('tf', TfidfVectorizer(\n"
     "            max_features=5000)),\n"
     "        ('clf', BernoulliNB())])\n"
     "}\n"
     "\n"
     "for name, model in models.items():\n"
     "    scores = cross_val_score(\n"
     "        model, X, y, cv=5,\n"
     "        scoring='accuracy')\n"
     "    print(f'{name}: '\n"
     "          f'{scores.mean():.4f}')")


# ━━━━━━━━━━━ S34: 참고문헌 & 추가 학습 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
y = title_slide_setup(s, "참고 논문 & 추가 학습 안내", BG_DEEP)

mtb(s, SAFE_L, y, CONTENT_W, 2.5, [
    ("참고 논문", 18, GOLD, True),
    ("", 4, PARCHMENT),
    ("1. McCallum & Nigam (1998)", 14, PARCHMENT, True),
    ("   'A Comparison of Event Models for Naive Bayes Text Classification'", 12, MUTED, False, True),
    ("   → Multinomial vs Multivariate Bernoulli 모델 비교", 12, MUTED),
    ("", 4, PARCHMENT),
    ("2. Rennie, Shih, Teevan, Karger (2003)", 14, PARCHMENT, True),
    ("   'Tackling the Poor Assumptions of Naive Bayes Text Classifiers'", 12, MUTED, False, True),
    ("   → NB의 한계와 개선 방법 (TF-IDF 변환, 길이 정규화 등)", 12, MUTED),
    ("", 4, PARCHMENT),
    ("3. Zhang (2004)", 14, PARCHMENT, True),
    ("   'The Optimality of Naive Bayes'", 12, MUTED, False, True),
    ("   → 독립 가정이 깨져도 NB가 잘 작동하는 이론적 분석", 12, MUTED),
])
y += 2.7

mtb(s, COL1_L, y, HALF_W, 2.5, [
    ("scikit-learn 공식 문서", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("• sklearn.naive_bayes 모듈", 13, PARCHMENT),
    ("  MultinomialNB, GaussianNB, BernoulliNB", 12, MUTED),
    ("• sklearn.feature_extraction.text", 13, PARCHMENT),
    ("  CountVectorizer, TfidfVectorizer", 12, MUTED),
    ("• sklearn.metrics", 13, PARCHMENT),
    ("  confusion_matrix, classification_report", 12, MUTED),
])

mtb(s, COL2_L, y, HALF_W, 2.5, [
    ("다음 학습 내용", 16, GOLD, True),
    ("", 4, PARCHMENT),
    ("• 10장: 재현율(Recall), 정밀도(Precision),", 13, PARCHMENT),
    ("  F1 Score", 13, PARCHMENT),
    ("• 11장: AUC (Area Under Curve)", 13, PARCHMENT),
    ("", 4, PARCHMENT),
    ("NLTK 공식 문서", 16, GOLD, True),
    ("• nltk.org — stopwords, tokenizers", 13, PARCHMENT),
    ("• 한국어 불용어:", 13, PARCHMENT),
    ("  www.ranks.nl/stopwords/korean", 12, MUTED),
])


# ━━━━━━━━━━━ S35: 마무리 ━━━━━━━━━━━
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DEEP)
draw_border(s)

tb(s, SAFE_L, 1.5, CONTENT_W, 0.4, "C H A P T E R  7",
   sz=14, color=MUTED, font="Consolas", align=PP_ALIGN.CENTER)
tb(s, SAFE_L, 2.2, CONTENT_W, 0.9, "나이브 베이즈",
   sz=48, color=GOLD, bold=True, italic=True, align=PP_ALIGN.CENTER)
hline(s, 3.3, SLIDE_W/2 - 2.5, 5.0)
tb(s, SAFE_L, 3.6, CONTENT_W, 0.6, "수고하셨습니다",
   sz=26, color=PARCHMENT, align=PP_ALIGN.CENTER)
mtb(s, SAFE_L, 4.8, CONTENT_W, 1.5, [
    ("다음 단계 안내", 16, GOLD, True, False),
    ("", 8, PARCHMENT),
    ("재현율(Recall), 정밀도(Precision), F1 Score → 10장", 14, MUTED),
    ("AUC (Area Under Curve) → 11장", 14, MUTED),
    ("", 8, PARCHMENT),
    ("과제: 베이즈 정리를 직접 손으로 풀어보기", 14, PARCHMENT),
    ("과제: GaussianNB, BernoulliNB(binarize=True) 비교 실험", 14, PARCHMENT),
])


# ══════════════════════════════════════════════════════
# Add slide numbers
# ══════════════════════════════════════════════════════
total = len(prs.slides)
for i, sl in enumerate(prs.slides):
    slide_num(sl, i + 1, total)

# Save
output_path = r"F:\minpodata\기계학습\기계학습\7장\7장_책버전.pptx"
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Total slides: {total}")
