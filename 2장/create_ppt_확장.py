"""2장 머신러닝을 위한 파이썬 - 확장 강의 PPT 생성 스크립트 (상세 버전)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = '맑은 고딕'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    total_w = sum(col_widths)
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(1), Inches(2.5), Inches(11.333), Inches(3), SECTION_BG, accent)
    add_accent_line(s, Inches(5.5), Inches(2.6), Inches(2.333), accent)
    add_text(s, Inches(1), Inches(2.8), Inches(11.333), Inches(0.5),
             f"SECTION {section_num}", font_size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.4), Inches(11.333), Inches(1),
             title, font_size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.4), Inches(11.333), Inches(0.6),
             subtitle, font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG, ACCENT_BLUE)
add_accent_line(s, Inches(4), Inches(1.5), Inches(5.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.8), Inches(11.333), Inches(0.5),
         "CHAPTER 2", font_size=22, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.5), Inches(11.333), Inches(1.2),
         "머신러닝을 위한 파이썬", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.333), Inches(0.5),
         "Python for Machine Learning", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.8),
         "파이썬 언어의 역사, 기초 문법, 자료구조, 함수/클래스 설계, 성능 최적화, 개발환경",
         font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
         "기계학습 | Machine Learning", font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 2: 목차 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, None, "목차 (Table of Contents)", "Chapter 2 전체 구성")
toc_left = [
    "2.1  왜 Python인가",
    "2.2  기초 문법 (변수, 연산자, 제어문)",
    "2.3  자료구조 심화 (list, tuple, dict, set)",
    "2.4  함수와 클래스 (sklearn 패턴)",
    "2.5  성능 최적화 (프로파일링, 벡터화, GIL)",
    "2.6  개발환경 (conda, Jupyter, IDE)",
]
toc_right = [
    "2.7   논문 리뷰 (5편)",
    "2.8   실습: Python 프로파일링",
    "2.9   실습: 자료구조 벤치마크",
    "2.10  실습: 함수형 프로그래밍 ML",
    "2.11  End-to-end ML 파이프라인",
    "2.12  핵심 요약 + 복습 질문",
]
add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
         "Part I: 이론", toc_left, ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
         "Part II: 실습 & 논문", toc_right, ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# 슬라이드 3: 학습 목표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, None, "학습 목표", "이 장을 마치면 할 수 있는 것들")
objectives = [
    "1. 파이썬이 ML 표준 언어가 된 역사적 배경을 설명할 수 있다",
    "2. 변수, 자료형, 연산자, 제어문 등 기초 문법을 활용할 수 있다",
    "3. list, tuple, dict, set의 시간 복잡도를 비교하고 적절히 선택할 수 있다",
    "4. 함수, 클래스, sklearn의 fit/predict/transform 패턴을 구현할 수 있다",
    "5. NumPy 벡터화, 프로파일링, GIL 우회 등 성능 최적화 기법을 적용할 수 있다",
    "6. conda/venv 가상환경, Jupyter Notebook 환경을 구축할 수 있다",
    "7. 순수 파이썬만으로 End-to-end ML 파이프라인을 구현할 수 있다",
]
add_bullet_list(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(5),
                objectives, font_size=18, color=LIGHT_GRAY, spacing=Pt(12))

# ============================================================
# SECTION 1: 왜 Python인가
# ============================================================
section_divider("왜 Python인가", "ML 표준 언어의 역사와 이유", 1, ACCENT_BLUE)

# 슬라이드 5: 파이썬 역사 타임라인
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 1, "ML 표준 언어가 된 역사", "1991년부터 현재까지의 이정표")
add_table_slide(s,
    ["연도", "이정표", "의의"],
    [
        ["1991", "Python 0.9 공개", "van Rossum의 'batteries included' 철학"],
        ["2001", "SciPy 프로젝트 시작", "과학 컴퓨팅 생태계의 씨앗"],
        ["2006", "NumPy 1.0 (Oliphant)", "고성능 ndarray - ML의 핵심 인프라"],
        ["2007", "IPython 발표", "인터랙티브 과학 컴퓨팅 환경 확립"],
        ["2010", "scikit-learn 0.1 공개", "일관된 fit/predict/transform API"],
        ["2014", "Jupyter Notebook 분리", "언어 독립적 리터릿 컴퓨팅 플랫폼"],
        ["2015", "TensorFlow 공개", "딥러닝 프레임워크 시대 개막"],
        ["2016", "PyTorch 공개", "동적 계산 그래프, 연구 친화적"],
        ["2020", "Raschka et al. 서베이", "파이썬 ML 생태계 전체 조망"],
    ],
    Inches(0.8), Inches(2.0), [1.2, 3.5, 7.0],
    row_height=0.5, font_size=12, header_font_size=13)

# 슬라이드 6: 파이썬이 ML 표준이 된 3가지 이유
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 1, "파이썬이 ML 표준이 된 3가지 이유")
add_card(s, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5),
         "1. 접착 언어 (Glue Language)", [
             "C/Fortran 고성능 라이브러리를",
             "파이썬 인터페이스로 감싸서 사용",
             "BLAS, LAPACK, cuDNN 등",
             "개발 생산성 + 실행 성능 동시 확보",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.5),
         "2. 풍부한 생태계", [
             "NumPy, Pandas, scikit-learn",
             "TensorFlow, PyTorch, matplotlib",
             "ML 파이프라인의 모든 단계 지원",
             "데이터 수집 → 전처리 → 모델링 → 배포",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.5),
         "3. 낮은 진입 장벽", [
             "동적 타이핑, 간결한 문법",
             "들여쓰기 기반 블록 구조",
             "비전공자도 빠르게 학습 가능",
             "TIOBE 인덱스 1위 (2024~)",
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# 슬라이드 7: R/Julia/MATLAB 비교
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 1, "Python vs R vs Julia vs MATLAB", "각 언어의 강점과 ML 적합성 비교")
add_table_slide(s,
    ["기준", "Python", "R", "Julia", "MATLAB"],
    [
        ["주력 분야", "범용+ML/DL", "통계 분석", "고성능 수치", "공학 시뮬레이션"],
        ["성능", "중간(NumPy보완)", "느림(Rcpp보완)", "빠름(C수준)", "중간"],
        ["DL 프레임워크", "TF, PyTorch", "제한적", "Flux.jl", "제한적"],
        ["커뮤니티", "매우 큼", "통계 중심 큼", "성장 중", "학교/기업 중심"],
        ["라이선스", "오픈소스(무료)", "오픈소스(무료)", "오픈소스(무료)", "상용(유료)"],
        ["산업 채택", "매우 높음", "중간", "낮음", "중간"],
    ],
    Inches(0.5), Inches(2.2), [2.0, 2.2, 2.2, 2.2, 2.6],
    row_height=0.5, font_size=12, header_font_size=13)
add_text(s, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
         "핵심: ML 전체 파이프라인을 하나의 언어로 커버하는 생태계의 폭에서 파이썬을 능가하는 언어는 없다",
         font_size=15, color=ACCENT_CYAN, bold=True)

# 슬라이드 8: The Zen of Python
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 1, "파이썬 설계 철학: The Zen of Python", "import this")
add_code_block(s, Inches(0.8), Inches(2.2), Inches(7), Inches(3.5), [
    "import this",
    "",
    "# Beautiful is better than ugly.",
    "#   → 가독성 최우선",
    "# Explicit is better than implicit.",
    "#   → 명시적 코드 선호",
    "# Simple is better than complex.",
    "#   → 단순한 해법 추구",
    "# There should be one obvious way to do it.",
    "#   → 일관된 API 설계",
], font_size=14)
add_card(s, Inches(8.3), Inches(2.2), Inches(4.4), Inches(3.5),
         "scikit-learn에 반영된 철학", [
             "fit() → 학습",
             "predict() → 예측",
             "transform() → 변환",
             "일관된 API = 'one obvious way'",
             "모든 모델이 동일 패턴 따름",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# SECTION 2: 기초 문법
# ============================================================
section_divider("기초 문법", "변수, 자료형, 연산자, 제어문", 2, ACCENT_GREEN)

# 슬라이드 10: 변수와 동적 타이핑
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "변수와 동적 타이핑 (Dynamic Typing)", "자료형 선언 없이 값에 따라 자동 결정")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(3.0), [
    "# 변수 할당 - 자료형 선언 불필요",
    "x = 42          # int",
    "y = 3.14        # float",
    "name = 'ML'     # str",
    "is_valid = True  # bool",
    "",
    "# 동적 타이핑: 같은 변수에 다른 타입 재할당",
    "x = '이제 문자열'  # int → str로 변경",
    "print(type(x))     # <class 'str'>",
], font_size=13)
add_card(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(1.8),
         "ML 변수 명명 관례", [
             "X = 특성 행렬 (대문자), y = 타겟 벡터 (소문자)",
             "X_train / X_test = 분할된 데이터",
             "n_samples, n_features = 데이터 차원",
         ], ACCENT_CYAN, ACCENT_CYAN)
add_card(s, Inches(7.0), Inches(4.3), Inches(5.8), Inches(1.5),
         "주의: 정적 vs 동적 타이핑", [
             "정적(C, Java): 컴파일 시 타입 오류 발견 → 안전",
             "동적(Python): 런타임 시 타입 오류 발견 → 유연",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 슬라이드 11: 자료형과 산술 연산자
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "자료형과 산술 연산자", "4가지 기본 자료형 + 7가지 연산자")
add_table_slide(s,
    ["자료형", "표기", "예시", "ML 맥락"],
    [
        ["정수", "int", "42", "클래스 레이블, 인덱스"],
        ["실수", "float", "3.14", "특성값, 가중치, 손실"],
        ["문자열", "str", "'hello'", "범주형 특성, 텍스트"],
        ["논리형", "bool", "True/False", "마스크, 조건 필터"],
    ],
    Inches(0.5), Inches(2.0), [2.0, 1.5, 2.0, 5.0],
    row_height=0.45, font_size=13)
add_table_slide(s,
    ["연산자", "의미", "예시", "결과"],
    [
        ["+", "덧셈", "3 + 2", "5"],
        ["-", "뺄셈", "3 - 2", "1"],
        ["*", "곱셈", "3 * 2", "6"],
        ["/", "나눗셈", "7 / 2", "3.5 (항상 float)"],
        ["//", "몫", "7 // 2", "3"],
        ["%", "나머지", "7 % 2", "1"],
        ["**", "거듭제곱", "3 ** 2", "9"],
    ],
    Inches(0.5), Inches(4.4), [1.5, 1.5, 2.0, 5.5],
    row_height=0.35, font_size=12)

# 슬라이드 12: float 정밀도와 ML
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "float 정밀도가 ML에 미치는 영향", "IEEE 754 부동소수점의 한계와 대처")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(2.5), [
    "# IEEE 754 부동소수점의 한계",
    "print(0.1 + 0.2)         # 0.30000000000000004",
    "print(0.1 + 0.2 == 0.3)  # False !!",
    "",
    "# 올바른 비교 방법",
    "import math",
    "print(math.isclose(0.1 + 0.2, 0.3))  # True",
], font_size=13)
add_card(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.5),
         "ML에서 정밀도가 중요한 이유", [
             "float32 정밀도 ≈ 10⁻⁷",
             "float64 정밀도 ≈ 10⁻¹⁶",
             "경사하강법: η=10⁻⁶일 때 float32에서 기울기 소실",
             "손실함수: log(p)에서 p→0이면 -∞ 발산",
             "  → np.clip(p, 1e-15, 1) 필요",
             "혼합 정밀도: float16 + float32로 GPU 최적화",
         ], ACCENT_RED, ACCENT_RED)
add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.8),
         "수식: float32 정밀도 ≈ 10⁻⁷,  float64 정밀도 ≈ 10⁻¹⁶",
         font_size=18, color=ACCENT_ORANGE, bold=True)

# 슬라이드 13: 조건문과 반복문
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "조건문과 반복문", "if-elif-else, for, while, enumerate")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.5), [
    "# 조건문: if-elif-else",
    "score = 85",
    "if score >= 90:",
    "    grade = 'A'",
    "elif score >= 80:",
    "    grade = 'B'",
    "else:",
    "    grade = 'C'",
    "",
    "# for 반복문",
    "for i in range(5):",
    "    print(i, end=' ')  # 0 1 2 3 4",
    "",
    "# enumerate: 인덱스 + 값 동시",
    "fruits = ['사과', '바나나', '체리']",
    "for idx, fruit in enumerate(fruits):",
    "    print(f'{idx}: {fruit}')",
], font_size=12)
add_code_block(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.5), [
    "# f-string: ML에서 자주 쓰는 출력 패턴",
    "epoch, loss, acc = 10, 0.0342, 0.9821",
    "print(f'Epoch {epoch:3d} | Loss: {loss:.4f}'",
    "      f' | Acc: {acc:.4f}')",
    "# Epoch  10 | Loss: 0.0342 | Acc: 0.9821",
], font_size=12)
add_card(s, Inches(6.8), Inches(5.0), Inches(5.8), Inches(1.5),
         "f-string 포매팅", [
             ":3d = 정수 3자리, :.4f = 소수점 4자리",
             ":.2% = 퍼센트 표시 (0.95 → 95.00%)",
         ], ACCENT_CYAN, ACCENT_CYAN)

# 슬라이드 14: 리스트 컴프리헨션
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 2, "리스트 컴프리헨션", "for문을 한 줄로 축약하는 파이썬 고유 문법")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4.0), [
    "# 기본: [표현식 for 변수 in iterable]",
    "squares = [x ** 2 for x in range(1, 6)]       # [1, 4, 9, 16, 25]",
    "",
    "# 조건 포함: [표현식 for 변수 in iterable if 조건]",
    "evens = [x for x in range(1, 21) if x % 2 == 0]  # [2, 4, ..., 20]",
    "",
    "# if-else 포함: [표현식1 if 조건 else 표현식2 for 변수 in iterable]",
    "labels = [1 if x > 0 else 0 for x in [-1, 3, -2, 5, 0]]  # [0, 1, 0, 1, 0]",
    "",
    "# ML 활용 예시",
    "# 데이터 전처리: 결측값 제거",
    "cleaned = [x for x in raw_data if x is not None]",
    "# 레이블 인코딩",
    "encoded = [1 if label == 'positive' else 0 for label in labels]",
], font_size=13)
add_text(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
         "ML 핵심 활용: 데이터 전처리, 특성 추출, 레이블 인코딩에서 매우 자주 사용",
         font_size=16, color=ACCENT_GREEN, bold=True)

# ============================================================
# SECTION 3: 자료구조 심화
# ============================================================
section_divider("자료구조 심화", "list, tuple, dict, set의 특성과 시간 복잡도", 3, ACCENT_ORANGE)

# 슬라이드 16: 4가지 자료구조 비교
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 3, "4가지 기본 자료구조 비교", "특성과 용도")
add_table_slide(s,
    ["자료구조", "기호", "변경 가능", "순서 보장", "중복 허용", "해시 기반"],
    [
        ["리스트 (list)", "[ ]", "O (mutable)", "O", "O", "X"],
        ["튜플 (tuple)", "( )", "X (immutable)", "O", "O", "X"],
        ["딕셔너리 (dict)", "{ }", "O (mutable)", "O (3.7+)", "키 X", "O"],
        ["세트 (set)", "{ }", "O (mutable)", "X", "X", "O"],
    ],
    Inches(0.5), Inches(2.0), [2.5, 1.2, 1.8, 1.8, 1.8, 1.8],
    row_height=0.55, font_size=13)
add_text(s, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
         "핵심: dict와 set은 해시 테이블 기반이므로 탐색이 O(1). 대용량 데이터에서 list의 O(n) 대비 압도적",
         font_size=16, color=ACCENT_ORANGE, bold=True)

# 슬라이드 17: 시간 복잡도 비교
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 3, "시간 복잡도 비교 (Big-O)", "연산별 성능 특성")
add_table_slide(s,
    ["연산", "list", "tuple", "dict", "set"],
    [
        ["인덱스 접근 [i]", "O(1)", "O(1)", "O(1)", "N/A"],
        ["탐색 (in)", "O(n)", "O(n)", "O(1)", "O(1)"],
        ["끝 삽입", "O(1)*", "N/A", "O(1)", "O(1)"],
        ["앞 삽입", "O(n)", "N/A", "N/A", "N/A"],
        ["끝 삭제", "O(1)", "N/A", "O(1)", "O(1)"],
        ["앞 삭제", "O(n)", "N/A", "N/A", "N/A"],
    ],
    Inches(0.8), Inches(2.0), [2.5, 1.8, 1.8, 1.8, 1.8],
    row_height=0.5, font_size=13)
add_text(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
         "*O(1) 분할상환(amortized): 리스트 용량 초과 시 내부적으로 배열을 2배 확장",
         font_size=14, color=DARK_GRAY)

# 슬라이드 18: 리스트 상세
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 3, "리스트 (list) 상세", "가변, 순서 보장, 인덱싱/슬라이싱")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(4.5), [
    "# 생성과 인덱싱",
    "data = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10]",
    "print(data[0])     # 1 (첫 번째)",
    "print(data[-1])    # 10 (마지막)",
    "",
    "# 슬라이싱: [시작:끝:간격]",
    "print(data[2:5])   # [3, 4, 5]",
    "print(data[::2])   # [1, 3, 5, 7, 9]",
    "print(data[::-1])  # [10, 9, ..., 1] 역순",
    "",
    "# 주요 메서드",
    "data.append(100)       # 끝에 추가 O(1)",
    "data.insert(0, 99)     # 앞에 삽입 O(n)",
    "data.remove(100)       # 값으로 삭제 O(n)",
    "removed = data.pop()   # 끝에서 제거 O(1)",
], font_size=12)
add_card(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.5),
         "슬라이싱 규칙", [
             "[시작:끝] → 끝 인덱스 미포함",
             "[::간격] → 간격만큼 건너뛰기",
             "[::-1] → 역순 (전체 뒤집기)",
             "ML에서 자주 사용:",
             "  X_train = data[:split]",
             "  X_test = data[split:]",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.5),
         "성능 주의", [
             "append → O(1), insert(0) → O(n)",
             "앞 삽입이 잦으면 collections.deque 사용",
         ], ACCENT_RED, ACCENT_RED)

# 슬라이드 19: 딕셔너리와 세트
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 3, "딕셔너리 (dict) & 세트 (set)", "해시 테이블 기반 O(1) 탐색")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5), [
    "# 딕셔너리: {키: 값}",
    "params = {",
    "    'n_estimators': 100,",
    "    'max_depth': 5,",
    "    'learning_rate': 0.1,",
    "}",
    "# model = XGBClassifier(**params)  # 언패킹",
], font_size=12)
add_code_block(s, Inches(0.6), Inches(5.0), Inches(5.8), Inches(2.0), [
    "# 세트: 중복 자동 제거, 집합 연산",
    "train_feats = {'age', 'income', 'score'}",
    "test_feats = {'age', 'income', 'name'}",
    "missing = train_feats - test_feats  # {'score'}",
], font_size=12)
add_card(s, Inches(6.8), Inches(2.2), Inches(6), Inches(4.8),
         "ML에서의 자료구조 활용 가이드", [
             "특성값 저장 (순서 중요) → list / tuple",
             "하이퍼파라미터 관리 → dict (**kwargs 전달)",
             "범주 레이블 매핑 → dict (문자열→정수 인코딩)",
             "고유 카테고리 추출 → set (중복 제거, O(1) 탐색)",
             "모델 설정 (불변) → tuple (실수 변경 방지)",
             "빠른 멤버십 검사 → set / dict",
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# SECTION 4: 함수와 클래스
# ============================================================
section_divider("함수와 클래스", "def, lambda, 클래스, sklearn 패턴", 4, ACCENT_PURPLE)

# 슬라이드 21: 함수 정의
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 4, "함수 정의와 고급 기능", "기본값, *args, **kwargs, 다중 반환")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(4.5), [
    "# 기본값 매개변수",
    "def train(data, epochs=10, lr=0.01):",
    "    print(f'학습: {epochs} 에폭, 학습률={lr}')",
    "",
    "train(data)              # 기본값 사용",
    "train(data, epochs=50)   # 에폭만 변경",
    "",
    "# *args: 가변 위치 인자 (튜플)",
    "def add_all(*args):",
    "    return sum(args)",
    "",
    "# **kwargs: 가변 키워드 인자 (딕셔너리)",
    "def create_model(**kwargs):",
    "    for key, val in kwargs.items():",
    "        print(f'  {key} = {val}')",
], font_size=12)
add_card(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.5),
         "여러 값 반환 (튜플)", [
             "def calc_stats(numbers):",
             "    total = sum(numbers)",
             "    avg = total / len(numbers)",
             "    return total, avg, min(numbers)",
             "",
             "total, avg, min_v = calc_stats([85,90,78])",
         ], ACCENT_CYAN, ACCENT_CYAN)
add_card(s, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.5),
         "ML에서 **kwargs 활용", [
             "params = {'n_estimators':100, 'max_depth':5}",
             "model = RandomForestClassifier(**params)",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 슬라이드 22: lambda와 고차 함수
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 4, "lambda 함수와 고차 함수", "sorted, map, filter에서 활용")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4.5), [
    "# lambda: 익명 함수 (한 줄 함수)",
    "square = lambda x: x ** 2",
    "print(square(5))  # 25",
    "",
    "# sorted() + lambda: 특성 중요도 정렬",
    "feature_importance = [('age', 0.35), ('income', 0.52), ('score', 0.13)]",
    "sorted_features = sorted(feature_importance, key=lambda x: x[1], reverse=True)",
    "# [('income', 0.52), ('age', 0.35), ('score', 0.13)]",
    "",
    "# map(): 데이터 변환  |  filter(): 데이터 필터링",
    "data = [1, -2, 3, -4, 5]",
    "positives = list(filter(lambda x: x > 0, data))   # [1, 3, 5]",
    "squared = list(map(lambda x: x ** 2, data))        # [1, 4, 9, 16, 25]",
], font_size=13)

# 슬라이드 23: sklearn fit/predict/transform 패턴
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 4, "sklearn의 fit/predict/transform 패턴", "Estimator API - 모든 모델의 일관된 인터페이스")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(7), Inches(4.5), [
    "class SimpleScaler:",
    "    def __init__(self, feature_range=(0, 1)):",
    "        self.feature_range = feature_range",
    "        self.min_ = None",
    "        self.max_ = None",
    "",
    "    def fit(self, X):  # 학습 데이터에서 통계량 학습",
    "        self.min_ = min(X)",
    "        self.max_ = max(X)",
    "        return self",
    "",
    "    def transform(self, X):  # 학습된 통계량으로 변환",
    "        range_ = self.max_ - self.min_",
    "        lo, hi = self.feature_range",
    "        return [(x-self.min_)/range_*(hi-lo)+lo",
    "                for x in X]",
], font_size=11)
add_table_slide(s,
    ["메서드", "역할", "호출 시점"],
    [
        ["fit(X)", "학습 데이터에서 파라미터 학습", "학습 시"],
        ["predict(X)", "학습된 모델로 예측 수행", "추론 시"],
        ["transform(X)", "학습된 파라미터로 데이터 변환", "학습/추론 시"],
        ["fit_transform(X)", "fit + transform 한 번에", "학습 시 (편의)"],
    ],
    Inches(7.8), Inches(2.2), [2.0, 2.0, 1.2],
    row_height=0.55, font_size=11, header_font_size=12)
add_text(s, Inches(7.8), Inches(5.0), Inches(5), Inches(1.0),
         "중요: fit()은 반드시 학습 데이터에서만!\n테스트 데이터에 fit() 호출 = 데이터 누출(data leakage)",
         font_size=14, color=ACCENT_RED, bold=True)

# 슬라이드 24: 클래스 상속과 다형성
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 4, "클래스 상속과 다형성", "BaseModel → 구체적 모델 구현 패턴")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(4.5), [
    "# 기본 모델 클래스 (추상 기반)",
    "class BaseModel:",
    "    def fit(self, X, y):",
    "        raise NotImplementedError",
    "    def predict(self, X):",
    "        raise NotImplementedError",
    "    def score(self, X, y):",
    "        preds = self.predict(X)",
    "        correct = sum(1 for p,t",
    "                      in zip(preds,y) if p==t)",
    "        return correct / len(y)",
    "",
    "# 구체적 모델 구현",
    "class MajorityClassifier(BaseModel):",
    "    def fit(self, X, y):",
    "        from collections import Counter",
    "        self.majority_ = Counter(y)",
    "            .most_common(1)[0][0]",
    "        return self",
], font_size=11)
add_card(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.5),
         "다형성 (Polymorphism)의 장점", [
             "동일한 인터페이스로 다른 모델 사용 가능",
             "",
             "model = MajorityClassifier()",
             "model.fit(None, [0,0,1,0,1,0])",
             "model.predict([1,2,3])  # [0, 0, 0]",
             "model.score(None, [0,0,1])  # 0.667",
             "",
             "sklearn의 모든 모델이 이 패턴:",
             "  LogisticRegression, SVC, RandomForest...",
             "  모두 fit() → predict() → score() 지원",
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# ============================================================
# SECTION 5: 성능 최적화
# ============================================================
section_divider("성능 최적화", "프로파일링, 벡터화, GIL, 메모리 관리", 5, ACCENT_RED)

# 슬라이드 26: 프로파일링
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "프로파일링: 병목 지점 찾기", "측정하지 않으면 최적화하지 마라")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(4.5), [
    "import time, cProfile",
    "",
    "# 방법 1: time.perf_counter()",
    "start = time.perf_counter()",
    "result = sum(x**2 for x in range(1_000_000))",
    "elapsed = time.perf_counter() - start",
    "print(f'실행 시간: {elapsed:.4f}초')",
    "",
    "# 방법 2: cProfile (함수별 상세)",
    "def my_function():",
    "    data = [x**2 for x in range(100_000)]",
    "    return sorted(data, reverse=True)",
    "cProfile.run('my_function()')",
    "",
    "# 방법 3: Jupyter %timeit",
    "# %timeit sum(range(1000))",
], font_size=12)
add_card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.5),
         "프로파일링 3단계", [
             "1단계: time.perf_counter()",
             "   → 간단한 구간 측정",
             "   → 가장 빠르게 병목 확인",
             "",
             "2단계: cProfile",
             "   → 함수별 호출 횟수, 누적 시간",
             "   → 어떤 함수가 느린지 정밀 분석",
             "",
             "3단계: line_profiler",
             "   → 줄 단위 실행 시간 분석",
             "   → 정확한 병목 지점 특정",
         ], ACCENT_RED, ACCENT_RED)

# 슬라이드 27: NumPy 벡터화
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "벡터화: NumPy의 위력", "순수 파이썬 루프 대비 수십~수백 배 빠름")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(2.5), [
    "import numpy as np",
    "n = 1_000_000",
    "data_list = list(range(n))",
    "data_array = np.arange(n)",
    "",
    "# 순수 파이썬 루프: ~0.15초",
    "result_loop = [x ** 2 for x in data_list]",
    "",
    "# NumPy 벡터화: ~0.002초 (약 75배 빠름!)",
    "result_numpy = data_array ** 2",
], font_size=12)
add_card(s, Inches(0.6), Inches(5.0), Inches(6.5), Inches(2.0),
         "벡터화가 빠른 3가지 이유", [
             "1. C 수준 내부 루프: 인터프리터 오버헤드 제거",
             "2. 연속 메모리 접근: CPU 캐시 효율 극대화",
             "3. SIMD 명령어: 하드웨어 수준 병렬 처리",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(7.5), Inches(2.2), Inches(5.3), Inches(4.8),
         "속도 비교", [
             "for loop     ~150ms",
             "list comp.   ~100ms",
             "map()        ~120ms",
             "NumPy         ~2ms",
             "",
             "NumPy는 for 루프 대비 ~75x 빠름",
             "",
             "핵심 교훈:",
             "  수치 연산에는 반드시 NumPy 사용!",
             "  순수 파이썬에서는 컴프리헨션이 최적",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# 슬라이드 28: GIL과 병렬 처리
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "GIL (Global Interpreter Lock)", "CPython의 전역 잠금과 ML에서의 대처")
add_card(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5),
         "GIL이란?", [
             "CPython의 전역 잠금 메커니즘",
             "한 번에 하나의 스레드만 파이썬 바이트코드 실행",
             "CPU 바운드: 멀티스레딩으로 성능 향상 불가",
             "I/O 바운드: 멀티스레딩 효과 있음 (GIL 해제)",
         ], ACCENT_RED, ACCENT_RED)
add_table_slide(s,
    ["전략", "도구", "적용 상황"],
    [
        ["멀티프로세싱", "multiprocessing, joblib", "CPU 병렬 학습"],
        ["NumPy 내부", "BLAS/MKL", "행렬 연산 자동 병렬"],
        ["외부 라이브러리", "TensorFlow, PyTorch", "GPU 연산 (GIL 무관)"],
        ["비동기 I/O", "asyncio", "데이터 로딩 파이프라인"],
    ],
    Inches(0.6), Inches(5.0), [2.5, 3.5, 5.5],
    row_height=0.5, font_size=12)
add_code_block(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.0), [
    "# scikit-learn의 n_jobs 파라미터",
    "# 내부적으로 멀티프로세싱 사용",
    "from sklearn.ensemble import RandomForestClassifier",
    "model = RandomForestClassifier(",
    "    n_estimators=100, n_jobs=-1  # 모든 코어",
    ")",
], font_size=12)

# 슬라이드 29: 메모리 관리
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 5, "메모리 관리", "자료구조별 메모리와 제너레이터")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(2.5), [
    "import sys",
    "",
    "# 자료구조별 메모리 사용량 (10,000개)",
    "list_mem  = sys.getsizeof(list(range(10000)))",
    "tuple_mem = sys.getsizeof(tuple(range(10000)))",
    "dict_mem  = sys.getsizeof({i:i for i in range(10000)})",
], font_size=12)
add_code_block(s, Inches(0.6), Inches(5.0), Inches(6), Inches(2.0), [
    "# 제너레이터: 메모리 효율적 대안",
    "gen = (x ** 2 for x in range(1_000_000))",
    "print(sys.getsizeof(gen))  # ~200 bytes!",
    "# 리스트는 ~8MB, 제너레이터는 ~200B",
], font_size=12)
add_card(s, Inches(6.5), Inches(2.2), Inches(6.3), Inches(4.8),
         "ML 메모리 관리 팁", [
             "list: ~87 KB (10K) → 튜플보다 약간 큼",
             "tuple: ~80 KB → immutable, 메모리 효율적",
             "dict: ~300+ KB → 해시 테이블 오버헤드",
             "set: ~500+ KB → 해시 테이블 오버헤드",
             "",
             "대용량 데이터 처리 전략:",
             "  1. 제너레이터 배치 로딩 (yield)",
             "  2. np.memmap (디스크 매핑 배열)",
             "  3. pd.read_csv(chunksize=1000)",
             "  4. float64 → float32 다운캐스팅",
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# ============================================================
# SECTION 6: 개발환경
# ============================================================
section_divider("개발환경", "가상환경, Jupyter, IDE, 디버깅", 6, ACCENT_CYAN)

# 슬라이드 31: 가상환경
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 6, "가상환경: conda vs venv", "프로젝트별 독립 환경 구축")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.0), [
    "# venv (Python 내장)",
    "python -m venv ml_env",
    "source ml_env/bin/activate     # Linux/Mac",
    "ml_env\\Scripts\\activate      # Windows",
    "pip install numpy pandas scikit-learn",
], font_size=12)
add_code_block(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.0), [
    "# conda (Anaconda/Miniconda)",
    "conda create -n ml_env python=3.11",
    "conda activate ml_env",
    "conda install numpy pandas scikit-learn",
], font_size=12)
add_table_slide(s,
    ["기준", "venv + pip", "conda"],
    [
        ["패키지 소스", "PyPI", "Anaconda + PyPI"],
        ["바이너리 관리", "제한적", "우수 (C 의존성 자동)"],
        ["GPU 설정", "수동", "간편 (cudatoolkit)"],
        ["추천", "가벼운 프로젝트", "ML/DL 프로젝트"],
    ],
    Inches(6.8), Inches(2.2), [1.8, 2.0, 2.0],
    row_height=0.5, font_size=12)

# 슬라이드 32: Jupyter vs IDE + 디버깅
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 6, "Jupyter Notebook vs IDE", "탐색적 분석 vs 프로덕션 코드")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
         "Jupyter Notebook", [
             "강점: 셀 단위 실행, 인라인 시각화, 마크다운",
             "적합: EDA, 프로토타이핑, 강의/발표",
             "매직: %timeit, %matplotlib inline",
             "한계: .ipynb JSON 형식 → Git diff 불편",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.5),
         "IDE (VS Code, PyCharm)", [
             "강점: 디버깅, 리팩토링, 자동완성, Git 통합",
             "적합: 프로덕션 코드, 패키지 개발",
             ".py 파일 → Git 친화적",
             "VS Code + Jupyter 확장 = 최고 조합",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         "실무 워크플로우 & 디버깅 기법", [
             "Jupyter에서 프로토타이핑 → .py 모듈로 리팩토링 → IDE에서 테스트/배포",
             "디버깅 4단계: print() → assert → pdb.set_trace() → IDE 디버거",
             "방어적 프로그래밍: assert len(X) == len(y), 'X와 y의 길이가 다릅니다!'",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# SECTION 7: 논문 리뷰
# ============================================================
section_divider("논문 리뷰", "5편의 핵심 논문 분석", 7, ACCENT_BLUE)

# 슬라이드 34: 논문 리뷰 (1/2)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 7, "핵심 논문 리뷰 (1/2)", "파이썬 생태계의 기초를 다진 논문들")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(2.0),
         "Van Rossum (2007) - Python Programming Language", [
             "USENIX Annual Technical Conference (Invited Talk)",
             "파이썬 창시자가 설명하는 설계 철학",
             "동적 타이핑, GIL, 확장성 아키텍처",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(2.0),
         "Millman & Aivazis (2011) - Python for Scientists", [
             "Computing in Science & Engineering, 13(2)",
             "SciPy 스택이 과학 컴퓨팅의 근간",
             "소프트웨어 엔지니어링 관행의 중요성",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(0.5), Inches(4.3), Inches(6), Inches(2.3),
         "Oliphant (2007) - Python for Scientific Computing", [
             "NumPy 창시자가 설명하는 기술적 기반",
             "ndarray: 연속 메모리 + stride 기반 레이아웃",
             "BLAS/LAPACK 활용, f2py/Cython 최적화",
             "ML 핵심: ndarray가 모든 ML 알고리즘의 인프라",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_card(s, Inches(6.8), Inches(4.3), Inches(6), Inches(2.3),
         "Perez & Granger (2007) - IPython", [
             "Jupyter Notebook의 전신 (커널-프론트엔드 분리)",
             "인터랙티브 컴퓨팅: 생각→코딩→실행→확인",
             "매직 명령어: %timeit, %debug, %matplotlib",
             "리터릿 컴퓨팅: 코드+결과+설명 통합",
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# 슬라이드 35: 논문 리뷰 (2/2)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 7, "필독 논문: Raschka et al. (2020)", "Machine Learning in Python - 생태계 서베이")
add_card(s, Inches(0.5), Inches(2.0), Inches(6), Inches(5.0),
         "Raschka et al. (2020) - ML in Python", [
             "Information, Vol. 11, No. 4, Article 193",
             "",
             "역사적 발전:",
             "  NumPy(2006) → scikit-learn(2010)",
             "  → TensorFlow(2015) → PyTorch(2016)",
             "",
             "프레임워크 비교:",
             "  scikit-learn: 일관된 API",
             "  TensorFlow: 분산 학습",
             "  PyTorch: 동적 그래프, 연구 친화적",
             "",
             "최신 트렌드:",
             "  AutoML, Interpretable ML (SHAP/LIME)",
             "  MLOps (MLflow/Kubeflow)",
         ], ACCENT_CYAN, ACCENT_CYAN)
add_card(s, Inches(6.8), Inches(2.0), Inches(6), Inches(5.0),
         "본 과목과의 연결", [
             "이 서베이 논문은 본 과목 전체를",
             "파이썬 생태계 맥락에서 이해하게 함",
             "",
             "데이터 처리 도구 진화:",
             "  Pandas → Dask(분산) → RAPIDS(GPU)",
             "",
             "다룰 알고리즘 계보:",
             "  선형회귀, 로지스틱회귀",
             "  의사결정나무, 랜덤포레스트",
             "  XGBoost, LightGBM",
             "  KMeans, PCA",
             "",
             "생태계 과제:",
             "  GIL, 프로덕션 성능, 라이브러리 호환",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# SECTION 8: 실습 - 프로파일링
# ============================================================
section_divider("실습: Python 프로파일링", "for 루프 vs 컴프리헨션 vs map vs NumPy 성능 비교", 8, ACCENT_GREEN)

# 슬라이드 37: 프로파일링 실습 코드 (1)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 8, "실습 코드: 성능 측정 유틸리티", "01_python_profiling.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4.5), [
    "import time, numpy as np, matplotlib.pyplot as plt",
    "",
    "def 시간측정(함수, *args, 반복횟수=5):",
    '    """함수의 실행 시간을 측정 (여러 번 실행 후 평균)"""',
    "    시간들 = []",
    "    for _ in range(반복횟수):",
    "        시작 = time.perf_counter()",
    "        함수(*args)",
    "        종료 = time.perf_counter()",
    "        시간들.append(종료 - 시작)",
    "    return np.mean(시간들)",
    "",
    "# 4가지 방법으로 제곱 연산 비교 (데이터 크기: 1,000,000)",
    "def for루프_제곱(data):     결과 = []; [결과.append(x**2) for x in data]; return 결과",
    "def 컴프리헨션_제곱(data):  return [x**2 for x in data]",
    "def map_제곱(data):        return list(map(lambda x: x**2, data))",
    "def numpy_제곱(data):      return data ** 2   # ← 핵심! C 수준 내부 루프",
], font_size=12)

# 슬라이드 38: 프로파일링 결과 요약
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 8, "프로파일링 실험 결과", "4가지 연산 + 3가지 실험")
add_table_slide(s,
    ["방법", "제곱 연산", "합산 연산", "필터링", "속도비(NumPy 대비)"],
    [
        ["for 루프", "~0.15s", "~0.05s", "~0.10s", "~75x 느림"],
        ["리스트 컴프리헨션", "~0.10s", "-", "~0.07s", "~50x 느림"],
        ["map()", "~0.12s", "-", "~0.09s", "~60x 느림"],
        ["NumPy 벡터화", "~0.002s", "~0.001s", "~0.003s", "기준 (1x)"],
    ],
    Inches(0.5), Inches(2.0), [2.5, 2.0, 2.0, 2.0, 2.5],
    row_height=0.55, font_size=13)
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         "핵심 교훈", [
             "1. 수치 연산에는 반드시 NumPy를 사용할 것 (수십~수백 배 차이)",
             "2. 순수 파이썬에서는 리스트 컴프리헨션이 가장 효율적",
             "3. 데이터가 클수록 NumPy의 이점이 극대화됨",
             "4. map()은 lambda 오버헤드로 컴프리헨션보다 느린 경우가 많음",
         ], ACCENT_GREEN, ACCENT_GREEN)

# ============================================================
# SECTION 9: 실습 - 자료구조 벤치마크
# ============================================================
section_divider("실습: 자료구조 벤치마크", "list, tuple, dict, set의 탐색/삽입/삭제 성능 비교", 9, ACCENT_ORANGE)

# 슬라이드 40: 벤치마크 실습 코드
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 9, "실습 코드: 탐색 성능 비교", "02_data_structures_benchmark.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4.5), [
    "# 탐색 성능 비교: O(n) vs O(1)의 실제 차이 측정",
    "크기들 = [1000, 5000, 10000, 50000, 100000, 500000]",
    "",
    "for 크기 in 크기들:",
    "    리스트_데이터 = list(range(크기))",
    "    딕셔너리_데이터 = {x: x for x in range(크기)}",
    "    세트_데이터 = set(range(크기))",
    "    탐색값들 = [random.randint(크기//2, 크기-1) for _ in range(100)]",
    "",
    "    # list: O(n) - 순차 탐색",
    "    def 리스트_탐색():",
    "        for v in 탐색값들: _ = v in 리스트_데이터",
    "",
    "    # dict/set: O(1) - 해시 테이블 즉시 조회",
    "    def 세트_탐색():",
    "        for v in 탐색값들: _ = v in 세트_데이터",
], font_size=12)

# 슬라이드 41: 벤치마크 결과
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 9, "벤치마크 결과 요약", "해시 기반 자료구조의 압도적 탐색 성능")
add_table_slide(s,
    ["연산", "list", "dict/set", "차이"],
    [
        ["탐색 (50만개)", "~0.5s", "~0.00005s", "~10,000x"],
        ["끝 삽입 (10만개)", "~0.005s (append)", "~0.006s", "유사"],
        ["앞 삽입 (1만개)", "~0.3s (insert)", "N/A", "list 느림"],
        ["끝 삭제 (1만개)", "~0.001s (pop)", "~0.001s (del)", "유사"],
        ["앞 삭제 (1만개)", "~0.05s (pop(0))", "N/A", "list 느림"],
    ],
    Inches(0.5), Inches(2.0), [3.0, 2.5, 2.5, 3.0],
    row_height=0.55, font_size=13)
add_card(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
         "실무 시사점", [
             "멤버십 검사(in)가 빈번하면: set 또는 dict 사용 (O(1) vs O(n))",
             "데이터 크기가 클수록: 해시 기반 자료구조의 이점이 극대화",
             "메모리: dict/set은 해시 테이블 오버헤드로 메모리를 더 많이 사용 (트레이드오프)",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# SECTION 10: 실습 - 함수형 프로그래밍
# ============================================================
section_divider("실습: 함수형 프로그래밍 ML", "map/filter/reduce, 데코레이터, 제너레이터", 10, ACCENT_PURPLE)

# 슬라이드 43: 함수형 데이터 파이프라인
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 10, "map/filter/reduce 데이터 파이프라인", "03_functional_programming_ml.py")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4.5), [
    "# 센서 데이터 전처리 파이프라인",
    "원시_온도데이터 = [round(random.gauss(25, 5), 1) for _ in range(100)]",
    "원시_온도데이터[10] = -99.0   # 센서 오류",
    "원시_온도데이터[75] = None    # 결측값",
    "",
    "# 단계 1: filter → 결측값(None) 제거",
    "유효 = list(filter(lambda x: x is not None, 원시_온도데이터))",
    "",
    "# 단계 2: filter → 이상값 제거 (0~50도 범위)",
    "정상 = list(filter(lambda x: 0 <= x <= 50, 유효))",
    "",
    "# 단계 3: map → 섭씨 → 화씨 변환",
    "화씨 = list(map(lambda c: round(c * 9/5 + 32, 1), 정상))",
    "",
    "# 단계 4: reduce → 통계량 계산",
    "from functools import reduce",
    "합계 = reduce(lambda acc, x: acc + x, 정상)",
], font_size=12)

# 슬라이드 44: 데코레이터와 제너레이터
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 10, "데코레이터 & 제너레이터 패턴", "함수 확장과 메모리 효율적 데이터 로딩")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.5), [
    "# 데코레이터: 함수 기능 확장",
    "def 타이밍(함수):",
    "    @functools.wraps(함수)",
    "    def 래퍼(*args, **kwargs):",
    "        시작 = time.perf_counter()",
    "        결과 = 함수(*args, **kwargs)",
    "        종료 = time.perf_counter()",
    "        print(f'{함수.__name__}: {종료-시작:.4f}초')",
    "        return 결과",
    "    return 래퍼",
    "",
    "@타이밍",
    "def 데이터_전처리(데이터):",
    "    # 자동으로 실행 시간 측정!",
    "    return [x for x in 데이터 if x > 0]",
], font_size=11)
add_code_block(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(3.0), [
    "# 제너레이터: 배치 데이터 로더",
    "def 배치_로더(데이터, 배치크기=32):",
    "    인덱스 = list(range(len(데이터)))",
    "    random.shuffle(인덱스)",
    "    for 시작 in range(0, len(인덱스), 배치크기):",
    "        배치 = 인덱스[시작:시작+배치크기]",
    "        yield [데이터[i] for i in 배치]",
    "",
    "# 사용: 메모리에 전체 로드하지 않음",
    "for 배치 in 배치_로더(데이터, 32):",
    "    학습(배치)  # 배치 단위 학습",
], font_size=11)
add_card(s, Inches(6.8), Inches(5.5), Inches(5.8), Inches(1.5),
         "@lru_cache: 메모이제이션", [
             "@functools.lru_cache(maxsize=128)",
             "def 피보나치(n): ...  # O(2^n) → O(n)",
         ], ACCENT_CYAN, ACCENT_CYAN)

# 슬라이드 45: 함수 합성 패턴
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 10, "함수 합성 (Function Composition)", "sklearn Pipeline의 함수형 버전")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(4.5), [
    "def 파이프라인(*함수들):",
    '    """함수를 순차 합성하는 파이프라인"""',
    "    def 합성_함수(입력):",
    "        결과 = 입력",
    "        for 함수 in 함수들:",
    "            결과 = 함수(결과)",
    "        return 결과",
    "    return 합성_함수",
    "",
    "# 전처리 함수들",
    "def 결측값_제거(d): return [x for x in d if x]",
    "def 이상값_제거(d): return [x for x in d if 0<=x<=50]",
    "def 표준화(d):",
    "    μ = sum(d)/len(d)",
    "    σ = (sum((x-μ)**2 for x in d)/len(d))**0.5",
    "    return [(x-μ)/σ for x in d]",
], font_size=11)
add_code_block(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.0), [
    "# 파이프라인 조합 & 실행",
    "ML_전처리 = 파이프라인(결측값_제거, 이상값_제거, 표준화)",
    "결과 = ML_전처리(원시_데이터)",
    "# 입력: 100개 → 결측 제거 → 이상치 제거 → 표준화",
], font_size=12)
add_card(s, Inches(7.0), Inches(4.5), Inches(5.8), Inches(2.5),
         "함수형 프로그래밍 핵심 정리", [
             "map(): 데이터 변환 (1:1 매핑)",
             "filter(): 데이터 필터링 (조건부 선택)",
             "reduce(): 집계 (리스트 → 단일 값)",
             "데코레이터: @타이밍, @로깅, @lru_cache",
             "제너레이터: yield 기반 배치 로딩, 메모리 절약",
             "함수 합성: sklearn Pipeline의 함수형 버전",
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# ============================================================
# SECTION 11: End-to-end ML 파이프라인
# ============================================================
section_divider("End-to-end ML 파이프라인", "순수 파이썬으로 KNN 분류기 구현", 11, ACCENT_BLUE)

# 슬라이드 47: 파이프라인 전체 구조
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 11, "순수 파이썬 KNN 분류기", "외부 라이브러리 없이 5단계 ML 파이프라인")
add_card(s, Inches(0.3), Inches(2.0), Inches(2.3), Inches(4.5),
         "1단계: 데이터 생성", [
             "합성 2D 데이터",
             "클래스 0: 중심 (2,2)",
             "클래스 1: 중심 (5,5)",
             "각 100개, 총 200개",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(2.8), Inches(2.0), Inches(2.3), Inches(4.5),
         "2단계: 전처리", [
             "Min-Max 정규화",
             "x' = (x-min)/(max-min)",
             "모든 특성을 0~1로",
             "스케일 통일",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(5.3), Inches(2.0), Inches(2.3), Inches(4.5),
         "3단계: 분할", [
             "학습/테스트 분할",
             "test_ratio = 0.2",
             "random.shuffle()",
             "80% 학습, 20% 테스트",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_card(s, Inches(7.8), Inches(2.0), Inches(2.7), Inches(4.5),
         "4단계: KNN 구현", [
             "유클리드 거리 계산",
             "k개 최근접 이웃 선택",
             "다수결 투표로 예측",
             "k=1,3,5,7,9 비교",
         ], ACCENT_PURPLE, ACCENT_PURPLE)
add_card(s, Inches(10.7), Inches(2.0), Inches(2.3), Inches(4.5),
         "5단계: 평가", [
             "정확도 (Accuracy)",
             "정밀도 (Precision)",
             "재현율 (Recall)",
             "F1-score",
         ], ACCENT_RED, ACCENT_RED)

# 슬라이드 48: 핵심 코드 - KNN
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 11, "핵심 코드: KNN 분류기 구현", "유클리드 거리 + 다수결 투표")
add_code_block(s, Inches(0.6), Inches(2.2), Inches(6), Inches(4.5), [
    "import math",
    "",
    "def euclidean_distance(a, b):",
    '    """유클리드 거리 계산"""',
    "    return math.sqrt(",
    "        sum((ai-bi)**2 for ai,bi in zip(a,b)))",
    "",
    "def knn_predict(X_train, y_train, x_query, k=5):",
    '    """KNN 예측: k개 최근접 이웃의 다수결"""',
    "    distances = [",
    "        (euclidean_distance(x_query, x_tr), label)",
    "        for x_tr, label",
    "        in zip(X_train, y_train)]",
    "    distances.sort(key=lambda x: x[0])",
    "    k_nearest = distances[:k]",
    "    votes = [label for _, label in k_nearest]",
    "    return max(set(votes), key=votes.count)",
], font_size=12)
add_card(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.5),
         "사용된 2장 개념 매핑", [
             "데이터 생성 → 함수, 리스트, for 루프",
             "전처리 → 리스트 컴프리헨션, min/max",
             "분할 → 슬라이싱, random.shuffle()",
             "KNN → lambda, sorted(), zip(), 튜플",
             "평가 → f-string, 산술 연산",
             "",
             "수식: d(x,y) = √Σ(xi - yi)²",
             "",
             "k값별 성능 비교:",
             "k=1: 노이즈에 민감 (과적합)",
             "k=5: 적절한 일반화",
             "k=9: 과도한 평활화 (과소적합)",
         ], ACCENT_BLUE, ACCENT_BLUE)

# ============================================================
# SECTION 12: 핵심 요약
# ============================================================
section_divider("핵심 요약 + 복습 질문", "2장 전체 정리", 12, ACCENT_CYAN)

# 슬라이드 50: 핵심 요약표
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 12, "핵심 요약표", "2장에서 배운 모든 개념")
add_table_slide(s,
    ["주제", "핵심 내용"],
    [
        ["왜 Python인가", "접착 언어 + 풍부한 생태계 + 낮은 진입 장벽"],
        ["자료형", "int, float, str, bool + 동적 타이핑"],
        ["리스트/튜플", "mutable vs immutable, 슬라이싱, O(1) 접근"],
        ["딕셔너리/세트", "해시 기반 O(1) 탐색, 하이퍼파라미터 관리"],
        ["함수", "def, lambda, *args, **kwargs, 다중 반환"],
        ["클래스", "fit/predict/transform 패턴 = sklearn API"],
        ["성능", "프로파일링 → NumPy 벡터화 → GIL 우회"],
        ["개발환경", "conda/venv + Jupyter(탐색) + IDE(프로덕션)"],
    ],
    Inches(0.5), Inches(2.0), [3.0, 8.5],
    row_height=0.5, font_size=13)

# 슬라이드 51: 수식 요약
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 12, "수식 요약", "2장에서 등장한 핵심 수식들")
add_card(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.0),
         "유클리드 거리 (Euclidean Distance)", [
             "d(x, y) = sqrt( Σ (xi - yi)² )",
             "KNN, K-Means 등 거리 기반 알고리즘의 기본",
         ], ACCENT_BLUE, ACCENT_BLUE)
add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.0),
         "Min-Max 정규화", [
             "x' = (x - x_min) / (x_max - x_min)",
             "모든 특성을 [0, 1] 범위로 변환",
         ], ACCENT_GREEN, ACCENT_GREEN)
add_card(s, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.0),
         "Z-score 표준화", [
             "z = (x - μ) / σ",
             "평균=0, 표준편차=1로 변환",
         ], ACCENT_ORANGE, ACCENT_ORANGE)
add_card(s, Inches(6.7), Inches(4.3), Inches(6.1), Inches(2.0),
         "정확도 (Accuracy)", [
             "Acc = (TP + TN) / (TP + TN + FP + FN)",
             "전체 예측 중 올바른 예측의 비율",
         ], ACCENT_PURPLE, ACCENT_PURPLE)

# 슬라이드 52: 복습 질문 (1/2)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 12, "복습 질문 (1/2)", "이론 + 실습 확인")
questions_1 = [
    "Q1. 동적 타이핑이란? 정적 타이핑 대비 ML 개발의 장단점?",
    "Q2. 0.1 + 0.2 == 0.3이 False인 이유? ML에서의 대처법?",
    "Q3. list.append()가 O(1)이고 insert(0)이 O(n)인 이유?",
    "Q4. 100만개 데이터에서 멤버십 검사: list vs set 선택 이유?",
    "Q5. 테스트 데이터에 fit()을 호출하면 안 되는 이유? (데이터 누출)",
]
add_bullet_list(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                questions_1, font_size=17, color=LIGHT_GRAY, spacing=Pt(14))

# 슬라이드 53: 복습 질문 (2/2)
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
slide_header(s, 12, "복습 질문 (2/2)", "코딩 + 심화")
questions_2 = [
    "Q6. [x**2 for x in range(10) if x%3==0]의 결과는?",
    "Q7. NumPy 벡터화가 for 루프보다 수십~수백 배 빠른 이유 3가지?",
    "Q8. GIL이란? ML에서 CPU 병렬 처리 방법?",
    "Q9. 제너레이터의 메모리 절약 원리와 ML 배치 로딩 활용?",
    "Q10. 결측값 제거 → 음수 0으로 클리핑 → 2배 스케일링을",
    "      map/filter와 리스트 컴프리헨션 각각으로 구현하시오",
]
add_bullet_list(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                questions_2, font_size=17, color=LIGHT_GRAY, spacing=Pt(14))

# ============================================================
# 슬라이드 54: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), SECTION_BG, ACCENT_BLUE)
add_accent_line(s, Inches(4), Inches(2.5), Inches(5.333), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.8), Inches(11.333), Inches(1),
         "Thank You", font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.0), Inches(11.333), Inches(0.5),
         "Chapter 2: 머신러닝을 위한 파이썬", font_size=22, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
         "다음 장: NumPy & Pandas 심화", font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 저장
# ============================================================
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "2장_머신러닝을_위한_파이썬_강의PPT_확장.pptx")
prs.save(out_path)
print(f"[완료] {out_path}")
print(f"[슬라이드 수] {len(prs.slides)}")
