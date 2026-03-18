"""2장 머신러닝을 위한 파이썬 - 강의 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    card = add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))


def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)


# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
# 상단 장식 라인
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
# 중앙 정보
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "기계학습 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "2장: 머신러닝을 위한 파이썬", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "Python for Machine Learning", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
         "핵심 키워드: Python 역사 · 동적 타이핑 · 자료구조 · 함수/클래스 · 벡터화 · GIL · 개발환경",
         font_size=14, color=DARK_GRAY)
# 하단 정보
add_shape(s, Inches(0), Inches(6.3), prs.slide_width, Inches(1.2), RGBColor(0x0A, 0x0A, 0x1A))
add_text(s, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
         "2026년 1학기", font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 2: 목차
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (Contents)")

toc = [
    ("01", "왜 Python인가 - ML 표준 언어의 역사", ACCENT_BLUE),
    ("02", "Python vs R / Julia / MATLAB 비교", ACCENT_CYAN),
    ("03", "기초 문법 - 변수, 자료형, 조건문, 반복문", ACCENT_GREEN),
    ("04", "자료구조 심화 - list / tuple / dict / set", ACCENT_PURPLE),
    ("05", "함수와 클래스 - sklearn 패턴", ACCENT_ORANGE),
    ("06", "성능 최적화 - 프로파일링과 벡터화", ACCENT_RED),
    ("07", "GIL과 메모리 관리", ACCENT_BLUE),
    ("08", "개발환경 - conda / venv / Jupyter / IDE", ACCENT_CYAN),
    ("09", "논문 리뷰 (5편) & 실습 소개", ACCENT_GREEN),
    ("10", "핵심 요약 & 복습 질문", ACCENT_PURPLE),
]
for i, (num, title, color) in enumerate(toc):
    y = Inches(2.0) + Inches(0.5) * i
    add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.38), color, radius=True)
    add_text(s, Inches(1.2), y, Inches(0.55), Inches(0.38), num,
             font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), y, Inches(8), Inches(0.38), title,
             font_size=17, color=WHITE)

# ============================================================
# 슬라이드 3: 왜 Python인가 - 타임라인
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "왜 Python인가?", "ML 표준 언어가 된 역사")

milestones = [
    ("1991", "Python 0.9\n공개", "van Rossum\n배터리 포함 철학"),
    ("2001", "SciPy\n프로젝트", "과학 컴퓨팅\n생태계의 씨앗"),
    ("2006", "NumPy 1.0\n(Oliphant)", "고성능 ndarray\nML 핵심 인프라"),
    ("2010", "scikit-learn\n0.1 공개", "fit/predict/\ntransform API"),
    ("2014", "Jupyter\nNotebook", "언어 독립적\n리터릿 컴퓨팅"),
    ("2015", "TensorFlow\n(Google)", "딥러닝 프레임워크\n시대 개막"),
    ("2016", "PyTorch\n(Facebook)", "동적 계산 그래프\n연구 친화적"),
    ("2020", "Raschka\net al. 서베이", "파이썬 ML 생태계\n전체 조망"),
]
# 타임라인 선
add_shape(s, Inches(0.6), Inches(3.6), Inches(12.0), Pt(3), ACCENT_BLUE)
for i, (year, event, desc) in enumerate(milestones):
    x = Inches(0.3) + Inches(1.55) * i
    # 점
    add_shape(s, x + Inches(0.5), Inches(3.45), Inches(0.25), Inches(0.25), ACCENT_CYAN, radius=True)
    add_text(s, x, Inches(2.5), Inches(1.4), Inches(0.4),
             year, font_size=13, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x - Inches(0.1), Inches(2.85), Inches(1.6), Inches(0.7),
             event, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x - Inches(0.1), Inches(4.0), Inches(1.6), Inches(1.2),
             desc, font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 세 가지 이유
reasons = [
    ("접착 언어 (Glue Language)", "C/Fortran 고성능 라이브러리\n(BLAS, LAPACK, cuDNN)를\n파이썬 인터페이스로 감싸서 사용\n→ 생산성 + 성능 동시 확보", ACCENT_BLUE),
    ("풍부한 생태계", "NumPy, Pandas, scikit-learn,\nTensorFlow, PyTorch, matplotlib\n→ ML 파이프라인의 모든 단계를\n하나의 언어로 커버", ACCENT_GREEN),
    ("낮은 진입 장벽", "동적 타이핑, 간결한 문법,\n들여쓰기 기반 블록 구조\n→ 비전공자도 빠르게 학습\n→ ML 교육의 표준 언어", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(reasons):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(5.3), Inches(3.9), Inches(2.0), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(5.4), Inches(3.5), Inches(0.4),
             title, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(5.8), Inches(3.5), Inches(1.4),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 4: Python vs R/Julia/MATLAB 비교 테이블
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "Python vs R / Julia / MATLAB", "언어 비교")

headers = ["기준", "Python", "R", "Julia", "MATLAB"]
rows = [
    ["주력 분야", "범용 + ML/DL", "통계 분석", "고성능 수치 계산", "공학 시뮬레이션"],
    ["성능", "중간 (NumPy 보완)", "느림 (Rcpp 보완)", "빠름 (C 수준)", "중간"],
    ["DL 프레임워크", "TF, PyTorch", "제한적", "Flux.jl", "제한적"],
    ["커뮤니티", "매우 큼", "통계 중심 큼", "성장 중", "학교/기업 중심"],
    ["라이선스", "오픈소스 (무료)", "오픈소스 (무료)", "오픈소스 (무료)", "상용 (유료)"],
    ["산업 채택", "매우 높음", "중간", "낮음", "중간"],
]
header_colors = [DARK_GRAY, ACCENT_BLUE, ACCENT_GREEN, ACCENT_PURPLE, ACCENT_ORANGE]

# 헤더
for j, (h, hc) in enumerate(zip(headers, header_colors)):
    x = Inches(0.5) + Inches(2.5) * j
    add_shape(s, x, Inches(2.3), Inches(2.35), Inches(0.6), hc)
    add_text(s, x, Inches(2.3), Inches(2.35), Inches(0.6),
             h, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 데이터 행
for i, row in enumerate(rows):
    y = Inches(3.0) + Inches(0.6) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.5) + Inches(2.5) * j
        add_shape(s, x, y, Inches(2.35), Inches(0.55), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        if j == 1:
            fc = ACCENT_BLUE
        add_text(s, x, y, Inches(2.35), Inches(0.55),
                 cell, font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 핵심 메시지
add_shape(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.8), Inches(6.75), Inches(11.8), Inches(0.5),
         "핵심: ML 전체 파이프라인(데이터 수집 - 전처리 - 모델링 - 배포)을 하나의 언어로 커버하는 생태계의 폭에서 파이썬을 능가하는 언어는 현재 없다.",
         font_size=13, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 5: 기초 문법 - 변수와 동적 타이핑
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "변수와 동적 타이핑", "Basic Syntax - Variables & Dynamic Typing")

# 동적 타이핑 설명
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "동적 타이핑 (Dynamic Typing)", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.8), [
    "변수 선언 시 자료형을 명시하지 않음",
    "대입되는 값에 따라 자료형 자동 결정",
    "같은 변수에 다른 타입 재할당 가능",
    "x = 42 (int) → x = \"문자열\" (str)",
    "type(x) 함수로 타입 확인",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 자료형 테이블
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "기본 자료형 (Data Types)", font_size=16, color=ACCENT_CYAN, bold=True)

dtypes = [
    ("int", "42", "클래스 레이블, 인덱스"),
    ("float", "3.14", "특성값, 가중치, 손실"),
    ("str", "\"hello\"", "범주형 특성, 텍스트"),
    ("bool", "True/False", "마스크, 조건 필터"),
]
for i, (dtype, ex, ml) in enumerate(dtypes):
    y = Inches(2.85) + Inches(0.38) * i
    add_text(s, Inches(7.1), y, Inches(1.0), Inches(0.35),
             dtype, font_size=13, color=ACCENT_GREEN, bold=True)
    add_text(s, Inches(8.1), y, Inches(1.2), Inches(0.35),
             ex, font_size=13, color=WHITE)
    add_text(s, Inches(9.5), y, Inches(2.8), Inches(0.35),
             ml, font_size=12, color=DARK_GRAY)

# ML 관례
add_shape(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(5.1), Inches(2.5), Inches(0.35),
         "ML 변수 이름 관례", font_size=14, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(3.5), Inches(5.1), Inches(9.0), Inches(0.8),
         "X = 특성 행렬 (대문자)  |  y = 타겟 벡터 (소문자)  |  X_train, X_test = 분할 데이터  |  n_samples, n_features = 차원 정보",
         font_size=13, color=LIGHT_GRAY)

# float 정밀도 주의
add_shape(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(1.0), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(6.25), Inches(2.5), Inches(0.35),
         "float 정밀도 주의!", font_size=14, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(3.5), Inches(6.25), Inches(9.0), Inches(0.8),
         "0.1 + 0.2 == 0.3 → False (IEEE 754)  |  math.isclose()로 비교  |  경사하강법 학습률 소실, log(p) 발산 → np.clip(p, 1e-15, 1)",
         font_size=13, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 6: 기초 문법 - 조건문/반복문/컴프리헨션
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "조건문, 반복문, 리스트 컴프리헨션", "Control Flow & List Comprehension")

# 조건문
add_shape(s, Inches(0.6), Inches(2.2), Inches(3.8), Inches(2.5), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(3.2), Inches(0.4),
         "조건문 (if-elif-else)", font_size=15, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(3.2), Inches(2.0), [
    "if score >= 90:",
    "    grade = 'A'",
    "elif score >= 80:",
    "    grade = 'B'",
    "else:",
    "    grade = 'C'",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(2))

# 반복문
add_shape(s, Inches(4.75), Inches(2.2), Inches(3.8), Inches(2.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(5.05), Inches(2.3), Inches(3.2), Inches(0.4),
         "반복문 (for / while)", font_size=15, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(5.05), Inches(2.7), Inches(3.2), Inches(2.0), [
    "for i in range(5): ...",
    "while a > 0: a -= 1",
    "",
    "enumerate: 인덱스+값 동시",
    "for idx, val in enumerate(lst):",
    "    print(f\"{idx}: {val}\")",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(2))

# 리스트 컴프리헨션
add_shape(s, Inches(8.9), Inches(2.2), Inches(3.8), Inches(2.5), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(9.2), Inches(2.3), Inches(3.2), Inches(0.4),
         "리스트 컴프리헨션", font_size=15, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(9.2), Inches(2.7), Inches(3.2), Inches(2.0), [
    "[x**2 for x in range(1,6)]",
    "→ [1, 4, 9, 16, 25]",
    "",
    "[x for x in data if x%2==0]",
    "→ 조건 필터링",
    "",
    "[1 if x>0 else 0 for x in data]",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(2))

# f-string
add_shape(s, Inches(0.6), Inches(5.0), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(5.1), Inches(5.2), Inches(0.4),
         "f-string 포매팅", font_size=15, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(5.5), Inches(5.2), Inches(1.6), [
    "f\"{name}의 정확도: {accuracy:.2%}\"",
    "→ 홍길동의 정확도: 95.43%",
    "",
    "ML 출력 패턴:",
    "f\"Epoch {epoch:3d} | Loss: {loss:.4f} | Acc: {acc:.4f}\"",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# 비교/논리 연산자
add_shape(s, Inches(6.8), Inches(5.0), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(7.1), Inches(5.1), Inches(5.2), Inches(0.4),
         "비교 & 논리 연산자", font_size=15, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(5.5), Inches(5.2), Inches(1.6), [
    "비교: ==, !=, >, <, >=, <=",
    "논리: and, or, not",
    "",
    "산술: +, -, *, /, //(몫), %(나머지), **(거듭제곱)",
    "주의: / 연산은 항상 float 반환",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 7: 자료구조 심화 - 4가지 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "자료구조 심화: 4가지 비교", "list / tuple / dict / set")

# 비교 테이블 헤더
th = ["자료구조", "기호", "변경 가능", "순서 보장", "중복 허용", "해시 기반"]
for j, h in enumerate(th):
    x = Inches(0.5) + Inches(2.05) * j
    add_shape(s, x, Inches(2.2), Inches(1.95), Inches(0.5), ACCENT_BLUE)
    add_text(s, x, Inches(2.2), Inches(1.95), Inches(0.5),
             h, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

ds_rows = [
    ["리스트 (list)", "[ ]", "O (mutable)", "O", "O", "X"],
    ["튜플 (tuple)", "( )", "X (immutable)", "O", "O", "X"],
    ["딕셔너리 (dict)", "{ }", "O (mutable)", "O (3.7+)", "키 X", "O"],
    ["세트 (set)", "{ }", "O (mutable)", "X", "X", "O"],
]
for i, row in enumerate(ds_rows):
    y = Inches(2.8) + Inches(0.5) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.5) + Inches(2.05) * j
        add_shape(s, x, y, Inches(1.95), Inches(0.45), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y, Inches(1.95), Inches(0.45),
                 cell, font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 시간 복잡도 테이블
add_text(s, Inches(0.6), Inches(5.0), Inches(5), Inches(0.4),
         "시간 복잡도 (Big-O) 비교", font_size=16, color=WHITE, bold=True)

tc_h = ["연산", "list", "tuple", "dict", "set"]
tc_rows = [
    ["인덱스 접근 [i]", "O(1)", "O(1)", "O(1)", "N/A"],
    ["탐색 (in)", "O(n)", "O(n)", "O(1)", "O(1)"],
    ["끝 삽입", "O(1)*", "N/A", "O(1)", "O(1)"],
    ["앞 삽입", "O(n)", "N/A", "N/A", "N/A"],
]
tc_header_colors = [DARK_GRAY, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE]
for j, (h, hc) in enumerate(zip(tc_h, tc_header_colors)):
    x = Inches(0.5) + Inches(2.5) * j
    add_shape(s, x, Inches(5.4), Inches(2.35), Inches(0.45), hc)
    add_text(s, x, Inches(5.4), Inches(2.35), Inches(0.45),
             h, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(tc_rows):
    y = Inches(5.9) + Inches(0.4) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.5) + Inches(2.5) * j
        add_shape(s, x, y, Inches(2.35), Inches(0.38), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        # Highlight O(1) in green, O(n) in red for search
        if i == 1 and cell == "O(1)":
            fc = ACCENT_GREEN
        elif i == 1 and cell == "O(n)":
            fc = ACCENT_RED
        add_text(s, x, y, Inches(2.35), Inches(0.38),
                 cell, font_size=11, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 8: 자료구조 - ML 활용 가이드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "자료구조: ML에서의 활용 가이드", "어떤 상황에 어떤 자료구조?")

guides = [
    ("list", "특성값 저장 (순서 중요)\n인덱스 접근 O(1)\nappend()로 데이터 수집\n슬라이싱으로 데이터 분할", ACCENT_BLUE,
     "data = [1, 2, 3, 4, 5]\ndata[2:5]  # [3, 4, 5]\ndata.append(6)"),
    ("tuple", "하이퍼파라미터 조합 (불변)\n모델 설정, shape 정보\n딕셔너리 키로 사용 가능\n실수로 변경 방지", ACCENT_CYAN,
     "config = ('relu', 0.001, 32)\nactivation, lr, batch = config"),
    ("dict", "하이퍼파라미터 관리\n범주 레이블 매핑\n이름으로 접근, 유연한 구조\n**kwargs 언패킹", ACCENT_GREEN,
     "params = {'lr': 0.01,\n  'epochs': 100}\nmodel = XGB(**params)"),
    ("set", "고유 카테고리 추출\n중복 자동 제거\n빠른 멤버십 검사 O(1)\n특성 비교 (집합 연산)", ACCENT_ORANGE,
     "train_feat - test_feat\n→ 누락된 특성 감지"),
]
for i, (title, desc, color, code) in enumerate(guides):
    x = Inches(0.3) + Inches(3.25) * i
    add_shape(s, x, Inches(2.2), Inches(3.0), Inches(2.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(3.0), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(3.0), Inches(0.5),
             title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(2.9), Inches(2.7), Inches(1.8),
             desc, font_size=12, color=LIGHT_GRAY)

    # 코드 예시 박스
    add_shape(s, x, Inches(5.2), Inches(3.0), Inches(2.0), RGBColor(0x15, 0x15, 0x30), color, radius=True)
    add_text(s, x + Inches(0.15), Inches(5.3), Inches(2.7), Inches(1.8),
             code, font_size=10, color=ACCENT_GREEN, font_name='Consolas')

# ============================================================
# 슬라이드 9: 함수 기본 - def, *args, **kwargs, lambda
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "함수: def, *args, **kwargs, lambda", "Functions & Lambda")

# 기본 함수
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "기본 함수 정의", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(1.8), [
    "def train(data, epochs=10, lr=0.01):",
    "    print(f\"학습: {epochs} 에폭\")",
    "",
    "기본값 매개변수로 유연한 호출",
    "여러 값 반환: return total, avg, min_v",
    "→ 튜플로 반환, 언패킹으로 받기",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# *args, **kwargs
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "*args와 **kwargs", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(1.8), [
    "*args: 가변 위치 인자 (튜플로 수집)",
    "  def add_all(*args): return sum(args)",
    "",
    "**kwargs: 가변 키워드 인자 (딕셔너리)",
    "  def create_model(**kwargs):",
    "    for k, v in kwargs.items(): ...",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# lambda
add_shape(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.5), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.4),
         "lambda 함수 - 익명 함수", font_size=16, color=ACCENT_PURPLE, bold=True)

lambda_items = [
    ("lambda 기본", "square = lambda x: x ** 2", ACCENT_CYAN),
    ("sorted + lambda", "sorted(features, key=lambda x: x[1], reverse=True)", ACCENT_GREEN),
    ("filter + lambda", "list(filter(lambda x: x > 0, data))", ACCENT_ORANGE),
    ("map + lambda", "list(map(lambda x: x ** 2, data))", ACCENT_RED),
]
for i, (name, code, color) in enumerate(lambda_items):
    y = Inches(5.4) + Inches(0.45) * i
    add_text(s, Inches(0.9), y, Inches(2.5), Inches(0.4),
             name, font_size=13, color=color, bold=True)
    add_text(s, Inches(3.5), y, Inches(9.0), Inches(0.4),
             code, font_size=13, color=LIGHT_GRAY, font_name='Consolas')

# ============================================================
# 슬라이드 10: sklearn의 fit/predict/transform 패턴
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "sklearn Estimator API 패턴", "fit / predict / transform")

# 패턴 설명
add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.3), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "scikit-learn의 일관된 인터페이스: 모든 모델/전처리기가 동일한 메서드 패턴을 따른다", font_size=15, color=WHITE)
add_text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.5),
         "\"There should be one obvious way to do it.\" — The Zen of Python → sklearn API 설계 철학", font_size=13, color=DARK_GRAY)

# 4개 메서드 카드
methods = [
    ("fit(X)", "학습 데이터에서\n파라미터 학습\n(평균, 분산, 가중치)", "학습 시", ACCENT_BLUE),
    ("predict(X)", "학습된 모델로\n예측 수행", "추론 시", ACCENT_GREEN),
    ("transform(X)", "학습된 파라미터로\n데이터 변환\n(전처리기)", "학습/추론 시", ACCENT_ORANGE),
    ("fit_transform(X)", "fit + transform을\n한 번에 수행\n(편의 메서드)", "학습 시", ACCENT_PURPLE),
]
for i, (name, desc, when, color) in enumerate(methods):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(3.8), Inches(2.9), Inches(2.5), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(3.8), Inches(2.9), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(3.85), Inches(2.9), Inches(0.5),
             name, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(4.5), Inches(2.5), Inches(1.2),
             desc, font_size=13, color=LIGHT_GRAY)
    add_text(s, x + Inches(0.2), Inches(5.7), Inches(2.5), Inches(0.4),
             when, font_size=12, color=color, bold=True)

# 중요 원칙 경고
add_shape(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.7), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.55),
         "중요 원칙: fit()은 반드시 학습 데이터에서만 호출! 테스트 데이터에 fit()을 호출하면 데이터 누출(data leakage)이 발생하여 모델 평가 왜곡",
         font_size=14, color=ACCENT_RED, bold=True)

# ============================================================
# 슬라이드 11: 클래스 상속과 다형성
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "클래스 상속과 다형성", "BaseModel 패턴")

# BaseModel
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(3.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "추상 기반 클래스: BaseModel", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(2.2), [
    "class BaseModel:",
    "    def fit(self, X, y):",
    "        raise NotImplementedError",
    "    def predict(self, X):",
    "        raise NotImplementedError",
    "    def score(self, X, y):",
    "        predictions = self.predict(X)",
    "        correct = sum(p==t for p,t in zip(...))",
    "        return correct / len(y)",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(2))

# MajorityClassifier
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(3.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "구체적 구현: MajorityClassifier", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(2.2), [
    "class MajorityClassifier(BaseModel):",
    "    def fit(self, X, y):",
    "        self.majority_ = Counter(y)",
    "            .most_common(1)[0][0]",
    "        return self",
    "    def predict(self, X):",
    "        return [self.majority_] * len(X)",
    "",
    "→ 다형성: 동일 인터페이스로 다른 모델",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(2))

# 핵심 개념 카드
concepts = [
    ("상속 (Inheritance)", "BaseModel → MajorityClassifier\n부모 클래스의 메서드/속성 물려받기\n공통 인터페이스 강제", ACCENT_CYAN),
    ("다형성 (Polymorphism)", "동일한 predict() 호출, 다른 동작\nfor model in models:\n    model.fit(X, y).predict(X_test)", ACCENT_PURPLE),
    ("캡슐화 (Encapsulation)", "학습된 파라미터: self.majority_\n밑줄 접미사(_)는 학습 후 생성\nsklearn 컨벤션", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(concepts):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(5.5), Inches(3.9), Inches(1.8), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(5.6), Inches(3.5), Inches(0.4),
             title, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(6.0), Inches(3.5), Inches(1.2),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 12: 성능 최적화 - 프로파일링
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "성능 최적화: 프로파일링", "\"측정하지 않으면 최적화하지 마라\"")

# 3가지 프로파일링 방법
profile_methods = [
    ("1. time.perf_counter()", "간단한 구간 측정", [
        "start = time.perf_counter()",
        "result = sum(x**2 for x in range(10**6))",
        "elapsed = time.perf_counter() - start",
        "print(f\"실행 시간: {elapsed:.4f}초\")",
    ], ACCENT_BLUE),
    ("2. cProfile", "함수별 상세 프로파일링", [
        "import cProfile",
        "def my_function():",
        "    data = [x**2 for x in range(100000)]",
        "    return sorted(data, reverse=True)",
        "cProfile.run('my_function()')",
    ], ACCENT_GREEN),
    ("3. Jupyter %timeit", "자동 반복 실행 후 평균", [
        "%timeit sum(range(1000))",
        "→ 자동으로 여러 번 실행 후 평균",
        "",
        "%%timeit",
        "→ 셀 전체 시간 측정",
    ], ACCENT_ORANGE),
]
for i, (title, desc, code, color) in enumerate(profile_methods):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(4.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(3.9), Inches(0.5),
             title, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.9), Inches(3.5), Inches(0.3),
             desc, font_size=13, color=color)
    add_bullet_list(s, x + Inches(0.2), Inches(3.3), Inches(3.5), Inches(3.5),
                    code, font_size=11, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 13: 벡터화의 위력
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "벡터화의 위력: Python Loop vs NumPy", "Vectorization")

# 속도 비교 시각화 (텍스트 기반)
add_shape(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.5), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
         "n = 1,000,000 데이터에서 제곱 연산 (x²) 성능 비교", font_size=16, color=WHITE, bold=True)

speed_items = [
    ("for 루프", "~150ms", Inches(9.0), ACCENT_RED),
    ("리스트 컴프리헨션", "~100ms", Inches(6.0), ACCENT_ORANGE),
    ("map() + lambda", "~110ms", Inches(6.6), ACCENT_PURPLE),
    ("NumPy 벡터화", "~2ms", Inches(0.15), ACCENT_GREEN),
]
for i, (name, time_str, bar_w, color) in enumerate(speed_items):
    y = Inches(2.85) + Inches(0.2) * i
    add_text(s, Inches(0.9), y, Inches(2.5), Inches(0.2),
             name, font_size=11, color=LIGHT_GRAY)
    add_shape(s, Inches(3.5), y + Inches(0.02), bar_w, Inches(0.15), color)
    add_text(s, Inches(3.5) + bar_w + Inches(0.1), y, Inches(1.0), Inches(0.2),
             time_str, font_size=10, color=color, bold=True)

# 벡터화가 빠른 3가지 이유
add_text(s, Inches(0.6), Inches(3.9), Inches(5), Inches(0.4),
         "벡터화가 빠른 3가지 이유", font_size=18, color=WHITE, bold=True)

reasons = [
    ("1", "C 수준 내부 루프", "파이썬 인터프리터의 오버헤드를\n완전히 제거. CPython의 바이트코드\n해석 과정을 건너뜀", ACCENT_BLUE),
    ("2", "연속 메모리 접근", "ndarray는 연속 메모리 블록에 저장.\nCPU 캐시 효율 극대화.\n리스트처럼 포인터 추적 불필요", ACCENT_GREEN),
    ("3", "SIMD 명령어", "하드웨어 수준 병렬 처리.\n하나의 CPU 명령어로 여러 데이터\n동시 처리 (Single Instruction, Multiple Data)", ACCENT_ORANGE),
]
for i, (num, title, desc, color) in enumerate(reasons):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(4.4), Inches(3.9), Inches(2.8), CARD_BG, color, radius=True)
    add_shape(s, x + Inches(0.15), Inches(4.5), Inches(0.4), Inches(0.4), color, radius=True)
    add_text(s, x + Inches(0.15), Inches(4.5), Inches(0.4), Inches(0.4),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.65), Inches(4.55), Inches(3.0), Inches(0.35),
             title, font_size=15, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(5.1), Inches(3.5), Inches(1.8),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 14: GIL과 메모리 관리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "GIL과 메모리 관리", "Global Interpreter Lock & Memory")

# GIL 설명
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "GIL (Global Interpreter Lock)", font_size=16, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.8), Inches(5.2), Inches(1.8), [
    "CPython의 전역 잠금 메커니즘",
    "한 번에 하나의 스레드만 바이트코드 실행",
    "",
    "CPU 바운드: 멀티스레딩 성능 향상 불가",
    "I/O 바운드: 멀티스레딩 효과 있음",
    "  (GIL이 I/O 대기 중 해제)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# ML 대처 방안
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "ML에서의 GIL 대처 방안", font_size=16, color=ACCENT_GREEN, bold=True)

strategies = [
    ("멀티프로세싱", "multiprocessing, joblib", "CPU 병렬 학습"),
    ("NumPy 내부 병렬화", "BLAS/MKL", "행렬 연산 자동 병렬화"),
    ("외부 라이브러리", "TF, PyTorch", "GPU 연산 (GIL 무관)"),
    ("n_jobs 파라미터", "sklearn", "RandomForest(n_jobs=-1)"),
]
for i, (strategy, tool, usage) in enumerate(strategies):
    y = Inches(2.85) + Inches(0.38) * i
    add_text(s, Inches(7.1), y, Inches(1.8), Inches(0.35),
             strategy, font_size=12, color=ACCENT_CYAN, bold=True)
    add_text(s, Inches(8.9), y, Inches(1.5), Inches(0.35),
             tool, font_size=11, color=WHITE)
    add_text(s, Inches(10.4), y, Inches(2.0), Inches(0.35),
             usage, font_size=11, color=DARK_GRAY)

# 메모리 관리
add_shape(s, Inches(0.6), Inches(5.0), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(0.9), Inches(5.1), Inches(5.2), Inches(0.4),
         "메모리 사용량 비교", font_size=16, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(0.9), Inches(5.5), Inches(5.2), Inches(1.6), [
    "list(range(10000)):   ~87 KB",
    "tuple(range(10000)):  ~80 KB (더 작음)",
    "dict (10000 항목):    ~300+ KB",
    "generator:            ~200 bytes (!)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 제너레이터 설명
add_shape(s, Inches(6.8), Inches(5.0), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(7.1), Inches(5.1), Inches(5.2), Inches(0.4),
         "제너레이터: 메모리 절약 비법", font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_list(s, Inches(7.1), Inches(5.5), Inches(5.2), Inches(1.6), [
    "gen = (x**2 for x in range(10**6))",
    "→ 리스트 대비 수만 배 메모리 절약!",
    "",
    "ML 활용:",
    "  배치 데이터 로딩, np.memmap,",
    "  pandas.read_csv(chunksize=...)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 15: 개발환경 - conda vs venv, Jupyter vs IDE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "개발환경: 가상환경 & 편집기", "conda / venv / Jupyter / IDE")

# conda vs venv 테이블
add_text(s, Inches(0.6), Inches(2.0), Inches(5), Inches(0.4),
         "가상환경: conda vs venv", font_size=16, color=WHITE, bold=True)

env_h = ["기준", "venv + pip", "conda"]
env_rows = [
    ["패키지 소스", "PyPI", "Anaconda + PyPI"],
    ["바이너리 관리", "제한적", "우수 (C/Fortran 자동 해결)"],
    ["GPU 설정", "수동", "conda install pytorch 간편"],
    ["환경 크기", "작음", "상대적으로 큼"],
    ["추천", "가벼운 프로젝트", "ML/DL 프로젝트"],
]
env_hc = [DARK_GRAY, ACCENT_BLUE, ACCENT_GREEN]
for j, (h, hc) in enumerate(zip(env_h, env_hc)):
    x = Inches(0.6) + Inches(2.0) * j
    add_shape(s, x, Inches(2.4), Inches(1.9), Inches(0.45), hc)
    add_text(s, x, Inches(2.4), Inches(1.9), Inches(0.45),
             h, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(env_rows):
    y = Inches(2.9) + Inches(0.4) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.6) + Inches(2.0) * j
        add_shape(s, x, y, Inches(1.9), Inches(0.38), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y, Inches(1.9), Inches(0.38),
                 cell, font_size=11, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# Jupyter vs IDE 테이블
add_text(s, Inches(6.8), Inches(2.0), Inches(5), Inches(0.4),
         "Jupyter Notebook vs IDE", font_size=16, color=WHITE, bold=True)

ide_h = ["기준", "Jupyter", "IDE (VS Code)"]
ide_rows = [
    ["강점", "탐색적 분석, 시각화", "디버깅, 리팩토링"],
    ["적합한 작업", "EDA, 프로토타이핑", "프로덕션, 패키지 개발"],
    ["매직 명령어", "%timeit, %matplotlib", "N/A"],
    ["버전 관리", ".ipynb (diff 불편)", ".py (Git 친화적)"],
]
ide_hc = [DARK_GRAY, ACCENT_ORANGE, ACCENT_PURPLE]
for j, (h, hc) in enumerate(zip(ide_h, ide_hc)):
    x = Inches(6.8) + Inches(2.0) * j
    add_shape(s, x, Inches(2.4), Inches(1.9), Inches(0.45), hc)
    add_text(s, x, Inches(2.4), Inches(1.9), Inches(0.45),
             h, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(ide_rows):
    y = Inches(2.9) + Inches(0.4) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(6.8) + Inches(2.0) * j
        add_shape(s, x, y, Inches(1.9), Inches(0.38), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y, Inches(1.9), Inches(0.38),
                 cell, font_size=11, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 실무 워크플로우
add_shape(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.8), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(5.25), Inches(11.5), Inches(0.35),
         "실무 워크플로우", font_size=14, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.35),
         "Jupyter에서 프로토타이핑  →  .py 모듈로 리팩토링  →  IDE에서 테스트/배포", font_size=14, color=WHITE)

# 디버깅 기법 카드
debug_methods = [
    ("print 디버깅", "f\"[DEBUG] type: {type(data)}\"", ACCENT_BLUE),
    ("assert 문", "assert len(X) == len(y)", ACCENT_GREEN),
    ("pdb 디버거", "import pdb; pdb.set_trace()", ACCENT_ORANGE),
    ("try-except", "예외 처리로 안전한 실행", ACCENT_RED),
]
for i, (name, code, color) in enumerate(debug_methods):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(6.2), Inches(2.9), Inches(1.1), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(6.25), Inches(2.6), Inches(0.35),
             name, font_size=12, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(6.6), Inches(2.6), Inches(0.5),
             code, font_size=10, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 16: 논문 리뷰 5편
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "핵심 논문 리뷰 (5편)", "Key Paper Reviews")

papers = [
    ("Van Rossum\n(2007)", "Python Programming\nLanguage",
     "파이썬 설계 철학\n동적 타이핑, GIL\nC 확장 인터페이스", ACCENT_BLUE),
    ("Millman &\nAivazis (2011)", "Python for Scientists\nand Engineers",
     "SciPy 스택 정립\n과학 컴퓨팅 생태계\n소프트웨어 공학 관행", ACCENT_CYAN),
    ("Oliphant\n(2007)", "Python for Scientific\nComputing",
     "NumPy ndarray 구조\nstride 기반 레이아웃\nBLAS/LAPACK 활용", ACCENT_GREEN),
    ("Perez &\nGranger (2007)", "IPython: Interactive\nScientific Computing",
     "Jupyter 전신\n매직 명령어, 커널 분리\n리터릿 컴퓨팅 비전", ACCENT_ORANGE),
    ("Raschka et al.\n(2020)", "Machine Learning\nin Python",
     "ML 생태계 서베이\nsklearn vs TF vs PyTorch\nAutoML, MLOps 트렌드", ACCENT_PURPLE),
]
for i, (author, title, contrib, color) in enumerate(papers):
    x = Inches(0.2) + Inches(2.6) * i
    add_shape(s, x, Inches(2.2), Inches(2.4), Inches(4.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(2.4), Inches(0.7), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(2.4), Inches(0.65),
             author, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.1), Inches(2.1), Inches(1.2),
             title, font_size=12, color=LIGHT_GRAY)
    add_accent_line(s, x + Inches(0.2), Inches(4.3), Inches(2.0), color)
    add_text(s, x + Inches(0.15), Inches(4.5), Inches(2.1), Inches(2.0),
             contrib, font_size=11, color=DARK_GRAY)

# ============================================================
# 슬라이드 17: 논문 리뷰 - Raschka et al. 상세
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "Raschka et al. (2020) 상세", "Machine Learning in Python - 필독 논문")

# 역사적 발전
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "역사적 발전 계보", font_size=15, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.9), Inches(2.75), Inches(5.2), Inches(1.2),
         "NumPy (2006)  →  scikit-learn (2010)  →  TensorFlow (2015)  →  PyTorch (2016)\n\n프레임워크 비교: sklearn(일관 API) vs TF(분산 학습) vs PyTorch(동적 그래프)",
         font_size=13, color=LIGHT_GRAY)

# 데이터 처리 도구 진화
add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "데이터 처리 도구 진화", font_size=15, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(7.1), Inches(2.75), Inches(5.2), Inches(1.2),
         "Pandas  →  Dask (분산 처리)  →  RAPIDS (GPU 가속)\n\n대용량 데이터 처리의 병목을 해결하는 방향으로 진화",
         font_size=13, color=LIGHT_GRAY)

# 최신 트렌드
trends = [
    ("AutoML", "Auto-sklearn, TPOT\n자동 알고리즘 선택\n하이퍼파라미터 튜닝", ACCENT_CYAN),
    ("Interpretable ML", "SHAP, LIME\n블랙박스 모델 해석\n공정성, 투명성", ACCENT_ORANGE),
    ("MLOps", "MLflow, Kubeflow\n모델 배포, 모니터링\n재현 가능한 파이프라인", ACCENT_PURPLE),
    ("Federated Learning", "분산 학습\n데이터 프라이버시 보호\n엣지 디바이스 학습", ACCENT_RED),
]
for i, (name, desc, color) in enumerate(trends):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(4.5), Inches(2.9), Inches(2.7), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(4.5), Inches(2.9), Inches(0.5), color, radius=True)
    add_text(s, x, Inches(4.55), Inches(2.9), Inches(0.45),
             name, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(5.15), Inches(2.6), Inches(1.8),
             desc, font_size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 18: 실습 소개
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "실습 소개 (3개 구현 소스)")

labs = [
    ("실습 1", "01_python_profiling.py", "파이썬 성능 프로파일링", [
        "for 루프 vs 리스트 컴프리헨션 vs map vs NumPy",
        "제곱, 합산, 필터링 연산 비교",
        "데이터 크기별 스케일링 테스트",
        "matplotlib 시각화로 결과 확인",
    ], ACCENT_BLUE),
    ("실습 2", "02_data_structure_benchmark.py", "자료구조 벤치마크", [
        "list/tuple/dict/set 탐색 성능 비교",
        "삽입(append vs insert) 성능 차이",
        "삭제 연산 성능 비교",
        "메모리 사용량 측정 및 시각화",
    ], ACCENT_GREEN),
    ("실습 3", "03_functional_ml.py", "함수형 프로그래밍 ML", [
        "map/filter/reduce를 활용한 데이터 전처리",
        "sklearn 패턴 직접 구현 (SimpleScaler)",
        "클래스 상속으로 모델 구현",
        "파이프라인 패턴 실습",
    ], ACCENT_ORANGE),
]
for i, (title, file, desc, items, color) in enumerate(labs):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(4.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(3.9), Inches(0.5),
             title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.9), Inches(3.5), Inches(0.3),
             file, font_size=11, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(0.3),
             desc, font_size=13, color=WHITE)
    add_bullet_list(s, x + Inches(0.2), Inches(3.6), Inches(3.5), Inches(3.0),
                    items, font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# 슬라이드 19: 핵심 요약 (6개 카드)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "핵심 요약 (Key Takeaways)")

summaries = [
    ("왜 Python인가", "접착 언어 + 풍부한 생태계 + 낮은 진입 장벽\nR/Julia/MATLAB 대비 ML 파이프라인 전체 커버", ACCENT_BLUE),
    ("기초 문법", "동적 타이핑, 4가지 기본 자료형\nfloat 정밀도 주의, f-string 포매팅", ACCENT_CYAN),
    ("자료구조", "list(순서) / tuple(불변) / dict(해시) / set(고유)\ndict/set은 O(1) 탐색, list는 O(n)", ACCENT_GREEN),
    ("함수와 클래스", "sklearn fit/predict/transform 패턴\n상속과 다형성으로 일관된 API 설계", ACCENT_ORANGE),
    ("성능 최적화", "프로파일링 → 벡터화(NumPy)\nfor 루프 대비 수십~수백 배 성능 향상", ACCENT_RED),
    ("개발환경", "conda(ML/DL) vs venv(가벼운 프로젝트)\nJupyter(탐색) → IDE(프로덕션)", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(summaries):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.2) + Inches(2.5) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.1), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
             title, font_size=17, color=color, bold=True)
    add_accent_line(s, x + Inches(0.2), y + Inches(0.6), Inches(2.0), color)
    add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.2),
             desc, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 20: 수식 요약
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "핵심 수식 정리")

formulas = [
    ("유클리드 거리", "d(x, y) = sqrt( sum( (xi - yi)^2 ) )", "벡터 간 거리 측정"),
    ("Min-Max 정규화", "x' = (x - x_min) / (x_max - x_min)", "값을 [0, 1] 범위로 변환"),
    ("Z-score 표준화", "z = (x - mu) / sigma", "평균 0, 분산 1로 표준화"),
    ("정확도", "Acc = (TP + TN) / (TP + TN + FP + FN)", "전체 중 맞춘 비율"),
    ("시간 복잡도", "list 탐색: O(n)  vs  dict/set 탐색: O(1)", "해시 테이블의 이점"),
    ("벡터화 속도비", "NumPy: ~2ms  vs  for 루프: ~150ms (75x)", "C 내부 루프 + SIMD"),
]
for i, (name, formula, meaning) in enumerate(formulas):
    y = Inches(2.2) + Inches(0.85) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    add_shape(s, Inches(0.6), y, Inches(12.1), Inches(0.75), bg, radius=True)
    add_text(s, Inches(0.9), y + Inches(0.05), Inches(2.5), Inches(0.65),
             name, font_size=14, color=ACCENT_CYAN, bold=True)
    add_text(s, Inches(3.5), y + Inches(0.05), Inches(4.5), Inches(0.65),
             formula, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.2), y + Inches(0.05), Inches(4.3), Inches(0.65),
             meaning, font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 21: 복습 질문 10개
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "복습 질문 (Review Questions)")

questions = [
    "Q1. 파이썬의 동적 타이핑이란 무엇이며, ML 개발에서 어떤 장단점이 있는가?",
    "Q2. 0.1 + 0.2 == 0.3이 False를 반환하는 이유를 IEEE 754와 ML 관점에서 설명하시오.",
    "Q3. 리스트의 append()가 O(1)이고 insert(0)가 O(n)인 이유를 내부 구조로 설명하시오.",
    "Q4. 100만 개 정수에서 멤버십 검사 시 list와 set 중 어떤 것을 사용해야 하는가?",
    "Q5. sklearn의 fit()/transform()에서 테스트 데이터에 fit()을 호출하면 안 되는 이유는?",
    "Q6. [x**2 for x in range(10) if x%3==0]의 실행 결과를 예측하고 동작 원리를 설명하시오.",
    "Q7. NumPy 벡터화가 파이썬 for 루프보다 수십~수백 배 빠른 이유 3가지를 제시하시오.",
    "Q8. GIL이란 무엇이며, ML 학습에서 CPU 병렬 처리가 필요할 때 어떻게 우회하는가?",
    "Q9. 제너레이터가 리스트 대비 메모리를 절약하는 원리와 ML 활용 사례를 기술하시오.",
    "Q10. map/filter와 리스트 컴프리헨션으로 전처리 파이프라인을 각각 구현하시오.",
]
for i, q in enumerate(questions):
    y = Inches(2.0) + Inches(0.52) * i
    color = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED,
             ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE][i]
    add_shape(s, Inches(0.8), y + Inches(0.05), Inches(0.08), Inches(0.3), color)
    add_text(s, Inches(1.1), y, Inches(11.5), Inches(0.5),
             q, font_size=13, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 22: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.5), prs.slide_width, Inches(1.0),
         "Thank You", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), prs.slide_width, Inches(0.6),
         "2장: 머신러닝을 위한 파이썬", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(4.2), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(4.5), prs.slide_width, Inches(0.5),
         "다음 장: 3장 - 판다스와 넘파이", font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ── 저장 ──
output_path = os.path.join(os.path.dirname(__file__), "2장_파이썬_기초_강의PPT.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
