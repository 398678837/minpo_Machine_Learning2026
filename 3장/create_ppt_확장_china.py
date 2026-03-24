"""第3章 Pandas与NumPy - 扩展讲义PPT生成脚本（详细版）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色调色板 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    """以卡片形式绘制表格的辅助函数"""
    total_w = sum(col_widths)
    # 表头
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    # 行
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    """章节分隔幻灯片"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), accent)
    add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), accent)
    add_text(s, Inches(0), Inches(2.0), prs.slide_width, Inches(0.5),
             f"SECTION {section_num}", font_size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, Inches(5.5), Inches(2.7), Inches(2.3), accent)
    add_text(s, Inches(0), Inches(3.0), prs.slide_width, Inches(1.0),
             title, font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.5),
             subtitle, font_size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    return s


# ============================================================
# 幻灯片 1: 封面
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "机器学习 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "第3章: Pandas与NumPy", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "Pandas & NumPy for Machine Learning", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.3), Inches(11), Inches(0.8),
         "[ 扩展详细版 ]", font_size=20, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
         "核心关键词: NumPy ndarray · Pandas DataFrame · 广播 · 向量化 · Tidy Data · 线性代数 · 数据预处理",
         font_size=14, color=DARK_GRAY)
add_text(s, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
         "研究生博士课程级别 | 机器学习 第3章", font_size=14, color=DARK_GRAY)

# ============================================================
# 幻灯片 2: 目录 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "目录 (Table of Contents) - 1/2")

toc_left = [
    "Section 1. 引言: 没有数据就没有ML",
    "    1.1 Python科学生态系统的层级结构",
    "    1.2 ML管道中的数据流",
    "    1.3 为什么选择Python?",
    "",
    "Section 2. NumPy深入: ndarray的世界",
    "    2.1 ndarray内部结构 (4个核心要素)",
    "    2.2 内存布局与Stride",
    "    2.3 C-order vs F-order / dtype系统",
    "    2.4 数组创建函数总结",
    "    2.5 数组索引与切片",
    "    2.6 广播的3条规则",
    "    2.7 广播代码示例",
    "    2.8 通用函数 (ufunc)",
    "    2.9 线性代数基础 / SVD / 特征值分解",
    "    2.10 正规方程与View vs Copy",
]
add_bullet_list(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5.0),
                toc_left, font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# ============================================================
# 幻灯片 3: 目录 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "目录 (Table of Contents) - 2/2")

toc_right = [
    "Section 3. Pandas深入: DataFrame的世界",
    "    3.1 DataFrame结构 / 索引 (loc vs iloc)",
    "    3.2 方法链 / GroupBy / Merge·Join",
    "    3.3 时间序列窗口函数 / Apply·Map",
    "",
    "Section 4. 数据预处理",
    "    4.1 缺失值处理5种方法 / 异常值检测3种方法",
    "    4.2 特征缩放3种方法 / 编码",
    "",
    "Section 5. Tidy Data概念 (Wickham 2014)",
    "Section 6. 向量化 vs 循环: 性能的科学",
    "Section 7. 论文综述整合",
    "Section 8. 实现代码详细解说 (3个实验)",
    "Section 9. 核心总结及复习题",
]
add_bullet_list(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5.0),
                toc_right, font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# ============================================================
# 幻灯片 4: 学习目标
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "学习目标 (Learning Objectives)")

objectives = [
    "1. 能够说明NumPy ndarray的内部内存结构(stride, dtype, C/F-order)，并应用于性能优化。",
    "2. 能够准确应用广播规则，对不同形状的数组进行运算。",
    "3. 能够用NumPy实现LU/QR/SVD/特征值分解，并说明在ML算法(PCA、线性回归)中的应用。",
    "4. 能够理解Pandas DataFrame的内部结构，用method chaining和pipe()构建数据管道。",
    "5. 能够比较分析缺失值处理5种方法、异常值检测3种方法、特征缩放3种方法。",
    "6. 能够应用Wickham的Tidy Data原则将messy data转换为tidy data。",
    "7. 能够从硬件层面(SIMD、缓存)解释向量化运算与Python循环的性能差异。",
]
add_bullet_list(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5.0),
                objectives, font_size=15, color=LIGHT_GRAY, spacing=Pt(8))

# ============================================================
# Section 1: 引言
# ============================================================
section_divider("引言: 没有数据就没有ML", "Python生态系统与ML管道中NumPy·Pandas的角色", 1, ACCENT_BLUE)

# ============================================================
# 幻灯片 5: Python科学生态系统的层级结构
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 1, "Python科学生态系统的层级结构", "基于Harris et al. (2020) Nature论文")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.2),
         "Layer 4: 应用库 (Applications)",
         ["scikit-learn, TensorFlow, PyTorch, Keras"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(0.6), Inches(3.3), Inches(5.5), Inches(1.3),
         "Layer 3: 领域特定库 (Domain-specific)",
         ["Pandas (表格数据)", "SciPy (科学算法), matplotlib (可视化)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.7), Inches(5.5), Inches(1.3),
         "Layer 2: NumPy (基础设施)",
         ["ndarray, ufunc, broadcasting", "linear algebra, FFT, random"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(0.6), Inches(6.1), Inches(5.5), Inches(0.9),
         "Layer 1: Python + C Extensions",
         ["BLAS, LAPACK, C/Fortran 编译代码"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(6.5), Inches(2.0), Inches(6.2), Inches(5.0),
         "核心要点",
         ["NumPy作为生态系统的基础设施(foundation infrastructure)",
          "Pandas DataFrame内部包装(wrapping)了NumPy数组",
          "scikit-learn的fit()/predict()接受NumPy数组作为输入",
          "TensorFlow/PyTorch张量也提供与NumPy数组的互转",
          "",
          "Andrew Ng: '数据为王 (Data is King)'",
          "ML项目中数据处理占50~80%时间 (Wickham, 2014)"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 幻灯片 6: ML管道中的数据流
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 1, "ML管道中的数据流", "Pandas -> NumPy -> scikit-learn 数据转换过程")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(5.0), [
    "import pandas as pd",
    "import numpy as np",
    "from sklearn.model_selection import train_test_split",
    "from sklearn.preprocessing import StandardScaler",
    "from sklearn.linear_model import LogisticRegression",
    "",
    "# 第1步: 用Pandas加载数据",
    "df = pd.read_csv('data.csv')",
    "",
    "# 第2步: 用Pandas探索和预处理",
    "df = df.dropna()",
    "df['new_feature'] = df['feature_a'] * df['feature_b']",
    "",
    "# 第3步: 转换为NumPy数组",
    "X = df[['feature_a','feature_b','new_feature']].to_numpy()",
    "y = df['target'].to_numpy()",
    "",
    "# 第4步: 用scikit-learn建模 (内部使用NumPy)",
    "X_train, X_test, y_train, y_test = train_test_split(X, y)",
    "scaler = StandardScaler()",
    "X_train_scaled = scaler.fit_transform(X_train)",
    "model = LogisticRegression()",
    "model.fit(X_train_scaled, y_train)",
], font_size=11)

add_card(s, Inches(8.5), Inches(2.0), Inches(4.2), Inches(5.0),
         "管道流程概要",
         ["数据收集 -> 数据探索(EDA)",
          "-> 数据预处理",
          "-> 特征工程",
          "-> 模型训练",
          "-> 模型评估 -> 部署",
          "",
          "Pandas: 负责加载、探索、预处理",
          "NumPy: 提供数值运算基础",
          "scikit-learn: 建模 (内部使用NumPy)",
          "",
          "实际项目中模型训练之前的所有步骤",
          "都属于数据处理 (占总时间的50-80%)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 7: 为什么选择Python?
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 1, "为什么选择Python?", "作为粘合语言(Glue Language)的角色")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "Python的核心策略",
         ["Python本身很慢，但性能关键部分由",
          "C/Fortran编写的库(BLAS, LAPACK)处理",
          "NumPy是这一策略的典型成功案例",
          "粘合语言(glue language): 连接各种库的角色"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "McKinney (2010)开发Pandas的动机",
         ["'Python中没有与R的data.frame对应的",
          "  高级数据结构。'",
          "NumPy ndarray: 擅长同质数值数据",
          "但对异质表格数据支持不足",
          "-> 为填补这一空白，Pandas诞生!"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.2),
         "NumPy + Pandas的角色分工",
         ["NumPy: 同质(homogeneous)数值数据的高速运算 -- ndarray将相同类型数据存储在连续内存中",
          "Pandas: 异质(heterogeneous)表格数据 -- DataFrame允许每列不同的dtype (字符串+数字+日期)",
          "这一组合在使Python成为数据科学主导语言中起到了决定性作用"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# Section 2: NumPy深入
# ============================================================
section_divider("NumPy深入: ndarray的世界", "ndarray内部结构、内存布局、广播、线性代数", 2, ACCENT_CYAN)

# ============================================================
# 幻灯片 8: ndarray内部结构 - 4个核心要素
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "ndarray内部结构: 4个核心要素", "基于Walt et al. (2011)")

add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(1.8),
         "data (数据指针)",
         ["指向数据存储内存缓冲区的指针",
          "实际数值数据存储在连续内存中",
          "与Python列表的根本区别"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(1.8),
         "dtype (数据类型)",
         ["每个元素的数据类型 (float64, int32等)",
          "所有元素类型相同 (同质的)",
          "决定内存大小和运算方式"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(1.8),
         "shape (形状)",
         ["表示每个维度大小的元组",
          "例: (2, 3) -> 2行3列",
          "ndim: 维度数, size: 总元素数"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.0), Inches(5.5), Inches(1.5),
         "strides (步幅)",
         ["在每个维度中到下一个元素的字节数",
          "例: (24, 8) -> 行移动24字节, 列移动8字节",
          "转置(transpose)只交换stride (无数据复制!)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_code_block(s, Inches(6.5), Inches(4.0), Inches(6.2), Inches(3.2), [
    "import numpy as np",
    "",
    "arr = np.array([[1.0, 2.0, 3.0],",
    "                [4.0, 5.0, 6.0]])",
    "",
    "print(f'dtype:   {arr.dtype}')    # float64",
    "print(f'shape:   {arr.shape}')    # (2, 3)",
    "print(f'strides: {arr.strides}')  # (24, 8)",
    "print(f'nbytes:  {arr.nbytes}')   # 48 = 2*3*8",
    "print(f'ndim:    {arr.ndim}')     # 2",
    "print(f'size:    {arr.size}')     # 6",
], font_size=11)

# ============================================================
# 幻灯片 9: 内存布局与Stride
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "内存布局与Stride", "ndarray vs Python列表的根本区别")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.0),
         "Python列表的内存结构",
         ["列表对象 --> [指针1, 指针2, ...]",
          "每个指针 -> 独立的PyObject (28+ bytes each)",
          "不连续内存: CPU缓存未命中频繁",
          "每次元素访问时都有类型检查开销"],
         title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "NumPy数组的内存结构",
         ["ndarray对象 --> [8B | 8B | 8B | ...]",
          "在连续内存块中存储同质类型数据",
          "CPU缓存命中率(hit rate)最大化",
          "可利用SIMD(单指令多数据)"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8),
         "Stride可视化说明 (2x3 float64, C-order)",
         ["内存: [1.0][2.0][3.0][4.0][5.0][6.0]",
          "字节:  0    8   16   24   32   40",
          "",
          "shape = (2, 3)",
          "strides = (24, 8)",
          "  行方向(axis=0): 24字节 = 3元素 x 8B/元素",
          "  列方向(axis=1):  8字节 = 1元素 x 8B/元素"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_code_block(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.8), [
    "a = np.array([[1, 2, 3, 4],",
    "              [5, 6, 7, 8],",
    "              [9, 10, 11, 12]], dtype=np.float64)",
    "",
    "print(f'shape: {a.shape}')      # (3, 4)",
    "print(f'strides: {a.strides}')  # (32, 8)",
    "",
    "# 转置(transpose)只改变stride",
    "# (无数据复制!)",
    "b = a.T",
    "print(f'转置 shape: {b.shape}')    # (4, 3)",
    "print(f'转置 strides: {b.strides}')# (8, 32)",
], font_size=11)

# ============================================================
# 幻灯片 10: C-order vs F-order, dtype系统
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "C-order vs F-order / dtype系统", "内存顺序与数据类型体系")

add_table_slide(s,
    headers=["顺序", "内存布局", "快速访问方向", "使用示例"],
    rows=[
        ["C-order (row-major)", "行优先存储", "最后一个轴 (列方向)", "NumPy默认, C"],
        ["F-order (col-major)", "列优先存储", "第一个轴 (行方向)", "Fortran, MATLAB, R"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 2.5, 2.5, 2.5],
    header_color=ACCENT_BLUE)

add_table_slide(s,
    headers=["类别", "dtype", "字节", "范围/用途"],
    rows=[
        ["整数", "int8 / int16", "1 / 2", "-128~127 / -32768~32767"],
        ["整数", "int32 / int64", "4 / 8", "一般整数 / 大整数"],
        ["无符号整数", "uint8", "1", "0~255 (图像像素)"],
        ["浮点数", "float16 / float32", "2 / 4", "深度学习推理 / 训练"],
        ["浮点数", "float64", "8", "科学计算默认值"],
        ["复数/布尔", "complex128 / bool_", "16 / 1", "信号处理 / True-False"],
    ],
    left=Inches(0.6), top=Inches(3.5), col_widths=[2.5, 2.5, 2.0, 3.0],
    header_color=ACCENT_PURPLE, row_height=0.4, font_size=12)

add_text(s, Inches(0.6), Inches(6.7), Inches(11), Inches(0.4),
         "性能启示: 行方向遍历时C-order，列方向遍历时F-order的CPU缓存效率更高。",
         font_size=14, color=ACCENT_ORANGE, bold=True)

# ============================================================
# 幻灯片 11: 数组创建函数总结
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "数组创建函数总结", "各种NumPy数组创建方法")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.2), [
    "# === 基本创建 ===",
    "np.array([1, 2, 3])",
    "np.array([[1, 2], [3, 4]], dtype=float)",
    "",
    "# === 初始化创建 ===",
    "np.zeros((3, 4))        # 用0填充的3x4",
    "np.ones((2, 3))         # 用1填充的2x3",
    "np.full((2, 3), 7)      # 用7填充的2x3",
    "np.empty((3, 3))        # 不初始化直接分配",
    "np.eye(4)               # 4x4单位矩阵",
], font_size=11)

add_code_block(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.2), [
    "# === 范围创建 ===",
    "np.arange(0, 10, 2)     # [0, 2, 4, 6, 8]",
    "np.linspace(0, 1, 5)    # [0, 0.25, 0.5, 0.75, 1]",
    "np.logspace(0, 3, 4)    # [1, 10, 100, 1000]",
    "",
    "# === 随机数创建 ===",
    "np.random.seed(42)",
    "np.random.rand(3, 4)          # U(0,1)",
    "np.random.randn(3, 4)         # N(0,1)",
    "np.random.randint(0, 10, (3,4))# 整数随机数",
], font_size=11)

add_card(s, Inches(0.6), Inches(4.5), Inches(12.0), Inches(2.8),
         "主要创建函数比较",
         ["zeros/ones/full/empty: 创建特定值初始化的数组 (传入shape元组)",
          "eye/diag: 创建单位矩阵、对角矩阵 (线性代数运算必需)",
          "arange: 整数范围 (类似Python range), linspace: 均匀分割 (常用于可视化)",
          "logspace: 对数尺度均匀分割 (学习率搜索等)",
          "random.rand/randn: 均匀分布/正态分布随机数 (ML权重初始化中使用)",
          "random.choice: 随机选择 (自助法、交叉验证等)"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 幻灯片 12: 数组索引与切片
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "数组索引与切片", "4种索引方法 (Harris et al., 2020)")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), [
    "arr = np.array([[1,  2,  3,  4],",
    "                [5,  6,  7,  8],",
    "                [9, 10, 11, 12]])",
    "",
    "# 1. 基本索引 (Basic Indexing)",
    "arr[0, 1]              # 2",
    "",
    "# 2. 切片 (Slicing) -- 返回视图(view)",
    "arr[0:2, 1:3]          # [[2, 3], [6, 7]]",
    "arr[:, 2]              # [3, 7, 11]",
    "",
    "# 3. 花式索引 (Fancy) -- 返回副本(copy)",
    "arr[[0, 2], [1, 3]]    # [2, 12]",
    "",
    "# 4. 布尔索引 (Boolean)",
    "arr[arr > 5]           # [6,7,8,9,10,11,12]",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.0),
         "View vs Copy 核心区别",
         ["切片 (a[1:4]): 返回View -> 共享内存",
          "转置 (a.T): 返回View -> 只交换stride",
          "花式索引 (a[[0,2]]): 返回Copy",
          "布尔索引 (a[mask]): 返回Copy",
          "显式复制: a.copy()始终返回Copy"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(7.0), Inches(4.3), Inches(5.6), Inches(2.5),
         "注意事项",
         ["修改View时原始数据也会改变!",
          "  b = a[1:4]; b[0]=99 -> a也改变!",
          "安全做法: c = a[1:4].copy()",
          "",
          "reshape(): 尽可能返回View",
          "flatten(): 始终返回Copy",
          "ravel(): 尽可能返回View"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 幻灯片 13: 广播的3条规则
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "广播(Broadcasting): 3条规则", "不同大小数组间运算的自动机制")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(1.3),
         "规则1: 对齐维度数",
         ["当两个数组的维度数(ndim)不同时，在维度数较少的数组shape前面添加1。",
          "例: shape (3,) -> (1, 3)   |   shape (4,) -> (1, 4)   |   shape (5,) -> (1, 1, 5)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(0.6), Inches(3.5), Inches(12.0), Inches(1.3),
         "规则2: 检查兼容性",
         ["在每个维度上，大小相同或其中一个为1时兼容(compatible)。",
          "不兼容时抛出ValueError。例: (3,) + (4,) -> 3 != 4 且都不是1 -> ERROR"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(5.0), Inches(12.0), Inches(1.3),
         "规则3: 虚拟扩展",
         ["大小为1的维度按另一个数组的大小虚拟扩展(stretch)。",
          "核心: 实际上不复制数据! 内部通过将stride设为0来实现内存效率。"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_text(s, Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5),
         "例: (3,1)+(1,4) -> 规则1不需要, 规则2兼容(3vs1, 1vs4), 规则3扩展 -> 结果shape (3,4)",
         font_size=15, color=ACCENT_BLUE, bold=True)

# ============================================================
# 幻灯片 14: 广播代码示例
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "广播代码示例", "在ML中的实际应用")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(5.0), [
    "# 示例1: 矩阵 + 向量",
    "A = np.array([[1,2,3],[4,5,6],[7,8,9]])",
    "b = np.array([10, 20, 30])",
    "print(A + b)  # 每行加上b",
    "",
    "# 示例2: 列向量 + 行向量 -> 矩阵 (外积模式)",
    "col = np.array([[1],[2],[3]])  # (3,1)",
    "row = np.array([10,20,30,40]) # (4,)-->(1,4)",
    "print(col + row)  # (3,4) 矩阵",
    "",
    "# 示例3: ML数据中心化(centering)",
    "X = np.random.randn(100, 5)",
    "mean = X.mean(axis=0)          # (5,)",
    "X_centered = X - mean          # (100,5)-(5,)",
    "",
    "# 示例4: Z-score标准化",
    "std = X.std(axis=0)",
    "X_standardized = (X - mean) / std",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "可视化说明: (3,3) + (3,)",
         ["数组 b: (3,) -> (1,3) -> (3,3)",
          "[a00 a01 a02]   [b0 b1 b2]",
          "[a10 a11 a12] + [b0 b1 b2]",
          "[a20 a21 a22]   [b0 b1 b2]",
          "",
          "-> 每行加上相同的向量b"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.2),
         "在ML中的核心应用",
         ["数据中心化: X - X.mean(axis=0)",
          "Z-score标准化: (X - mean) / std",
          "批量归一化(Batch Normalization)",
          "加权求和(weighted sum): X @ w + b"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 15: 通用函数 (ufunc)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "通用函数 (ufunc)", "向量化的逐元素运算")

add_table_slide(s,
    headers=["分类", "函数示例", "说明"],
    rows=[
        ["数学函数", "np.sin, np.cos, np.exp, np.log, np.sqrt", "逐元素数学运算"],
        ["比较函数", "np.greater, np.equal, np.logical_and", "逐元素比较"],
        ["算术函数", "np.add, np.subtract, np.multiply", "四则运算"],
        ["聚合函数", "np.sum, np.prod, np.min, np.max", "按轴聚合"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 5.0, 4.5],
    header_color=ACCENT_BLUE, row_height=0.5)

add_code_block(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.8), [
    "A = np.array([[1, 2, 3],",
    "              [4, 5, 6]])",
    "",
    "print(np.sum(A))          # 21 (总和)",
    "print(np.sum(A, axis=0))  # [5, 7, 9] (列方向)",
    "print(np.sum(A, axis=1))  # [6, 15] (行方向)",
    "print(np.mean(A, axis=0)) # [2.5, 3.5, 4.5]",
], font_size=12)

add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.8),
         "理解axis",
         ["axis=None: 对所有元素聚合",
          "axis=0: 沿行方向缩减 (按列聚合)",
          "axis=1: 沿列方向缩减 (按行聚合)",
          "",
          "ufunc的核心: 基于C循环的向量化运算",
          "相比Python循环提升100~1000倍性能"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 16: 线性代数基础
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "线性代数基本运算", "ML的数学基础")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "A = np.array([[1,2],[3,4]], dtype=float)",
    "B = np.array([[5,6],[7,8]], dtype=float)",
    "",
    "print(A * B)                # 逐元素乘法",
    "print(A @ B)                # 矩阵乘法",
    "print(A.T)                  # 转置",
    "print(np.linalg.det(A))     # 行列式",
    "print(np.linalg.inv(A))     # 逆矩阵",
    "print(np.trace(A))          # 迹",
    "print(np.linalg.norm(A,'fro'))# Frobenius范数",
], font_size=11)

add_table_slide(s,
    headers=["分解方法", "公式", "ML应用", "NumPy函数"],
    rows=[
        ["LU分解", "A = P * L * U", "线性方程组, 行列式", "scipy.linalg.lu()"],
        ["QR分解", "A = Q * R (Q^T Q = I)", "最小二乘法稳定求解", "np.linalg.qr()"],
        ["SVD", "A = U * S * V^T", "PCA, 推荐系统, LSA", "np.linalg.svd()"],
        ["特征值分解", "Av = lambda*v", "PCA, 谱聚类", "np.linalg.eigh()"],
    ],
    left=Inches(0.6), top=Inches(5.2), col_widths=[2.0, 2.5, 3.5, 2.5],
    header_color=ACCENT_PURPLE, row_height=0.45, font_size=12)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.8),
         "主要矩阵运算概要",
         ["* (逐元素) vs @ (矩阵乘): 必须区分!",
          "A.T: 转置 (View, 只交换stride)",
          "det(A): 行列式 (判断可逆性)",
          "inv(A): 逆矩阵 (需注意数值问题)",
          "trace(A): 迹 = 特征值之和",
          "norm(A): L1, L2, Frobenius范数"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 17: SVD与特征值分解
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "SVD(奇异值分解)与特征值分解", "PCA降维的数学基础")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "# SVD: A = U * Sigma * V^T",
    "A = np.array([[1,2,0],[0,1,1],[2,0,1],",
    "              [1,1,1],[3,2,1]], dtype=float)",
    "",
    "U, s, Vt = np.linalg.svd(A, full_matrices=False)",
    "",
    "# 降维: 只使用前k个奇异值",
    "k = 2",
    "A_approx = U[:,:k] @ np.diag(s[:k]) @ Vt[:k,:]",
    "",
    "# 方差解释力",
    "energy = s**2 / np.sum(s**2) * 100",
    "print(f'累积解释比例: {np.cumsum(energy)}')",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(5.2), Inches(6.0), Inches(2.0), [
    "# 特征值分解: Av = lambda*v",
    "cov = np.array([[4,2,1],[2,3,1],[1,1,2]],",
    "               dtype=float)",
    "eigenvalues, eigenvectors = np.linalg.eigh(cov)",
    "variance_ratio = eigenvalues/np.sum(eigenvalues)*100",
    "print(f'方差解释比例: {variance_ratio}')",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "SVD在ML中的应用",
         ["PCA降维: 用前k个奇异值近似数据",
          "推荐系统: 用户-物品矩阵的潜在因子",
          "潜在语义分析(LSA): 文档-词矩阵",
          "数据压缩: 计算信息保留比例",
          "方差解释力 = s^2 / sum(s^2) * 100"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.4),
         "特征值分解在ML中的应用",
         ["PCA: 协方差矩阵的特征值分解",
          "谱聚类: 拉普拉斯矩阵",
          "特征值 = 该方向的方差量",
          "eigh(): 对称矩阵专用 (更快更稳定)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 18: 正规方程
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "正规方程 (Normal Equation)", "线性回归的解析解: beta = (X^T X)^{-1} X^T y")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(4.8), [
    "np.random.seed(42)",
    "n_samples = 50",
    "X_raw = np.random.uniform(0, 10, (n_samples, 1))",
    "noise = np.random.normal(0, 1, (n_samples, 1))",
    "y = 3 * X_raw + 2 + noise  # 真值: 斜率=3, 截距=2",
    "",
    "X = np.hstack([np.ones((n_samples, 1)), X_raw])",
    "",
    "# 方法1: 逆矩阵 (数值不稳定 -- 不推荐)",
    "beta_inv = np.linalg.inv(X.T @ X) @ X.T @ y",
    "",
    "# 方法2: solve (基于LU分解, 数值稳定)",
    "beta_solve = np.linalg.solve(X.T @ X, X.T @ y)",
    "",
    "# 方法3: lstsq (基于SVD, 最稳定, 推荐!)",
    "beta_lstsq, _, _, _ = np.linalg.lstsq(X, y, rcond=None)",
    "",
    "print(f'真值: 截距=2, 斜率=3')",
    "print(f'估计: 截距={beta_lstsq[0,0]:.4f}, 斜率={beta_lstsq[1,0]:.4f}')",
], font_size=11)

add_card(s, Inches(8.0), Inches(2.0), Inches(4.6), Inches(5.0),
         "3种求解方法比较",
         ["方法1: inv(X^T X) @ X^T @ y",
          "  -> 条件数大的矩阵数值不稳定",
          "",
          "方法2: solve(X^T X, X^T @ y)",
          "  -> 基于LU分解, 更稳定",
          "",
          "方法3: lstsq(X, y) [推荐!]",
          "  -> 基于SVD, 最稳定",
          "  -> 可处理秩不足矩阵",
          "",
          "实际推荐: 始终使用lstsq()",
          "逆矩阵因条件数问题不推荐"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 幻灯片 19: 线性代数与ML关系总结
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "线性代数与ML关系总结", "各分解在ML算法中的应用")

add_table_slide(s,
    headers=["线性代数运算", "ML应用", "NumPy函数"],
    rows=[
        ["SVD", "PCA, 推荐系统, LSA", "np.linalg.svd()"],
        ["特征值分解", "PCA, 谱聚类", "np.linalg.eigh()"],
        ["QR分解", "最小二乘法稳定求解", "np.linalg.qr()"],
        ["LU分解", "线性方程组高效求解", "scipy.linalg.lu()"],
        ["正规方程", "线性回归", "np.linalg.lstsq()"],
        ["矩阵乘法", "神经网络前向传播", "@ 运算符"],
        ["范数", "L1, L2正则化", "np.linalg.norm()"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.0, 4.5, 3.5],
    header_color=ACCENT_PURPLE, row_height=0.5)

add_card(s, Inches(0.6), Inches(5.8), Inches(12.0), Inches(1.5),
         "核心信息",
         ["线性代数是ML算法的数学引擎。NumPy是这个引擎的实现。",
          "PCA = SVD或特征值分解 | 线性回归 = 正规方程 | 神经网络 = 矩阵乘法的链式运算",
          "所有ML算法的内部都存在线性代数运算。"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 20: View vs Copy
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "View vs Copy详解", "内存共享与数据安全性")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.5), [
    "a = np.array([1, 2, 3, 4, 5])",
    "b = a[1:4]        # 视图(View)",
    "b[0] = 99",
    "print(a)           # [1, 99, 3, 4, 5] -- 原始数据也改变!",
    "",
    "c = a[1:4].copy()  # 副本(Copy)",
    "c[0] = 100",
    "print(a)           # [1, 99, 3, 4, 5] -- 原始不变",
], font_size=12)

add_table_slide(s,
    headers=["操作", "返回类型", "内存共享"],
    rows=[
        ["切片 a[1:4]", "View", "O"],
        ["转置 a.T", "View", "O"],
        ["reshape()", "View (可能时)", "O"],
        ["ravel()", "View (可能时)", "O"],
        ["花式索引 a[[0,2]]", "Copy", "X"],
        ["布尔索引 a[mask]", "Copy", "X"],
        ["flatten()", "Copy", "X"],
        ["copy()", "Copy", "X"],
    ],
    left=Inches(0.6), top=Inches(4.8), col_widths=[3.5, 3.0, 2.0],
    header_color=ACCENT_ORANGE, row_height=0.35, font_size=12)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.2),
         "View vs Copy判断标准",
         ["View: 与原始数据共享内存, 效率高",
          "  -> 修改会影响原始数据 (注意!)",
          "  -> 切片、转置、reshape等",
          "",
          "Copy: 复制到独立内存",
          "  -> 与原始独立, 修改安全",
          "  -> 花式索引、布尔索引等",
          "",
          "检查: np.shares_memory(a, b)",
          "安全复制: 显式调用a.copy()",
          "",
          "实践提示: 修改前确认是View还是Copy!"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# Section 3: Pandas深入
# ============================================================
section_divider("Pandas深入: DataFrame的世界", "DataFrame结构、索引、方法链、GroupBy、Merge/Join", 3, ACCENT_ORANGE)

# ============================================================
# 幻灯片 21: DataFrame内部结构
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "DataFrame内部结构", "McKinney (2010)设计的异质表格数据结构")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "import pandas as pd",
    "import numpy as np",
    "",
    "df = pd.DataFrame({",
    "    '姓名': ['张三', '李四', '王五'],",
    "    '年龄': [28, 35, 42],",
    "    '身高':   [175.5, 162.3, 180.1],",
    "    '学生': [True, False, False]",
    "})",
    "",
    "print(df.dtypes)",
    "# 姓名    object  /  年龄    int64",
    "# 身高    float64  /  学生      bool",
], font_size=11)

add_table_slide(s,
    headers=["属性", "说明", "示例"],
    rows=[
        ["df.shape", "(行数, 列数)", "(3, 4)"],
        ["df.dtypes", "每列的数据类型", "int64, float64, ..."],
        ["df.index", "行索引", "RangeIndex(0, 3)"],
        ["df.columns", "列名", "Index(['姓名','年龄',...])"],
        ["df.values", "转换为NumPy数组", "ndarray"],
        ["df.info()", "整体信息", "行/列数, dtype, 内存"],
        ["df.describe()", "描述性统计", "均值, 标准差, 四分位数"],
    ],
    left=Inches(7.0), top=Inches(2.0), col_widths=[1.8, 2.2, 2.2],
    header_color=ACCENT_ORANGE, row_height=0.4, font_size=11)

add_card(s, Inches(0.6), Inches(5.3), Inches(12.0), Inches(1.8),
         "DataFrame的核心特征",
         ["每列作为独立的NumPy数组(或Extension Array)存储 -> 每列可以有不同dtype",
          "通过行索引(Index)实现自动数据对齐(automatic alignment)",
          "系统化的缺失值(NaN)处理: isna(), fillna(), dropna()等一致的API"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# 幻灯片 22: 索引 loc vs iloc
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "索引: loc vs iloc", "基于标签 vs 基于整数位置")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), [
    "# 1. 列选择",
    "df['年龄']                  # 返回Series",
    "df[['姓名', '年龄']]         # 返回DataFrame",
    "",
    "# 2. loc - 基于标签的索引",
    "df.loc[0, '姓名']           # '张三'",
    "df.loc[0:1, ['姓名','年龄']] # 行0~1 (包含末尾!)",
    "",
    "# 3. iloc - 基于整数位置的索引",
    "df.iloc[0, 0]              # '张三'",
    "df.iloc[0:2, 0:2]          # 行0~1 (不包含末尾!)",
    "",
    "# 4. 布尔索引",
    "df[df['年龄'] > 30]         # 年龄 > 30的行",
    "df.query('年龄 > 30 and 身高 > 170')",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "loc vs iloc核心区别",
         ["loc: 通过标签(名称)访问",
          "  -> 切片时包含末尾索引",
          "  -> df.loc[0:2] -> 行0, 1, 2 (3个)",
          "",
          "iloc: 通过整数位置访问",
          "  -> 切片时不包含末尾索引",
          "  -> df.iloc[0:2] -> 行0, 1 (2个)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.2),
         "query()方法",
         ["用字符串表达式进行行过滤",
          "可读性好，适合链式调用",
          "内部使用numexpr引擎 (快速)",
          "例: df.query('年龄 > 30 and 身高 > 170')"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 23: 方法链 (Method Chaining)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "方法链 (Method Chaining)", "无需中间变量连续应用多个转换")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.5), [
    "# 不好的例子: 滥用中间变量",
    "df2 = df.dropna()",
    "df3 = df2[df2['年龄'] > 25]",
    "df4 = df3.assign(年龄段=df3['年龄']//10*10)",
    "result = df4.sort_values('身高', ascending=False)",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.5), [
    "# 好的例子: 方法链",
    "result = (df",
    "    .dropna()",
    "    .query('年龄 > 25')",
    "    .assign(年龄段=lambda x: x['年龄']//10*10)",
    "    .sort_values('身高', ascending=False)",
    ")",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.2),
         "pipe()整合自定义函数",
         ["def 去除异常值(df, col, n_std=3):",
          "    mean, std = df[col].mean(), df[col].std()",
          "    return df[abs(df[col]-mean) <= n_std*std]",
          "",
          "result = (df",
          "    .pipe(去除异常值, '身高')",
          "    .assign(BMI=lambda x: x['身高']/100))"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.5), Inches(5.6), Inches(2.7),
         "方法链的优点",
         ["可读性: 数据转换过程从上到下阅读",
          "调试: 注释各步骤查看中间结果",
          "可重现性: 整个转换为一个表达式",
          "pipe(): 自定义函数也可集成到链中",
          "assign(): 添加新列 (用lambda引用当前df)"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 24: GroupBy - 分割-应用-合并
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "GroupBy: 分割-应用-合并 (Split-Apply-Combine)", "McKinney (2010)强调的核心模式")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(5.0), [
    "np.random.seed(42)",
    "df = pd.DataFrame({",
    "    '部门': np.random.choice(['开发','市场','销售'], 100),",
    "    '职级': np.random.choice(['员工','主管','经理'], 100),",
    "    '年薪': np.random.normal(5000,1000,100).astype(int),",
    "    '绩效': np.random.uniform(0,100,100).round(1)",
    "})",
    "",
    "# agg: 分组聚合",
    "df.groupby('部门')['年薪'].agg(['mean','std','min','max'])",
    "",
    "# transform: 保持原始大小 (分组Z-score)",
    "df['年薪_zscore'] = df.groupby('部门')['年薪'].transform(",
    "    lambda x: (x - x.mean()) / x.std())",
    "",
    "# apply: 高自由度的分组操作",
    "def 前N名(group, n=3):",
    "    return group.nlargest(n, '绩效')",
    "df.groupby('部门').apply(前N名, n=2)",
], font_size=11)

add_card(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(5.0),
         "GroupBy的3种模式",
         ["1. agg (聚合): 组 -> 标量",
          "   分组平均、标准差、最小、最大",
          "   结果: 组数个行",
          "",
          "2. transform (转换): 组 -> 原始大小",
          "   分组Z-score、累积和、比例",
          "   结果: 与原始相同大小 (非常有用!)",
          "",
          "3. apply (自由形式): 最灵活",
          "   分组提取前N个",
          "   应用自定义函数",
          "",
          "Split-Apply-Combine模式:",
          "分割(groupby) -> 应用(agg/transform/apply)",
          "-> 合并(自动)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 25: Merge与Join
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "Merge与Join", "与SQL JOIN相同的表连接操作")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.2), [
    "# 准备两个表",
    "员工 = pd.DataFrame({",
    "    '工号': [1, 2, 3, 4],",
    "    '姓名': ['张三','李四','王五','赵六'],",
    "    '部门代码': ['D01','D02','D01','D03']",
    "})",
    "部门 = pd.DataFrame({",
    "    '代码': ['D01','D02','D03'],",
    "    '部门名': ['开发部','市场部','销售部']",
    "})",
    "",
    "# 与SQL JOIN相同",
    "result = pd.merge(员工, 部门,",
    "    left_on='部门代码', right_on='代码', how='left')",
], font_size=11)

add_table_slide(s,
    headers=["Join类型", "SQL对应", "说明"],
    rows=[
        ["inner", "INNER JOIN", "只保留两侧都有键的行"],
        ["left", "LEFT JOIN", "以左表为基准"],
        ["right", "RIGHT JOIN", "以右表为基准"],
        ["outer", "FULL OUTER JOIN", "包含两侧所有行"],
    ],
    left=Inches(0.6), top=Inches(5.5), col_widths=[2.0, 3.0, 5.0],
    header_color=ACCENT_ORANGE, row_height=0.4, font_size=12)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(3.2),
         "Merge核心参数",
         ["on: 两侧表的键列名相同时",
          "left_on / right_on: 键列名不同时",
          "how: 'inner','left','right','outer'",
          "",
          "concat: 行方向/列方向简单拼接",
          "  pd.concat([df1, df2], axis=0)",
          "",
          "join: 基于索引的连接",
          "  df1.join(df2, how='left')"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# 幻灯片 26: Apply/Map与时间序列窗口函数
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "Apply/Map与时间序列窗口函数", "自定义函数应用与移动平均")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.5), [
    "# apply: 对DataFrame的每列/行应用函数",
    "df.apply(np.mean, axis=0)  # 按列求均值",
    "df.apply(np.mean, axis=1)  # 按行求均值",
    "",
    "# map: 对Series的每个元素应用函数/字典",
    "df['Sex'] = df['Sex'].map({'male':0, 'female':1})",
    "",
    "# applymap: 对DataFrame所有元素应用函数",
    "df.applymap(lambda x: round(x, 2))",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.5), [
    "# 时间序列窗口函数",
    "df['MA_7'] = df['股价'].rolling(window=7).mean()",
    "df['MA_30'] = df['股价'].rolling(window=30).mean()",
    "df['波动率'] = df['股价'].rolling(window=20).std()",
    "",
    "# 累积窗口",
    "df['累积均值'] = df['股价'].expanding().mean()",
    "",
    "# 指数加权移动平均 (EWM)",
    "df['EWM_12'] = df['股价'].ewm(span=12).mean()",
], font_size=11)

add_table_slide(s,
    headers=["窗口函数", "说明", "使用示例"],
    rows=[
        ["rolling(n)", "固定大小n的窗口", "移动平均, 移动标准差"],
        ["expanding()", "从开始到当前的累积", "累积最大值, 累积均值"],
        ["ewm(span=n)", "指数加权", "MACD, 指数移动平均"],
    ],
    left=Inches(7.0), top=Inches(2.0), col_widths=[1.6, 2.2, 2.4],
    header_color=ACCENT_ORANGE, row_height=0.45, font_size=11)

add_card(s, Inches(7.0), Inches(4.0), Inches(5.6), Inches(3.2),
         "实践提示: 移动平均交叉(MA Crossover)",
         ["金叉: 短期MA(7日)上穿长期MA(30日)",
          "  -> 买入信号",
          "死叉: 下穿 -> 卖出信号",
          "",
          "MACD = EWM(12) - EWM(26)",
          "技术分析的核心指标",
          "",
          "McKinney(2010)基于金融数据分析的",
          "经验内置于Pandas"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 27: 数据透视表与MultiIndex
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "数据透视表与MultiIndex", "对应SQL GROUP BY + CASE WHEN的数据汇总")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "# 基本数据透视表",
    "透视 = pd.pivot_table(",
    "    销售,",
    "    values='销售额',",
    "    index='地区',       # 行基准",
    "    columns='产品',      # 列基准",
    "    aggfunc='mean',     # 聚合函数",
    "    margins=True         # 添加合计行/列",
    ")",
    "",
    "# 交叉表 (Crosstab) - 频次统计",
    "pd.crosstab(销售['地区'], 销售['产品'], margins=True)",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(5.2), Inches(6.0), Inches(2.0), [
    "# MultiIndex (多级索引)",
    "df_multi.loc['北京']                # 北京的所有季度",
    "df_multi.loc[('北京', '第一季度')]     # 特定行",
    "df_multi.xs('第一季度', level='季度')  # 选择特定级别",
    "df_multi.unstack(level='季度')      # 将季度移到列",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.2),
         "Pandas内存优化4种策略",
         ["1. 数值类型降级:",
          "   int64 -> int32/int16 (节约50~75%)",
          "",
          "2. Categorical类型转换:",
          "   对重复字符串有效 (节约90~99%)",
          "   100万行 x 3个唯一值: 64MB -> 1MB",
          "",
          "3. Sparse类型:",
          "   大部分为0的数据 (节约90%+)",
          "",
          "4. 分块读取 (chunksize):",
          "   chunks = pd.read_csv('huge.csv',",
          "       chunksize=100000)",
          "   在RAM限制内处理"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 28: SQL vs Pandas比较
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "SQL vs Pandas比较", "为了解SQL的人提供的Pandas映射")

add_table_slide(s,
    headers=["SQL", "Pandas", "备注"],
    rows=[
        ["SELECT col FROM t", "df['col'] 或 df[['col']]", "列选择"],
        ["WHERE cond", "df[df['col']>v] / df.query()", "行过滤"],
        ["GROUP BY col", "df.groupby('col')", "分组聚合"],
        ["ORDER BY col", "df.sort_values('col')", "排序"],
        ["JOIN t1, t2 ON", "pd.merge(t1, t2, on=)", "表连接"],
        ["DISTINCT", "df.drop_duplicates()", "去重"],
        ["COUNT(*)", "df.shape[0] / len(df)", "行数"],
        ["LIMIT n", "df.head(n)", "前n个"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.5, 4.0, 3.5],
    header_color=ACCENT_BLUE, row_height=0.45, font_size=12)

add_card(s, Inches(0.6), Inches(6.2), Inches(12.0), Inches(1.0),
         "核心信息",
         ["有SQL经验的话Pandas非常直观! 几乎可以1:1对应，方法链使分析更加灵活。"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# Section 4: 数据预处理
# ============================================================
section_divider("数据预处理", "缺失值处理、异常值检测、特征缩放、编码", 4, ACCENT_GREEN)

# ============================================================
# 幻灯片 29: 缺失值处理5种策略
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 4, "缺失值处理5种策略", "根据情况选择合适的策略很重要")

add_table_slide(s,
    headers=["策略", "方法", "优点", "缺点", "适用场景"],
    rows=[
        ["删除行", "dropna()", "简单", "数据损失", "缺失 < 5%"],
        ["均值替代", "fillna(mean)", "保持分布", "低估方差", "正态分布, MCAR"],
        ["中位数替代", "fillna(median)", "对异常值鲁棒", "低估方差", "偏斜分布"],
        ["插值", "interpolate()", "保持连续性", "外推风险", "时间序列数据"],
        ["分组替代", "groupby.transform", "反映组特性", "实现复杂", "组间差异大时"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[1.5, 2.0, 1.8, 1.8, 2.2],
    header_color=ACCENT_GREEN, row_height=0.5, font_size=11)

add_code_block(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.5), [
    "df = pd.DataFrame({'A': [1, np.nan, 3, np.nan, 5], 'B': [10,20,np.nan,40,50], '分组': ['X','X','Y','Y','Y']})",
    "",
    "# 各策略比较",
    "print('删除:',   df['A'].dropna().values)                              # 删除NaN",
    "print('均值:',   df['A'].fillna(df['A'].mean()).values)               # 用整体均值替代",
    "print('中位数:', df['A'].fillna(df['A'].median()).values)             # 用整体中位数替代",
    "print('插值:',   df['A'].interpolate().values)                        # 线性插值",
    "print('分组:', df.groupby('分组')['A'].transform(lambda x: x.fillna(x.mean())).values)  # 用组均值替代",
], font_size=10)

# ============================================================
# 幻灯片 30: 异常值检测3种方法
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 4, "异常值检测3种方法", "IQR, Z-score, 可视化方法")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "IQR方法 (四分位距)",
         ["异常值范围: Q1 - 1.5*IQR < x < Q3 + 1.5*IQR",
          "IQR = Q3 - Q1 (四分位距)",
          "不需要正态分布假设",
          "与箱线图的标准相同"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "Z-score方法",
         ["z = (x - mu) / sigma",
          "|z| > 3 时判定为异常值",
          "需要正态分布假设",
          "对均值和标准差敏感"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_code_block(s, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.5), [
    "def iqr_outlier(series):",
    "    Q1, Q3 = series.quantile([0.25, 0.75])",
    "    IQR = Q3 - Q1",
    "    lower = Q1 - 1.5 * IQR",
    "    upper = Q3 + 1.5 * IQR",
    "    return (series < lower) | (series > upper)",
], font_size=11)

add_code_block(s, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.5), [
    "def zscore_outlier(series, threshold=3):",
    "    z = (series - series.mean()) / series.std()",
    "    return abs(z) > threshold",
    "",
    "# 使用示例",
    "mask = iqr_outlier(df['年薪'])",
    "df_clean = df[~mask]  # 去除异常值",
], font_size=11)

# ============================================================
# 幻灯片 31: 特征缩放3种方法
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 4, "特征缩放3种方法", "StandardScaler, MinMaxScaler, RobustScaler")

add_table_slide(s,
    headers=["缩放器", "公式", "特性", "使用场景"],
    rows=[
        ["StandardScaler", "z = (x-mu)/sigma", "均值0, 方差1", "正态分布, SVM, 逻辑回归"],
        ["MinMaxScaler", "z = (x-min)/(max-min)", "[0,1]范围", "神经网络, 图像"],
        ["RobustScaler", "z = (x-Q2)/(Q3-Q1)", "基于中位数/IQR", "存在异常值时"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 3.0, 2.5, 3.0],
    header_color=ACCENT_GREEN, row_height=0.55, font_size=12)

add_code_block(s, Inches(0.6), Inches(4.0), Inches(6.0), Inches(2.0), [
    "from sklearn.preprocessing import (StandardScaler,",
    "    MinMaxScaler, RobustScaler)",
    "",
    "# 重要: fit只在train上, transform在train/test上都做",
    "scaler = StandardScaler()",
    "X_train_scaled = scaler.fit_transform(X_train)",
    "X_test_scaled = scaler.transform(X_test)  # 不用fit!",
], font_size=12)

add_card(s, Inches(7.0), Inches(4.0), Inches(5.6), Inches(3.2),
         "Data Leakage警告!",
         ["对全部数据做fit_transform会导致",
          "测试数据的信息泄露到训练中!",
          "",
          "正确做法:",
          "  1. scaler.fit_transform(X_train) -- 在train上fit",
          "  2. scaler.transform(X_test) -- 在test上只transform",
          "",
          "错误做法:",
          "  scaler.fit_transform(X_all) -- 在全部数据上fit",
          "  -> 发生Data Leakage!"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 幻灯片 32: 编码
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 4, "分类变量编码", "Label Encoding, One-Hot Encoding, Ordinal Encoding")

add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(2.5),
         "Label Encoding",
         ["为每个类别分配整数",
          "例: male->0, female->1",
          "适合树模型",
          "可能产生顺序关系 (注意)",
          "df['Sex'].map({'male':0,'female':1})"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(2.5),
         "One-Hot Encoding",
         ["将每个类别转为二进制向量",
          "例: [1,0,0], [0,1,0], [0,0,1]",
          "无顺序关系 (安全)",
          "唯一值多时维度爆炸问题",
          "pd.get_dummies(df, columns=['地区'])"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(2.5),
         "Ordinal Encoding",
         ["为有序类别分配整数",
          "例: 小学->1, 初中->2, 高中->3",
          "保留顺序信息",
          "不适合名义变量",
          "OrdinalEncoder(categories=...)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.4),
         "编码选择指南",
         ["二元类别 (2个值): Label Encoding -> male/female, yes/no",
          "名义类别 (无序): One-Hot Encoding -> 地区, 颜色, 产品类型",
          "有序类别 (有序): Ordinal Encoding -> 学历, 等级, 满意度",
          "唯一值非常多 (100+): 考虑Target Encoding, Hash Encoding",
          "",
          "注意: 对名义变量使用Label Encoding会导致模型学习错误的顺序关系!"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# Section 5: Tidy Data概念
# ============================================================
section_divider("Tidy Data概念", "Wickham (2014)的整洁数据3原则", 5, ACCENT_PURPLE)

# ============================================================
# 幻灯片 33: Wickham的Tidy Data 3原则
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 5, "Wickham的Tidy Data 3原则", "Journal of Statistical Software, 2014")

add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(2.0),
         "原则1: 变量 = 列",
         ["每个变量(variable)",
          "构成一列(column)。",
          "例: 姓名、年龄、分数等",
          "各自表示为独立的列"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(2.0),
         "原则2: 观测 = 行",
         ["每个观测(observation)",
          "构成一行(row)。",
          "例: 一个学生的信息",
          "表示为一行"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(2.0),
         "原则3: 观测单位 = 表",
         ["每种观测单位类型",
          "构成一个表。",
          "例: 学生表、课程表",
          "分为独立的表"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.3), Inches(12.0), Inches(3.0),
         "Messy Data的5种类型 (Wickham 2014)",
         ["(1) 列标题是值而非变量名  ->  用melt()转换",
          "(2) 一列中存储了多个变量  ->  用str.split() + assign()分离",
          "(3) 变量同时存储在行和列中  ->  melt() + pivot() 组合",
          "(4) 一个表中混合了多个观测单位  ->  分离表",
          "(5) 一个观测单位分散在多个表中  ->  用merge()合并",
          "",
          "scikit-learn的输入格式 '行=样本, 列=特征' 与Tidy Data原则完全一致!",
          "Tidy Data -> groupby/agg自然运行, seaborn可视化兼容, 直接投入ML管道"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 34: melt/pivot转换
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 5, "melt/pivot: Messy -> Tidy转换", "Wide Format <-> Long Format转换")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), [
    "# Messy: 列标题是值的情况",
    "messy = pd.DataFrame({",
    "    '姓名': ['张三', '李四'],",
    "    '2023_语文': [85, 92],",
    "    '2023_数学': [90, 88],",
    "    '2024_语文': [88, 95],",
    "    '2024_数学': [92, 90]",
    "})",
    "",
    "# 转换为Tidy",
    "tidy = (messy",
    "    .melt(id_vars='姓名',",
    "          var_name='科目_年份',",
    "          value_name='分数')",
    "    .assign(",
    "        年份=lambda x: x['科目_年份'].str.split('_').str[0],",
    "        科目=lambda x: x['科目_年份'].str.split('_').str[1]",
    "    )",
    "    .drop(columns='科目_年份')",
    ")",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "melt() -- Wide -> Long",
         ["id_vars: 保留的列 (标识符)",
          "value_vars: 要融化的列 (省略则全部其余)",
          "var_name: 变量名列的名称",
          "value_name: 值列的名称",
          "",
          "结果: 宽表 -> 长表"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.4),
         "pivot() / pivot_table() -- Long -> Wide",
         ["index: 成为行的列",
          "columns: 成为列的列",
          "values: 成为值的列",
          "aggfunc: 聚合函数 (仅pivot_table)",
          "",
          "结果: 长表 -> 宽表"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# Section 6: 向量化 vs 循环
# ============================================================
section_divider("向量化 vs 循环: 性能的科学", "相比Python循环提升100~1000倍速度的原理", 6, ACCENT_RED)

# ============================================================
# 幻灯片 35: 向量化 vs 循环性能比较
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 6, "向量化 vs 循环: 性能比较", "Harris et al. (2020) Nature论文确认的100~1000倍速度提升")

add_table_slide(s,
    headers=["因素", "Python循环", "NumPy向量化"],
    rows=[
        ["执行引擎", "Python解释器", "C/Fortran编译代码"],
        ["类型检查", "每次运算", "只做一次"],
        ["内存访问", "不连续 (指针追踪)", "连续 (缓存友好)"],
        ["SIMD利用", "不可能", "可能 (SSE, AVX)"],
        ["GIL", "持有", "可释放"],
        ["BLAS/LAPACK", "未使用", "使用 (Intel MKL等)"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 3.5, 3.5],
    header_color=ACCENT_RED, row_height=0.45, font_size=12)

add_code_block(s, Inches(0.6), Inches(5.0), Inches(6.0), Inches(2.3), [
    "n = 1_000_000",
    "a = np.random.randn(n); b = np.random.randn(n)",
    "",
    "# Python循环: ~0.3秒",
    "c_loop = [a[i]+b[i] for i in range(n)]",
    "",
    "# NumPy向量化: ~0.001秒 (快300倍!)",
    "c_vec = a + b",
], font_size=12)

add_card(s, Inches(7.0), Inches(5.0), Inches(5.6), Inches(2.3),
         "速度提升原因总结",
         ["1. C编译代码: 消除Python解释器开销",
          "2. 连续内存: 最大化CPU L1/L2缓存命中率",
          "3. SIMD: 单条指令同时处理多个数据",
          "4. 省略类型检查: 同质dtype只需检查一次",
          "5. BLAS/LAPACK: 使用优化的数值库"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 36: SIMD与缓存的作用
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 6, "SIMD与CPU缓存的作用", "向量化性能的硬件原理")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "SIMD (Single Instruction, Multiple Data)",
         ["一条CPU指令同时处理多个数据",
          "AVX-256: 一次运算8个float32",
          "AVX-512: 一次运算16个float32",
          "NumPy内部利用SIMD指令",
          "Python循环中无法利用SIMD"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "CPU缓存效率",
         ["NumPy数组: 连续内存 -> L1/L2缓存命中率高",
          "Python列表: 指针间接访问 -> 缓存未命中频繁",
          "",
          "L1缓存: ~1ns (最快, 32~64KB)",
          "L2缓存: ~4ns (256KB~1MB)",
          "主内存: ~100ns (非常慢)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.5),
         "实战基准测试: 成对欧氏距离",
         ["200x300 成对距离计算结果:",
          "  Python 三重循环: ~15秒",
          "  NumPy 向量化: ~0.01秒",
          "  Numba JIT: ~0.02秒",
          "  Numba parallel: ~0.005秒",
          "向量化比循环快约500~1000倍!"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.5),
         "性能优化选择指南",
         ["第1选择: NumPy向量化 (最简单且最快)",
          "第2选择: Numba JIT (难以向量化的复杂循环)",
          "第3选择: Numba parallel (利用多核)",
          "",
          "能向量化就用NumPy!",
          "复杂条件分支/动态规划才用Numba"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 37: 向量化代码示例 - 欧氏距离
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 6, "向量化实战: 成对欧氏距离", "循环 vs 广播比较")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.8), [
    "# Python循环: O(n*m*d) 三重循环",
    "def dist_loop(X, Y):",
    "    n, d = X.shape; m = Y.shape[0]",
    "    D = np.zeros((n, m))",
    "    for i in range(n):",
    "        for j in range(m):",
    "            s = 0",
    "            for k in range(d):",
    "                s += (X[i,k] - Y[j,k])**2",
    "            D[i,j] = np.sqrt(s)",
    "    return D",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(5.0), Inches(6.0), Inches(2.2), [
    "# NumPy广播: 只需一行!",
    "def dist_vectorized(X, Y):",
    "    return np.sqrt(np.sum(",
    "        (X[:,np.newaxis,:] - Y[np.newaxis,:,:])**2,",
    "        axis=2))",
    "",
    "# KNN、K-Means等基于距离的算法必需!",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.2),
         "广播距离计算解析",
         ["X: shape (n, d)  ->  (n, 1, d)",
          "Y: shape (m, d)  ->  (1, m, d)",
          "",
          "X[:,np.newaxis,:] - Y[np.newaxis,:,:]",
          "  -> shape (n, m, d)  (差向量)",
          "",
          "**2  -> 平方",
          "np.sum(..., axis=2)  -> 求和 -> (n, m)",
          "np.sqrt(...)  -> 欧氏距离",
          "",
          "整个过程在C层面执行!",
          "内存: 需要O(n*m*d) (大数据注意)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# Section 7: 论文综述整合
# ============================================================
section_divider("论文综述整合", "Harris(2020), McKinney(2010), Walt(2011), Wickham(2014), Reback(2020)", 7, ACCENT_PURPLE)

# ============================================================
# 幻灯片 38: 论文概览
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 7, "第3章相关论文概览", "5篇核心论文及学术背景")

add_table_slide(s,
    headers=["#", "论文", "核心贡献", "引用数"],
    rows=[
        ["1", "McKinney (2010)", "Pandas DataFrame, Split-Apply-Combine", "10,000+"],
        ["2", "Walt et al. (2011)", "ndarray内部结构 (stride, dtype) 文档化", "5,000+"],
        ["3", "Harris et al. (2020)", "NumPy Nature论文, 数组编程范式", "8,000+"],
        ["4", "Reback et al. (2020)", "Pandas官方引用文档, 全部功能概览", "15,000+"],
        ["5", "Wickham (2014)", "Tidy Data 3原则, messy data 5种类型", "7,000+"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[0.5, 3.0, 5.5, 1.5],
    header_color=ACCENT_PURPLE, row_height=0.55, font_size=12)

add_card(s, Inches(0.6), Inches(5.2), Inches(12.0), Inches(2.0),
         "论文间关联关系 (按时间顺序)",
         ["McKinney (2010): 'ndarray对异质数据不够用' -> 开发Pandas",
          "Walt et al. (2011): ndarray技术设计详细文档化 (stride, ufunc, broadcast)",
          "Wickham (2014): 数据整理原则体系化 -> 影响了Pandas melt()/pivot_table()的设计",
          "Harris et al. (2020): NumPy 15年发展总结, 获Nature正式认可"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# 幻灯片 39: Harris (2020) - NumPy Nature论文
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 7, "Harris et al. (2020) -- NumPy Nature论文", '"Array programming with NumPy", Nature, 585, 357-362')

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "背景与意义",
         ["软件库论文在Nature上发表 = 证明NumPy的学术影响力",
          "总结了NumPy 15年的发展历程",
          "正式认可NumPy在现代科学计算中的核心角色",
          "正式定义了数组编程(Array Programming)范式"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "4个核心贡献",
         ["1. 数组编程范式: 用数组整体运算替代循环",
          "2. Python科学生态系统层级结构可视化",
          "3. NumPy API协议: __array_ufunc__, __array_function__",
          "4. GPU/分布式环境扩展: JAX, CuPy, Dask Array兼容"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.4),
         "核心引文与ML课程相关性",
         ['"数组编程是一种对整个数组应用运算而非对单个元素进行循环的编程范式。',
          ' 这种范式同时提升了代码的简洁性、可读性和性能。"',
          "",
          "ML课程相关性: 本课程使用的几乎所有库(scikit-learn, TensorFlow, PyTorch)都以NumPy为基础。",
          "本论文系统地解释了该基础设施的设计原则。"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 40: McKinney (2010) & Wickham (2014)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 7, "McKinney (2010) & Wickham (2014)", "Pandas的诞生与Tidy Data原则")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
         "McKinney (2010) -- Pandas的诞生",
         ['"Data Structures for Statistical Computing in Python"',
          "SciPy 2010",
          "",
          "背景: 在AQR Capital进行金融数据分析时",
          "发现Python缺少与R的data.frame对应的结构",
          "",
          "核心贡献:",
          "  1. 引入Series和DataFrame (基于索引的自动对齐)",
          "  2. 系统化的缺失值(NaN)处理API",
          "  3. Split-Apply-Combine模式 (groupby)",
          "  4. 关系数据连接 (merge, join)",
          "",
          '引文: "Python中没有与R的data.frame对应的',
          '高级数据结构。"'],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
         "Wickham (2014) -- Tidy Data原则",
         ['"Tidy Data", J. of Statistical Software, 59(10)',
          "",
          "背景: 数据科学家50~80%的工作是数据清洗",
          "缺乏系统理论 -> 原则体系化",
          "",
          "核心贡献:",
          "  1. 整洁数据的3原则 (变量=列, 观测=行, 观测单位=表)",
          "  2. 混乱数据的5种类型分类",
          "  3. 转换工具: melt(wide->long), pivot(long->wide)",
          "",
          "ML相关性:",
          '  scikit-learn的 "行=样本, 列=特征" 格式',
          "  = 与Tidy Data原则完全一致!",
          '引文: "整洁的数据可以复用相同的分析代码"'],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# Section 8: 实现代码详细解说
# ============================================================
section_divider("实现代码详细解说", "3个实验代码: NumPy线性代数、Pandas管道、向量化基准测试", 8, ACCENT_CYAN)

# ============================================================
# 幻灯片 41: 实验1 - NumPy线性代数
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "实验1: NumPy线性代数", "01_numpy_linear_algebra.py - SVD降维")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(5.0), [
    "# 使用SVD进行降维",
    "import numpy as np",
    "",
    "A = np.array([[1,2,0],[0,1,1],[2,0,1],",
    "              [1,1,1],[3,2,1]], dtype=float)",
    "",
    "U, s, Vt = np.linalg.svd(A, full_matrices=False)",
    "",
    "# 用前k个奇异值近似",
    "k = 2",
    "A_approx = U[:,:k] @ np.diag(s[:k]) @ Vt[:k,:]",
    "",
    "# 能量(方差解释力)分析",
    "energy = s**2 / np.sum(s**2) * 100",
    "print(f'奇异值: {s}')",
    "print(f'能量比例: {energy}')",
    "print(f'前{k}个累积: {np.sum(energy[:k]):.1f}%')",
], font_size=11)

add_card(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(5.0),
         "代码解析",
         ["SVD: 将A分解为 U * Sigma * V^T",
          "",
          "奇异值(s)的大小 = 该方向的数据方差量",
          "只使用前k个 = 最小化信息损失 + 降维",
          "-> 这就是PCA的数学基础!",
          "",
          "full_matrices=False: 截断SVD",
          "  U: (m, k), s: (k,), Vt: (k, n)",
          "",
          "能量比例 = s^2 / sum(s^2) * 100",
          "  每个奇异值解释总方差的百分比",
          "",
          "实践: 选择使累积能量>95%的k"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# 幻灯片 42: 实验1 (续) - 正规方程
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "实验1 (续): 正规方程求解比较", "逆矩阵 vs solve vs lstsq")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(4.5), [
    "# 数据生成: y = 3x + 2 + noise",
    "np.random.seed(42)",
    "X_raw = np.random.uniform(0, 10, (50, 1))",
    "y = 3 * X_raw + 2 + np.random.normal(0, 1, (50, 1))",
    "X = np.hstack([np.ones((50, 1)), X_raw])",
    "",
    "# 方法1: 逆矩阵 (数值不稳定)",
    "beta_inv = np.linalg.inv(X.T @ X) @ X.T @ y",
    "",
    "# 方法2: solve (基于LU分解, 稳定)",
    "beta_solve = np.linalg.solve(X.T @ X, X.T @ y)",
    "",
    "# 方法3: lstsq (基于SVD, 最稳定, 推荐!)",
    "beta_lstsq, _, _, _ = np.linalg.lstsq(X, y, rcond=None)",
    "",
    "print(f'真值: 截距=2, 斜率=3')",
    "print(f'lstsq: 截距={beta_lstsq[0,0]:.4f}, 斜率={beta_lstsq[1,0]:.4f}')",
], font_size=11)

add_card(s, Inches(8.0), Inches(2.0), Inches(4.6), Inches(5.0),
         "建议",
         ["实际中始终使用np.linalg.lstsq()!",
          "",
          "逆矩阵计算的问题:",
          "  - 条件数(condition number)大的矩阵",
          "    数值不稳定",
          "  - 无法处理奇异矩阵(singular matrix)",
          "",
          "lstsq的优点:",
          "  - 基于SVD最稳定",
          "  - 可处理秩不足(rank-deficient)矩阵",
          "  - 同时返回残差、秩、奇异值",
          "",
          "scikit-learn的LinearRegression也",
          "内部使用lstsq!"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 幻灯片 43: 实验2 - Pandas数据管道
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "实验2: Pandas数据管道", "02_pandas_data_pipeline.py - 缺失值策略比较")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(5.0), [
    "import pandas as pd; import numpy as np",
    "np.random.seed(42); n = 500",
    "df = pd.DataFrame({",
    "    '日期': pd.date_range('2024-01-01', periods=n, freq='6H'),",
    "    '客户ID': np.random.choice(['C001','C002','C003'], n),",
    "    '金额': np.random.exponential(50000, n).astype(int)+5000,",
    "    '数量': np.random.randint(1, 10, size=n)})",
    "",
    "# 人工插入缺失值 (10%)",
    "mask = np.random.random(n) < 0.1",
    "df.loc[mask, '金额'] = np.nan",
    "",
    "# 5种策略比较",
    "strategies = {",
    "    '删除': df['金额'].dropna(),",
    "    '均值': df['金额'].fillna(df['金额'].mean()),",
    "    '中位数': df['金额'].fillna(df['金额'].median()),",
    "    '0替代': df['金额'].fillna(0),",
    "    '插值': df['金额'].interpolate()}",
], font_size=10)

add_card(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(5.0),
         "代码解析与结果",
         ["用均值和标准差比较各策略的影响:",
          "",
          "删除: 保持原始统计量 (数据损失)",
          "均值: 均值相同, 低估方差",
          "中位数: 对异常值鲁棒, 低估方差",
          "0替代: 均值和方差都扭曲 (不推荐)",
          "插值: 在时间序列中保持连续性",
          "",
          "GroupBy + Transform模式:",
          "  分组Z-score归一化",
          "  分组累积求和 (cumsum)",
          "  用组中位数替代缺失值",
          "",
          "最推荐: 根据情况选择合适的策略!"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 44: 实验2 (续) - GroupBy Transform
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "实验2 (续): GroupBy + Transform模式", "分组归一化与缺失值处理")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.5), [
    "# 按客户分组Z-score归一化",
    "df['金额_zscore'] = df.groupby('客户ID')['金额'].transform(",
    "    lambda x: (x - x.mean()) / x.std())",
    "",
    "# 按客户分组累积求和",
    "df['累积金额'] = df.groupby('客户ID')['金额'].cumsum()",
], font_size=12)

add_code_block(s, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.5), [
    "# 用组中位数替代缺失值",
    "df['金额_filled'] = df.groupby('客户ID')['金额'].transform(",
    "    lambda x: x.fillna(x.median()))",
    "",
    "# 这个模式是ML预处理中最强大的工具!",
    "# 分组替代比整体均值更准确",
], font_size=12)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.2),
         "transform的核心特性",
         ["transform vs agg 区别:",
          "  agg: 组 -> 标量 (行数减少)",
          "  transform: 组 -> 原始大小 (行数不变)",
          "",
          "transform应用模式:",
          "  1. 分组Z-score归一化",
          "  2. 分组比例计算",
          "  3. 分组缺失值替代",
          "  4. 分组累积和/移动平均",
          "",
          "为什么有用?",
          "  可直接作为新列添加到原始DataFrame",
          "  分组统计比整体统计更准确的替代",
          "  ML中特征工程的核心工具"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 45: 实验3 - 向量化基准测试
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "实验3: 向量化基准测试", "03_vectorization_benchmark.py - 距离计算性能比较")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(5.0), [
    "def dist_loop(X, Y):",
    "    n, d = X.shape; m = Y.shape[0]",
    "    D = np.zeros((n, m))",
    "    for i in range(n):",
    "        for j in range(m):",
    "            s = 0",
    "            for k in range(d):",
    "                s += (X[i,k]-Y[j,k])**2",
    "            D[i,j] = np.sqrt(s)",
    "    return D",
    "",
    "def dist_vectorized(X, Y):",
    "    return np.sqrt(np.sum(",
    "        (X[:,np.newaxis,:]-Y[np.newaxis,:,:])**2, axis=2))",
    "",
    "X = np.random.randn(200, 10)",
    "Y = np.random.randn(300, 10)",
    "# 循环: ~15秒 / 向量化: ~0.01秒 / 速度: ~1000倍",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.0),
         "基准测试结果及解析",
         ["200x300 成对距离 (10维):",
          "",
          "方法            | 时间      | 速度比",
          "Python循环      | ~15秒    | 1x",
          "NumPy向量化     | ~0.01秒  | 1,500x",
          "Numba JIT       | ~0.02秒  | 750x",
          "Numba parallel  | ~0.005秒 | 3,000x",
          "",
          "结果一致: np.allclose(D_loop, D_vec) = True",
          "",
          "KNN、K-Means等基于距离的算法中",
          "向量化不是选择而是必须!",
          "scikit-learn内部也使用向量化"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# Section 9: 核心总结
# ============================================================
section_divider("核心总结与复习", "第3章核心概念整理与复习题", 9, ACCENT_BLUE)

# ============================================================
# 幻灯片 46: 核心总结表
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "核心总结表", "第3章学习的所有概念整理")

add_table_slide(s,
    headers=["概念", "核心要点"],
    rows=[
        ["ndarray", "连续内存, stride, 同质dtype, 向量化运算的基础"],
        ["广播", "不同大小数组间运算自动化, stride=0虚拟扩展"],
        ["DataFrame", "异质表结构, 按列独立数组, 基于标签索引"],
        ["方法链", ".assign().query().sort_values() 模式"],
        ["GroupBy", "分割-应用-合并, agg / transform / apply"],
        ["Tidy Data", "变量=列, 观测=行, 观测单位=表 (Wickham 2014)"],
        ["向量化", "C循环, 连续内存, SIMD带来100~1000倍性能提升"],
        ["缺失值策略", "删除 / 均值 / 中位数 / 插值 / 分组 -- 根据情况选择"],
        ["缩放", "Standard / MinMax / Robust -- 只在train上fit, 注意Leakage"],
        ["SVD", "A = U Sigma V^T, PCA的数学基础, 降维的核心"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 9.5],
    header_color=ACCENT_BLUE, row_height=0.45, font_size=12)

# ============================================================
# 幻灯片 47: 公式整理
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "第3章核心公式整理", "NumPy与ML中使用的主要公式")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.3),
         "线性代数分解",
         ["SVD: A = U * Sigma * V^T",
          "特征值分解: Av = lambda * v",
          "LU分解: A = P * L * U",
          "QR分解: A = Q * R, Q^T Q = I"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.3),
         "正规方程与回归",
         ["正规方程: beta = (X^T X)^{-1} X^T y",
          "最小二乘法: min ||Xbeta - y||^2",
          "条件数: kappa(A) = ||A|| * ||A^{-1}||",
          "推荐: np.linalg.lstsq(X, y)"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.3),
         "特征缩放",
         ["StandardScaler: z = (x - mu) / sigma",
          "MinMaxScaler: z = (x - x_min) / (x_max - x_min)",
          "RobustScaler: z = (x - Q2) / (Q3 - Q1)",
          "Data Leakage: 只在train上fit, test上只transform!"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.3),
         "异常值检测",
         ["IQR: Q1 - 1.5*IQR < x < Q3 + 1.5*IQR",
          "Z-score: |z| = |(x - mu)/sigma| > 3",
          "方差解释力: s^2 / sum(s^2) * 100",
          "选择使累积能量>95%的k"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 48: 大规模数据处理
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "大规模数据处理工具", "超越Pandas的局限")

add_table_slide(s,
    headers=["库", "特点", "适用场景"],
    rows=[
        ["Dask", "Pandas API兼容, 延迟执行, 分布式处理", "中大规模 (数~数百GB)"],
        ["Polars", "基于Rust, 延迟执行, 非常快", "单机大规模"],
        ["Vaex", "Out-of-core, 内存映射, 10亿行", "探索性分析"],
        ["PySpark", "集群分布式处理", "超大规模 (TB以上)"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.0, 5.0, 4.0],
    header_color=ACCENT_BLUE, row_height=0.5, font_size=12)

add_code_block(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.5), [
    "# Dask示例 -- 与Pandas API几乎相同!",
    "import dask.dataframe as dd",
    "ddf = dd.read_csv('large_*.csv')",
    "result = ddf.groupby('col').mean().compute()",
    "",
    "# np.einsum: 爱因斯坦求和约定",
    "C = np.einsum('ik,kj->ij', A, B)  # == A @ B",
    "# 批量矩阵乘法 (深度学习注意力)",
    "scores = np.einsum('hid,hjd->hij', Q, K)",
], font_size=11)

add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.5),
         "实际数据管道: Titanic示例",
         ["1. pd.read_csv() -> DataFrame加载",
          "2. groupby + transform分组替代缺失值",
          "3. map()进行分类编码",
          "4. .values转换为NumPy",
          "5. train_test_split进行分割",
          "6. StandardScaler: 只在train上fit!",
          "7. LogisticRegression训练和评估"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 49: np.einsum与Numba
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "高级工具: np.einsum与Numba JIT", "张量运算与JIT编译加速")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "# np.einsum: 爱因斯坦求和约定",
    "A = np.random.randn(3, 4)",
    "B = np.random.randn(4, 5)",
    "",
    "# 矩阵乘法: C_ij = sum_k A_ik * B_kj",
    "C = np.einsum('ik,kj->ij', A, B)",
    "",
    "# 迹(Trace): sum_i A_ii",
    "trace = np.einsum('ii->', np.eye(3))",
    "",
    "# 批量矩阵乘法 (深度学习注意力核心!)",
    "Q = np.random.randn(8, 64, 32)  # (heads, seq, d_k)",
    "K = np.random.randn(8, 64, 32)",
    "scores = np.einsum('hid,hjd->hij', Q, K)",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(5.2), Inches(6.0), Inches(2.0), [
    "# Numba JIT: 将Python循环加速到C速度",
    "from numba import jit",
    "@jit(nopython=True)",
    "def fast_func(X, Y):",
    "    # 用于难以向量化的复杂循环",
    "    ...",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "np.einsum应用案例",
         ["矩阵乘法: 'ik,kj->ij' == A @ B",
          "迹: 'ii->' == np.trace(A)",
          "外积: 'i,j->ij' == np.outer(a, b)",
          "批量矩阵乘法: 'bij,bjk->bik'",
          "Self-Attention: 'hid,hjd->hij'",
          "可读性好且内部优化也很好!"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.4),
         "Numba选择指南",
         ["能向量化 -> NumPy最好",
          "难以向量化的复杂循环 -> Numba JIT",
          "需要并行化 -> Numba parallel (prange)",
          "",
          "Numba JIT: 首次调用时编译 (warm-up)",
          "之后调用达到C级速度!"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 幻灯片 50: 复习题 (1/3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "复习题 (1/3)", "NumPy相关问题")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(5.2),
         "NumPy核心问题 1~4",
         ["Q1. 说明NumPy ndarray的stride概念，并解释转置(transpose)运算时stride如何变化。",
          "    -> stride = 每个维度到下一个元素的字节数。转置只交换stride (无数据复制!)",
          "",
          "Q2. 描述广播的三条规则，并推导(3,1) + (1,4)运算的结果shape。",
          "    -> 规则1: 对齐维度, 规则2: 检查兼容性(3vs1, 1vs4), 规则3: 虚拟扩展 -> 结果 (3,4)",
          "",
          "Q3. 从内存布局角度说明C-order和F-order的区别，",
          "    并说明行方向遍历时哪个更高效及其原因。",
          "    -> C-order: 行优先存储, 行方向遍历时连续内存访问使缓存效率高",
          "",
          "Q4. 说明Harris et al.(2020)的Nature论文中提出的Python科学生态系统层级结构，",
          "    并解释NumPy位于基础层的原因。",
          "    -> NumPy是Pandas、SciPy、scikit-learn等的基础。ndarray提供标准数组接口。"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 幻灯片 51: 复习题 (2/3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "复习题 (2/3)", "Pandas及数据预处理相关问题")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(5.2),
         "Pandas/预处理核心问题 5~7",
         ["Q5. 说明Pandas的loc和iloc的区别，并解释切片时末尾索引包含与否不同的原因。",
          "    -> loc: 基于标签 (包含末尾) / iloc: 基于整数位置 (不包含末尾, Python惯例)",
          "",
          "Q6. 说明McKinney(2010)开发Pandas的动机，并描述NumPy ndarray不适合表格数据的原因。",
          "    -> Python缺少与R的data.frame对应的结构。ndarray只支持同质类型，",
          "       对于字符串+数字+日期混合的异质表格数据处理不足。",
          "",
          "Q7. 比较缺失值处理5种策略的优缺点，并说明对于时间序列数据",
          "    最合适的策略及其理由。",
          "    -> 时间序列: 插值(interpolate)最合适。保持时间连续性，反映前后值的趋势。",
          "       删除行会断开时间序列，均值替代忽略了时间模式。"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 52: 复习题 (3/3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "复习题 (3/3)", "Tidy Data、向量化、Data Leakage相关问题")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(5.2),
         "进阶问题 8~10",
         ["Q8. 描述Wickham(2014)的Tidy Data 3原则，并编写将列标题为值的messy data",
          "    转换为tidy data的Pandas代码。",
          "    -> 原则: 变量=列, 观测=行, 观测单位=表。pd.melt(messy, id_vars='姓名', ...)",
          "",
          "Q9. 从硬件层面(SIMD, CPU缓存)解释向量化运算比Python循环",
          "    快100~1000倍的原因。",
          "    -> SIMD: 一条指令同时处理多个数据 (AVX-256: 8个float32)",
          "       CPU缓存: 连续内存访问最大化L1/L2缓存命中率",
          "       + C编译代码, 省略类型检查, GIL释放",
          "",
          "Q10. 什么是Data Leakage，说明错误使用StandardScaler时如何发生。",
          "    -> 测试数据的信息泄露到训练中的现象。",
          "       错误: scaler.fit_transform(X_all) -> 测试的均值/标准差反映到训练中",
          "       正确: scaler.fit_transform(X_train) 后 scaler.transform(X_test)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 53: 第3章核心关键词总结
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "第3章核心关键词总结", "一页看完第3章的全部内容")

add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(2.2),
         "NumPy",
         ["ndarray, stride, dtype",
          "C-order/F-order",
          "广播3条规则",
          "ufunc, 向量化",
          "SVD, 特征值分解, 正规方程"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(2.2),
         "Pandas",
         ["DataFrame, Series",
          "loc/iloc, query",
          "方法链, pipe()",
          "GroupBy: agg/transform/apply",
          "merge/join, pivot/melt"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(2.2),
         "数据预处理",
         ["缺失值5种策略",
          "异常值: IQR, Z-score",
          "缩放: Standard/MinMax/Robust",
          "编码: Label/OneHot/Ordinal",
          "Data Leakage防止"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.5), Inches(3.8), Inches(2.2),
         "Tidy Data",
         ["Wickham 2014",
          "变量=列, 观测=行",
          "观测单位=表",
          "melt(wide->long)",
          "pivot(long->wide)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(4.7), Inches(4.5), Inches(3.8), Inches(2.2),
         "向量化 vs 循环",
         ["100~1000倍速度提升",
          "SIMD, CPU缓存",
          "C编译, 省略类型检查",
          "GIL释放, BLAS/LAPACK",
          "Numba JIT替代方案"],
         title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(8.8), Inches(4.5), Inches(3.8), Inches(2.2),
         "5篇论文",
         ["McKinney (2010): Pandas",
          "Walt et al. (2011): ndarray",
          "Harris et al. (2020): NumPy Nature",
          "Reback et al. (2020): Pandas官方",
          "Wickham (2014): Tidy Data"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 幻灯片 54: View vs Copy / 内存优化整理
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "内存相关整理: View/Copy & 优化", "NumPy与Pandas的内存效率策略")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "NumPy View vs Copy整理",
         ["View (共享内存): 切片, 转置, reshape, ravel",
          "Copy (独立内存): 花式索引, 布尔索引, flatten, copy()",
          "检查: np.shares_memory(a, b)",
          "注意: 修改View时原始数据也会改变!"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "Pandas内存优化",
         ["降级: int64->int32 (节约50~75%)",
          "Categorical: object->category (节约90~99%)",
          "Sparse: 大部分为0的数据 (节约90%+)",
          "分块读取: chunksize参数 (在RAM限制内)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.4),
         "np.memmap: 内存映射文件",
         ["处理大于RAM的数据时使用",
          "fp = np.memmap('data.dat', dtype='float32', mode='w+', shape=(1000000, 100))",
          "只将需要的部分加载到内存: chunk = fp_read[500:600]",
          "scikit-learn的部分算法也支持memmap"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 幻灯片 55: 实践提示 & 下一章预览
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "实践提示 & 下一章预览", "第3章学到的内容如何在第4章之后使用")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "实践提示 Top 5",
         ["1. 始终使用np.linalg.lstsq() (不推荐逆矩阵)",
          "2. StandardScaler: 只在train上fit! (防止Data Leakage)",
          "3. 方法链 + pipe()构建易读的管道",
          "4. GroupBy transform进行分组预处理 (比整体均值更准确)",
          "5. 优先向量化, 无法向量化时才用Numba JIT"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "第3章的位置: 在ML管道中",
         ["第1章: ML概述 (定义, 学习类型)",
          "第2章: ML项目 (完整管道)",
          "第3章: 数据处理 (NumPy + Pandas) [当前]",
          "第4章~: 具体ML算法",
          "",
          "第3章是所有后续章节的数据处理基础!"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.4),
         "参考文献",
         ["1. McKinney, W. (2010). 'Data Structures for Statistical Computing in Python.' SciPy 2010",
          "2. Walt, S. et al. (2011). 'The NumPy Array.' Computing in Science & Engineering, 13(2)",
          "3. Harris, C.R. et al. (2020). 'Array programming with NumPy.' Nature, 585, 357-362",
          "4. Reback, J. et al. (2020). 'pandas-dev/pandas: Pandas.' Zenodo",
          "5. Wickham, H. (2014). 'Tidy Data.' Journal of Statistical Software, 59(10)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 幻灯片 56: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.2), prs.slide_width, Inches(0.5),
         "THANK YOU", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(3.0), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(3.3), prs.slide_width, Inches(0.6),
         "第3章: Pandas与NumPy (扩展详细版)", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.4),
         "NumPy ndarray | Pandas DataFrame | Broadcasting | Vectorization | Tidy Data | Linear Algebra",
         font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.0), prs.slide_width, Inches(0.4),
         "如有任何问题，请随时提问!", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# PPT保存
# ============================================================
save_path = os.path.join(os.path.dirname(__file__), "3장_판다스와_넘파이_강의PPT_확장_china.pptx")
prs.save(save_path)
print(f"PPT saved: {save_path}")
print(f"Total slides: {len(prs.slides)}")
