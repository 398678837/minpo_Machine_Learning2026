"""3장 판다스와 넘파이 - 확장 강의 PPT 생성 스크립트 (상세 버전)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)


def add_bg(slide, color=DARK_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name; p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = '맑은 고딕'; p.space_after = spacing; p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, radius=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN; p.font.name = 'Consolas'; p.space_after = Pt(2)

def add_table_slide(slide, headers, rows, left, top, col_widths, header_color=ACCENT_BLUE,
                    row_height=0.5, font_size=13, header_font_size=14):
    """테이블을 카드 형태로 그리는 헬퍼"""
    total_w = sum(col_widths)
    # 헤더
    cx = left
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        add_shape(slide, cx, top, Inches(w), Inches(0.5), header_color)
        add_text(slide, cx, top, Inches(w), Inches(0.5),
                 h, font_size=header_font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += Inches(w)
    # 행
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + Inches(row_height) * i
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
        cx = left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_shape(slide, cx, y, Inches(w), Inches(row_height), bg)
            fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
            add_text(slide, cx, y, Inches(w), Inches(row_height),
                     cell, font_size=font_size, color=fc, bold=(j==0), align=PP_ALIGN.CENTER)
            cx += Inches(w)

def section_divider(title, subtitle, section_num, accent=ACCENT_BLUE):
    """섹션 구분 슬라이드"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, SECTION_BG)
    add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), accent)
    add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), accent)
    add_text(s, Inches(0), Inches(2.0), prs.slide_width, Inches(0.5),
             f"SECTION {section_num}", font_size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_accent_line(s, Inches(5.5), Inches(2.7), Inches(2.3), accent)
    add_text(s, Inches(0), Inches(3.0), prs.slide_width, Inches(1.0),
             title, font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.5),
             subtitle, font_size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    return s


# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "기계학습 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "3장: 판다스와 넘파이", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "Pandas & NumPy for Machine Learning", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.3), Inches(11), Inches(0.8),
         "[ 확장 상세 버전 ]", font_size=20, color=ACCENT_ORANGE, bold=True)
add_text(s, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
         "핵심 키워드: NumPy ndarray · Pandas DataFrame · 브로드캐스팅 · 벡터화 · Tidy Data · 선형대수 · 데이터 전처리",
         font_size=14, color=DARK_GRAY)
add_text(s, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
         "대학원 박사과정 수준 | 기계학습 3장", font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 2: 목차 (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (Table of Contents) - 1/2")

toc_left = [
    "Section 1. 도입: 데이터 없이는 ML도 없다",
    "    1.1 Python 과학 생태계의 레이어 구조",
    "    1.2 ML 파이프라인에서의 데이터 흐름",
    "    1.3 왜 파이썬인가?",
    "",
    "Section 2. NumPy 심화: ndarray의 세계",
    "    2.1 ndarray 내부 구조 (4가지 핵심 요소)",
    "    2.2 메모리 레이아웃과 Stride",
    "    2.3 C-order vs F-order / dtype 시스템",
    "    2.4 배열 생성 함수 총정리",
    "    2.5 배열 인덱싱과 슬라이싱",
    "    2.6 브로드캐스팅 3가지 규칙",
    "    2.7 브로드캐스팅 코드 예제",
    "    2.8 유니버설 함수 (ufunc)",
    "    2.9 선형대수 기본 / SVD / 고유값 분해",
    "    2.10 정규방정식과 View vs Copy",
]
add_bullet_list(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5.0),
                toc_left, font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# ============================================================
# 슬라이드 3: 목차 (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (Table of Contents) - 2/2")

toc_right = [
    "Section 3. Pandas 심화: DataFrame의 세계",
    "    3.1 DataFrame 구조 / 인덱싱 (loc vs iloc)",
    "    3.2 메서드 체이닝 / GroupBy / Merge·Join",
    "    3.3 시계열 윈도우 함수 / Apply·Map",
    "",
    "Section 4. 데이터 전처리",
    "    4.1 결측치 처리 5가지 / 이상치 탐지 3가지",
    "    4.2 피처 스케일링 3가지 / 인코딩",
    "",
    "Section 5. Tidy Data 개념 (Wickham 2014)",
    "Section 6. 벡터화 vs 루프: 성능의 과학",
    "Section 7. 논문 리뷰 통합",
    "Section 8. 구현 코드 상세 해설 (3개 실습)",
    "Section 9. 핵심 요약 및 복습 질문",
]
add_bullet_list(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5.0),
                toc_right, font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# ============================================================
# 슬라이드 4: 학습 목표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "학습 목표 (Learning Objectives)")

objectives = [
    "1. NumPy ndarray의 내부 메모리 구조(stride, dtype, C/F-order)를 설명하고, 성능 최적화에 활용할 수 있다.",
    "2. 브로드캐스팅 규칙을 정확히 적용하여 다양한 형태의 배열 간 연산을 수행할 수 있다.",
    "3. LU/QR/SVD/고유값 분해를 NumPy로 구현하고, ML 알고리즘(PCA, 선형회귀)에서의 활용을 설명할 수 있다.",
    "4. Pandas DataFrame의 내부 구조를 이해하고, method chaining과 pipe()로 데이터 파이프라인을 구성할 수 있다.",
    "5. 결측치 처리 5가지, 이상치 탐지 3가지, 피처 스케일링 3가지를 비교 분석할 수 있다.",
    "6. Wickham의 Tidy Data 원칙을 적용하여 messy data를 tidy data로 변환할 수 있다.",
    "7. 벡터화 연산과 Python 루프의 성능 차이를 하드웨어 수준(SIMD, 캐시)에서 설명할 수 있다.",
]
add_bullet_list(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5.0),
                objectives, font_size=15, color=LIGHT_GRAY, spacing=Pt(8))

# ============================================================
# Section 1: 도입
# ============================================================
section_divider("도입: 데이터 없이는 ML도 없다", "Python 생태계와 ML 파이프라인에서의 NumPy·Pandas 역할", 1, ACCENT_BLUE)

# ============================================================
# 슬라이드 5: Python 과학 생태계의 레이어 구조
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 1, "Python 과학 생태계의 레이어 구조", "Harris et al. (2020) Nature 논문 기반")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.2),
         "Layer 4: 응용 라이브러리 (Applications)",
         ["scikit-learn, TensorFlow, PyTorch, Keras"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(0.6), Inches(3.3), Inches(5.5), Inches(1.3),
         "Layer 3: 도메인별 라이브러리 (Domain-specific)",
         ["Pandas (표 형식 데이터)", "SciPy (과학 알고리즘), matplotlib (시각화)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.7), Inches(5.5), Inches(1.3),
         "Layer 2: NumPy (기반 인프라)",
         ["ndarray, ufunc, broadcasting", "linear algebra, FFT, random"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(0.6), Inches(6.1), Inches(5.5), Inches(0.9),
         "Layer 1: Python + C Extensions",
         ["BLAS, LAPACK, C/Fortran 컴파일 코드"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(6.5), Inches(2.0), Inches(6.2), Inches(5.0),
         "핵심 포인트",
         ["NumPy는 생태계의 기반 인프라(foundation infrastructure)로 기능",
          "Pandas DataFrame은 내부적으로 NumPy 배열을 감싸고(wrapping) 있음",
          "scikit-learn의 fit()/predict()는 NumPy 배열을 입력으로 받음",
          "TensorFlow/PyTorch 텐서도 NumPy 배열과 상호 변환 기본 제공",
          "",
          "Andrew Ng: '데이터가 왕이다 (Data is King)'",
          "ML 프로젝트에서 데이터 처리에 50~80% 시간 소비 (Wickham, 2014)"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 슬라이드 6: ML 파이프라인에서의 데이터 흐름
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 1, "ML 파이프라인에서의 데이터 흐름", "Pandas -> NumPy -> scikit-learn 데이터 변환 과정")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(5.0), [
    "import pandas as pd",
    "import numpy as np",
    "from sklearn.model_selection import train_test_split",
    "from sklearn.preprocessing import StandardScaler",
    "from sklearn.linear_model import LogisticRegression",
    "",
    "# 1단계: Pandas로 데이터 로딩",
    "df = pd.read_csv('data.csv')",
    "",
    "# 2단계: Pandas로 탐색 및 전처리",
    "df = df.dropna()",
    "df['new_feature'] = df['feature_a'] * df['feature_b']",
    "",
    "# 3단계: NumPy 배열로 변환",
    "X = df[['feature_a','feature_b','new_feature']].to_numpy()",
    "y = df['target'].to_numpy()",
    "",
    "# 4단계: scikit-learn으로 모델링 (내부적으로 NumPy)",
    "X_train, X_test, y_train, y_test = train_test_split(X, y)",
    "scaler = StandardScaler()",
    "X_train_scaled = scaler.fit_transform(X_train)",
    "model = LogisticRegression()",
    "model.fit(X_train_scaled, y_train)",
], font_size=11)

add_card(s, Inches(8.5), Inches(2.0), Inches(4.2), Inches(5.0),
         "파이프라인 흐름 요약",
         ["데이터 수집 -> 데이터 탐색(EDA)",
          "-> 데이터 전처리",
          "-> 피처 엔지니어링",
          "-> 모델 학습",
          "-> 모델 평가 -> 배포",
          "",
          "Pandas: 로딩, 탐색, 전처리 담당",
          "NumPy: 수치 연산의 기반 제공",
          "scikit-learn: 모델링 (내부 NumPy)",
          "",
          "실무에서 모델 학습 이전의 모든 단계가",
          "데이터 처리에 해당 (전체 시간의 50-80%)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 7: 왜 파이썬인가?
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 1, "왜 파이썬인가?", "접착 언어(Glue Language)로서의 역할")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "파이썬의 핵심 전략",
         ["파이썬 자체는 느리지만, 성능이 중요한 부분은",
          "C/Fortran으로 작성된 라이브러리(BLAS, LAPACK)가 처리",
          "NumPy는 이 전략의 대표적 성공 사례",
          "접착 언어(glue language): 다양한 라이브러리를 연결하는 역할"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "McKinney (2010)의 Pandas 개발 동기",
         ["'파이썬에는 R의 data.frame에 대응하는",
          "  고수준 데이터 구조가 없었다.'",
          "NumPy ndarray: 동질적 수치 데이터에 탁월",
          "하지만 이질적 표 형식 데이터는 부족",
          "-> 이 간극을 메우기 위해 Pandas 탄생!"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.2),
         "NumPy + Pandas의 역할 분담",
         ["NumPy: 동질적(homogeneous) 수치 데이터의 고속 연산 -- ndarray는 연속 메모리에 동일 타입 데이터 저장",
          "Pandas: 이질적(heterogeneous) 표 형식 데이터 -- DataFrame은 열마다 다른 dtype 허용 (문자열+숫자+날짜)",
          "이 조합이 파이썬을 데이터 과학의 지배적 언어로 만드는 데 결정적 역할"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# Section 2: NumPy 심화
# ============================================================
section_divider("NumPy 심화: ndarray의 세계", "ndarray 내부 구조, 메모리 레이아웃, 브로드캐스팅, 선형대수", 2, ACCENT_CYAN)

# ============================================================
# 슬라이드 8: ndarray 내부 구조 - 4가지 핵심 요소
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "ndarray 내부 구조: 4가지 핵심 요소", "Walt et al. (2011) 기반")

add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(1.8),
         "data (데이터 포인터)",
         ["데이터가 저장된 메모리 버퍼의 포인터",
          "실제 수치 데이터가 연속 메모리에 저장",
          "Python 리스트와의 근본적 차이"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(1.8),
         "dtype (데이터 타입)",
         ["각 원소의 데이터 타입 (float64, int32 등)",
          "모든 원소가 동일한 타입 (동질적)",
          "메모리 크기와 연산 방식 결정"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(1.8),
         "shape (형태)",
         ["각 차원의 크기를 나타내는 튜플",
          "예: (2, 3) -> 2행 3열",
          "ndim: 차원 수, size: 전체 원소 수"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.0), Inches(5.5), Inches(1.5),
         "strides (보폭)",
         ["각 차원에서 다음 원소까지의 바이트 수",
          "예: (24, 8) -> 행 이동 시 24바이트, 열 이동 시 8바이트",
          "전치(transpose)는 stride만 교환 (데이터 복사 없음!)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_code_block(s, Inches(6.5), Inches(4.0), Inches(6.2), Inches(3.2), [
    "import numpy as np",
    "",
    "arr = np.array([[1.0, 2.0, 3.0],",
    "                [4.0, 5.0, 6.0]])",
    "",
    "print(f'dtype:   {arr.dtype}')    # float64",
    "print(f'shape:   {arr.shape}')    # (2, 3)",
    "print(f'strides: {arr.strides}')  # (24, 8)",
    "print(f'nbytes:  {arr.nbytes}')   # 48 = 2*3*8",
    "print(f'ndim:    {arr.ndim}')     # 2",
    "print(f'size:    {arr.size}')     # 6",
], font_size=11)

# ============================================================
# 슬라이드 9: 메모리 레이아웃과 Stride
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "메모리 레이아웃과 Stride", "ndarray vs Python 리스트의 근본적 차이")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.0),
         "Python 리스트의 메모리 구조",
         ["리스트 객체 --> [포인터1, 포인터2, ...]",
          "각 포인터 -> 독립적인 PyObject (28+ bytes each)",
          "불연속 메모리: CPU 캐시 미스 빈번",
          "매 원소 접근 시 타입 체크 오버헤드"],
         title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.0),
         "NumPy 배열의 메모리 구조",
         ["ndarray 객체 --> [8B | 8B | 8B | ...]",
          "연속 메모리 블록에 동질적 타입 데이터 저장",
          "CPU 캐시 적중률(hit rate) 극대화",
          "SIMD(단일 명령 다중 데이터) 활용 가능"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8),
         "Stride 시각적 설명 (2x3 float64, C-order)",
         ["메모리: [1.0][2.0][3.0][4.0][5.0][6.0]",
          "바이트:  0    8   16   24   32   40",
          "",
          "shape = (2, 3)",
          "strides = (24, 8)",
          "  행 방향(axis=0): 24바이트 = 3원소 x 8B/원소",
          "  열 방향(axis=1):  8바이트 = 1원소 x 8B/원소"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_code_block(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.8), [
    "a = np.array([[1, 2, 3, 4],",
    "              [5, 6, 7, 8],",
    "              [9, 10, 11, 12]], dtype=np.float64)",
    "",
    "print(f'shape: {a.shape}')      # (3, 4)",
    "print(f'strides: {a.strides}')  # (32, 8)",
    "",
    "# 전치(transpose)는 stride만 바꾼다",
    "# (데이터 복사 없음!)",
    "b = a.T",
    "print(f'전치 shape: {b.shape}')    # (4, 3)",
    "print(f'전치 strides: {b.strides}')# (8, 32)",
], font_size=11)

# ============================================================
# 슬라이드 10: C-order vs F-order, dtype 시스템
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "C-order vs F-order / dtype 시스템", "메모리 순서와 데이터 타입 체계")

add_table_slide(s,
    headers=["순서", "메모리 레이아웃", "빠른 접근 방향", "사용 예"],
    rows=[
        ["C-order (row-major)", "행 우선 저장", "마지막 축 (열 방향)", "NumPy 기본, C"],
        ["F-order (col-major)", "열 우선 저장", "첫번째 축 (행 방향)", "Fortran, MATLAB, R"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 2.5, 2.5, 2.5],
    header_color=ACCENT_BLUE)

add_table_slide(s,
    headers=["카테고리", "dtype", "바이트", "범위/용도"],
    rows=[
        ["정수", "int8 / int16", "1 / 2", "-128~127 / -32768~32767"],
        ["정수", "int32 / int64", "4 / 8", "일반 정수 / 큰 정수"],
        ["부호없는 정수", "uint8", "1", "0~255 (이미지 픽셀)"],
        ["부동소수점", "float16 / float32", "2 / 4", "딥러닝 추론 / 학습"],
        ["부동소수점", "float64", "8", "과학 계산 기본값"],
        ["복소수/불리언", "complex128 / bool_", "16 / 1", "신호처리 / True-False"],
    ],
    left=Inches(0.6), top=Inches(3.5), col_widths=[2.5, 2.5, 2.0, 3.0],
    header_color=ACCENT_PURPLE, row_height=0.4, font_size=12)

add_text(s, Inches(0.6), Inches(6.7), Inches(11), Inches(0.4),
         "성능 시사점: 행 방향 순회 시 C-order가, 열 방향 순회 시 F-order가 CPU 캐시 효율이 높다.",
         font_size=14, color=ACCENT_ORANGE, bold=True)

# ============================================================
# 슬라이드 11: 배열 생성 함수 총정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "배열 생성 함수 총정리", "다양한 NumPy 배열 생성 방법")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.2), [
    "# === 기본 생성 ===",
    "np.array([1, 2, 3])",
    "np.array([[1, 2], [3, 4]], dtype=float)",
    "",
    "# === 초기화 생성 ===",
    "np.zeros((3, 4))        # 0으로 채운 3x4",
    "np.ones((2, 3))         # 1로 채운 2x3",
    "np.full((2, 3), 7)      # 7로 채운 2x3",
    "np.empty((3, 3))        # 초기화 없이 할당",
    "np.eye(4)               # 4x4 단위행렬",
], font_size=11)

add_code_block(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.2), [
    "# === 범위 생성 ===",
    "np.arange(0, 10, 2)     # [0, 2, 4, 6, 8]",
    "np.linspace(0, 1, 5)    # [0, 0.25, 0.5, 0.75, 1]",
    "np.logspace(0, 3, 4)    # [1, 10, 100, 1000]",
    "",
    "# === 난수 생성 ===",
    "np.random.seed(42)",
    "np.random.rand(3, 4)          # U(0,1)",
    "np.random.randn(3, 4)         # N(0,1)",
    "np.random.randint(0, 10, (3,4))# 정수 난수",
], font_size=11)

add_card(s, Inches(0.6), Inches(4.5), Inches(12.0), Inches(2.8),
         "주요 생성 함수 비교",
         ["zeros/ones/full/empty: 특정 값으로 초기화된 배열 생성 (shape 튜플 전달)",
          "eye/diag: 단위행렬, 대각행렬 생성 (선형대수 연산에 필수)",
          "arange: 정수 범위 (Python range와 유사), linspace: 균등 분할 (시각화에 자주 사용)",
          "logspace: 로그 스케일 균등 분할 (학습률 탐색 등)",
          "random.rand/randn: 균등분포/정규분포 난수 (ML 가중치 초기화에 활용)",
          "random.choice: 임의 선택 (부트스트랩, 교차 검증 등)"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 슬라이드 12: 배열 인덱싱과 슬라이싱
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "배열 인덱싱과 슬라이싱", "4가지 인덱싱 방법 (Harris et al., 2020)")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), [
    "arr = np.array([[1,  2,  3,  4],",
    "                [5,  6,  7,  8],",
    "                [9, 10, 11, 12]])",
    "",
    "# 1. 기본 인덱싱 (Basic Indexing)",
    "arr[0, 1]              # 2",
    "",
    "# 2. 슬라이싱 (Slicing) -- 뷰(view) 반환",
    "arr[0:2, 1:3]          # [[2, 3], [6, 7]]",
    "arr[:, 2]              # [3, 7, 11]",
    "",
    "# 3. 팬시 인덱싱 (Fancy) -- 카피(copy) 반환",
    "arr[[0, 2], [1, 3]]    # [2, 12]",
    "",
    "# 4. 불리언 인덱싱 (Boolean)",
    "arr[arr > 5]           # [6,7,8,9,10,11,12]",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.0),
         "View vs Copy 핵심 차이",
         ["슬라이싱 (a[1:4]): View 반환 -> 메모리 공유",
          "전치 (a.T): View 반환 -> stride만 교환",
          "팬시 인덱싱 (a[[0,2]]): Copy 반환",
          "불리언 인덱싱 (a[mask]): Copy 반환",
          "명시적 복사: a.copy()로 항상 Copy"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(7.0), Inches(4.3), Inches(5.6), Inches(2.5),
         "주의사항",
         ["View 수정 시 원본도 변경됨!",
          "  b = a[1:4]; b[0]=99 -> a도 변경!",
          "안전하게: c = a[1:4].copy()",
          "",
          "reshape(): 가능하면 View 반환",
          "flatten(): 항상 Copy 반환",
          "ravel(): 가능하면 View 반환"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 13: 브로드캐스팅 3가지 규칙
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "브로드캐스팅(Broadcasting): 3가지 규칙", "크기가 다른 배열 간 연산의 자동 메커니즘")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(1.3),
         "규칙 1: 차원 수 맞추기",
         ["두 배열의 차원 수(ndim)가 다르면, 차원 수가 적은 배열의 shape 앞에 1을 추가한다.",
          "예: shape (3,) -> (1, 3)   |   shape (4,) -> (1, 4)   |   shape (5,) -> (1, 1, 5)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(0.6), Inches(3.5), Inches(12.0), Inches(1.3),
         "규칙 2: 호환성 확인",
         ["각 차원에서 크기가 같거나, 둘 중 하나가 1이면 호환(compatible)된다.",
          "호환되지 않으면 ValueError 발생.  예: (3,) + (4,) -> 3 != 4 이고 둘 다 1이 아님 -> ERROR"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(5.0), Inches(12.0), Inches(1.3),
         "규칙 3: 가상 확장",
         ["크기가 1인 차원은 다른 배열의 크기에 맞게 가상으로 확장(stretch)된다.",
          "핵심: 실제로 데이터를 복사하지 않는다! 내부적으로 stride를 0으로 설정하여 메모리 효율적."],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_text(s, Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5),
         "예: (3,1)+(1,4) -> 규칙1 불필요, 규칙2 호환(3vs1, 1vs4), 규칙3 확장 -> 결과 shape (3,4)",
         font_size=15, color=ACCENT_BLUE, bold=True)

# ============================================================
# 슬라이드 14: 브로드캐스팅 코드 예제
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "브로드캐스팅 코드 예제", "ML에서의 실전 활용")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(5.0), [
    "# 예제 1: 행렬 + 벡터",
    "A = np.array([[1,2,3],[4,5,6],[7,8,9]])",
    "b = np.array([10, 20, 30])",
    "print(A + b)  # 각 행에 b가 더해짐",
    "",
    "# 예제 2: 열벡터 + 행벡터 -> 행렬 (외적 패턴)",
    "col = np.array([[1],[2],[3]])  # (3,1)",
    "row = np.array([10,20,30,40]) # (4,)-->(1,4)",
    "print(col + row)  # (3,4) 행렬",
    "",
    "# 예제 3: ML 데이터 중심화(centering)",
    "X = np.random.randn(100, 5)",
    "mean = X.mean(axis=0)          # (5,)",
    "X_centered = X - mean          # (100,5)-(5,)",
    "",
    "# 예제 4: Z-score 표준화",
    "std = X.std(axis=0)",
    "X_standardized = (X - mean) / std",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "시각적 설명: (3,3) + (3,)",
         ["배열 b: (3,) -> (1,3) -> (3,3)",
          "[a00 a01 a02]   [b0 b1 b2]",
          "[a10 a11 a12] + [b0 b1 b2]",
          "[a20 a21 a22]   [b0 b1 b2]",
          "",
          "-> 각 행에 같은 벡터 b가 더해짐"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.2),
         "ML에서의 핵심 활용",
         ["데이터 중심화: X - X.mean(axis=0)",
          "Z-score 표준화: (X - mean) / std",
          "배치 정규화(Batch Normalization)",
          "가중합(weighted sum): X @ w + b"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 15: 유니버설 함수 (ufunc)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "유니버설 함수 (ufunc)", "벡터화된 원소별 연산")

add_table_slide(s,
    headers=["분류", "함수 예시", "설명"],
    rows=[
        ["수학 함수", "np.sin, np.cos, np.exp, np.log, np.sqrt", "원소별 수학 연산"],
        ["비교 함수", "np.greater, np.equal, np.logical_and", "원소별 비교"],
        ["산술 함수", "np.add, np.subtract, np.multiply", "사칙연산"],
        ["집계 함수", "np.sum, np.prod, np.min, np.max", "축별 집계"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 5.0, 4.5],
    header_color=ACCENT_BLUE, row_height=0.5)

add_code_block(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.8), [
    "A = np.array([[1, 2, 3],",
    "              [4, 5, 6]])",
    "",
    "print(np.sum(A))          # 21 (전체 합)",
    "print(np.sum(A, axis=0))  # [5, 7, 9] (열방향)",
    "print(np.sum(A, axis=1))  # [6, 15] (행방향)",
    "print(np.mean(A, axis=0)) # [2.5, 3.5, 4.5]",
], font_size=12)

add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.8),
         "axis 이해하기",
         ["axis=None: 전체 원소에 대해 집계",
          "axis=0: 행 방향으로 축소 (열별 집계)",
          "axis=1: 열 방향으로 축소 (행별 집계)",
          "",
          "ufunc의 핵심: C 루프 기반 벡터화 연산",
          "Python 루프 대비 100~1000배 성능 향상"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 16: 선형대수 기본
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "선형대수 기본 연산", "ML의 수학적 기반")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "A = np.array([[1,2],[3,4]], dtype=float)",
    "B = np.array([[5,6],[7,8]], dtype=float)",
    "",
    "print(A * B)                # 원소별 곱셈",
    "print(A @ B)                # 행렬곱",
    "print(A.T)                  # 전치",
    "print(np.linalg.det(A))     # 행렬식",
    "print(np.linalg.inv(A))     # 역행렬",
    "print(np.trace(A))          # 대각합",
    "print(np.linalg.norm(A,'fro'))# 프로베니우스 노름",
], font_size=11)

add_table_slide(s,
    headers=["분해 방법", "수식", "ML 활용", "NumPy 함수"],
    rows=[
        ["LU 분해", "A = P * L * U", "연립방정식, 행렬식", "scipy.linalg.lu()"],
        ["QR 분해", "A = Q * R (Q^T Q = I)", "최소제곱법 안정 풀이", "np.linalg.qr()"],
        ["SVD", "A = U * S * V^T", "PCA, 추천시스템, LSA", "np.linalg.svd()"],
        ["고유값 분해", "Av = lambda*v", "PCA, 스펙트럴 클러스터링", "np.linalg.eigh()"],
    ],
    left=Inches(0.6), top=Inches(5.2), col_widths=[2.0, 2.5, 3.5, 2.5],
    header_color=ACCENT_PURPLE, row_height=0.45, font_size=12)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.8),
         "주요 행렬 연산 요약",
         ["* (원소별) vs @ (행렬곱): 반드시 구분!",
          "A.T: 전치 (View, stride만 교환)",
          "det(A): 행렬식 (가역 여부 판단)",
          "inv(A): 역행렬 (수치적 주의 필요)",
          "trace(A): 대각합 = 고유값의 합",
          "norm(A): L1, L2, Frobenius 노름"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 17: SVD와 고유값 분해
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "SVD(특이값 분해)와 고유값 분해", "PCA 차원 축소의 수학적 기반")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "# SVD: A = U * Sigma * V^T",
    "A = np.array([[1,2,0],[0,1,1],[2,0,1],",
    "              [1,1,1],[3,2,1]], dtype=float)",
    "",
    "U, s, Vt = np.linalg.svd(A, full_matrices=False)",
    "",
    "# 차원 축소: 상위 k개 특이값만 사용",
    "k = 2",
    "A_approx = U[:,:k] @ np.diag(s[:k]) @ Vt[:k,:]",
    "",
    "# 분산 설명력",
    "energy = s**2 / np.sum(s**2) * 100",
    "print(f'누적 설명 비율: {np.cumsum(energy)}')",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(5.2), Inches(6.0), Inches(2.0), [
    "# 고유값 분해: Av = lambda*v",
    "cov = np.array([[4,2,1],[2,3,1],[1,1,2]],",
    "               dtype=float)",
    "eigenvalues, eigenvectors = np.linalg.eigh(cov)",
    "variance_ratio = eigenvalues/np.sum(eigenvalues)*100",
    "print(f'분산 설명 비율: {variance_ratio}')",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "SVD의 ML 활용",
         ["PCA 차원 축소: 상위 k개 특이값으로 데이터 근사",
          "추천 시스템: 사용자-아이템 행렬의 잠재 요인",
          "잠재 의미 분석(LSA): 문서-단어 행렬",
          "데이터 압축: 정보 보존 비율 계산",
          "분산 설명력 = s^2 / sum(s^2) * 100"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.4),
         "고유값 분해의 ML 활용",
         ["PCA: 공분산 행렬의 고유값 분해",
          "스펙트럴 클러스터링: 라플라시안 행렬",
          "고유값 = 해당 방향의 분산량",
          "eigh(): 대칭 행렬 전용 (더 빠르고 안정)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 18: 정규방정식
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "정규방정식 (Normal Equation)", "선형회귀의 해석적 풀이: beta = (X^T X)^{-1} X^T y")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(4.8), [
    "np.random.seed(42)",
    "n_samples = 50",
    "X_raw = np.random.uniform(0, 10, (n_samples, 1))",
    "noise = np.random.normal(0, 1, (n_samples, 1))",
    "y = 3 * X_raw + 2 + noise  # 참값: 기울기=3, 절편=2",
    "",
    "X = np.hstack([np.ones((n_samples, 1)), X_raw])",
    "",
    "# 방법 1: 역행렬 (수치적으로 불안정 -- 비추천)",
    "beta_inv = np.linalg.inv(X.T @ X) @ X.T @ y",
    "",
    "# 방법 2: solve (LU 분해 기반, 수치적으로 안정)",
    "beta_solve = np.linalg.solve(X.T @ X, X.T @ y)",
    "",
    "# 방법 3: lstsq (SVD 기반, 가장 안정적, 권장!)",
    "beta_lstsq, _, _, _ = np.linalg.lstsq(X, y, rcond=None)",
    "",
    "print(f'참값: 절편=2, 기울기=3')",
    "print(f'추정: 절편={beta_lstsq[0,0]:.4f}, 기울기={beta_lstsq[1,0]:.4f}')",
], font_size=11)

add_card(s, Inches(8.0), Inches(2.0), Inches(4.6), Inches(5.0),
         "3가지 풀이 방법 비교",
         ["방법 1: inv(X^T X) @ X^T @ y",
          "  -> 조건수 큰 행렬에서 수치적 불안정",
          "",
          "방법 2: solve(X^T X, X^T @ y)",
          "  -> LU 분해 기반, 더 안정적",
          "",
          "방법 3: lstsq(X, y) [권장!]",
          "  -> SVD 기반, 가장 안정적",
          "  -> 랭크 부족 행렬도 처리 가능",
          "",
          "실무 권장: 항상 lstsq() 사용",
          "역행렬은 조건수 문제로 비추천"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 19: 선형대수와 ML 관계 총정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "선형대수와 ML 관계 총정리", "각 분해가 ML 알고리즘에서 어떻게 활용되는가")

add_table_slide(s,
    headers=["선형대수 연산", "ML 활용", "NumPy 함수"],
    rows=[
        ["SVD", "PCA, 추천 시스템, LSA", "np.linalg.svd()"],
        ["고유값 분해", "PCA, 스펙트럴 클러스터링", "np.linalg.eigh()"],
        ["QR 분해", "최소제곱법 안정적 풀이", "np.linalg.qr()"],
        ["LU 분해", "연립방정식 효율적 풀이", "scipy.linalg.lu()"],
        ["정규방정식", "선형 회귀", "np.linalg.lstsq()"],
        ["행렬곱", "신경망 순전파", "@ 연산자"],
        ["노름", "L1, L2 규제 (정규화)", "np.linalg.norm()"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.0, 4.5, 3.5],
    header_color=ACCENT_PURPLE, row_height=0.5)

add_card(s, Inches(0.6), Inches(5.8), Inches(12.0), Inches(1.5),
         "핵심 메시지",
         ["선형대수는 ML 알고리즘의 수학적 엔진이다. NumPy는 이 엔진의 구현체이다.",
          "PCA = SVD 또는 고유값 분해 | 선형회귀 = 정규방정식 | 신경망 = 행렬곱의 연쇄",
          "모든 ML 알고리즘의 내부에는 선형대수 연산이 존재한다."],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 20: View vs Copy
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 2, "View vs Copy 상세", "메모리 공유와 데이터 안전성")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.5), [
    "a = np.array([1, 2, 3, 4, 5])",
    "b = a[1:4]        # 뷰(View)",
    "b[0] = 99",
    "print(a)           # [1, 99, 3, 4, 5] -- 원본도 변경!",
    "",
    "c = a[1:4].copy()  # 카피(Copy)",
    "c[0] = 100",
    "print(a)           # [1, 99, 3, 4, 5] -- 원본 불변",
], font_size=12)

add_table_slide(s,
    headers=["연산", "반환 유형", "메모리 공유"],
    rows=[
        ["슬라이싱 a[1:4]", "View", "O"],
        ["전치 a.T", "View", "O"],
        ["reshape()", "View (가능 시)", "O"],
        ["ravel()", "View (가능 시)", "O"],
        ["팬시 인덱싱 a[[0,2]]", "Copy", "X"],
        ["불리언 인덱싱 a[mask]", "Copy", "X"],
        ["flatten()", "Copy", "X"],
        ["copy()", "Copy", "X"],
    ],
    left=Inches(0.6), top=Inches(4.8), col_widths=[3.5, 3.0, 2.0],
    header_color=ACCENT_ORANGE, row_height=0.35, font_size=12)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.2),
         "View vs Copy 판단 기준",
         ["View: 원본과 메모리를 공유하므로 효율적",
          "  -> 수정하면 원본도 변경됨 (주의!)",
          "  -> 슬라이싱, 전치, reshape 등",
          "",
          "Copy: 독립적인 메모리에 데이터 복사",
          "  -> 원본과 독립적, 수정해도 안전",
          "  -> 팬시 인덱싱, 불리언 인덱싱 등",
          "",
          "확인: np.shares_memory(a, b)",
          "안전한 복사: a.copy() 명시적 호출",
          "",
          "실무 팁: View인지 Copy인지 확인 후 수정!"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# Section 3: Pandas 심화
# ============================================================
section_divider("Pandas 심화: DataFrame의 세계", "DataFrame 구조, 인덱싱, 메서드 체이닝, GroupBy, Merge/Join", 3, ACCENT_ORANGE)

# ============================================================
# 슬라이드 21: DataFrame 내부 구조
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "DataFrame 내부 구조", "McKinney (2010)가 설계한 이질적 표 형식 데이터 구조")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "import pandas as pd",
    "import numpy as np",
    "",
    "df = pd.DataFrame({",
    "    '이름': ['김철수', '이영희', '박민수'],",
    "    '나이': [28, 35, 42],",
    "    '키':   [175.5, 162.3, 180.1],",
    "    '학생': [True, False, False]",
    "})",
    "",
    "print(df.dtypes)",
    "# 이름    object  /  나이    int64",
    "# 키    float64  /  학생      bool",
], font_size=11)

add_table_slide(s,
    headers=["속성", "설명", "예시"],
    rows=[
        ["df.shape", "(행 수, 열 수)", "(3, 4)"],
        ["df.dtypes", "각 열의 데이터 타입", "int64, float64, ..."],
        ["df.index", "행 인덱스", "RangeIndex(0, 3)"],
        ["df.columns", "열 이름", "Index(['이름','나이',...])"],
        ["df.values", "NumPy 배열 변환", "ndarray"],
        ["df.info()", "전반적 정보", "행/열수, dtype, 메모리"],
        ["df.describe()", "기술통계량", "평균, 표준편차, 사분위수"],
    ],
    left=Inches(7.0), top=Inches(2.0), col_widths=[1.8, 2.2, 2.2],
    header_color=ACCENT_ORANGE, row_height=0.4, font_size=11)

add_card(s, Inches(0.6), Inches(5.3), Inches(12.0), Inches(1.8),
         "DataFrame의 핵심 특징",
         ["각 열은 독립적인 NumPy 배열(또는 Extension Array)로 저장 -> 열마다 다른 dtype 가능",
          "행 인덱스(Index)를 통한 자동 데이터 정렬(automatic alignment)",
          "결측치(NaN)의 체계적 처리: isna(), fillna(), dropna() 등의 일관된 API"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# 슬라이드 22: 인덱싱 loc vs iloc
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "인덱싱: loc vs iloc", "라벨 기반 vs 정수 위치 기반")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), [
    "# 1. 열 선택",
    "df['나이']                  # Series 반환",
    "df[['이름', '나이']]         # DataFrame 반환",
    "",
    "# 2. loc - 라벨 기반 인덱싱",
    "df.loc[0, '이름']           # '김철수'",
    "df.loc[0:1, ['이름','나이']] # 행 0~1 (끝 포함!)",
    "",
    "# 3. iloc - 정수 위치 기반 인덱싱",
    "df.iloc[0, 0]              # '김철수'",
    "df.iloc[0:2, 0:2]          # 행 0~1 (끝 미포함!)",
    "",
    "# 4. 불리언 인덱싱",
    "df[df['나이'] > 30]         # 나이 > 30인 행",
    "df.query('나이 > 30 and 키 > 170')",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "loc vs iloc 핵심 차이",
         ["loc: 라벨(이름)으로 접근",
          "  -> 슬라이싱 시 끝 인덱스 포함",
          "  -> df.loc[0:2] -> 행 0, 1, 2 (3개)",
          "",
          "iloc: 정수 위치로 접근",
          "  -> 슬라이싱 시 끝 인덱스 미포함",
          "  -> df.iloc[0:2] -> 행 0, 1 (2개)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.2),
         "query() 메서드",
         ["문자열 표현식으로 행 필터링",
          "가독성이 좋고 체이닝에 적합",
          "내부적으로 numexpr 엔진 사용 (빠름)",
          "예: df.query('나이 > 30 and 키 > 170')"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 23: 메서드 체이닝 (Method Chaining)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "메서드 체이닝 (Method Chaining)", "중간 변수 없이 여러 변환을 연속 적용")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.5), [
    "# 나쁜 예: 중간 변수 남발",
    "df2 = df.dropna()",
    "df3 = df2[df2['나이'] > 25]",
    "df4 = df3.assign(나이대=df3['나이']//10*10)",
    "result = df4.sort_values('키', ascending=False)",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.5), [
    "# 좋은 예: 메서드 체이닝",
    "result = (df",
    "    .dropna()",
    "    .query('나이 > 25')",
    "    .assign(나이대=lambda x: x['나이']//10*10)",
    "    .sort_values('키', ascending=False)",
    ")",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.2),
         "pipe() 커스텀 함수 통합",
         ["def 이상치_제거(df, col, n_std=3):",
          "    mean, std = df[col].mean(), df[col].std()",
          "    return df[abs(df[col]-mean) <= n_std*std]",
          "",
          "result = (df",
          "    .pipe(이상치_제거, '키')",
          "    .assign(BMI=lambda x: x['키']/100))"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.5), Inches(5.6), Inches(2.7),
         "메서드 체이닝의 장점",
         ["가독성: 데이터 변환 과정이 위에서 아래로 읽힘",
          "디버깅: 각 단계를 주석 처리하여 중간 결과 확인",
          "재현성: 전체 변환이 하나의 표현식",
          "pipe(): 커스텀 함수도 체이닝에 통합 가능",
          "assign(): 새 열 추가 (lambda로 현재 df 참조)"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 24: GroupBy - 분할-적용-결합
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "GroupBy: 분할-적용-결합 (Split-Apply-Combine)", "McKinney (2010)가 강조한 핵심 패턴")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(5.0), [
    "np.random.seed(42)",
    "df = pd.DataFrame({",
    "    '부서': np.random.choice(['개발','마케팅','영업'], 100),",
    "    '직급': np.random.choice(['사원','대리','과장'], 100),",
    "    '연봉': np.random.normal(5000,1000,100).astype(int),",
    "    '성과': np.random.uniform(0,100,100).round(1)",
    "})",
    "",
    "# agg: 그룹별 집계",
    "df.groupby('부서')['연봉'].agg(['mean','std','min','max'])",
    "",
    "# transform: 원본 크기 유지 (그룹별 Z-score)",
    "df['연봉_zscore'] = df.groupby('부서')['연봉'].transform(",
    "    lambda x: (x - x.mean()) / x.std())",
    "",
    "# apply: 자유도 높은 그룹별 연산",
    "def 상위N(group, n=3):",
    "    return group.nlargest(n, '성과')",
    "df.groupby('부서').apply(상위N, n=2)",
], font_size=11)

add_card(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(5.0),
         "GroupBy 3가지 패턴",
         ["1. agg (집계): 그룹 -> 스칼라",
          "   그룹별 평균, 표준편차, 최소, 최대",
          "   결과: 그룹 수만큼의 행",
          "",
          "2. transform (변환): 그룹 -> 원본 크기",
          "   그룹별 Z-score, 누적합, 비율",
          "   결과: 원본과 같은 크기 (매우 유용!)",
          "",
          "3. apply (자유형): 가장 유연",
          "   그룹별 상위 N개 추출",
          "   커스텀 함수 적용",
          "",
          "Split-Apply-Combine 패턴:",
          "분할(groupby) -> 적용(agg/transform/apply)",
          "-> 결합(자동)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 25: Merge와 Join
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "Merge와 Join", "SQL JOIN과 동일한 테이블 결합 연산")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.2), [
    "# 두 테이블 준비",
    "직원 = pd.DataFrame({",
    "    '사번': [1, 2, 3, 4],",
    "    '이름': ['김철수','이영희','박민수','정수진'],",
    "    '부서코드': ['D01','D02','D01','D03']",
    "})",
    "부서 = pd.DataFrame({",
    "    '코드': ['D01','D02','D03'],",
    "    '부서명': ['개발팀','마케팅팀','영업팀']",
    "})",
    "",
    "# SQL JOIN과 동일",
    "result = pd.merge(직원, 부서,",
    "    left_on='부서코드', right_on='코드', how='left')",
], font_size=11)

add_table_slide(s,
    headers=["Join 유형", "SQL 대응", "설명"],
    rows=[
        ["inner", "INNER JOIN", "양쪽 모두 키 있는 행만"],
        ["left", "LEFT JOIN", "왼쪽 테이블 기준"],
        ["right", "RIGHT JOIN", "오른쪽 테이블 기준"],
        ["outer", "FULL OUTER JOIN", "양쪽 모두 포함"],
    ],
    left=Inches(0.6), top=Inches(5.5), col_widths=[2.0, 3.0, 5.0],
    header_color=ACCENT_ORANGE, row_height=0.4, font_size=12)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(3.2),
         "Merge 핵심 파라미터",
         ["on: 양쪽 테이블의 키 열이 같을 때",
          "left_on / right_on: 키 열 이름이 다를 때",
          "how: 'inner','left','right','outer'",
          "",
          "concat: 행 방향/열 방향 단순 결합",
          "  pd.concat([df1, df2], axis=0)",
          "",
          "join: 인덱스 기반 결합",
          "  df1.join(df2, how='left')"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# 슬라이드 26: Apply/Map과 시계열 윈도우 함수
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "Apply/Map과 시계열 윈도우 함수", "커스텀 함수 적용과 이동 평균")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.5), [
    "# apply: DataFrame의 각 열/행에 함수 적용",
    "df.apply(np.mean, axis=0)  # 열별 평균",
    "df.apply(np.mean, axis=1)  # 행별 평균",
    "",
    "# map: Series의 각 원소에 함수/딕셔너리 적용",
    "df['Sex'] = df['Sex'].map({'male':0, 'female':1})",
    "",
    "# applymap: DataFrame 전체 원소에 함수 적용",
    "df.applymap(lambda x: round(x, 2))",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.5), [
    "# 시계열 윈도우 함수",
    "df['MA_7'] = df['주가'].rolling(window=7).mean()",
    "df['MA_30'] = df['주가'].rolling(window=30).mean()",
    "df['변동성'] = df['주가'].rolling(window=20).std()",
    "",
    "# 누적 윈도우",
    "df['누적평균'] = df['주가'].expanding().mean()",
    "",
    "# 지수 가중 이동 평균 (EWM)",
    "df['EWM_12'] = df['주가'].ewm(span=12).mean()",
], font_size=11)

add_table_slide(s,
    headers=["윈도우 함수", "설명", "사용 예시"],
    rows=[
        ["rolling(n)", "고정 크기 n 윈도우", "이동평균, 이동표준편차"],
        ["expanding()", "시작~현재 누적", "누적 최대, 누적 평균"],
        ["ewm(span=n)", "지수 가중", "MACD, 지수이동평균"],
    ],
    left=Inches(7.0), top=Inches(2.0), col_widths=[1.6, 2.2, 2.4],
    header_color=ACCENT_ORANGE, row_height=0.45, font_size=11)

add_card(s, Inches(7.0), Inches(4.0), Inches(5.6), Inches(3.2),
         "실무 팁: 이동평균 교차(MA Crossover)",
         ["골든 크로스: 단기 MA(7일)가 장기 MA(30일)를",
          "  상향 돌파 -> 매수 신호",
          "데드 크로스: 하향 돌파 -> 매도 신호",
          "",
          "MACD = EWM(12) - EWM(26)",
          "기술적 분석의 핵심 지표",
          "",
          "McKinney(2010)가 금융 데이터 분석에서의",
          "경험을 바탕으로 Pandas에 기본 내장"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 27: 피벗 테이블과 MultiIndex
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "피벗 테이블과 MultiIndex", "SQL GROUP BY + CASE WHEN에 대응하는 데이터 요약")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "# 기본 피벗 테이블",
    "피벗 = pd.pivot_table(",
    "    매출,",
    "    values='매출액',",
    "    index='지역',       # 행 기준",
    "    columns='제품',      # 열 기준",
    "    aggfunc='mean',     # 집계 함수",
    "    margins=True         # 합계 행/열 추가",
    ")",
    "",
    "# 교차표 (Crosstab) - 빈도 집계",
    "pd.crosstab(매출['지역'], 매출['제품'], margins=True)",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(5.2), Inches(6.0), Inches(2.0), [
    "# MultiIndex (다중 인덱스)",
    "df_multi.loc['서울']                # 서울의 모든 분기",
    "df_multi.loc[('서울', '1분기')]     # 특정 행",
    "df_multi.xs('1분기', level='분기')  # 특정 레벨 선택",
    "df_multi.unstack(level='분기')      # 분기를 열로 이동",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.2),
         "Pandas 메모리 최적화 4가지 전략",
         ["1. 수치형 다운캐스팅:",
          "   int64 -> int32/int16 (50~75% 절약)",
          "",
          "2. Categorical 타입 변환:",
          "   반복 문자열에 효과적 (90~99% 절약)",
          "   100만행 x 3개 고유값: 64MB -> 1MB",
          "",
          "3. Sparse 타입:",
          "   대부분 0인 데이터 (90%+ 절약)",
          "",
          "4. 청크 단위 읽기 (chunksize):",
          "   chunks = pd.read_csv('huge.csv',",
          "       chunksize=100000)",
          "   RAM 제한 내에서 처리 가능"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 28: SQL vs Pandas 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 3, "SQL vs Pandas 비교", "SQL을 아는 사람을 위한 Pandas 매핑")

add_table_slide(s,
    headers=["SQL", "Pandas", "비고"],
    rows=[
        ["SELECT col FROM t", "df['col'] 또는 df[['col']]", "열 선택"],
        ["WHERE cond", "df[df['col']>v] / df.query()", "행 필터링"],
        ["GROUP BY col", "df.groupby('col')", "그룹별 집계"],
        ["ORDER BY col", "df.sort_values('col')", "정렬"],
        ["JOIN t1, t2 ON", "pd.merge(t1, t2, on=)", "테이블 결합"],
        ["DISTINCT", "df.drop_duplicates()", "중복 제거"],
        ["COUNT(*)", "df.shape[0] / len(df)", "행 수"],
        ["LIMIT n", "df.head(n)", "상위 n개"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[3.5, 4.0, 3.5],
    header_color=ACCENT_BLUE, row_height=0.45, font_size=12)

add_card(s, Inches(0.6), Inches(6.2), Inches(12.0), Inches(1.0),
         "핵심 메시지",
         ["SQL 경험이 있다면 Pandas는 매우 직관적! 거의 1:1 대응이 가능하며, 메서드 체이닝으로 더 유연한 분석이 가능하다."],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# Section 4: 데이터 전처리
# ============================================================
section_divider("데이터 전처리", "결측치 처리, 이상치 탐지, 피처 스케일링, 인코딩", 4, ACCENT_GREEN)

# ============================================================
# 슬라이드 29: 결측치 처리 5가지 전략
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 4, "결측치 처리 5가지 전략", "상황에 따라 적절한 전략 선택이 중요")

add_table_slide(s,
    headers=["전략", "방법", "장점", "단점", "적합한 상황"],
    rows=[
        ["행 삭제", "dropna()", "단순함", "데이터 손실", "결측 < 5%"],
        ["평균 대체", "fillna(mean)", "분포 유지", "분산 과소추정", "정규분포, MCAR"],
        ["중앙값 대체", "fillna(median)", "이상치 강건", "분산 과소추정", "편향 분포"],
        ["보간", "interpolate()", "연속성 유지", "외삽 위험", "시계열 데이터"],
        ["그룹별 대체", "groupby.transform", "그룹 특성 반영", "구현 복잡", "그룹 간 차이 큰 경우"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[1.5, 2.0, 1.8, 1.8, 2.2],
    header_color=ACCENT_GREEN, row_height=0.5, font_size=11)

add_code_block(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.5), [
    "df = pd.DataFrame({'A': [1, np.nan, 3, np.nan, 5], 'B': [10,20,np.nan,40,50], '그룹': ['X','X','Y','Y','Y']})",
    "",
    "# 전략별 비교",
    "print('삭제:',   df['A'].dropna().values)                              # NaN 제거",
    "print('평균:',   df['A'].fillna(df['A'].mean()).values)               # 전체 평균으로 대체",
    "print('중앙값:', df['A'].fillna(df['A'].median()).values)             # 전체 중앙값으로 대체",
    "print('보간:',   df['A'].interpolate().values)                        # 선형 보간",
    "print('그룹별:', df.groupby('그룹')['A'].transform(lambda x: x.fillna(x.mean())).values)  # 그룹 평균으로 대체",
], font_size=10)

# ============================================================
# 슬라이드 30: 이상치 탐지 3가지 방법
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 4, "이상치 탐지 3가지 방법", "IQR, Z-score, 시각적 방법")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "IQR 방법 (사분위수 범위)",
         ["이상치 범위: Q1 - 1.5*IQR < x < Q3 + 1.5*IQR",
          "IQR = Q3 - Q1 (사분위수 범위)",
          "정규분포 가정 불필요",
          "박스플롯과 동일한 기준"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "Z-score 방법",
         ["z = (x - mu) / sigma",
          "|z| > 3 이면 이상치로 판단",
          "정규분포 가정 필요",
          "평균과 표준편차에 민감"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_code_block(s, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.5), [
    "def iqr_이상치(series):",
    "    Q1, Q3 = series.quantile([0.25, 0.75])",
    "    IQR = Q3 - Q1",
    "    lower = Q1 - 1.5 * IQR",
    "    upper = Q3 + 1.5 * IQR",
    "    return (series < lower) | (series > upper)",
], font_size=11)

add_code_block(s, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.5), [
    "def zscore_이상치(series, threshold=3):",
    "    z = (series - series.mean()) / series.std()",
    "    return abs(z) > threshold",
    "",
    "# 사용 예",
    "mask = iqr_이상치(df['연봉'])",
    "df_clean = df[~mask]  # 이상치 제거",
], font_size=11)

# ============================================================
# 슬라이드 31: 피처 스케일링 3가지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 4, "피처 스케일링 3가지 방법", "StandardScaler, MinMaxScaler, RobustScaler")

add_table_slide(s,
    headers=["스케일러", "수식", "특성", "사용 상황"],
    rows=[
        ["StandardScaler", "z = (x-mu)/sigma", "평균0, 분산1", "정규분포, SVM, 로지스틱"],
        ["MinMaxScaler", "z = (x-min)/(max-min)", "[0,1] 범위", "신경망, 이미지"],
        ["RobustScaler", "z = (x-Q2)/(Q3-Q1)", "중앙값/IQR 기반", "이상치 존재 시"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 3.0, 2.5, 3.0],
    header_color=ACCENT_GREEN, row_height=0.55, font_size=12)

add_code_block(s, Inches(0.6), Inches(4.0), Inches(6.0), Inches(2.0), [
    "from sklearn.preprocessing import (StandardScaler,",
    "    MinMaxScaler, RobustScaler)",
    "",
    "# 중요: fit은 train에만, transform은 train/test 모두",
    "scaler = StandardScaler()",
    "X_train_scaled = scaler.fit_transform(X_train)",
    "X_test_scaled = scaler.transform(X_test)  # fit 없이!",
], font_size=12)

add_card(s, Inches(7.0), Inches(4.0), Inches(5.6), Inches(3.2),
         "Data Leakage 경고!",
         ["fit_transform을 전체 데이터에 적용하면",
          "테스트 데이터의 정보가 훈련에 유출됨!",
          "",
          "올바른 방법:",
          "  1. scaler.fit_transform(X_train) -- train에 fit",
          "  2. scaler.transform(X_test) -- test에 transform만",
          "",
          "잘못된 방법:",
          "  scaler.fit_transform(X_all) -- 전체 데이터에 fit",
          "  -> Data Leakage 발생!"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 32: 인코딩
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 4, "범주형 변수 인코딩", "Label Encoding, One-Hot Encoding, Ordinal Encoding")

add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(2.5),
         "Label Encoding",
         ["각 범주에 정수 할당",
          "예: male->0, female->1",
          "트리 모델에 적합",
          "순서 관계가 생길 수 있음 (주의)",
          "df['Sex'].map({'male':0,'female':1})"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(2.5),
         "One-Hot Encoding",
         ["각 범주를 이진 벡터로 변환",
          "예: [1,0,0], [0,1,0], [0,0,1]",
          "순서 관계 없음 (안전)",
          "고유값 많으면 차원 폭발 문제",
          "pd.get_dummies(df, columns=['지역'])"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(2.5),
         "Ordinal Encoding",
         ["순서가 있는 범주에 정수 할당",
          "예: 초등->1, 중등->2, 고등->3",
          "순서 정보를 보존",
          "명목형 변수에는 부적절",
          "OrdinalEncoder(categories=...)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.4),
         "인코딩 선택 가이드",
         ["이진 범주 (2개 값): Label Encoding -> male/female, yes/no",
          "명목형 범주 (순서 없음): One-Hot Encoding -> 지역, 색상, 제품종류",
          "순서형 범주 (순서 있음): Ordinal Encoding -> 학력, 등급, 만족도",
          "고유값 매우 많음 (100+): Target Encoding, Hash Encoding 고려",
          "",
          "주의: Label Encoding을 명목형에 사용하면 모델이 잘못된 순서 관계를 학습할 수 있음!"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# Section 5: Tidy Data 개념
# ============================================================
section_divider("Tidy Data 개념", "Wickham (2014)의 깔끔한 데이터 3원칙", 5, ACCENT_PURPLE)

# ============================================================
# 슬라이드 33: Wickham의 Tidy Data 3원칙
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 5, "Wickham의 Tidy Data 3원칙", "Journal of Statistical Software, 2014")

add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(2.0),
         "원칙 1: 변수 = 열",
         ["각 변수(variable)는",
          "하나의 열(column)을 구성한다.",
          "예: 이름, 나이, 점수 등",
          "각각 별도의 열로 표현"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(2.0),
         "원칙 2: 관측 = 행",
         ["각 관측(observation)은",
          "하나의 행(row)을 구성한다.",
          "예: 한 학생의 정보는",
          "하나의 행으로 표현"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(2.0),
         "원칙 3: 관측단위 = 테이블",
         ["각 관측 단위 유형은",
          "하나의 테이블을 구성한다.",
          "예: 학생 테이블, 과목 테이블",
          "별도의 테이블로 분리"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.3), Inches(12.0), Inches(3.0),
         "Messy Data의 5가지 유형 (Wickham 2014)",
         ["(1) 열 헤더가 변수명이 아닌 값인 경우  ->  melt()로 변환",
          "(2) 하나의 열에 여러 변수가 저장된 경우  ->  str.split() + assign()으로 분리",
          "(3) 변수가 행과 열 모두에 걸쳐 저장된 경우  ->  melt() + pivot() 조합",
          "(4) 하나의 테이블에 여러 관측 단위가 혼재된 경우  ->  테이블 분리",
          "(5) 하나의 관측 단위가 여러 테이블에 분산된 경우  ->  merge()로 결합",
          "",
          "scikit-learn 입력 형식 '행=샘플, 열=특성'은 Tidy Data 원칙과 정확히 일치!",
          "Tidy Data -> groupby/agg 자연스럽게 동작, seaborn 시각화 호환, ML 파이프라인 바로 투입"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 34: melt/pivot 변환
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 5, "melt/pivot: Messy -> Tidy 변환", "Wide Format <-> Long Format 변환")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), [
    "# Messy: 열 헤더가 값인 경우",
    "messy = pd.DataFrame({",
    "    '이름': ['김철수', '이영희'],",
    "    '2023_국어': [85, 92],",
    "    '2023_수학': [90, 88],",
    "    '2024_국어': [88, 95],",
    "    '2024_수학': [92, 90]",
    "})",
    "",
    "# Tidy로 변환",
    "tidy = (messy",
    "    .melt(id_vars='이름',",
    "          var_name='과목_연도',",
    "          value_name='점수')",
    "    .assign(",
    "        연도=lambda x: x['과목_연도'].str.split('_').str[0],",
    "        과목=lambda x: x['과목_연도'].str.split('_').str[1]",
    "    )",
    "    .drop(columns='과목_연도')",
    ")",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "melt() -- Wide -> Long",
         ["id_vars: 유지할 열 (식별자)",
          "value_vars: 녹일 열 (생략 시 나머지 전부)",
          "var_name: 변수명 열 이름",
          "value_name: 값 열 이름",
          "",
          "결과: 넓은 테이블 -> 긴 테이블"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.4),
         "pivot() / pivot_table() -- Long -> Wide",
         ["index: 행이 될 열",
          "columns: 열이 될 열",
          "values: 값이 될 열",
          "aggfunc: 집계 함수 (pivot_table만)",
          "",
          "결과: 긴 테이블 -> 넓은 테이블"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# Section 6: 벡터화 vs 루프
# ============================================================
section_divider("벡터화 vs 루프: 성능의 과학", "Python 루프 대비 100~1000배 속도 향상의 원리", 6, ACCENT_RED)

# ============================================================
# 슬라이드 35: 벡터화 vs 루프 성능 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 6, "벡터화 vs 루프: 성능 비교", "Harris et al. (2020) Nature 논문에서 확인된 100~1000배 속도 향상")

add_table_slide(s,
    headers=["요인", "Python 루프", "NumPy 벡터화"],
    rows=[
        ["실행 엔진", "Python 인터프리터", "C/Fortran 컴파일 코드"],
        ["타입 체크", "매 연산마다", "한 번만"],
        ["메모리 접근", "불연속 (포인터 추적)", "연속 (캐시 친화적)"],
        ["SIMD 활용", "불가", "가능 (SSE, AVX)"],
        ["GIL", "보유", "해제 가능"],
        ["BLAS/LAPACK", "미사용", "활용 (Intel MKL 등)"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 3.5, 3.5],
    header_color=ACCENT_RED, row_height=0.45, font_size=12)

add_code_block(s, Inches(0.6), Inches(5.0), Inches(6.0), Inches(2.3), [
    "n = 1_000_000",
    "a = np.random.randn(n); b = np.random.randn(n)",
    "",
    "# Python 루프: ~0.3초",
    "c_loop = [a[i]+b[i] for i in range(n)]",
    "",
    "# NumPy 벡터화: ~0.001초 (300배 빠름!)",
    "c_vec = a + b",
], font_size=12)

add_card(s, Inches(7.0), Inches(5.0), Inches(5.6), Inches(2.3),
         "속도 향상의 원인 요약",
         ["1. C 컴파일 코드: Python 인터프리터 오버헤드 제거",
          "2. 연속 메모리: CPU L1/L2 캐시 적중률 극대화",
          "3. SIMD: 단일 명령으로 여러 데이터 동시 처리",
          "4. 타입 체크 생략: 동질적 dtype으로 한 번만 확인",
          "5. BLAS/LAPACK: 최적화된 수치 라이브러리 활용"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 36: SIMD와 캐시의 역할
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 6, "SIMD와 CPU 캐시의 역할", "벡터화 성능의 하드웨어 원리")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "SIMD (Single Instruction, Multiple Data)",
         ["하나의 CPU 명령어로 여러 데이터를 동시 처리",
          "AVX-256: 8개의 float32를 한 번에 연산",
          "AVX-512: 16개의 float32를 한 번에 연산",
          "NumPy는 내부적으로 SIMD 명령어 활용",
          "Python 루프에서는 SIMD 활용 불가능"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "CPU 캐시 효율",
         ["NumPy 배열: 연속 메모리 -> L1/L2 캐시 적중률 높음",
          "Python 리스트: 포인터 간접 접근 -> 캐시 미스 빈번",
          "",
          "L1 캐시: ~1ns (가장 빠름, 32~64KB)",
          "L2 캐시: ~4ns (256KB~1MB)",
          "메인 메모리: ~100ns (매우 느림)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.5),
         "실전 벤치마크: 쌍별 유클리드 거리",
         ["200x300 쌍별 거리 계산 결과:",
          "  Python 3중 루프: ~15초",
          "  NumPy 벡터화: ~0.01초",
          "  Numba JIT: ~0.02초",
          "  Numba parallel: ~0.005초",
          "벡터화는 루프 대비 약 500~1000배 빠름!"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.5),
         "성능 최적화 선택 가이드",
         ["1순위: NumPy 벡터화 (가장 간단하고 빠름)",
          "2순위: Numba JIT (벡터화 어려운 복잡한 루프)",
          "3순위: Numba parallel (멀티코어 활용)",
          "",
          "벡터화 가능하면 NumPy가 최선!",
          "복잡한 조건분기/동적 프로그래밍만 Numba 사용"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 37: 벡터화 코드 예제 - 유클리드 거리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 6, "벡터화 실전: 쌍별 유클리드 거리", "루프 vs 브로드캐스팅 비교")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.8), [
    "# Python 루프: O(n*m*d) 3중 루프",
    "def 거리_루프(X, Y):",
    "    n, d = X.shape; m = Y.shape[0]",
    "    D = np.zeros((n, m))",
    "    for i in range(n):",
    "        for j in range(m):",
    "            s = 0",
    "            for k in range(d):",
    "                s += (X[i,k] - Y[j,k])**2",
    "            D[i,j] = np.sqrt(s)",
    "    return D",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(5.0), Inches(6.0), Inches(2.2), [
    "# NumPy 브로드캐스팅: 단 한 줄!",
    "def 거리_벡터화(X, Y):",
    "    return np.sqrt(np.sum(",
    "        (X[:,np.newaxis,:] - Y[np.newaxis,:,:])**2,",
    "        axis=2))",
    "",
    "# KNN, K-Means 등 거리 기반 알고리즘에 필수!",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.2),
         "브로드캐스팅 거리 계산 해부",
         ["X: shape (n, d)  ->  (n, 1, d)",
          "Y: shape (m, d)  ->  (1, m, d)",
          "",
          "X[:,np.newaxis,:] - Y[np.newaxis,:,:]",
          "  -> shape (n, m, d)  (차이 벡터)",
          "",
          "**2  -> 제곱",
          "np.sum(..., axis=2)  -> 합산 -> (n, m)",
          "np.sqrt(...)  -> 유클리드 거리",
          "",
          "전체 과정이 C 레벨에서 실행!",
          "메모리: O(n*m*d) 필요 (큰 데이터 주의)"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# Section 7: 논문 리뷰 통합
# ============================================================
section_divider("논문 리뷰 통합", "Harris(2020), McKinney(2010), Walt(2011), Wickham(2014), Reback(2020)", 7, ACCENT_PURPLE)

# ============================================================
# 슬라이드 38: 논문 개관
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 7, "3장 관련 논문 개관", "5편의 핵심 논문과 학문적 맥락")

add_table_slide(s,
    headers=["#", "논문", "핵심 기여", "인용수"],
    rows=[
        ["1", "McKinney (2010)", "Pandas DataFrame, Split-Apply-Combine", "10,000+"],
        ["2", "Walt et al. (2011)", "ndarray 내부 구조 (stride, dtype) 문서화", "5,000+"],
        ["3", "Harris et al. (2020)", "NumPy Nature 논문, 배열 프로그래밍 패러다임", "8,000+"],
        ["4", "Reback et al. (2020)", "Pandas 공식 인용 문서, 전체 기능 개관", "15,000+"],
        ["5", "Wickham (2014)", "Tidy Data 3원칙, messy data 5가지 유형", "7,000+"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[0.5, 3.0, 5.5, 1.5],
    header_color=ACCENT_PURPLE, row_height=0.55, font_size=12)

add_card(s, Inches(0.6), Inches(5.2), Inches(12.0), Inches(2.0),
         "논문 간 연결 관계 (시간순)",
         ["McKinney (2010): 'ndarray는 이질적 데이터에 부족' -> Pandas 개발",
          "Walt et al. (2011): ndarray의 기술적 설계 상세 문서화 (stride, ufunc, broadcast)",
          "Wickham (2014): 데이터 정리의 원칙 체계화 -> Pandas melt()/pivot_table() 설계에 영향",
          "Harris et al. (2020): NumPy 15년간 발전 총정리, Nature에서 공식 인정"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# 슬라이드 39: Harris (2020) - NumPy Nature 논문
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 7, "Harris et al. (2020) -- NumPy Nature 논문", '"Array programming with NumPy", Nature, 585, 357-362')

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "배경과 의의",
         ["Nature에 SW 라이브러리 논문 게재 = NumPy의 학술적 영향력 증명",
          "15년간의 NumPy 발전을 총정리",
          "현대 과학 컴퓨팅에서 NumPy의 핵심 역할 공식 인정",
          "배열 프로그래밍(Array Programming) 패러다임 공식 정의"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "핵심 기여 4가지",
         ["1. 배열 프로그래밍 패러다임: 반복문 대신 배열 전체 연산",
          "2. 파이썬 과학 생태계 레이어 구조 시각화",
          "3. NumPy API 프로토콜: __array_ufunc__, __array_function__",
          "4. GPU/분산 환경 확장: JAX, CuPy, Dask Array 호환"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.4),
         "핵심 인용과 ML 과목 관련성",
         ['"배열 프로그래밍은 개별 원소에 대한 반복문 대신, 배열 전체에 연산을 적용하는 프로그래밍 패러다임이다.',
          ' 이 패러다임은 코드의 간결성, 가독성, 성능을 동시에 향상시킨다."',
          "",
          "ML 과목 관련성: 본 과목에서 사용하는 거의 모든 라이브러리(scikit-learn, TensorFlow, PyTorch)의 기반이 NumPy.",
          "이 논문은 그 기반 인프라의 설계 원칙을 체계적으로 설명한다."],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 40: McKinney (2010) & Wickham (2014)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 7, "McKinney (2010) & Wickham (2014)", "Pandas의 탄생과 Tidy Data 원칙")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
         "McKinney (2010) -- Pandas의 탄생",
         ['"Data Structures for Statistical Computing in Python"',
          "SciPy 2010",
          "",
          "배경: AQR Capital에서 금융 데이터 분석 중",
          "파이썬에 R의 data.frame 대응 구조 부재 절감",
          "",
          "핵심 기여:",
          "  1. Series와 DataFrame 도입 (인덱스 기반 자동 정렬)",
          "  2. 결측치(NaN) 체계적 처리 API",
          "  3. Split-Apply-Combine 패턴 (groupby)",
          "  4. 관계형 데이터 결합 (merge, join)",
          "",
          '인용: "파이썬에 R의 data.frame에 대응하는',
          '고수준 데이터 구조가 없었다."'],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
         "Wickham (2014) -- Tidy Data 원칙",
         ['"Tidy Data", J. of Statistical Software, 59(10)',
          "",
          "배경: 데이터 과학자 업무의 50~80%가 데이터 정리",
          "체계적 이론 부재 -> 원칙 체계화",
          "",
          "핵심 기여:",
          "  1. 깔끔한 데이터의 3원칙 (변수=열, 관측=행, 관측단위=테이블)",
          "  2. 지저분한 데이터 5가지 유형 분류",
          "  3. 변환 도구: melt(wide->long), pivot(long->wide)",
          "",
          "ML 관련성:",
          '  scikit-learn의 "행=샘플, 열=특성" 형식',
          "  = Tidy Data 원칙과 정확히 일치!",
          '인용: "깔끔한 데이터면 동일 분석 코드를 재사용 가능"'],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# Section 8: 구현 코드 상세 해설
# ============================================================
section_divider("구현 코드 상세 해설", "3개 실습 코드: NumPy 선형대수, Pandas 파이프라인, 벡터화 벤치마크", 8, ACCENT_CYAN)

# ============================================================
# 슬라이드 41: 실습 1 - NumPy 선형대수
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "실습 1: NumPy 선형대수", "01_numpy_linear_algebra.py - SVD 차원 축소")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(5.0), [
    "# SVD를 이용한 차원 축소",
    "import numpy as np",
    "",
    "A = np.array([[1,2,0],[0,1,1],[2,0,1],",
    "              [1,1,1],[3,2,1]], dtype=float)",
    "",
    "U, s, Vt = np.linalg.svd(A, full_matrices=False)",
    "",
    "# 상위 k개 특이값으로 근사",
    "k = 2",
    "A_approx = U[:,:k] @ np.diag(s[:k]) @ Vt[:k,:]",
    "",
    "# 에너지(분산 설명력) 분석",
    "energy = s**2 / np.sum(s**2) * 100",
    "print(f'특이값: {s}')",
    "print(f'에너지 비율: {energy}')",
    "print(f'상위 {k}개 누적: {np.sum(energy[:k]):.1f}%')",
], font_size=11)

add_card(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(5.0),
         "코드 해석",
         ["SVD: A = U * Sigma * V^T 로 분해",
          "",
          "특이값(s)의 크기 = 해당 방향의 데이터 분산량",
          "상위 k개만 사용 = 정보 손실 최소화 + 차원 축소",
          "-> 이것이 PCA의 수학적 기반!",
          "",
          "full_matrices=False: 축소된 SVD",
          "  U: (m, k), s: (k,), Vt: (k, n)",
          "",
          "에너지 비율 = s^2 / sum(s^2) * 100",
          "  각 특이값이 전체 분산의 몇 %를 설명하는지",
          "",
          "실무: 누적 에너지 95% 이상 보존하는 k 선택"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

# ============================================================
# 슬라이드 42: 실습 1 (계속) - 정규방정식
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "실습 1 (계속): 정규방정식 풀이 비교", "역행렬 vs solve vs lstsq")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(4.5), [
    "# 데이터 생성: y = 3x + 2 + noise",
    "np.random.seed(42)",
    "X_raw = np.random.uniform(0, 10, (50, 1))",
    "y = 3 * X_raw + 2 + np.random.normal(0, 1, (50, 1))",
    "X = np.hstack([np.ones((50, 1)), X_raw])",
    "",
    "# 방법 1: 역행렬 (수치적으로 불안정)",
    "beta_inv = np.linalg.inv(X.T @ X) @ X.T @ y",
    "",
    "# 방법 2: solve (LU 분해 기반, 안정적)",
    "beta_solve = np.linalg.solve(X.T @ X, X.T @ y)",
    "",
    "# 방법 3: lstsq (SVD 기반, 가장 안정적, 권장!)",
    "beta_lstsq, _, _, _ = np.linalg.lstsq(X, y, rcond=None)",
    "",
    "print(f'참값: 절편=2, 기울기=3')",
    "print(f'lstsq: 절편={beta_lstsq[0,0]:.4f}, 기울기={beta_lstsq[1,0]:.4f}')",
], font_size=11)

add_card(s, Inches(8.0), Inches(2.0), Inches(4.6), Inches(5.0),
         "권장 사항",
         ["실무에서는 항상 np.linalg.lstsq() 사용!",
          "",
          "역행렬 계산의 문제점:",
          "  - 조건수(condition number)가 큰 행렬에서",
          "    수치적 불안정",
          "  - 특이 행렬(singular matrix) 처리 불가",
          "",
          "lstsq의 장점:",
          "  - SVD 기반으로 가장 안정적",
          "  - 랭크 부족(rank-deficient) 행렬도 처리",
          "  - 잔차, 랭크, 특이값도 함께 반환",
          "",
          "scikit-learn의 LinearRegression도",
          "내부적으로 lstsq 사용!"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# 슬라이드 43: 실습 2 - Pandas 데이터 파이프라인
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "실습 2: Pandas 데이터 파이프라인", "02_pandas_data_pipeline.py - 결측치 전략 비교")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(5.0), [
    "import pandas as pd; import numpy as np",
    "np.random.seed(42); n = 500",
    "df = pd.DataFrame({",
    "    '날짜': pd.date_range('2024-01-01', periods=n, freq='6H'),",
    "    '고객ID': np.random.choice(['C001','C002','C003'], n),",
    "    '금액': np.random.exponential(50000, n).astype(int)+5000,",
    "    '수량': np.random.randint(1, 10, size=n)})",
    "",
    "# 인위적 결측치 삽입 (10%)",
    "mask = np.random.random(n) < 0.1",
    "df.loc[mask, '금액'] = np.nan",
    "",
    "# 5가지 전략 비교",
    "strategies = {",
    "    '삭제': df['금액'].dropna(),",
    "    '평균': df['금액'].fillna(df['금액'].mean()),",
    "    '중앙값': df['금액'].fillna(df['금액'].median()),",
    "    '0대체': df['금액'].fillna(0),",
    "    '보간': df['금액'].interpolate()}",
], font_size=10)

add_card(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(5.0),
         "코드 해석 및 결과",
         ["각 전략의 영향을 평균과 표준편차로 비교:",
          "",
          "삭제: 원본 통계량 유지 (데이터 손실)",
          "평균: 평균 동일, 분산 과소 추정",
          "중앙값: 이상치에 강건, 분산 과소 추정",
          "0대체: 평균과 분산 모두 왜곡 (비추천)",
          "보간: 시계열에서 연속성 유지",
          "",
          "GroupBy + Transform 패턴:",
          "  그룹별 Z-score 정규화",
          "  그룹별 누적 합계 (cumsum)",
          "  그룹별 결측치를 그룹 중앙값으로 대체",
          "",
          "가장 권장: 상황에 맞는 전략 선택!"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 44: 실습 2 (계속) - GroupBy Transform
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "실습 2 (계속): GroupBy + Transform 패턴", "그룹별 정규화와 결측치 처리")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.5), [
    "# 고객별 Z-score 정규화",
    "df['금액_zscore'] = df.groupby('고객ID')['금액'].transform(",
    "    lambda x: (x - x.mean()) / x.std())",
    "",
    "# 고객별 누적 합계",
    "df['누적금액'] = df.groupby('고객ID')['금액'].cumsum()",
], font_size=12)

add_code_block(s, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.5), [
    "# 고객별 결측치를 그룹 중앙값으로 대체",
    "df['금액_filled'] = df.groupby('고객ID')['금액'].transform(",
    "    lambda x: x.fillna(x.median()))",
    "",
    "# 이 패턴은 ML 전처리에서 가장 강력한 도구!",
    "# 전체 평균보다 그룹별 대체가 더 정확함",
], font_size=12)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.2),
         "transform의 핵심 특성",
         ["transform vs agg 차이:",
          "  agg: 그룹 -> 스칼라 (행 수 축소)",
          "  transform: 그룹 -> 원본 크기 (행 수 유지)",
          "",
          "transform 활용 패턴:",
          "  1. 그룹별 Z-score 정규화",
          "  2. 그룹별 비율 계산",
          "  3. 그룹별 결측치 대체",
          "  4. 그룹별 누적합/이동평균",
          "",
          "왜 유용한가?",
          "  원본 DataFrame에 바로 새 열로 추가 가능",
          "  전체 통계보다 그룹 통계가 더 정확한 대체",
          "  ML에서 피처 엔지니어링의 핵심 도구"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 45: 실습 3 - 벡터화 벤치마크
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 8, "실습 3: 벡터화 벤치마크", "03_vectorization_benchmark.py - 거리 계산 성능 비교")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(5.0), [
    "def 거리_루프(X, Y):",
    "    n, d = X.shape; m = Y.shape[0]",
    "    D = np.zeros((n, m))",
    "    for i in range(n):",
    "        for j in range(m):",
    "            s = 0",
    "            for k in range(d):",
    "                s += (X[i,k]-Y[j,k])**2",
    "            D[i,j] = np.sqrt(s)",
    "    return D",
    "",
    "def 거리_벡터화(X, Y):",
    "    return np.sqrt(np.sum(",
    "        (X[:,np.newaxis,:]-Y[np.newaxis,:,:])**2, axis=2))",
    "",
    "X = np.random.randn(200, 10)",
    "Y = np.random.randn(300, 10)",
    "# 루프: ~15초 / 벡터화: ~0.01초 / 속도: ~1000배",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(5.0),
         "벤치마크 결과 및 해석",
         ["200x300 쌍별 거리 (10차원):",
          "",
          "방법            | 시간      | 속도비",
          "Python 루프     | ~15초     | 1x",
          "NumPy 벡터화    | ~0.01초   | 1,500x",
          "Numba JIT       | ~0.02초   | 750x",
          "Numba parallel  | ~0.005초  | 3,000x",
          "",
          "결과 일치: np.allclose(D_loop, D_vec) = True",
          "",
          "KNN, K-Means 등 거리 기반 알고리즘에서",
          "벡터화는 선택이 아닌 필수!",
          "scikit-learn도 내부적으로 벡터화 사용"],
         title_color=ACCENT_RED, border=ACCENT_RED)

# ============================================================
# Section 9: 핵심 요약
# ============================================================
section_divider("핵심 요약 및 복습", "3장의 핵심 개념 정리와 복습 질문", 9, ACCENT_BLUE)

# ============================================================
# 슬라이드 46: 핵심 요약표
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "핵심 요약표", "3장에서 배운 모든 개념 정리")

add_table_slide(s,
    headers=["개념", "핵심 포인트"],
    rows=[
        ["ndarray", "연속 메모리, stride, 동질적 dtype, 벡터화 연산의 기반"],
        ["브로드캐스팅", "크기 다른 배열 간 연산 자동화, stride=0 가상 확장"],
        ["DataFrame", "이질적 표 구조, 열별 독립 배열, 라벨 기반 인덱싱"],
        ["메서드 체이닝", ".assign().query().sort_values() 패턴"],
        ["GroupBy", "분할-적용-결합, agg / transform / apply"],
        ["Tidy Data", "변수=열, 관측=행, 관측단위=테이블 (Wickham 2014)"],
        ["벡터화", "C 루프, 연속 메모리, SIMD로 100~1000배 성능 향상"],
        ["결측치 전략", "삭제 / 평균 / 중앙값 / 보간 / 그룹별 -- 상황에 맞게 선택"],
        ["스케일링", "Standard / MinMax / Robust -- train에만 fit, Leakage 주의"],
        ["SVD", "A = U Sigma V^T, PCA의 수학적 기반, 차원 축소의 핵심"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.5, 9.5],
    header_color=ACCENT_BLUE, row_height=0.45, font_size=12)

# ============================================================
# 슬라이드 47: 수식 정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "3장 핵심 수식 정리", "NumPy와 ML에서 사용되는 주요 수식")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.3),
         "선형대수 분해",
         ["SVD: A = U * Sigma * V^T",
          "고유값 분해: Av = lambda * v",
          "LU 분해: A = P * L * U",
          "QR 분해: A = Q * R, Q^T Q = I"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.3),
         "정규방정식과 회귀",
         ["정규방정식: beta = (X^T X)^{-1} X^T y",
          "최소제곱법: min ||Xbeta - y||^2",
          "조건수: kappa(A) = ||A|| * ||A^{-1}||",
          "권장: np.linalg.lstsq(X, y)"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.3),
         "피처 스케일링",
         ["StandardScaler: z = (x - mu) / sigma",
          "MinMaxScaler: z = (x - x_min) / (x_max - x_min)",
          "RobustScaler: z = (x - Q2) / (Q3 - Q1)",
          "Data Leakage: train에만 fit, test에는 transform만!"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.3),
         "이상치 탐지",
         ["IQR: Q1 - 1.5*IQR < x < Q3 + 1.5*IQR",
          "Z-score: |z| = |(x - mu)/sigma| > 3",
          "분산 설명력: s^2 / sum(s^2) * 100",
          "누적 에너지 95% 이상 보존하는 k 선택"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 48: 대용량 데이터 처리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "대용량 데이터 처리 도구", "Pandas의 한계를 넘어서")

add_table_slide(s,
    headers=["라이브러리", "특징", "적합한 상황"],
    rows=[
        ["Dask", "Pandas API 호환, 지연 실행, 분산 처리", "중-대규모 (수~수백 GB)"],
        ["Polars", "Rust 기반, 지연 실행, 매우 빠름", "단일 머신 대규모"],
        ["Vaex", "Out-of-core, 메모리 매핑, 10억 행", "탐색적 분석"],
        ["PySpark", "클러스터 분산 처리", "초대규모 (TB 이상)"],
    ],
    left=Inches(0.6), top=Inches(2.0), col_widths=[2.0, 5.0, 4.0],
    header_color=ACCENT_BLUE, row_height=0.5, font_size=12)

add_code_block(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.5), [
    "# Dask 예시 -- Pandas API와 거의 동일!",
    "import dask.dataframe as dd",
    "ddf = dd.read_csv('large_*.csv')",
    "result = ddf.groupby('col').mean().compute()",
    "",
    "# np.einsum: 아인슈타인 합산 표기법",
    "C = np.einsum('ik,kj->ij', A, B)  # == A @ B",
    "# 배치 행렬곱 (딥러닝 어텐션)",
    "scores = np.einsum('hid,hjd->hij', Q, K)",
], font_size=11)

add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.5),
         "실무 데이터 파이프라인: Titanic 예제",
         ["1. pd.read_csv() -> DataFrame 로딩",
          "2. groupby + transform으로 결측치 그룹별 대체",
          "3. map()으로 범주형 인코딩",
          "4. .values로 NumPy 변환",
          "5. train_test_split으로 분할",
          "6. StandardScaler: train에만 fit!",
          "7. LogisticRegression으로 학습 및 평가"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 49: np.einsum과 Numba
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "고급 도구: np.einsum과 Numba JIT", "텐서 연산과 JIT 컴파일 가속")

add_code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.0), [
    "# np.einsum: 아인슈타인 합산 표기법",
    "A = np.random.randn(3, 4)",
    "B = np.random.randn(4, 5)",
    "",
    "# 행렬곱: C_ij = sum_k A_ik * B_kj",
    "C = np.einsum('ik,kj->ij', A, B)",
    "",
    "# 대각합(Trace): sum_i A_ii",
    "trace = np.einsum('ii->', np.eye(3))",
    "",
    "# 배치 행렬곱 (딥러닝 어텐션 핵심!)",
    "Q = np.random.randn(8, 64, 32)  # (heads, seq, d_k)",
    "K = np.random.randn(8, 64, 32)",
    "scores = np.einsum('hid,hjd->hij', Q, K)",
], font_size=11)

add_code_block(s, Inches(0.6), Inches(5.2), Inches(6.0), Inches(2.0), [
    "# Numba JIT: Python 루프를 C 속도로 가속",
    "from numba import jit",
    "@jit(nopython=True)",
    "def fast_func(X, Y):",
    "    # 벡터화 어려운 복잡한 루프에 적용",
    "    ...",
], font_size=11)

add_card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.5),
         "np.einsum 활용 사례",
         ["행렬곱: 'ik,kj->ij' == A @ B",
          "대각합: 'ii->' == np.trace(A)",
          "외적: 'i,j->ij' == np.outer(a, b)",
          "배치 행렬곱: 'bij,bjk->bik'",
          "Self-Attention: 'hid,hjd->hij'",
          "가독성 뛰어나고 내부 최적화도 우수!"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(7.0), Inches(4.8), Inches(5.6), Inches(2.4),
         "Numba 선택 가이드",
         ["벡터화 가능 -> NumPy가 최선",
          "벡터화 어려운 복잡한 루프 -> Numba JIT",
          "병렬화 필요 -> Numba parallel (prange)",
          "",
          "Numba JIT: 첫 호출 시 컴파일 (warm-up)",
          "이후 호출은 C 수준 속도!"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# ============================================================
# 슬라이드 50: 복습 질문 (1/3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "복습 질문 (1/3)", "NumPy 관련 질문")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(5.2),
         "NumPy 핵심 질문 1~4",
         ["Q1. NumPy ndarray의 stride 개념을 설명하고, 전치(transpose) 연산 시 stride가 어떻게 변하는지 설명하시오.",
          "    -> stride = 각 차원에서 다음 원소까지의 바이트 수. 전치는 stride를 교환 (데이터 복사 없음!)",
          "",
          "Q2. 브로드캐스팅의 세 가지 규칙을 기술하고, (3,1) + (1,4) 연산의 결과 shape을 유도하시오.",
          "    -> 규칙1: 차원 맞추기, 규칙2: 호환성 확인(3vs1, 1vs4), 규칙3: 가상 확장 -> 결과 (3,4)",
          "",
          "Q3. C-order와 F-order의 차이점을 메모리 레이아웃 관점에서 설명하고,",
          "    행 방향 순회 시 어느 것이 더 효율적인지 이유와 함께 설명하시오.",
          "    -> C-order: 행 우선 저장, 행 방향 순회 시 연속 메모리 접근으로 캐시 효율 높음",
          "",
          "Q4. Harris et al.(2020)의 Nature 논문에서 제시한 파이썬 과학 생태계의 레이어 구조를 설명하고,",
          "    NumPy가 기반 계층에 위치하는 이유를 설명하시오.",
          "    -> NumPy는 Pandas, SciPy, scikit-learn 등의 기반. ndarray가 표준 배열 인터페이스 제공."],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 슬라이드 51: 복습 질문 (2/3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "복습 질문 (2/3)", "Pandas 및 데이터 전처리 관련 질문")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(5.2),
         "Pandas/전처리 핵심 질문 5~7",
         ["Q5. Pandas의 loc과 iloc의 차이를 설명하고, 슬라이싱 시 끝 인덱스 포함 여부가 다른 이유를 설명하시오.",
          "    -> loc: 라벨 기반 (끝 포함) / iloc: 정수 위치 기반 (끝 미포함, Python 관례)",
          "",
          "Q6. McKinney(2010)가 Pandas를 개발한 동기를 설명하고, NumPy ndarray가 표 형식 데이터에 부적합한 이유를 기술하시오.",
          "    -> 파이썬에 R의 data.frame 대응 구조 부재. ndarray는 동질적 타입만 지원하므로",
          "       문자열+숫자+날짜가 혼합된 이질적 표 형식 데이터 처리에 부족.",
          "",
          "Q7. 결측치 처리 5가지 전략의 장단점을 비교하고, 시계열 데이터에서 가장 적합한 전략은",
          "    무엇인지 근거와 함께 설명하시오.",
          "    -> 시계열: 보간(interpolate)이 적합. 시간적 연속성을 유지하며 전후 값의 추세를 반영.",
          "       행 삭제는 시간 순서 끊김, 평균 대체는 시간적 패턴 무시하는 문제."],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 52: 복습 질문 (3/3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "복습 질문 (3/3)", "Tidy Data, 벡터화, Data Leakage 관련 질문")

add_card(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(5.2),
         "심화 질문 8~10",
         ["Q8. Wickham(2014)의 Tidy Data 3원칙을 기술하고, 열 헤더가 값인 messy data를",
          "    tidy data로 변환하는 Pandas 코드를 작성하시오.",
          "    -> 원칙: 변수=열, 관측=행, 관측단위=테이블. pd.melt(messy, id_vars='이름', ...)",
          "",
          "Q9. 벡터화 연산이 Python 루프보다 100~1000배 빠른 이유를",
          "    하드웨어 수준(SIMD, CPU 캐시)에서 설명하시오.",
          "    -> SIMD: 하나의 명령으로 여러 데이터 동시 처리 (AVX-256: 8개 float32)",
          "       CPU 캐시: 연속 메모리 접근으로 L1/L2 캐시 적중률 극대화",
          "       + C 컴파일 코드, 타입 체크 생략, GIL 해제",
          "",
          "Q10. Data Leakage란 무엇이며, StandardScaler를 잘못 사용했을 때 어떻게 발생하는지 설명하시오.",
          "    -> 테스트 데이터의 정보가 훈련에 유출되는 현상.",
          "       잘못: scaler.fit_transform(X_all) -> 테스트의 평균/표준편차가 훈련에 반영",
          "       올바름: scaler.fit_transform(X_train) 후 scaler.transform(X_test)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 53: 3장 핵심 키워드 총정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "3장 핵심 키워드 총정리", "한 페이지로 보는 3장의 모든 것")

add_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(2.2),
         "NumPy",
         ["ndarray, stride, dtype",
          "C-order/F-order",
          "브로드캐스팅 3규칙",
          "ufunc, 벡터화",
          "SVD, 고유값 분해, 정규방정식"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(4.7), Inches(2.0), Inches(3.8), Inches(2.2),
         "Pandas",
         ["DataFrame, Series",
          "loc/iloc, query",
          "메서드 체이닝, pipe()",
          "GroupBy: agg/transform/apply",
          "merge/join, pivot/melt"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(2.2),
         "데이터 전처리",
         ["결측치 5가지 전략",
          "이상치: IQR, Z-score",
          "스케일링: Standard/MinMax/Robust",
          "인코딩: Label/OneHot/Ordinal",
          "Data Leakage 방지"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.5), Inches(3.8), Inches(2.2),
         "Tidy Data",
         ["Wickham 2014",
          "변수=열, 관측=행",
          "관측단위=테이블",
          "melt(wide->long)",
          "pivot(long->wide)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

add_card(s, Inches(4.7), Inches(4.5), Inches(3.8), Inches(2.2),
         "벡터화 vs 루프",
         ["100~1000배 속도 향상",
          "SIMD, CPU 캐시",
          "C 컴파일, 타입 체크 생략",
          "GIL 해제, BLAS/LAPACK",
          "Numba JIT 대안"],
         title_color=ACCENT_RED, border=ACCENT_RED)

add_card(s, Inches(8.8), Inches(4.5), Inches(3.8), Inches(2.2),
         "논문 5편",
         ["McKinney (2010): Pandas",
          "Walt et al. (2011): ndarray",
          "Harris et al. (2020): NumPy Nature",
          "Reback et al. (2020): Pandas 공식",
          "Wickham (2014): Tidy Data"],
         title_color=ACCENT_BLUE, border=ACCENT_BLUE)

# ============================================================
# 슬라이드 54: View vs Copy / 메모리 최적화 정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "메모리 관련 정리: View/Copy & 최적화", "NumPy와 Pandas의 메모리 효율성 전략")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "NumPy View vs Copy 정리",
         ["View (메모리 공유): 슬라이싱, 전치, reshape, ravel",
          "Copy (독립 메모리): 팬시인덱싱, 불리언인덱싱, flatten, copy()",
          "확인: np.shares_memory(a, b)",
          "주의: View 수정 시 원본도 변경됨!"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "Pandas 메모리 최적화",
         ["다운캐스팅: int64->int32 (50~75% 절약)",
          "Categorical: object->category (90~99% 절약)",
          "Sparse: 대부분 0인 데이터 (90%+ 절약)",
          "청크 읽기: chunksize 파라미터 (RAM 제한 내)"],
         title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.4),
         "np.memmap: 메모리 매핑 파일",
         ["RAM보다 큰 데이터를 다룰 때 활용",
          "fp = np.memmap('data.dat', dtype='float32', mode='w+', shape=(1000000, 100))",
          "필요한 부분만 메모리에 로드: chunk = fp_read[500:600]",
          "scikit-learn의 일부 알고리즘도 memmap 지원"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# ============================================================
# 슬라이드 55: 실무 팁 & 다음 장 미리보기
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, 9, "실무 팁 & 다음 장 미리보기", "3장에서 배운 것을 4장 이후에 어떻게 활용하는가")

add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "실무 팁 Top 5",
         ["1. 항상 np.linalg.lstsq() 사용 (역행렬 비추천)",
          "2. StandardScaler: train에만 fit! (Data Leakage 방지)",
          "3. 메서드 체이닝 + pipe()로 읽기 쉬운 파이프라인 구성",
          "4. GroupBy transform으로 그룹별 전처리 (평균보다 정확)",
          "5. 벡터화 우선, 벡터화 불가 시에만 Numba JIT"],
         title_color=ACCENT_CYAN, border=ACCENT_CYAN)

add_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
         "3장의 위치: ML 파이프라인에서",
         ["1장: ML 개요 (정의, 학습 유형)",
          "2장: ML 프로젝트 (파이프라인 전체)",
          "3장: 데이터 처리 (NumPy + Pandas) [현재]",
          "4장~: 구체적 ML 알고리즘",
          "",
          "3장은 모든 이후 장의 데이터 처리 기반!"],
         title_color=ACCENT_GREEN, border=ACCENT_GREEN)

add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(2.4),
         "참고 문헌",
         ["1. McKinney, W. (2010). 'Data Structures for Statistical Computing in Python.' SciPy 2010",
          "2. Walt, S. et al. (2011). 'The NumPy Array.' Computing in Science & Engineering, 13(2)",
          "3. Harris, C.R. et al. (2020). 'Array programming with NumPy.' Nature, 585, 357-362",
          "4. Reback, J. et al. (2020). 'pandas-dev/pandas: Pandas.' Zenodo",
          "5. Wickham, H. (2014). 'Tidy Data.' Journal of Statistical Software, 59(10)"],
         title_color=ACCENT_PURPLE, border=ACCENT_PURPLE)

# ============================================================
# 슬라이드 56: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_shape(s, Inches(0), Inches(7.2), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.2), prs.slide_width, Inches(0.5),
         "THANK YOU", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(3.0), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(3.3), prs.slide_width, Inches(0.6),
         "3장: 판다스와 넘파이 (확장 상세 버전)", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.4),
         "NumPy ndarray | Pandas DataFrame | Broadcasting | Vectorization | Tidy Data | Linear Algebra",
         font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.0), prs.slide_width, Inches(0.4),
         "질문이 있으시면 언제든지 질문해 주세요!", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# PPT 저장
# ============================================================
save_path = os.path.join(os.path.dirname(__file__), "3장_판다스와_넘파이_강의PPT_확장.pptx")
prs.save(save_path)
print(f"PPT 저장 완료: {save_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
