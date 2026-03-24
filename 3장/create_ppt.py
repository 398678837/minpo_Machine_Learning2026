"""3장 판다스와 넘파이 - 강의 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xA0, 0x6C, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DARK_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SECTION_BG = RGBColor(0x10, 0x10, 0x28)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_CYAN, border=None):
    bc = border if border else CARD_BG
    card = add_shape(slide, left, top, width, height, CARD_BG, bc, radius=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=15, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))


def slide_header(slide, section_num, title, subtitle=""):
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(1.2), ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(2), Inches(0.4),
             f"SECTION {section_num}" if section_num else "", font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.6),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=DARK_GRAY)


# ============================================================
# 슬라이드 1: 표지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         "기계학습 (Machine Learning)", font_size=20, color=ACCENT_CYAN, bold=True)
add_accent_line(s, Inches(1), Inches(1.85), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "3장: 판다스(Pandas)와 넘파이(NumPy)", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "Data Processing with Pandas & NumPy", font_size=24, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
         "핵심 키워드: ndarray · DataFrame · 브로드캐스팅 · 벡터화 · 선형대수 · GroupBy · Tidy Data · 데이터 전처리",
         font_size=14, color=DARK_GRAY)
add_shape(s, Inches(0), Inches(6.3), prs.slide_width, Inches(1.2), RGBColor(0x0A, 0x0A, 0x1A))
add_text(s, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
         "2026년 1학기", font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 2: 목차
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "목차 (Contents)")

toc = [
    ("01", "도입: 파이썬 과학 생태계와 ML 데이터 흐름", ACCENT_BLUE),
    ("02", "NumPy ndarray 내부 구조 (data, dtype, shape, strides)", ACCENT_CYAN),
    ("03", "C-order vs F-order, dtype 시스템", ACCENT_GREEN),
    ("04", "배열 생성 함수 총정리", ACCENT_PURPLE),
    ("05", "인덱싱 4가지 (기본, 슬라이싱, 팬시, 불리언)", ACCENT_ORANGE),
    ("06", "브로드캐스팅 규칙과 시각적 이해", ACCENT_RED),
    ("07", "유니버설 함수 (ufunc) + 집계 함수", ACCENT_BLUE),
    ("08", "선형대수 연산 (LU, QR, SVD, 고유값)", ACCENT_CYAN),
    ("09", "뷰(View) vs 카피(Copy)", ACCENT_GREEN),
    ("10", "Pandas DataFrame & 인덱싱 + GroupBy/Merge", ACCENT_PURPLE),
    ("11", "데이터 전처리 + Tidy Data + 벡터화 성능", ACCENT_ORANGE),
    ("12", "논문 리뷰 · 실습 · 핵심 요약 · 복습 질문", ACCENT_RED),
]
for i, (num, title, color) in enumerate(toc):
    y = Inches(2.0) + Inches(0.44) * i
    add_shape(s, Inches(1.2), y, Inches(0.55), Inches(0.35), color, radius=True)
    add_text(s, Inches(1.2), y, Inches(0.55), Inches(0.35), num,
             font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), y, Inches(10), Inches(0.35), title,
             font_size=16, color=WHITE)

# ============================================================
# 슬라이드 3: 도입 - 파이썬 과학 생태계 레이어 구조
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "01", "도입: 파이썬 과학 생태계", "Python Scientific Ecosystem & ML Data Flow")

# 레이어 구조
layers = [
    ("Applications", "scikit-learn, TensorFlow, PyTorch, Keras", ACCENT_PURPLE, Inches(2.2)),
    ("Domain-specific", "Pandas (표 형식), SciPy (과학), matplotlib (시각화)", ACCENT_ORANGE, Inches(3.1)),
    ("NumPy", "ndarray, ufunc, broadcasting, linear algebra, FFT, random", ACCENT_BLUE, Inches(4.0)),
    ("Python + C Extensions", "CPython 인터프리터, BLAS, LAPACK", ACCENT_GREEN, Inches(4.9)),
]
for label, desc, color, y_pos in layers:
    add_shape(s, Inches(0.6), y_pos, Inches(5.5), Inches(0.75), CARD_BG, color, radius=True)
    add_text(s, Inches(0.8), y_pos + Inches(0.02), Inches(2.0), Inches(0.35),
             label, font_size=14, color=color, bold=True)
    add_text(s, Inches(0.8), y_pos + Inches(0.35), Inches(5.0), Inches(0.35),
             desc, font_size=11, color=LIGHT_GRAY)

# 화살표들 (아래로)
for i in range(3):
    y = Inches(2.95) + Inches(0.9) * i
    add_shape(s, Inches(3.2), y, Inches(0.15), Inches(0.15), DARK_GRAY)

# ML 파이프라인 데이터 흐름
add_shape(s, Inches(6.6), Inches(2.2), Inches(6.1), Inches(4.8), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(6.9), Inches(2.3), Inches(5.5), Inches(0.4),
         "ML 파이프라인 데이터 흐름", font_size=16, color=ACCENT_CYAN, bold=True)

flow_steps = [
    ("1. Pandas", "pd.read_csv() 로 데이터 로딩", ACCENT_ORANGE),
    ("2. Pandas", "탐색(EDA) + 전처리 (결측치, 인코딩)", ACCENT_ORANGE),
    ("3. NumPy", ".to_numpy() 로 배열 변환", ACCENT_BLUE),
    ("4. sklearn", "StandardScaler, train_test_split", ACCENT_GREEN),
    ("5. 모델링", "model.fit(X_train, y_train) - 내부 NumPy", ACCENT_PURPLE),
]
for i, (step, desc, color) in enumerate(flow_steps):
    y = Inches(2.9) + Inches(0.75) * i
    add_shape(s, Inches(6.9), y, Inches(1.5), Inches(0.55), color, radius=True)
    add_text(s, Inches(6.9), y + Inches(0.05), Inches(1.5), Inches(0.45),
             step, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.6), y + Inches(0.05), Inches(3.8), Inches(0.45),
             desc, font_size=12, color=LIGHT_GRAY)

# 하단 인용
add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
         "\"데이터 과학자의 업무 시간 중 50~80%가 데이터 처리에 소비된다\" - Wickham (2014)",
         font_size=13, color=DARK_GRAY)
add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
         "Harris et al. (2020) Nature: NumPy는 파이썬 과학 생태계의 기반 인프라(foundation infrastructure)",
         font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 4: NumPy ndarray 내부 구조
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "02", "NumPy ndarray 내부 구조", "4가지 핵심 구성 요소 (Walt et al., 2011)")

# 4가지 요소 카드
elements = [
    ("data", "데이터가 저장된\n메모리 버퍼 포인터", "연속 메모리 블록에\n동질적 타입 저장", ACCENT_BLUE),
    ("dtype", "각 원소의\n데이터 타입", "float64, int32 등\n타입 고정 → 성능", ACCENT_CYAN),
    ("shape", "각 차원의\n크기 튜플", "(2, 3) → 2행 3열\n(100, 5) → 100샘플", ACCENT_GREEN),
    ("strides", "다음 원소까지의\n바이트 수", "(24, 8) → 행: 24B\n열: 8B (float64)", ACCENT_ORANGE),
]
for i, (name, desc, detail, color) in enumerate(elements):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(2.2), Inches(2.9), Inches(2.5), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(2.9), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(2.9), Inches(0.45),
             name, font_size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.9),
             desc, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.8), Inches(2.5), Inches(0.8),
             detail, font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 메모리 레이아웃 비교
add_shape(s, Inches(0.6), Inches(5.0), Inches(5.8), Inches(2.2), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(5.1), Inches(5.2), Inches(0.4),
         "Python 리스트 메모리 구조", font_size=14, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(0.9), Inches(5.5), Inches(5.2), Inches(1.5), [
    "리스트 객체 → [포인터1, 포인터2, ...]",
    "각 포인터 → PyObject (28+ bytes)",
    "불연속 메모리, 타입 체크 매번 → 느림",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

add_shape(s, Inches(6.8), Inches(5.0), Inches(5.8), Inches(2.2), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(5.1), Inches(5.2), Inches(0.4),
         "NumPy ndarray 메모리 구조", font_size=14, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(7.1), Inches(5.5), Inches(5.2), Inches(1.5), [
    "ndarray → [8B | 8B | 8B | ...] 연속 블록",
    "동질적 타입, 타입 체크 한 번만",
    "CPU 캐시 친화적, SIMD 활용 가능 → 빠름",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 5: C-order vs F-order + dtype 시스템
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "03", "C-order vs F-order & dtype 시스템")

# C vs F 비교 테이블
headers_cf = ["구분", "C-order (row-major)", "F-order (column-major)"]
rows_cf = [
    ["메모리 레이아웃", "행 우선 저장", "열 우선 저장"],
    ["빠른 접근 방향", "마지막 축 (열 방향)", "첫 번째 축 (행 방향)"],
    ["사용 예", "NumPy 기본값, C 언어", "Fortran, MATLAB, R"],
    ["strides (2x3, int64)", "(24, 8)", "(8, 16)"],
]

for j, h in enumerate(headers_cf):
    x = Inches(0.6) + Inches(3.9) * j
    add_shape(s, x, Inches(2.2), Inches(3.7), Inches(0.5), ACCENT_BLUE)
    add_text(s, x, Inches(2.2), Inches(3.7), Inches(0.5),
             h, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows_cf):
    y = Inches(2.8) + Inches(0.5) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.6) + Inches(3.9) * j
        add_shape(s, x, y, Inches(3.7), Inches(0.45), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y, Inches(3.7), Inches(0.45),
                 cell, font_size=13, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# dtype 시스템
add_text(s, Inches(0.6), Inches(5.0), Inches(5), Inches(0.4),
         "dtype 시스템", font_size=18, color=WHITE, bold=True)

dtype_headers = ["카테고리", "dtype", "바이트", "ML 활용"]
dtype_rows = [
    ["정수", "int8 ~ int64", "1 ~ 8", "인덱스, 레이블"],
    ["부호 없는 정수", "uint8", "1", "이미지 픽셀 (0~255)"],
    ["부동소수점", "float16 / float32", "2 / 4", "딥러닝 추론 / 학습"],
    ["부동소수점", "float64", "8", "과학 계산 (NumPy 기본)"],
    ["불리언", "bool_", "1", "마스크, 불리언 인덱싱"],
]

for j, h in enumerate(dtype_headers):
    x = Inches(0.6) + Inches(3.1) * j
    add_shape(s, x, Inches(5.4), Inches(2.9), Inches(0.45), ACCENT_GREEN)
    add_text(s, x, Inches(5.4), Inches(2.9), Inches(0.45),
             h, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(dtype_rows):
    y = Inches(5.9) + Inches(0.38) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        x = Inches(0.6) + Inches(3.1) * j
        add_shape(s, x, y, Inches(2.9), Inches(0.35), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x, y, Inches(2.9), Inches(0.35),
                 cell, font_size=11, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 6: 배열 생성 함수 총정리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "04", "배열 생성 함수 총정리", "Array Creation Functions")

# 3개 카드
add_card(s, Inches(0.4), Inches(2.2), Inches(3.9), Inches(4.8),
         "기본 / 초기화 생성", [
             "np.array([1, 2, 3])",
             "np.zeros((3, 4))        # 0으로 채움",
             "np.ones((2, 3))          # 1로 채움",
             "np.full((2, 3), 7)       # 7로 채움",
             "np.empty((3, 3))        # 미초기화",
             "np.eye(4)                  # 단위행렬",
             "np.diag([1, 2, 3])      # 대각행렬",
         ], ACCENT_BLUE, ACCENT_BLUE)

add_card(s, Inches(4.6), Inches(2.2), Inches(3.9), Inches(4.8),
         "범위 생성", [
             "np.arange(0, 10, 2)",
             "  → [0, 2, 4, 6, 8]",
             "",
             "np.linspace(0, 1, 5)",
             "  → [0, 0.25, 0.5, 0.75, 1.0]",
             "",
             "np.logspace(0, 3, 4)",
             "  → [1, 10, 100, 1000]",
         ], ACCENT_GREEN, ACCENT_GREEN)

add_card(s, Inches(8.8), Inches(2.2), Inches(3.9), Inches(4.8),
         "난수 생성", [
             "np.random.seed(42)",
             "np.random.rand(3, 4)       # U(0,1)",
             "np.random.randn(3, 4)     # N(0,1)",
             "np.random.randint(0, 10, (3,4))",
             "np.random.normal(0, 1, (100,))",
             "np.random.choice(['a','b'], 5)",
         ], ACCENT_ORANGE, ACCENT_ORANGE)

# ============================================================
# 슬라이드 7: 인덱싱 4가지
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "05", "인덱싱 4가지", "Basic, Slicing, Fancy, Boolean (Harris et al., 2020)")

idx_types = [
    ("1. 기본 인덱싱", "arr[0, 1]  →  2\n\n단일 원소 접근\n차원별 정수 인덱스", ACCENT_BLUE),
    ("2. 슬라이싱", "arr[0:2, 1:3]  →  [[2,3],[6,7]]\narr[:, 2]  →  [3, 7, 11]\n\n뷰(View) 반환 → 메모리 공유", ACCENT_CYAN),
    ("3. 팬시 인덱싱", "arr[[0,2], [1,3]]  →  [2, 12]\n\n배열로 인덱싱\n카피(Copy) 반환", ACCENT_GREEN),
    ("4. 불리언 인덱싱", "arr[arr > 5]\n  →  [6, 7, 8, 9, 10, 11, 12]\n\n조건 마스크 활용\n카피(Copy) 반환", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(idx_types):
    x = Inches(0.3) + Inches(3.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.0), Inches(4.0), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(3.0), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(3.0), Inches(0.45),
             title, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.9), Inches(2.6), Inches(3.0),
             desc, font_size=12, color=LIGHT_GRAY)

# 하단 배열 예시
add_shape(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7), CARD_BG, radius=True)
add_text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.6),
         "arr = np.array([[1,2,3,4], [5,6,7,8], [9,10,11,12]])    # shape (3, 4)",
         font_size=14, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 8: 브로드캐스팅
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "06", "브로드캐스팅 (Broadcasting)", "크기 다른 배열 간 자동 연산 (Walt et al., 2011)")

# 3가지 규칙
rules = [
    ("규칙 1", "차원 수가 다르면, 적은 쪽의\nshape 앞에 1을 추가", ACCENT_BLUE),
    ("규칙 2", "각 차원에서 크기가 같거나\n둘 중 하나가 1이면 호환", ACCENT_GREEN),
    ("규칙 3", "크기 1인 차원을 다른 배열\n크기에 맞게 가상 확장", ACCENT_ORANGE),
]
for i, (rule, desc, color) in enumerate(rules):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(1.3), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), Inches(2.3), Inches(1.2), Inches(0.4),
             rule, font_size=15, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.7), Inches(3.5), Inches(0.7),
             desc, font_size=13, color=LIGHT_GRAY)

# 시각적 예제들
add_shape(s, Inches(0.4), Inches(3.8), Inches(6.0), Inches(3.2), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.7), Inches(3.9), Inches(5.4), Inches(0.4),
         "예제 1: (3,3) + (3,)  →  (3,3)", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(0.7), Inches(4.3), Inches(5.4), Inches(2.5),
         "[a00 a01 a02]     [b0 b1 b2]     [a00+b0  a01+b1  a02+b2]\n"
         "[a10 a11 a12]  +  [b0 b1 b2]  =  [a10+b0  a11+b1  a12+b2]\n"
         "[a20 a21 a22]     [b0 b1 b2]     [a20+b0  a21+b1  a22+b2]",
         font_size=11, color=LIGHT_GRAY, font_name='Consolas')

add_shape(s, Inches(6.7), Inches(3.8), Inches(6.0), Inches(1.5), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.0), Inches(3.9), Inches(5.4), Inches(0.4),
         "예제 2: (3,1) + (1,4)  →  (3,4)", font_size=14, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(7.0), Inches(4.3), Inches(5.4), Inches(0.9),
         "[a0]                         [a0+b0  a0+b1  a0+b2  a0+b3]\n"
         "[a1]  +  [b0 b1 b2 b3]  =   [a1+b0  a1+b1  a1+b2  a1+b3]\n"
         "[a2]                         [a2+b0  a2+b1  a2+b2  a2+b3]",
         font_size=11, color=LIGHT_GRAY, font_name='Consolas')

# ML 활용 예
add_shape(s, Inches(6.7), Inches(5.5), Inches(6.0), Inches(1.5), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.0), Inches(5.6), Inches(5.4), Inches(0.4),
         "ML 활용: 데이터 중심화 & Z-score 표준화", font_size=14, color=ACCENT_ORANGE, bold=True)
add_bullet_list(s, Inches(7.0), Inches(6.0), Inches(5.4), Inches(0.9), [
    "X_centered = X - X.mean(axis=0)    # (100,5) - (5,)",
    "X_std = (X - mean) / std              # 브로드캐스팅 2회",
    "핵심: 실제 데이터 복사 없이 stride=0으로 가상 확장",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 9: ufunc + 집계 함수
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "07", "유니버설 함수 (ufunc) + 집계 함수", "Vectorized Element-wise Operations")

# ufunc 테이블
ufunc_headers = ["분류", "함수 예시", "설명"]
ufunc_rows = [
    ["수학 함수", "np.sin, np.cos, np.exp, np.log, np.sqrt", "원소별 수학 연산"],
    ["비교 함수", "np.greater, np.equal, np.logical_and", "원소별 비교 (불리언 반환)"],
    ["산술 함수", "np.add, np.subtract, np.multiply", "사칙연산 (연산자 오버로딩)"],
    ["집계 함수", "np.sum, np.prod, np.min, np.max, np.mean", "축(axis)별 집계"],
]

for j, h in enumerate(ufunc_headers):
    widths = [Inches(2.0), Inches(6.5), Inches(3.5)]
    x_positions = [Inches(0.6), Inches(2.6), Inches(9.1)]
    add_shape(s, x_positions[j], Inches(2.2), widths[j], Inches(0.5), ACCENT_BLUE)
    add_text(s, x_positions[j], Inches(2.2), widths[j], Inches(0.5),
             h, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(ufunc_rows):
    y = Inches(2.8) + Inches(0.55) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        widths = [Inches(2.0), Inches(6.5), Inches(3.5)]
        x_positions = [Inches(0.6), Inches(2.6), Inches(9.1)]
        add_shape(s, x_positions[j], y, widths[j], Inches(0.5), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x_positions[j], y, widths[j], Inches(0.5),
                 cell, font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# axis 시각적 설명
add_shape(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4),
         "axis 파라미터 이해: np.sum(A, axis=?)", font_size=16, color=ACCENT_GREEN, bold=True)

axis_items = [
    ("axis=None", "전체 합: 21", "A = [[1,2,3],\n      [4,5,6]]", ACCENT_BLUE),
    ("axis=0", "열 방향 합:\n[5, 7, 9]", "행을 따라 축소\n(위→아래 접기)", ACCENT_CYAN),
    ("axis=1", "행 방향 합:\n[6, 15]", "열을 따라 축소\n(왼→오 접기)", ACCENT_ORANGE),
]
for i, (ax, result, desc, color) in enumerate(axis_items):
    x = Inches(0.9) + Inches(4.0) * i
    add_shape(s, x, Inches(5.8), Inches(3.6), Inches(1.2), RGBColor(0x2D, 0x2D, 0x45), color, radius=True)
    add_text(s, x + Inches(0.15), Inches(5.85), Inches(1.2), Inches(0.35),
             ax, font_size=14, color=color, bold=True)
    add_text(s, x + Inches(1.4), Inches(5.85), Inches(2.0), Inches(0.5),
             result, font_size=12, color=WHITE)
    add_text(s, x + Inches(0.15), Inches(6.4), Inches(3.3), Inches(0.5),
             desc, font_size=10, color=DARK_GRAY)

# ============================================================
# 슬라이드 10: 선형대수 연산 + ML 활용
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "08", "선형대수 연산과 ML 활용", "Linear Algebra for Machine Learning")

# 4대 분해
decomps = [
    ("LU 분해", "A = P·L·U\n하삼각 × 상삼각", "연립방정식, 행렬식\n역행렬 효율 계산", "scipy.linalg.lu()", ACCENT_BLUE),
    ("QR 분해", "A = Q·R\nQ^T Q = I (직교)", "최소제곱법의\n수치 안정적 풀이", "np.linalg.qr()", ACCENT_CYAN),
    ("SVD", "A = U·Σ·V^T\n특이값 분해", "PCA, 추천시스템\nLSA, 데이터 압축", "np.linalg.svd()", ACCENT_GREEN),
    ("고유값 분해", "Av = λv\n고유값·고유벡터", "PCA (공분산행렬)\n스펙트럴 클러스터링", "np.linalg.eigh()", ACCENT_ORANGE),
]
for i, (name, formula, ml_use, func, color) in enumerate(decomps):
    x = Inches(0.3) + Inches(3.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.0), Inches(2.8), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(2.3), Inches(2.7), Inches(0.4),
             name, font_size=15, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(2.7), Inches(2.7), Inches(0.8),
             formula, font_size=12, color=WHITE)
    add_accent_line(s, x + Inches(0.2), Inches(3.5), Inches(2.5), color)
    add_text(s, x + Inches(0.15), Inches(3.6), Inches(2.7), Inches(0.7),
             ml_use, font_size=12, color=LIGHT_GRAY)
    add_text(s, x + Inches(0.15), Inches(4.4), Inches(2.7), Inches(0.4),
             func, font_size=10, color=DARK_GRAY)

# ML 관계 테이블
linalg_headers = ["선형대수 연산", "ML 활용", "NumPy 함수"]
linalg_rows = [
    ["SVD", "PCA, 추천 시스템, LSA", "np.linalg.svd()"],
    ["고유값 분해", "PCA, 스펙트럴 클러스터링", "np.linalg.eigh()"],
    ["정규방정식", "선형 회귀 (β = (X^TX)^-1 X^Ty)", "np.linalg.lstsq()"],
    ["행렬곱", "신경망 순전파", "@ 연산자"],
    ["노름", "L1, L2 규제", "np.linalg.norm()"],
]

for j, h in enumerate(linalg_headers):
    widths2 = [Inches(2.5), Inches(5.0), Inches(4.3)]
    x_pos2 = [Inches(0.6), Inches(3.1), Inches(8.1)]
    add_shape(s, x_pos2[j], Inches(5.2), widths2[j], Inches(0.45), ACCENT_PURPLE)
    add_text(s, x_pos2[j], Inches(5.2), widths2[j], Inches(0.45),
             h, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(linalg_rows):
    y = Inches(5.7) + Inches(0.38) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        widths2 = [Inches(2.5), Inches(5.0), Inches(4.3)]
        x_pos2 = [Inches(0.6), Inches(3.1), Inches(8.1)]
        add_shape(s, x_pos2[j], y, widths2[j], Inches(0.35), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x_pos2[j], y, widths2[j], Inches(0.35),
                 cell, font_size=11, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 11: 뷰(View) vs 카피(Copy)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "09", "뷰(View) vs 카피(Copy)", "메모리 공유 여부에 따른 동작 차이")

# 설명
add_shape(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.4),
         "뷰 (View) - 메모리 공유", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(1.3), [
    "b = a[1:4]    # 슬라이싱 → 뷰",
    "b[0] = 99     # 원본도 변경됨!",
    "효율적이지만 의도치 않은 변경 주의",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

add_shape(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.0), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
         "카피 (Copy) - 독립 메모리", font_size=16, color=ACCENT_RED, bold=True)
add_bullet_list(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(1.3), [
    "c = a[1:4].copy()  # 명시적 복사",
    "c[0] = 100   # 원본 불변",
    "메모리 추가 사용, 안전한 수정 가능",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 테이블
vc_headers = ["연산", "반환 유형", "메모리 공유"]
vc_rows = [
    ["슬라이싱  a[1:4]", "뷰", "O"],
    ["전치  a.T", "뷰", "O"],
    ["reshape()  (가능 시)", "뷰", "O"],
    ["ravel()  (가능 시)", "뷰", "O"],
    ["팬시 인덱싱  a[[0,2]]", "카피", "X"],
    ["불리언 인덱싱  a[mask]", "카피", "X"],
    [".copy()", "카피", "X"],
    [".flatten()", "카피", "X"],
]

for j, h in enumerate(vc_headers):
    widths3 = [Inches(5.0), Inches(3.5), Inches(3.5)]
    x_pos3 = [Inches(0.6), Inches(5.6), Inches(9.1)]
    add_shape(s, x_pos3[j], Inches(4.5), widths3[j], Inches(0.45), ACCENT_BLUE)
    add_text(s, x_pos3[j], Inches(4.5), widths3[j], Inches(0.45),
             h, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(vc_rows):
    y = Inches(5.0) + Inches(0.34) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        widths3 = [Inches(5.0), Inches(3.5), Inches(3.5)]
        x_pos3 = [Inches(0.6), Inches(5.6), Inches(9.1)]
        add_shape(s, x_pos3[j], y, widths3[j], Inches(0.3), bg)
        if j == 1:
            fc = ACCENT_GREEN if cell == "뷰" else ACCENT_RED
        elif j == 2:
            fc = ACCENT_GREEN if cell == "O" else ACCENT_RED
        else:
            fc = LIGHT_GRAY
        add_text(s, x_pos3[j], y, widths3[j], Inches(0.3),
                 cell, font_size=11, color=fc, bold=(j > 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 12: Pandas DataFrame 내부 구조
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "Pandas DataFrame 내부 구조", "McKinney (2010): 이질적 표 형식 데이터 처리")

# 설명
add_shape(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(2.3), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(2.3), Inches(5.4), Inches(0.4),
         "DataFrame = 열별 독립 NumPy 배열", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_list(s, Inches(0.9), Inches(2.7), Inches(5.4), Inches(1.6), [
    "각 열이 서로 다른 dtype 가능 (이질적)",
    "내부: 열마다 독립적 NumPy 배열 (또는 ExtensionArray)",
    "행 인덱스(Index) + 열 이름(Columns) → 라벨 기반 접근",
    "NumPy ndarray와 달리 문자열+숫자+날짜 혼합 가능",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# 핵심 속성 테이블
add_shape(s, Inches(6.9), Inches(2.2), Inches(5.8), Inches(2.3), CARD_BG, ACCENT_ORANGE, radius=True)
add_text(s, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.4),
         "핵심 속성", font_size=16, color=ACCENT_ORANGE, bold=True)

attrs = [
    ("df.shape", "(행 수, 열 수)"),
    ("df.dtypes", "각 열의 데이터 타입"),
    ("df.index", "행 인덱스"),
    ("df.columns", "열 이름"),
    ("df.info()", "전반적 정보 (행/열, dtype, 메모리)"),
    ("df.describe()", "기술통계량 (평균, 표준편차, 사분위수)"),
]
for i, (attr, desc) in enumerate(attrs):
    y = Inches(2.8) + Inches(0.28) * i
    add_text(s, Inches(7.2), y, Inches(2.2), Inches(0.25),
             attr, font_size=11, color=ACCENT_CYAN, bold=True)
    add_text(s, Inches(9.5), y, Inches(3.0), Inches(0.25),
             desc, font_size=11, color=LIGHT_GRAY)

# 코드 예시
add_shape(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.3), CARD_BG, radius=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.4),
         "DataFrame 생성 예시", font_size=14, color=WHITE, bold=True)
add_text(s, Inches(0.9), Inches(5.3), Inches(7.0), Inches(1.6),
         "df = pd.DataFrame({\n"
         "    '이름': ['김철수', '이영희', '박민수'],\n"
         "    '나이': [28, 35, 42],\n"
         "    '키': [175.5, 162.3, 180.1],\n"
         "    '학생': [True, False, False]\n"
         "})",
         font_size=12, color=LIGHT_GRAY, font_name='Consolas')

add_text(s, Inches(8.2), Inches(5.3), Inches(4.2), Inches(1.6),
         "df.dtypes:\n"
         "  이름     object\n"
         "  나이      int64\n"
         "  키      float64\n"
         "  학생       bool",
         font_size=12, color=ACCENT_GREEN, font_name='Consolas')

# ============================================================
# 슬라이드 13: Pandas 인덱싱 4가지 + 메서드 체이닝
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "Pandas 인덱싱 & 메서드 체이닝")

# 4가지 인덱싱
pd_idx = [
    ("1. 열 선택", "df['나이']           # Series\ndf[['이름', '나이']]   # DataFrame", ACCENT_BLUE),
    ("2. loc (라벨)", "df.loc[0, '이름']    # 라벨 기반\ndf.loc[0:1, ['이름']]  # 끝 포함", ACCENT_CYAN),
    ("3. iloc (정수)", "df.iloc[0, 0]        # 정수 위치\ndf.iloc[0:2, 0:2]   # 끝 미포함", ACCENT_GREEN),
    ("4. 불리언", "df[df['나이'] > 30]\ndf.query('나이 > 30 and 키 > 170')", ACCENT_ORANGE),
]
for i, (title, code, color) in enumerate(pd_idx):
    x = Inches(0.3) + Inches(3.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.0), Inches(1.8), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(2.3), Inches(2.7), Inches(0.35),
             title, font_size=13, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(2.7), Inches(2.7), Inches(1.2),
             code, font_size=10, color=LIGHT_GRAY, font_name='Consolas')

# 메서드 체이닝
add_shape(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(4.4), Inches(5.2), Inches(0.4),
         "메서드 체이닝 (Method Chaining)", font_size=16, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(0.9), Inches(4.8), Inches(5.2), Inches(2.1),
         "# 중간 변수 없이 연속 변환\n"
         "result = (df\n"
         "    .dropna()\n"
         "    .query('나이 > 25')\n"
         "    .assign(나이대=lambda x: x['나이']//10*10)\n"
         "    .sort_values('키', ascending=False)\n"
         ")",
         font_size=12, color=LIGHT_GRAY, font_name='Consolas')

# pipe() 설명
add_shape(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.8), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(7.1), Inches(4.4), Inches(5.2), Inches(0.4),
         "pipe() - 커스텀 함수 체이닝", font_size=16, color=ACCENT_RED, bold=True)
add_text(s, Inches(7.1), Inches(4.8), Inches(5.2), Inches(2.1),
         "def 이상치_제거(df, col, n_std=3):\n"
         "    mean = df[col].mean()\n"
         "    std = df[col].std()\n"
         "    return df[abs(df[col]-mean) <= n_std*std]\n"
         "\n"
         "result = (df\n"
         "    .pipe(이상치_제거, '키')\n"
         "    .assign(BMI=lambda x: x['키']/100)\n"
         ")",
         font_size=11, color=LIGHT_GRAY, font_name='Consolas')

# ============================================================
# 슬라이드 14: GroupBy: 분할-적용-결합
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "GroupBy: 분할-적용-결합", "Split-Apply-Combine (McKinney, 2010)")

# 3가지 패턴
gb_patterns = [
    ("agg", "그룹별 집계", [
        "df.groupby('부서')['연봉']",
        "  .agg(['mean', 'std', 'min', 'max'])",
        "",
        "여러 통계량을 한 번에 계산",
        "결과: 그룹 수 x 집계함수 수",
    ], ACCENT_BLUE),
    ("transform", "원본 크기 유지", [
        "df['연봉_zscore'] = ",
        "  df.groupby('부서')['연봉']",
        "  .transform(",
        "    lambda x: (x-x.mean())/x.std())",
        "결과 크기 = 원본 크기 (브로드캐스트 패턴)",
    ], ACCENT_GREEN),
    ("apply", "자유도 높은 연산", [
        "def 상위N(group, n=3):",
        "    return group.nlargest(n, '성과')",
        "",
        "df.groupby('부서')",
        "  .apply(상위N, n=2)",
    ], ACCENT_ORANGE),
]
for i, (name, desc, items, color) in enumerate(gb_patterns):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(4.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(0.9), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(3.9), Inches(0.4),
             name, font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.65), Inches(3.9), Inches(0.4),
             desc, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_bullet_list(s, x + Inches(0.2), Inches(3.3), Inches(3.5), Inches(3.5),
                    items, font_size=11, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 15: Merge/Join
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "Merge와 Join", "SQL 스타일 데이터 결합")

# 4가지 Join 유형
join_headers = ["Join 유형", "SQL 대응", "설명", "결과 행 수"]
join_rows = [
    ["inner", "INNER JOIN", "양쪽 모두 키 있는 행만", "교집합"],
    ["left", "LEFT JOIN", "왼쪽 테이블 기준", "왼쪽 전체"],
    ["right", "RIGHT JOIN", "오른쪽 테이블 기준", "오른쪽 전체"],
    ["outer", "FULL OUTER JOIN", "양쪽 모두 포함", "합집합"],
]

for j, h in enumerate(join_headers):
    widths_j = [Inches(2.5), Inches(3.0), Inches(3.5), Inches(2.8)]
    x_j = [Inches(0.6), Inches(3.1), Inches(6.1), Inches(9.6)]
    add_shape(s, x_j[j], Inches(2.2), widths_j[j], Inches(0.55), ACCENT_BLUE)
    add_text(s, x_j[j], Inches(2.2), widths_j[j], Inches(0.55),
             h, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(join_rows):
    y = Inches(2.85) + Inches(0.55) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        widths_j = [Inches(2.5), Inches(3.0), Inches(3.5), Inches(2.8)]
        x_j = [Inches(0.6), Inches(3.1), Inches(6.1), Inches(9.6)]
        add_shape(s, x_j[j], y, widths_j[j], Inches(0.5), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x_j[j], y, widths_j[j], Inches(0.5),
                 cell, font_size=13, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 코드 예시
add_shape(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(2.1), CARD_BG, ACCENT_CYAN, radius=True)
add_text(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.4),
         "Merge 코드 예시", font_size=14, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(0.9), Inches(5.6), Inches(5.5), Inches(1.4),
         "직원 = pd.DataFrame({\n"
         "    '사번': [1, 2, 3, 4],\n"
         "    '이름': ['김철수','이영희','박민수','정수진'],\n"
         "    '부서코드': ['D01','D02','D01','D03']\n"
         "})",
         font_size=11, color=LIGHT_GRAY, font_name='Consolas')
add_text(s, Inches(6.8), Inches(5.6), Inches(5.5), Inches(1.4),
         "부서 = pd.DataFrame({\n"
         "    '코드': ['D01','D02','D03'],\n"
         "    '부서명': ['개발팀','마케팅팀','영업팀']\n"
         "})\n"
         "result = pd.merge(직원, 부서,\n"
         "    left_on='부서코드', right_on='코드', how='left')",
         font_size=11, color=LIGHT_GRAY, font_name='Consolas')

# ============================================================
# 슬라이드 16: Window Functions + pivot_table
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "10", "윈도우 함수 & 피벗 테이블", "Rolling, Expanding, EWM, pivot_table")

# 3가지 윈도우
win_funcs = [
    ("rolling(n)", "고정 크기 n 윈도우", [
        "df['MA_7'] = df['주가'].rolling(7).mean()",
        "7일/30일 이동평균 → 골든/데드 크로스",
        "이동 표준편차 → 변동성 추정",
    ], ACCENT_BLUE),
    ("expanding()", "시작 ~ 현재 누적", [
        "df['누적평균'] = df['주가'].expanding().mean()",
        "df['누적최대'] = df['주가'].expanding().max()",
        "역대 최고가, 누적 수익률 계산",
    ], ACCENT_GREEN),
    ("ewm(span=n)", "지수 가중 이동평균", [
        "df['EWM_12'] = df['주가'].ewm(span=12).mean()",
        "최근 데이터에 더 큰 가중치",
        "MACD = EWM_12 - EWM_26",
    ], ACCENT_ORANGE),
]
for i, (name, desc, items, color) in enumerate(win_funcs):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(2.5), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(2.3), Inches(3.6), Inches(0.35),
             name, font_size=15, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(2.6), Inches(3.6), Inches(0.3),
             desc, font_size=12, color=WHITE)
    add_bullet_list(s, x + Inches(0.15), Inches(2.95), Inches(3.6), Inches(1.5),
                    items, font_size=11, color=LIGHT_GRAY, spacing=Pt(3))

# pivot_table
add_shape(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2), CARD_BG, ACCENT_PURPLE, radius=True)
add_text(s, Inches(0.9), Inches(5.1), Inches(5.0), Inches(0.4),
         "pivot_table - SQL GROUP BY + CASE WHEN", font_size=15, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(0.9), Inches(5.5), Inches(5.5), Inches(1.5),
         "피벗 = pd.pivot_table(\n"
         "    매출,\n"
         "    values='매출액',\n"
         "    index='지역',         # 행 기준\n"
         "    columns='제품',       # 열 기준\n"
         "    aggfunc='mean',      # 집계 함수\n"
         "    margins=True          # 합계 행/열\n"
         ")",
         font_size=11, color=LIGHT_GRAY, font_name='Consolas')
add_bullet_list(s, Inches(6.8), Inches(5.5), Inches(5.5), Inches(1.5), [
    "다중 집계: aggfunc={'매출액': ['mean','sum'], '수량': 'sum'}",
    "교차표: pd.crosstab(매출['지역'], 매출['제품'])",
    "",
    "Tidy Data 관점: pivot_table 결과 = wide format",
    "분석 후 melt()로 long format(tidy)으로 재변환 가능",
], font_size=11, color=LIGHT_GRAY, spacing=Pt(3))

# ============================================================
# 슬라이드 17: 데이터 전처리 - 결측치, 이상치, 스케일링
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "11", "데이터 전처리", "결측치 5가지, 이상치 3가지, 스케일링 3가지")

# 결측치 처리 5가지
missing_headers = ["전략", "방법", "장점", "적합한 상황"]
missing_rows = [
    ["행 삭제", "dropna()", "단순함", "결측 < 5%"],
    ["평균 대체", "fillna(mean)", "분포 유지", "정규분포, MCAR"],
    ["중앙값 대체", "fillna(median)", "이상치 강건", "편향 분포"],
    ["보간", "interpolate()", "연속성 유지", "시계열 데이터"],
    ["그룹별 대체", "groupby.transform", "그룹 특성 반영", "그룹 간 차이 큼"],
]

add_text(s, Inches(0.6), Inches(2.0), Inches(5), Inches(0.3),
         "결측치 처리 5가지", font_size=15, color=ACCENT_BLUE, bold=True)

for j, h in enumerate(missing_headers):
    w_m = [Inches(1.8), Inches(2.5), Inches(2.0), Inches(2.5)]
    x_m = [Inches(0.6), Inches(2.4), Inches(4.9), Inches(6.9)]
    add_shape(s, x_m[j], Inches(2.3), w_m[j], Inches(0.35), ACCENT_BLUE)
    add_text(s, x_m[j], Inches(2.3), w_m[j], Inches(0.35),
             h, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(missing_rows):
    y = Inches(2.7) + Inches(0.32) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        w_m = [Inches(1.8), Inches(2.5), Inches(2.0), Inches(2.5)]
        x_m = [Inches(0.6), Inches(2.4), Inches(4.9), Inches(6.9)]
        add_shape(s, x_m[j], y, w_m[j], Inches(0.28), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x_m[j], y, w_m[j], Inches(0.28),
                 cell, font_size=10, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 이상치 탐지
add_text(s, Inches(0.6), Inches(4.4), Inches(5), Inches(0.3),
         "이상치 탐지 3가지", font_size=15, color=ACCENT_RED, bold=True)

outlier_cards = [
    ("IQR 방법", "Q1 - 1.5*IQR < x < Q3 + 1.5*IQR\n상자 수염 그림의 경계", ACCENT_RED),
    ("Z-score", "|z| = |x - mu| / sigma > 3\n정규분포 가정", ACCENT_ORANGE),
    ("Isolation Forest", "이상치는 고립되기 쉬움\n트리 기반 비지도 학습", ACCENT_PURPLE),
]
for i, (name, desc, color) in enumerate(outlier_cards):
    x = Inches(0.4) + Inches(3.2) * i
    add_shape(s, x, Inches(4.7), Inches(3.0), Inches(1.0), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(4.75), Inches(1.2), Inches(0.3),
             name, font_size=12, color=color, bold=True)
    add_text(s, x + Inches(0.15), Inches(5.05), Inches(2.7), Inches(0.6),
             desc, font_size=10, color=LIGHT_GRAY)

# 스케일링 3가지
add_text(s, Inches(0.6), Inches(5.9), Inches(5), Inches(0.3),
         "피처 스케일링 3가지", font_size=15, color=ACCENT_GREEN, bold=True)

scale_headers = ["스케일러", "수식", "특성", "사용 상황"]
scale_rows = [
    ["StandardScaler", "z = (x - mu) / sigma", "평균 0, 분산 1", "SVM, 로지스틱"],
    ["MinMaxScaler", "z = (x - min) / (max - min)", "[0, 1] 범위", "신경망, 이미지"],
    ["RobustScaler", "z = (x - Q2) / (Q3 - Q1)", "중앙값/IQR 기반", "이상치 존재 시"],
]

for j, h in enumerate(scale_headers):
    w_s = [Inches(2.2), Inches(3.8), Inches(2.5), Inches(3.3)]
    x_s = [Inches(0.6), Inches(2.8), Inches(6.6), Inches(9.1)]
    add_shape(s, x_s[j], Inches(6.2), w_s[j], Inches(0.35), ACCENT_GREEN)
    add_text(s, x_s[j], Inches(6.2), w_s[j], Inches(0.35),
             h, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(scale_rows):
    y = Inches(6.6) + Inches(0.3) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        w_s = [Inches(2.2), Inches(3.8), Inches(2.5), Inches(3.3)]
        x_s = [Inches(0.6), Inches(2.8), Inches(6.6), Inches(9.1)]
        add_shape(s, x_s[j], y, w_s[j], Inches(0.27), bg)
        fc = ACCENT_CYAN if j == 0 else LIGHT_GRAY
        add_text(s, x_s[j], y, w_s[j], Inches(0.27),
                 cell, font_size=10, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 18: Tidy Data 원칙
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "11", "Tidy Data 원칙", "Wickham (2014) - Journal of Statistical Software")

# 3원칙
tidy_principles = [
    ("원칙 1", "각 변수는 하나의 열을 구성", "변수 = 열 (Column)", ACCENT_BLUE),
    ("원칙 2", "각 관측은 하나의 행을 구성", "관측 = 행 (Row)", ACCENT_GREEN),
    ("원칙 3", "각 관측 단위 유형은 하나의 테이블", "관측단위 = 테이블", ACCENT_ORANGE),
]
for i, (name, desc, short, color) in enumerate(tidy_principles):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(1.2), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.15), Inches(2.3), Inches(1.2), Inches(0.35),
             name, font_size=15, color=color, bold=True)
    add_text(s, x + Inches(1.4), Inches(2.3), Inches(2.3), Inches(0.35),
             short, font_size=13, color=WHITE, bold=True)
    add_text(s, x + Inches(0.15), Inches(2.7), Inches(3.6), Inches(0.5),
             desc, font_size=13, color=LIGHT_GRAY)

# Messy → Tidy 변환 예시
add_shape(s, Inches(0.6), Inches(3.7), Inches(5.8), Inches(3.3), CARD_BG, ACCENT_RED, radius=True)
add_text(s, Inches(0.9), Inches(3.8), Inches(5.2), Inches(0.4),
         "Messy Data (열 헤더가 값)", font_size=14, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.9), Inches(4.2), Inches(5.2), Inches(2.6),
         "messy = pd.DataFrame({\n"
         "    '이름': ['김철수', '이영희'],\n"
         "    '2023_국어': [85, 92],\n"
         "    '2023_수학': [90, 88],\n"
         "    '2024_국어': [88, 95],\n"
         "    '2024_수학': [92, 90]\n"
         "})",
         font_size=11, color=LIGHT_GRAY, font_name='Consolas')

add_shape(s, Inches(6.8), Inches(3.7), Inches(5.8), Inches(3.3), CARD_BG, ACCENT_GREEN, radius=True)
add_text(s, Inches(7.1), Inches(3.8), Inches(5.2), Inches(0.4),
         "Tidy Data (melt 변환)", font_size=14, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(7.1), Inches(4.2), Inches(5.2), Inches(2.6),
         "tidy = (messy\n"
         "    .melt(id_vars='이름',\n"
         "          var_name='과목_연도',\n"
         "          value_name='점수')\n"
         "    .assign(\n"
         "        연도=lambda x: x['과목_연도'].str[:4],\n"
         "        과목=lambda x: x['과목_연도'].str[5:])\n"
         "    .drop(columns='과목_연도')\n"
         ")",
         font_size=11, color=LIGHT_GRAY, font_name='Consolas')

# 실무적 이점
add_shape(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35), CARD_BG, radius=True)
add_text(s, Inches(0.9), Inches(7.05), Inches(11.5), Inches(0.35),
         "Tidy Data 이점: groupby/agg 자연 동작 | seaborn 호환 | ML 파이프라인 직접 투입 | scikit-learn: 행=샘플, 열=특성",
         font_size=11, color=DARK_GRAY)

# ============================================================
# 슬라이드 19: 벡터화 vs 루프 성능 비교
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "11", "벡터화 vs 루프: 성능의 과학", "왜 NumPy가 100~1000배 빠른가? (Harris et al., 2020)")

# 비교 테이블
perf_headers = ["요인", "Python 루프", "NumPy 벡터화"]
perf_rows = [
    ["실행 엔진", "Python 인터프리터", "C/Fortran 컴파일 코드"],
    ["타입 체크", "매 연산마다", "한 번만"],
    ["메모리 접근", "불연속 (포인터 추적)", "연속 (캐시 친화적)"],
    ["SIMD 활용", "불가", "가능 (SSE, AVX)"],
    ["GIL", "보유", "해제 가능"],
    ["BLAS/LAPACK", "미사용", "활용 (Intel MKL 등)"],
]

for j, h in enumerate(perf_headers):
    w_p = [Inches(2.5), Inches(4.0), Inches(4.5)]
    x_p = [Inches(0.6), Inches(3.1), Inches(7.1)]
    color_h = DARK_GRAY if j == 0 else (ACCENT_RED if j == 1 else ACCENT_GREEN)
    add_shape(s, x_p[j], Inches(2.2), w_p[j], Inches(0.5), color_h)
    add_text(s, x_p[j], Inches(2.2), w_p[j], Inches(0.5),
             h, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(perf_rows):
    y = Inches(2.8) + Inches(0.48) * i
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, cell in enumerate(row):
        w_p = [Inches(2.5), Inches(4.0), Inches(4.5)]
        x_p = [Inches(0.6), Inches(3.1), Inches(7.1)]
        add_shape(s, x_p[j], y, w_p[j], Inches(0.43), bg)
        fc = ACCENT_CYAN if j == 0 else (ACCENT_RED if j == 1 else ACCENT_GREEN)
        add_text(s, x_p[j], y, w_p[j], Inches(0.43),
                 cell, font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)

# 성능 수치
add_shape(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.4), CARD_BG, ACCENT_BLUE, radius=True)
add_text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.4),
         "벤치마크: 100만 원소 덧셈 (a + b)", font_size=15, color=ACCENT_BLUE, bold=True)

bench_items = [
    ("Python 루프", "for i in range(n): c[i] = a[i]+b[i]", "~300ms", "1x", ACCENT_RED),
    ("List Comprehension", "[a[i]+b[i] for i in range(n)]", "~200ms", "1.5x", ACCENT_ORANGE),
    ("NumPy 벡터화", "c = a + b", "~0.3ms", "1000x", ACCENT_GREEN),
]
for i, (method, code, time_val, speedup, color) in enumerate(bench_items):
    x = Inches(0.9) + Inches(4.0) * i
    add_text(s, x, Inches(6.3), Inches(1.8), Inches(0.3),
             method, font_size=12, color=color, bold=True)
    add_text(s, x + Inches(1.9), Inches(6.3), Inches(1.0), Inches(0.3),
             time_val, font_size=12, color=WHITE, bold=True)
    add_text(s, x + Inches(2.9), Inches(6.3), Inches(0.8), Inches(0.3),
             speedup, font_size=12, color=color, bold=True)

# ============================================================
# 슬라이드 20: 논문 리뷰 5편
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "12", "핵심 논문 리뷰 (5편)", "Key Paper Reviews")

papers = [
    ("McKinney\n(2010)", "Data Structures for\nStatistical Computing\nin Python",
     "Pandas 탄생\nDataFrame, GroupBy\nSplit-Apply-Combine", ACCENT_BLUE),
    ("Walt et al.\n(2011)", "The NumPy Array:\nEfficient Numerical\nComputation",
     "ndarray 아키텍처\nstride, ufunc\n브로드캐스팅 문서화", ACCENT_CYAN),
    ("Harris et al.\n(2020)", "Array programming\nwith NumPy\n(Nature)",
     "배열 프로그래밍 패러다임\n생태계 레이어 시각화\nNumPy API 프로토콜", ACCENT_GREEN),
    ("Reback et al.\n(2020)", "pandas-dev/pandas\n(Zenodo)",
     "Pandas 공식 인용\nApache Arrow 백엔드\nCSV~Parquet IO 지원", ACCENT_ORANGE),
    ("Wickham\n(2014)", "Tidy Data\n(J. Stat. Software)",
     "Tidy Data 3원칙\nMessy Data 5유형\nmelt/pivot 체계화", ACCENT_RED),
]
for i, (author, title, contrib, color) in enumerate(papers):
    x = Inches(0.2) + Inches(2.6) * i
    add_shape(s, x, Inches(2.2), Inches(2.4), Inches(4.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(2.4), Inches(0.7), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(2.4), Inches(0.6),
             author, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.1), Inches(2.1), Inches(1.3),
             title, font_size=11, color=LIGHT_GRAY)
    add_accent_line(s, x + Inches(0.2), Inches(4.4), Inches(2.0), color)
    add_text(s, x + Inches(0.15), Inches(4.6), Inches(2.1), Inches(2.0),
             contrib, font_size=11, color=DARK_GRAY)

# ============================================================
# 슬라이드 21: 실습 소개
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "12", "실습 소개 (3개 구현 소스)")

labs = [
    ("실습 1", "01_numpy_linear_algebra.py", "NumPy 선형대수", [
        "LU / QR / SVD 분해 구현",
        "고유값 분해와 PCA 연결",
        "정규방정식 3가지 풀이 비교",
        "  inv(비추천) vs solve vs lstsq(권장)",
        "SVD 차원 축소: 정보 보존 비율 분석",
    ], ACCENT_BLUE),
    ("실습 2", "02_pandas_pipeline.py", "Pandas 데이터 파이프라인", [
        "메서드 체이닝 + pipe() 실전 패턴",
        "결측치 5가지 전략 비교 실험",
        "GroupBy + transform (그룹별 Z-score)",
        "Rolling 윈도우 함수 활용",
        "Tidy Data 변환 (melt/pivot)",
    ], ACCENT_GREEN),
    ("실습 3", "03_vectorization_benchmark.py", "벡터화 벤치마크", [
        "쌍별 유클리드 거리 3중 루프 vs 벡터화",
        "행렬곱 성능 비교",
        "기술통계량 계산 벤치마크",
        "결과: 벡터화 500~1000배 속도 향상",
        "KNN, K-Means 등 거리 알고리즘 필수",
    ], ACCENT_ORANGE),
]
for i, (title, file, desc, items, color) in enumerate(labs):
    x = Inches(0.4) + Inches(4.2) * i
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(4.8), CARD_BG, color, radius=True)
    add_shape(s, x, Inches(2.2), Inches(3.9), Inches(0.55), color, radius=True)
    add_text(s, x, Inches(2.25), Inches(3.9), Inches(0.5),
             title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.9), Inches(3.5), Inches(0.3),
             file, font_size=11, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(0.3),
             desc, font_size=13, color=WHITE)
    add_bullet_list(s, x + Inches(0.2), Inches(3.6), Inches(3.5), Inches(3.0),
                    items, font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

# ============================================================
# 슬라이드 22: 핵심 요약
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "핵심 요약 (Key Takeaways)")

summaries = [
    ("ndarray", "연속 메모리, stride, 동질적 dtype\n벡터화 연산의 기반 인프라", ACCENT_BLUE),
    ("브로드캐스팅", "크기 다른 배열 간 자동 연산\nstride=0 가상 확장, 메모리 효율", ACCENT_CYAN),
    ("DataFrame", "이질적 표 구조, 열별 독립 배열\nloc/iloc 라벨/정수 인덱싱", ACCENT_GREEN),
    ("GroupBy", "분할-적용-결합 패턴\nagg / transform / apply", ACCENT_ORANGE),
    ("Tidy Data", "변수=열, 관측=행 (Wickham 2014)\nmelt/pivot으로 변환", ACCENT_RED),
    ("벡터화 성능", "C 루프 + SIMD + 캐시 = 1000배\n모든 ML 코드의 기본 원칙", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(summaries):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.2) + Inches(2.5) * row
    add_shape(s, x, y, Inches(3.9), Inches(2.1), CARD_BG, color, radius=True)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
             title, font_size=17, color=color, bold=True)
    add_accent_line(s, x + Inches(0.2), y + Inches(0.6), Inches(2.0), color)
    add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.2),
             desc, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 23: 복습 질문 10개
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_header(s, "", "복습 질문 (Review Questions)")

questions = [
    "Q1. NumPy ndarray의 stride 개념을 설명하고, 전치(transpose) 시 stride 변화를 설명하시오.",
    "Q2. 브로드캐스팅의 3가지 규칙을 기술하고, (3,1) + (1,4) 결과 shape을 유도하시오.",
    "Q3. C-order와 F-order의 차이를 메모리 레이아웃 관점에서 설명하시오.",
    "Q4. Pandas loc과 iloc의 차이를 설명하고, 슬라이싱 끝 인덱스 포함 여부를 비교하시오.",
    "Q5. McKinney(2010)가 Pandas를 개발한 동기와 NumPy ndarray의 한계를 기술하시오.",
    "Q6. 결측치 처리 5가지 전략의 장단점을 비교하고, 시계열에 적합한 전략을 설명하시오.",
    "Q7. Wickham(2014)의 Tidy Data 3원칙과 messy→tidy 변환 Pandas 코드를 작성하시오.",
    "Q8. 벡터화 연산이 Python 루프보다 빠른 이유를 하드웨어(SIMD, 캐시)에서 설명하시오.",
    "Q9. Harris et al.(2020)의 파이썬 과학 생태계 레이어 구조를 설명하시오.",
    "Q10. Data Leakage란 무엇이며, StandardScaler 잘못 사용 시 어떻게 발생하는지 설명하시오.",
]
for i, q in enumerate(questions):
    y = Inches(2.0) + Inches(0.52) * i
    color = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED,
             ACCENT_PURPLE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE][i]
    add_shape(s, Inches(0.8), y + Inches(0.05), Inches(0.08), Inches(0.3), color)
    add_text(s, Inches(1.1), y, Inches(11.5), Inches(0.5),
             q, font_size=13, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 24: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, SECTION_BG)
add_shape(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
add_text(s, Inches(0), Inches(2.5), prs.slide_width, Inches(1.0),
         "Thank You", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), prs.slide_width, Inches(0.6),
         "3장: 판다스(Pandas)와 넘파이(NumPy)", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_accent_line(s, Inches(5.5), Inches(4.2), Inches(2.3), ACCENT_BLUE)
add_text(s, Inches(0), Inches(4.5), prs.slide_width, Inches(0.5),
         "다음 장: 4장 - 선형회귀", font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ── 저장 ──
output_path = os.path.join(os.path.dirname(__file__), "3장_판다스와_넘파이_강의PPT.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
